"""Compile accepted E03H references into one editable hybrid reference pack."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

from src.presentation_agent.magic_layer.e02h_candidate_compiler import _draw_candidate, _write_backplate_assets
from src.presentation_agent.magic_layer.e02h_hybrid_object_graph_builder import SLIDE_H_PX, SLIDE_W_PX
from src.presentation_agent.magic_layer.e03h_candidate_compiler import render_e03h_candidate_preview


SLIDE_W_IN = 16.0
SLIDE_H_IN = 9.0


def compile_e03h_reference_pack(payloads: list[dict[str, Any]], output_dir: str | Path) -> dict[str, Any]:
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    preview_paths: list[Path] = []
    for payload in payloads:
        asset_dir = output / "pack_backplate_assets" / payload["reference_id"]
        _write_backplate_assets(payload, asset_dir)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _draw_candidate(slide, payload, asset_dir)
        preview_dir = output / "pack_previews" / payload["reference_id"]
        preview_manifest = render_e03h_candidate_preview(payload, preview_dir)
        preview_paths.append(Path(preview_manifest["rendered_candidate"]))
    pptx_path = output / "editable_hybrid_reference_pack.pptx"
    prs.save(pptx_path)
    inventory = _audit_pack(pptx_path)
    contact_sheet = output / "editable_hybrid_reference_pack_contact_sheet.png"
    _contact_sheet(preview_paths, contact_sheet, "E03H Editable Hybrid Reference Pack")
    return {
        "schema_name": "editable_hybrid_reference_pack_render_manifest",
        "status": "passed" if inventory["semantic_raster_violation_count"] == 0 and inventory["full_slide_raster_count"] == 0 else "failed",
        "pptx_path": pptx_path.as_posix(),
        "pptx_exists": pptx_path.exists(),
        "contact_sheet": contact_sheet.as_posix(),
        "rendered_contact_sheet": contact_sheet.as_posix(),
        "slide_count": len(payloads),
        "semantic_raster_violation_count": inventory["semantic_raster_violation_count"],
        "full_slide_raster_count": inventory["full_slide_raster_count"],
        "screenshot_slide_count": 0,
        "text_count": inventory["text_count"],
        "media_count": inventory["media_count"],
        "native_chart_count": inventory["native_chart_count"],
        "native_table_count": inventory["native_table_count"],
        "canva_parity_claimed": False,
    }


def _audit_pack(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    objects = []
    full_slide_raster = 0
    semantic_raster = 0
    for slide_index, slide in enumerate(prs.slides):
        for z, shape in enumerate(slide.shapes):
            shape_type = str(shape.shape_type)
            is_picture = "PICTURE" in shape_type or int(getattr(shape.shape_type, "value", -999)) == 13
            name = shape.name or f"slide_{slide_index + 1}_shape_{z}"
            bbox = {
                "x": round(float(shape.left / prs.slide_width), 6),
                "y": round(float(shape.top / prs.slide_height), 6),
                "w": round(float(shape.width / prs.slide_width), 6),
                "h": round(float(shape.height / prs.slide_height), 6),
            }
            if is_picture and bbox["w"] >= 0.92 and bbox["h"] >= 0.92:
                full_slide_raster += 1
            if is_picture and any(token in name.lower() for token in ("title", "text", "icon", "footer", "card", "chart", "table", "matrix")):
                semantic_raster += 1
            objects.append(
                {
                    "slide_index": slide_index,
                    "shape_name": name,
                    "shape_type": shape_type,
                    "z_order": z,
                    "bbox_norm": bbox,
                    "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
                    "has_chart": bool(getattr(shape, "has_chart", False)),
                    "has_table": bool(getattr(shape, "has_table", False)),
                    "is_picture": is_picture,
                    "text": shape.text if getattr(shape, "has_text_frame", False) else "",
                }
            )
    return {
        "schema_name": "pptx_pack_inventory",
        "status": "passed" if full_slide_raster == 0 and semantic_raster == 0 else "failed",
        "pptx_path": pptx_path.as_posix(),
        "slide_count": len(prs.slides),
        "object_count": len(objects),
        "shape_count": sum(1 for row in objects if not row["is_picture"]),
        "text_count": sum(1 for row in objects if row["has_text_frame"] and row["text"].strip()),
        "media_count": sum(1 for row in objects if row["is_picture"]),
        "native_chart_count": sum(1 for row in objects if row["has_chart"]),
        "native_table_count": sum(1 for row in objects if row["has_table"]),
        "semantic_raster_violation_count": semantic_raster,
        "full_slide_raster_count": full_slide_raster,
        "screenshot_slide_count": 0,
        "objects": objects,
    }


def _contact_sheet(paths: list[Path], output: Path, title: str) -> None:
    thumbs = []
    for path in paths:
        with Image.open(path).convert("RGB") as image:
            thumbs.append(image.resize((320, 180)))
    cols = 3
    rows = max(1, (len(thumbs) + cols - 1) // cols)
    sheet = Image.new("RGB", (cols * 340 + 20, rows * 225 + 60), "#041826")
    draw = ImageDraw.Draw(sheet)
    draw.text((18, 16), title, fill="#F3F7FA", font=_font(14))
    for idx, thumb in enumerate(thumbs):
        x = 20 + (idx % cols) * 340
        y = 55 + (idx // cols) * 225
        sheet.paste(thumb, (x, y))
        draw.text((x, y + 184), paths[idx].parent.name, fill="#F3A51A", font=_font(7))
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)


def _font(size: int) -> ImageFont.ImageFont:
    try:
        return ImageFont.truetype("arial.ttf", size * 2)
    except OSError:
        return ImageFont.load_default()
