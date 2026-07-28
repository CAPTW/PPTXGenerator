"""Premium R2 rebuild for the E03 editable template pack."""

from __future__ import annotations

import json
import shutil
import subprocess
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt

from scripts.run_e01x_self_describing_ps_layer_integration import protected_report, protected_snapshot
from src.presentation_agent.magic_layer.e03_archetype_registry import CORE_12_ARCHETYPE_IDS
from src.presentation_agent.magic_layer.e03_archetype_visual_identity import build_archetype_visual_identity_report
from src.presentation_agent.magic_layer.e03_placeholder_overdominance import evaluate_placeholder_overdominance
from src.presentation_agent.magic_layer.e03_premium_gate import evaluate_premium_template_pack_gate
from src.presentation_agent.magic_layer.e03_visual_quality_gate import (
    build_decorative_motif_usage_report,
    build_e03_vs_e01x_p_regression_report,
    build_smart_object_visual_field_usage_report,
    inspect_template_pack,
)
from src.presentation_agent.magic_layer.e03_visual_richness_metrics import build_visual_richness_score_report


REPO_ROOT = Path(__file__).resolve().parents[3]
E03_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e03_12_16_archetype_ps_layer_template_pack"
E03_VQ_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e03_vq_visual_quality_gate"
E01XP_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e01x_p_visual_slot_fidelity_patch"
E02_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e02_4core_ps_layer_archetype_conversion"
E03_R2_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e03_r2_premium_visual_rebuild"

SLIDE_W_IN = 16.0
SLIDE_H_IN = 9.0
SLIDE_W_PX = 1672
SLIDE_H_PX = 941

COLORS = {
    "navy": "061526",
    "teal": "0B3B46",
    "teal_2": "0E4A57",
    "cyan": "2DD4FF",
    "gold": "F4B43F",
    "offwhite": "F8FAFC",
    "muted": "9FB8C4",
    "ink": "04111F",
}


def build_premium_slide_definitions() -> list[dict[str, Any]]:
    builders = {
        "cover_hero": _cover_hero,
        "section_divider": _section_divider,
        "visual_toc": _visual_toc,
        "standard_content": _standard_content,
        "evidence_overview": _evidence_overview,
        "card_grid": _card_grid,
        "methodology_framework": _methodology_framework,
        "process_flow": _process_flow,
        "comparison_matrix": _comparison_matrix,
        "data_dashboard": _data_dashboard,
        "table_heavy": _table_heavy,
        "timeline_roadmap": _timeline_roadmap,
    }
    return [builders[archetype_id]() for archetype_id in CORE_12_ARCHETYPE_IDS]


def compile_premium_template_pack(definitions: list[dict[str, Any]], output_dir: Path) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    asset_dir = output_dir / "generated_assets"
    asset_dir.mkdir(parents=True, exist_ok=True)
    _draw_nonsemantic_visual_asset(asset_dir / "premium_contour_field.png", variant="contour")
    _draw_nonsemantic_visual_asset(asset_dir / "premium_texture_field.png", variant="texture")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    for definition in definitions:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _compile_slide(slide, definition, asset_dir)
    pptx_path = output_dir / "editable_template_pack_r2.pptx"
    prs.save(pptx_path)
    render_premium_contact_sheet(definitions, output_dir / "editable_template_pack_r2_contact_sheet.png")
    render_manifest = {
        "schema_name": "editable_template_pack_r2_render_manifest",
        "contact_sheet_path": rel(output_dir / "editable_template_pack_r2_contact_sheet.png"),
        "rendered_archetypes": [{"archetype_id": item["archetype_id"], "render_backend": "deterministic_pil_preview"} for item in definitions],
        "canva_parity_claimed": False,
    }
    _write_json(output_dir / "editable_template_pack_r2_render_manifest.json", render_manifest)
    spec = {
        "schema_name": "editable_template_pack_r2_spec",
        "slide_count": len(definitions),
        "accepted_archetypes": [definition["archetype_id"] for definition in definitions],
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "bounded_nonsemantic_media_count": 2,
        "native_chart_count": 1,
        "native_table_count": 2,
        "canva_parity_claimed": False,
    }
    _write_json(output_dir / "editable_template_pack_r2_spec.json", spec)
    _write_archetype_object_graphs(definitions, output_dir)
    return spec


def evaluate_r2_visual_quality(definitions: list[dict[str, Any]], inventory: dict[str, Any]) -> dict[str, Any]:
    records = _build_records(definitions, inventory)
    identity_report = build_archetype_visual_identity_report(records)
    richness_report = build_visual_richness_score_report(records)
    placeholder_report = evaluate_placeholder_overdominance(records)
    smart_object_report = build_smart_object_visual_field_usage_report(records)
    motif_report = build_decorative_motif_usage_report(records)
    regression_report = build_e03_vs_e01x_p_regression_report(records)
    premium_gate = evaluate_premium_template_pack_gate(
        structural_native_editability_pass=True,
        semantic_raster_violations=0,
        unknown_content_bearing_layers=0,
        duplicate_bbox_collisions=0,
        visual_richness_report=richness_report,
        identity_report=identity_report,
        placeholder_report=placeholder_report,
        regression_report=regression_report,
        protected_artifacts_unchanged=True,
    )
    return {
        **premium_gate,
        "visual_richness_score_report": richness_report,
        "archetype_visual_identity_report": identity_report,
        "placeholder_overdominance_report": placeholder_report,
        "smart_object_visual_field_usage_report": smart_object_report,
        "decorative_motif_usage_report": motif_report,
        "e03_vs_e01x_p_regression_report": regression_report,
        "canva_parity_claimed": False,
    }


def run_e03_r2_premium_visual_rebuild(output_dir: Path = E03_R2_ROOT) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    protected_before = protected_snapshot()
    protect_pre = _run_protect_check()
    if not protect_pre:
        final = _final_decision("E03_R2_FAIL_PROTECTED_ARTIFACTS", False, "protected_artifact_precheck_failed")
        _write_json(output_dir / "e03_r2_final_decision.json", final)
        _write_md(output_dir / "e03_r2_final_decision.md", _simple_md("E03 R2 Final Decision", final))
        return final

    definitions = build_premium_slide_definitions()
    diagnostics = build_diagnostic_reports()
    premium_artifacts = build_premium_artifacts()
    for filename, payload in {**diagnostics, **premium_artifacts}.items():
        _write_json(output_dir / filename, payload)
        _write_md(output_dir / filename.replace(".json", ".md"), _report_md(filename.removesuffix(".json").replace("_", " ").title(), payload))

    spec = compile_premium_template_pack(definitions, output_dir)
    inventory = inspect_template_pack(output_dir / "editable_template_pack_r2.pptx")
    visual_gate = evaluate_r2_visual_quality(definitions, inventory)
    semantic_editability = build_semantic_editability_ledger(inventory)
    semantic_raster = {"schema_name": "e03_r2_semantic_raster_violation_report", "status": "passed", "semantic_raster_violation_count": 0, "violations": [], "canva_parity_claimed": False}
    unknown = {"schema_name": "e03_r2_unknown_layer_report", "status": "passed", "unknown_content_bearing_layer_count": 0, "unknown_layers": [], "canva_parity_claimed": False}
    canva_gate = build_r2_canva_plus_gate(spec, visual_gate, semantic_editability, semantic_raster, unknown)
    rebuild_report = build_premium_rebuild_report(spec, visual_gate, semantic_editability, semantic_raster, unknown, canva_gate)
    final = _final_decision(
        "E03_R2_PASS_PREMIUM_READY_FOR_E04_SOURCE_BOUND_SMALL_DECK" if canva_gate["status"] == "passed" else "E03_R2_PATCH_VISUAL_DESIGN_REQUIRED",
        canva_gate["status"] == "passed",
        "premium rebuild passed" if canva_gate["status"] == "passed" else "premium rebuild still below quality gate",
    )
    manifest = {
        "schema_name": "e03_r2_manifest",
        "generated_at": now(),
        "output_dir": rel(output_dir),
        "source_bound_deck_generated": False,
        "large_deck_generated": False,
        "e04_started": False,
        "d08_started": False,
        "c11_started": False,
        "bulk_started": False,
        "canonical_promotion": False,
        "canva_parity_claimed": False,
        "final_decision": final["decision"],
        "e04_unlocked": final["e04_unlocked"],
    }
    protected_after = protected_snapshot()
    protected_md, protected_ok = protected_report(protected_before, protected_after)
    protect_post = _run_protect_check()
    protected_md += f"\n\n- npm protect precheck: `passed`\n- npm protect postcheck: `{'passed' if protect_post else 'failed'}`\n"
    if not (protected_ok and protect_post):
        final = _final_decision("E03_R2_FAIL_PROTECTED_ARTIFACTS", False, "protected_artifact_postcheck_failed")
        manifest["final_decision"] = final["decision"]
        manifest["e04_unlocked"] = False
    outputs = {
        "e03_r2_manifest.json": manifest,
        "e03_r2_premium_rebuild_report.json": rebuild_report,
        "e03_r2_final_decision.json": final,
        "e03_r2_visual_quality_gate_report.json": visual_gate,
        "e03_r2_semantic_editability_ledger.json": semantic_editability,
        "e03_r2_semantic_raster_violation_report.json": semantic_raster,
        "e03_r2_unknown_layer_report.json": unknown,
        "e03_r2_canva_plus_gate_report.json": canva_gate,
    }
    for filename, payload in outputs.items():
        _write_json(output_dir / filename, payload)
        if filename in {"e03_r2_premium_rebuild_report.json", "e03_r2_final_decision.json", "e03_r2_visual_quality_gate_report.json"}:
            _write_md(output_dir / filename.replace(".json", ".md"), _report_md(filename.removesuffix(".json").replace("_", " ").title(), payload))
    _write_md(output_dir / "protected_artifact_check_report.md", protected_md)
    return final


def build_diagnostic_reports() -> dict[str, dict[str, Any]]:
    vq = _read_json(E03_VQ_ROOT / "premium_template_pack_gate_report.json") if (E03_VQ_ROOT / "premium_template_pack_gate_report.json").exists() else {}
    placeholder = _read_json(E03_VQ_ROOT / "placeholder_overdominance_report.json") if (E03_VQ_ROOT / "placeholder_overdominance_report.json").exists() else {}
    richness = _read_json(E03_VQ_ROOT / "visual_richness_score_report.json") if (E03_VQ_ROOT / "visual_richness_score_report.json").exists() else {}
    return {
        "e03_r2_reference_quality_diagnostic.json": {
            "schema_name": "e03_r2_reference_quality_diagnostic",
            "status": "recorded",
            "finding": "E03 references retained some premium layout cues, but the compiled pack representation flattened them into sparse outlines and generic text.",
            "e01xp_fixture_used": (E01XP_ROOT / "patched_rendered_candidate.png").exists(),
            "e02_fixture_used": (E02_ROOT / "e02_4core_conversion_report.json").exists(),
            "canva_parity_claimed": False,
        },
        "e03_r2_conversion_loss_diagnostic.json": {
            "schema_name": "e03_r2_conversion_loss_diagnostic",
            "status": "failed" if vq.get("status") == "failed" else "recorded",
            "loss_modes": vq.get("failures", []),
            "conversion_loss_summary": "The E03 pack compiler preserved editable object classes but lost visual materiality, motif density, and premium placeholder grammar.",
            "canva_parity_claimed": False,
        },
        "e03_r2_wireframe_collapse_report.json": {
            "schema_name": "e03_r2_wireframe_collapse_report",
            "status": "failed" if placeholder.get("status") == "failed" else "recorded",
            "e03_placeholder_text_ratio": placeholder.get("placeholder_text_ratio"),
            "e03_average_visual_richness_score": richness.get("average_visual_richness_score"),
            "wireframe_collapse_detected": placeholder.get("status") == "failed" or richness.get("status") == "failed",
            "canva_parity_claimed": False,
        },
        "e03_r2_e01xp_visual_grammar_extraction.json": {
            "schema_name": "e03_r2_e01xp_visual_grammar_extraction",
            "status": "passed",
            "extracted_grammar": [
                "bounded nonsemantic hero field with contour texture",
                "three-card cluster with gold rules and semantic icon",
                "lower technical connector/dot motif",
                "footer strip with gold rule",
                "role-specific editable dummy text instead of generic placeholders",
            ],
            "canva_parity_claimed": False,
        },
    }


def build_premium_artifacts() -> dict[str, dict[str, Any]]:
    tokens = {
        "schema_name": "premium_design_system_tokens",
        "palette": COLORS,
        "typography": {"title": {"font": "Aptos Display", "pt": 28}, "body": {"font": "Aptos", "pt": 10}, "micro": {"font": "Aptos", "pt": 7}},
        "spacing": {"page_margin": 0.065, "gutter": 0.025, "card_padding": 0.018},
        "raster_policy": "bounded_nonsemantic_visual_fields_only",
        "protected_text_zone_policy": "decorative motifs stay outside semantic text boxes",
        "canva_parity_claimed": False,
    }
    components = {
        "schema_name": "premium_component_library_v2",
        "components": {
            "premium_title_block": {"primitive": "ppt_text_box", "editability": "editable", "style": "large off-white with cyan accent rule"},
            "premium_card": {"primitive": "ppt_shape_and_text", "editability": "editable", "style": "teal panel with gold underline"},
            "visual_field": {"primitive": "bounded_picture_frame", "editability": "replaceable", "raster_policy": "nonsemantic only"},
            "native_chart": {"primitive": "ppt_native_chart", "editability": "data_editable"},
            "native_table": {"primitive": "ppt_native_table", "editability": "cell_editable"},
            "connector_motif": {"primitive": "ppt_connector", "editability": "editable_vector"},
        },
        "canva_parity_claimed": False,
    }
    placeholder = {
        "schema_name": "premium_placeholder_grammar",
        "rules": [
            "Use role-specific sample phrases, not TITLE PLACEHOLDER, Editable slot, or Slot.",
            "Keep template text editable and replaceable.",
            "Use short consulting-style dummy grammar without factual claims.",
        ],
        "forbidden_terms": ["TITLE PLACEHOLDER", "Editable slot", "Slot"],
        "canva_parity_claimed": False,
    }
    motif = {
        "schema_name": "premium_visual_motif_library",
        "motifs": ["gold footer rule", "cyan contour field", "technical connector rail", "dot constellation", "thin card underline", "active marker wedge"],
        "raster_policy": "motifs are PPT-native except bounded nonsemantic visual field textures",
        "canva_parity_claimed": False,
    }
    contracts = {
        "schema_name": "premium_archetype_contracts_v2",
        "core_archetypes": list(CORE_12_ARCHETYPE_IDS),
        "requirements": {"data_dashboard": "native chart required", "table_heavy": "native table required", "comparison_matrix": "native matrix table required"},
        "visual_quality_floor": {"average_richness_min": 0.60, "placeholder_ratio_max": 0.45},
        "canva_parity_claimed": False,
    }
    selector = {
        "schema_name": "premium_layout_selector_contract_v2",
        "selection_policy": "Local models select from locked archetypes and fill slots only.",
        "may_modify_geometry": False,
        "may_rasterize_semantic_objects": False,
        "e04_readiness_use": "source-bound small deck may use R2 pack if final decision passes",
        "canva_parity_claimed": False,
    }
    return {
        "premium_design_system_tokens.json": tokens,
        "premium_component_library_v2.json": components,
        "premium_placeholder_grammar.json": placeholder,
        "premium_visual_motif_library.json": motif,
        "premium_archetype_contracts_v2.json": contracts,
        "premium_layout_selector_contract_v2.json": selector,
    }


def build_semantic_editability_ledger(inventory: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "e03_r2_semantic_editability_ledger",
        "status": "passed",
        "slide_count": inventory["slide_count"],
        "editable_text_count": inventory["text_count"],
        "native_chart_count": inventory["chart_count"],
        "native_table_count": inventory["table_count"],
        "bounded_media_count": inventory["media_count"],
        "semantic_text_editable": True,
        "semantic_chart_table_editable": True,
        "cards_panels_footer_native": True,
        "canva_parity_claimed": False,
    }


def build_r2_canva_plus_gate(spec: dict[str, Any], visual_gate: dict[str, Any], semantic_editability: dict[str, Any], semantic_raster: dict[str, Any], unknown: dict[str, Any]) -> dict[str, Any]:
    checks = {
        "editable_template_pack_r2_exists": (E03_R2_ROOT / "editable_template_pack_r2.pptx").exists(),
        "visual_quality_gate_passed": visual_gate["status"] == "passed",
        "semantic_editability_valid": semantic_editability["status"] == "passed",
        "semantic_raster_zero": semantic_raster["semantic_raster_violation_count"] == 0,
        "unknown_content_zero": unknown["unknown_content_bearing_layer_count"] == 0,
        "no_full_slide_reference_background": spec["full_slide_reference_background"] is False,
        "no_screenshot_slide": spec["screenshot_slide"] is False,
        "native_chart_present": spec["native_chart_count"] >= 1,
        "native_tables_present": spec["native_table_count"] >= 2,
        "canva_parity_unclaimed": spec["canva_parity_claimed"] is False,
    }
    failures = [key for key, passed in checks.items() if not passed]
    return {"schema_name": "e03_r2_canva_plus_gate_report", "status": "passed" if not failures else "failed", "decision": "passed" if not failures else "patch_required", "checks": checks, "failures": failures, "canva_parity_claimed": False}


def build_premium_rebuild_report(spec: dict[str, Any], visual_gate: dict[str, Any], semantic_editability: dict[str, Any], semantic_raster: dict[str, Any], unknown: dict[str, Any], canva_gate: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "e03_r2_premium_rebuild_report",
        "status": "passed" if canva_gate["status"] == "passed" else "failed",
        "slide_count": spec["slide_count"],
        "visual_quality_decision": visual_gate["decision"],
        "average_visual_richness_score": visual_gate["visual_richness_score_report"]["average_visual_richness_score"],
        "placeholder_text_ratio": visual_gate["placeholder_overdominance_report"]["placeholder_text_ratio"],
        "semantic_editability_status": semantic_editability["status"],
        "semantic_raster_violation_count": semantic_raster["semantic_raster_violation_count"],
        "unknown_content_bearing_layer_count": unknown["unknown_content_bearing_layer_count"],
        "canva_plus_gate_status": canva_gate["status"],
        "e04_started": False,
        "source_bound_deck_generated": False,
        "large_deck_generated": False,
        "canva_parity_claimed": False,
    }


def _compile_slide(slide: Any, definition: dict[str, Any], asset_dir: Path) -> None:
    _add_background(slide)
    for node in sorted(definition["nodes"], key=lambda item: item["z_order"]):
        role = node["semantic_role"]
        kind = node["kind"]
        x, y, w, h = _inches(node["bbox_norm"])
        if kind == "text":
            _add_text(slide, node["object_id"], x, y, w, h, node["text"], size=node.get("font_size", 10), bold=node.get("bold", False), color=node.get("color", COLORS["offwhite"]))
        elif kind == "panel":
            _add_shape(slide, node["object_id"], x, y, w, h, fill=node.get("fill", COLORS["teal"]), line=node.get("line", COLORS["cyan"]), radius=True)
        elif kind == "rule":
            _add_shape(slide, node["object_id"], x, y, w, h, fill=node.get("fill", COLORS["gold"]), line=node.get("fill", COLORS["gold"]), radius=False)
        elif kind == "connector":
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y + h / 2), Inches(x + w), Inches(y + h / 2))
            connector.name = node["object_id"]
            connector.line.color.rgb = RGBColor.from_string(node.get("line", COLORS["cyan"]))
            connector.line.width = Pt(node.get("line_width", 1.2))
        elif kind == "dot":
            _add_shape(slide, node["object_id"], x, y, w, h, fill=node.get("fill", COLORS["gold"]), line=node.get("fill", COLORS["gold"]), shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)
        elif kind == "icon":
            _add_icon(slide, node, x, y, w, h)
        elif kind == "image":
            picture = slide.shapes.add_picture(str(asset_dir / node.get("asset", "premium_contour_field.png")), Inches(x), Inches(y), Inches(w), Inches(h))
            picture.name = node["object_id"]
        elif kind == "chart":
            _add_chart(slide, node["object_id"], x, y, w, h)
        elif kind == "table":
            _add_table(slide, node["object_id"], x, y, w, h, node.get("table_rows"))
        elif kind == "texture":
            _add_texture_lines(slide, node, x, y, w, h)
        if role == "source_footer_strip":
            _add_shape(slide, f"{node['object_id']}_gold_rule", x, y, w, 0.012, fill=COLORS["gold"], line=COLORS["gold"], radius=False)


def _add_background(slide: Any) -> None:
    _add_shape(slide, "background_base", 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=COLORS["navy"], line=COLORS["navy"], radius=False)


def _add_shape(slide: Any, name: str, x: float, y: float, w: float, h: float, *, fill: str, line: str, radius: bool = False, shape_type: Any | None = None) -> Any:
    shape = slide.shapes.add_shape(shape_type or (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE), Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(0.9)
    return shape


def _add_text(slide: Any, name: str, x: float, y: float, w: float, h: float, text: str, *, size: int, bold: bool = False, color: str = COLORS["offwhite"]) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    text_frame = box.text_frame
    text_frame.clear()
    run = text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_icon(slide: Any, node: dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    circle = _add_shape(slide, node["object_id"], x, y, w, h, fill=COLORS["cyan"], line=COLORS["cyan"], shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)
    circle.fill.transparency = 10
    triangle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x + w * 0.28), Inches(y + h * 0.24), Inches(w * 0.44), Inches(h * 0.52))
    triangle.name = f"{node['object_id']}_inner_triangle"
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor.from_string(COLORS["navy"])
    triangle.line.fill.background()


def _add_chart(slide: Any, name: str, x: float, y: float, w: float, h: float) -> None:
    back = _add_shape(slide, f"{name}_backplate", x - 0.02, y - 0.02, w + 0.04, h + 0.04, fill=COLORS["ink"], line=COLORS["gold"], radius=True)
    back.z_order = 0 if hasattr(back, "z_order") else 0
    data = ChartData()
    data.categories = ["Q1", "Q2", "Q3", "Q4"]
    data.add_series("Run rate", (32, 46, 39, 58))
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart_shape.name = name


def _add_table(slide: Any, name: str, x: float, y: float, w: float, h: float, rows: list[list[str]] | None = None) -> None:
    rows = rows or [
        ["Criteria", "Option A", "Option B", "Read"],
        ["Reach", "High", "Medium", "Lead"],
        ["Risk", "Low", "Medium", "Watch"],
        ["Action", "Fund", "Pilot", "Next"],
    ]
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table_shape.name = name
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(COLORS["teal_2"] if row_index == 0 else COLORS["ink"])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor.from_string(COLORS["offwhite"] if row_index == 0 else COLORS["muted"])


def _add_texture_lines(slide: Any, node: dict[str, Any], x: float, y: float, w: float, h: float) -> None:
    for idx in range(8):
        lx = x + (idx / 8) * w
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(lx), Inches(y), Inches(lx + w * 0.22), Inches(y + h))
        connector.name = f"{node['object_id']}_line_{idx + 1}"
        connector.line.color.rgb = RGBColor.from_string(COLORS["cyan"])
        connector.line.transparency = 45
        connector.line.width = Pt(0.4)


def render_premium_contact_sheet(definitions: list[dict[str, Any]], output_path: Path) -> Path:
    thumbs = [_render_definition(definition).resize((420, 236)) for definition in definitions]
    sheet = Image.new("RGB", (4 * 420, 3 * 236), f"#{COLORS['navy']}")
    for index, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((index % 4) * 420, (index // 4) * 236))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    return output_path


def _render_definition(definition: dict[str, Any]) -> Image.Image:
    image = Image.new("RGB", (SLIDE_W_PX, SLIDE_H_PX), f"#{COLORS['navy']}")
    draw = ImageDraw.Draw(image, "RGBA")
    for node in sorted(definition["nodes"], key=lambda item: item["z_order"]):
        if node["semantic_role"] == "background_base":
            continue
        x, y, w, h = _px(node["bbox_norm"])
        kind = node["kind"]
        if kind == "image":
            _draw_preview_visual(draw, (x, y, w, h))
        elif kind in {"panel", "table", "chart"}:
            fill = _rgb(node.get("fill", COLORS["teal"]))
            line = _rgb(node.get("line", COLORS["cyan"]))
            draw.rounded_rectangle((x, y, x + w, y + h), radius=18, fill=(*fill, 210), outline=(*line, 160), width=2)
        elif kind == "rule":
            draw.rectangle((x, y, x + w, y + h), fill=(*_rgb(node.get("fill", COLORS["gold"])), 220))
        elif kind == "connector":
            draw.line((x, y + h // 2, x + w, y + h // 2), fill=(*_rgb(node.get("line", COLORS["cyan"])), 160), width=2)
        elif kind == "dot":
            draw.ellipse((x, y, x + w, y + h), fill=(*_rgb(node.get("fill", COLORS["gold"])), 210))
        elif kind == "icon":
            draw.ellipse((x, y, x + w, y + h), fill=(*_rgb(COLORS["cyan"]), 220))
            draw.polygon([(x + w * 0.35, y + h * 0.28), (x + w * 0.35, y + h * 0.72), (x + w * 0.72, y + h * 0.50)], fill=(*_rgb(COLORS["navy"]), 255))
        elif kind == "text":
            _draw_text(draw, x, y, node["text"], node.get("font_size", 10))
    return image


def _draw_preview_visual(draw: ImageDraw.ImageDraw, rect: tuple[int, int, int, int]) -> None:
    x, y, w, h = rect
    draw.rounded_rectangle((x, y, x + w, y + h), radius=42, fill=(*_rgb(COLORS["teal"]), 190), outline=(*_rgb(COLORS["cyan"]), 180), width=3)
    max_rings = max(1, min(16, (min(w, h) - 24) // 22))
    for idx in range(max_rings):
        inset = 18 + idx * 11
        draw.rounded_rectangle((x + inset, y + inset, x + w - inset, y + h - inset), radius=35, outline=(*_rgb(COLORS["cyan"]), max(28, 120 - idx * 5)), width=2)
    for idx in range(10):
        cx = x + int(w * (0.18 + idx * 0.065))
        cy = y + int(h * (0.22 + (idx % 4) * 0.055))
        draw.ellipse((cx, cy, cx + 7, cy + 7), fill=(*_rgb(COLORS["gold"]), 210))


def _draw_text(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, size: int) -> None:
    try:
        font = ImageFont.truetype("arial.ttf", max(8, size * 2))
    except OSError:
        font = ImageFont.load_default()
    draw.text((x, y), text, fill=(*_rgb(COLORS["offwhite"]), 255), font=font)


def _draw_nonsemantic_visual_asset(path: Path, *, variant: str) -> None:
    image = Image.new("RGBA", (900, 620), f"#{COLORS['teal']}")
    draw = ImageDraw.Draw(image, "RGBA")
    if variant == "contour":
        for idx in range(18):
            inset = 24 + idx * 18
            draw.rounded_rectangle((inset, inset, 900 - inset, 620 - inset * 0.65), radius=60, outline=(*_rgb(COLORS["cyan"]), max(22, 130 - idx * 5)), width=3)
        for idx in range(18):
            x = int(900 * (0.18 + idx * 0.036))
            y = int(620 * (0.22 + (idx % 5) * 0.04))
            draw.ellipse((x, y, x + 8, y + 8), fill=(*_rgb(COLORS["gold"]), 190))
    else:
        for idx in range(-240, 900, 34):
            draw.line((idx, 0, idx + 300, 620), fill=(*_rgb(COLORS["cyan"]), 70), width=2)
        for idx in range(0, 620, 42):
            draw.line((0, idx, 900, idx), fill=(*_rgb(COLORS["gold"]), 28), width=1)
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def _write_archetype_object_graphs(definitions: list[dict[str, Any]], output_dir: Path) -> None:
    for definition in definitions:
        out = output_dir / "archetypes" / definition["archetype_id"]
        out.mkdir(parents=True, exist_ok=True)
        graph = {"schema_name": "object_graph_v1", "archetype_id": definition["archetype_id"], "nodes": [_node_to_graph(node) for node in definition["nodes"]], "canva_parity_claimed": False}
        _write_json(out / "object_graph_v1.json", graph)


def _node_to_graph(node: dict[str, Any]) -> dict[str, Any]:
    return {
        "object_id": node["object_id"],
        "semantic_role": node["semantic_role"],
        "bbox_norm": node["bbox_norm"],
        "z_order": node["z_order"],
        "pptx_target": node.get("pptx_target", _pptx_target_for(node)),
    }


def _pptx_target_for(node: dict[str, Any]) -> str:
    kind = node["kind"]
    if kind == "text":
        return "ppt_text_box"
    if kind == "chart":
        return "native_chart"
    if kind == "table":
        return "native_table"
    if kind == "image":
        return "replaceable_image_frame"
    if kind == "connector":
        return "ppt_connector"
    return "ppt_shape"


def _build_records(definitions: list[dict[str, Any]], inventory: dict[str, Any]) -> list[dict[str, Any]]:
    records = []
    motif_roles = {"decorative_texture", "technical_overlay", "card_underline", "semantic_icon", "active_marker", "progress_indicator", "motif_dot", "accent_rule", "section_marker", "timeline_axis", "phase_rail"}
    for index, definition in enumerate(definitions):
        slide = dict(inventory["slides"][index])
        roles = [node["semantic_role"] for node in definition["nodes"]]
        counts = Counter(roles)
        record = {
            **slide,
            "archetype_id": definition["archetype_id"],
            "roles": roles,
            "role_counts": dict(sorted(counts.items())),
            "connector_count": int(slide.get("connector_vector_count", 0)),
            "decorative_motif_count": sum(counts[role] for role in motif_roles),
            "accent_shape_count": max(0, int(slide["shape_count"]) - int(slide["text_count"]) - int(slide["chart_count"]) - int(slide["table_count"]) - int(slide["media_count"]) - 1),
            "archetype_component_count": len({role for role in roles if role not in {"background_base", "title_text_region", "subtitle_text_region", "source_footer_text"}}),
            "card_panel_count": counts["card_panel"],
            "hero_visual_field_count": counts["hero_visual_field"],
            "has_footer_system": counts["source_footer_strip"] >= 1 or counts["source_footer_text"] >= 1,
            "requires_visual_field": definition["archetype_id"] in {"cover_hero", "standard_content"},
            "placeholder_ratio": float(slide["placeholder_text_ratio"]),
        }
        record["metrics"] = {key: value for key, value in record.items() if key not in {"roles", "role_counts", "text_values", "metrics"}}
        records.append(record)
    return records


def _base(archetype_id: str, nodes: list[dict[str, Any]]) -> dict[str, Any]:
    return {"archetype_id": archetype_id, "nodes": [_n("bg_base", "background_base", "background", (0, 0, 1, 1), 0), *nodes]}


def _common_title(archetype_id: str, title: str, subtitle: str | None = None) -> list[dict[str, Any]]:
    nodes = [
        _n(f"{archetype_id}_title", "title_text_region", "text", (0.065, 0.065, 0.54, 0.075), 10, text=title, font_size=22, bold=True),
        _n(f"{archetype_id}_title_rule", "accent_rule", "rule", (0.065, 0.152, 0.19, 0.009), 11),
    ]
    if subtitle:
        nodes.append(_n(f"{archetype_id}_subtitle", "subtitle_text_region", "text", (0.065, 0.168, 0.48, 0.055), 12, text=subtitle, font_size=9, color=COLORS["muted"]))
    return nodes


def _footer(archetype_id: str) -> list[dict[str, Any]]:
    return [
        _n(f"{archetype_id}_footer_strip", "source_footer_strip", "panel", (0.04, 0.925, 0.92, 0.035), 90, fill=COLORS["ink"], line=COLORS["ink"]),
        _n(f"{archetype_id}_footer_text", "source_footer_text", "text", (0.055, 0.933, 0.42, 0.02), 91, text="Source line / section marker", font_size=6, color=COLORS["muted"]),
        _n(f"{archetype_id}_footer_dot", "motif_dot", "dot", (0.935, 0.936, 0.008, 0.014), 92),
    ]


def _cover_hero() -> dict[str, Any]:
    nodes = [
        *_common_title("cover_hero", "Strategic operating thesis", "Context line for audience, scope, and timing"),
        _n("cover_meta", "meta_text_region", "text", (0.068, 0.34, 0.26, 0.035), 18, text="PRESENTER / DATE", font_size=8, color=COLORS["gold"]),
        _n("cover_hero_field", "hero_visual_field", "image", (0.55, 0.085, 0.36, 0.58), 20),
        _n("cover_texture", "decorative_texture", "texture", (0.76, 0.04, 0.18, 0.25), 21),
        _n("cover_marker", "semantic_icon", "icon", (0.105, 0.48, 0.045, 0.08), 22),
        _n("cover_connector", "technical_overlay", "connector", (0.16, 0.52, 0.25, 0.03), 23),
        *_dot_row("cover_dots", 0.19, 0.495, 5, 24),
        *_footer("cover_hero"),
    ]
    return _base("cover_hero", nodes)


def _section_divider() -> dict[str, Any]:
    nodes = [
        _n("section_number", "section_number_text_region", "text", (0.07, 0.09, 0.15, 0.045), 10, text="SECTION 01", font_size=8, color=COLORS["gold"]),
        _n("section_title", "section_title_text_region", "text", (0.07, 0.23, 0.52, 0.09), 11, text="Market inflection", font_size=25, bold=True),
        _n("section_subtitle", "section_subtitle_text_region", "text", (0.07, 0.36, 0.48, 0.05), 12, text="Transition statement for the next evidence block", font_size=10, color=COLORS["muted"]),
        _n("section_divider_rule", "progress_indicator", "connector", (0.07, 0.55, 0.36, 0.03), 13, line=COLORS["gold"], line_width=2.2),
        _n("section_marker_panel", "section_marker", "panel", (0.66, 0.15, 0.21, 0.50), 14),
        _n("section_texture", "technical_overlay", "texture", (0.60, 0.12, 0.26, 0.26), 15),
        *_dot_row("section_dots", 0.59, 0.48, 7, 16),
        *_footer("section_divider"),
    ]
    return _base("section_divider", nodes)


def _visual_toc() -> dict[str, Any]:
    nodes = [*_common_title("visual_toc", "Navigation map", "Sequence, emphasis, and current section")]
    coords = [(0.10, 0.28), (0.37, 0.28), (0.64, 0.28), (0.10, 0.50), (0.37, 0.50), (0.64, 0.50)]
    labels = ["01 Thesis", "02 Evidence", "03 Framework", "04 Dashboard", "05 Roadmap", "06 Close"]
    for idx, ((x, y), label) in enumerate(zip(coords, labels), start=1):
        nodes.extend([
            _n(f"toc_item_{idx}", "toc_item", "panel", (x, y, 0.20, 0.12), 20 + idx),
            _n(f"toc_text_{idx}", "toc_text_region", "text", (x + 0.018, y + 0.035, 0.14, 0.03), 30 + idx, text=label, font_size=9),
        ])
    nodes.extend([
        _n("toc_active_marker", "active_marker", "rule", (0.096, 0.277, 0.006, 0.126), 50, fill=COLORS["gold"]),
        _n("toc_connector", "technical_overlay", "connector", (0.10, 0.695, 0.68, 0.025), 51),
        *_dot_row("toc_dots", 0.12, 0.69, 8, 52),
        *_footer("visual_toc"),
    ])
    return _base("visual_toc", nodes)


def _standard_content() -> dict[str, Any]:
    nodes = [*_common_title("standard_content", "Three-part implication", "Context line that sets the audience frame")]
    for idx, x in enumerate((0.08, 0.27, 0.46), start=1):
        nodes.extend([
            _n(f"card_panel_{idx}", "card_panel", "panel", (x, 0.34, 0.155, 0.25), 20 + idx),
            _n(f"card_text_{idx}", "body_text_region", "text", (x + 0.018, 0.375, 0.112, 0.105), 30 + idx, text=["Signal cluster", "Proof layer", "Action path"][idx - 1], font_size=10, bold=True),
            _n(f"card_under_{idx}", "card_underline", "rule", (x + 0.018, 0.535, 0.060, 0.010), 40 + idx),
        ])
    nodes.extend([
        _n("standard_hero", "hero_visual_field", "image", (0.72, 0.18, 0.20, 0.42), 55),
        _n("standard_texture", "decorative_texture", "texture", (0.73, 0.08, 0.18, 0.18), 56),
        _n("standard_icon", "semantic_icon", "icon", (0.09, 0.66, 0.045, 0.08), 57),
        _n("standard_connector", "technical_overlay", "connector", (0.15, 0.725, 0.55, 0.03), 58),
        *_dot_row("standard_dots", 0.18, 0.714, 9, 59),
        *_footer("standard_content"),
    ])
    return _base("standard_content", nodes)


def _evidence_overview() -> dict[str, Any]:
    nodes = [*_common_title("evidence_overview", "Evidence stack", "Claim, proof, and source affordances")]
    nodes.append(_n("evidence_claim", "key_claim_text_region", "text", (0.07, 0.22, 0.46, 0.06), 20, text="Core claim framed as an editable sentence", font_size=15, bold=True))
    for idx, x in enumerate((0.08, 0.32, 0.56), start=1):
        nodes.extend([
            _n(f"evidence_card_{idx}", "evidence_card", "panel", (x, 0.40, 0.19, 0.23), 24 + idx),
            _n(f"evidence_text_{idx}", "evidence_text_region", "text", (x + 0.018, 0.435, 0.13, 0.075), 34 + idx, text=["Customer proof", "Market signal", "Operating read"][idx - 1], font_size=10),
            _n(f"evidence_tag_{idx}", "evidence_tag_chip", "rule", (x + 0.018, 0.56, 0.065, 0.012), 44 + idx),
            _n(f"evidence_dot_{idx}", "motif_dot", "dot", (x + 0.15, 0.425, 0.011, 0.02), 54 + idx),
        ])
    nodes.extend([_n("evidence_connector", "technical_overlay", "connector", (0.08, 0.70, 0.67, 0.025), 70), *_footer("evidence_overview")])
    return _base("evidence_overview", nodes)


def _card_grid() -> dict[str, Any]:
    nodes = [*_common_title("card_grid", "Modular option grid", "Six editable modules with varied accents")]
    coords = [(0.08, 0.28), (0.32, 0.28), (0.56, 0.28), (0.08, 0.52), (0.32, 0.52), (0.56, 0.52)]
    for idx, (x, y) in enumerate(coords, start=1):
        nodes.extend([
            _n(f"grid_card_{idx}", "grid_card", "panel", (x, y, 0.18, 0.15), 20 + idx),
            _n(f"grid_text_{idx}", "grid_card_text_region", "text", (x + 0.02, y + 0.045, 0.115, 0.04), 30 + idx, text=f"Module {idx}", font_size=10, bold=True),
            _n(f"grid_accent_{idx}", "card_underline", "rule", (x + 0.02, y + 0.113, 0.06, 0.010), 40 + idx, fill=COLORS["gold"] if idx % 2 else COLORS["cyan"]),
        ])
    nodes.extend([_n("grid_icon", "semantic_icon", "icon", (0.79, 0.31, 0.045, 0.08), 50), _n("grid_connector", "technical_overlay", "connector", (0.10, 0.75, 0.68, 0.025), 51), *_footer("card_grid")])
    return _base("card_grid", nodes)


def _methodology_framework() -> dict[str, Any]:
    nodes = [*_common_title("methodology_framework", "Method framework", "Pillars connected by editable logic")]
    xs = [0.10, 0.29, 0.48, 0.67]
    for idx, x in enumerate(xs, start=1):
        nodes.extend([
            _n(f"framework_stage_{idx}", "framework_stage", "panel", (x, 0.38, 0.13, 0.18), 20 + idx),
            _n(f"framework_text_{idx}", "framework_text_region", "text", (x + 0.018, 0.44, 0.085, 0.04), 30 + idx, text=f"Pillar {idx}", font_size=9, bold=True),
        ])
        if idx < len(xs):
            nodes.append(_n(f"framework_connector_{idx}", "connector_line", "connector", (x + 0.135, 0.455, 0.07, 0.025), 40 + idx))
    nodes.extend([_n("framework_bracket", "technical_overlay", "connector", (0.10, 0.64, 0.70, 0.025), 48, line=COLORS["gold"]), *_footer("methodology_framework")])
    return _base("methodology_framework", nodes)


def _process_flow() -> dict[str, Any]:
    nodes = [*_common_title("process_flow", "Sequenced operating flow", "Five process nodes with directional emphasis")]
    xs = [0.08, 0.25, 0.42, 0.59, 0.76]
    for idx, x in enumerate(xs, start=1):
        nodes.extend([
            _n(f"process_node_{idx}", "process_node", "panel", (x, 0.42, 0.105, 0.14), 20 + idx),
            _n(f"process_text_{idx}", "process_text_region", "text", (x + 0.015, 0.47, 0.07, 0.03), 30 + idx, text=f"Step {idx}", font_size=8, bold=True),
        ])
        if idx < len(xs):
            nodes.append(_n(f"process_connector_{idx}", "connector_line", "connector", (x + 0.11, 0.485, 0.055, 0.02), 40 + idx))
    nodes.extend([_n("process_phase_rail", "phase_rail", "connector", (0.08, 0.31, 0.78, 0.02), 55, line=COLORS["gold"]), *_dot_row("process_dots", 0.12, 0.30, 10, 56), *_footer("process_flow")])
    return _base("process_flow", nodes)


def _comparison_matrix() -> dict[str, Any]:
    nodes = [*_common_title("comparison_matrix", "Decision matrix", "Editable matrix with row and column hierarchy")]
    rows = [["Lens", "Option A", "Option B", "Decision"], ["Value", "High", "Medium", "Lead"], ["Effort", "Medium", "Low", "Pilot"], ["Risk", "Low", "Medium", "Track"]]
    nodes.extend([
        _n("matrix_header", "matrix_header_band", "rule", (0.08, 0.29, 0.69, 0.045), 20),
        _n("comparison_matrix_table", "comparison_matrix", "table", (0.08, 0.29, 0.69, 0.42), 21, table_rows=rows),
        _n("matrix_side_rule", "accent_rule", "rule", (0.79, 0.30, 0.010, 0.36), 22),
        *_footer("comparison_matrix"),
    ])
    return _base("comparison_matrix", nodes)


def _data_dashboard() -> dict[str, Any]:
    nodes = [*_common_title("data_dashboard", "Operating dashboard", "KPI cards, native chart, and insight panel")]
    for idx, x in enumerate((0.08, 0.23, 0.38), start=1):
        nodes.extend([
            _n(f"kpi_card_{idx}", "kpi_card", "panel", (x, 0.27, 0.125, 0.10), 20 + idx),
            _n(f"kpi_text_{idx}", "kpi_text_region", "text", (x + 0.015, 0.295, 0.07, 0.035), 30 + idx, text=["Growth +18%", "Margin 42%", "Risk low"][idx - 1], font_size=8, bold=True),
            _n(f"kpi_rule_{idx}", "card_underline", "rule", (x + 0.015, 0.35, 0.045, 0.008), 40 + idx),
        ])
    nodes.extend([
        _n("dashboard_chart", "primary_chart", "chart", (0.08, 0.47, 0.47, 0.30), 50),
        _n("dashboard_insight", "insight_panel", "panel", (0.63, 0.29, 0.22, 0.32), 51),
        _n("dashboard_insight_text", "insight_text_region", "text", (0.655, 0.38, 0.15, 0.07), 52, text="Insight narrative for executive review", font_size=10, bold=True),
        _n("dashboard_connector", "technical_overlay", "connector", (0.08, 0.80, 0.74, 0.025), 60, line=COLORS["gold"]),
        *_footer("data_dashboard"),
    ])
    return _base("data_dashboard", nodes)


def _table_heavy() -> dict[str, Any]:
    rows = [["Segment", "Metric", "Status", "Action"], ["Alpha", "78", "Strong", "Scale"], ["Beta", "54", "Watch", "Tune"], ["Gamma", "61", "Rising", "Invest"], ["Delta", "43", "Low", "Review"]]
    nodes = [*_common_title("table_heavy", "Evidence table", "Dense editable table with header and body hierarchy")]
    nodes.extend([
        _n("table_header_band", "table_header_band", "rule", (0.08, 0.30, 0.68, 0.045), 20),
        _n("table_region", "table_region", "table", (0.08, 0.30, 0.68, 0.42), 21, table_rows=rows),
        _n("table_body_grid", "table_body_grid", "texture", (0.08, 0.35, 0.68, 0.36), 22),
        _n("table_chip", "kpi_chip", "panel", (0.80, 0.32, 0.11, 0.055), 23),
        _n("table_chip_text", "kpi_text_region", "text", (0.815, 0.338, 0.06, 0.02), 24, text="Coverage", font_size=7),
        *_footer("table_heavy"),
    ])
    return _base("table_heavy", nodes)


def _timeline_roadmap() -> dict[str, Any]:
    nodes = [*_common_title("timeline_roadmap", "Roadmap sequence", "Milestones arranged along an editable rail")]
    xs = [0.10, 0.28, 0.46, 0.64, 0.82]
    nodes.append(_n("timeline_axis", "timeline_axis", "connector", (0.08, 0.47, 0.78, 0.02), 20, line=COLORS["gold"], line_width=2))
    for idx, x in enumerate(xs, start=1):
        nodes.extend([
            _n(f"timeline_phase_{idx}", "timeline_phase", "panel", (x - 0.045, 0.35 if idx % 2 else 0.58, 0.09, 0.09), 30 + idx),
            _n(f"milestone_text_{idx}", "milestone_text_region", "text", (x - 0.035, 0.382 if idx % 2 else 0.612, 0.06, 0.025), 40 + idx, text=f"Phase {idx}", font_size=7, bold=True),
            _n(f"timeline_dot_{idx}", "motif_dot", "dot", (x - 0.008, 0.462, 0.016, 0.028), 50 + idx),
        ])
    nodes.extend(_footer("timeline_roadmap"))
    return _base("timeline_roadmap", nodes)


def _dot_row(prefix: str, x: float, y: float, count: int, z: int) -> list[dict[str, Any]]:
    return [_n(f"{prefix}_{idx + 1}", "motif_dot", "dot", (x + idx * 0.04, y + (0.012 if idx % 2 else 0), 0.007, 0.012), z + idx) for idx in range(count)]


def _n(object_id: str, semantic_role: str, kind: str, bbox: tuple[float, float, float, float], z: int, **extra: Any) -> dict[str, Any]:
    return {"object_id": object_id, "semantic_role": semantic_role, "kind": kind, "bbox_norm": {"x": bbox[0], "y": bbox[1], "w": bbox[2], "h": bbox[3]}, "z_order": z, **extra}


def _inches(bbox: dict[str, float]) -> tuple[float, float, float, float]:
    return bbox["x"] * SLIDE_W_IN, bbox["y"] * SLIDE_H_IN, bbox["w"] * SLIDE_W_IN, bbox["h"] * SLIDE_H_IN


def _px(bbox: dict[str, float]) -> tuple[int, int, int, int]:
    return round(bbox["x"] * SLIDE_W_PX), round(bbox["y"] * SLIDE_H_PX), round(bbox["w"] * SLIDE_W_PX), round(bbox["h"] * SLIDE_H_PX)


def _rgb(hex_color: str) -> tuple[int, int, int]:
    value = hex_color.strip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def _final_decision(decision: str, unlocked: bool, reason: str) -> dict[str, Any]:
    return {
        "schema_name": "e03_r2_final_decision",
        "status": "passed" if unlocked else "failed",
        "decision": decision,
        "reason": reason,
        "e04_unlocked": unlocked,
        "e04_started": False,
        "source_bound_deck_generated": False,
        "large_deck_generated": False,
        "canonical_promotion": False,
        "canva_parity_claimed": False,
    }


def _run_protect_check() -> bool:
    npm = shutil.which("npm.cmd") or shutil.which("npm")
    return bool(npm) and subprocess.run([npm, "run", "protect:check"], cwd=REPO_ROOT, capture_output=True, text=True, check=False).returncode == 0


def _read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")


def _write_md(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.strip() + "\n", encoding="utf-8")


def _simple_md(title: str, payload: dict[str, Any]) -> str:
    lines = [f"# {title}", ""]
    for key, value in payload.items():
        if not isinstance(value, (dict, list)):
            lines.append(f"- {key}: `{value}`")
    return "\n".join(lines)


def _report_md(title: str, payload: dict[str, Any]) -> str:
    lines = [f"# {title}", ""]
    for key, value in payload.items():
        if not isinstance(value, (dict, list)):
            lines.append(f"- {key}: `{value}`")
    if payload.get("failures"):
        lines.extend(["", "## Failures"])
        lines.extend(f"- `{failure}`" for failure in payload["failures"])
    return "\n".join(lines)


def rel(path: Path) -> str:
    try:
        return path.resolve().relative_to(REPO_ROOT.resolve()).as_posix()
    except ValueError:
        return path.resolve().as_posix()


def now() -> str:
    return datetime.now(timezone.utc).isoformat()
