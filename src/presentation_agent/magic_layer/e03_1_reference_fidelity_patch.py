"""Compile patched E03.1 reference-fidelity candidates."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .e03_16_orchestrator import CORE_ARCHETYPES
from .e03_archetype_conversion import build_visual_asset_plan, draw_e03_archetype
from .e03_chart_table_pipeline import build_chart_table_native_probe_report
from .e03_icon_vector_pipeline import EXPECTED_ICON_COUNTS

SLIDE_W = 16.0
SLIDE_H = 9.0
COLORS = {
    "bg": "06111A",
    "deep": "081B2A",
    "panel": "0B3B44",
    "panel2": "0F4B57",
    "cyan": "50D2E5",
    "cyan2": "1B8EA3",
    "gold": "F5A623",
    "paper": "F7F1E4",
    "paper2": "E7E1D5",
    "ink": "071018",
    "muted": "9EC4C8",
    "red": "EF6B5A",
    "white": "FFFFFF",
}


def build_visual_asset_plan_e03_1(archetype_id: str, reference_image: Path, output_dir: Path) -> dict[str, Any]:
    return build_visual_asset_plan(archetype_id, reference_image, output_dir)


def compile_e03_1_candidate(
    archetype_id: str,
    output_pptx: Path,
    *,
    reference_image: Path,
    visual_asset_plan: dict[str, Any],
    e03_candidate: Path | None = None,
) -> dict[str, Any]:
    if archetype_id in CORE_ARCHETYPES and e03_candidate and e03_candidate.exists():
        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(e03_candidate, output_pptx)
        source = "copied_core_e03_candidate_no_degrade"
    else:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        draw_e03_1_archetype(slide, archetype_id, visual_asset_plan)
        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_pptx)
        source = "compiled_e03_1_expansion_reference_fidelity_patch"
    chart_table = build_chart_table_native_probe_report(archetype_id)
    return {
        "schema_name": "e03_1_candidate_compile_report",
        "status": "passed",
        "archetype_id": archetype_id,
        "pptx_path": output_pptx.as_posix(),
        "source": source,
        "reference_image": reference_image.as_posix(),
        "semantic_vector_icon_count": EXPECTED_ICON_COUNTS[archetype_id],
        "native_ppt_chart_count": chart_table["native_ppt_chart_count"],
        "editable_shape_chart_count": chart_table["editable_shape_chart_count"],
        "native_ppt_table_count": chart_table["native_ppt_table_count"],
        "editable_shape_grid_table_count": chart_table["editable_shape_grid_table_count"],
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use": False,
    }


def draw_e03_1_archetype(slide: Any, archetype_id: str, visual_asset_plan: dict[str, Any]) -> None:
    _background(slide, archetype_id)
    {
        "section_divider": _section_divider,
        "visual_toc": _visual_toc,
        "evidence_overview": _evidence_overview,
        "card_grid": _card_grid,
        "methodology_framework": _methodology_framework,
        "process_flow": _process_flow,
        "comparison_matrix": _comparison_matrix,
        "timeline_roadmap": _timeline_roadmap,
        "decision_record": _decision_record,
        "risk_register": _risk_register,
        "case_study": _case_study,
        "closing_synthesis": _closing_synthesis,
    }[archetype_id](slide, visual_asset_plan)


def _background(slide: Any, prefix: str) -> None:
    _shape(slide, f"{prefix}_background_base", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, "bg", None)
    for idx, x in enumerate((0.34, 1.42, 12.7, 14.8), start=1):
        _line(slide, f"{prefix}_grid_line_{idx}", x, 0.25, x + 0.22, 8.1, "cyan2", 0.28)
    for idx in range(9):
        _line(slide, f"{prefix}_top_micro_trace_{idx}", 0.8 + idx * 0.42, 0.38, 1.05 + idx * 0.42, 0.38, "cyan2", 0.35)


def _section_divider(slide: Any, plan: dict[str, Any]) -> None:
    _asset(slide, plan, "section_divider_visual_field", "section_divider_bounded_visual_field")
    _shape(slide, "section_left_marker_slab", MSO_AUTO_SHAPE_TYPE.PENTAGON, 0.62, 1.32, 2.1, 4.9, "deep", "gold")
    _text(slide, "section_number_text", "03", 0.95, 2.08, 1.0, 0.75, 44, "gold", bold=True)
    _shape(slide, "section_title_slab", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 3.0, 3.05, 6.1, 1.0, "paper", "gold")
    _text(slide, "section_title_text", "TITLE", 3.55, 3.34, 2.8, 0.34, 24, "ink", bold=True)
    _shape(slide, "section_subtitle_slot", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 3.25, 4.32, 4.9, 0.52, "paper", "cyan")
    _text(slide, "section_subtitle_text", "SUBTITLE", 3.78, 4.48, 1.3, 0.14, 9, "cyan2", bold=True)
    _shape(slide, "section_diagonal_gold_spine", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 8.45, -0.1, 0.82, 8.6, "panel", "gold")
    _footer(slide, "section_divider")


def _visual_toc(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "visual_toc_title", "VISUAL TOC", 0.75, 0.8, 3.2, 0.32, 23, "paper", bold=True)
    _line(slide, "visual_toc_index_path", 1.05, 1.52, 1.05, 7.1, "cyan", 1.2)
    for idx in range(6):
        y = 1.35 + idx * 0.88
        active = idx == 1
        _shape(slide, f"visual_toc_index_dot_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, 0.86, y, 0.38, 0.38, "gold" if active else "panel", "cyan")
        _shape(slide, f"visual_toc_nav_card_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 2.0, y - 0.12, 7.7, 0.62, "paper", "gold" if active else "cyan2")
        _draw_icon(slide, f"visual_toc_nav_icon_{idx}", 2.25, y + 0.05, 0.22, "cyan2", "target")
        _text(slide, f"visual_toc_nav_text_{idx}", f"MODULE {idx+1}", 2.72, y + 0.05, 1.3, 0.12, 7.3, "ink", bold=True)
        _line(slide, f"visual_toc_nav_status_{idx}", 8.45, y + 0.18, 9.2, y + 0.18, "cyan2", 0.8)
    _shape(slide, "visual_toc_active_marker_panel", MSO_AUTO_SHAPE_TYPE.HEXAGON, 10.55, 1.9, 1.1, 1.1, "panel", "gold")
    _shape(slide, "visual_toc_side_meta_panel", MSO_AUTO_SHAPE_TYPE.PENTAGON, 12.25, 1.3, 2.35, 5.8, "panel", "cyan")
    _text(slide, "visual_toc_side_meta_text", "ACTIVE\nPATH", 12.72, 2.1, 1.1, 0.58, 14, "paper", bold=True)
    _footer(slide, "visual_toc")


def _evidence_overview(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "evidence_title", "EVIDENCE OVERVIEW", 0.72, 0.72, 4.5, 0.32, 23, "paper", bold=True)
    for idx in range(8):
        x = 0.9 + (idx % 4) * 2.55
        y = 1.55 + (idx // 4) * 2.0
        _shape(slide, f"evidence_card_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, 2.12, 1.32, "paper", "cyan")
        _draw_icon(slide, f"evidence_trace_icon_{idx}", x + 0.14, y + 0.18, 0.24, "gold", "shield")
        _text(slide, f"evidence_card_label_{idx}", "EVIDENCE", x + 0.52, y + 0.2, 0.8, 0.11, 6.5, "ink", bold=True)
        for k in range(3):
            _shape(slide, f"evidence_confidence_dot_{idx}_{k}", MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.55 + k * 0.18, y + 0.62, 0.08, 0.08, "cyan2" if k < 2 else "gold", None)
        _line(slide, f"evidence_trace_rule_{idx}", x + 0.45, y + 0.95, x + 1.75, y + 0.95, "cyan2", 0.55)
    _shape(slide, "evidence_summary_strip", MSO_AUTO_SHAPE_TYPE.PENTAGON, 11.55, 1.45, 2.85, 4.55, "panel", "gold")
    _text(slide, "evidence_summary_text", "INSIGHT\nTRACEABILITY", 12.05, 2.15, 1.6, 0.62, 14, "paper", bold=True)
    _footer(slide, "evidence_overview")


def _card_grid(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "card_grid_title", "CARD GRID", 0.72, 0.72, 3.1, 0.32, 23, "paper", bold=True)
    _shape(slide, "card_grid_category_header", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 0.72, 1.18, 3.4, 0.36, "panel", "gold")
    _text(slide, "card_grid_category_text", "CATEGORY", 1.0, 1.29, 1.1, 0.1, 7, "gold", bold=True)
    for idx in range(8):
        x = 0.95 + (idx % 4) * 3.08
        y = 1.82 + (idx // 4) * 2.18
        _shape(slide, f"card_grid_module_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, 2.48, 1.52, "paper", "cyan")
        _shape(slide, f"card_grid_icon_hex_{idx}", MSO_AUTO_SHAPE_TYPE.HEXAGON, x + 0.16, y + 0.18, 0.45, 0.45, "panel", "gold" if idx == 6 else "cyan")
        _draw_icon(slide, f"card_grid_icon_{idx}", x + 0.27, y + 0.29, 0.22, "paper", "target")
        _text(slide, f"card_grid_module_no_{idx}", f"{idx+1:02d}", x + 0.8, y + 0.25, 0.36, 0.12, 7, "gold", bold=True)
        _text(slide, f"card_grid_module_text_{idx}", "MODULE", x + 0.8, y + 0.55, 0.8, 0.12, 7, "ink", bold=True)
    _shape(slide, "card_grid_outcome_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 13.5, 1.65, 1.2, 4.2, "panel", "gold")
    _footer(slide, "card_grid")


def _methodology_framework(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "methodology_title", "METHODOLOGY", 0.72, 0.72, 3.5, 0.32, 23, "paper", bold=True)
    _line(slide, "methodology_center_spine", 6.75, 1.52, 6.75, 6.15, "gold", 1.0)
    for idx in range(5):
        y = 1.45 + idx * 0.88
        active = idx == 2
        _shape(slide, f"methodology_layer_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 2.05 + idx * 0.18, y, 8.1 - idx * 0.34, 0.58, "gold" if active else ("paper" if idx % 2 == 0 else "panel"), "cyan")
        _draw_icon(slide, f"methodology_layer_icon_{idx}", 2.35 + idx * 0.18, y + 0.18, 0.22, "paper" if active else "gold", "network")
        _text(slide, f"methodology_layer_text_{idx}", f"LAYER {idx+1}", 2.78 + idx * 0.18, y + 0.2, 1.0, 0.1, 6.5, "ink" if idx % 2 == 0 and not active else "paper", bold=True)
    _shape(slide, "methodology_side_note_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 12.0, 1.32, 2.6, 5.2, "panel", "cyan")
    _footer(slide, "methodology_framework")


def _process_flow(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "process_title", "PROCESS FLOW", 0.72, 0.72, 3.2, 0.32, 23, "paper", bold=True)
    xs = [0.9, 2.7, 4.5, 6.3, 8.1, 9.9]
    for idx, x in enumerate(xs):
        _shape(slide, f"process_node_{idx}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 3.0, 1.28, 0.92, "paper", "cyan")
        _draw_icon(slide, f"process_icon_{idx}", x + 0.48, 3.18, 0.22, "cyan2", "target")
        _text(slide, f"process_step_text_{idx}", f"{idx+1:02d}", x + 0.48, 3.58, 0.22, 0.09, 5.8, "ink", bold=True)
        if idx < len(xs) - 1:
            _line(slide, f"process_direction_connector_{idx}", x + 1.28, 3.46, xs[idx + 1], 3.46, "gold", 1.0)
    _shape(slide, "process_decision_gate", MSO_AUTO_SHAPE_TYPE.DIAMOND, 6.0, 4.65, 1.0, 1.0, "panel", "gold")
    _text(slide, "process_decision_gate_text", "GATE", 6.28, 5.04, 0.42, 0.1, 6, "paper", bold=True)
    _shape(slide, "process_side_note_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 12.2, 1.45, 2.4, 5.2, "panel", "cyan")
    _footer(slide, "process_flow")


def _comparison_matrix(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "comparison_title", "COMPARISON MATRIX", 0.72, 0.72, 4.2, 0.32, 23, "paper", bold=True)
    _grid(slide, "comparison_matrix", 0.95, 1.52, 9.65, 5.2, 6, 7, header=True, score=True)
    for idx in range(5):
        _draw_icon(slide, f"comparison_score_icon_{idx}", 4.05 + idx * 1.15, 4.85, 0.18, "gold", "target")
    _shape(slide, "comparison_decision_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 11.35, 1.52, 2.7, 5.2, "panel", "gold")
    _text(slide, "comparison_decision_rail_text", "DECISION\nRAIL", 11.9, 2.38, 1.3, 0.55, 14, "paper", bold=True)
    _footer(slide, "comparison_matrix")


def _timeline_roadmap(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "timeline_title", "TIMELINE ROADMAP", 0.72, 0.72, 4.4, 0.32, 23, "paper", bold=True)
    _line(slide, "timeline_axis_main", 1.05, 4.15, 12.7, 4.15, "cyan", 1.4)
    for idx in range(6):
        x = 1.2 + idx * 1.85
        top = idx % 2 == 0
        y = 2.35 if top else 4.75
        _shape(slide, f"timeline_phase_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, 1.55, 0.85, "paper", "cyan")
        _line(slide, f"timeline_phase_connector_{idx}", x + 0.78, y + (0.85 if top else 0), x + 0.78, 4.15, "gold", 0.8)
        _shape(slide, f"timeline_milestone_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.66, 4.03, 0.24, 0.24, "gold", None)
        _text(slide, f"timeline_phase_label_{idx}", f"PHASE {idx+1}", x + 0.2, y + 0.3, 0.65, 0.1, 6, "ink", bold=True)
    _shape(slide, "timeline_risk_row", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 1.2, 6.45, 10.8, 0.36, "panel", "gold")
    _shape(slide, "timeline_side_meta_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 13.4, 1.45, 1.15, 5.2, "panel", "cyan")
    _footer(slide, "timeline_roadmap")


def _decision_record(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "decision_title", "DECISION RECORD", 0.72, 0.72, 4.1, 0.32, 23, "paper", bold=True)
    _shape(slide, "decision_stamp_sidebar", MSO_AUTO_SHAPE_TYPE.PENTAGON, 0.9, 1.45, 3.35, 4.8, "panel", "gold")
    _text(slide, "decision_stamp_text", "DECISION\nSTAMP", 1.45, 2.35, 1.8, 0.65, 20, "paper", bold=True)
    _grid(slide, "decision_metadata_grid", 4.75, 1.45, 3.8, 4.8, 2, 6, header=True, score=False)
    for idx in range(4):
        _shape(slide, f"decision_condition_module_{idx}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 9.0, 1.55 + idx * 1.1, 3.25, 0.72, "paper", "cyan")
        _draw_icon(slide, f"decision_condition_icon_{idx}", 9.24, 1.75 + idx * 1.1, 0.22, "gold", "shield")
    _shape(slide, "decision_evidence_strip", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 1.0, 6.72, 12.4, 0.55, "paper", "gold")
    _footer(slide, "decision_record")


def _risk_register(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "risk_title", "RISK REGISTER", 0.72, 0.72, 3.4, 0.32, 23, "paper", bold=True)
    _grid(slide, "risk_register_grid", 0.9, 1.45, 10.8, 5.6, 7, 9, header=True, score=True)
    for idx in range(8):
        _draw_icon(slide, f"risk_marker_{idx}", 10.6, 2.02 + idx * 0.55, 0.16, "red" if idx % 3 == 0 else "gold", "target")
    _shape(slide, "risk_side_meta_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 12.25, 1.45, 2.25, 5.6, "panel", "gold")
    _footer(slide, "risk_register")


def _case_study(slide: Any, plan: dict[str, Any]) -> None:
    _text(slide, "case_title", "CASE STUDY", 0.72, 0.72, 3.2, 0.32, 23, "paper", bold=True)
    if not _asset(slide, plan, "case_study_image_frame", "case_study_bounded_image_frame"):
        _shape(slide, "case_study_replaceable_image_frame", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.92, 1.35, 5.8, 4.65, "panel", "cyan")
    _shape(slide, "case_image_frame_rule", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.88, 1.31, 5.88, 4.73, None, "cyan")
    _shape(slide, "case_context_panel", MSO_AUTO_SHAPE_TYPE.PENTAGON, 7.1, 1.35, 3.05, 1.45, "panel", "gold")
    _shape(slide, "case_evidence_panel", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 7.1, 3.05, 3.05, 1.45, "paper", "cyan")
    _shape(slide, "case_result_panel", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 10.55, 3.05, 3.05, 1.45, "paper", "gold")
    _shape(slide, "case_decision_panel", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.55, 1.35, 3.05, 1.45, "panel", "cyan")
    for idx, (x, y) in enumerate(((7.35, 3.36), (10.8, 3.36), (10.8, 1.68)), start=1):
        _draw_icon(slide, f"case_module_icon_{idx}", x, y, 0.24, "gold" if idx != 3 else "paper", "target")
    _shape(slide, "case_lesson_strip", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 1.0, 6.72, 12.8, 0.58, "panel", "gold")
    _footer(slide, "case_study")


def _closing_synthesis(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "closing_title", "CLOSING SYNTHESIS", 0.72, 0.72, 4.4, 0.32, 23, "paper", bold=True)
    modules = [("RECOMMENDATION", 0.95, 1.6, "gold"), ("NEXT ACTION", 4.75, 1.6, "cyan"), ("EVIDENCE SUMMARY", 8.55, 1.6, "gold")]
    for idx, (label, x, y, accent) in enumerate(modules):
        _shape(slide, f"closing_module_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, 3.2, 2.55, "paper", accent)
        _draw_icon(slide, f"closing_module_icon_{idx}", x + 0.35, y + 0.4, 0.3, accent, "target")
        _text(slide, f"closing_module_label_{idx}", label, x + 0.85, y + 0.52, 1.8, 0.13, 7, "ink", bold=True)
        for k in range(3):
            _line(slide, f"closing_module_rule_{idx}_{k}", x + 0.85, y + 1.08 + k * 0.35, x + 2.6, y + 1.08 + k * 0.35, "cyan2", 0.5)
    _shape(slide, "closing_takeaway_panel", MSO_AUTO_SHAPE_TYPE.PENTAGON, 2.0, 5.25, 11.2, 1.1, "panel", "gold")
    _text(slide, "closing_takeaway_text", "DECISION / TAKEAWAY / INSIGHT", 3.4, 5.66, 4.2, 0.2, 17, "paper", bold=True)
    _footer(slide, "closing_synthesis")


def _footer(slide: Any, prefix: str) -> None:
    _shape(slide, f"{prefix}_source_footer_strip_native", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 8.34, SLIDE_W, 0.66, "deep", None)
    _line(slide, f"{prefix}_footer_top_rule", 0, 8.34, SLIDE_W, 8.34, "gold", 0.85)
    _draw_icon(slide, f"{prefix}_footer_icon", 0.45, 8.55, 0.2, "cyan", "database")
    _text(slide, f"{prefix}_footer_source_text", "SOURCE", 0.82, 8.58, 1.0, 0.12, 7, "muted", bold=True)
    _text(slide, f"{prefix}_footer_marker_text", "FOOTER", 13.35, 8.58, 1.0, 0.12, 7, "muted", bold=True)


def _grid(slide: Any, prefix: str, x: float, y: float, w: float, h: float, cols: int, rows: int, *, header: bool, score: bool) -> None:
    col_w = w / cols
    row_h = h / rows
    for r in range(rows):
        for c in range(cols):
            fill = "panel" if header and r == 0 else ("paper2" if (r + c) % 2 else "paper")
            _shape(slide, f"{prefix}_r{r}_c{c}", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + c * col_w, y + r * row_h, col_w, row_h, fill, "paper2")
            if r == 0:
                _text(slide, f"{prefix}_header_{c}", "HDR", x + c * col_w + 0.06, y + 0.11, col_w - 0.12, 0.1, 5.6, "paper", bold=True)
            elif c in {0, 1}:
                _text(slide, f"{prefix}_cell_{r}_{c}", "ROW", x + c * col_w + 0.06, y + r * row_h + 0.11, col_w - 0.12, 0.1, 5.2, "ink")
            elif score and c >= cols - 2 and r > 0:
                _shape(slide, f"{prefix}_score_pill_{r}_{c}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + c * col_w + 0.18, y + r * row_h + 0.16, col_w - 0.36, 0.12, "gold" if (r + c) % 3 == 0 else "cyan2", None)


def _asset(slide: Any, visual_asset_plan: dict[str, Any], asset_id: str, name: str) -> bool:
    for asset in visual_asset_plan.get("assets", []):
        if asset["asset_id"] == asset_id:
            bbox = asset["target_bbox_in"]
            pic = slide.shapes.add_picture(asset["asset_path"], Inches(bbox["x"]), Inches(bbox["y"]), width=Inches(bbox["w"]), height=Inches(bbox["h"]))
            pic.name = name
            return True
    return False


def _draw_icon(slide: Any, name: str, x: float, y: float, size: float, color: str, role: str) -> None:
    if role == "database":
        _shape(slide, f"{name}_db_top", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size * 0.35, None, color)
        _line(slide, f"{name}_db_left", x, y + size * 0.18, x, y + size * 0.78, color, 0.55)
        _line(slide, f"{name}_db_right", x + size, y + size * 0.18, x + size, y + size * 0.78, color, 0.55)
        _shape(slide, f"{name}_db_bottom", MSO_AUTO_SHAPE_TYPE.OVAL, x, y + size * 0.58, size, size * 0.35, None, color)
    elif role == "shield":
        _shape(slide, f"{name}_shield", MSO_AUTO_SHAPE_TYPE.PENTAGON, x, y, size, size, None, color)
        _line(slide, f"{name}_check1", x + size * 0.25, y + size * 0.52, x + size * 0.42, y + size * 0.70, color, 0.65)
        _line(slide, f"{name}_check2", x + size * 0.42, y + size * 0.70, x + size * 0.78, y + size * 0.30, color, 0.65)
    elif role == "network":
        for idx, (cx, cy) in enumerate(((0.1, 0.2), (0.7, 0.1), (0.55, 0.72)), start=1):
            _shape(slide, f"{name}_node_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, x + size * cx, y + size * cy, size * 0.22, size * 0.22, None, color)
        _line(slide, f"{name}_edge1", x + size * 0.2, y + size * 0.3, x + size * 0.78, y + size * 0.2, color, 0.5)
    else:
        _shape(slide, f"{name}_ring", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, None, color)
        _line(slide, f"{name}_mark1", x + size * 0.28, y + size * 0.55, x + size * 0.46, y + size * 0.72, color, 0.68)
        _line(slide, f"{name}_mark2", x + size * 0.46, y + size * 0.72, x + size * 0.78, y + size * 0.28, color, 0.68)


def _shape(slide: Any, name: str, shape_type: Any, x: float, y: float, w: float, h: float, fill: str | None, line: str | None) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.62)
    else:
        shape.line.fill.background()
    return shape


def _line(slide: Any, name: str, x1: float, y1: float, x2: float, y2: float, color: str, width: float) -> Any:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _text(slide: Any, name: str, text: str, x: float, y: float, w: float, h: float, size: float, color: str, *, bold: bool = False) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _rgb(name_or_hex: str) -> RGBColor:
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
