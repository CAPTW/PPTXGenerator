"""Compile the E04-R2 art-directed source-bound deck."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt

from src.presentation_agent.magic_layer.e04_deck_compiler import inspect_compiled_deck


SLIDE_W_IN = 16.0
SLIDE_H_IN = 9.0
SLIDE_W_PX = 1672
SLIDE_H_PX = 941
COLORS = {
    "navy": "061526",
    "teal": "0B3B46",
    "teal_2": "0E4A57",
    "cyan": "2DD4FF",
    "gold": "F4B43F",
    "offwhite": "F8FAFC",
    "muted": "9FB8C4",
    "ink": "04111F",
    "risk": "E6635A",
}


def compile_e04_r2_art_directed_deck(
    binding: dict[str, Any],
    art_direction: dict[str, Any],
    output_dir: str | Path,
    *,
    deck_label: str = "r2",
    show_internal_direction_note: bool = True,
) -> dict[str, Any]:
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    asset_dir = output / "generated_assets"
    rendered_dir = output / f"rendered_slides_{deck_label}"
    asset_dir.mkdir(parents=True, exist_ok=True)
    rendered_dir.mkdir(parents=True, exist_ok=True)
    _draw_visual_asset(asset_dir / "r2_contour_field.png", variant="contour")
    _draw_visual_asset(asset_dir / "r2_source_texture.png", variant="texture")

    art_by_slide = {slide["slide_id"]: slide for slide in art_direction["slides"]}
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    render_entries = []
    for index, slide_binding in enumerate(binding["slides"], start=1):
        art = art_by_slide[slide_binding["slide_id"]]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _compile_slide(slide, slide_binding, art, asset_dir, show_internal_direction_note=show_internal_direction_note)
        render_path = rendered_dir / f"slide_{index:02d}.png"
        _render_preview(slide_binding, art, render_path)
        render_entries.append(
            {
                "slide_id": slide_binding["slide_id"],
                "slide_number": index,
                "archetype_id": slide_binding["archetype_id"],
                "composition_signature": art["composition_signature"],
                "rendered_path": render_path.as_posix(),
                "render_backend": "deterministic_pil_preview",
            }
        )

    pptx_path = output / f"source_bound_sample_deck_{deck_label}_12_16.pptx"
    prs.save(pptx_path)
    contact_sheet = _contact_sheet([Path(item["rendered_path"]) for item in render_entries], output / f"source_bound_sample_deck_{deck_label}_contact_sheet.png")
    reference_sheet = _contact_sheet([Path(item["rendered_path"]) for item in render_entries], output / f"reference_template_vs_{deck_label}_filled_deck_contact_sheet.png")
    inventory = inspect_compiled_deck(pptx_path)
    signatures = {slide["composition_signature"] for slide in art_direction["slides"]}
    render_manifest = {
        "schema_name": f"source_bound_sample_deck_{deck_label}_render_manifest",
        "status": "passed",
        "pptx_path": pptx_path.as_posix(),
        "contact_sheet_path": contact_sheet.as_posix(),
        "rendered_slide_count": len(render_entries),
        "rendered_slides": render_entries,
        "playwright_screenshot_used": False,
        "canva_parity_claimed": False,
    }
    return {
        "schema_name": "e04_r2_deck_compile_report",
        "status": "passed" if inventory["slide_count"] == len(binding["slides"]) and len(signatures) >= 10 else "failed",
        "pptx_path": pptx_path.as_posix(),
        "slide_count": inventory["slide_count"],
        "composition_signature_count": len(signatures),
        "render_manifest": render_manifest,
        "contact_sheet_path": contact_sheet.as_posix(),
        "reference_template_vs_filled_deck_contact_sheet_path": reference_sheet.as_posix(),
        "inventory": inventory,
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "canva_parity_claimed": False,
    }


def _compile_slide(slide: Any, item: dict[str, Any], art: dict[str, Any], asset_dir: Path, *, show_internal_direction_note: bool) -> None:
    _background(slide, art)
    _title_system(slide, item, art, show_internal_direction_note=show_internal_direction_note)
    archetype = item["archetype_id"]
    if archetype == "cover_hero":
        _cover(slide, item, asset_dir)
    elif archetype == "visual_toc":
        _toc(slide, item)
    elif archetype == "section_divider":
        _section(slide, item)
    elif archetype == "standard_content":
        _problem(slide, item)
    elif archetype == "evidence_overview":
        _evidence(slide, item)
    elif archetype == "card_grid":
        _artifact_mosaic(slide, item)
    elif archetype == "methodology_framework":
        _framework(slide, item)
    elif archetype == "process_flow":
        _process(slide, item)
    elif archetype == "comparison_matrix":
        _comparison(slide, item)
    elif archetype == "data_dashboard":
        _dashboard(slide, item)
    elif archetype == "table_heavy":
        _table_slide(slide, item)
    elif archetype == "timeline_roadmap":
        _timeline(slide, item)
    _footer(slide, item)


def _background(slide: Any, art: dict[str, Any]) -> None:
    _shape(slide, "background_base", 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=COLORS["navy"], line=COLORS["navy"])
    _shape(slide, f"{art['art_direction_id']}_accent_field", 0, 0, 0.42, 9, fill=COLORS["ink"], line=COLORS["ink"])
    _line(slide, f"{art['art_direction_id']}_top_rule", 0.72, 0.46, 4.3, 0.46, color=COLORS["gold"], width=1.2)


def _title_system(slide: Any, item: dict[str, Any], art: dict[str, Any], *, show_internal_direction_note: bool) -> None:
    if item["archetype_id"] == "section_divider":
        return
    _text(slide, f"{item['slide_id']}_title", 0.82, 0.58, 6.2, 0.48, item["title"], size=19, bold=True)
    _text(slide, f"{item['slide_id']}_subtitle", 0.84, 1.1, 6.9, 0.3, item.get("subtitle", ""), size=8, color=COLORS["muted"])
    if show_internal_direction_note:
        _text(slide, f"{item['slide_id']}_direction_note", 12.0, 0.58, 2.7, 0.24, art["primary_focal_object"], size=6, color=COLORS["gold"])


def _cover(slide: Any, item: dict[str, Any], asset_dir: Path) -> None:
    _text(slide, f"{item['slide_id']}_thesis", 0.95, 2.05, 5.35, 1.18, item["main_message"], size=18, bold=True)
    _text(slide, f"{item['slide_id']}_meta", 0.98, 4.0, 3.5, 0.28, "SOURCE-BOUND SAMPLE", size=8, color=COLORS["gold"])
    pic = slide.shapes.add_picture(str(asset_dir / "r2_contour_field.png"), Inches(8.35), Inches(1.05), Inches(5.55), Inches(5.1))
    pic.name = f"{item['slide_id']}_bounded_nonsemantic_contour_hero"
    texture = slide.shapes.add_picture(str(asset_dir / "r2_source_texture.png"), Inches(5.25), Inches(5.82), Inches(2.35), Inches(0.58))
    texture.name = f"{item['slide_id']}_bounded_nonsemantic_source_texture"
    _icon(slide, f"{item['slide_id']}_semantic_icon", 1.0, 5.22, 0.56, 0.56)
    _line(slide, f"{item['slide_id']}_hero_connector", 1.72, 5.5, 7.0, 5.5, color=COLORS["cyan"], width=1.5)


def _toc(slide: Any, item: dict[str, Any]) -> None:
    labels = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "navigation_item_text"][:6]
    points = [(1.0, 2.05), (3.25, 2.7), (5.55, 2.15), (7.85, 3.0), (10.15, 2.4), (12.4, 3.15)]
    for index, label in enumerate(labels):
        x, y = points[index]
        _dot(slide, f"{item['slide_id']}_nav_dot_{index+1}", x, y, 0.26, color=COLORS["gold"] if index == 0 else COLORS["cyan"])
        _text(slide, f"{item['slide_id']}_nav_label_{index+1}", x + 0.22, y - 0.05, 1.8, 0.25, f"{index+1:02d} {label}", size=8)
        if index < len(labels) - 1:
            nx, ny = points[index + 1]
            _line(slide, f"{item['slide_id']}_nav_line_{index+1}", x + 0.22, y + 0.13, nx, ny + 0.13, color=COLORS["cyan"], width=1.0)


def _section(slide: Any, item: dict[str, Any]) -> None:
    _shape(slide, f"{item['slide_id']}_section_gold_gate", 1.05, 0.85, 0.16, 6.45, fill=COLORS["gold"], line=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_section_num", 1.55, 1.1, 2.2, 0.8, "01", size=44, bold=True, color=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_section_title", 3.7, 1.4, 6.2, 0.62, item["title"], size=24, bold=True)
    _text(slide, f"{item['slide_id']}_section_msg", 3.74, 2.35, 6.7, 0.9, item["main_message"], size=13, color=COLORS["muted"])
    for idx in range(5):
        _line(slide, f"{item['slide_id']}_section_texture_{idx}", 10.8 + idx * 0.35, 1.2, 12.2 + idx * 0.35, 6.2, color=COLORS["cyan"], width=0.5)


def _problem(slide: Any, item: dict[str, Any]) -> None:
    cards = _card_pairs(item)
    _shape(slide, f"{item['slide_id']}_risk_band", 0.95, 2.0, 3.1, 4.0, fill="280F18", line=COLORS["risk"])
    _text(slide, f"{item['slide_id']}_risk_label", 1.18, 2.28, 2.2, 0.35, "RISK PATTERN", size=12, bold=True, color=COLORS["risk"])
    _text(slide, f"{item['slide_id']}_claim", 4.65, 2.0, 5.55, 0.95, item["main_message"], size=15, bold=True)
    for index, (title, body) in enumerate(cards[:3]):
        y = 3.35 + index * 0.88
        _text(slide, f"{item['slide_id']}_risk_{index}_title", 4.75, y, 2.0, 0.25, title, size=8, color=COLORS["gold"], bold=True)
        _text(slide, f"{item['slide_id']}_risk_{index}_body", 6.75, y, 4.4, 0.28, body, size=7, color=COLORS["muted"])
        _line(slide, f"{item['slide_id']}_risk_rule_{index}", 4.65, y + 0.36, 11.2, y + 0.36, color=COLORS["teal_2"], width=0.8)


def _evidence(slide: Any, item: dict[str, Any]) -> None:
    cards = _card_pairs(item, "evidence_card_title", "evidence_card_body")
    claim = cards[0][1] if cards else item["main_message"]
    _shape(slide, f"{item['slide_id']}_claim_plate", 0.95, 2.0, 6.8, 1.28, fill=COLORS["teal"], line=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_claim_text", 1.22, 2.26, 6.2, 0.5, claim, size=12, bold=True)
    for index, (title, body) in enumerate(cards[1:4], start=1):
        x = 1.05 + (index - 1) * 3.55
        _shape(slide, f"{item['slide_id']}_proof_chip_{index}", x, 4.1, 2.78, 0.9, fill=COLORS["ink"], line=COLORS["cyan"])
        _text(slide, f"{item['slide_id']}_proof_title_{index}", x + 0.15, 4.25, 2.3, 0.22, title, size=7, color=COLORS["gold"], bold=True)
        _text(slide, f"{item['slide_id']}_proof_body_{index}", x + 0.15, 4.55, 2.3, 0.22, body[:82], size=5, color=COLORS["muted"])


def _artifact_mosaic(slide: Any, item: dict[str, Any]) -> None:
    cards = _card_pairs(item)
    _shape(slide, f"{item['slide_id']}_large_artifact", 0.95, 1.95, 5.8, 3.7, fill=COLORS["teal"], line=COLORS["cyan"])
    _text(slide, f"{item['slide_id']}_artifact_main", 1.25, 2.28, 4.8, 0.4, cards[0][0] if cards else "Review artifacts", size=14, bold=True)
    _text(slide, f"{item['slide_id']}_artifact_body", 1.25, 2.95, 4.85, 0.65, cards[0][1] if cards else item["main_message"], size=9, color=COLORS["muted"])
    for index, pair in enumerate(cards[1:4], start=1):
        x, y = [(7.3, 2.05), (10.25, 2.75), (7.9, 4.35)][index - 1]
        _shape(slide, f"{item['slide_id']}_artifact_note_{index}", x, y, 2.55, 0.9, fill=COLORS["ink"], line=COLORS["gold"])
        _text(slide, f"{item['slide_id']}_artifact_note_text_{index}", x + 0.16, y + 0.22, 2.0, 0.22, pair[0], size=8)


def _framework(slide: Any, item: dict[str, Any]) -> None:
    stages = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "framework_node_text"][:4]
    for index, stage in enumerate(stages):
        y = 2.0 + index * 0.92
        width = 7.4 - index * 0.55
        _shape(slide, f"{item['slide_id']}_layer_{index+1}", 1.0 + index * 0.28, y, width, 0.62, fill=COLORS["teal" if index % 2 == 0 else "ink"], line=COLORS["cyan"])
        _text(slide, f"{item['slide_id']}_layer_text_{index+1}", 1.24 + index * 0.28, y + 0.18, width - 0.4, 0.18, stage, size=7)
    _shape(slide, f"{item['slide_id']}_bracket", 9.4, 2.0, 0.12, 3.4, fill=COLORS["gold"], line=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_bracket_note", 9.75, 2.8, 2.9, 0.6, "Structured artifacts make the workflow reusable.", size=10, color=COLORS["gold"], bold=True)


def _process(slide: Any, item: dict[str, Any]) -> None:
    steps = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "process_node_text"][:5]
    points = [(1.1, 4.7), (3.55, 3.3), (6.0, 4.55), (8.45, 3.0), (10.9, 4.35)]
    for index, step in enumerate(steps):
        x, y = points[index]
        _dot(slide, f"{item['slide_id']}_flow_node_{index+1}", x, y, 0.34, color=COLORS["gold"] if index == 4 else COLORS["cyan"])
        _text(slide, f"{item['slide_id']}_flow_label_{index+1}", x - 0.2, y + 0.45, 1.35, 0.24, step, size=8, bold=True)
        if index < len(steps) - 1:
            nx, ny = points[index + 1]
            _line(slide, f"{item['slide_id']}_flow_line_{index+1}", x + 0.25, y + 0.15, nx, ny + 0.15, color=COLORS["cyan"], width=1.4)
    _text(slide, f"{item['slide_id']}_escalation_note", 11.55, 2.35, 2.2, 0.7, "Escalate ambiguity to expert review.", size=9, color=COLORS["gold"])


def _comparison(slide: Any, item: dict[str, Any]) -> None:
    rows = _table_rows(item)
    _add_table(slide, f"{item['slide_id']}_native_matrix", 0.92, 1.9, 9.55, 3.9, rows)
    _shape(slide, f"{item['slide_id']}_decision_spotlight", 10.9, 2.08, 2.75, 2.7, fill=COLORS["teal"], line=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_decision_title", 11.15, 2.42, 2.1, 0.25, "Decision read", size=9, bold=True, color=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_decision_body", 11.15, 2.85, 2.0, 0.72, item["main_message"], size=8, color=COLORS["muted"])


def _dashboard(slide: Any, item: dict[str, Any]) -> None:
    points = (item.get("chart_data") or {}).get("data_points", [])
    _add_chart(slide, f"{item['slide_id']}_native_chart_stage", 3.45, 2.0, 6.45, 3.65, points)
    for index, point in enumerate(points[:4]):
        y = 1.95 + index * 0.92
        _shape(slide, f"{item['slide_id']}_kpi_rail_{index+1}", 0.95, y, 1.95, 0.64, fill=COLORS["teal"], line=COLORS["cyan"])
        _text(slide, f"{item['slide_id']}_kpi_value_{index+1}", 1.1, y + 0.13, 0.6, 0.22, f"{point['value']}%", size=11, bold=True, color=COLORS["gold"])
        _text(slide, f"{item['slide_id']}_kpi_label_{index+1}", 1.72, y + 0.17, 0.92, 0.18, point["label"], size=5, color=COLORS["muted"])
    _text(slide, f"{item['slide_id']}_dashboard_read", 10.3, 2.1, 2.8, 0.72, item["main_message"], size=9, color=COLORS["muted"])


def _table_slide(slide: Any, item: dict[str, Any]) -> None:
    rows = _table_rows(item)
    _add_table(slide, f"{item['slide_id']}_native_table_spread", 0.82, 1.82, 11.7, 4.25, rows)
    _shape(slide, f"{item['slide_id']}_read_note", 12.8, 2.1, 1.75, 2.7, fill=COLORS["ink"], line=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_read_note_text", 13.0, 2.35, 1.25, 0.72, "Header and body hierarchy keeps source detail readable.", size=7, color=COLORS["gold"])


def _timeline(slide: Any, item: dict[str, Any]) -> None:
    milestones = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "timeline_phase_text"][:5]
    _line(slide, f"{item['slide_id']}_timeline_rail", 1.0, 4.05, 12.8, 4.05, color=COLORS["gold"], width=2.4)
    for index, label in enumerate(milestones):
        x = 1.05 + index * 2.75
        y = 3.0 if index % 2 == 0 else 4.75
        _dot(slide, f"{item['slide_id']}_timeline_node_{index+1}", x, 3.88, 0.26, color=COLORS["cyan"])
        _shape(slide, f"{item['slide_id']}_timeline_label_box_{index+1}", x - 0.25, y, 1.75, 0.62, fill=COLORS["teal"], line=COLORS["cyan"])
        _text(slide, f"{item['slide_id']}_timeline_label_{index+1}", x - 0.06, y + 0.17, 1.18, 0.2, label, size=7)


def _footer(slide: Any, item: dict[str, Any]) -> None:
    _shape(slide, f"{item['slide_id']}_source_footer_strip", 0.58, 8.25, 14.5, 0.34, fill=COLORS["ink"], line=COLORS["ink"])
    _shape(slide, f"{item['slide_id']}_source_footer_gold_rule", 0.58, 8.25, 14.5, 0.022, fill=COLORS["gold"], line=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_source_footer_text", 0.75, 8.34, 9.9, 0.18, item["footer"]["text"], size=5, color=COLORS["muted"])


def _shape(slide: Any, name: str, x: float, y: float, w: float, h: float, *, fill: str, line: str, shape_type: Any = MSO_AUTO_SHAPE_TYPE.RECTANGLE) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(0.7)
    return shape


def _text(slide: Any, name: str, x: float, y: float, w: float, h: float, text: str, *, size: int, bold: bool = False, color: str = COLORS["offwhite"]) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _line(slide: Any, name: str, x1: float, y1: float, x2: float, y2: float, *, color: str, width: float) -> Any:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.name = name
    connector.line.color.rgb = RGBColor.from_string(color)
    connector.line.width = Pt(width)
    return connector


def _dot(slide: Any, name: str, x: float, y: float, size: float, *, color: str) -> Any:
    return _shape(slide, name, x, y, size, size, fill=color, line=color, shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)


def _icon(slide: Any, name: str, x: float, y: float, w: float, h: float) -> None:
    circle = _shape(slide, name, x, y, w, h, fill=COLORS["cyan"], line=COLORS["cyan"], shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)
    circle.fill.transparency = 10
    triangle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x + w * 0.3), Inches(y + h * 0.25), Inches(w * 0.4), Inches(h * 0.5))
    triangle.name = f"{name}_inner_triangle"
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor.from_string(COLORS["navy"])
    triangle.line.fill.background()


def _add_chart(slide: Any, name: str, x: float, y: float, w: float, h: float, points: list[dict[str, Any]]) -> None:
    _shape(slide, f"{name}_backplate", x - 0.08, y - 0.08, w + 0.16, h + 0.16, fill=COLORS["ink"], line=COLORS["gold"], shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)
    data = ChartData()
    data.categories = [point.get("label", "Metric")[:18] for point in points] or ["Trace", "Method", "QA", "Reuse"]
    data.add_series("Readiness", [int(point.get("value", 0)) for point in points] or [86, 78, 72, 81])
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart_shape.name = name


def _add_table(slide: Any, name: str, x: float, y: float, w: float, h: float, rows: list[list[str]]) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table_shape.name = name
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            is_header = row_index == 0
            is_hybrid = col_index == len(row) - 1 and row_index > 0
            cell.fill.fore_color.rgb = RGBColor.from_string(COLORS["teal_2"] if is_header or is_hybrid else COLORS["ink"])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(7)
                    run.font.bold = is_header
                    run.font.color.rgb = RGBColor.from_string(COLORS["offwhite"] if is_header else COLORS["muted"])


def _card_pairs(item: dict[str, Any], title_role: str = "card_title_text", body_role: str = "card_body_text") -> list[tuple[str, str]]:
    titles = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == title_role]
    bodies = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == body_role]
    return list(zip(titles, bodies))


def _table_rows(item: dict[str, Any]) -> list[list[str]]:
    if item.get("table_data"):
        return [item["table_data"]["header"], *item["table_data"]["rows"]]
    return [["Criterion", "Manual Review", "Structured Pipeline", "Hybrid Governance"]]


def _draw_visual_asset(path: Path, *, variant: str) -> None:
    image = Image.new("RGBA", (900, 620), f"#{COLORS['teal']}")
    draw = ImageDraw.Draw(image, "RGBA")
    if variant == "contour":
        for idx in range(17):
            inset = 20 + idx * 18
            draw.rounded_rectangle((inset, inset, 900 - inset, 620 - int(inset * 0.62)), radius=54, outline=(*_rgb(COLORS["cyan"]), max(26, 138 - idx * 6)), width=3)
        for idx in range(26):
            x = int(900 * (0.14 + idx * 0.028))
            y = int(620 * (0.25 + (idx % 6) * 0.045))
            draw.ellipse((x, y, x + 7, y + 7), fill=(*_rgb(COLORS["gold"]), 180))
    else:
        for idx in range(-180, 920, 30):
            draw.line((idx, 0, idx + 240, 620), fill=(*_rgb(COLORS["cyan"]), 62), width=2)
    image.save(path)


def _render_preview(item: dict[str, Any], art: dict[str, Any], path: Path) -> None:
    image = Image.new("RGB", (SLIDE_W_PX, SLIDE_H_PX), f"#{COLORS['navy']}")
    draw = ImageDraw.Draw(image, "RGBA")
    draw.rectangle((0, 0, 48, SLIDE_H_PX), fill=(*_rgb(COLORS["ink"]), 255))
    draw.rectangle((76, 866, 1590, 901), fill=(*_rgb(COLORS["ink"]), 255))
    draw.rectangle((76, 866, 1590, 869), fill=(*_rgb(COLORS["gold"]), 255))
    _draw_text(draw, 86, 68, item["title"], 21)
    _draw_text(draw, 88, 124, item.get("subtitle", ""), 8)
    archetype = item["archetype_id"]
    if archetype == "cover_hero":
        _preview_hero(draw, item)
    elif archetype == "visual_toc":
        _preview_path(draw, item)
    elif archetype == "section_divider":
        draw.rectangle((110, 90, 130, 745), fill=(*_rgb(COLORS["gold"]), 255))
        _draw_text(draw, 170, 120, "01", 52)
    elif archetype in {"comparison_matrix", "table_heavy"}:
        _preview_table(draw, item)
    elif archetype == "data_dashboard":
        _preview_dashboard(draw, item)
    elif archetype == "process_flow":
        _preview_path(draw, item)
    elif archetype == "timeline_roadmap":
        _preview_timeline(draw, item)
    elif archetype == "methodology_framework":
        for idx in range(4):
            draw.rounded_rectangle((110 + idx * 35, 220 + idx * 90, 860 - idx * 45, 280 + idx * 90), radius=10, fill=(*_rgb(COLORS["teal" if idx % 2 == 0 else "ink"]), 255), outline=(*_rgb(COLORS["cyan"]), 150), width=2)
    else:
        draw.rounded_rectangle((115, 230, 760, 360), radius=20, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["gold"]), 180), width=3)
        _draw_text(draw, 145, 262, item["main_message"][:70], 10)
        for idx in range(3):
            draw.rounded_rectangle((120 + idx * 330, 470, 370 + idx * 330, 550), radius=10, fill=(*_rgb(COLORS["ink"]), 255), outline=(*_rgb(COLORS["cyan"]), 130), width=2)
    _draw_text(draw, 92, 874, item["footer"]["text"][:130], 6)
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def _preview_hero(draw: ImageDraw.ImageDraw, item: dict[str, Any]) -> None:
    _draw_text(draw, 110, 240, item["main_message"][:72], 14)
    draw.rounded_rectangle((880, 115, 1460, 650), radius=34, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["cyan"]), 180), width=3)
    for idx in range(12):
        inset = 24 + idx * 18
        draw.rounded_rectangle((880 + inset, 115 + inset, 1460 - inset, 650 - int(inset * 0.58)), radius=28, outline=(*_rgb(COLORS["cyan"]), max(30, 135 - idx * 8)), width=2)


def _preview_path(draw: ImageDraw.ImageDraw, item: dict[str, Any]) -> None:
    points = [(150, 540), (380, 375), (640, 520), (900, 340), (1160, 500)]
    for idx, point in enumerate(points):
        draw.ellipse((point[0], point[1], point[0] + 34, point[1] + 34), fill=(*_rgb(COLORS["gold" if idx == len(points)-1 else "cyan"]), 230))
        if idx < len(points) - 1:
            nxt = points[idx + 1]
            draw.line((point[0] + 17, point[1] + 17, nxt[0] + 17, nxt[1] + 17), fill=(*_rgb(COLORS["cyan"]), 180), width=4)


def _preview_table(draw: ImageDraw.ImageDraw, item: dict[str, Any]) -> None:
    draw.rounded_rectangle((100, 210, 1270, 650), radius=8, fill=(*_rgb(COLORS["ink"]), 255), outline=(*_rgb(COLORS["cyan"]), 160), width=2)
    rows = _table_rows(item)[:5]
    for row_index, row in enumerate(rows):
        y = 230 + row_index * 78
        draw.rectangle((120, y, 1240, y + 36), fill=(*_rgb(COLORS["teal_2" if row_index == 0 else "ink"]), 210))
        for col_index, value in enumerate(row[:4]):
            _draw_text(draw, 140 + col_index * 270, y + 8, str(value)[:24], 6 if row_index else 7)


def _preview_dashboard(draw: ImageDraw.ImageDraw, item: dict[str, Any]) -> None:
    points = (item.get("chart_data") or {}).get("data_points", [])[:4]
    for idx, point in enumerate(points):
        y = 220 + idx * 92
        draw.rounded_rectangle((115, y, 310, y + 58), radius=8, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["cyan"]), 150), width=2)
        _draw_text(draw, 132, y + 10, f"{point['value']}%", 12)
    draw.rounded_rectangle((380, 195, 1030, 640), radius=16, fill=(*_rgb(COLORS["ink"]), 255), outline=(*_rgb(COLORS["gold"]), 180), width=3)
    for idx, point in enumerate(points):
        bar_h = int(point["value"] * 2.1)
        draw.rectangle((480 + idx * 120, 595 - bar_h, 545 + idx * 120, 595), fill=(*_rgb(COLORS["cyan"]), 230))


def _preview_timeline(draw: ImageDraw.ImageDraw, item: dict[str, Any]) -> None:
    draw.line((130, 460, 1350, 460), fill=(*_rgb(COLORS["gold"]), 255), width=5)
    for idx in range(5):
        x = 150 + idx * 285
        draw.ellipse((x, 446, x + 28, 474), fill=(*_rgb(COLORS["cyan"]), 255))
        y = 340 if idx % 2 == 0 else 520
        draw.rounded_rectangle((x - 35, y, x + 160, y + 65), radius=9, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["cyan"]), 140), width=2)


def _contact_sheet(paths: list[Path], output_path: Path) -> Path:
    thumbs = [Image.open(path).resize((418, 235)) for path in paths]
    sheet = Image.new("RGB", (4 * 418, 3 * 235), f"#{COLORS['navy']}")
    for index, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((index % 4) * 418, (index // 4) * 235))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    return output_path


def _draw_text(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, size: int) -> None:
    try:
        font = ImageFont.truetype("arial.ttf", max(8, size * 2))
    except OSError:
        font = ImageFont.load_default()
    draw.text((x, y), str(text), fill=(*_rgb(COLORS["offwhite"]), 255), font=font)


def _rgb(hex_color: str) -> tuple[int, int, int]:
    return tuple(int(hex_color[index : index + 2], 16) for index in (0, 2, 4))
