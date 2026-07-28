"""Compile the E01.1 semantic component candidate PPTX."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt

SLIDE_W = 16.0
SLIDE_H = 9.0


def build_editable_candidate_spec_e01_1(component_graph: dict[str, Any], text_lift: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "editable_candidate_spec_e01_1",
        "slide_size": {"width_in": SLIDE_W, "height_in": SLIDE_H},
        "component_count": component_graph["component_count"],
        "editable_text_region_count": text_lift["editable_text_region_count"],
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "allowed_raster_targets": ["hero_photo_field", "thumbnail_callout_image_frames"],
        "canva_parity_claimed": False,
    }


def compile_e01_1_candidate(
    *,
    reference_image: Path,
    component_graph: dict[str, Any],
    oracle: dict[str, Any],
    output_pptx: Path,
) -> dict[str, Any]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    assets_dir = output_pptx.parent / "assets"
    assets_dir.mkdir(exist_ok=True)
    crops = _create_crops(reference_image, assets_dir)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_background(slide)
    slide.shapes.add_picture(str(crops["hero"]), Inches(0), Inches(0), Inches(9.75), Inches(7.0)).name = "hero_photo_field_bounded_reference_crop"
    _add_technical_overlay(slide)
    _add_thumbnail_callouts(slide, crops, oracle["thumbnail_callouts"])
    _add_checklist(slide, oracle)
    _add_bottom_action_bar(slide, oracle)
    _add_source_footer(slide)

    prs.save(output_pptx)
    return {
        "schema_name": "e01_1_candidate_compile_report",
        "status": "passed" if output_pptx.exists() else "failed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 1,
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "bounded_raster_layer_count": 4,
        "canva_parity_claimed": False,
    }


def _create_crops(reference_image: Path, assets_dir: Path) -> dict[str, Path]:
    with Image.open(reference_image) as image:
        image = image.convert("RGB")
        w, h = image.size
        crops = {
            "hero": image.crop((0, 0, int(w * 0.61), int(h * 0.735))),
            "thumb_1": image.crop((int(w * 0.225), int(h * 0.585), int(w * 0.335), int(h * 0.72))),
            "thumb_2": image.crop((int(w * 0.352), int(h * 0.585), int(w * 0.462), int(h * 0.72))),
            "thumb_3": image.crop((int(w * 0.48), int(h * 0.585), int(w * 0.59), int(h * 0.72))),
        }
    paths: dict[str, Path] = {}
    for key, crop in crops.items():
        path = assets_dir / f"{key}.png"
        crop.save(path)
        paths[key] = path
    return paths


def _add_background(slide: Any) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    shape.name = "background_base_native_shape"
    _fill(shape, "061526")
    shape.line.fill.background()


def _add_checklist(slide: Any, oracle: dict[str, Any]) -> None:
    panel_x, panel_y, panel_w, panel_h = 9.72, 0.35, 5.95, 7.05
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(panel_x), Inches(panel_y), Inches(panel_w), Inches(panel_h))
    panel.name = "checklist_system_panel_native"
    _fill(panel, "071A27", transparency=8)
    _line(panel, "24C5D8", 1.4)

    _add_text(slide, "checklist_title_text", oracle["title"], panel_x + 0.9, panel_y + 0.28, 4.2, 0.42, 22, "55D7E5", bold=True, align_center=True)
    _add_line(slide, panel_x + 0.2, panel_y + 0.72, panel_x + panel_w - 0.2, panel_y + 0.72, "147B8C", 0.9)

    row_x = panel_x + 0.18
    row_w = panel_w - 0.36
    row_h = 1.15
    row_gap = 0.105
    row_y0 = panel_y + 0.83
    for idx, step in enumerate(oracle["steps"]):
        y = row_y0 + idx * (row_h + row_gap)
        row = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(row_x), Inches(y), Inches(row_w), Inches(row_h))
        row.name = f"step_{step['index']}_card_panel_native"
        _fill(row, "0B2A34", transparency=4)
        _line(row, "176F80", 0.8)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(row_x + 0.14), Inches(y + 0.16), Inches(0.78), Inches(0.78))
        circle.name = f"step_{step['index']}_icon_circle_vector"
        _fill(circle, "092535", transparency=35)
        _line(circle, "37D7E8", 1.4)
        _add_icon_mark(slide, step["icon_role"], row_x + 0.36, y + 0.37, "6DE7F3")
        _add_text(slide, f"step_{step['index']}_number_text", step["number"], row_x + 1.48, y + 0.25, 0.65, 0.55, 28, "36C7D9", bold=True)
        _add_line(slide, row_x + 2.18, y + 0.25, row_x + 2.18, y + 0.9, "1D8E9D", 1.0)
        _add_text(slide, f"step_{step['index']}_heading_text", step["heading"], row_x + 2.38, y + 0.2, 2.65, 0.28, 14, "F8FAFC", bold=True)
        _add_text(slide, f"step_{step['index']}_body_text", step["body"], row_x + 2.38, y + 0.52, 2.55, 0.42, 11, "F4F7FA")
        chevron = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, Inches(row_x + row_w - 0.38), Inches(y + 0.43), Inches(0.22), Inches(0.28))
        chevron.name = f"step_{step['index']}_chevron_vector"
        _fill(chevron, "37D7E8")
        chevron.line.fill.background()


def _add_thumbnail_callouts(slide: Any, crops: dict[str, Path], callouts: list[dict[str, Any]]) -> None:
    positions = [(3.55, 5.63), (5.5, 5.63), (7.45, 5.63)]
    for idx, callout in enumerate(callouts, start=1):
        x, y = positions[idx - 1]
        pic = slide.shapes.add_picture(str(crops[f"thumb_{idx}"]), Inches(x), Inches(y), Inches(1.65), Inches(1.2))
        pic.name = f"thumbnail_{idx}_bounded_image_frame"
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.03), Inches(y - 0.08), Inches(1.72), Inches(1.35))
        frame.name = f"thumbnail_{idx}_circular_frame_vector"
        frame.fill.background()
        _line(frame, "38D9EA", 1.5)
        _add_text(slide, f"thumbnail_{idx}_caption_text", callout["label"], x + 0.03, y + 1.33, 1.7, 0.24, 9, "F8FAFC", bold=True, align_center=True)


def _add_bottom_action_bar(slide: Any, oracle: dict[str, Any]) -> None:
    bar_y = 7.58
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(bar_y), Inches(SLIDE_W), Inches(1.26))
    bar.name = "bottom_action_bar_container_native"
    _fill(bar, "06111A", transparency=2)
    bar.line.fill.background()
    _add_line(slide, 0.0, bar_y, SLIDE_W, bar_y, "F59E1B", 1.1)
    _add_line(slide, 0.55, 8.75, 14.4, 8.75, "B87518", 0.9)

    x0 = 0.75
    cell_w = 2.82
    for idx, action in enumerate(oracle["actions"]):
        x = x0 + idx * cell_w
        if idx > 0:
            _add_line(slide, x - 0.22, 7.82, x - 0.22, 8.55, "1C3E47", 0.8)
        _add_action_icon(slide, action["icon_role"], x, 7.85)
        _add_text(slide, f"action_{action['index']}_label_top_text", action["label_top"], x + 0.78, 7.85, 1.72, 0.28, 11, "F5A623", bold=True)
        _add_text(slide, f"action_{action['index']}_label_bottom_text", action["label_bottom"], x + 0.78, 8.16, 1.95, 0.28, 10, "F5A623", bold=True)


def _add_source_footer(slide: Any) -> None:
    _add_line(slide, 0.4, 8.88, 15.2, 8.88, "174B58", 0.7)
    _add_text(slide, "source_footer_text_editable", "SOURCE / FOOTER SLOT", 0.42, 8.82, 2.4, 0.16, 5, "8BC8D2")


def _add_technical_overlay(slide: Any) -> None:
    for offset, size in ((0, 1.65), (0.18, 1.28), (0.36, 0.92), (0.54, 0.56)):
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.25 + offset), Inches(0.55 + offset), Inches(size), Inches(size))
        circle.name = f"technical_overlay_radar_ring_{size}"
        circle.fill.background()
        _line(circle, "0DAEC3", 0.7)
    _add_line(slide, 0.35, 0.55, 4.65, 0.55, "0DAEC3", 0.5)
    _add_line(slide, 2.9, 6.25, 8.9, 6.25, "0DAEC3", 0.6)
    for x, y in ((4.65, 1.35), (6.08, 0.78), (8.58, 6.35), (3.3, 6.25)):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.06), Inches(0.06))
        dot.name = "technical_overlay_dot_vector"
        _fill(dot, "23CAE0")
        dot.line.fill.background()


def _add_icon_mark(slide: Any, role: str, x: float, y: float, color: str) -> None:
    _add_line(slide, x, y + 0.22, x + 0.36, y + 0.22, color, 1.4)
    _add_line(slide, x + 0.18, y, x + 0.18, y + 0.42, color, 1.4)
    if "shield" in role:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(x - 0.06), Inches(y - 0.04), Inches(0.5), Inches(0.48))
    elif "document" in role or "clipboard" in role:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x - 0.03), Inches(y - 0.02), Inches(0.43), Inches(0.46))
    else:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.04), Inches(y), Inches(0.45), Inches(0.45))
    shape.name = f"{role}_vector_icon_mark"
    shape.fill.background()
    _line(shape, color, 1.1)


def _add_action_icon(slide: Any, role: str, x: float, y: float) -> None:
    color = "F5A623"
    if "warning" in role:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y + 0.02), Inches(0.54), Inches(0.48))
    elif "lock" in role:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(y + 0.12), Inches(0.45), Inches(0.36))
    elif "shield" in role:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(x), Inches(y), Inches(0.58), Inches(0.58))
    elif "chat" in role:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + 0.08), Inches(0.62), Inches(0.44))
    else:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y + 0.05), Inches(0.58), Inches(0.5))
    icon.name = f"{role}_bottom_action_vector_icon"
    icon.fill.background()
    _line(icon, color, 2.0)


def _add_text(
    slide: Any,
    name: str,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: str,
    *,
    bold: bool = False,
    align_center: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.01)
    frame.margin_right = Inches(0.01)
    frame.margin_top = Inches(0.0)
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align_center:
        paragraph.alignment = 2
    run = paragraph.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _add_line(slide: Any, x1: float, y1: float, x2: float, y2: float, color: str, width_pt: float) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = "native_connector_line"
    line.line.color.rgb = RGBColor.from_string(color)
    line.line.width = Pt(width_pt)


def _fill(shape: Any, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.fill.transparency = transparency


def _line(shape: Any, color: str, width_pt: float) -> None:
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(width_pt)

