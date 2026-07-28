"""Shape object creation for E06.2 contract compilation."""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt

from src.presentation_agent.magic_layer.e06_2_contract_object_factory import contract_shape_name


def add_contract_shape(slide: Any, obj: dict[str, Any]) -> Any:
    bbox = obj["bbox_emu"]
    if obj.get("object_type") == "line" or str(obj.get("shape_type", "")).startswith("LINE"):
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            int(bbox["x"]),
            int(bbox["y"]),
            int(bbox["x"] + bbox["w"]),
            int(bbox["y"] + bbox["h"]),
        )
        shape.name = contract_shape_name(obj)
        shape.line.color.rgb = RGBColor(64, 82, 101)
        shape.line.width = Pt(0.6)
        return shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(bbox["x"]), int(bbox["y"]), max(1, int(bbox["w"])), max(1, int(bbox["h"])))
    shape.name = contract_shape_name(obj)
    fill, line = _colors_for(obj)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)
    shape.line.width = Pt(0.35)
    return shape


def add_contract_placeholder(slide: Any, obj: dict[str, Any]) -> Any:
    shape = add_contract_shape(slide, obj)
    shape.fill.background()
    shape.line.fill.background()
    return shape


def _colors_for(obj: dict[str, Any]) -> tuple[tuple[int, int, int], tuple[int, int, int]]:
    object_type = obj.get("object_type")
    if object_type == "icon_background":
        return (255, 255, 255), (203, 213, 225)
    if object_type == "source_footer":
        return (15, 23, 42), (242, 169, 0)
    if object_type == "card_region":
        return (13, 27, 40), (48, 72, 94)
    if object_type == "table_region":
        return (18, 28, 42), (90, 105, 128)
    if object_type == "chart_region":
        return (14, 32, 46), (40, 215, 232)
    if object_type == "image_field":
        return (20, 36, 50), (56, 217, 158)
    return (7, 16, 24), (32, 48, 64)
