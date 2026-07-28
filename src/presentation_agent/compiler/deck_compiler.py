"""Compile slide_blueprint content through editable template layouts into a final PPTX."""

from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.image import ImagePart
from pptx.util import Inches, Pt

from ..generator_contracts import validateDeckAssemblyPlan, validateEditableTemplateSpec
from .blueprint_adapter import load_valid_slide_blueprints
from .content_binding_policy import (
    DEFAULT_TEMPLATE_CONTRACTS_DIR,
    apply_content_binding_policy,
    load_template_contracts,
)
from .decorative_budget import resolve_decorative_budget, update_budget_decisions_after_render
from .deterministic_fallbacks import (
    choose_title_body_slots,
    normalize_card_blocks,
    normalize_chart_data,
    normalize_table_data,
)
from .icon_resolver import IconResolver, IconResolution, SVG_CONTENT_TYPE
from .layout_matcher import build_deck_assembly_plan
from .template_spec_selector import (
    DEFAULT_FINAL_TEMPLATE_SPEC_PATH as SELECTOR_FINAL_TEMPLATE_SPEC_PATH,
    load_explicit_template_spec,
    select_template_spec,
)


DEFAULT_BLUEPRINT_PATH = Path("outputs/slide_blueprint.json")
DEFAULT_TEMPLATE_SPEC_PATH = Path("outputs/editable_template_spec.json")
DEFAULT_FINAL_TEMPLATE_SPEC_PATH = SELECTOR_FINAL_TEMPLATE_SPEC_PATH
DEFAULT_ASSEMBLY_PLAN_PATH = Path("outputs/deck_assembly_plan.json")
DEFAULT_OUTPUT_PATH = Path("outputs/final_deck.pptx")
DEFAULT_MANIFEST_PATH = Path("outputs/final_deck_manifest.json")
DEFAULT_PHOTO_PLACEHOLDER_MANIFEST = Path("assets/photo_placeholders/photo_placeholder_manifest.json")
PRIORITY_MASTER_ARCHETYPES = {
    "creative_cover",
    "visual_table_of_contents",
    "section_divider",
    "research_overview",
    "methodology_framework",
    "data_table_appendix",
}


def compile_final_deck(
    slide_blueprints: dict[str, Any] | list[dict[str, Any]],
    editable_template_spec: dict[str, Any],
    deck_assembly_plan: dict[str, Any] | None,
    output_path: str | Path,
    *,
    template_contracts_dir: str | Path = DEFAULT_TEMPLATE_CONTRACTS_DIR,
) -> dict[str, Any]:
    validateEditableTemplateSpec(editable_template_spec)
    slides = _normalize_slides(slide_blueprints)
    assembly_plan = deck_assembly_plan or build_deck_assembly_plan(slides, editable_template_spec)
    validateDeckAssemblyPlan(assembly_plan)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    canvas = editable_template_spec["canvas"]
    presentation.slide_width = Inches(float(canvas["width"]))
    presentation.slide_height = Inches(float(canvas["height"]))
    blank_layout = presentation.slide_layouts[6]
    style = _DeckStyle(editable_template_spec)
    layouts = {layout["layout_id"]: layout for layout in editable_template_spec["layouts"]}
    slide_by_id = {str(slide.get("slide_id") or slide.get("id") or f"slide-{index:03d}"): slide for index, slide in enumerate(slides, start=1)}
    template_contracts = load_template_contracts(template_contracts_dir)
    icon_resolver = IconResolver()

    compiled: list[dict[str, Any]] = []
    icon_reports: list[dict[str, Any]] = []
    for binding in assembly_plan["slide_layout_bindings"]:
        slide_id = binding["slide_id"]
        source_slide = slide_by_id.get(slide_id)
        layout = layouts.get(binding["selected_layout_id"])
        if source_slide is None or layout is None:
            continue
        ppt_slide = presentation.slides.add_slide(blank_layout)
        render_warnings: list[dict[str, Any]] = []
        contract = template_contracts.get(str(layout.get("archetype_id") or ""))
        render_source, binding_policy = apply_content_binding_policy(source_slide, layout, binding, contract)
        render_warnings.extend(binding_policy.get("overflow_warnings") or [])
        render_binding = dict(binding)
        render_binding["content_binding_policy"] = binding_policy
        render_binding["text_first_contract"] = True
        render_binding["presentation_role"] = (contract or {}).get("presentation_role")
        render_binding["reading_order"] = (contract or {}).get("reading_order") or []
        render_binding["content_capacity"] = (contract or {}).get("content_capacity") or {}
        budget_plan = resolve_decorative_budget(
            layout,
            contract,
            binding=render_binding,
            slot_capacity_status=binding_policy.get("slot_capacity_status") if isinstance(binding_policy, dict) else {},
        )
        render_binding["decorative_budget_plan"] = budget_plan
        render_binding["decorative_budget"] = {
            "max_shape_count_target": budget_plan["max_shape_count_target"],
            "max_ornament_density": budget_plan["allowed_ornament_density"],
            "max_background_coverage": budget_plan["max_background_coverage"],
        }
        slide_style = style.for_tone(str(binding.get("selected_tone_variant") or assembly_plan.get("selected_tone_variant") or "creative"))
        slide_style = slide_style.for_layout(layout, render_binding)
        icon_report = icon_resolver.empty_report()
        _render_slide(
            ppt_slide,
            render_source,
            layout,
            render_binding,
            slide_style,
            render_warnings,
            icon_resolver=icon_resolver,
            icon_report=icon_report,
        )
        icon_reports.append(icon_report)
        _add_notes(ppt_slide, render_source)
        semantic_notes_slots = _preserve_semantic_source_in_notes(ppt_slide, source_slide)
        compiled.append(
            {
                "slide_id": slide_id,
                "layout_id": layout["layout_id"],
                "archetype_id": layout["archetype_id"],
                "pptx_index": len(compiled),
                "contract_used": binding_policy.get("contract_used"),
                "content_trimmed": binding_policy.get("content_trimmed"),
                "content_split": binding_policy.get("content_split"),
                "content_moved_to_notes": binding_policy.get("content_moved_to_notes"),
                "semantic_source_preserved_in_notes": bool(semantic_notes_slots),
                "semantic_source_notes_slots": semantic_notes_slots,
                "appendix_created": binding_policy.get("appendix_created"),
                "overflow_warnings": binding_policy.get("overflow_warnings"),
                "source_anchor_preserved": binding_policy.get("source_anchor_preserved"),
                "slot_capacity_status": binding_policy.get("slot_capacity_status"),
                "overflow_actions": binding_policy.get("overflow_actions"),
                "selected_tone_variant": slide_style.tone_variant,
                "ornament_density_mode": slide_style.ornament_density,
                "tone_tokens": slide_style.report_tokens(),
                "decorative_budget_plan": update_budget_decisions_after_render(
                    budget_plan,
                    shape_count=len(ppt_slide.shapes),
                ),
                "icons_used": icon_report["icons_used"],
                "missing_icons": icon_report["missing_icons"],
                "unresolved_icon_roles": icon_report["unresolved_icon_roles"],
                "icon_asset_paths": icon_report["icon_asset_paths"],
                "icon_family": icon_report["icon_family"],
                "warnings": render_warnings,
            }
        )

    presentation.save(output)
    return {
        "schema_name": "final_deck_manifest",
        "schema_version": "1.0",
        "pptx_path": _display_path(output),
        "slide_count": len(compiled),
        "source_blueprint_path": _display_path(DEFAULT_BLUEPRINT_PATH),
        "source_template_spec_path": str((assembly_plan.get("template_spec_source") or {}).get("path") or _display_path(DEFAULT_TEMPLATE_SPEC_PATH)),
        "source_assembly_plan_path": _display_path(DEFAULT_ASSEMBLY_PLAN_PATH),
        "template_contracts_dir": _display_path(Path(template_contracts_dir)),
        "template_contracts_loaded": len(template_contracts),
        "photo_placeholder_manifest_path": _display_path(DEFAULT_PHOTO_PLACEHOLDER_MANIFEST),
        "icon_report": icon_resolver.merge_reports(icon_reports),
        "compiled_slides": compiled,
        "editable_policy": {
            "text_boxes_are_editable": True,
            "tables_are_editable": True,
            "charts_are_editable_or_shape_group_mvp": True,
            "photo_placeholders_only_in_declared_image_slots": True,
            "no_gpt_image_reference_images_inserted": True,
            "no_full_slide_raster_images": True,
        },
    }


def compile_final_deck_from_files(
    *,
    slide_blueprint_path: str | Path = DEFAULT_BLUEPRINT_PATH,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC_PATH,
    final_template_spec_path: str | Path = DEFAULT_FINAL_TEMPLATE_SPEC_PATH,
    prefer_final_template_spec: bool = True,
    assembly_plan_path: str | Path = DEFAULT_ASSEMBLY_PLAN_PATH,
    output_path: str | Path = DEFAULT_OUTPUT_PATH,
    manifest_path: str | Path = DEFAULT_MANIFEST_PATH,
    template_contracts_dir: str | Path = DEFAULT_TEMPLATE_CONTRACTS_DIR,
) -> Path:
    slide_blueprints, selected_slide_blueprint_path = load_valid_slide_blueprints(slide_blueprint_path)
    if Path(template_spec_path) == DEFAULT_TEMPLATE_SPEC_PATH:
        selection = select_template_spec(
            base_template_spec_path=template_spec_path,
            final_template_spec_path=final_template_spec_path,
            prefer_final=prefer_final_template_spec,
        )
    else:
        selection = load_explicit_template_spec(template_spec_path)
    template_spec = selection.spec
    assembly_plan = _load_json(assembly_plan_path) if Path(assembly_plan_path).exists() else None
    if assembly_plan is None or not _assembly_plan_matches_template_source(assembly_plan, selection.source):
        assembly_plan = build_deck_assembly_plan(slide_blueprints, template_spec, template_spec_source=selection.source)
        Path(assembly_plan_path).write_text(json.dumps(assembly_plan, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    manifest = compile_final_deck(
        slide_blueprints,
        template_spec,
        assembly_plan,
        output_path,
        template_contracts_dir=template_contracts_dir,
    )
    manifest["source_blueprint_path"] = _display_path(selected_slide_blueprint_path)
    manifest_file = Path(manifest_path)
    manifest_file.parent.mkdir(parents=True, exist_ok=True)
    manifest_file.write_text(json.dumps(manifest, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    return Path(output_path)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Compile outputs/final_deck.pptx from blueprints and editable template layouts.")
    parser.add_argument("--slide-blueprint", type=Path, default=DEFAULT_BLUEPRINT_PATH)
    parser.add_argument("--template-spec", type=Path, default=DEFAULT_TEMPLATE_SPEC_PATH)
    parser.add_argument("--final-template-spec", type=Path, default=DEFAULT_FINAL_TEMPLATE_SPEC_PATH)
    parser.add_argument("--use-base-template-spec", action="store_true")
    parser.add_argument("--assembly-plan", type=Path, default=DEFAULT_ASSEMBLY_PLAN_PATH)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT_PATH)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST_PATH)
    parser.add_argument("--template-contracts", type=Path, default=DEFAULT_TEMPLATE_CONTRACTS_DIR)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        output = compile_final_deck_from_files(
            slide_blueprint_path=args.slide_blueprint,
            template_spec_path=args.template_spec,
            final_template_spec_path=args.final_template_spec,
            prefer_final_template_spec=not args.use_base_template_spec,
            assembly_plan_path=args.assembly_plan,
            output_path=args.output,
            manifest_path=args.manifest,
            template_contracts_dir=args.template_contracts,
        )
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"BUILD_DECK_FAILED {exc}")
        return 1
    print(f"WROTE {output}")
    return 0


def _ornament_density_tokens(raw: Any) -> dict[str, dict[str, Any]]:
    defaults: dict[str, dict[str, Any]] = {
        "low": {
            "mode": "low",
            "grid_columns": 4,
            "guide_line_count": 4,
            "connector_cluster_count": 1,
            "micro_node_count": 8,
            "blueprint_arc_count": 1,
            "max_extra_shapes": 18,
            "line_weight_pt": 0.12,
        },
        "medium": {
            "mode": "medium",
            "grid_columns": 8,
            "guide_line_count": 10,
            "connector_cluster_count": 2,
            "micro_node_count": 18,
            "blueprint_arc_count": 3,
            "max_extra_shapes": 42,
            "line_weight_pt": 0.18,
        },
        "high": {
            "mode": "high",
            "grid_columns": 12,
            "guide_line_count": 18,
            "connector_cluster_count": 3,
            "micro_node_count": 30,
            "blueprint_arc_count": 5,
            "max_extra_shapes": 72,
            "line_weight_pt": 0.24,
        },
    }
    if isinstance(raw, dict):
        for mode in ("low", "medium", "high"):
            payload = raw.get(mode)
            if isinstance(payload, dict):
                merged = dict(defaults[mode])
                merged.update(payload)
                merged["mode"] = mode
                defaults[mode] = merged
    return defaults


def _normalize_ornament_density(value: Any, default: str = "medium") -> str:
    normalized = str(value or default).strip().lower().replace("_", "-")
    if normalized in {"medium-low", "medium_low"}:
        return "medium"
    if normalized in {"low", "medium", "high"}:
        return normalized
    return default


def _cap_ornament_density(current: str, maximum: str) -> str:
    order = {"low": 0, "medium": 1, "high": 2}
    current_mode = _normalize_ornament_density(current, "medium")
    max_mode = _normalize_ornament_density(maximum, "medium")
    return current_mode if order[current_mode] <= order[max_mode] else max_mode


def _effective_ornament_density(current: str, layout: dict[str, Any], binding: dict[str, Any]) -> str:
    requested = _normalize_ornament_density(binding.get("ornament_density"), current)
    family_id = str(layout.get("layout_family_id") or binding.get("layout_family_id") or "")
    layout_id = str(layout.get("layout_id") or binding.get("selected_layout_id") or "")
    slide_type = str(binding.get("slide_type") or layout.get("archetype_id") or "")
    tone = str(binding.get("selected_tone_variant") or "").lower()
    component_density = str(binding.get("component_density") or "").lower()

    high_allowed = (
        family_id in {"expressive_cover_divider", "visual_toc_navigation"}
        or slide_type in {"cover", "creative_cover", "section_divider", "visual_toc"}
        or layout_id.endswith("creative_cover")
        or (tone == "creative" and component_density not in {"data_dense", "evidence_dense"})
    )
    body_should_stay_medium = component_density in {"data_dense", "evidence_dense"} and family_id not in {
        "expressive_cover_divider",
        "visual_toc_navigation",
    }
    if body_should_stay_medium:
        return "medium" if requested == "high" else requested
    if high_allowed:
        return "high"
    return requested


class _DeckStyle:
    def __init__(self, spec: dict[str, Any], tone_variant: str = "base") -> None:
        self.canvas_width = float(spec["canvas"]["width"])
        self.canvas_height = float(spec["canvas"]["height"])
        self.tone_variant = tone_variant
        self.colors = dict(spec["tokens"]["colors"])
        self.typography = dict(spec["tokens"]["typography"])
        self.spacing = spec["tokens"].get("spacing", {})
        self.ornament_density_tokens = _ornament_density_tokens(spec["tokens"].get("ornament_density"))
        self.tone_variants = ((spec.get("tokens") or {}).get("typography") or {}).get("tone_variants") or {}
        self.footer_style = "standard"
        self.card_style = "standard"
        self.chart_table_style = "standard"
        self.section_style = "standard"
        self.ornament_density = "medium"
        self.image_frame_style = "diagonal"
        self.dark_hero = False
        self.index_rail = False
        self.expressive_cover_identity = False
        self.fractured_geometry = False
        self.archetype_id = ""
        self.layout_family_id = ""
        self.priority_master_identity = False
        self.layout_reference_source = "unspecified"
        self.identity_lock = "standard"
        self.reference_geometry: dict[str, Any] = {}
        self.minimum_visual_features: dict[str, Any] = {}
        self.component_grammar: dict[str, Any] = {}
        self.refinement_patch_directives: list[dict[str, Any]] = []
        self.deck_scale = ""
        self.text_first_contract = False
        self.decorative_budget: dict[str, Any] = {}
        self.decorative_budget_plan: dict[str, Any] = {}
        self.protected_text_zones: list[dict[str, Any]] = []
        self.protected_table_chart_zones: list[dict[str, Any]] = []
        self.card_chrome_simplified = False
        self.footer_chrome_simplified = False
        self.background_density_reduced = False
        self._apply_tone_variant(tone_variant)

    def for_tone(self, tone_variant: str) -> "_DeckStyle":
        normalized = str(tone_variant or "creative").strip().lower()
        if normalized not in {"academic", "professional", "creative"}:
            normalized = "creative"
        return _DeckStyle(
            {
                "canvas": {"width": self.canvas_width, "height": self.canvas_height},
                "tokens": {
                    "colors": self.colors,
                    "typography": {**self.typography, "tone_variants": self.tone_variants},
                    "spacing": self.spacing,
                    "ornament_density": self.ornament_density_tokens,
                },
            },
            normalized,
        )

    def color(self, token: str, fallback: str = "#111827") -> RGBColor:
        return _hex_to_rgb(self.colors.get(token, fallback))

    def font_size(self, token: str, fallback: float) -> Pt:
        return Pt(float(self.typography.get(token, {}).get("size_pt", fallback)))

    def font_family(self, token: str) -> str:
        return str(self.typography.get(token, {}).get("font_family", "Aptos"))

    def for_layout(self, layout: dict[str, Any], binding: dict[str, Any]) -> "_DeckStyle":
        layout_style = self.for_tone(str(binding.get("selected_tone_variant") or self.tone_variant or "creative"))
        layout_id = str(layout.get("layout_id") or "")
        family_id = str(layout.get("layout_family_id") or binding.get("layout_family_id") or "")
        archetype_id = str(layout.get("archetype_id") or binding.get("slide_type") or "")
        rhythm = str(binding.get("section_rhythm_role") or "")
        component_bindings = layout.get("component_bindings") or {}
        layout_style.archetype_id = archetype_id
        layout_style.layout_family_id = family_id
        layout_style.priority_master_identity = archetype_id in PRIORITY_MASTER_ARCHETYPES
        layout_style.layout_reference_source = str(binding.get("layout_reference_source") or "template_spec")
        layout_style.identity_lock = str(binding.get("identity_lock") or "standard")
        layout_style.reference_geometry = dict(binding.get("reference_geometry") or {})
        layout_style.minimum_visual_features = dict(binding.get("minimum_visual_features") or {})
        layout_style.component_grammar = dict(binding.get("component_grammar") or {})
        layout_style.refinement_patch_directives = [
            dict(item)
            for item in (binding.get("refinement_patch_directives") or [])
            if isinstance(item, dict)
        ]
        layout_style.deck_scale = str(binding.get("deck_scale") or "")
        layout_style.text_first_contract = bool(binding.get("text_first_contract"))
        layout_style.decorative_budget = dict(binding.get("decorative_budget") or {})
        layout_style.decorative_budget_plan = dict(binding.get("decorative_budget_plan") or {})
        layout_style.protected_text_zones = list(layout_style.decorative_budget_plan.get("protected_text_zones") or [])
        layout_style.protected_table_chart_zones = list(layout_style.decorative_budget_plan.get("protected_table_chart_zones") or [])
        budget_decisions = layout_style.decorative_budget_plan.get("decisions") if isinstance(layout_style.decorative_budget_plan.get("decisions"), dict) else {}
        layout_style.card_chrome_simplified = bool(budget_decisions.get("card_chrome_simplified"))
        layout_style.footer_chrome_simplified = bool(budget_decisions.get("footer_chrome_simplified"))
        layout_style.background_density_reduced = bool(layout_style.text_first_contract or budget_decisions.get("background_density_reduced"))
        if layout_style.identity_lock == "high":
            layout_style.footer_style = "citation_dense"
            if layout_style.minimum_visual_features.get("index_presence"):
                layout_style.index_rail = True
            if layout_style.minimum_visual_features.get("image_mask_presence"):
                layout_style.image_frame_style = "medallion_topology_frame"
            if layout_style.layout_reference_source == "layout_ref":
                layout_style.ornament_density = "high" if layout_style.minimum_visual_features.get("ornament_presence") else layout_style.ornament_density
        if family_id == "expressive_cover_divider" or "expressive_cover_identity" in component_bindings.values():
            layout_style.expressive_cover_identity = True
            layout_style.index_rail = True
            layout_style.fractured_geometry = True
            layout_style.footer_style = "editorial_dense"
            layout_style.card_style = "layered_editorial_card"
            layout_style.ornament_density = "high"
        if layout_id.endswith("creative_cover") and rhythm != "divider":
            layout_style.dark_hero = True
            layout_style.ornament_density = "high"
            layout_style.footer_style = "editorial_dense"
            layout_style.image_frame_style = "medallion_topology_frame"
            layout_style.colors.update(
                {
                    "background": "#071A33",
                    "surface": "#0B2446",
                    "surface_alt": "#102E54",
                    "text": "#F7F0DC",
                    "muted_text": "#D8CFAE",
                    "line": "#D0AF37",
                    "grid": "#24405E",
                    "accent": self.colors.get("gold_accent", "#D0AF37"),
                    "accent_secondary": self.colors.get("teal", "#0FA3A3"),
                }
            )
            layout_style._scale_type({"title": -10, "subtitle": -1, "body": -1})
        if "section_divider" in layout_id or rhythm == "divider":
            layout_style.section_style = "oversized_creative_marker"
            layout_style.index_rail = True
            layout_style.fractured_geometry = True
            layout_style.colors.update(
                {
                    "background": "#FFF9EA",
                    "surface": "#FFFFFF",
                    "surface_alt": "#F8EED2",
                    "accent": self.colors.get("gold_accent", "#D0AF37"),
                    "accent_secondary": self.colors.get("teal", "#0FA3A3"),
                    "line": "#D6B666",
                    "grid": "#EADAB6",
                    "text": self.colors.get("primary_navy", "#071A33"),
                }
            )
        if family_id in {"visual_toc_navigation", "expressive_cover_divider"}:
            layout_style.index_rail = True
        if family_id == "evidence_overview":
            layout_style.footer_style = "citation_dense"
            layout_style.card_style = "evidence_card"
        if family_id in {"kpi_dashboard", "table_appendix"}:
            layout_style.chart_table_style = "expressive_modular"
            layout_style.footer_style = "citation_dense"
        if archetype_id == "visual_table_of_contents":
            layout_style.index_rail = True
            layout_style.card_style = "layered_editorial_card"
            layout_style.ornament_density = "high"
        if archetype_id == "research_overview":
            layout_style.footer_style = "citation_dense"
            layout_style.card_style = "evidence_card"
            layout_style.ornament_density = "medium"
        if archetype_id == "methodology_framework":
            layout_style.footer_style = "citation_dense"
            layout_style.card_style = "evidence_card"
            layout_style.ornament_density = "medium"
        if archetype_id == "data_table_appendix":
            layout_style.footer_style = "citation_dense"
            layout_style.chart_table_style = "dense_academic"
            layout_style.ornament_density = "medium"
        layout_style.ornament_density = _effective_ornament_density(layout_style.ornament_density, layout, binding)
        if layout_style.text_first_contract:
            layout_style.ornament_density = _cap_ornament_density(
                layout_style.ornament_density,
                str(
                    layout_style.decorative_budget_plan.get("allowed_ornament_density")
                    or layout_style.decorative_budget.get("max_ornament_density")
                    or "medium"
                ),
            )
            if layout_style.archetype_id not in {"creative_cover", "section_divider", "visual_table_of_contents"}:
                layout_style.ornament_density = _cap_ornament_density(layout_style.ornament_density, "low")
            layout_style.footer_style = "citation_dense"
        return layout_style

    def report_tokens(self) -> dict[str, Any]:
        density_settings = self.ornament_density_tokens.get(self.ornament_density, {})
        return {
            "tone_variant": self.tone_variant,
            "palette": {key: self.colors.get(key) for key in ("background", "surface", "surface_alt", "accent", "accent_secondary", "text", "line", "grid")},
            "typography": {
                key: self.typography.get(key, {}).get("size_pt")
                for key in ("title", "subtitle", "body", "label", "footer")
                if isinstance(self.typography.get(key), dict)
            },
            "footer_style": self.footer_style,
            "card_style": self.card_style,
            "chart_table_style": self.chart_table_style,
            "section_style": self.section_style,
            "background_ornament_intensity": self.ornament_density,
            "ornament_density_mode": self.ornament_density,
            "ornament_density_tokens": density_settings,
            "image_frame_style": self.image_frame_style,
            "dark_hero": self.dark_hero,
            "index_rail": self.index_rail,
            "expressive_cover_identity": self.expressive_cover_identity,
            "fractured_geometry": self.fractured_geometry,
            "archetype_id": self.archetype_id,
            "priority_master_identity": self.priority_master_identity,
            "layout_reference_source": self.layout_reference_source,
            "identity_lock": self.identity_lock,
            "minimum_visual_features": self.minimum_visual_features,
            "refinement_patch_directives": self.refinement_patch_directives,
            "text_first_contract": self.text_first_contract,
            "decorative_budget": self.decorative_budget,
            "decorative_budget_plan": self.decorative_budget_plan,
        }

    def _apply_tone_variant(self, tone_variant: str) -> None:
        if tone_variant == "academic":
            self.colors.update(
                {
                    "background": self.colors.get("paper", "#FBFAF6"),
                    "surface": "#FFFFFF",
                    "surface_alt": "#F8F5ED",
                    "accent": self.colors.get("gold_accent", "#B8872D"),
                    "accent_secondary": self.colors.get("academic_blue", "#1E3A5F"),
                    "line": "#C9B98E",
                    "grid": "#E8DEC6",
                }
            )
            self._scale_type({"title": -2, "subtitle": -1, "body": -1, "label": -1, "footer": 0})
            self.footer_style = "citation_dense"
            self.card_style = "evidence_card"
            self.chart_table_style = "dense_academic"
            self.section_style = "restrained_navy_band"
            self.ornament_density = "low"
            self.image_frame_style = "thin_caption_frame"
        elif tone_variant == "professional":
            self.colors.update(
                {
                    "background": "#F6F8FB",
                    "surface": "#FFFFFF",
                    "surface_alt": "#EAF2F5",
                    "accent": self.colors.get("primary_navy", "#12355B"),
                    "accent_secondary": self.colors.get("teal", "#0F766E"),
                    "line": "#A8B4C2",
                    "grid": "#D7DEE8",
                }
            )
            self._scale_type({"title": 0, "subtitle": 0, "body": 0, "label": 0, "footer": 0})
            self.footer_style = "crisp_dashboard"
            self.card_style = "decision_panel"
            self.chart_table_style = "high_contrast_dashboard"
            self.section_style = "navy_teal_band"
            self.ornament_density = "medium"
            self.image_frame_style = "clean_diagonal_frame"
        elif tone_variant == "creative":
            self.colors.update(
                {
                    "background": "#FFF8E8",
                    "surface": "#FFFFFF",
                    "surface_alt": "#F6EED8",
                    "accent": self.colors.get("gold_accent", "#D49A2A"),
                    "accent_secondary": self.colors.get("teal", "#0F766E"),
                    "line": "#D8B86A",
                    "grid": "#EDDDB1",
                }
            )
            self._scale_type({"title": 4, "subtitle": 1, "body": 0, "label": 1, "footer": 0})
            self.footer_style = "editorial_dense"
            self.card_style = "layered_editorial_card"
            self.chart_table_style = "expressive_modular"
            self.section_style = "oversized_creative_marker"
            self.ornament_density = "high"
            self.image_frame_style = "diagonal_circular_frame"

    def _scale_type(self, adjustments: dict[str, float]) -> None:
        for token, delta in adjustments.items():
            payload = self.typography.get(token)
            if isinstance(payload, dict) and isinstance(payload.get("size_pt"), (int, float)):
                updated = dict(payload)
                updated["size_pt"] = max(6, float(payload["size_pt"]) + delta)
                self.typography[token] = updated


def _render_slide(
    slide: Any,
    source: dict[str, Any],
    layout: dict[str, Any],
    binding: dict[str, Any],
    style: _DeckStyle,
    warnings: list[dict[str, Any]],
    *,
    icon_resolver: IconResolver | None = None,
    icon_report: dict[str, Any] | None = None,
) -> None:
    _draw_background(slide, style)
    _draw_adapter_component_primitives(slide, layout, style)
    if style.expressive_cover_identity:
        _draw_expressive_identity_layer(slide, layout, binding, style)
    _draw_footer_rule(slide, style)
    if style.index_rail:
        _draw_index_rail(slide, source, binding, style)
    if _is_section_marker_layout(layout, binding):
        _draw_oversized_section_marker(slide, source, binding, style)
    if style.priority_master_identity:
        _draw_priority_master_motifs(slide, layout, binding, style)
    placement = choose_title_body_slots(layout)
    binding.setdefault("slot_bindings", {})
    if placement["title"] and "title" not in binding["slot_bindings"]:
        binding["slot_bindings"][placement["title"]] = "title"
        warnings.append(_render_warning("TITLE_SLOT_FALLBACK_BOUND", source, f"Title was bound to {placement['title']} deterministically."))
    if placement["body"] and placement["body"] not in binding["slot_bindings"]:
        binding["slot_bindings"][placement["body"]] = "content_blocks"
        warnings.append(_render_warning("BODY_SLOT_FALLBACK_BOUND", source, f"Body content was bound to {placement['body']} deterministically."))
    for slot in layout["slots"]:
        source_ref = binding.get("slot_bindings", {}).get(slot["slot_id"])
        _render_slot(slide, source, slot, source_ref, style, warnings)
        _render_slot_icon(slide, source, slot, style, warnings, icon_resolver, icon_report)
    if _render_missing_citation_footer(slide, source, layout, style):
        warnings.append(
            _render_warning(
                "CITATION_FOOTER_FALLBACK_RENDERED",
                source,
                "The selected layout omitted a footer slot; citations were preserved in a deterministic editable footer.",
            )
        )
    for component in _iter_layout_icon_components(layout):
        _render_slot_icon(slide, source, component, style, warnings, icon_resolver, icon_report)


def _draw_adapter_component_primitives(slide: Any, layout: dict[str, Any], style: _DeckStyle) -> None:
    geometry_strategy = layout.get("geometry_strategy") or {}
    primitives = geometry_strategy.get("adapter_component_primitives") or []
    for primitive in sorted(
        (item for item in primitives if isinstance(item, dict) and item.get("render_before_slots")),
        key=lambda item: int(item.get("z_index") or 0),
    ):
        component_type = _normalize_key(primitive.get("source_component_type"))
        if component_type in {"group", "text_box", "footer_component", "image_frame", "chart_frame", "table_frame"}:
            continue
        bounds = primitive.get("bounds") or {}
        if not all(key in bounds for key in ("x", "y", "w", "h")):
            continue
        x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
        if (component_type == "line" and w <= 0 and h <= 0) or (component_type != "line" and (w <= 0 or h <= 0)):
            continue
        source_component_id = str(primitive.get("source_component_id") or primitive.get("primitive_id") or "component")
        if component_type == "line":
            shape = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x),
                Inches(y),
                Inches(x + w),
                Inches(y + h),
            )
            shape.line.color.rgb = style.color("grid", "#D5DEEB")
            shape.line.width = Pt(0.6)
        else:
            auto_shape_type = {
                "rounded_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                "freeform_polygon": MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
                "svg_icon": MSO_AUTO_SHAPE_TYPE.OVAL,
                "svg_layer": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            }.get(component_type, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
            shape = slide.shapes.add_shape(auto_shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
            if component_type in {"svg_icon", "svg_layer"}:
                shape.fill.background()
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = style.color("surface_alt", "#F6F8FB")
            shape.line.color.rgb = style.color("grid", "#D5DEEB")
            shape.line.width = Pt(0.5)
        try:
            shape.name = f"Adapter Primitive {source_component_id}"
        except (AttributeError, ValueError):
            pass


def _draw_priority_master_motifs(slide: Any, layout: dict[str, Any], binding: dict[str, Any], style: _DeckStyle) -> None:
    archetype_id = str(layout.get("archetype_id") or binding.get("slide_type") or style.archetype_id)
    if style.identity_lock == "high":
        _draw_reference_geometry_lock(slide, style)
    if archetype_id == "visual_table_of_contents":
        _draw_visual_toc_reference_grid(slide, style)
    elif archetype_id == "research_overview" and not style.text_first_contract:
        _draw_research_overview_reference_frame(slide, style)
    elif archetype_id == "methodology_framework" and not style.text_first_contract:
        _draw_methodology_reference_scaffold(slide, style)
    elif archetype_id == "data_table_appendix":
        _draw_appendix_table_reference_frame(slide, style)
    elif archetype_id == "creative_cover" and style.minimum_visual_features.get("image_mask_presence"):
        _draw_reference_image_mask(slide, style)
    elif archetype_id == "section_divider" and style.minimum_visual_features.get("index_presence"):
        _draw_reference_section_index(slide, style)
    _draw_refinement_patch_directives(slide, style)


def _draw_refinement_patch_directives(slide: Any, style: _DeckStyle) -> None:
    patch_types = {
        str(directive.get("patch_type") or "")
        for directive in style.refinement_patch_directives
        if isinstance(directive, dict)
    }
    if not patch_types:
        return
    if style.text_first_contract:
        patch_types = {patch for patch in patch_types if patch in {"strengthen_footer_strip", "add_index_markers", "adjust_photo_mask_geometry"}}
        if not patch_types:
            return
    geometry = style.reference_geometry or {}
    content_zone = _reference_zone(geometry.get("content_zone"), style) or {"x": 0.82, "y": 1.28, "w": 11.55, "h": 4.82}
    visual_zone = _reference_zone(geometry.get("visual_zone"), style) or content_zone
    footer_zone = _reference_zone(geometry.get("footer_zone"), style) or {"x": 0.55, "y": 6.72, "w": 12.2, "h": 0.44}
    index_zone = _reference_zone(geometry.get("index_zone"), style) or {"x": 0.24, "y": 0.72, "w": 0.46, "h": 5.58}
    image_zone = _reference_zone(geometry.get("image_mask_zone"), style) or {"x": 7.45, "y": 1.18, "w": 3.85, "h": 4.25}

    if "increase_dark_panel_ratio" in patch_types and not style.dark_hero:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
            Inches(max(0.3, visual_zone["x"] + visual_zone["w"] * 0.58)),
            Inches(max(0.45, visual_zone["y"] - 0.18)),
            Inches(min(4.1, visual_zone["w"] * 0.42)),
            Inches(min(5.9, visual_zone["h"] + 0.4)),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = style.color("surface_alt", "#102E54")
        panel.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
        panel.line.width = Pt(0.7)
    if "strengthen_footer_strip" in patch_types and not style.footer_chrome_simplified:
        _draw_reference_footer_system(slide, footer_zone, style)
        _draw_reference_footer_system(slide, {**footer_zone, "y": max(0.0, footer_zone["y"] - 0.08), "h": max(0.28, footer_zone["h"] * 0.72)}, style)
    if "add_topology_ornaments" in patch_types and not style.text_first_contract:
        _draw_dense_topology_field(
            slide,
            style,
            x0=max(0.55, content_zone["x"]),
            y0=max(0.52, content_zone["y"] - 0.36),
            w=min(11.8, content_zone["w"]),
            h=min(5.4, content_zone["h"] + 0.62),
            step=0.30,
        )
        _draw_topology_network(slide, style, origin_x=max(0.58, content_zone["x"] + content_zone["w"] - 2.4), origin_y=max(0.75, content_zone["y"] - 0.1), scale=0.78)
    if "add_index_markers" in patch_types:
        _draw_reference_index_zone(slide, index_zone, style)
    if "adjust_photo_mask_geometry" in patch_types:
        _draw_reference_image_mask(slide, style, image_zone)
    if ("increase_card_density" in patch_types or "reduce_white_space" in patch_types) and (not style.text_first_contract or style.archetype_id == "research_overview"):
        _draw_reference_card_scaffold(slide, content_zone, style)
    if "increase_table_chrome" in patch_types:
        _draw_reference_table_grid(slide, content_zone, style)
        _add_data_module_chrome(slide, content_zone, style, "APPENDIX")
    if "improve_chart_module_density" in patch_types:
        _add_data_module_chrome(slide, visual_zone, style, "CHART")


def _draw_reference_geometry_lock(slide: Any, style: _DeckStyle) -> None:
    geometry = style.reference_geometry or {}
    visual_zone = _reference_zone(geometry.get("visual_zone"), style)
    content_zone = _reference_zone(geometry.get("content_zone"), style)
    footer_zone = _reference_zone(geometry.get("footer_zone"), style)
    index_zone = _reference_zone(geometry.get("index_zone"), style)
    image_zone = _reference_zone(geometry.get("image_mask_zone"), style)
    if visual_zone and not style.text_first_contract:
        _draw_reference_zone_frame(slide, visual_zone, style, "accent_secondary", fill=False, weight=0.45)
    if content_zone and str(style.minimum_visual_features.get("card_density") or "") in {"medium", "high"} and not style.text_first_contract:
        _draw_reference_card_scaffold(slide, content_zone, style)
    if (footer_zone or style.minimum_visual_features.get("footer_presence")) and not style.footer_chrome_simplified:
        zone = footer_zone or {"x": 0.55, "y": 6.78, "w": 12.15, "h": 0.35}
        _draw_reference_footer_system(slide, zone, style)
    if index_zone or style.minimum_visual_features.get("index_presence"):
        zone = index_zone or {"x": 0.22, "y": 0.62, "w": 0.42, "h": 5.65}
        _draw_reference_index_zone(slide, zone, style)
    if image_zone or style.minimum_visual_features.get("image_mask_presence"):
        zone = image_zone or {"x": 7.45, "y": 1.18, "w": 3.85, "h": 4.25}
        _draw_reference_image_mask(slide, style, zone)
    if style.minimum_visual_features.get("table_density") == "high" and not style.text_first_contract:
        zone = content_zone or visual_zone or {"x": 0.75, "y": 1.3, "w": 11.8, "h": 4.8}
        _draw_reference_table_grid(slide, zone, style)


def _reference_zone(zone: Any, style: _DeckStyle) -> dict[str, float] | None:
    if not isinstance(zone, dict):
        return None
    try:
        x = float(zone.get("x"))
        y = float(zone.get("y"))
        w = float(zone.get("w"))
        h = float(zone.get("h"))
    except (TypeError, ValueError):
        return None
    if str(zone.get("coordinate_system") or "").startswith("normalized"):
        return {
            "x": x * style.canvas_width,
            "y": y * style.canvas_height,
            "w": w * style.canvas_width,
            "h": h * style.canvas_height,
        }
    return {"x": x, "y": y, "w": w, "h": h}


def _draw_reference_zone_frame(
    slide: Any,
    zone: dict[str, float],
    style: _DeckStyle,
    color_token: str,
    *,
    fill: bool,
    weight: float,
) -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(zone["x"]), Inches(zone["y"]), Inches(zone["w"]), Inches(zone["h"]))
    if fill:
        rect.fill.solid()
        rect.fill.fore_color.rgb = style.color("surface_alt", "#EEF2F7")
        rect.fill.transparency = 35
    else:
        rect.fill.background()
    rect.line.color.rgb = style.color(color_token, "#0F766E")
    rect.line.width = Pt(weight)


def _draw_reference_card_scaffold(slide: Any, zone: dict[str, float], style: _DeckStyle) -> None:
    cols = 4 if zone["w"] > 7.5 else 3
    gap = 0.12
    card_w = max(0.5, (zone["w"] - gap * (cols - 1)) / cols)
    card_h = min(1.25, max(0.5, zone["h"] * 0.55))
    y = zone["y"] + zone["h"] * 0.18
    for index in range(cols):
        x = zone["x"] + index * (card_w + gap)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = style.color("surface", "#FFFFFF")
        card.fill.transparency = 8
        card.line.color.rgb = style.color("line", "#CBD5E1")
        card.line.width = Pt(0.55)
        rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.12), Inches(y + 0.16), Inches(card_w - 0.24), Inches(0.04))
        rail.fill.solid()
        rail.fill.fore_color.rgb = style.color("accent" if index % 2 == 0 else "accent_secondary", "#D0AF37")
        rail.line.fill.background()


def _draw_reference_footer_system(slide: Any, zone: dict[str, float], style: _DeckStyle) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(zone["x"]), Inches(zone["y"]), Inches(min(0.16, zone["w"])), Inches(zone["h"]))
    band.fill.solid()
    band.fill.fore_color.rgb = style.color("accent", "#D0AF37")
    band.line.fill.background()
    rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(zone["x"]), Inches(zone["y"] + zone["h"] * 0.18), Inches(zone["x"] + zone["w"]), Inches(zone["y"] + zone["h"] * 0.18))
    rule.line.color.rgb = style.color("line", "#CBD5E1")
    rule.line.width = Pt(0.55)
    for index in range(5):
        tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(zone["x"] + 0.42 + index * 0.32), Inches(zone["y"] + zone["h"] * 0.55), Inches(0.18), Inches(0.025))
        tick.fill.solid()
        tick.fill.fore_color.rgb = style.color("grid", "#E2E8F0")
        tick.line.fill.background()


def _draw_reference_index_zone(slide: Any, zone: dict[str, float], style: _DeckStyle) -> None:
    rail_x = zone["x"] + min(0.08, zone["w"] * 0.25)
    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(rail_x), Inches(zone["y"]), Inches(0.045), Inches(zone["h"]))
    rail.fill.solid()
    rail.fill.fore_color.rgb = style.color("accent_secondary", "#0FA3A3")
    rail.line.fill.background()
    for index in range(4):
        y = zone["y"] + zone["h"] * (0.14 + index * 0.22)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(rail_x - 0.045), Inches(y), Inches(0.135), Inches(0.135))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent" if index == 0 else "surface", "#D0AF37")
        dot.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
        dot.line.width = Pt(0.45)


def _draw_reference_image_mask(slide: Any, style: _DeckStyle, zone: dict[str, float] | None = None) -> None:
    zone = zone or {"x": 7.8, "y": 1.1, "w": 3.3, "h": 4.3}
    mask = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(zone["x"]), Inches(zone["y"]), Inches(zone["w"]), Inches(zone["h"]))
    mask.fill.background()
    mask.line.color.rgb = style.color("accent", "#D0AF37")
    mask.line.width = Pt(1.2)
    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(zone["x"] + 0.16), Inches(zone["y"] + 0.16), Inches(max(0.1, zone["w"] - 0.32)), Inches(max(0.1, zone["h"] - 0.32)))
    inner.fill.solid()
    inner.fill.fore_color.rgb = style.color("surface_alt", "#EEF2F7")
    inner.fill.transparency = 35
    inner.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
    inner.line.width = Pt(0.55)


def _draw_reference_section_index(slide: Any, style: _DeckStyle) -> None:
    for index, (x, y) in enumerate(((0.9, 0.82), (1.28, 1.08), (1.66, 0.92))):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.28), Inches(0.08))
        box.fill.solid()
        box.fill.fore_color.rgb = style.color("accent" if index == 0 else "accent_secondary", "#D0AF37")
        box.line.fill.background()


def _draw_reference_table_grid(slide: Any, zone: dict[str, float], style: _DeckStyle) -> None:
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(zone["x"]), Inches(zone["y"]), Inches(zone["w"]), Inches(0.32))
    header.fill.solid()
    header.fill.fore_color.rgb = style.color("surface_alt", "#EEF2F7")
    header.line.color.rgb = style.color("line", "#CBD5E1")
    header.line.width = Pt(0.35)
    rows = 7
    cols = 5
    for row in range(rows + 1):
        y = zone["y"] + 0.32 + row * max(0.28, (zone["h"] - 0.32) / rows)
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(zone["x"]), Inches(y), Inches(zone["x"] + zone["w"]), Inches(y))
        rule.line.color.rgb = style.color("grid", "#E2E8F0")
        rule.line.width = Pt(0.25)
    for col in range(cols + 1):
        x = zone["x"] + col * zone["w"] / cols
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(zone["y"]), Inches(x), Inches(zone["y"] + zone["h"]))
        rule.line.color.rgb = style.color("grid", "#E2E8F0")
        rule.line.width = Pt(0.22)


def _draw_visual_toc_reference_grid(slide: Any, style: _DeckStyle) -> None:
    for index, x in enumerate((1.05, 3.88, 6.71, 9.54), start=1):
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(1.62), Inches(2.1), Inches(0.08))
        marker.fill.solid()
        marker.fill.fore_color.rgb = style.color("accent" if index % 2 else "accent_secondary", "#D49A2A")
        marker.line.fill.background()
        _add_textbox(slide, f"{index:02d}", {"x": x + 0.02, "y": 1.28, "w": 0.34, "h": 0.18}, style, "footer", "accent")
    rail_y = 5.95
    for index in range(9):
        x = 1.35 + index * 1.17
        tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(rail_y), Inches(0.34), Inches(0.055))
        tick.fill.solid()
        tick.fill.fore_color.rgb = style.color("accent_secondary" if index % 2 else "accent", "#0F766E")
        tick.line.fill.background()


def _draw_research_overview_reference_frame(slide: Any, style: _DeckStyle) -> None:
    for y in (1.45, 2.6, 3.75, 4.9):
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.82), Inches(y), Inches(12.35), Inches(y + 0.04))
        rule.line.color.rgb = style.color("grid", "#E8DEC6")
        rule.line.width = Pt(0.22)
    for index in range(5):
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.45 + index * 0.14), Inches(1.06), Inches(0.055), Inches(0.055))
        chip.fill.solid()
        chip.fill.fore_color.rgb = style.color("accent" if index % 2 else "accent_secondary", "#D0AF37")
        chip.line.fill.background()
    if style.text_first_contract:
        for row in range(2):
            for col in range(5):
                x = 7.68 + col * 0.42
                y = 5.28 + row * 0.22
                bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.28), Inches(0.035))
                bar.fill.solid()
                bar.fill.fore_color.rgb = style.color("accent" if (row + col) % 2 else "accent_secondary", "#D0AF37")
                bar.line.fill.background()


def _draw_methodology_reference_scaffold(slide: Any, style: _DeckStyle) -> None:
    cx, cy = 6.85, 3.1
    for scale in (1.0, 0.72, 0.44):
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - scale), Inches(cy - scale), Inches(scale * 2), Inches(scale * 2))
        ring.fill.background()
        ring.line.color.rgb = style.color("line", "#C9B98E")
        ring.line.width = Pt(0.38)
    for x1, y1, x2, y2 in ((4.1, 3.1, 9.6, 3.1), (6.85, 1.15, 6.85, 5.05), (4.9, 1.72, 8.8, 4.48), (8.8, 1.72, 4.9, 4.48)):
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        rule.line.color.rgb = style.color("grid", "#E8DEC6")
        rule.line.width = Pt(0.24)


def _draw_appendix_table_reference_frame(slide: Any, style: _DeckStyle) -> None:
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.75), Inches(0.42), Inches(1.35), Inches(0.28))
    tag.fill.solid()
    tag.fill.fore_color.rgb = style.color("accent_secondary", "#1E4A5F")
    tag.line.fill.background()
    _add_textbox(slide, "APPENDIX", {"x": 10.93, "y": 0.49, "w": 0.9, "h": 0.12}, style, "footer", "surface")
    for index in range(7):
        y = 1.25 + index * 0.68
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.72), Inches(y), Inches(12.62), Inches(y))
        rule.line.color.rgb = style.color("grid", "#E8DEC6")
        rule.line.width = Pt(0.18)


def _draw_background(slide: Any, style: _DeckStyle) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(style.canvas_width), Inches(style.canvas_height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = style.color("background", "#F8FAFC")
    shape.line.fill.background()
    rule_width = 0.28 if style.text_first_contract else 0.45 if style.ornament_density == "low" else 0.75 if style.ornament_density == "medium" else 1.05
    for y in (0.18, 7.22):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.58), Inches(y), Inches(12.75), Inches(y))
        line.line.color.rgb = style.color("grid", "#E2E8F0")
        line.line.width = Pt(rule_width)
    grid_x = (0.62, 12.7) if style.text_first_contract else (3.0, 6.0, 9.0, 12.0) if style.ornament_density == "low" else (1.5, 3.0, 4.5, 6.0, 7.5, 9.0, 10.5, 12.0)
    for x in grid_x:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(0.55), Inches(x), Inches(6.65))
        line.line.color.rgb = style.color("grid", "#E2E8F0")
        line.line.width = Pt(0.12 if style.ornament_density == "low" else 0.18)
    corner_dots = ((0.7, 0.58), (12.48, 0.58), (0.7, 6.58), (12.48, 6.58)) if not style.text_first_contract else ((0.7, 0.58), (12.48, 0.58))
    for x, y in corner_dots:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.045), Inches(0.045))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent_secondary", "#0F766E")
        dot.line.fill.background()
    if style.ornament_density == "high" and not style.text_first_contract:
        for x, y, size in ((10.9, 0.62, 0.36), (11.45, 1.02, 0.18), (0.94, 5.92, 0.24)):
            ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
            ring.fill.background()
            ring.line.color.rgb = style.color("accent", "#D49A2A")
            ring.line.width = Pt(1.2)
        diagonal = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.7), Inches(0.35), Inches(12.55), Inches(2.35))
        diagonal.line.color.rgb = style.color("accent_secondary", "#0F766E")
        diagonal.line.width = Pt(1.0)
        _draw_topology_network(slide, style, origin_x=0.35 if style.dark_hero else 9.6, origin_y=0.7 if style.dark_hero else 0.9, scale=1.0 if style.dark_hero else 0.65)
    if style.dark_hero:
        rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.52), Inches(style.canvas_height))
        rail.fill.solid()
        rail.fill.fore_color.rgb = style.color("surface", "#0B2446")
        rail.line.fill.background()
        if not style.text_first_contract:
            _draw_topology_network(slide, style, origin_x=7.95, origin_y=1.2, scale=1.35)


def _draw_expressive_identity_layer(slide: Any, layout: dict[str, Any], binding: dict[str, Any], style: _DeckStyle) -> None:
    layout_id = str(layout.get("layout_id") or "")
    if style.dark_hero:
        _draw_dark_header_band(slide, style)
        if not style.text_first_contract:
            _draw_dense_topology_field(slide, style, x0=0.82, y0=0.32, w=11.85, h=6.0, step=0.24)
            _draw_cover_topology_lattice(slide, style)
            _draw_cover_blueprint_guides(slide, style)
            _draw_topology_network(slide, style, origin_x=0.55, origin_y=0.88, scale=1.4)
            _draw_topology_network(slide, style, origin_x=6.95, origin_y=1.15, scale=1.55)
            _draw_high_density_cover_ornaments(slide, style)
        else:
            _draw_topology_network(slide, style, origin_x=10.25, origin_y=0.88, scale=0.58)
        _draw_fractured_geometry_panel(slide, style, dark=True)
        _draw_creative_section_tab(slide, style, "01", x=0.78, y=0.11)
    elif "section_divider" in layout_id or str(binding.get("slide_type") or "") == "section_divider":
        if not style.text_first_contract:
            _draw_section_contour_field(slide, style)
            _draw_high_density_section_ornaments(slide, style)
        _draw_dark_header_band(slide, style, compact=True)
        _draw_fractured_geometry_panel(slide, style, dark=False)
        _draw_creative_section_tab(slide, style, "02", x=0.2, y=0.08)
        if style.text_first_contract:
            for index in range(6):
                tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.15 + index * 0.28), Inches(5.86), Inches(0.18), Inches(0.035))
                tick.fill.solid()
                tick.fill.fore_color.rgb = style.color("accent" if index % 2 else "accent_secondary", "#D0AF37")
                tick.line.fill.background()


def _draw_dark_header_band(slide: Any, style: _DeckStyle, *, compact: bool = False) -> None:
    h = 0.3 if compact else 0.38
    w = 2.1 if compact else 2.55
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.08), Inches(w), Inches(h))
    band.fill.solid()
    band.fill.fore_color.rgb = style.color("surface", "#0B2446") if style.dark_hero else style.color("accent_secondary", "#0FA3A3")
    band.line.fill.background()
    tab = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(0.36), Inches(0.08), Inches(0.38), Inches(h))
    tab.fill.solid()
    tab.fill.fore_color.rgb = style.color("accent", "#D0AF37")
    tab.line.fill.background()


def _draw_dense_topology_field(slide: Any, style: _DeckStyle, *, x0: float, y0: float, w: float, h: float, step: float) -> None:
    x = x0
    row = 0
    while x <= x0 + w:
        y = y0 + (0.08 if row % 2 else 0)
        while y <= y0 + h:
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.018), Inches(0.018))
            dot.fill.solid()
            dot.fill.fore_color.rgb = style.color("accent" if style.dark_hero and (row + int(y * 10)) % 7 == 0 else "grid", "#24405E")
            dot.line.fill.background()
            y += step
        row += 1
        x += step


def _draw_cover_topology_lattice(slide: Any, style: _DeckStyle) -> None:
    points = [
        (5.95, 1.28), (6.55, 2.05), (7.05, 1.62), (7.55, 2.58),
        (8.1, 1.34), (8.7, 2.22), (9.18, 1.62), (9.72, 2.78),
        (6.25, 3.42), (7.08, 4.35), (8.0, 3.86), (8.8, 4.72),
        (9.55, 3.72), (10.18, 4.42), (10.82, 3.2),
    ]
    edges = [
        (0, 1), (0, 2), (1, 2), (2, 3), (2, 4), (3, 5), (4, 5),
        (5, 6), (5, 10), (6, 7), (8, 9), (9, 10), (10, 11), (10, 12),
        (11, 13), (12, 13), (13, 14), (7, 14), (3, 8), (4, 10),
    ]
    for start_idx, end_idx in edges:
        start = points[start_idx]
        end = points[end_idx]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start[0]), Inches(start[1]), Inches(end[0]), Inches(end[1]))
        line.line.color.rgb = style.color("accent", "#D0AF37")
        line.line.width = Pt(0.55)
    for index, (x, y) in enumerate(points):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.035), Inches(y - 0.035), Inches(0.07), Inches(0.07))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent" if index % 3 else "accent_secondary", "#D0AF37")
        dot.line.fill.background()


def _draw_cover_blueprint_guides(slide: Any, style: _DeckStyle) -> None:
    """Add editable guide-line and connector density matching the dark design-board hero."""

    mesh_x = 4.65
    mesh_y = 0.42
    mesh_w = 7.85
    mesh_h = 5.95
    for index in range(44):
        x = mesh_x + index * 0.18
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(mesh_y), Inches(x + 0.18), Inches(mesh_y + mesh_h))
        line.line.color.rgb = style.color("accent_secondary" if index % 8 == 0 else "grid", "#0FA3A3")
        line.line.width = Pt(0.22 if index % 8 == 0 else 0.14)
    for index in range(34):
        y = mesh_y + index * 0.18
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(mesh_x), Inches(y), Inches(mesh_x + mesh_w), Inches(y + 0.06))
        line.line.color.rgb = style.color("accent" if index % 7 == 0 else "grid", "#D0AF37")
        line.line.width = Pt(0.2 if index % 7 == 0 else 0.12)
    for index in range(32):
        x = mesh_x + (index % 8) * 0.92 + 0.12
        y = mesh_y + (index // 8) * 1.08 + 0.28
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.035), Inches(0.035))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent" if index % 3 == 0 else "accent_secondary", "#D0AF37")
        dot.line.fill.background()
    for index in range(15):
        x = 0.85 + index * 0.74
        guide = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(0.45), Inches(x + 1.8), Inches(6.35))
        guide.line.color.rgb = style.color("accent_secondary" if index % 4 == 0 else "accent", "#0FA3A3")
        guide.line.width = Pt(0.42 if index % 4 == 0 else 0.28)
    for index in range(13):
        y = 0.7 + index * 0.42
        guide = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.95), Inches(y), Inches(12.15), Inches(y + 0.32))
        guide.line.color.rgb = style.color("accent" if index % 3 == 0 else "grid", "#D0AF37")
        guide.line.width = Pt(0.32 if index % 3 == 0 else 0.24)
    for index in range(18):
        y = 0.54 + index * 0.29
        start = 4.55 + (index % 4) * 0.14
        end = 12.45 - (index % 5) * 0.1
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start), Inches(y), Inches(end), Inches(y + 0.22))
        rule.line.color.rgb = style.color("accent_secondary" if index % 3 == 0 else "grid", "#0FA3A3")
        rule.line.width = Pt(0.34 if index % 3 == 0 else 0.26)
    for index in range(15):
        x = 4.8 + index * 0.5
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(0.5), Inches(x + 0.55), Inches(6.05))
        rule.line.color.rgb = style.color("grid", "#24405E")
        rule.line.width = Pt(0.28)
    for index in range(15):
        x = 5.1 + index * 0.5
        y = 0.85 + (index % 5) * 0.58
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.045), Inches(y - 0.045), Inches(0.09), Inches(0.09))
        node.fill.solid()
        node.fill.fore_color.rgb = style.color("accent_secondary" if index % 2 else "accent", "#0FA3A3")
        node.line.fill.background()
        tick = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.09), Inches(y), Inches(x + 0.52), Inches(y + 0.16))
        tick.line.color.rgb = style.color("accent", "#D0AF37")
        tick.line.width = Pt(0.48)
    for index in range(9):
        y = 1.05 + index * 0.48
        arc = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(7.15), Inches(y), Inches(11.6), Inches(y + 0.36))
        arc.line.color.rgb = style.color("accent" if index % 2 else "accent_secondary", "#D0AF37")
        arc.line.width = Pt(0.32)
    for index, scale in enumerate((1.0, 0.78, 0.56, 0.36)):
        ring = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(9.55 - 0.95 * scale),
            Inches(3.65 - 0.95 * scale),
            Inches(1.9 * scale),
            Inches(1.9 * scale),
        )
        ring.fill.background()
        ring.line.color.rgb = style.color("accent_secondary" if index % 2 else "accent", "#D0AF37")
        ring.line.width = Pt(0.45)


def _draw_high_density_cover_ornaments(slide: Any, style: _DeckStyle) -> None:
    if style.ornament_density != "high":
        return
    settings = style.ornament_density_tokens.get("high", {})
    max_shapes = max(24, min(int(settings.get("max_extra_shapes") or 72), 86))
    line_weight = float(settings.get("line_weight_pt") or 0.24)
    guide_count = min(int(settings.get("guide_line_count") or 18), 24)
    node_count = min(int(settings.get("micro_node_count") or 30), 36)
    arc_count = min(int(settings.get("blueprint_arc_count") or 5), 7)
    shape_budget = max_shapes

    for index in range(min(guide_count, shape_budget)):
        x = 0.85 + (index % 9) * 0.86
        y = 0.62 + (index // 9) * 2.25
        rail = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + 2.15), Inches(y + 0.46))
        rail.line.color.rgb = style.color("accent_secondary" if index % 3 else "accent", "#0FA3A3")
        rail.line.width = Pt(line_weight + (0.08 if index % 3 == 0 else 0))
        shape_budget -= 1
    for index in range(min(node_count, shape_budget)):
        x = 0.72 + (index % 12) * 0.94
        y = 0.84 + (index // 12) * 1.58 + (0.18 if index % 2 else 0)
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.04), Inches(0.04))
        node.fill.solid()
        node.fill.fore_color.rgb = style.color("accent" if index % 4 == 0 else "accent_secondary", "#D0AF37")
        node.line.fill.background()
        shape_budget -= 1
    for index in range(min(arc_count, shape_budget)):
        y = 1.05 + index * 0.58
        arc = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(8.4), Inches(y), Inches(12.05), Inches(y + 0.32))
        arc.line.color.rgb = style.color("grid" if index % 2 else "accent_secondary", "#24405E")
        arc.line.width = Pt(line_weight)
        shape_budget -= 1
    for index, (x, y, w, h) in enumerate(((8.3, 0.8, 3.25, 0.34), (8.85, 5.05, 2.85, 0.26), (5.2, 5.8, 3.6, 0.2))):
        if shape_budget <= 0:
            break
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(x), Inches(y), Inches(w), Inches(h))
        band.fill.solid()
        band.fill.fore_color.rgb = style.color("surface_alt", "#102E54")
        band.line.color.rgb = style.color("accent_secondary" if index % 2 else "accent", "#0FA3A3")
        band.line.width = Pt(0.35)
        shape_budget -= 1


def _draw_high_density_section_ornaments(slide: Any, style: _DeckStyle) -> None:
    if style.ornament_density != "high":
        return
    settings = style.ornament_density_tokens.get("high", {})
    line_weight = float(settings.get("line_weight_pt") or 0.24)
    guide_count = min(int(settings.get("guide_line_count") or 18), 18)
    for index in range(guide_count):
        x = 0.55 + (index % 6) * 1.04
        y = 0.72 + (index // 6) * 1.65
        rail = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + 2.4), Inches(y + 0.18))
        rail.line.color.rgb = style.color("line" if index % 2 else "accent_secondary", "#D6B666")
        rail.line.width = Pt(line_weight)
    for index in range(18):
        x = 6.9 + (index % 6) * 0.74
        y = 0.88 + (index // 6) * 0.66
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.035), Inches(0.035))
        node.fill.solid()
        node.fill.fore_color.rgb = style.color("accent" if index % 3 == 0 else "accent_secondary", "#D0AF37")
        node.line.fill.background()


def _draw_fractured_geometry_panel(slide: Any, style: _DeckStyle, *, dark: bool) -> None:
    panel_color = "surface_alt" if dark else "surface"
    x, y, w, h = (8.45, 0.52, 3.75, 5.45) if dark else (9.55, 0.42, 2.55, 5.75)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = style.color(panel_color, "#102E54")
    panel.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
    panel.line.width = Pt(0.45 if style.text_first_contract else 0.85)
    if style.text_first_contract:
        return
    for idx, scale in enumerate((1.0, 0.78, 0.56, 0.34)):
        ring = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x + w * 0.42 - 0.98 * scale),
            Inches(y + h * 0.44 - 0.98 * scale),
            Inches(1.96 * scale),
            Inches(1.96 * scale),
        )
        ring.fill.background()
        ring.line.color.rgb = style.color("accent" if idx % 2 == 0 else "accent_secondary", "#D0AF37")
        ring.line.width = Pt(0.8 if idx == 0 else 0.45)
    for offset in (0.0, 0.34, 0.68):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x - 0.28 + offset),
            Inches(y + h - 0.2),
            Inches(x + w - 0.12 + offset),
            Inches(y + 0.2),
        )
        line.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
        line.line.width = Pt(0.55)


def _draw_section_contour_field(slide: Any, style: _DeckStyle) -> None:
    for index in range(7):
        y = 0.95 + index * 0.38
        start_x = 0.38 + index * 0.08
        end_x = 4.15 + index * 0.16
        line = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(start_x), Inches(y), Inches(end_x), Inches(y + 0.52))
        line.line.color.rgb = style.color("line", "#D6B666")
        line.line.width = Pt(0.35)
    _draw_dense_topology_field(slide, style, x0=6.85, y0=0.6, w=4.8, h=2.2, step=0.28)


def _draw_creative_section_tab(slide: Any, style: _DeckStyle, label: str, *, x: float, y: float) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.42), Inches(0.28))
    box.fill.solid()
    box.fill.fore_color.rgb = style.color("accent", "#D0AF37")
    box.line.fill.background()
    if not style.text_first_contract:
        _add_textbox(slide, label, {"x": x + 0.09, "y": y + 0.055, "w": 0.25, "h": 0.12}, style, "footer", "surface" if style.dark_hero else "text")


def _is_section_marker_layout(layout: dict[str, Any], binding: dict[str, Any]) -> bool:
    layout_id = str(layout.get("layout_id") or "")
    rhythm = str(binding.get("section_rhythm_role") or "")
    slide_type = str(binding.get("slide_type") or "")
    return "section_divider" in layout_id or rhythm == "divider" or slide_type == "section_divider"


def _draw_topology_network(slide: Any, style: _DeckStyle, *, origin_x: float, origin_y: float, scale: float) -> None:
    points = [
        (0.0, 0.5),
        (0.55, 0.0),
        (1.1, 0.45),
        (1.65, 0.18),
        (2.1, 0.78),
        (0.38, 1.18),
        (1.35, 1.28),
        (2.34, 1.42),
    ]
    edges = [(0, 1), (1, 2), (2, 3), (2, 6), (3, 4), (0, 5), (5, 6), (6, 7), (4, 7), (1, 6)]
    if style.text_first_contract:
        points = [(0.0, 0.5), (0.55, 0.0), (1.1, 0.45), (1.65, 0.18), (0.38, 1.18), (1.35, 1.28), (2.34, 1.42)]
        edges = [(0, 1), (1, 2), (2, 3), (4, 5), (5, 6)]
    for start_idx, end_idx in edges:
        start = points[start_idx]
        end = points[end_idx]
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(origin_x + start[0] * scale),
            Inches(origin_y + start[1] * scale),
            Inches(origin_x + end[0] * scale),
            Inches(origin_y + end[1] * scale),
        )
        line.line.color.rgb = style.color("accent", "#D0AF37")
        line.line.width = Pt(0.45 if style.dark_hero else 0.32)
    for x, y in points:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(origin_x + x * scale - 0.025), Inches(origin_y + y * scale - 0.025), Inches(0.05), Inches(0.05))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent", "#D0AF37")
        dot.line.fill.background()


def _draw_index_rail(slide: Any, source: dict[str, Any], binding: dict[str, Any], style: _DeckStyle) -> None:
    x = 0.18
    top = 0.58
    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(top), Inches(0.055), Inches(5.85))
    rail.fill.solid()
    rail.fill.fore_color.rgb = style.color("accent_secondary", "#0FA3A3")
    rail.line.fill.background()
    dot_count = 3 if style.text_first_contract else 5
    for index in range(dot_count):
        y = top + 0.35 + index * 1.1
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.055), Inches(y), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent" if index == 0 else "surface", "#D0AF37")
        dot.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
        dot.line.width = Pt(0.55)
    if not style.text_first_contract:
        _add_textbox(slide, "INDEX", {"x": 0.12, "y": 0.18, "w": 0.5, "h": 0.22}, style, "footer", "muted_text")


def _draw_oversized_section_marker(slide: Any, source: dict[str, Any], binding: dict[str, Any], style: _DeckStyle) -> None:
    raw = str(source.get("section_id") or binding.get("section_rhythm_role") or "01")
    match = "".join(ch for ch in raw if ch.isdigit())[-2:] or "01"
    _add_textbox(slide, match.zfill(2), {"x": 5.45, "y": 1.6, "w": 1.75, "h": 1.25}, style, "title", "accent")
    rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.95), Inches(2.28), Inches(7.68), Inches(2.28))
    rule.line.color.rgb = style.color("accent", "#D0AF37")
    rule.line.width = Pt(1.2)
    seal = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.82), Inches(1.82), Inches(0.72), Inches(0.72))
    seal.fill.background()
    seal.line.color.rgb = style.color("accent", "#D0AF37")
    seal.line.width = Pt(1.0)
    _add_textbox(slide, "A", {"x": 8.03, "y": 2.0, "w": 0.28, "h": 0.24}, style, "label", "accent_secondary")


def _draw_footer_rule(slide: Any, style: _DeckStyle) -> None:
    band_w = 0.09 if style.footer_style == "citation_dense" else 0.12 if style.footer_style == "crisp_dashboard" else 0.18
    band_h = 0.12 if style.footer_chrome_simplified else 0.22 if style.footer_style == "citation_dense" else 0.26 if style.footer_style == "crisp_dashboard" else 0.34
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(6.82), Inches(band_w), Inches(band_h))
    band.fill.solid()
    band.fill.fore_color.rgb = style.color("accent", "#2563EB")
    band.line.fill.background()
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.6), Inches(6.78), Inches(12.7), Inches(6.78))
    line.line.color.rgb = style.color("line", "#CBD5E1")
    line.line.width = Pt(0.65 if style.footer_style == "citation_dense" else 0.85 if style.footer_style == "crisp_dashboard" else 1.1)
    if not style.footer_chrome_simplified:
        subline = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.74), Inches(6.86), Inches(12.6), Inches(6.86))
        subline.line.color.rgb = style.color("grid", "#E2E8F0")
        subline.line.width = Pt(0.35)


def _render_slot(
    slide: Any,
    source: dict[str, Any],
    slot: dict[str, Any],
    source_ref: str | None,
    style: _DeckStyle,
    warnings: list[dict[str, Any]],
) -> None:
    slot_type = slot["slot_type"]
    if slot_type in {"text", "footer"}:
        _render_text_slot(slide, source, slot, source_ref, style)
    elif slot_type == "content":
        _render_content_slot(slide, source, slot, source_ref, style, warnings)
    elif slot_type == "table":
        _render_table_slot(slide, source, slot, source_ref, style, warnings)
    elif slot_type == "chart":
        _render_chart_slot(slide, source, slot, source_ref, style, warnings)
    elif slot_type == "image":
        _render_image_slot(slide, source, slot, style)
    else:
        _render_shape_slot(slide, slot, style)


def _render_slot_icon(
    slide: Any,
    source: dict[str, Any],
    slot: dict[str, Any],
    style: _DeckStyle,
    warnings: list[dict[str, Any]],
    icon_resolver: IconResolver | None,
    icon_report: dict[str, Any] | None,
) -> None:
    if icon_resolver is None or icon_report is None or not _component_has_icon(slot):
        return
    context = {
        "slide_id": _slide_id(source),
        "slot_id": str(slot.get("slot_id") or ""),
        "slot_type": str(slot.get("slot_type") or ""),
        "component_id": str(slot.get("component_id") or ""),
    }
    if _icon_is_content_replacement_risk(slot):
        warnings.append(
            _render_warning(
                "ICON_CONTENT_SLOT_REPLACEMENT_FORBIDDEN",
                source,
                f"Icon on required content slot `{context['slot_id']}` was skipped; icons cannot replace title/body/table/chart content.",
            )
        )
        return
    missing_before = len(icon_report["missing_icons"])
    unresolved_before = len(icon_report["unresolved_icon_roles"])
    resolution = icon_resolver.resolve(slot, report=icon_report, context=context)
    _append_icon_resolution_warnings(source, warnings, icon_report, missing_before, unresolved_before)
    if resolution is None:
        return
    bounds = _icon_render_bounds(slot, resolution)
    color_hex = _style_color_hex(style, resolution.color_token, "#111827")
    temp_svg: Path | None = None
    try:
        temp_svg = icon_resolver.materialize_svg(resolution, color_hex=color_hex, stroke_width=resolution.stroke_width)
        placement = _add_svg_icon_picture(slide, temp_svg, bounds, resolution)
    except (OSError, ValueError, AttributeError) as exc:
        warnings.append(
            _render_warning(
                "ICON_SVG_INSERT_FAILED",
                source,
                f"SVG icon `{resolution.icon_id}` could not be inserted and no raster fallback was used: {exc}",
            )
        )
        return
    finally:
        if temp_svg is not None:
            try:
                temp_svg.unlink()
            except OSError:
                pass
    icon_resolver.record_used(
        resolution,
        icon_report,
        context=context,
        color_hex=color_hex,
        bounds=bounds,
        relationship_id=placement.get("relationship_id"),
        object_id=placement.get("object_id"),
        object_name=placement.get("object_name"),
        object_description=placement.get("object_description"),
    )


def _component_has_icon(component: dict[str, Any]) -> bool:
    return any(
        key in component and component.get(key) not in (None, "")
        for key in (
            "icon_role",
            "icon_id",
            "icon_family",
            "icon_size",
            "icon_color_token",
            "icon_stroke_width",
            "icon_position",
        )
    )


def _iter_layout_icon_components(layout: dict[str, Any]) -> list[dict[str, Any]]:
    components: list[dict[str, Any]] = []
    for item in layout.get("components") or []:
        if isinstance(item, dict) and _component_has_icon(item):
            components.append(item)
    for slot in layout.get("slots") or []:
        if not isinstance(slot, dict):
            continue
        for item in slot.get("components") or []:
            if isinstance(item, dict) and _component_has_icon(item):
                component = dict(item)
                component.setdefault("slot_id", slot.get("slot_id"))
                component.setdefault("slot_type", slot.get("slot_type"))
                component.setdefault("bounds", slot.get("bounds"))
                components.append(component)
    return components


def _icon_is_content_replacement_risk(component: dict[str, Any]) -> bool:
    slot_id = str(component.get("slot_id") or "").strip().lower()
    slot_type = str(component.get("slot_type") or "").strip().lower()
    component_id = str(component.get("component_id") or "").strip().lower()
    if component_id in {
        "footer",
        "citation_strip",
        "kpi_card",
        "evidence_chip",
        "section_marker",
        "index_rail",
        "chart_header",
        "table_header",
        "insight_callout",
    }:
        return False
    if isinstance(component.get("icon_position"), dict):
        return False
    if slot_id in {"title", "subtitle", "body"} and slot_type == "text":
        return True
    if slot_id in {"table", "chart"} or slot_type in {"table", "chart"}:
        return True
    return False


def _append_icon_resolution_warnings(
    source: dict[str, Any],
    warnings: list[dict[str, Any]],
    icon_report: dict[str, Any],
    missing_before: int,
    unresolved_before: int,
) -> None:
    for entry in icon_report["missing_icons"][missing_before:]:
        warnings.append(_render_warning(str(entry.get("code") or "ICON_MISSING"), source, str(entry.get("message") or "Icon was not resolved.")))
    for entry in icon_report["unresolved_icon_roles"][unresolved_before:]:
        warnings.append(_render_warning(str(entry.get("code") or "ICON_ROLE_UNRESOLVED"), source, str(entry.get("message") or "Icon role was not resolved.")))


def _icon_render_bounds(component: dict[str, Any], resolution: IconResolution) -> dict[str, float]:
    bounds = component.get("bounds") if isinstance(component.get("bounds"), dict) else {}
    x = _float_or_default(bounds.get("x"), 0.0)
    y = _float_or_default(bounds.get("y"), 0.0)
    w = _float_or_default(bounds.get("w"), resolution.size_in)
    h = _float_or_default(bounds.get("h"), resolution.size_in)
    size = resolution.size_in
    position = component.get("icon_position")
    if isinstance(position, dict):
        explicit_w = _float_or_none(position.get("w"))
        explicit_h = _float_or_none(position.get("h"))
        explicit_size = _float_or_none(position.get("size"))
        size_w = explicit_w or explicit_size or size
        size_h = explicit_h or explicit_size or size
        if _float_or_none(position.get("x")) is not None and _float_or_none(position.get("y")) is not None:
            return {
                "x": _float_or_default(position.get("x"), x),
                "y": _float_or_default(position.get("y"), y),
                "w": size_w,
                "h": size_h,
            }
        anchor = str(position.get("anchor") or "top_left")
    else:
        anchor = str(position or "top_left")
        size_w = size
        size_h = size
    inset = 0.12
    if anchor in {"top_right", "right"}:
        return {"x": x + max(0.0, w - size_w - inset), "y": y + inset, "w": size_w, "h": size_h}
    if anchor in {"bottom_left"}:
        return {"x": x + inset, "y": y + max(0.0, h - size_h - inset), "w": size_w, "h": size_h}
    if anchor in {"bottom_right"}:
        return {"x": x + max(0.0, w - size_w - inset), "y": y + max(0.0, h - size_h - inset), "w": size_w, "h": size_h}
    if anchor in {"center", "middle"}:
        return {"x": x + max(0.0, (w - size_w) / 2), "y": y + max(0.0, (h - size_h) / 2), "w": size_w, "h": size_h}
    return {"x": x + inset, "y": y + inset, "w": size_w, "h": size_h}


def _add_svg_icon_picture(slide: Any, svg_path: Path, bounds: dict[str, float], resolution: IconResolution) -> dict[str, Any]:
    image_part = ImagePart(
        slide.part.package.next_image_partname("svg"),
        SVG_CONTENT_TYPE,
        slide.part.package,
        svg_path.read_bytes(),
        svg_path.name,
    )
    r_id = slide.part.relate_to(image_part, RT.IMAGE)
    shape_id = slide.shapes._next_shape_id
    safe_id = resolution.icon_id.replace("&", "and").replace("<", "").replace(">", "")
    object_name = f"SVG Icon {safe_id}"
    object_description = f"svg-icon:{safe_id}"
    slide.shapes._grpSp.add_pic(
        shape_id,
        object_name,
        object_description,
        r_id,
        Inches(bounds["x"]),
        Inches(bounds["y"]),
        Inches(bounds["w"]),
        Inches(bounds["h"]),
    )
    slide.shapes._recalculate_extents()
    return {
        "relationship_id": r_id,
        "object_id": shape_id,
        "object_name": object_name,
        "object_description": object_description,
    }


def _style_color_hex(style: _DeckStyle, token: str, fallback: str) -> str:
    aliases = {
        "gold": "accent",
        "cyan": "accent_secondary",
        "dark_teal": "accent_secondary",
        "panel_teal": "accent_secondary",
        "off_white": "surface",
        "grid_line": "grid",
        "primary_navy": "text",
    }
    for candidate in (token, aliases.get(token, ""), "accent", "text"):
        value = style.colors.get(candidate)
        if isinstance(value, str) and value.startswith("#") and len(value) == 7:
            return value
    return fallback


def _float_or_none(value: Any) -> float | None:
    try:
        if value in (None, ""):
            return None
        return float(value)
    except (TypeError, ValueError):
        return None


def _float_or_default(value: Any, default: float) -> float:
    parsed = _float_or_none(value)
    return default if parsed is None else parsed


def _render_text_slot(slide: Any, source: dict[str, Any], slot: dict[str, Any], source_ref: str | None, style: _DeckStyle) -> None:
    text = _text_for_source(source, source_ref, slot["slot_id"])
    token = "title" if slot["slot_id"] == "title" else "subtitle" if slot["slot_id"] == "subtitle" else "body"
    bounds = _dark_hero_text_bounds(slot["slot_id"], slot["bounds"], style)
    if slot["slot_id"] == "footer" or str(slot.get("component_id") or "") in {"dense_footer", "citation_strip", "citation_micro_footer"}:
        _draw_citation_footer_system(slide, bounds, style)
        _add_textbox(slide, text, _inset(bounds, 0.26, 0.06), style, "footer", "muted_text" if style.dark_hero else "text")
    elif slot.get("component_id") == "card":
        _add_panel(slide, bounds, style)
        _add_textbox(slide, text, _inset(bounds, 0.18, 0.12), style, token, "text")
    else:
        _add_textbox(slide, text, bounds, style, token, "text")


def _draw_citation_footer_system(slide: Any, bounds: dict[str, Any], style: _DeckStyle) -> None:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    footer_h = min(max(h, 0.24), 0.28) if style.footer_chrome_simplified else max(h, 0.34 if style.footer_style == "citation_dense" else 0.3)
    y = min(y, style.canvas_height - footer_h - 0.12)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(footer_h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = style.color("surface_alt" if not style.dark_hero else "surface", "#F8F5ED")
    strip.line.color.rgb = style.color("line", "#C9B98E")
    strip.line.width = Pt(0.45)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(min(0.16, w)), Inches(footer_h))
    accent.fill.solid()
    accent.fill.fore_color.rgb = style.color("accent", "#D0AF37")
    accent.line.fill.background()
    if not style.footer_chrome_simplified:
        tag_w = min(1.2, max(0.7, w * 0.12))
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.24), Inches(y + 0.07), Inches(tag_w), Inches(max(0.12, footer_h - 0.14)))
        tag.fill.solid()
        tag.fill.fore_color.rgb = style.color("accent_secondary", "#0FA3A3")
        tag.line.fill.background()
        _add_textbox(slide, "SOURCE", {"x": x + 0.34, "y": y + 0.11, "w": tag_w - 0.18, "h": 0.12}, style, "footer", "surface")
        for index in range(6):
            marker_x = x + w - 2.15 + index * 0.28
            tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(marker_x), Inches(y + footer_h * 0.34), Inches(0.12), Inches(footer_h * 0.32))
            tick.fill.solid()
            tick.fill.fore_color.rgb = style.color("accent" if index % 2 == 0 else "grid", "#D0AF37")
            tick.line.fill.background()
        top = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.22), Inches(y + 0.04), Inches(x + w - 0.22), Inches(y + 0.04))
        top.line.color.rgb = style.color("grid", "#E8DEC6")
        top.line.width = Pt(0.32)


def _render_missing_citation_footer(
    slide: Any,
    source: dict[str, Any],
    layout: dict[str, Any],
    style: _DeckStyle,
) -> bool:
    citations = source.get("citations") or []
    if not citations:
        return False
    if any(
        _normalize_key(slot.get("slot_id")) in {"footer", "citation_strip"}
        or _normalize_key(slot.get("component_id")) in {"dense_footer", "citation_strip", "citation_micro_footer"}
        for slot in layout.get("slots") or []
        if isinstance(slot, dict)
    ):
        return False
    text = _text_for_source(source, "citations", "footer")
    if not text.strip():
        return False
    bounds = {
        "x": 0.5,
        "y": max(0.12, style.canvas_height - 0.54),
        "w": max(0.5, style.canvas_width - 1.0),
        "h": 0.3,
    }
    _draw_citation_footer_system(slide, bounds, style)
    _add_textbox(
        slide,
        text,
        _inset(bounds, 0.26, 0.06),
        style,
        "footer",
        "muted_text" if style.dark_hero else "text",
    )
    return True


def _dark_hero_text_bounds(slot_id: str, bounds: dict[str, Any], style: _DeckStyle) -> dict[str, Any]:
    if not style.dark_hero:
        return bounds
    updated = dict(bounds)
    if slot_id == "title":
        updated["y"] = max(0.8, float(updated["y"]) - 0.28)
        updated["h"] = max(float(updated["h"]), 2.0)
        updated["w"] = max(float(updated["w"]), 4.85)
    elif slot_id == "subtitle":
        updated["y"] = max(float(updated["y"]), 4.25)
        updated["h"] = max(float(updated["h"]), 0.7)
    return updated


def _render_content_slot(
    slide: Any,
    source: dict[str, Any],
    slot: dict[str, Any],
    source_ref: str | None,
    style: _DeckStyle,
    warnings: list[dict[str, Any]],
) -> None:
    bounds = slot["bounds"]
    blocks = _content_blocks_for_slot(source, slot["slot_id"], source_ref)
    component_id = str(slot.get("component_id") or "")
    if component_id == "creative_section_tab":
        _render_academic_seal(slide, bounds, style)
    elif component_id == "oversized_section_number":
        _render_oversized_number_slot(slide, source, bounds, style)
    else:
        draws_structured_content = (
            slot["slot_id"] == "cards"
            or slot["slot_id"] in {"index_navigation", "progress_markers", "metric_panels"}
            or slot.get("component_id") in {"radial_process", "curved_timeline", "connector_style"}
            or slot["slot_id"] in {"diagram", "process", "timeline"}
        )
        if not style.text_first_contract or draws_structured_content:
            _add_panel(slide, bounds, style, component_id=slot.get("component_id", "card"))
    if component_id in {"creative_section_tab", "oversized_section_number"}:
        return
    if slot["slot_id"] == "cards":
        blocks, card_warnings = normalize_card_blocks(blocks, slide_id=_slide_id(source))
        warnings.extend(card_warnings)
        _render_cards(slide, blocks, bounds, style)
    elif slot["slot_id"] in {"index_navigation", "progress_markers"} or component_id == "index_navigation":
        _render_navigation_module(slide, blocks, bounds, style, compact=slot["slot_id"] == "progress_markers")
    elif slot["slot_id"] == "metric_panels":
        _render_kpi_cards(slide, source, bounds, style)
    elif slot.get("component_id") in {"radial_process", "curved_timeline", "connector_style"} or slot["slot_id"] in {"diagram", "process", "timeline"}:
        _render_process_module(slide, blocks, bounds, style, variant=str(slot.get("component_id") or slot["slot_id"]))
    else:
        text = _blocks_to_text(blocks) or _text_for_source(source, source_ref, slot["slot_id"])
        _add_textbox(slide, text, _inset(bounds, 0.16, 0.12), style, "body", "text", bullet=True)


def _render_navigation_module(slide: Any, blocks: list[dict[str, Any]], bounds: dict[str, Any], style: _DeckStyle, *, compact: bool) -> None:
    labels = [_block_title(block, index) for index, block in enumerate(blocks[:4])] or ["Context", "Method", "Evidence", "Decision"]
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    if compact:
        step = w / max(1, len(labels))
        baseline = y + h * 0.48
        for index, label in enumerate(labels):
            cx = x + step * (index + 0.5)
            tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.12), Inches(baseline - 0.12), Inches(0.24), Inches(0.24))
            tick.fill.solid()
            tick.fill.fore_color.rgb = style.color("accent" if index == 0 else "surface", "#D0AF37")
            tick.line.color.rgb = style.color("accent_secondary", "#0F766E")
            tick.line.width = Pt(0.65)
            _add_textbox(slide, str(index + 1), {"x": cx - 0.035, "y": baseline - 0.065, "w": 0.08, "h": 0.1}, style, "footer", "text")
            _add_textbox(slide, label.splitlines()[0], {"x": cx - 0.46, "y": baseline + 0.18, "w": 0.92, "h": 0.18}, style, "footer", "text")
        return
    gap = 0.22
    card_w = (w - gap * (len(labels) - 1)) / max(1, len(labels))
    for index, label in enumerate(labels):
        card_bounds = {"x": x + index * (card_w + gap), "y": y, "w": card_w, "h": h}
        _add_panel(slide, card_bounds, style, component_id="layered_card")
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(card_bounds["x"] + 0.2), Inches(card_bounds["y"] + 0.24), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = style.color("accent" if index % 2 == 0 else "accent_secondary", "#D0AF37")
        badge.line.fill.background()
        _add_textbox(slide, f"{index + 1}", {"x": card_bounds["x"] + 0.31, "y": card_bounds["y"] + 0.325, "w": 0.12, "h": 0.12}, style, "footer", "surface")
        _add_textbox(slide, label, _inset({**card_bounds, "y": card_bounds["y"] + 0.48, "h": max(0.1, card_bounds["h"] - 0.48)}, 0.18, 0.08), style, "body", "text")


def _render_cards(slide: Any, blocks: list[dict[str, Any]], bounds: dict[str, Any], style: _DeckStyle) -> None:
    items = blocks or [{"content": "Card item"}]
    columns = 3 if len(items) >= 3 else max(1, len(items))
    rows = (len(items) + columns - 1) // columns
    gap = 0.18
    card_w = (float(bounds["w"]) - gap * (columns - 1)) / columns
    card_h = (float(bounds["h"]) - gap * (rows - 1)) / max(1, rows)
    for index, block in enumerate(items):
        col = index % columns
        row = index // columns
        card_bounds = {"x": float(bounds["x"]) + col * (card_w + gap), "y": float(bounds["y"]) + row * (card_h + gap), "w": card_w, "h": card_h}
        _add_panel(slide, card_bounds, style, component_id="card")
        _add_textbox(slide, _block_title(block, index), _inset(card_bounds, 0.18, 0.14), style, "body", "text")


def _render_kpi_cards(slide: Any, source: dict[str, Any], bounds: dict[str, Any], style: _DeckStyle) -> None:
    series = (source.get("chart_data") or {}).get("series") or []
    labels = [str(item.get("name") or f"KPI {index + 1}") for index, item in enumerate(series[:3])] or ["KPI", "TREND", "DELTA"]
    gap = 0.16
    card_h = (float(bounds["h"]) - gap * (len(labels) - 1)) / len(labels)
    for index, label in enumerate(labels):
        card_bounds = {"x": bounds["x"], "y": float(bounds["y"]) + index * (card_h + gap), "w": bounds["w"], "h": card_h}
        _add_panel(slide, card_bounds, style, component_id="kpi_card")
        marker = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(float(card_bounds["x"]) + 0.16),
            Inches(float(card_bounds["y"]) + 0.14),
            Inches(0.11),
            Inches(0.11),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = style.color("accent", "#2563EB")
        marker.line.fill.background()
        _add_textbox(slide, label, _inset(card_bounds, 0.32, 0.08), style, "label", "muted_text")


def _render_process_module(slide: Any, blocks: list[dict[str, Any]], bounds: dict[str, Any], style: _DeckStyle, *, variant: str) -> None:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    labels = [_block_title(block, index) for index, block in enumerate(blocks[:5])] or ["Frame", "Review", "Decide", "Learn"]
    if style.text_first_contract:
        if style.archetype_id == "literature_map":
            return
        text = _blocks_to_text(blocks) or "\n".join(f"- {label}" for label in labels)
        _add_textbox(slide, text, _inset(bounds, 0.16, 0.12), style, "body", "text", bullet=True)
        return
    if "radial" in variant:
        center_x = x + w * 0.5
        center_y = y + h * 0.5
        radius = min(w, h) * 0.32
        prev: tuple[float, float] | None = None
        first: tuple[float, float] | None = None
        for index, label in enumerate(labels[:6]):
            angle = (-90 + index * 360 / max(1, len(labels[:6]))) * 3.14159 / 180
            px = center_x + radius * math.cos(angle)
            py = center_y + radius * math.sin(angle)
            if prev is not None and not style.text_first_contract:
                line = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(prev[0]), Inches(prev[1]), Inches(px), Inches(py))
                line.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
                line.line.width = Pt(0.85)
            first = first or (px, py)
            prev = (px, py)
            node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(px - 0.18), Inches(py - 0.18), Inches(0.36), Inches(0.36))
            node.fill.solid()
            node.fill.fore_color.rgb = style.color("surface", "#FFFFFF")
            node.line.color.rgb = style.color("accent", "#D0AF37")
            node.line.width = Pt(1.0)
            _add_textbox(slide, str(index + 1), {"x": px - 0.07, "y": py - 0.09, "w": 0.14, "h": 0.14}, style, "footer", "accent_secondary")
            if style.text_first_contract:
                continue
            _add_textbox(slide, label, {"x": px - 0.52, "y": py + 0.24, "w": 1.05, "h": 0.34}, style, "footer", "text")
        if prev and first and not style.text_first_contract:
            close = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(prev[0]), Inches(prev[1]), Inches(first[0]), Inches(first[1]))
            close.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
            close.line.width = Pt(0.6)
    else:
        step_gap = w / max(1, len(labels))
        baseline = y + h * 0.48
        for index, label in enumerate(labels[:6]):
            cx = x + step_gap * (index + 0.5)
            if index > 0 and not style.text_first_contract:
                connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx - step_gap + 0.26), Inches(baseline), Inches(cx - 0.26), Inches(baseline))
                connector.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
                connector.line.width = Pt(1.0)
            node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.22), Inches(baseline - 0.22), Inches(0.44), Inches(0.44))
            node.fill.solid()
            node.fill.fore_color.rgb = style.color("surface", "#FFFFFF")
            node.line.color.rgb = style.color("accent", "#D0AF37")
            node.line.width = Pt(0.9)
            _add_textbox(slide, label, {"x": cx - 0.55, "y": baseline + 0.32, "w": 1.1, "h": 0.4}, style, "footer", "text")


def _render_table_slot(
    slide: Any,
    source: dict[str, Any],
    slot: dict[str, Any],
    source_ref: str | None,
    style: _DeckStyle,
    warnings: list[dict[str, Any]],
) -> None:
    data = _table_data(source, source_ref)
    data, table_warnings = normalize_table_data(
        data,
        slide_id=_slide_id(source),
        fallback_rows=[[str(block.get("slot", "")), _block_text(block)] for block in source.get("content_blocks") or [] if isinstance(block, dict)],
    )
    warnings.extend(table_warnings)
    rows = data["rows"]
    headers = data["headers"]
    row_count = max(1, len(rows) + 1)
    col_count = max(1, len(headers))
    _add_panel(slide, slot["bounds"], style, fill_token="surface_alt", component_id="table_frame")
    _add_data_module_chrome(slide, slot["bounds"], style, "TABLE")
    table_bounds = _data_module_inner_bounds(slot["bounds"])
    shape = slide.shapes.add_table(row_count, col_count, Inches(table_bounds["x"]), Inches(table_bounds["y"]), Inches(table_bounds["w"]), Inches(table_bounds["h"]))
    table = shape.table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        _set_cell_text(cell, header, style, is_header=True)
    for row_index, row in enumerate(rows, start=1):
        for col in range(col_count):
            _set_cell_text(table.cell(row_index, col), str(row[col] if col < len(row) else ""), style, is_header=False)
    _add_table_rules(slide, table_bounds, row_count, style)


def _render_chart_slot(
    slide: Any,
    source: dict[str, Any],
    slot: dict[str, Any],
    source_ref: str | None,
    style: _DeckStyle,
    warnings: list[dict[str, Any]],
) -> None:
    raw_chart = source.get("chart_data") if source_ref == "chart_data" or source.get("chart_data") else None
    chart_payload, chart_warnings = normalize_chart_data(raw_chart, slide_id=_slide_id(source))
    warnings.extend(chart_warnings)
    _add_panel(slide, slot["bounds"], style, fill_token="surface_alt", component_id="chart_frame")
    _add_data_module_chrome(slide, slot["bounds"], style, "CHART")
    chart_bounds = _data_module_inner_bounds(slot["bounds"])
    chart_data = CategoryChartData()
    chart_data.categories = chart_payload["categories"]
    for series in chart_payload["series"]:
        chart_data.add_series(str(series.get("name") or "Series"), tuple(series.get("values") or [1, 2, 3]))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(chart_bounds["x"]),
        Inches(chart_bounds["y"]),
        Inches(chart_bounds["w"]),
        Inches(chart_bounds["h"]),
        chart_data,
    )
    _add_chart_ornaments(slide, slot["bounds"], style)


def _render_image_slot(slide: Any, source: dict[str, Any], slot: dict[str, Any], style: _DeckStyle) -> None:
    bounds = slot["bounds"]
    component_id = str(slot.get("component_id") or "image_frame")
    slot_id = str(slot.get("slot_id") or "")
    _add_panel(slide, bounds, style, fill_token="surface_alt", component_id=component_id)
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    picture_bounds = _photo_picture_bounds(bounds, component_id, style)
    if _should_insert_photo_placeholder(slot_id, component_id, style):
        image_path = _photo_placeholder_path(slot_id, component_id, style)
        if image_path is not None:
            slide.shapes.add_picture(
                str(image_path),
                Inches(picture_bounds["x"]),
                Inches(picture_bounds["y"]),
                width=Inches(picture_bounds["w"]),
                height=Inches(picture_bounds["h"]),
            )
    _draw_photo_frame_fallback(slide, bounds, style)
    if component_id in {"diagonal_photo_panel", "diagonal_hero_panel", "fractured_geometry_panel", "diagonal_image_frame"} and not style.text_first_contract:
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(x + w * 0.58), Inches(y), Inches(w * 0.3), Inches(h))
        band.fill.solid()
        band.fill.fore_color.rgb = style.color("surface", "#FFFFFF") if not style.dark_hero else style.color("surface_alt", "#102E54")
        band.line.color.rgb = style.color("accent_secondary", "#0F766E")
        band.line.width = Pt(0.8)
        for idx in range(3):
            diagonal = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + w * (0.1 + idx * 0.15)),
                Inches(y + h - 0.1),
                Inches(x + w * (0.68 + idx * 0.08)),
                Inches(y + 0.12),
            )
            diagonal.line.color.rgb = style.color("accent_secondary", "#0FA3A3")
            diagonal.line.width = Pt(0.45)


def _should_insert_photo_placeholder(slot_id: str, component_id: str, style: _DeckStyle) -> bool:
    if style.deck_scale == "template_preview":
        return False
    if slot_id in {"hero_image", "section_image", "photo_grid", "photo_frame", "image_frame"}:
        return True
    if component_id in {"diagonal_photo_panel", "diagonal_hero_panel", "fractured_geometry_panel", "diagonal_image_frame"}:
        return True
    return style.archetype_id in {"creative_cover", "section_divider", "photo_caption_grid"} and "image" in slot_id


def _photo_placeholder_path(slot_id: str, component_id: str, style: _DeckStyle) -> Path | None:
    manifest_path = DEFAULT_PHOTO_PLACEHOLDER_MANIFEST
    if not manifest_path.exists():
        return None
    try:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    key = "field_grid"
    if style.archetype_id == "creative_cover" or component_id in {"diagonal_hero_panel", "diagonal_photo_panel"}:
        key = "architecture_neutral"
    elif style.archetype_id == "section_divider" or component_id == "fractured_geometry_panel" or slot_id == "section_image":
        key = "section_geometry_neutral"
    elif style.archetype_id == "photo_caption_grid" or slot_id == "photo_grid":
        key = "field_grid"
    for item in manifest.get("placeholders") or []:
        if isinstance(item, dict) and item.get("id") == key and isinstance(item.get("path"), str):
            path = Path(item["path"])
            if not path.is_absolute():
                path = Path.cwd() / path
            return path if path.exists() else None
    return None


def _photo_picture_bounds(bounds: dict[str, Any], component_id: str, style: _DeckStyle) -> dict[str, float]:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    inset_x = min(0.16, w * 0.06)
    inset_y = min(0.16, h * 0.06)
    if component_id in {"diagonal_photo_panel", "diagonal_hero_panel", "fractured_geometry_panel", "diagonal_image_frame"}:
        inset_x = min(0.22, w * 0.08)
        inset_y = min(0.20, h * 0.07)
    if style.image_frame_style == "medallion_topology_frame":
        inset_x = min(0.24, w * 0.08)
        inset_y = min(0.24, h * 0.08)
    return {
        "x": x + inset_x,
        "y": y + inset_y,
        "w": max(0.1, w - inset_x * 2),
        "h": max(0.1, h - inset_y * 2),
    }


def _draw_photo_frame_fallback(slide: Any, bounds: dict[str, Any], style: _DeckStyle) -> None:
    if style.text_first_contract:
        return
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    cx = x + w * 0.5
    cy = y + h * 0.5
    radius = min(w, h) * 0.34
    for index, scale in enumerate((1.0, 0.72, 0.48)):
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - radius * scale), Inches(cy - radius * scale), Inches(radius * scale * 2), Inches(radius * scale * 2))
        ring.fill.background()
        ring.line.color.rgb = style.color("accent" if index == 0 else "accent_secondary", "#D0AF37")
        ring.line.width = Pt(1.0 if index == 0 else 0.55)
    for angle_idx, (dx, dy) in enumerate(((-0.48, -0.12), (-0.18, 0.32), (0.22, -0.28), (0.54, 0.18))):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx + dx * radius - 0.035), Inches(cy + dy * radius - 0.035), Inches(0.07), Inches(0.07))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent" if angle_idx % 2 else "accent_secondary", "#0FA3A3")
        dot.line.fill.background()


def _render_shape_slot(slide: Any, slot: dict[str, Any], style: _DeckStyle) -> None:
    if slot.get("component_id") == "oversized_section_number":
        _render_oversized_number_slot(slide, {"section_id": "01"}, slot["bounds"], style)
    elif slot.get("component_id") == "creative_section_tab":
        _render_academic_seal(slide, slot["bounds"], style)
    elif slot.get("component_id") == "section_marker":
        _add_section_band(slide, slot["bounds"], style)
    else:
        _add_panel(slide, slot["bounds"], style, fill_token="accent")


def _render_oversized_number_slot(slide: Any, source: dict[str, Any], bounds: dict[str, Any], style: _DeckStyle) -> None:
    raw = str(source.get("section_id") or "01")
    number = "".join(ch for ch in raw if ch.isdigit())[-2:] or "01"
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    _add_textbox(slide, number.zfill(2), {"x": x, "y": y - 0.1, "w": w, "h": max(h, 0.9)}, style, "title", "accent")
    rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + w + 0.08), Inches(y + h * 0.54), Inches(x + w + 0.78), Inches(y + h * 0.54))
    rule.line.color.rgb = style.color("accent", "#D0AF37")
    rule.line.width = Pt(1.15)


def _render_academic_seal(slide: Any, bounds: dict[str, Any], style: _DeckStyle) -> None:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    size = max(0.42, min(w, h, 0.72))
    cx = x + w * 0.5
    cy = y + h * 0.5
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - size / 2), Inches(cy - size / 2), Inches(size), Inches(size))
    outer.fill.background()
    outer.line.color.rgb = style.color("accent", "#D0AF37")
    outer.line.width = Pt(1.0)
    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - size * 0.31), Inches(cy - size * 0.31), Inches(size * 0.62), Inches(size * 0.62))
    inner.fill.solid()
    inner.fill.fore_color.rgb = style.color("surface", "#FFFFFF")
    inner.line.color.rgb = style.color("line", "#D6B666")
    inner.line.width = Pt(0.55)
    _add_textbox(slide, "A", {"x": cx - 0.08, "y": cy - 0.1, "w": 0.18, "h": 0.18}, style, "label", "accent_secondary")


def _add_panel(slide: Any, bounds: dict[str, Any], style: _DeckStyle, fill_token: str = "surface", component_id: str = "card") -> Any:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    chrome_components = {"card", "kpi_card", "chart_frame", "table_frame", "layered_card", "premium_kpi_card", "chart_module", "thin_grid_table", "glass_overlay_card", "diagonal_hero_panel", "fractured_geometry_panel"}
    text_first_simple = (
        style.text_first_contract
        and component_id in {"card", "kpi_card", "layered_card", "premium_kpi_card", "glass_overlay_card"}
    )
    simplify_chrome = style.card_chrome_simplified and component_id in chrome_components
    if component_id in chrome_components and not text_first_simple and not simplify_chrome:
        shadow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x + 0.045),
            Inches(y + 0.055),
            Inches(w),
            Inches(h),
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = style.color("grid", "#E2E8F0")
        shadow.line.fill.background()
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(float(bounds["x"])), Inches(float(bounds["y"])), Inches(float(bounds["w"])), Inches(float(bounds["h"])))
    shape.fill.solid()
    fill_choice = fill_token
    if component_id in {"glass_overlay_card", "diagonal_hero_panel", "fractured_geometry_panel"}:
        fill_choice = "surface_alt"
    elif style.text_first_contract and style.layout_family_id == "evidence_overview" and component_id in {"card", "cards", "layered_card", "radial_process"}:
        fill_choice = "grid"
    elif style.card_style == "evidence_card" and component_id in {"card", "cards", "layered_card", "kpi_card", "premium_kpi_card"}:
        fill_choice = "surface_alt"
    shape.fill.fore_color.rgb = style.color(fill_choice, "#FFFFFF")
    shape.line.color.rgb = style.color("line", "#CBD5E1")
    line_width = 0.35 if simplify_chrome else 0.65 if style.card_style == "evidence_card" else 0.85 if style.card_style == "decision_panel" else 1.05
    shape.line.width = Pt(line_width)
    if component_id in chrome_components and not text_first_simple and not simplify_chrome:
        rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(min(0.08, w)), Inches(h))
        rail.fill.solid()
        rail.fill.fore_color.rgb = style.color("accent_secondary" if component_id in {"card", "layered_card"} else "accent", "#2563EB")
        rail.line.fill.background()
        header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.18), Inches(y + 0.12), Inches(max(0.1, w - 0.36)), Inches(0.035))
        header.fill.solid()
        header.fill.fore_color.rgb = style.color("accent" if component_id in {"chart_frame", "table_frame", "chart_module", "thin_grid_table"} else "grid", "#D0AF37")
        header.line.fill.background()
    if style.card_style == "layered_editorial_card" and component_id in {"card", "layered_card", "kpi_card", "premium_kpi_card", "glass_overlay_card"} and not text_first_simple:
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + w - min(0.14, w)), Inches(y), Inches(min(0.14, w)), Inches(min(0.5, h)))
        accent.fill.solid()
        accent.fill.fore_color.rgb = style.color("accent", "#D49A2A")
        accent.line.fill.background()
    if style.card_style == "evidence_card" and component_id in {"card", "cards", "layered_card"} and not text_first_simple:
        source_chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + w - min(0.72, w * 0.28)), Inches(y + 0.1), Inches(min(0.56, w * 0.22)), Inches(0.16))
        source_chip.fill.solid()
        source_chip.fill.fore_color.rgb = style.color("accent_secondary", "#1E4A5F")
        source_chip.line.fill.background()
        for marker_index in range(3):
            marker = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(x + 0.18 + marker_index * 0.16),
                Inches(y + h - 0.24),
                Inches(0.055),
                Inches(0.055),
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = style.color("accent", "#D0AF37")
            marker.line.fill.background()
    if component_id in {"card", "kpi_card", "layered_card", "premium_kpi_card", "glass_overlay_card"} and not simplify_chrome and not text_first_simple:
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.16), Inches(y + 0.16), Inches(x + w - 0.16), Inches(y + 0.16))
        rule.line.color.rgb = style.color("grid", "#E2E8F0")
        rule.line.width = Pt(0.35 if style.card_style == "evidence_card" else 0.5 if style.card_style == "decision_panel" else 0.8)
        seal = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.15), Inches(y + 0.18), Inches(0.16), Inches(0.16))
        seal.fill.solid()
        seal.fill.fore_color.rgb = style.color("surface_alt", "#F6EED8")
        seal.line.color.rgb = style.color("accent", "#D0AF37")
        seal.line.width = Pt(0.5)
    return shape


def _add_section_band(slide: Any, bounds: dict[str, Any], style: _DeckStyle) -> None:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    band.fill.solid()
    band.fill.fore_color.rgb = style.color("accent_secondary" if style.section_style == "restrained_navy_band" else "accent", "#111827")
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(min(0.16, w)), Inches(h))
    accent.fill.solid()
    accent.fill.fore_color.rgb = style.color("accent" if style.section_style != "oversized_creative_marker" else "text", "#2563EB")
    accent.line.fill.background()


def _add_chart_ornaments(slide: Any, bounds: dict[str, Any], style: _DeckStyle) -> None:
    if style.text_first_contract:
        return
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    for index in range(1, 4):
        rule_y = y + h * index / 4
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.2), Inches(rule_y), Inches(x + w - 0.2), Inches(rule_y))
        rule.line.color.rgb = style.color("grid", "#E2E8F0")
        rule.line.width = Pt(0.2 if style.chart_table_style == "dense_academic" else 0.35 if style.chart_table_style == "high_contrast_dashboard" else 0.5)
    for index in range(4):
        chip_x = x + w - 1.45 + index * 0.28
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(chip_x), Inches(y + 0.16), Inches(0.11), Inches(0.11))
        chip.fill.solid()
        chip.fill.fore_color.rgb = style.color("accent" if index % 2 == 0 else "accent_secondary", "#D0AF37")
        chip.line.fill.background()


def _add_table_rules(slide: Any, bounds: dict[str, Any], row_count: int, style: _DeckStyle) -> None:
    if style.text_first_contract:
        return
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    for index in range(1, min(row_count + 1, 8)):
        rule_y = y + h * index / max(1, row_count)
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(rule_y), Inches(x + w), Inches(rule_y))
        rule.line.color.rgb = style.color("grid", "#E2E8F0")
        rule.line.width = Pt(0.2 if style.chart_table_style == "dense_academic" else 0.35 if style.chart_table_style == "high_contrast_dashboard" else 0.5)
    for index in range(1, 5):
        rule_x = x + w * index / 5
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(rule_x), Inches(y), Inches(rule_x), Inches(y + h))
        rule.line.color.rgb = style.color("grid", "#E2E8F0")
        rule.line.width = Pt(0.18 if style.chart_table_style == "dense_academic" else 0.28)


def _add_data_module_chrome(slide: Any, bounds: dict[str, Any], style: _DeckStyle, label: str) -> None:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    if style.card_chrome_simplified or style.text_first_contract:
        header_h = min(0.22, max(0.16, h * 0.08))
        header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.08), Inches(y + 0.08), Inches(max(0.2, w - 0.16)), Inches(header_h))
        header.fill.solid()
        header.fill.fore_color.rgb = style.color("surface_alt", "#F8F5ED")
        header.line.fill.background()
        _add_textbox(slide, label, {"x": x + 0.22, "y": y + 0.115, "w": 1.05, "h": 0.11}, style, "footer", "muted_text")
        return
    header_h = min(0.36, max(0.24, h * 0.12))
    footer_h = min(0.3, max(0.18, h * 0.09))
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.08), Inches(y + 0.08), Inches(max(0.2, w - 0.16)), Inches(header_h))
    header.fill.solid()
    header.fill.fore_color.rgb = style.color("accent_secondary" if style.chart_table_style != "dense_academic" else "surface_alt", "#0FA3A3")
    header.line.fill.background()
    _add_textbox(slide, label, {"x": x + 0.22, "y": y + 0.14, "w": 1.05, "h": 0.13}, style, "footer", "surface" if style.chart_table_style != "dense_academic" else "text")
    source = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.08), Inches(y + h - footer_h - 0.08), Inches(max(0.2, w - 0.16)), Inches(footer_h))
    source.fill.solid()
    source.fill.fore_color.rgb = style.color("surface_alt", "#F8F5ED")
    source.line.color.rgb = style.color("line", "#C9B98E")
    source.line.width = Pt(0.35)
    _add_textbox(slide, "evidence source strip", {"x": x + 0.22, "y": y + h - footer_h - 0.03, "w": min(2.2, w - 0.4), "h": 0.12}, style, "footer", "muted_text")
    for index in range(5):
        marker_x = x + w - 1.6 + index * 0.24
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(marker_x), Inches(y + h - footer_h + 0.01), Inches(0.1), Inches(0.1))
        marker.fill.solid()
        marker.fill.fore_color.rgb = style.color("accent" if index % 2 == 0 else "accent_secondary", "#D0AF37")
        marker.line.fill.background()
    for index in range(3):
        rail_x = x + 0.18 + index * max(0.7, w / 5)
        rail = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(rail_x), Inches(y + header_h + 0.2), Inches(rail_x + w * 0.18), Inches(y + h - footer_h - 0.14))
        rail.line.color.rgb = style.color("grid", "#E2E8F0")
        rail.line.width = Pt(0.22)


def _data_module_inner_bounds(bounds: dict[str, Any]) -> dict[str, float]:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    top = min(0.48, max(0.34, h * 0.14))
    bottom = min(0.38, max(0.24, h * 0.11))
    return {
        "x": x + 0.18,
        "y": y + top,
        "w": max(0.5, w - 0.36),
        "h": max(0.5, h - top - bottom),
    }


def _add_textbox(
    slide: Any,
    text: str,
    bounds: dict[str, Any],
    style: _DeckStyle,
    typography_token: str,
    color_token: str,
    *,
    bullet: bool = False,
) -> Any:
    box = slide.shapes.add_textbox(Inches(float(bounds["x"])), Inches(float(bounds["y"])), Inches(float(bounds["w"])), Inches(float(bounds["h"])))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    lines = str(text).splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.level = 0
        if bullet and index > 0:
            line = line.lstrip("- ")
        run = paragraph.add_run()
        run.text = line
        run.font.name = style.font_family(typography_token)
        run.font.size = style.font_size(typography_token, 13)
        run.font.color.rgb = style.color(color_token, "#111827")
        run.font.bold = typography_token in {"title", "kpi"}
    return box


def _set_cell_text(cell: Any, text: str, style: _DeckStyle, *, is_header: bool) -> None:
    cell.text = str(text)
    cell.fill.solid()
    if is_header and style.chart_table_style in {"dense_academic", "expressive_modular", "high_contrast_dashboard"}:
        cell.fill.fore_color.rgb = style.color("accent_secondary", "#0F766E")
    elif style.chart_table_style in {"dense_academic", "expressive_modular"}:
        cell.fill.fore_color.rgb = style.color("grid", "#E8DEC6")
    else:
        cell.fill.fore_color.rgb = style.color("surface", "#FFFFFF")
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = style.font_family("label" if is_header else "body")
            run.font.size = style.font_size("label" if is_header else "body", 10)
            run.font.bold = is_header
            run.font.color.rgb = style.color("surface" if is_header and style.chart_table_style in {"dense_academic", "expressive_modular", "high_contrast_dashboard"} else "text", "#111827")


def _add_notes(slide: Any, source: dict[str, Any]) -> None:
    notes_text = source.get("speaker_notes")
    if not notes_text:
        return
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = str(notes_text)
    except Exception:
        return


def _preserve_semantic_source_in_notes(slide: Any, source: dict[str, Any]) -> list[str]:
    visible_text = "\n".join(
        str(shape.text or "")
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_text = str(notes_frame.text or "")
    except Exception:
        return []

    semantic_items: list[tuple[str, str]] = []
    subtitle = str(source.get("subtitle") or "").strip()
    if subtitle:
        semantic_items.append(("subtitle", subtitle))
    for index, block in enumerate(source.get("content_blocks") or [], start=1):
        if not isinstance(block, dict):
            continue
        slot_id = str(block.get("slot") or block.get("block_id") or f"content_{index}")
        text = _canonical_block_text(block.get("content"))
        if text:
            semantic_items.append((slot_id, text))

    searchable = _normalize_semantic_text(f"{visible_text}\n{notes_text}")
    missing: list[tuple[str, str]] = []
    for slot_id, value in semantic_items:
        if _normalize_semantic_text(value) not in searchable:
            missing.append((slot_id, value))
    if not missing:
        return []

    appendix = ["Canonical semantic source:"]
    for slot_id, value in missing:
        appendix.append(f"[{slot_id}]")
        appendix.append(value)
    preserved = "\n".join(appendix)
    notes_frame.text = f"{notes_text.rstrip()}\n\n{preserved}".strip()
    return [slot_id for slot_id, _value in missing]


def _canonical_block_text(content: Any) -> str:
    if isinstance(content, list):
        return "\n".join(str(item) for item in content if str(item).strip())
    if isinstance(content, dict):
        return "\n".join(f"{key}: {value}" for key, value in content.items())
    return str(content or "").strip()


def _normalize_semantic_text(value: Any) -> str:
    return " ".join(str(value or "").split()).casefold()


def _text_for_source(source: dict[str, Any], source_ref: str | None, slot_id: str) -> str:
    if source_ref == "title" or slot_id == "title":
        return str(source.get("title") or "")
    if source_ref == "subtitle" or slot_id == "subtitle":
        return str(source.get("subtitle") or "")
    if source_ref == "citations" or slot_id == "footer":
        citations = source.get("citations") or []
        if citations:
            return " | ".join(str(item.get("label") or item.get("source") or "") for item in citations if isinstance(item, dict))
        return str(source.get("section_id") or "")
    if source_ref and source_ref.startswith("content_blocks."):
        block_id = source_ref.split(".", 1)[1]
        for block in source.get("content_blocks") or []:
            if isinstance(block, dict) and str(block.get("block_id")) == block_id:
                return _block_text(block)
    return _blocks_to_text(source.get("content_blocks") or [])


def _content_blocks_for_slot(source: dict[str, Any], slot_id: str, source_ref: str | None) -> list[dict[str, Any]]:
    blocks = [block for block in source.get("content_blocks") or [] if isinstance(block, dict)]
    if source_ref and source_ref.startswith("content_blocks."):
        block_id = source_ref.split(".", 1)[1]
        return [block for block in blocks if str(block.get("block_id")) == block_id]
    matching = [block for block in blocks if _normalize_key(block.get("slot")) == _normalize_key(slot_id)]
    if source.get("strict_slot_content"):
        return matching
    return matching or blocks


def _blocks_to_text(blocks: list[dict[str, Any]]) -> str:
    return "\n".join(_block_text(block) for block in blocks if isinstance(block, dict)).strip()


def _block_text(block: dict[str, Any]) -> str:
    content = block.get("content")
    if isinstance(content, list):
        return "\n".join(f"- {item}" for item in content)
    if isinstance(content, dict):
        return "\n".join(f"{_display_key(key)}: {value}" for key, value in content.items())
    return str(content or "")


def _block_title(block: dict[str, Any], index: int) -> str:
    raw_title = block.get("title") or block.get("label")
    title = _display_key(raw_title) if raw_title else _fallback_block_title(block.get("slot"), index)
    body = _block_text(block)
    return f"{title}\n{body}" if body and str(title) not in body else str(title)


def _table_data(source: dict[str, Any], source_ref: str | None) -> dict[str, Any]:
    data = source.get("table_data") if source_ref == "table_data" or source.get("table_data") else None
    if isinstance(data, dict):
        rows = data.get("rows") or []
        headers = data.get("headers")
        if headers is None and rows:
            headers = [f"Col {index + 1}" for index in range(len(rows[0]))]
        return {"headers": headers or ["A", "B"], "rows": rows or [["", ""]]}
    blocks = source.get("content_blocks") or []
    rows = [[_fallback_block_title(block.get("slot"), index), _block_text(block)] for index, block in enumerate(blocks) if isinstance(block, dict)]
    return {"headers": ["Item", "Detail"], "rows": rows or [["Item", "Detail"]]}


RAW_SEMANTIC_DISPLAY_KEYS = {
    "cards": "Evidence modules",
    "timeline_steps": "Milestones",
    "source_notes": "Source context",
    "gap_visual": "Evidence gap",
    "process_visual": "Process logic",
    "diagram": "System logic",
    "method_steps": "Method steps",
    "side_notes": "Context notes",
    "caption_cards": "Evidence captions",
    "photo_grid": "Field visuals",
    "index_navigation": "Navigation",
    "progress_markers": "Progress markers",
    "metric_panels": "Key signals",
    "data_table": "Evidence table",
    "chart": "Evidence chart",
    "table": "Evidence table",
}


def _display_key(value: Any) -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    normalized = _normalize_key(text)
    if normalized in RAW_SEMANTIC_DISPLAY_KEYS:
        return RAW_SEMANTIC_DISPLAY_KEYS[normalized]
    if text == normalized or "_" in text:
        return text.replace("_", " ").replace("-", " ").title()
    return text


def _fallback_block_title(slot: Any, index: int) -> str:
    display = _display_key(slot)
    if display:
        return display
    return f"Insight {index + 1}"


def _inset(bounds: dict[str, Any], dx: float, dy: float) -> dict[str, float]:
    return {
        "x": float(bounds["x"]) + dx,
        "y": float(bounds["y"]) + dy,
        "w": max(0.1, float(bounds["w"]) - dx * 2),
        "h": max(0.1, float(bounds["h"]) - dy * 2),
    }


def _normalize_slides(slide_blueprints: dict[str, Any] | list[dict[str, Any]]) -> list[dict[str, Any]]:
    if isinstance(slide_blueprints, list):
        return [slide for slide in slide_blueprints if isinstance(slide, dict)]
    if isinstance(slide_blueprints, dict) and isinstance(slide_blueprints.get("slides"), list):
        return [slide for slide in slide_blueprints["slides"] if isinstance(slide, dict)]
    if isinstance(slide_blueprints, dict) and isinstance(slide_blueprints.get("slide_blueprints"), list):
        return [slide for slide in slide_blueprints["slide_blueprints"] if isinstance(slide, dict)]
    if isinstance(slide_blueprints, dict):
        return [slide_blueprints]
    raise ValueError("slide_blueprint must be an object or array")


def _load_json(path: str | Path) -> Any:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _assembly_plan_matches_template_source(assembly_plan: dict[str, Any], template_spec_source: dict[str, Any]) -> bool:
    source = assembly_plan.get("template_spec_source")
    if not isinstance(source, dict):
        return False
    return _same_path(source.get("path"), template_spec_source.get("path")) and source.get("selection") == template_spec_source.get("selection")


def _normalize_key(value: Any) -> str:
    return str(value or "").strip().lower().replace("-", "_").replace(" ", "_")


def _hex_to_rgb(value: str) -> RGBColor:
    text = value.strip().lstrip("#")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _display_path(path: Path) -> str:
    return str(path.as_posix())


def _same_path(left: Any, right: Any) -> bool:
    if not left or not right:
        return False
    try:
        return Path(str(left)).resolve() == Path(str(right)).resolve()
    except OSError:
        return str(left).replace("\\", "/") == str(right).replace("\\", "/")


def _slide_id(source: dict[str, Any]) -> str:
    return str(source.get("slide_id") or source.get("id") or "slide")


def _render_warning(code: str, source: dict[str, Any], message: str) -> dict[str, Any]:
    return {"code": code, "slide_id": _slide_id(source), "severity": "warning", "message": message}


if __name__ == "__main__":
    raise SystemExit(main())
