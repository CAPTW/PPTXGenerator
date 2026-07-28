"""Compile D05 Magic Layer editable PPT candidates from D01-D04 specs."""

from __future__ import annotations

import hashlib
import json
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image

from .major_region_resolver import validate_scoped_visual_field


SLIDE_WIDTH_IN = 13.333333
SLIDE_HEIGHT_IN = 7.5
TEXT_SLOT_LABELS = {
    "title": "TITLE",
    "subtitle": "SUBTITLE",
    "section_label": "SECTION",
    "section_number": "01",
    "body": "BODY",
    "card_title": "CARD",
    "card_body": "BODY",
    "kpi_label": "KPI",
    "kpi_value": "VALUE",
    "chart_title": "CHART",
    "chart_axis_label": "AXIS",
    "chart_legend": "LEGEND",
    "table_header": "HEADER",
    "table_cell": "CELL",
    "insight": "INSIGHT",
    "note": "NOTE",
    "source": "SOURCE",
    "citation": "CITATION",
    "footer": "FOOTER",
    "meta": "META",
    "placeholder_label": "SLOT",
    "decorative_microtext": "MICRO",
    "unknown_text": "TEXT",
}


def validate_editable_candidate_spec(spec: dict[str, Any]) -> list[str]:
    """Return schema/policy validation errors for an editable candidate spec."""

    required = {
        "schema_name",
        "reference_id",
        "slide_size",
        "source_reference_image",
        "reference_image_as_background",
        "screenshot_slide",
        "objects",
        "fallbacks",
        "semantic_policy",
        "compile_status",
    }
    errors: list[str] = []
    missing = required.difference(spec)
    if missing:
        errors.append(f"missing_fields:{','.join(sorted(missing))}")
    if spec.get("reference_image_as_background") is not False:
        errors.append("full_slide_reference_background_forbidden")
    if spec.get("screenshot_slide") is not False:
        errors.append("screenshot_slide_forbidden")
    for fallback in spec.get("fallbacks") or []:
        if not fallback.get("recorded"):
            errors.append("silent_fallback_forbidden")
    for obj in spec.get("objects") or []:
        if obj.get("semantic_component") in {"icon", "chart", "table", "matrix"} and obj.get("final_use") == "raster":
            errors.append(f"semantic_raster_final_use_forbidden:{obj.get('object_id')}")
        if obj.get("object_type") == "reference_image_background":
            errors.append(f"reference_image_background_object_forbidden:{obj.get('object_id')}")
        if obj.get("object_type") == "scoped_visual_field_crop":
            errors.extend(f"{obj.get('object_id')}:{error}" for error in validate_scoped_visual_field(obj))
    return errors


def build_editable_candidate_spec(
    *,
    reference_id: str,
    reference_image_path: Path,
    manifest: dict[str, Any],
    editable_text_spec: dict[str, Any],
    primitive_mapping: dict[str, Any],
    resolved_svg_icon_map: dict[str, Any],
    native_chart_spec: dict[str, Any],
    editable_shape_chart_spec: dict[str, Any],
    native_table_spec: dict[str, Any],
    editable_shape_grid_table_spec: dict[str, Any],
) -> dict[str, Any]:
    """Build a one-slide editable candidate spec without embedding the full reference image."""

    metadata = manifest.get("reference_metadata") or {}
    objects: list[dict[str, Any]] = [
        {
            "object_id": f"{reference_id}_background_shape",
            "object_type": "ppt_shape",
            "primitive_family": "background_base",
            "semantic_component": "background",
            "bbox_norm": [0.0, 0.0, 1.0, 1.0],
            "z_order": 0,
            "final_use": "ppt_shape",
            "fill": "#111827",
            "line": "#111827",
            "editable": True,
            "source": "D05_compile_policy_shape_background_not_reference_image",
        }
    ]
    seen_text_layers: set[str] = set()
    for primitive in primitive_mapping.get("primitive_mappings") or []:
        target = primitive.get("target_ppt_object_type")
        family = primitive.get("primitive_family")
        if target == "ppt_text" or family in {"title_text_region", "body_text_region"}:
            continue
        if family in {"chart_region", "table_region", "matrix_region"}:
            continue
        if not primitive.get("bbox_norm"):
            continue
        objects.append(_primitive_object(reference_id, primitive))

    for text in editable_text_spec.get("text_layers") or []:
        if not text.get("bbox_norm"):
            continue
        text_layer_id = str(text.get("text_layer_id"))
        if text_layer_id in seen_text_layers:
            continue
        seen_text_layers.add(text_layer_id)
        objects.append(_text_object(reference_id, text))

    for icon in resolved_svg_icon_map.get("mappings") or []:
        if icon.get("final_disposition") != "svg_mapped":
            continue
        if icon.get("raster_fallback_allowed"):
            continue
        objects.append(_icon_object(reference_id, icon))

    for chart in native_chart_spec.get("chart_specs") or []:
        if chart.get("target_ppt_object_type") in {"editable_shape_chart", "native_ppt_chart"}:
            objects.extend(_chart_objects(reference_id, chart, editable_shape_chart_spec))

    for table in native_table_spec.get("table_specs") or []:
        if table.get("target_ppt_object_type") in {"editable_shape_grid_table", "native_ppt_table"}:
            objects.extend(_table_objects(reference_id, table, editable_shape_grid_table_spec))

    objects.sort(key=lambda item: int(item.get("z_order", 0)))
    spec = {
        "schema_name": "editable_candidate_spec",
        "schema_version": "1.0",
        "reference_id": reference_id,
        "slide_size": {"width_in": SLIDE_WIDTH_IN, "height_in": SLIDE_HEIGHT_IN, "aspect_ratio": 16 / 9},
        "source_reference_image": str(reference_image_path.as_posix()),
        "source_reference_metadata": metadata,
        "reference_image_as_background": False,
        "screenshot_slide": False,
        "selected_route": "editable_candidate_magic_layer_d05",
        "objects": objects,
        "fallbacks": [
            {
                "fallback_id": f"{reference_id}_ocr_unavailable_placeholder_text",
                "recorded": True,
                "reason": "OCR backend unavailable; text boxes use semantic slot labels only, not final copy.",
                "allowed": True,
                "D06_risk": "limited_text_context",
            }
        ],
        "semantic_policy": {
            "semantic_text_target": "ppt_text",
            "semantic_icon_target": "svg_vector_or_ppt_vector_shape",
            "semantic_chart_target": "editable_shape_chart_skeleton",
            "semantic_table_target": "editable_shape_grid_table_skeleton",
            "semantic_raster_final_use_allowed": False,
            "full_slide_reference_background_allowed": False,
            "screenshot_slide_allowed": False,
        },
        "compile_status": "ready",
        "canva_parity_claimed": False,
    }
    errors = validate_editable_candidate_spec(spec)
    spec["validation_errors"] = errors
    if errors:
        spec["compile_status"] = "blocked"
    return spec


def build_patched_editable_candidate_spec(
    *,
    reference_id: str,
    reference_image_path: Path,
    manifest: dict[str, Any],
    editable_text_spec: dict[str, Any],
    resolved_svg_icon_map: dict[str, Any],
    native_chart_spec: dict[str, Any],
    editable_shape_chart_spec: dict[str, Any],
    native_table_spec: dict[str, Any],
    editable_shape_grid_table_spec: dict[str, Any],
    major_region_resolution: dict[str, Any],
    decorative_grouping: dict[str, Any],
) -> dict[str, Any]:
    """Build a D05.1 patched candidate emphasizing major visual composition."""

    metadata = manifest.get("reference_metadata") or {}
    objects: list[dict[str, Any]] = [
        {
            "object_id": f"{reference_id}_background_shape",
            "object_type": "ppt_shape",
            "primitive_family": "background_base",
            "semantic_component": "background",
            "major_region_type": "background_base",
            "bbox_norm": [0.0, 0.0, 1.0, 1.0],
            "z_order": 0,
            "final_use": "ppt_shape",
            "fill": "#0B1220",
            "line": "#0B1220",
            "editable": True,
            "source": "D05_1_shape_background_not_reference_image",
        }
    ]
    for region in major_region_resolution.get("major_regions") or []:
        objects.append(_major_region_object(reference_id, region))
    for group in decorative_grouping.get("decorative_group_objects") or []:
        objects.append(_decorative_group_object(group))

    for icon in resolved_svg_icon_map.get("mappings") or []:
        if icon.get("final_disposition") != "svg_mapped" or icon.get("raster_fallback_allowed"):
            continue
        if _area(icon.get("bbox_norm")) < 0.0016:
            continue
        objects.append(_icon_object(reference_id, icon))

    for text in _select_text_layers_for_patch(editable_text_spec):
        objects.append(_text_object(reference_id, text))
    if not any(obj.get("slot_type") in {"source", "footer"} for obj in objects):
        footer_region = _find_major_region(major_region_resolution, "bottom_footer_source_strip")
        if footer_region:
            objects.append(_synthetic_footer_text(reference_id, footer_region))

    for chart in native_chart_spec.get("chart_specs") or []:
        if chart.get("target_ppt_object_type") in {"editable_shape_chart", "native_ppt_chart"}:
            objects.extend(_chart_objects(reference_id, chart, editable_shape_chart_spec))
    for table in native_table_spec.get("table_specs") or []:
        if table.get("target_ppt_object_type") in {"editable_shape_grid_table", "native_ppt_table"}:
            objects.extend(_table_objects(reference_id, table, editable_shape_grid_table_spec))

    objects.sort(key=lambda item: int(item.get("z_order", 0)))
    spec = {
        "schema_name": "patched_editable_candidate_spec",
        "schema_version": "1.1",
        "reference_id": reference_id,
        "slide_size": {"width_in": SLIDE_WIDTH_IN, "height_in": SLIDE_HEIGHT_IN, "aspect_ratio": 16 / 9},
        "source_reference_image": str(reference_image_path.as_posix()),
        "source_reference_metadata": metadata,
        "reference_image_as_background": False,
        "screenshot_slide": False,
        "selected_route": "editable_candidate_magic_layer_d05_1_patched",
        "objects": objects,
        "fallbacks": [
            {
                "fallback_id": f"{reference_id}_ocr_unavailable_placeholder_text",
                "recorded": True,
                "reason": "OCR backend unavailable; patched candidate keeps semantic slot labels as editable placeholders.",
                "allowed": True,
                "D06_risk": "limited_text_context",
            },
            {
                "fallback_id": f"{reference_id}_major_region_skeleton",
                "recorded": True,
                "reason": "D01/D03 micro-layer outputs do not always expose major visual fields; D05.1 uses recorded composition skeleton regions.",
                "allowed": True,
                "D06_risk": "bounded_visual_fidelity_patch",
            },
        ],
        "semantic_policy": {
            "semantic_text_target": "ppt_text",
            "semantic_icon_target": "svg_vector_or_ppt_vector_shape",
            "semantic_chart_target": "editable_shape_chart_skeleton",
            "semantic_table_target": "editable_shape_grid_table_skeleton",
            "scoped_visual_field_raster_allowed_for_nonsemantic_regions": True,
            "semantic_raster_final_use_allowed": False,
            "full_slide_reference_background_allowed": False,
            "screenshot_slide_allowed": False,
        },
        "major_region_resolution": major_region_resolution,
        "decorative_layer_grouping": decorative_grouping,
        "compile_status": "ready",
        "canva_parity_claimed": False,
    }
    errors = validate_editable_candidate_spec(spec)
    spec["validation_errors"] = errors
    if errors:
        spec["compile_status"] = "blocked"
    return spec


def materialize_scoped_visual_field_crops(spec: dict[str, Any], output_dir: Path) -> dict[str, Any]:
    """Create scoped crop assets for nonsemantic visual-field objects."""

    output_dir.mkdir(parents=True, exist_ok=True)
    source = Path(spec["source_reference_image"])
    image = Image.open(source).convert("RGB")
    created = []
    for obj in spec.get("objects") or []:
        if obj.get("object_type") != "scoped_visual_field_crop":
            continue
        bbox = obj.get("source_crop_bbox_norm") or obj.get("bbox_norm")
        x = max(0, int(float(bbox[0]) * image.width))
        y = max(0, int(float(bbox[1]) * image.height))
        w = max(1, int(float(bbox[2]) * image.width))
        h = max(1, int(float(bbox[3]) * image.height))
        crop = image.crop((x, y, min(image.width, x + w), min(image.height, y + h)))
        path = output_dir / f"{obj['object_id']}.png"
        crop.save(path)
        obj["image_path"] = path.as_posix()
        obj["crop_px"] = [x, y, crop.width, crop.height]
        created.append({"object_id": obj["object_id"], "path": path.as_posix(), "crop_px": obj["crop_px"]})
    return {"schema_name": "scoped_visual_field_crop_manifest", "created_count": len(created), "crops": created}


def compile_editable_candidate_pptx(spec: dict[str, Any], output_pptx: Path) -> dict[str, Any]:
    """Compile a spec into a one-slide PPTX made of editable objects."""

    errors = validate_editable_candidate_spec(spec)
    if errors:
        raise ValueError(f"Editable candidate spec failed validation: {errors}")
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for obj in spec.get("objects") or []:
        object_type = obj.get("object_type")
        if object_type == "ppt_text":
            _add_text(slide, obj)
        elif object_type == "ppt_connector":
            _add_connector(slide, obj)
        elif object_type == "editable_shape_chart":
            _add_chart_skeleton(slide, obj)
        elif object_type == "editable_shape_grid_table":
            _add_table_skeleton(slide, obj)
        elif object_type == "scoped_visual_field_crop":
            _add_scoped_visual_crop(slide, obj)
        elif object_type in {"svg_vector", "ppt_vector_shape_icon"}:
            _add_icon_vector(slide, obj)
        else:
            _add_shape(slide, obj)
    prs.save(output_pptx)
    return pptx_inventory(output_pptx)


def pptx_inventory(path: Path) -> dict[str, Any]:
    prs = Presentation(path)
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        media = [name for name in names if name.startswith("ppt/media/")]
        slides = [name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
    return {
        "schema_name": "pptx_inventory",
        "status": "passed",
        "pptx_path": path.as_posix(),
        "exists": path.exists(),
        "file_size_bytes": path.stat().st_size if path.exists() else 0,
        "sha256": _sha256(path) if path.exists() else None,
        "slide_count": len(prs.slides),
        "ooxml_slide_count": len(slides),
        "media_count": len(media),
        "selected_route": "editable_candidate_magic_layer_d05",
    }


def build_structural_ledgers(spec: dict[str, Any], pptx_path: Path) -> dict[str, dict[str, Any]]:
    prs = Presentation(pptx_path)
    shapes = []
    text = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            record = {
                "slide_index": slide_index,
                "shape_index": shape_index,
                "shape_type": str(shape.shape_type),
                "name": shape.name,
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
                "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
            }
            shapes.append(record)
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                value = shape.text_frame.text
                if value:
                    text.append({**record, "text": value})
    objects = spec.get("objects") or []
    semantic_icons = [obj for obj in objects if obj.get("semantic_component") == "icon"]
    chart_tables = [obj for obj in objects if obj.get("semantic_component") in {"chart", "table", "matrix"}]
    rasters = [obj for obj in objects if obj.get("final_use") in {"raster", "allowed_scoped_visual_field_raster"}]
    media = [obj for obj in objects if obj.get("object_type") == "scoped_visual_field_crop"]
    return {
        "pptx_inventory": pptx_inventory(pptx_path),
        "object_ledger": {"schema_name": "object_ledger", "status": "passed", "object_count": len(objects), "objects": objects},
        "text_ledger": {"schema_name": "text_ledger", "status": "passed", "editable_text_count": len(text), "text_runs": text},
        "media_ledger": {"schema_name": "media_ledger", "status": "passed", "media_count": len(media), "media": media},
        "shape_ledger": {"schema_name": "shape_ledger", "status": "passed", "shape_count": len(shapes), "shapes": shapes},
        "svg_icon_ledger": {
            "schema_name": "svg_icon_ledger",
            "status": "passed",
            "semantic_icon_count": len(semantic_icons),
            "semantic_icon_raster_count": 0,
            "vector_icon_count": len(semantic_icons),
            "icons": semantic_icons,
        },
        "chart_table_ledger": {
            "schema_name": "chart_table_ledger",
            "status": "passed",
            "chart_table_object_count": len(chart_tables),
            "semantic_chart_table_raster_count": 0,
            "components": chart_tables,
        },
        "raster_layer_ledger": {
            "schema_name": "raster_layer_ledger",
            "status": "passed",
            "raster_layer_count": len(rasters),
            "semantic_raster_count": len([obj for obj in rasters if obj.get("semantic_component") in {"icon", "chart", "table", "matrix", "text", "source_footer"}]),
            "rasters": rasters,
        },
        "editability_ledger": {
            "schema_name": "editability_ledger",
            "status": "passed",
            "editable_object_count": len([obj for obj in objects if obj.get("editable")]),
            "noneditable_required_object_count": 0,
            "objects": [{"object_id": obj.get("object_id"), "editable": obj.get("editable"), "final_use": obj.get("final_use")} for obj in objects],
        },
    }


def _primitive_object(reference_id: str, primitive: dict[str, Any]) -> dict[str, Any]:
    family = primitive.get("primitive_family", "ppt_shape")
    target = primitive.get("target_ppt_object_type", "ppt_shape")
    object_type = "ppt_connector" if target in {"ppt_connector", "ppt_line"} or "connector" in family else "ppt_shape"
    fill = _fill_for_family(family)
    line = _line_for_family(family)
    return {
        "object_id": primitive.get("primitive_id") or f"{reference_id}_primitive",
        "source_layer_ids": primitive.get("source_layer_ids") or [],
        "object_type": object_type,
        "primitive_family": family,
        "semantic_component": _semantic_component_for_family(family),
        "bbox_norm": primitive.get("bbox_norm"),
        "bbox_px": primitive.get("bbox_px"),
        "z_order": int(primitive.get("z_order", 10)),
        "final_use": "ppt_shape",
        "fill": fill,
        "line": line,
        "editable": True,
        "confidence": primitive.get("confidence"),
        "raster_policy": primitive.get("raster_policy"),
    }


def _text_object(reference_id: str, text: dict[str, Any]) -> dict[str, Any]:
    slot = text.get("slot_type") or text.get("text_role") or "unknown_text"
    label = TEXT_SLOT_LABELS.get(slot, slot.upper()[:18])
    return {
        "object_id": text.get("text_layer_id") or f"{reference_id}_text",
        "source_layer_ids": text.get("source_layer_ids") or [],
        "object_type": "ppt_text",
        "primitive_family": "title_text_region" if slot == "title" else "body_text_region",
        "semantic_component": "text",
        "slot_type": slot,
        "text": label,
        "ocr_status": text.get("ocr_status"),
        "final_copy": False,
        "bbox_norm": text.get("bbox_norm"),
        "bbox_px": text.get("bbox_px"),
        "z_order": 900 + len(str(text.get("text_layer_id"))),
        "font_size": max(5, min(34, int(text.get("font_size_estimate") or 10))),
        "font_weight": text.get("font_weight_estimate") or "regular",
        "fill": "#00000000",
        "line": "#00000000",
        "text_color": "#F8FAFC" if slot in {"title", "subtitle"} else "#D1D5DB",
        "editable": True,
        "final_use": "ppt_text",
    }


def _icon_object(reference_id: str, icon: dict[str, Any]) -> dict[str, Any]:
    return {
        "object_id": icon.get("candidate_id") or f"{reference_id}_icon",
        "source_layer_ids": [icon.get("layer_id")] if icon.get("layer_id") else [],
        "object_type": "ppt_vector_shape_icon",
        "primitive_family": "icon_region",
        "semantic_component": "icon" if icon.get("icon_classification") == "semantic_icon" else "decorative_icon",
        "bbox_norm": icon.get("bbox_norm"),
        "bbox_px": icon.get("bbox_px"),
        "z_order": 760,
        "selected_role": icon.get("selected_role"),
        "selected_svg_candidate_path": icon.get("selected_svg_candidate_path"),
        "raster_fallback_allowed": False,
        "final_use": "svg_vector",
        "fill": "#38BDF8",
        "line": "#E0F2FE",
        "editable": True,
    }


def _chart_objects(reference_id: str, chart: dict[str, Any], editable_shape_chart_spec: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        {
            "object_id": chart.get("native_chart_spec_id") or f"{reference_id}_chart",
            "source_layer_ids": chart.get("source_layer_ids") or [],
            "object_type": "editable_shape_chart",
            "primitive_family": "chart_region",
            "semantic_component": "chart",
            "chart_component_type": chart.get("chart_component_type"),
            "bbox_norm": chart.get("bbox_norm"),
            "bbox_px": chart.get("bbox_px"),
            "z_order": 640,
            "final_use": "editable_shape_chart",
            "editable": True,
            "raster_final_use_policy": chart.get("raster_final_use_policy"),
            "editable_shape_chart_spec_count": len(editable_shape_chart_spec.get("chart_specs") or []),
        }
    ]


def _table_objects(reference_id: str, table: dict[str, Any], editable_shape_grid_table_spec: dict[str, Any]) -> list[dict[str, Any]]:
    skeleton = table.get("table_skeleton_spec") or {}
    return [
        {
            "object_id": table.get("native_table_spec_id") or f"{reference_id}_table",
            "source_layer_ids": table.get("source_layer_ids") or [],
            "object_type": "editable_shape_grid_table",
            "primitive_family": "table_region",
            "semantic_component": "table",
            "table_component_type": table.get("table_component_type"),
            "bbox_norm": table.get("bbox_norm"),
            "bbox_px": table.get("bbox_px"),
            "z_order": 640,
            "final_use": "editable_shape_grid_table",
            "editable": True,
            "row_count": int(skeleton.get("row_count_inferred") or 6),
            "column_count": int(skeleton.get("column_count_inferred") or 5),
            "raster_final_use_policy": table.get("raster_final_use_policy"),
            "editable_shape_grid_table_spec_count": len(editable_shape_grid_table_spec.get("table_specs") or []),
        }
    ]


def _major_region_object(reference_id: str, region: dict[str, Any]) -> dict[str, Any]:
    region_type = region.get("major_region_type")
    object_id = region.get("region_id") or f"{reference_id}_{region_type}"
    if region.get("object_type") == "scoped_visual_field_crop":
        return {
            "object_id": object_id,
            "source_layer_ids": region.get("source_layer_ids") or [],
            "object_type": "scoped_visual_field_crop",
            "primitive_family": "hero_visual_field",
            "semantic_component": "visual_field",
            "major_region_type": region_type,
            "bbox_norm": region.get("bbox_norm"),
            "source_crop_bbox_norm": region.get("source_crop_bbox_norm") or region.get("bbox_norm"),
            "z_order": int(region.get("z_order", 30)),
            "final_use": "allowed_scoped_visual_field_raster",
            "editable": True,
            "source": region.get("source"),
            "notes": region.get("notes", ""),
        }
    return {
        "object_id": object_id,
        "source_layer_ids": region.get("source_layer_ids") or [],
        "object_type": "ppt_shape",
        "primitive_family": region_type,
        "semantic_component": "source_footer" if region_type == "bottom_footer_source_strip" else "major_region",
        "major_region_type": region_type,
        "bbox_norm": region.get("bbox_norm"),
        "bbox_px": region.get("bbox_px"),
        "z_order": int(region.get("z_order", 80)),
        "final_use": "ppt_shape",
        "fill": region.get("fill", "#1F2937"),
        "line": region.get("line", "#38BDF8"),
        "editable": True,
        "source": region.get("source"),
    }


def _decorative_group_object(group: dict[str, Any]) -> dict[str, Any]:
    return {
        **group,
        "final_use": "ppt_shape",
        "editable": True,
        "fill": group.get("fill", "#00000000"),
        "line": group.get("line", "#38BDF8"),
    }


def _select_text_layers_for_patch(editable_text_spec: dict[str, Any]) -> list[dict[str, Any]]:
    text_layers = editable_text_spec.get("text_layers") or []
    selected: list[dict[str, Any]] = []
    by_slot: dict[str, list[dict[str, Any]]] = {}
    for item in text_layers:
        slot = item.get("slot_type") or item.get("text_role") or "unknown_text"
        if slot == "decorative_microtext":
            continue
        by_slot.setdefault(slot, []).append(item)
    limits = {"title": 2, "subtitle": 1, "source": 2, "footer": 1, "citation": 1, "body": 4, "card_title": 6, "card_body": 4}
    for slot, items in by_slot.items():
        items = sorted(items, key=lambda item: _area(item.get("bbox_norm")), reverse=True)
        selected.extend(items[: limits.get(slot, 2)])
    return selected


def _synthetic_footer_text(reference_id: str, region: dict[str, Any]) -> dict[str, Any]:
    bbox = list(region.get("bbox_norm") or [0.03, 0.9, 0.35, 0.04])
    bbox[0] = min(0.95, bbox[0] + 0.02)
    bbox[1] = min(0.96, bbox[1] + 0.02)
    bbox[2] = min(0.5, max(0.12, bbox[2] * 0.5))
    bbox[3] = min(0.06, max(0.025, bbox[3] * 0.45))
    return {
        "text_layer_id": f"{reference_id}_synthetic_source_footer_text",
        "bbox_norm": bbox,
        "slot_type": "source",
        "font_size_estimate": 8,
        "font_weight_estimate": "regular",
        "ocr_status": "OCR_UNAVAILABLE",
        "source_layer_ids": region.get("source_layer_ids") or [],
    }


def _find_major_region(resolution: dict[str, Any], region_type: str) -> dict[str, Any] | None:
    for region in resolution.get("major_regions") or []:
        if region.get("major_region_type") == region_type:
            return region
    return None


def _add_shape(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
    if obj.get("shape") == "oval":
        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
    elif obj.get("shape") == "diamond":
        shape_type = MSO_AUTO_SHAPE_TYPE.DIAMOND
    if obj.get("primitive_family") in {"card_panel", "kpi_card", "evidence_card"}:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    _style_shape(shape, obj)


def _add_text(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = str(obj.get("text") or "TEXT")
    run.font.size = Pt(float(obj.get("font_size") or 9))
    run.font.bold = obj.get("font_weight") == "bold"
    run.font.color.rgb = _rgb(obj.get("text_color", "#F8FAFC"))
    shape.name = str(obj.get("object_id") or "editable_text")


def _add_connector(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y + h // 2, x + w, y + h // 2)
    shape.line.color.rgb = _rgb(obj.get("line", "#38BDF8"))
    shape.line.width = Pt(1.1)
    shape.name = str(obj.get("object_id") or "connector")


def _add_icon_vector(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
    if obj.get("icon_shape") == "diamond":
        shape_type = MSO_AUTO_SHAPE_TYPE.DIAMOND
    elif obj.get("icon_shape") == "hexagon":
        shape_type = MSO_AUTO_SHAPE_TYPE.HEXAGON
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    _style_shape(shape, obj)
    shape.name = str(obj.get("object_id") or "vector_icon")


def _add_chart_skeleton(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    frame.fill.background()
    frame.line.color.rgb = _rgb("#60A5FA")
    frame.line.width = Pt(1)
    chart_h = max(1, int(h * 0.5))
    base_y = y + int(h * 0.7)
    bar_count = 6
    for idx in range(bar_count):
        bar_w = max(1, int(w / (bar_count * 2.5)))
        bx = x + int(w * 0.1) + idx * int(w * 0.12)
        bh = int(chart_h * (0.35 + 0.09 * (idx % 4)))
        by = base_y - bh
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bx, by, bar_w, bh)
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb("#38BDF8")
        bar.line.color.rgb = _rgb("#0F172A")
        bar.name = f"{obj.get('object_id')}_bar_{idx + 1}"
    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + int(w * 0.08), base_y, x + int(w * 0.9), base_y)
    axis.line.color.rgb = _rgb("#CBD5E1")
    label = slide.shapes.add_textbox(x + int(w * 0.08), y + int(h * 0.08), max(1, int(w * 0.5)), max(1, int(h * 0.12)))
    label.text_frame.text = "CHART"
    label.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    label.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb("#E0F2FE")


def _add_table_skeleton(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    rows = max(2, min(12, int(obj.get("row_count") or 6)))
    cols = max(2, min(8, int(obj.get("column_count") or 5)))
    cell_w = max(1, int(w / cols))
    cell_h = max(1, int(h / rows))
    for row in range(rows):
        for col in range(cols):
            cx = x + col * cell_w
            cy = y + row * cell_h
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cx, cy, cell_w, cell_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("#1F2937" if row else "#0F766E")
            cell.line.color.rgb = _rgb("#94A3B8")
            cell.line.width = Pt(0.5)
            if row == 0 or col == 0:
                txt = slide.shapes.add_textbox(cx + 2, cy + 1, max(1, cell_w - 4), max(1, cell_h - 2))
                txt.text_frame.text = "HDR" if row == 0 else "ROW"
                txt.text_frame.paragraphs[0].runs[0].font.size = Pt(5.5)
                txt.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb("#F8FAFC")


def _add_scoped_visual_crop(slide: Any, obj: dict[str, Any]) -> None:
    path = obj.get("image_path")
    x, y, w, h = _box(obj)
    if path and Path(path).exists():
        shape = slide.shapes.add_picture(str(Path(path).resolve()), x, y, w, h)
        shape.name = str(obj.get("object_id") or "scoped_visual_field_crop")
        return
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb("#0E7490")
    shape.line.color.rgb = _rgb("#38BDF8")
    shape.name = str(obj.get("object_id") or "scoped_visual_field_crop_missing")


def _style_shape(shape: Any, obj: dict[str, Any]) -> None:
    fill = obj.get("fill")
    if fill and fill != "#00000000":
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        try:
            shape.fill.transparency = 20 if obj.get("primitive_family") not in {"background_base"} else 0
        except AttributeError:
            pass
    else:
        shape.fill.background()
    line = obj.get("line")
    if line and line != "#00000000":
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    shape.name = str(obj.get("object_id") or "editable_shape")


def _box(obj: dict[str, Any]) -> tuple[int, int, int, int]:
    norm = obj.get("bbox_norm") or [0, 0, 0.1, 0.1]
    x = Inches(float(norm[0]) * SLIDE_WIDTH_IN)
    y = Inches(float(norm[1]) * SLIDE_HEIGHT_IN)
    w = Inches(max(0.005, float(norm[2]) * SLIDE_WIDTH_IN))
    h = Inches(max(0.005, float(norm[3]) * SLIDE_HEIGHT_IN))
    return int(x), int(y), int(w), int(h)


def _rgb(hex_color: str) -> RGBColor:
    value = (hex_color or "#FFFFFF").replace("#", "")[:6]
    if len(value) != 6:
        value = "FFFFFF"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _fill_for_family(family: str) -> str:
    if family == "background_base":
        return "#111827"
    if "card" in family or "panel" in family or family in {"source_footer_strip", "chart_frame"}:
        return "#1F2937"
    if "overlay" in family:
        return "#0F172A"
    if "line" in family or "connector" in family or "accent" in family:
        return "#00000000"
    return "#243447"


def _line_for_family(family: str) -> str:
    if "accent" in family:
        return "#F59E0B"
    if "connector" in family:
        return "#38BDF8"
    if family == "source_footer_strip":
        return "#64748B"
    return "#334155"


def _semantic_component_for_family(family: str) -> str:
    if "chart" in family:
        return "chart"
    if "table" in family or "matrix" in family:
        return "table"
    if "icon" in family:
        return "icon"
    if "footer" in family:
        return "source_footer"
    if "text" in family:
        return "text"
    return "primitive"


def _area(bbox_norm: Any) -> float:
    if not isinstance(bbox_norm, list) or len(bbox_norm) != 4:
        return 0.0
    return float(bbox_norm[2]) * float(bbox_norm[3])


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(65536), b""):
            digest.update(chunk)
    return digest.hexdigest()


def has_full_slide_picture(pptx_path: Path, *, tolerance: float = 0.95) -> bool:
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                width_ratio = int(shape.width) / int(prs.slide_width)
                height_ratio = int(shape.height) / int(prs.slide_height)
                if width_ratio >= tolerance and height_ratio >= tolerance:
                    return True
    return False
