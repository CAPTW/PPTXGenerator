"""Image policy QA for final editable PPTX decks."""

from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_PPTX_PATH = Path("outputs/final_deck.pptx")
DEFAULT_TEMPLATE_SPEC = Path("outputs/editable_template_spec.final.json")
DEFAULT_ASSEMBLY_PLAN = Path("outputs/deck_assembly_plan.json")
DEFAULT_TEMPLATE_IMAGE_MANIFEST = Path("outputs/template_images/template_image_manifest.json")
DEFAULT_JSON_REPORT = Path("outputs/final_deck_image_policy_report.json")
DEFAULT_MD_REPORT = Path("outputs/final_deck_image_policy_report.md")

EMU_PER_INCH = 914400
FULL_SLIDE_AREA_RATIO_THRESHOLD = 0.88
LARGE_PICTURE_AREA_RATIO_THRESHOLD = 0.62
SVG_ICON_AREA_RATIO_THRESHOLD = 0.02
ALLOWED_IMAGE_SLOT_IDS = {
    "image_frame",
    "hero_image",
    "section_image",
    "case_study_image",
    "photo_frame",
    "photo_grid",
    "diagonal_photo_panel",
    "diagonal_image_frame",
}


def build_final_deck_image_policy_report(
    *,
    pptx_path: str | Path = DEFAULT_PPTX_PATH,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    deck_assembly_plan_path: str | Path = DEFAULT_ASSEMBLY_PLAN,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    full_slide_area_ratio_threshold: float = FULL_SLIDE_AREA_RATIO_THRESHOLD,
    large_picture_area_ratio_threshold: float = LARGE_PICTURE_AREA_RATIO_THRESHOLD,
) -> dict[str, Any]:
    pptx_file = Path(pptx_path)
    template_spec = _load_json(template_spec_path) if Path(template_spec_path).exists() else {}
    assembly_plan = _load_json(deck_assembly_plan_path) if Path(deck_assembly_plan_path).exists() else {}
    template_image_manifest = _load_json(template_image_manifest_path) if Path(template_image_manifest_path).exists() else {}

    deck = Presentation(pptx_file)
    slide_w = deck.slide_width / EMU_PER_INCH
    slide_h = deck.slide_height / EMU_PER_INCH
    slide_area = max(0.01, slide_w * slide_h)
    layouts = {str(layout.get("layout_id")): layout for layout in template_spec.get("layouts") or [] if isinstance(layout, dict)}
    bindings = assembly_plan.get("slide_layout_bindings") or []
    reference_hashes = _template_reference_hashes(template_image_manifest)
    reference_names = _template_reference_names(template_image_manifest)

    slide_reports: list[dict[str, Any]] = []
    findings: list[dict[str, Any]] = []
    counts = {
        "picture_shape_count": 0,
        "full_slide_picture_count": 0,
        "reference_template_image_embedded_count": 0,
        "undeclared_picture_shape_count": 0,
        "allowed_photo_frame_picture_count": 0,
        "allowed_svg_icon_picture_count": 0,
    }

    for slide_index, slide in enumerate(deck.slides, start=1):
        binding = bindings[slide_index - 1] if slide_index - 1 < len(bindings) and isinstance(bindings[slide_index - 1], dict) else {}
        layout_id = str(binding.get("selected_layout_id") or binding.get("layout_id") or "")
        layout = layouts.get(layout_id, {})
        allowed_slots = _allowed_image_slots(layout)
        slide_picture_reports: list[dict[str, Any]] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            counts["picture_shape_count"] += 1
            bounds = _shape_bounds(shape)
            area_ratio = round((bounds["w"] * bounds["h"]) / slide_area, 6)
            matched_slot = _matching_allowed_slot(bounds, allowed_slots)
            image_info = _shape_image_info(shape)
            is_svg_icon = _is_allowed_svg_icon_shape(shape, image_info, area_ratio)
            is_allowed_photo_frame = matched_slot is not None
            if is_allowed_photo_frame:
                counts["allowed_photo_frame_picture_count"] += 1
            elif is_svg_icon:
                counts["allowed_svg_icon_picture_count"] += 1
            else:
                counts["undeclared_picture_shape_count"] += 1

            blob_hash = _shape_image_sha256(shape)
            is_reference = blob_hash in reference_hashes
            if is_reference:
                counts["reference_template_image_embedded_count"] += 1

            is_full_slide = _is_full_slide_picture(bounds, slide_w, slide_h, area_ratio, full_slide_area_ratio_threshold)
            if is_full_slide:
                counts["full_slide_picture_count"] += 1

            picture_report = {
                "shape_index": shape_index,
                "bounds": bounds,
                "area_ratio": area_ratio,
                "matched_slot_id": matched_slot.get("slot_id") if matched_slot else None,
                "is_allowed_photo_frame": is_allowed_photo_frame,
                "is_allowed_svg_icon": is_svg_icon,
                "is_full_slide_picture": is_full_slide,
                "is_template_reference_image": is_reference,
                "image_sha256": blob_hash,
                "image_content_type": image_info.get("content_type"),
                "image_partname": image_info.get("partname"),
            }
            slide_picture_reports.append(picture_report)

            if is_reference:
                findings.append(
                    _finding(
                        "TEMPLATE_REFERENCE_IMAGE_EMBEDDED",
                        "severe",
                        "A template reference PNG from outputs/template_images is embedded in the final deck.",
                        slide_index,
                        picture_report,
                    )
                )
            if is_full_slide:
                findings.append(
                    _finding(
                        "FULL_SLIDE_PICTURE_FORBIDDEN",
                        "severe",
                        "A picture shape covers most of the slide and is treated as a forbidden raster background.",
                        slide_index,
                        picture_report,
                    )
                )
            if not is_allowed_photo_frame and not is_svg_icon:
                findings.append(
                    _finding(
                        "UNDECLARED_PICTURE_SHAPE",
                        "severe",
                        "A picture shape is not contained in a declared image_frame, hero_image, or case_study_image slot.",
                        slide_index,
                        picture_report,
                    )
                )
            elif area_ratio >= large_picture_area_ratio_threshold and matched_slot is None:
                findings.append(
                    _finding(
                        "LARGE_PICTURE_WITHOUT_PHOTO_FRAME",
                        "severe",
                        "A picture covers a large area without an explicit allowed photo-frame slot.",
                        slide_index,
                        picture_report,
                    )
                )

        slide_reports.append(
            {
                "slide_index": slide_index,
                "layout_id": layout_id,
                "allowed_image_slots": [_slot_summary(slot) for slot in allowed_slots],
                "picture_shape_count": len(slide_picture_reports),
                "pictures": slide_picture_reports,
            }
        )

    severe_count = sum(1 for finding in findings if finding["severity"] == "severe")
    report = {
        "schema_name": "final_deck_image_policy_report",
        "schema_version": "1.0",
        "pptx_path": _display_path(pptx_file),
        "template_spec_path": _display_path(Path(template_spec_path)),
        "deck_assembly_plan_path": _display_path(Path(deck_assembly_plan_path)),
        "template_image_manifest_path": _display_path(Path(template_image_manifest_path)),
        "status": "failed" if severe_count else "passed",
        "qa_blocks_deck_generation": bool(severe_count),
        "thresholds": {
            "full_slide_area_ratio": full_slide_area_ratio_threshold,
            "large_picture_area_ratio": large_picture_area_ratio_threshold,
        },
        "reference_template_image_names": sorted(reference_names),
        "reference_template_image_hash_count": len(reference_hashes),
        "slide_count": len(deck.slides),
        "findings_summary": {
            "total": len(findings),
            "severe": severe_count,
            "warning": sum(1 for finding in findings if finding["severity"] == "warning"),
        },
        "slides": slide_reports,
        "findings": findings,
        **counts,
    }
    return report


def build_final_deck_image_policy_report_from_files(
    *,
    pptx_path: str | Path = DEFAULT_PPTX_PATH,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    deck_assembly_plan_path: str | Path = DEFAULT_ASSEMBLY_PLAN,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
) -> Path:
    report = build_final_deck_image_policy_report(
        pptx_path=pptx_path,
        template_spec_path=template_spec_path,
        deck_assembly_plan_path=deck_assembly_plan_path,
        template_image_manifest_path=template_image_manifest_path,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return json_path


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Validate final deck picture shapes against editable template image policy.")
    parser.add_argument("--pptx", type=Path, default=DEFAULT_PPTX_PATH)
    parser.add_argument("--template-spec", type=Path, default=DEFAULT_TEMPLATE_SPEC)
    parser.add_argument("--assembly-plan", type=Path, default=DEFAULT_ASSEMBLY_PLAN)
    parser.add_argument("--template-image-manifest", type=Path, default=DEFAULT_TEMPLATE_IMAGE_MANIFEST)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        output = build_final_deck_image_policy_report_from_files(
            pptx_path=args.pptx,
            template_spec_path=args.template_spec,
            deck_assembly_plan_path=args.assembly_plan,
            template_image_manifest_path=args.template_image_manifest,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
        report = _load_json(output)
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"FINAL_DECK_IMAGE_POLICY_FAILED {exc}")
        return 1
    print(f"WROTE {output}")
    if report.get("qa_blocks_deck_generation"):
        print("FINAL_DECK_IMAGE_POLICY failed")
        return 1
    print("FINAL_DECK_IMAGE_POLICY passed")
    return 0


def _allowed_image_slots(layout: dict[str, Any]) -> list[dict[str, Any]]:
    slots = []
    for slot in layout.get("slots") or []:
        if not isinstance(slot, dict):
            continue
        slot_id = str(slot.get("slot_id") or "")
        slot_type = str(slot.get("slot_type") or "")
        if slot_type == "image" or _normalize(slot_id) in ALLOWED_IMAGE_SLOT_IDS:
            bounds = slot.get("bounds")
            if isinstance(bounds, dict):
                slots.append({"slot_id": slot_id, "bounds": _bounds_to_float(bounds)})
    return slots


def _matching_allowed_slot(bounds: dict[str, float], allowed_slots: list[dict[str, Any]]) -> dict[str, Any] | None:
    for slot in allowed_slots:
        slot_bounds = slot["bounds"]
        if _center_inside(bounds, slot_bounds) or _intersection_over_min_area(bounds, slot_bounds) >= 0.72:
            return slot
    return None


def _is_full_slide_picture(bounds: dict[str, float], slide_w: float, slide_h: float, area_ratio: float, threshold: float) -> bool:
    nearly_origin = bounds["x"] <= 0.2 and bounds["y"] <= 0.2
    nearly_full_extent = bounds["w"] >= slide_w * 0.9 and bounds["h"] >= slide_h * 0.9
    return area_ratio >= threshold and nearly_origin and nearly_full_extent


def _template_reference_hashes(template_image_manifest: dict[str, Any]) -> set[str]:
    hashes: set[str] = set()
    for path in _template_reference_paths(template_image_manifest):
        try:
            hashes.add(hashlib.sha256(path.read_bytes()).hexdigest())
        except OSError:
            continue
    return hashes


def _template_reference_names(template_image_manifest: dict[str, Any]) -> set[str]:
    return {path.name for path in _template_reference_paths(template_image_manifest)}


def _template_reference_paths(template_image_manifest: dict[str, Any]) -> list[Path]:
    paths = []
    for record in template_image_manifest.get("images") or []:
        if isinstance(record, dict) and isinstance(record.get("image_output_path"), str):
            paths.append(Path(record["image_output_path"]))
    return paths


def _shape_image_sha256(shape: Any) -> str:
    blob = _shape_image_blob(shape)
    return hashlib.sha256(blob).hexdigest()


def _shape_image_blob(shape: Any) -> bytes:
    element = getattr(shape, "_element", None)
    r_id = getattr(element, "blip_rId", None)
    if r_id:
        try:
            part = shape.part.related_part(r_id)
            blob = getattr(part, "blob", b"")
            if isinstance(blob, bytes):
                return blob
        except (KeyError, AttributeError, ValueError):
            pass
    try:
        image = getattr(shape, "image", None)
        blob = getattr(image, "blob", b"") if image is not None else b""
        return blob if isinstance(blob, bytes) else b""
    except (AttributeError, ValueError):
        return b""


def _shape_image_info(shape: Any) -> dict[str, Any]:
    element = getattr(shape, "_element", None)
    r_id = getattr(element, "blip_rId", None)
    if not r_id:
        return {}
    try:
        part = shape.part.related_part(r_id)
    except (KeyError, AttributeError, ValueError):
        return {}
    return {
        "rId": r_id,
        "content_type": str(getattr(part, "content_type", "") or ""),
        "partname": str(getattr(part, "partname", "") or ""),
    }


def _is_allowed_svg_icon_shape(shape: Any, image_info: dict[str, Any], area_ratio: float) -> bool:
    if image_info.get("content_type") != "image/svg+xml":
        return False
    if area_ratio > SVG_ICON_AREA_RATIO_THRESHOLD:
        return False
    name = str(getattr(shape, "name", "") or "")
    return name.startswith("SVG Icon ")


def _shape_bounds(shape: Any) -> dict[str, float]:
    return {
        "x": round(shape.left / EMU_PER_INCH, 4),
        "y": round(shape.top / EMU_PER_INCH, 4),
        "w": round(shape.width / EMU_PER_INCH, 4),
        "h": round(shape.height / EMU_PER_INCH, 4),
    }


def _bounds_to_float(bounds: dict[str, Any]) -> dict[str, float]:
    return {key: float(bounds[key]) for key in ("x", "y", "w", "h")}


def _center_inside(inner: dict[str, float], outer: dict[str, float]) -> bool:
    center_x = inner["x"] + inner["w"] / 2
    center_y = inner["y"] + inner["h"] / 2
    return outer["x"] <= center_x <= outer["x"] + outer["w"] and outer["y"] <= center_y <= outer["y"] + outer["h"]


def _intersection_over_min_area(first: dict[str, float], second: dict[str, float]) -> float:
    left = max(first["x"], second["x"])
    top = max(first["y"], second["y"])
    right = min(first["x"] + first["w"], second["x"] + second["w"])
    bottom = min(first["y"] + first["h"], second["y"] + second["h"])
    if right <= left or bottom <= top:
        return 0.0
    intersection = (right - left) * (bottom - top)
    min_area = max(0.01, min(first["w"] * first["h"], second["w"] * second["h"]))
    return intersection / min_area


def _slot_summary(slot: dict[str, Any]) -> dict[str, Any]:
    return {"slot_id": slot["slot_id"], "bounds": slot["bounds"]}


def _finding(code: str, severity: str, message: str, slide_index: int, details: dict[str, Any]) -> dict[str, Any]:
    return {
        "code": code,
        "severity": severity,
        "message": message,
        "slide_index": slide_index,
        "details": details,
    }


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Final Deck Image Policy Report",
        "",
        f"Status: `{report['status']}`",
        f"Picture shapes: `{report['picture_shape_count']}`",
        f"Allowed photo-frame pictures: `{report['allowed_photo_frame_picture_count']}`",
        f"Allowed SVG icon pictures: `{report.get('allowed_svg_icon_picture_count', 0)}`",
        f"Undeclared pictures: `{report['undeclared_picture_shape_count']}`",
        f"Full-slide pictures: `{report['full_slide_picture_count']}`",
        f"Template reference images embedded: `{report['reference_template_image_embedded_count']}`",
        "",
        "| Slide | Layout | Pictures | Findings |",
        "|---:|---|---:|---|",
    ]
    findings_by_slide: dict[int, list[str]] = {}
    for finding in report["findings"]:
        findings_by_slide.setdefault(int(finding["slide_index"]), []).append(str(finding["code"]))
    for slide in report["slides"]:
        findings = ", ".join(findings_by_slide.get(int(slide["slide_index"]), [])) or "none"
        lines.append(f"| {slide['slide_index']} | `{slide['layout_id']}` | {slide['picture_shape_count']} | {findings} |")
    return "\n".join(lines) + "\n"


def _load_json(path: str | Path) -> Any:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _display_path(path: Path) -> str:
    return str(path.as_posix())


def _normalize(value: str) -> str:
    return value.strip().lower()


if __name__ == "__main__":
    raise SystemExit(main())
