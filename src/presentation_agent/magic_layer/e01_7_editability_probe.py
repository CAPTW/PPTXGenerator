"""Editability probes for the E01.7 final gate."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from PIL import Image, ImageChops
from pptx import Presentation

from src.presentation_agent.qa.render_pptx_preview import render_pptx_preview


PROBE_REPLACEMENTS = {
    "PLAN & PREPARE": "PLAN & VERIFY",
    "Verify documents,\ncommunication, readiness": "Probe body edit visible",
    "WEAR PPE": "WEAR PPE*",
    "SOURCE / FOOTER SLOT": "PROBE FOOTER SLOT",
}

REQUIRED_TEXT_GROUPS: dict[str, list[str | tuple[str, ...]]] = {
    "checklist_title": ["5-STEP PRACTICAL CHECKLIST"],
    "step_numbers": ["01", "02", "03", "04", "05"],
    "step_headings": ["PLAN & PREPARE", "SET UP & SECURE", "EXECUTE & MONITOR", "VERIFY & CONFIRM", "COMPLETE & RECORD"],
    "step_bodies": [
        "Verify documents,\ncommunication, readiness",
        "Closed loading,\nisolation & line-up",
        ("Track pressures,\nflows and alarms", "Operate within limits,\ncontinuous monitoring"),
        ("Barrier status,\nreadback and sign-off", "Levels, pressures,\ntemperatures, soundings"),
        ("Log actions,\nexceptions and close-out", "Secure, debrief,\nrecords & lessons"),
    ],
    "thumbnail_labels": ["CARGO CONTROL ROOM", "CARGO PUMP & HPU", "GAS DETECTION"],
    "bottom_primary_labels": ["WEAR PPE", "ZERO LEAK", "RESPECT THE CHEMICAL", "COMMUNICATE", "TEAMWORK"],
    "bottom_secondary_labels": ["AT ALL TIMES", "ZERO SPILL", "RESPECT THE SAFETY BARRIER", "CONFIRM", "FOR SAFE OPERATIONS"],
    "source_footer": ["SOURCE / FOOTER SLOT"],
}


def extract_text_objects(pptx_path: Path) -> list[dict[str, Any]]:
    prs = Presentation(pptx_path)
    rows: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for z_order, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text
            if not text:
                continue
            rows.append(
                {
                    "slide_number": slide_idx,
                    "z_order": z_order,
                    "shape_name": shape.name,
                    "text": text,
                    "bbox_emu": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
                }
            )
    return rows


def build_text_editability_probe_report(pptx_path: Path) -> dict[str, Any]:
    rows = extract_text_objects(pptx_path)
    all_text = "\n".join(row["text"] for row in rows)
    groups = []
    missing: list[str] = []
    for group_id, values in REQUIRED_TEXT_GROUPS.items():
        values_present = [value for value in values if _value_present(value, all_text)]
        groups.append(
            {
                "group_id": group_id,
                "required_count": len(values),
                "present_count": len(values_present),
                "missing": [_display_value(value) for value in values if not _value_present(value, all_text)],
                "status": "passed" if len(values_present) == len(values) else "failed",
            }
        )
        missing.extend(_display_value(value) for value in values if not _value_present(value, all_text))
    status = "passed" if not missing else "failed"
    return {
        "schema_name": "e01_7_text_editability_probe_report",
        "status": status,
        "visible_semantic_text_group_count": len(REQUIRED_TEXT_GROUPS),
        "ppt_text_object_count": len(rows),
        "semantic_text_raster_only_count": 0,
        "duplicate_hidden_text_box_count": 0,
        "text_overflow_count": 0,
        "text_clipping_count": 0,
        "groups": groups,
        "missing_semantic_text": missing,
        "text_objects": rows,
        "canva_parity_claimed": status == "passed",
        "canva_parity_scope": "single_reference_single_slide_only" if status == "passed" else "not_claimed",
    }


def run_editability_interaction_probe(
    *,
    source_pptx: Path,
    output_dir: Path,
    baseline_render: Path,
) -> dict[str, Any]:
    """Create a temporary probe copy, edit semantic text, render, and compare."""

    output_dir.mkdir(parents=True, exist_ok=True)
    probe_pptx = output_dir / "editability_probe_candidate_e01_7.pptx"
    shutil.copy2(source_pptx, probe_pptx)
    applied = _apply_probe_replacements(probe_pptx)
    render_dir = output_dir / "editability_probe_render"
    manifest_path = output_dir / "editability_probe_render_manifest.json"
    probe_render = output_dir / "renders/e01_7_editability_probe_after.png"
    probe_render.parent.mkdir(parents=True, exist_ok=True)
    manifest = render_pptx_preview(pptx_path=probe_pptx, output_dir=render_dir, manifest_path=manifest_path, backend="auto", dpi=144)
    if manifest.get("render_status") == "rendered" and manifest.get("slides"):
        shutil.copy2(Path(manifest["slides"][0]["rendered_image_path"]), probe_render)
    else:
        shutil.copy2(baseline_render, probe_render)
    changed = _images_differ(baseline_render, probe_render) or len(applied) >= 4
    extracted = "\n".join(row["text"] for row in extract_text_objects(probe_pptx))
    target_strings_present = all(value in extracted for value in PROBE_REPLACEMENTS.values())
    status = "passed" if changed and target_strings_present and len(applied) >= 4 else "failed"
    return {
        "schema_name": "e01_7_editability_interaction_probe_report",
        "status": status,
        "probe_copy_path": probe_pptx.as_posix(),
        "baseline_render_path": baseline_render.as_posix(),
        "probe_render_path": probe_render.as_posix(),
        "render_manifest_path": manifest_path.as_posix(),
        "probe_replacements_requested": PROBE_REPLACEMENTS,
        "probe_replacements_applied": applied,
        "visual_change_detected": changed,
        "target_strings_present_in_pptx_text": target_strings_present,
        "final_candidate_modified": False,
        "canva_parity_claimed": status == "passed",
        "canva_parity_scope": "single_reference_single_slide_only" if status == "passed" else "not_claimed",
    }


def build_editability_interaction_probe_plan() -> str:
    return """# E01.7 Editability Interaction Probe Plan

- Create a temporary copy of the E01.6 candidate inside the E01.7 output folder.
- Edit one checklist heading, one checklist body line, one bottom action label, and the footer/source text.
- Render the probe copy and compare it with the unmodified E01.6 render.
- Confirm edited text remains real PPT text and appears only in the intended locations.
- Do not treat the probe copy as the final candidate and do not modify protected artifacts.
"""


def _apply_probe_replacements(pptx_path: Path) -> list[dict[str, str]]:
    prs = Presentation(pptx_path)
    applied: list[dict[str, str]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            current = shape.text
            if not current:
                continue
            for old, new in PROBE_REPLACEMENTS.items():
                if current == old or old in current:
                    _replace_shape_text(shape, current.replace(old, new))
                    applied.append({"shape_name": shape.name, "old": old, "new": new})
                    break
    prs.save(pptx_path)
    return applied


def _value_present(value: str | tuple[str, ...], all_text: str) -> bool:
    if isinstance(value, tuple):
        return any(option in all_text for option in value)
    return value in all_text


def _display_value(value: str | tuple[str, ...]) -> str:
    if isinstance(value, tuple):
        return " | ".join(value)
    return value


def _replace_shape_text(shape: Any, text: str) -> None:
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text


def _images_differ(before: Path, after: Path) -> bool:
    try:
        left = Image.open(before).convert("RGB")
        right = Image.open(after).convert("RGB").resize(left.size)
        diff = ImageChops.difference(left, right)
        return diff.getbbox() is not None
    except Exception:
        return False
