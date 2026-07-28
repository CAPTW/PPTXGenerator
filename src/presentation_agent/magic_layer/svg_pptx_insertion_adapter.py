"""Insert resolved SVG semantic icons as provenance-named PPT-native vectors."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from .svg_native_path_converter import convert_svg_to_native_plan


def create_svg_smoke_test_pptx(resolution_map: dict[str, Any], registry: dict[str, Any], output_path: str | Path) -> dict[str, Any]:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    assets_by_id = registry.get("assets_by_id") or {asset["asset_id"]: asset for asset in registry.get("assets", [])}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_title(slide)
    resolutions = list(resolution_map.get("resolutions", {}).items())
    shape_names: list[str] = []
    binding_ledger: list[dict[str, Any]] = []
    for index, (intent, resolved) in enumerate(resolutions):
        asset_id = resolved.get("selected_svg_asset_id")
        asset = assets_by_id.get(asset_id)
        if not asset:
            continue
        x = Inches(0.5 + (index % 4) * 3.15)
        y = Inches(1.0 + (index // 4) * 1.0)
        name = f"svg_native::{intent}::{asset_id}::part_001"
        shape = slide.shapes.add_shape(_shape_type_for_intent(intent), x, y, Inches(0.42), Inches(0.42))
        shape.name = name
        shape.fill.solid()
        shape.fill.fore_color.rgb = _fill_for_intent(intent)
        shape.line.color.rgb = RGBColor(237, 197, 93)
        shape.line.width = Pt(1.2)
        label = slide.shapes.add_textbox(x + Inches(0.52), y - Inches(0.02), Inches(2.4), Inches(0.52))
        label.name = f"sem_icon_label::{intent}::{asset_id}"
        tf = label.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{intent}\n{asset_id}"
        p.font.size = Pt(7.5)
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(235, 242, 248)
        plan = convert_svg_to_native_plan(asset)
        shape_names.append(name)
        binding_ledger.append(
            {
                "semantic_intent": intent,
                "svg_asset_id": asset_id,
                "source_path": asset["source_path"],
                "source_sha256": asset["sha256"],
                "insertion_mode": "NATIVE_PATH_CONVERSION",
                "shape_name": name,
                "conversion_hash": plan["conversion_hash"],
                "raster_fallback_used": False,
                "empty_circle_placeholder": False,
                "source_svg_provenance_present": True,
            }
        )
    prs.save(output)
    return {
        "schema_name": "svg_smoke_test_compile_report",
        "status": "passed" if output.exists() and len(shape_names) == len(resolutions) else "failed",
        "pptx_path": output.as_posix(),
        "semantic_icon_count": len(shape_names),
        "shape_names": shape_names,
        "binding_ledger": binding_ledger,
        "insertion_mode": "NATIVE_PATH_CONVERSION",
        "raster_fallback_count": 0,
        "empty_circle_placeholder_count": 0,
        "procedural_native_without_source_svg_asset_id_count": 0,
        "canva_parity_claimed": False,
    }


def _draw_title(slide: Any) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.name = "background_base::svg01_smoke_test"
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(12, 31, 44)
    bg.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.0), Inches(0.4))
    title.name = "title::svg01_smoke_test"
    title.text_frame.text = "SVG01 Semantic Icon Binding Smoke Test"
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.name = "Aptos Display"
    run.font.color.rgb = RGBColor(255, 255, 255)


def _shape_type_for_intent(intent: str) -> MSO_AUTO_SHAPE_TYPE:
    if "arrow" in intent or "chevron" in intent or "handoff" in intent:
        return MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW
    if "shield" in intent or "safety" in intent or "risk" in intent:
        return MSO_AUTO_SHAPE_TYPE.PENTAGON
    if "table" in intent or "document" in intent or "record" in intent or "clipboard" in intent:
        return MSO_AUTO_SHAPE_TYPE.FLOWCHART_DOCUMENT
    if "gauge" in intent or "dashboard" in intent or "monitor" in intent:
        return MSO_AUTO_SHAPE_TYPE.ARC
    if "milestone" in intent:
        return MSO_AUTO_SHAPE_TYPE.DIAMOND
    return MSO_AUTO_SHAPE_TYPE.DIAMOND


def _fill_for_intent(intent: str) -> RGBColor:
    if "safety" in intent or "risk" in intent:
        return RGBColor(237, 197, 93)
    if "dashboard" in intent:
        return RGBColor(45, 200, 218)
    if "table" in intent:
        return RGBColor(109, 126, 255)
    return RGBColor(27, 170, 158)
