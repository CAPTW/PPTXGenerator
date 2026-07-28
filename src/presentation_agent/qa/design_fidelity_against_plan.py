"""QA final decks against the design production plan."""

from __future__ import annotations

import argparse
import json
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .final_deck_image_policy import build_final_deck_image_policy_report


DEFAULT_PPTX = Path("outputs/final_deck_large_premium.pptx")
DEFAULT_LARGE_REPORT = Path("outputs/large_premium_deck_report.json")
DEFAULT_ASSEMBLY_PLAN = Path("outputs/deck_assembly_plan_large_premium.json")
DEFAULT_TEMPLATE_SPEC = Path("outputs/editable_template_spec.final.json")
DEFAULT_VISUAL_TARGETS = Path("outputs/design_planning/visual_fidelity_targets.json")
DEFAULT_LAYOUT_FAMILY_PLAN = Path("outputs/design_planning/layout_family_plan.json")
DEFAULT_COMPONENT_TRANSLATION_PLAN = Path("outputs/design_planning/component_translation_plan.json")
DEFAULT_TEMPLATE_VISUAL_DIFF = Path("outputs/template_visual_diff_report.json")
DEFAULT_TEMPLATE_IMAGE_MANIFEST = Path("outputs/template_images/template_image_manifest.json")
DEFAULT_RENDERED_PNG_DIR = Path("outputs/final_deck_large_premium_preview_png")
DEFAULT_JSON_REPORT = Path("outputs/design_fidelity_report.json")
DEFAULT_MD_REPORT = Path("outputs/design_fidelity_report.md")

SMOKE_PHRASES = (
    "deterministic slot filling",
    "large deck smoke synthetic source",
    "Large Editable Template Smoke Deck",
    "Image frame",
)


def build_design_fidelity_report(
    *,
    pptx_path: str | Path = DEFAULT_PPTX,
    large_deck_report_path: str | Path = DEFAULT_LARGE_REPORT,
    deck_assembly_plan_path: str | Path = DEFAULT_ASSEMBLY_PLAN,
    visual_fidelity_targets_path: str | Path = DEFAULT_VISUAL_TARGETS,
    layout_family_plan_path: str | Path = DEFAULT_LAYOUT_FAMILY_PLAN,
    component_translation_plan_path: str | Path = DEFAULT_COMPONENT_TRANSLATION_PLAN,
    template_visual_diff_report_path: str | Path = DEFAULT_TEMPLATE_VISUAL_DIFF,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    rendered_png_dir: str | Path = DEFAULT_RENDERED_PNG_DIR,
) -> dict[str, Any]:
    pptx = Path(pptx_path)
    findings: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []

    assembly_plan = _load_optional_json(Path(deck_assembly_plan_path), findings, "DECK_ASSEMBLY_PLAN_MISSING")
    layout_family_plan = _load_optional_json(Path(layout_family_plan_path), findings, "LAYOUT_FAMILY_PLAN_MISSING")
    component_plan = _load_optional_json(Path(component_translation_plan_path), findings, "COMPONENT_TRANSLATION_PLAN_MISSING")
    visual_targets = _load_optional_json(Path(visual_fidelity_targets_path), findings, "VISUAL_FIDELITY_TARGETS_MISSING")
    large_report = _load_optional_json(Path(large_deck_report_path), warnings, "LARGE_DECK_REPORT_MISSING", severity="warning")
    visual_diff = _load_optional_json(Path(template_visual_diff_report_path), warnings, "TEMPLATE_VISUAL_DIFF_MISSING", severity="warning")
    rendered_preview_summary = _rendered_preview_summary(Path(rendered_png_dir), warnings)

    pptx_stats = _pptx_stats(pptx, findings)
    visible_text = pptx_stats.pop("visible_text", [])
    smoke_hits = _smoke_phrase_hits(visible_text)
    for phrase in smoke_hits:
        findings.append(_finding("SMOKE_TEXT_PRESENT", "severe", f"Smoke/test phrase is visible in the final deck: {phrase}", details={"phrase": phrase}))

    image_policy = _image_policy(
        pptx_path=pptx,
        template_spec_path=Path(template_spec_path),
        assembly_plan_path=Path(deck_assembly_plan_path),
        template_image_manifest_path=Path(template_image_manifest_path),
        findings=findings,
    )

    layout_coverage = _layout_family_coverage(assembly_plan, layout_family_plan, visual_targets, findings, warnings)
    component_coverage = _component_translation_coverage(assembly_plan, component_plan, findings, warnings)
    density_status = _visual_density_status(assembly_plan, visual_targets, visual_diff, rendered_preview_summary, pptx_stats, findings, warnings)
    tone_status = _tone_expression_status(assembly_plan, findings, warnings)
    production_status = _production_plan_compliance(assembly_plan, large_report, visual_targets, findings, warnings)

    if image_policy.get("reference_template_image_embedded_count", 0) > 0:
        findings.append(_finding("DESIGN_REFERENCE_IMAGE_EMBEDDED", "severe", "A template/design reference image is embedded in the final deck."))
    if image_policy.get("full_slide_picture_count", 0) > 0:
        findings.append(_finding("FULL_SLIDE_RASTER_BACKGROUND", "severe", "A full-slide raster background is present."))

    severe_count = sum(1 for finding in findings if finding["severity"] == "severe")
    report = {
        "schema_name": "design_fidelity_against_plan_report",
        "schema_version": "1.0",
        "status": "passed" if severe_count == 0 else "failed",
        "pptx_path": _display_path(pptx),
        "large_deck_report_path": _display_path(Path(large_deck_report_path)),
        "deck_assembly_plan_path": _display_path(Path(deck_assembly_plan_path)),
        "visual_fidelity_targets_path": _display_path(Path(visual_fidelity_targets_path)),
        "layout_family_plan_path": _display_path(Path(layout_family_plan_path)),
        "component_translation_plan_path": _display_path(Path(component_translation_plan_path)),
        "template_visual_diff_report_path": _display_path(Path(template_visual_diff_report_path)) if Path(template_visual_diff_report_path).exists() else None,
        "rendered_png_dir": _display_path(Path(rendered_png_dir)) if Path(rendered_png_dir).exists() else None,
        "layout_family_coverage": layout_coverage,
        "component_translation_coverage": component_coverage,
        "visual_density_targets": density_status,
        "tone_expression": tone_status,
        "production_plan_compliance": production_status,
        "image_policy": {
            "status": image_policy.get("status"),
            "picture_shape_count": image_policy.get("picture_shape_count"),
            "full_slide_picture_count": image_policy.get("full_slide_picture_count"),
            "reference_template_image_embedded_count": image_policy.get("reference_template_image_embedded_count"),
            "undeclared_picture_shape_count": image_policy.get("undeclared_picture_shape_count"),
            "allowed_photo_frame_picture_count": image_policy.get("allowed_photo_frame_picture_count"),
        },
        "pptx_stats": pptx_stats,
        "smoke_phrase_hits": smoke_hits,
        "findings_summary": {
            "total": len(findings) + len(warnings),
            "severe": severe_count,
            "warning": len(warnings),
        },
        "findings": findings,
        "warnings": warnings,
    }
    return report


def build_design_fidelity_report_from_files(
    *,
    pptx_path: str | Path = DEFAULT_PPTX,
    large_deck_report_path: str | Path = DEFAULT_LARGE_REPORT,
    deck_assembly_plan_path: str | Path = DEFAULT_ASSEMBLY_PLAN,
    visual_fidelity_targets_path: str | Path = DEFAULT_VISUAL_TARGETS,
    layout_family_plan_path: str | Path = DEFAULT_LAYOUT_FAMILY_PLAN,
    component_translation_plan_path: str | Path = DEFAULT_COMPONENT_TRANSLATION_PLAN,
    template_visual_diff_report_path: str | Path = DEFAULT_TEMPLATE_VISUAL_DIFF,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    rendered_png_dir: str | Path = DEFAULT_RENDERED_PNG_DIR,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
) -> Path:
    report = build_design_fidelity_report(
        pptx_path=pptx_path,
        large_deck_report_path=large_deck_report_path,
        deck_assembly_plan_path=deck_assembly_plan_path,
        visual_fidelity_targets_path=visual_fidelity_targets_path,
        layout_family_plan_path=layout_family_plan_path,
        component_translation_plan_path=component_translation_plan_path,
        template_visual_diff_report_path=template_visual_diff_report_path,
        template_spec_path=template_spec_path,
        template_image_manifest_path=template_image_manifest_path,
        rendered_png_dir=rendered_png_dir,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return json_path


def _layout_family_coverage(
    assembly_plan: dict[str, Any],
    layout_family_plan: dict[str, Any],
    visual_targets: dict[str, Any],
    findings: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> dict[str, Any]:
    bindings = _bindings(assembly_plan)
    family_counts = Counter(str(binding.get("layout_family_id") or "unrecorded") for binding in bindings)
    families = [family for family in layout_family_plan.get("families") or [] if isinstance(family, dict)]
    required_for_large = {
        "expressive_cover_divider",
        "visual_toc_navigation",
        "evidence_overview",
        "problem_research_gap",
        "methodology_framework",
        "technical_flow_process",
        "comparison_matrix",
        "kpi_dashboard",
        "table_appendix",
    }
    used_families = {family for family in family_counts if family and family != "unrecorded"}
    missing_required = sorted(required_for_large - used_families) if assembly_plan.get("deck_scale") == "large" else []
    minimum = int(((visual_targets.get("targets") or {}).get("minimum_distinct_layout_families_used")) or 1)
    max_family_ratio = max((count / max(1, len(bindings)) for count in family_counts.values()), default=0.0)

    if not bindings:
        findings.append(_finding("NO_LAYOUT_BINDINGS", "severe", "Deck assembly plan has no slide layout bindings."))
    if len(used_families) <= 1:
        findings.append(_finding("NO_LAYOUT_FAMILY_DIVERSITY", "severe", "Deck uses one or zero layout families."))
    if len(used_families) < min(minimum, len(bindings)):
        findings.append(_finding("LAYOUT_FAMILY_COVERAGE_LOW", "severe", "Deck does not meet minimum layout-family diversity.", details={"used": len(used_families), "minimum": minimum}))
    if missing_required:
        findings.append(_finding("LARGE_DECK_REQUIRED_FAMILIES_MISSING", "severe", "Large deck is missing required production-plan layout families.", details={"missing_families": missing_required}))
    if max_family_ratio > 0.34:
        findings.append(_finding("LAYOUT_FAMILY_OVERUSED", "severe", "One layout family dominates the deck.", details={"family_counts": dict(family_counts)}))
    elif max_family_ratio > 0.25:
        warnings.append(_finding("LAYOUT_FAMILY_REPETITION_RISK", "warning", "One layout family is used heavily; human review should check repetition.", details={"family_counts": dict(family_counts)}))

    return {
        "status": "passed" if not missing_required and len(used_families) >= min(minimum, len(bindings)) and max_family_ratio <= 0.34 else "failed",
        "layout_family_count": len(families),
        "used_layout_family_count": len(used_families),
        "minimum_required_family_count": minimum,
        "family_counts": dict(sorted(family_counts.items())),
        "missing_required_families_for_large_deck": missing_required,
        "max_family_ratio": round(max_family_ratio, 4),
    }


def _component_translation_coverage(
    assembly_plan: dict[str, Any],
    component_plan: dict[str, Any],
    findings: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> dict[str, Any]:
    plan_families = {
        str(item.get("component_family") or "")
        for item in component_plan.get("components") or []
        if isinstance(item, dict)
    }
    used_components: Counter[str] = Counter()
    needs: Counter[str] = Counter()
    for binding in _bindings(assembly_plan):
        for component_id in (binding.get("component_bindings") or {}).values():
            if component_id:
                used_components[str(component_id)] += 1
        for need in binding.get("content_needs") or []:
            needs[str(need)] += 1
    used_keys = set(used_components)
    coverage = {
        "footer_citation_modules_present": bool({"dense_footer", "citation_strip", "footer_system"} & used_keys),
        "card_system_present": bool({"layered_card", "cards", "premium_kpi_card", "kpi_cards"} & used_keys),
        "table_chart_modules_present": bool({"thin_grid_table", "table_modules", "chart_module", "chart_modules", "comparison_matrix"} & used_keys),
        "process_timeline_modules_present_when_needed": needs["process"] == 0 and needs["timeline"] == 0 or bool({"radial_process", "curved_timeline", "process_arrows", "timeline_blocks"} & used_keys),
        "section_index_navigation_present": bool({"index_navigation", "section_tabs", "oversized_section_number", "section_marker"} & used_keys),
    }
    if not plan_families:
        findings.append(_finding("COMPONENT_TRANSLATION_PLAN_EMPTY", "severe", "Component translation plan is missing or empty."))
    if not coverage["footer_citation_modules_present"]:
        findings.append(_finding("NO_FOOTER_CITATION_SYSTEM", "severe", "Deck does not use footer/citation components from the production plan."))
    if not coverage["card_system_present"]:
        findings.append(_finding("CARD_SYSTEM_MISSING", "severe", "Deck does not use the production-plan card system."))
    if needs["table"] or needs["chart"]:
        if not coverage["table_chart_modules_present"]:
            findings.append(_finding("DATA_MODULE_SYSTEM_MISSING", "severe", "Deck needs table/chart modules but does not use translated data components."))
    if not coverage["process_timeline_modules_present_when_needed"]:
        findings.append(_finding("PROCESS_TIMELINE_SYSTEM_MISSING", "severe", "Deck needs process/timeline modules but does not use translated process components."))
    if not coverage["section_index_navigation_present"]:
        warnings.append(_finding("SECTION_NAVIGATION_COMPONENTS_LOW", "warning", "Section/index/navigation component use is low or not explicit."))

    return {
        "status": "passed" if not any(code in {"NO_FOOTER_CITATION_SYSTEM", "CARD_SYSTEM_MISSING", "DATA_MODULE_SYSTEM_MISSING", "PROCESS_TIMELINE_SYSTEM_MISSING"} for code in _finding_codes(findings)) else "failed",
        "component_plan_family_count": len(plan_families),
        "used_component_count": sum(used_components.values()),
        "used_component_families": dict(sorted(used_components.items())),
        "content_need_counts": dict(sorted(needs.items())),
        **coverage,
    }


def _visual_density_status(
    assembly_plan: dict[str, Any],
    visual_targets: dict[str, Any],
    visual_diff: dict[str, Any],
    rendered_preview_summary: dict[str, Any],
    pptx_stats: dict[str, Any],
    findings: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> dict[str, Any]:
    bindings = _bindings(assembly_plan)
    component_density = Counter(str(binding.get("component_density") or "unrecorded") for binding in bindings)
    ornament_density = Counter(str(binding.get("ornament_density") or "unrecorded") for binding in bindings)
    sparse = [
        binding.get("slide_id")
        for binding in bindings
        if binding.get("component_density") in {"low", "none"} and binding.get("slide_character") not in {"cover", "section_divider", "closing"}
    ]
    dense_without_simplification = [
        binding.get("slide_id")
        for binding in bindings
        if binding.get("component_density") in {"data_dense", "evidence_dense"} and binding.get("ornament_density") not in {"low", "medium-low"}
    ]
    compliance = assembly_plan.get("visual_fidelity_target_compliance") or {}
    if compliance.get("generic_layout_ratio", 0) > ((visual_targets.get("targets") or {}).get("maximum_generic_layout_ratio") or 0.0):
        findings.append(_finding("GENERIC_WHITE_CARD_RATIO_TOO_HIGH", "severe", "Generic layout ratio exceeds visual fidelity target."))
    if len(sparse) > max(2, len(bindings) * 0.18):
        warnings.append(_finding("SLIDES_POTENTIALLY_TOO_SPARSE", "warning", "Several non-divider slides have low component density.", details={"slide_ids": sparse[:10]}))
    if dense_without_simplification:
        warnings.append(_finding("DENSE_SLIDES_NEED_ORNAMENT_SIMPLIFICATION_REVIEW", "warning", "Dense data/evidence slides should use reduced ornament density.", details={"slide_ids": dense_without_simplification[:10]}))
    if int(pptx_stats.get("shape_count") or 0) <= max(1, int(pptx_stats.get("slide_count") or 0)) * 4:
        warnings.append(_finding("OBJECT_DENSITY_LOW", "warning", "PPTX has low editable object count for a premium production-plan deck."))

    return {
        "status": "passed",
        "component_density_counts": dict(sorted(component_density.items())),
        "ornament_density_counts": dict(sorted(ornament_density.items())),
        "potentially_sparse_slide_ids": sparse,
        "dense_slides_without_simplified_ornaments": dense_without_simplification,
        "generic_layout_ratio": compliance.get("generic_layout_ratio"),
        "background_ornament_target": (visual_targets.get("targets") or {}).get("minimum_non_white_background_ornament_occupancy"),
        "footer_citation_target": (visual_targets.get("targets") or {}).get("footer_citation_presence_ratio"),
        "section_navigation_target": (visual_targets.get("targets") or {}).get("section_navigation_presence_ratio"),
        "template_visual_diff_status": visual_diff.get("status") or visual_diff.get("render_status"),
        "rendered_preview_summary": rendered_preview_summary,
    }


def _tone_expression_status(
    assembly_plan: dict[str, Any],
    findings: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> dict[str, Any]:
    bindings = _bindings(assembly_plan)
    tones = Counter(str(binding.get("selected_tone_variant") or assembly_plan.get("selected_tone_variant") or "") for binding in bindings)
    if not tones or set(tones) <= {""}:
        findings.append(_finding("NO_TONE_VARIANT", "severe", "No tone variant is selected in the assembly plan."))
    component_density_by_tone: dict[str, Counter[str]] = {}
    ornament_density_by_tone: dict[str, Counter[str]] = {}
    for binding in bindings:
        tone = str(binding.get("selected_tone_variant") or assembly_plan.get("selected_tone_variant") or "unrecorded")
        component_density_by_tone.setdefault(tone, Counter())[str(binding.get("component_density") or "unrecorded")] += 1
        ornament_density_by_tone.setdefault(tone, Counter())[str(binding.get("ornament_density") or "unrecorded")] += 1
    if len(tones) < 2:
        warnings.append(_finding("TONE_VARIETY_LOW", "warning", "Deck uses one tone variant; this can be acceptable but should be intentional."))
    return {
        "status": "passed" if tones and set(tones) != {""} else "failed",
        "selected_tone_counts": dict(sorted(tones.items())),
        "component_density_by_tone": {key: dict(value) for key, value in sorted(component_density_by_tone.items())},
        "ornament_density_by_tone": {key: dict(value) for key, value in sorted(ornament_density_by_tone.items())},
        "tone_affects_visible_component_choices": len(component_density_by_tone) >= 2 or len(ornament_density_by_tone) >= 2,
    }


def _production_plan_compliance(
    assembly_plan: dict[str, Any],
    large_report: dict[str, Any],
    visual_targets: dict[str, Any],
    findings: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> dict[str, Any]:
    bindings = _bindings(assembly_plan)
    fallback_count = sum(1 for binding in bindings if binding.get("fallback_used"))
    slide_count = max(1, len(bindings))
    fallback_ratio = fallback_count / slide_count
    max_fallback_ratio = float((visual_targets.get("targets") or {}).get("maximum_fallback_ratio") or 0.08)
    generic_bindings = [
        binding
        for binding in bindings
        if str(binding.get("selected_layout_id") or "").endswith("-mvp") or "standard-content-mvp" in str(binding.get("selected_layout_id") or "")
    ]
    rhythm_counts = Counter(str(binding.get("section_rhythm_role") or "unrecorded") for binding in bindings)
    provenance = large_report.get("template_spec_provenance") or {}
    if assembly_plan.get("production_plan_used") is not True:
        findings.append(_finding("DESIGN_PRODUCTION_PLAN_NOT_USED", "severe", "Deck assembly plan does not record production_plan_used=true."))
    if provenance and provenance.get("source") != "design_board_production_plan":
        findings.append(_finding("SPEC_NOT_PRODUCTION_PLAN_DERIVED", "severe", "Large deck report does not show a production-plan-derived final spec."))
    if provenance and provenance.get("extraction_source") != "actual":
        findings.append(_finding("SCHEMA_SAMPLE_EXTRACTION_FALLBACK_USED", "severe", "Template spec provenance does not record actual extraction source."))
    if generic_bindings:
        findings.append(_finding("GENERIC_SMOKE_LAYOUT_USED", "severe", "Generic smoke/MVP layout was selected.", details={"slide_ids": [item.get("slide_id") for item in generic_bindings]}))
    if fallback_ratio > max_fallback_ratio:
        findings.append(_finding("FALLBACK_RATIO_EXCEEDS_PLAN", "severe", "Fallback count exceeds visual fidelity target.", details={"fallback_count": fallback_count, "fallback_ratio": round(fallback_ratio, 4), "maximum": max_fallback_ratio}))
    missing_rhythm = [role for role in ("divider", "overview", "evidence", "analysis", "implication") if rhythm_counts[role] == 0 and assembly_plan.get("deck_scale") == "large"]
    if missing_rhythm:
        warnings.append(_finding("SECTION_RHYTHM_INCOMPLETE", "warning", "Large-deck rhythm roles are not all represented.", details={"missing_roles": missing_rhythm}))
    return {
        "status": "passed" if assembly_plan.get("production_plan_used") is True and not generic_bindings and fallback_ratio <= max_fallback_ratio else "failed",
        "production_plan_used": assembly_plan.get("production_plan_used") is True,
        "deck_scale": assembly_plan.get("deck_scale"),
        "fallback_count": fallback_count,
        "fallback_ratio": round(fallback_ratio, 4),
        "maximum_fallback_ratio": max_fallback_ratio,
        "generic_layout_count": len(generic_bindings),
        "section_rhythm_counts": dict(sorted(rhythm_counts.items())),
        "template_spec_provenance": provenance,
    }


def _image_policy(
    *,
    pptx_path: Path,
    template_spec_path: Path,
    assembly_plan_path: Path,
    template_image_manifest_path: Path,
    findings: list[dict[str, Any]],
) -> dict[str, Any]:
    try:
        return build_final_deck_image_policy_report(
            pptx_path=pptx_path,
            template_spec_path=template_spec_path,
            deck_assembly_plan_path=assembly_plan_path,
            template_image_manifest_path=template_image_manifest_path,
        )
    except Exception as exc:  # noqa: BLE001 - QA should report, not crash before writing.
        findings.append(_finding("IMAGE_POLICY_UNAVAILABLE", "severe", f"Image policy QA could not run: {type(exc).__name__}: {exc}"))
        return {"status": "failed", "picture_shape_count": None, "full_slide_picture_count": None, "reference_template_image_embedded_count": None}


def _pptx_stats(path: Path, findings: list[dict[str, Any]]) -> dict[str, Any]:
    if not path.exists():
        findings.append(_finding("FINAL_DECK_MISSING", "severe", f"Final deck is missing: {_display_path(path)}"))
        return {"slide_count": 0, "shape_count": 0, "text_shape_count": 0, "table_count": 0, "chart_count": 0, "picture_shape_count": 0, "visible_text": []}
    deck = Presentation(path)
    stats = {"slide_count": len(deck.slides), "shape_count": 0, "text_shape_count": 0, "table_count": 0, "chart_count": 0, "picture_shape_count": 0, "visible_text": []}
    for slide in deck.slides:
        for shape in slide.shapes:
            stats["shape_count"] += 1
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
                if text:
                    stats["text_shape_count"] += 1
                    stats["visible_text"].append(text)
            if getattr(shape, "has_table", False):
                stats["table_count"] += 1
            if getattr(shape, "has_chart", False):
                stats["chart_count"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                stats["picture_shape_count"] += 1
    return stats


def _rendered_preview_summary(rendered_png_dir: Path, warnings: list[dict[str, Any]]) -> dict[str, Any]:
    if not rendered_png_dir.exists():
        warnings.append(_finding("RENDERED_PREVIEWS_MISSING", "warning", "Rendered final-deck PNG previews were not found; visual density checks used assembly metadata only."))
        return {"status": "missing", "directory": _display_path(rendered_png_dir), "rendered_slide_count": 0, "paths": []}
    pngs = sorted(rendered_png_dir.glob("*.png"))
    summary: dict[str, Any] = {
        "status": "available" if pngs else "empty",
        "directory": _display_path(rendered_png_dir),
        "rendered_slide_count": len(pngs),
        "paths": [_display_path(path) for path in pngs[:8]],
    }
    try:
        from PIL import Image, ImageStat  # type: ignore

        brightness_values: list[float] = []
        for path in pngs[:12]:
            with Image.open(path) as image:
                gray = image.convert("L")
                stat = ImageStat.Stat(gray)
                brightness_values.append(float(stat.mean[0]) / 255.0)
        if brightness_values:
            summary["average_brightness_sample"] = round(sum(brightness_values) / len(brightness_values), 4)
            summary["brightness_range_sample"] = [round(min(brightness_values), 4), round(max(brightness_values), 4)]
    except Exception as exc:  # noqa: BLE001 - optional rendered metrics should not block QA.
        warnings.append(_finding("RENDERED_PREVIEW_METRICS_SKIPPED", "warning", f"Rendered PNG metrics could not be computed: {type(exc).__name__}: {exc}"))
    return summary


def _bindings(assembly_plan: dict[str, Any]) -> list[dict[str, Any]]:
    return [binding for binding in assembly_plan.get("slide_layout_bindings") or [] if isinstance(binding, dict)]


def _smoke_phrase_hits(visible_text: list[str]) -> list[str]:
    text = "\n".join(visible_text).lower()
    return [phrase for phrase in SMOKE_PHRASES if phrase.lower() in text]


def _finding_codes(findings: list[dict[str, Any]]) -> set[str]:
    return {str(finding.get("code") or "") for finding in findings}


def _load_optional_json(path: Path, findings: list[dict[str, Any]], code: str, *, severity: str = "severe") -> dict[str, Any]:
    if not path.exists():
        findings.append(_finding(code, severity, f"Required artifact is missing: {_display_path(path)}"))
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        findings.append(_finding(code, severity, f"Artifact could not be read: {type(exc).__name__}: {exc}", details={"path": _display_path(path)}))
        return {}


def _finding(code: str, severity: str, message: str, *, details: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = {"code": code, "severity": severity, "message": message}
    if details:
        payload["details"] = details
    return payload


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Design Fidelity Against Plan Report",
        "",
        f"Status: `{report['status']}`",
        f"Deck: `{report['pptx_path']}`",
        f"Assembly plan: `{report['deck_assembly_plan_path']}`",
        f"Layout family coverage: `{report['layout_family_coverage']['status']}`",
        f"Component translation coverage: `{report['component_translation_coverage']['status']}`",
        f"Tone expression: `{report['tone_expression']['status']}`",
        f"Production plan compliance: `{report['production_plan_compliance']['status']}`",
        f"Image policy: `{report['image_policy'].get('status')}`",
        "",
        "## Key Metrics",
        "",
        f"- Used layout families: `{report['layout_family_coverage']['used_layout_family_count']}`",
        f"- Fallback ratio: `{report['production_plan_compliance']['fallback_ratio']}`",
        f"- Generic layout ratio: `{report['visual_density_targets'].get('generic_layout_ratio')}`",
        f"- Smoke phrase hits: `{len(report['smoke_phrase_hits'])}`",
        "",
        "## Severe Findings",
        "",
    ]
    severe = [finding for finding in report["findings"] if finding["severity"] == "severe"]
    if severe:
        for finding in severe:
            lines.append(f"- `{finding['code']}`: {finding['message']}")
    else:
        lines.append("- None")
    lines.extend(["", "## Warnings", ""])
    if report["warnings"]:
        for warning in report["warnings"]:
            lines.append(f"- `{warning['code']}`: {warning['message']}")
    else:
        lines.append("- None")
    return "\n".join(lines) + "\n"


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Compare a final deck against design production plan fidelity targets.")
    parser.add_argument("--pptx", type=Path, default=DEFAULT_PPTX)
    parser.add_argument("--large-deck-report", type=Path, default=DEFAULT_LARGE_REPORT)
    parser.add_argument("--deck-assembly-plan", type=Path, default=DEFAULT_ASSEMBLY_PLAN)
    parser.add_argument("--visual-fidelity-targets", type=Path, default=DEFAULT_VISUAL_TARGETS)
    parser.add_argument("--layout-family-plan", type=Path, default=DEFAULT_LAYOUT_FAMILY_PLAN)
    parser.add_argument("--component-translation-plan", type=Path, default=DEFAULT_COMPONENT_TRANSLATION_PLAN)
    parser.add_argument("--template-visual-diff-report", type=Path, default=DEFAULT_TEMPLATE_VISUAL_DIFF)
    parser.add_argument("--template-spec", type=Path, default=DEFAULT_TEMPLATE_SPEC)
    parser.add_argument("--template-image-manifest", type=Path, default=DEFAULT_TEMPLATE_IMAGE_MANIFEST)
    parser.add_argument("--rendered-png-dir", type=Path, default=DEFAULT_RENDERED_PNG_DIR)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        output = build_design_fidelity_report_from_files(
            pptx_path=args.pptx,
            large_deck_report_path=args.large_deck_report,
            deck_assembly_plan_path=args.deck_assembly_plan,
            visual_fidelity_targets_path=args.visual_fidelity_targets,
            layout_family_plan_path=args.layout_family_plan,
            component_translation_plan_path=args.component_translation_plan,
            template_visual_diff_report_path=args.template_visual_diff_report,
            template_spec_path=args.template_spec,
            template_image_manifest_path=args.template_image_manifest,
            rendered_png_dir=args.rendered_png_dir,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
        report = json.loads(output.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"DESIGN_FIDELITY_QA_FAILED {type(exc).__name__}: {exc}")
        return 1
    print(f"WROTE {output}")
    if report.get("status") != "passed":
        print("DESIGN_FIDELITY_QA failed")
        for finding in report.get("findings", []):
            if finding.get("severity") == "severe":
                print(f"DESIGN_FIDELITY_FAILURE {finding.get('code')}: {finding.get('message')}")
        return 1
    print("DESIGN_FIDELITY_QA passed")
    return 0


def _display_path(path: Path) -> str:
    return path.as_posix()


if __name__ == "__main__":
    raise SystemExit(main())
