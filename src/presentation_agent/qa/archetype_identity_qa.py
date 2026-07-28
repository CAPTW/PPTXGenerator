"""Classify archetype identity drift in golden template masters."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .render_pptx_preview import render_pptx_preview


DEFAULT_GOLDEN_PPTX = Path("outputs/golden_template_masters.pptx")
DEFAULT_CROPS_DIR = Path("outputs/template_design_board/crops")
DEFAULT_MASTER_SPECS_DIR = Path("outputs/golden_master_specs")
DEFAULT_FIDELITY_REPORT = Path("outputs/golden_template_fidelity_report.json")
DEFAULT_RENDER_DIR = Path("outputs/archetype_identity_png")
DEFAULT_RENDER_MANIFEST = Path("outputs/archetype_identity_render_manifest.json")
DEFAULT_JSON_REPORT = Path("outputs/archetype_identity_report.json")
DEFAULT_MD_REPORT = Path("outputs/archetype_identity_report.md")

DEBUG_PREVIEW_STRINGS = (
    "SAFE MARGINS",
    "ARCHETYPE:",
    "DENSITY:",
    "FOOTER SYSTEM",
    "IMAGE FRAME ONLY",
    "PREVIEW WARNINGS",
    "layout-board-",
)

SEVERE_DRIFT_RULES = {
    ("section_divider", "content_page"): "section_divider rendered as content page",
    ("visual_table_of_contents", "navigation_list_page"): "visual_table_of_contents rendered as navigation/list page",
    ("creative_cover", "wireframe_info_page"): "creative_cover rendered as wireframe/info page",
    ("methodology_framework", "generic_card_page"): "methodology framework rendered as generic card page",
    ("data_table_appendix", "table_without_appendix_identity"): "data table appendix rendered without appendix table identity",
    ("data_table_appendix", "content_page"): "data table appendix rendered without appendix table identity",
}


def build_archetype_identity_report_from_files(
    *,
    golden_pptx_path: str | Path = DEFAULT_GOLDEN_PPTX,
    crops_dir: str | Path = DEFAULT_CROPS_DIR,
    master_specs_dir: str | Path = DEFAULT_MASTER_SPECS_DIR,
    fidelity_report_path: str | Path = DEFAULT_FIDELITY_REPORT,
    render_dir: str | Path = DEFAULT_RENDER_DIR,
    render_manifest_path: str | Path = DEFAULT_RENDER_MANIFEST,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
    renderer: str = "auto",
) -> dict[str, Any]:
    report = build_archetype_identity_report(
        golden_pptx_path=golden_pptx_path,
        crops_dir=crops_dir,
        master_specs_dir=master_specs_dir,
        fidelity_report_path=fidelity_report_path,
        render_dir=render_dir,
        render_manifest_path=render_manifest_path,
        renderer=renderer,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return report


def build_archetype_identity_report(
    *,
    golden_pptx_path: str | Path,
    crops_dir: str | Path,
    master_specs_dir: str | Path,
    fidelity_report_path: str | Path,
    render_dir: str | Path,
    render_manifest_path: str | Path,
    renderer: str = "auto",
) -> dict[str, Any]:
    pptx_path = Path(golden_pptx_path)
    if not pptx_path.exists():
        raise ValueError(f"golden template masters PPTX is missing: {pptx_path.as_posix()}")
    specs = _load_master_specs(Path(master_specs_dir))
    fidelity = _load_json(fidelity_report_path)
    render_paths = _render_paths_from_fidelity(fidelity)
    if renderer != "none" and not _render_paths_exist(render_paths):
        render_report = render_pptx_preview(
            pptx_path=pptx_path,
            output_dir=render_dir,
            report_path=render_manifest_path,
            backend=renderer,
        )
        render_paths = _render_paths_from_render_report(render_report)
    records = _slide_records(fidelity, specs)
    deck = Presentation(pptx_path)

    slide_results: list[dict[str, Any]] = []
    severe_findings: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    crops = _crop_paths_by_archetype(Path(crops_dir))

    for index, record in enumerate(records, start=1):
        slide_number = int(record.get("slide_number") or index)
        intended = str(record.get("archetype_id") or "")
        slide = deck.slides[slide_number - 1]
        spec = specs.get(intended, {})
        rendered_path = render_paths.get(slide_number) or _rendered_path_from_record(record)
        features = _slide_features(slide, rendered_path)
        classification = _classify_slide(intended, spec, features)
        finding = _severe_finding_for(intended, classification["rendered_visual_family"], slide_number)
        if finding:
            severe_findings.append(finding)
        elif not classification["identity_preserved"]:
            warnings.append(
                _finding(
                    "ARCHETYPE_IDENTITY_DRIFT",
                    "warning",
                    f"{intended} rendered as {classification['rendered_visual_family']}.",
                    slide_number=slide_number,
                    archetype_id=intended,
                )
            )
        slide_results.append(
            {
                "slide_number": slide_number,
                "intended_archetype": intended,
                "rendered_visual_family": classification["rendered_visual_family"],
                "identity_preserved": classification["identity_preserved"],
                "drift_target": None if classification["identity_preserved"] else classification["rendered_visual_family"],
                "production_master_like": classification["production_master_like"],
                "inspection_preview_like": classification["inspection_preview_like"],
                "recommended_fix": classification["recommended_fix"],
                "reference_crop_path": _display_path(crops[intended]) if intended in crops else None,
                "rendered_image_path": _display_path(rendered_path) if rendered_path else None,
                "evidence": classification["evidence"],
            }
        )

    status = "failed" if severe_findings else "passed" if not warnings else "issues_reported"
    return {
        "schema_name": "archetype_identity_report",
        "schema_version": "1.0",
        "status": status,
        "golden_pptx_path": _display_path(pptx_path),
        "crops_dir": _display_path(Path(crops_dir)),
        "golden_master_specs_dir": _display_path(Path(master_specs_dir)),
        "golden_template_fidelity_report_path": _display_path(Path(fidelity_report_path)),
        "slide_count": len(slide_results),
        "findings_summary": {
            "severe": len(severe_findings),
            "warning": len(warnings),
            "total": len(severe_findings) + len(warnings),
        },
        "severe_findings": severe_findings,
        "warnings": warnings,
        "slides": slide_results,
    }


def _slide_records(fidelity: dict[str, Any], specs: dict[str, dict[str, Any]]) -> list[dict[str, Any]]:
    comparisons = [item for item in fidelity.get("layout_comparisons") or [] if isinstance(item, dict)]
    if comparisons:
        return sorted(comparisons, key=lambda item: int(item.get("slide_number") or 0))
    result = []
    for index, archetype_id in enumerate(sorted(specs), start=1):
        result.append({"slide_number": index, "archetype_id": archetype_id})
    return result


def _classify_slide(intended: str, spec: dict[str, Any], features: dict[str, Any]) -> dict[str, Any]:
    text = features["text_joined"].lower()
    inspection_preview_like = _inspection_preview_like(text)
    production_master_like = (
        not inspection_preview_like
        and features["shape_count"] >= 35
        and features["picture_count"] == 0
        and features["text_shape_count"] >= 2
    )
    family = _rendered_visual_family(intended, spec, features)
    identity_preserved = _expected_family(intended) == family and production_master_like
    if intended == "visual_table_of_contents" and family == "visual_toc_navigation":
        identity_preserved = production_master_like
    if intended == "data_table_appendix" and family == "appendix_table":
        identity_preserved = production_master_like
    return {
        "rendered_visual_family": family,
        "identity_preserved": identity_preserved,
        "production_master_like": production_master_like,
        "inspection_preview_like": inspection_preview_like,
        "recommended_fix": "No action." if identity_preserved else _recommended_fix(intended, family),
        "evidence": {
            "shape_count": features["shape_count"],
            "text_shape_count": features["text_shape_count"],
            "line_count": features["line_count"],
            "table_count": features["table_count"],
            "chart_count": features["chart_count"],
            "picture_count": features["picture_count"],
            "dark_area_ratio": features.get("dark_area_ratio"),
            "non_white_ratio": features.get("non_white_ratio"),
            "key_text": features["text_samples"],
        },
    }


def _rendered_visual_family(intended: str, spec: dict[str, Any], features: dict[str, Any]) -> str:
    text = features["text_joined"].lower()
    table_count = features["table_count"]
    chart_count = features["chart_count"]
    line_count = features["line_count"]
    shape_count = features["shape_count"]
    dark_ratio = float(features.get("dark_area_ratio") or 0.0)
    has_index = "index" in text
    has_table_label = "table" in text
    has_chart_label = "chart" in text
    has_appendix = "appendix" in text
    has_section_number = _has_two_digit_marker(features["text_samples"])
    content_slot_ids = _content_slot_ids(spec)

    if intended == "creative_cover":
        if dark_ratio >= 0.35 and shape_count >= 100 and _has_image_zone(spec):
            return "creative_hero_cover"
        return "wireframe_info_page"
    if intended == "section_divider":
        if has_section_number and shape_count >= 90 and _has_image_zone(spec):
            return "section_divider"
        return "content_page"
    if intended == "visual_table_of_contents":
        if has_index and shape_count >= 45 and line_count >= 12 and {"index_navigation", "progress_markers"} & content_slot_ids:
            return "visual_toc_navigation"
        return "navigation_list_page"
    if intended == "methodology_framework":
        if "diagram" in content_slot_ids and {"side_notes", "method_steps"} <= content_slot_ids and line_count >= 12:
            return "methodology_framework"
        return "generic_card_page"
    if intended == "data_table_appendix":
        if table_count > 0 and has_appendix and has_table_label:
            return "appendix_table"
        if table_count > 0:
            return "table_without_appendix_identity"
        return "content_page"
    if table_count > 0:
        return "comparison_matrix" if intended == "comparison_matrix" else "table_layout"
    if chart_count > 0 or has_chart_label:
        return "kpi_dashboard"
    if intended == "photo_caption_grid" and _has_image_zone(spec):
        return "photo_caption_grid"
    if intended in {"technical_flow_chart", "work_support_sequence", "circular_process", "timeline_roadmap"}:
        if line_count >= 10:
            return _expected_family(intended)
        return "generic_card_page"
    if intended in {"research_gap", "literature_map", "concept_relationship_venn", "three_level_explanation"}:
        if line_count >= 10:
            return _expected_family(intended)
        return "generic_card_page"
    if intended in {"research_overview", "problem_statement"}:
        if shape_count >= 40:
            return _expected_family(intended)
        return "generic_card_page"
    return _expected_family(intended)


def _expected_family(archetype_id: str) -> str:
    mapping = {
        "creative_cover": "creative_hero_cover",
        "visual_table_of_contents": "visual_toc_navigation",
        "section_divider": "section_divider",
        "research_overview": "evidence_card_overview",
        "problem_statement": "problem_callout_layout",
        "research_gap": "research_gap_diagram",
        "literature_map": "literature_map_diagram",
        "methodology_framework": "methodology_framework",
        "technical_flow_chart": "technical_flow_diagram",
        "work_support_sequence": "work_support_sequence",
        "photo_caption_grid": "photo_caption_grid",
        "comparison_matrix": "comparison_matrix",
        "concept_relationship_venn": "relationship_diagram",
        "three_level_explanation": "three_level_explanation",
        "circular_process": "circular_process",
        "kpi_donut_chart": "kpi_dashboard",
        "timeline_roadmap": "timeline_roadmap",
        "data_table_appendix": "appendix_table",
    }
    return mapping.get(archetype_id, "production_master")


def _recommended_fix(intended: str, family: str) -> str:
    fixes = {
        "creative_cover": "Restore the dark hero field, diagonal photo mask, and dense topology ornaments from the canonical cover master spec.",
        "section_divider": "Rebuild the divider with the oversized section number, fractured geometry panel, and minimal section thesis content.",
        "visual_table_of_contents": "Use the visual navigation archetype: modular numbered navigation cards and progress markers, not a plain list.",
        "methodology_framework": "Restore the central method diagram plus side-note and method-step rails from the canonical master spec.",
        "data_table_appendix": "Render a dominant editable appendix table with appendix title, table chrome, and compact source strip.",
    }
    return fixes.get(intended, f"Adjust the golden master so it classifies as {_expected_family(intended)} instead of {family}.")


def _severe_finding_for(intended: str, family: str, slide_number: int) -> dict[str, Any] | None:
    message = SEVERE_DRIFT_RULES.get((intended, family))
    if not message:
        return None
    return _finding(
        "ARCHETYPE_IDENTITY_SEVERE_DRIFT",
        "severe",
        message,
        slide_number=slide_number,
        archetype_id=intended,
        details={"rendered_visual_family": family},
    )


def _slide_features(slide: Any, rendered_path: Path | None) -> dict[str, Any]:
    text_samples: list[str] = []
    features = {
        "shape_count": 0,
        "text_shape_count": 0,
        "table_count": 0,
        "chart_count": 0,
        "picture_count": 0,
        "line_count": 0,
        "text_joined": "",
        "text_samples": text_samples,
        "dark_area_ratio": None,
        "non_white_ratio": None,
    }
    for shape in slide.shapes:
        features["shape_count"] += 1
        if getattr(shape, "has_text_frame", False):
            text = str(shape.text or "").strip()
            if text:
                features["text_shape_count"] += 1
                text_samples.append(text.replace("\n", " | ")[:120])
        if getattr(shape, "has_table", False):
            features["table_count"] += 1
        if getattr(shape, "has_chart", False):
            features["chart_count"] += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            features["picture_count"] += 1
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            features["line_count"] += 1
    features["text_joined"] = "\n".join(text_samples)
    image_features = _image_features(rendered_path)
    features.update(image_features)
    return features


def _image_features(path: Path | None) -> dict[str, float | None]:
    if path is None or not path.exists():
        return {"dark_area_ratio": None, "non_white_ratio": None}
    try:
        image = Image.open(path).convert("RGB").resize((180, 100))
    except OSError:
        return {"dark_area_ratio": None, "non_white_ratio": None}
    raw = image.tobytes()
    pixels = zip(raw[0::3], raw[1::3], raw[2::3])
    total = max(1, image.width * image.height)
    dark = 0
    non_white = 0
    for r, g, b in pixels:
        if (r + g + b) / 3 < 80:
            dark += 1
        if min(255 - r, 255 - g, 255 - b) > 12:
            non_white += 1
    stat = ImageStat.Stat(image)
    return {
        "dark_area_ratio": round(dark / total, 6),
        "non_white_ratio": round(non_white / total, 6),
        "mean_luma": round(sum(stat.mean) / 3, 4),
    }


def _content_slot_ids(spec: dict[str, Any]) -> set[str]:
    slots = []
    geometry = spec.get("title_zone_geometry")
    if isinstance(geometry, dict):
        slots.extend(geometry.get("slots") or [])
    fixture = spec.get("dummy_content_fixture") if isinstance(spec.get("dummy_content_fixture"), dict) else {}
    for block in fixture.get("content_blocks") or []:
        if isinstance(block, dict) and block.get("slot"):
            slots.append({"slot_id": block["slot"]})
    return {str(slot.get("slot_id") or "") for slot in slots if isinstance(slot, dict)}


def _has_image_zone(spec: dict[str, Any]) -> bool:
    zone = spec.get("image_photo_zone_geometry")
    return isinstance(zone, dict) and bool(zone.get("slots"))


def _has_two_digit_marker(text_samples: list[str]) -> bool:
    for text in text_samples:
        stripped = text.strip()
        if len(stripped) <= 3 and any(ch.isdigit() for ch in stripped):
            return True
    return False


def _inspection_preview_like(text: str) -> bool:
    return any(marker.lower() in text for marker in DEBUG_PREVIEW_STRINGS)


def _load_master_specs(directory: Path) -> dict[str, dict[str, Any]]:
    if not directory.exists():
        raise ValueError(f"golden master specs directory is missing: {directory.as_posix()}")
    specs: dict[str, dict[str, Any]] = {}
    for path in sorted(directory.glob("*.json")):
        payload = _load_json(path)
        if isinstance(payload, dict) and payload.get("archetype_id"):
            specs[str(payload["archetype_id"])] = payload
    if not specs:
        raise ValueError(f"no golden master specs found in {directory.as_posix()}")
    return specs


def _crop_paths_by_archetype(crops_dir: Path) -> dict[str, Path]:
    result: dict[str, Path] = {}
    if not crops_dir.exists():
        return result
    for path in crops_dir.glob("slide_thumbnail_*_*.png"):
        stem = path.stem
        parts = stem.split("_", 3)
        if len(parts) == 4:
            result[parts[3]] = path
    return result


def _render_paths_from_fidelity(fidelity: dict[str, Any]) -> dict[int, Path]:
    result: dict[int, Path] = {}
    for record in fidelity.get("layout_comparisons") or []:
        if isinstance(record, dict) and record.get("rendered_image_path"):
            result[int(record.get("slide_number") or len(result) + 1)] = Path(str(record["rendered_image_path"]))
    return result


def _rendered_path_from_record(record: dict[str, Any]) -> Path | None:
    if record.get("rendered_image_path"):
        return Path(str(record["rendered_image_path"]))
    return None


def _render_paths_from_render_report(render_report: dict[str, Any]) -> dict[int, Path]:
    result: dict[int, Path] = {}
    for record in render_report.get("slides") or []:
        if isinstance(record, dict) and record.get("rendered_image_path"):
            result[int(record.get("slide_index") or len(result) + 1)] = Path(str(record["rendered_image_path"]))
    for index, path in enumerate(render_report.get("output_paths") or [], start=1):
        result.setdefault(index, Path(str(path)))
    return result


def _render_paths_exist(paths: dict[int, Path]) -> bool:
    return bool(paths) and all(path.exists() for path in paths.values())


def _finding(
    code: str,
    severity: str,
    message: str,
    *,
    slide_number: int | None = None,
    archetype_id: str | None = None,
    details: dict[str, Any] | None = None,
) -> dict[str, Any]:
    payload: dict[str, Any] = {"code": code, "severity": severity, "message": message}
    if slide_number is not None:
        payload["slide_number"] = slide_number
    if archetype_id is not None:
        payload["archetype_id"] = archetype_id
    if details is not None:
        payload["details"] = details
    return payload


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Archetype Identity QA Report",
        "",
        f"Status: `{report['status']}`",
        f"Deck: `{report['golden_pptx_path']}`",
        f"Slide count: `{report['slide_count']}`",
        f"Findings: `{report['findings_summary']['total']}` total, `{report['findings_summary']['severe']}` severe",
        "",
        "## Slide Classification",
        "",
        "| Slide | Intended | Rendered family | Preserved | Drift target | Production | Inspection | Recommended fix |",
        "|---:|---|---|---|---|---|---|---|",
    ]
    for slide in report.get("slides") or []:
        lines.append(
            "| {slide_number} | `{intended_archetype}` | `{rendered_visual_family}` | `{identity_preserved}` | `{drift_target}` | `{production_master_like}` | `{inspection_preview_like}` | {recommended_fix} |".format(
                **{**slide, "drift_target": slide.get("drift_target") or ""}
            )
        )
    lines.extend(["", "## Severe Findings", ""])
    if report.get("severe_findings"):
        for finding in report["severe_findings"]:
            lines.append(f"- Slide {finding.get('slide_number')}: `{finding['code']}` {finding['message']}")
    else:
        lines.append("- None")
    lines.extend(["", "## Warnings", ""])
    if report.get("warnings"):
        for finding in report["warnings"]:
            lines.append(f"- Slide {finding.get('slide_number')}: `{finding['code']}` {finding['message']}")
    else:
        lines.append("- None")
    return "\n".join(lines) + "\n"


def _load_json(path: str | Path) -> Any:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _display_path(path: Path) -> str:
    return path.as_posix()


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Classify archetype identity drift in golden template masters.")
    parser.add_argument("--pptx", type=Path, default=DEFAULT_GOLDEN_PPTX)
    parser.add_argument("--crops-dir", type=Path, default=DEFAULT_CROPS_DIR)
    parser.add_argument("--master-specs-dir", type=Path, default=DEFAULT_MASTER_SPECS_DIR)
    parser.add_argument("--fidelity-report", type=Path, default=DEFAULT_FIDELITY_REPORT)
    parser.add_argument("--render-dir", type=Path, default=DEFAULT_RENDER_DIR)
    parser.add_argument("--render-manifest", type=Path, default=DEFAULT_RENDER_MANIFEST)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    parser.add_argument("--renderer", choices=("auto", "powerpoint_com", "libreoffice", "none"), default="auto")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        report = build_archetype_identity_report_from_files(
            golden_pptx_path=args.pptx,
            crops_dir=args.crops_dir,
            master_specs_dir=args.master_specs_dir,
            fidelity_report_path=args.fidelity_report,
            render_dir=args.render_dir,
            render_manifest_path=args.render_manifest,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
            renderer=args.renderer,
        )
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        print(f"ARCHETYPE_IDENTITY_QA_FAILED {exc}")
        return 1
    print(f"WROTE {DEFAULT_JSON_REPORT}")
    print(f"ARCHETYPE_IDENTITY_QA {report['status']}")
    return 1 if report["findings_summary"]["severe"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
