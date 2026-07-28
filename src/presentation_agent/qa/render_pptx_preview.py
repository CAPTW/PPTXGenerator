"""Render PPTX files to per-slide PNG previews using local QA backends."""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation


DEFAULT_PPTX_PATH = Path("outputs/template_preview.pptx")
DEFAULT_TEMPLATE_OUTPUT_DIR = Path("outputs/template_preview_png")
DEFAULT_FINAL_OUTPUT_DIR = Path("outputs/final_deck_preview_png")
DEFAULT_MANIFEST_PATH = Path("outputs/render_preview_manifest.json")
EMU_PER_INCH = 914400


def render_pptx_preview(
    *,
    pptx_path: str | Path = DEFAULT_PPTX_PATH,
    output_dir: str | Path = DEFAULT_TEMPLATE_OUTPUT_DIR,
    manifest_path: str | Path | None = None,
    report_path: str | Path | None = None,
    backend: str = "auto",
    renderer: str | None = None,
    dpi: int = 144,
) -> dict[str, Any]:
    """Render a PPTX into PNG previews, or write an explicit skip manifest."""

    selected_backend = renderer or backend
    pptx_file = Path(pptx_path)
    output = Path(output_dir)
    manifest = Path(manifest_path or report_path or DEFAULT_MANIFEST_PATH)
    output.mkdir(parents=True, exist_ok=True)
    _clear_pngs(output)

    slide_count = len(Presentation(pptx_file).slides)
    warnings: list[dict[str, Any]] = []
    errors: list[dict[str, Any]] = []
    attempts: list[dict[str, Any]] = []

    if selected_backend == "none":
        warnings.append(_warning("RENDERER_UNAVAILABLE", "Renderer disabled by request."))
        report = _manifest(
            pptx_file=pptx_file,
            output_dir=output,
            backend="skip",
            backend_path=None,
            render_status="skipped",
            slide_count=slide_count,
            slides=[],
            warnings=warnings,
            errors=errors,
            attempts=attempts,
        )
        _write_json(report, manifest)
        return report

    rendered_slides: list[dict[str, Any]] = []
    rendered_backend = "skip"
    rendered_backend_path: str | None = None

    for candidate in _candidate_backends(selected_backend):
        availability = _backend_availability(candidate)
        attempts.append(availability)
        if not availability["available"]:
            warnings.append(_warning("RENDERER_UNAVAILABLE", availability["reason"], backend=candidate))
            continue
        try:
            if candidate == "powerpoint_com":
                rendered_slides = _render_with_powerpoint_com(pptx_file, output, dpi=dpi)
            elif candidate == "libreoffice":
                rendered_slides = _render_with_libreoffice(pptx_file, output, Path(str(availability["path"])), dpi=dpi)
            else:
                warnings.append(_warning("RENDERER_UNSUPPORTED", f"Unsupported renderer backend: {candidate}", backend=candidate))
                continue
            rendered_backend = candidate
            rendered_backend_path = availability.get("path")
            break
        except Exception as exc:
            errors.append(_error("RENDER_BACKEND_FAILED", f"{candidate} render failed.", backend=candidate, details={"error": str(exc)}))
            _clear_pngs(output)
            rendered_slides = []

    if rendered_slides:
        render_status = "rendered"
        if len(rendered_slides) != slide_count:
            warnings.append(
                _warning(
                    "RENDER_SLIDE_COUNT_MISMATCH",
                    "Rendered PNG count does not match PPTX slide count.",
                    details={"expected": slide_count, "actual": len(rendered_slides)},
                )
            )
    else:
        render_status = "skipped"
        if not warnings and not errors:
            warnings.append(_warning("RENDERER_UNAVAILABLE", "No local renderer backend was available."))

    report = _manifest(
        pptx_file=pptx_file,
        output_dir=output,
        backend=rendered_backend,
        backend_path=rendered_backend_path,
        render_status=render_status,
        slide_count=slide_count,
        slides=rendered_slides,
        warnings=warnings,
        errors=errors,
        attempts=attempts,
    )
    _write_json(report, manifest)
    return report


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Render a PPTX to per-slide PNG previews using local QA backends.")
    parser.add_argument("--pptx", type=Path, default=DEFAULT_PPTX_PATH)
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_TEMPLATE_OUTPUT_DIR)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST_PATH)
    parser.add_argument("--report", type=Path, default=None, help="Compatibility alias for --manifest.")
    parser.add_argument("--backend", choices=("auto", "powerpoint_com", "libreoffice", "none"), default="auto")
    parser.add_argument("--renderer", choices=("auto", "powerpoint_com", "libreoffice", "none"), default=None)
    parser.add_argument("--dpi", type=int, default=144)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        report = render_pptx_preview(
            pptx_path=args.pptx,
            output_dir=args.output_dir,
            manifest_path=args.report or args.manifest,
            backend=args.backend,
            renderer=args.renderer,
            dpi=args.dpi,
        )
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        print(f"RENDER_PREVIEW_FAILED {exc}")
        return 1
    print(
        f"RENDER_PREVIEW status={report['render_status']} backend={report['backend']} "
        f"slides={report['slide_count']} rendered={report['rendered_slide_count']}"
    )
    return 0


def _candidate_backends(backend: str) -> list[str]:
    if backend == "auto":
        return ["powerpoint_com", "libreoffice"]
    return [backend]


def _backend_availability(backend: str) -> dict[str, Any]:
    if backend == "powerpoint_com":
        if os.name != "nt":
            return _availability(backend, False, None, "PowerPoint COM rendering is available only on Windows.")
        try:
            import pythoncom  # noqa: F401
            import win32com.client  # noqa: F401
        except ImportError:
            return _availability(backend, False, None, "pywin32 is not installed, so PowerPoint COM rendering is unavailable.")
        return _availability(backend, True, "Microsoft PowerPoint COM", None)
    if backend == "libreoffice":
        soffice = _find_soffice()
        if soffice is None:
            return _availability(backend, False, None, "LibreOffice soffice executable was not found locally.")
        return _availability(backend, True, str(soffice), None)
    if backend == "none":
        return _availability(backend, False, None, "Renderer disabled by request.")
    return _availability(backend, False, None, f"Unsupported renderer backend: {backend}")


def _render_with_powerpoint_com(pptx_file: Path, output_dir: Path, *, dpi: int) -> list[dict[str, Any]]:
    import pythoncom
    import win32com.client

    presentation_probe = Presentation(pptx_file)
    width_px = max(1, round(presentation_probe.slide_width / EMU_PER_INCH * dpi))
    height_px = max(1, round(presentation_probe.slide_height / EMU_PER_INCH * dpi))

    pythoncom.CoInitialize()
    app = None
    deck = None
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        deck = app.Presentations.Open(str(pptx_file.resolve()), WithWindow=False)
        slides: list[dict[str, Any]] = []
        for index in range(1, int(deck.Slides.Count) + 1):
            output_path = output_dir / f"slide-{index:03d}.png"
            deck.Slides(index).Export(str(output_path.resolve()), "PNG", width_px, height_px)
            slides.append(_slide_record(index, output_path, width_px, height_px, "powerpoint_com"))
        return slides
    finally:
        if deck is not None:
            deck.Close()
        if app is not None:
            app.Quit()
        pythoncom.CoUninitialize()


def _render_with_libreoffice(pptx_file: Path, output_dir: Path, soffice_path: Path, *, dpi: int) -> list[dict[str, Any]]:
    with tempfile.TemporaryDirectory() as temp_dir_text:
        temp_dir = Path(temp_dir_text)
        command = [
            str(soffice_path),
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(temp_dir),
            str(pptx_file.resolve()),
        ]
        result = subprocess.run(command, cwd=pptx_file.parent, capture_output=True, text=True, check=False, timeout=120)
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stdout} {result.stderr}".strip())
        pdf_candidates = sorted(temp_dir.glob("*.pdf"))
        if not pdf_candidates:
            raise RuntimeError("LibreOffice did not produce a PDF file.")
        pdf_path = pdf_candidates[0]
        document = fitz.open(pdf_path)
        slides: list[dict[str, Any]] = []
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            output_path = output_dir / f"slide-{index:03d}.png"
            pixmap.save(output_path)
            slides.append(_slide_record(index, output_path, pixmap.width, pixmap.height, "libreoffice"))
        document.close()
        return slides


def _find_soffice() -> Path | None:
    env_path = os.environ.get("LIBREOFFICE_PATH")
    candidates = [Path(env_path)] if env_path else []
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            candidates.append(Path(found))
    candidates.extend(
        [
            Path("C:/Program Files/LibreOffice/program/soffice.exe"),
            Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
        ]
    )
    for candidate in candidates:
        if str(candidate) and candidate.is_file():
            return candidate
    return None


def _manifest(
    *,
    pptx_file: Path,
    output_dir: Path,
    backend: str,
    backend_path: str | None,
    render_status: str,
    slide_count: int,
    slides: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
    errors: list[dict[str, Any]],
    attempts: list[dict[str, Any]],
) -> dict[str, Any]:
    output_paths = [str(record["rendered_image_path"]) for record in slides]
    findings = [*_as_findings(warnings, "warning"), *_as_findings(errors, "error")]
    return {
        "schema_name": "render_preview_manifest",
        "schema_version": "1.0",
        "input_pptx": _display_path(pptx_file),
        "pptx_path": _display_path(pptx_file),
        "output_dir": _display_path(output_dir),
        "backend": backend,
        "render_status": render_status,
        "rendered_slide_count": len(slides),
        "output_paths": output_paths,
        "warnings": warnings,
        "errors": errors,
        "backend_attempts": attempts,
        "slide_count": slide_count,
        "slides": slides,
        "status": render_status,
        "renderer": {
            "name": backend,
            "available": render_status == "rendered",
            "path": backend_path,
            "reason": None if render_status == "rendered" else "No local renderer completed successfully.",
        },
        "findings": findings,
    }


def _slide_record(index: int, output_path: Path, width_px: int, height_px: int, backend: str) -> dict[str, Any]:
    return {
        "slide_index": index,
        "rendered_image_path": _display_path(output_path),
        "width_px": width_px,
        "height_px": height_px,
        "backend": backend,
    }


def _availability(backend: str, available: bool, path: str | None, reason: str | None) -> dict[str, Any]:
    return {"backend": backend, "available": available, "path": path, "reason": reason}


def _warning(code: str, message: str, *, backend: str | None = None, details: dict[str, Any] | None = None) -> dict[str, Any]:
    payload: dict[str, Any] = {"code": code, "severity": "warning", "message": message}
    if backend:
        payload["backend"] = backend
    if details:
        payload["details"] = details
    return payload


def _error(code: str, message: str, *, backend: str | None = None, details: dict[str, Any] | None = None) -> dict[str, Any]:
    payload: dict[str, Any] = {"code": code, "severity": "error", "message": message}
    if backend:
        payload["backend"] = backend
    if details:
        payload["details"] = details
    return payload


def _as_findings(records: list[dict[str, Any]], default_severity: str) -> list[dict[str, Any]]:
    findings = []
    for record in records:
        finding = dict(record)
        finding.setdefault("severity", default_severity)
        findings.append(finding)
    return findings


def _clear_pngs(output_dir: Path) -> None:
    for path in output_dir.glob("*.png"):
        path.unlink()


def _write_json(payload: dict[str, Any], path: str | Path) -> None:
    destination = Path(path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_text(json.dumps(payload, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")


def _display_path(path: Path) -> str:
    return str(path.as_posix())


if __name__ == "__main__":
    raise SystemExit(main())
