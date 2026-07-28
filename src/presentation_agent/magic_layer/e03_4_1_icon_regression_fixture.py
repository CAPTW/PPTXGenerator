"""PowerPoint-rendered v7.1 icon regression fixture."""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.image import ImagePart
from pptx.util import Inches, Pt

SVG_CONTENT_TYPE = "image/svg+xml"
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
RENDER_DPI = 144
SIZE_ROWS = ((16, 16 / RENDER_DPI), (24, 24 / RENDER_DPI), (32, 32 / RENDER_DPI))


def build_icon_regression_fixture_v7_1(
    curated_v7_1_manifest: dict[str, Any],
    output_dir: Path,
    *,
    render: bool = True,
) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "magic_layer_v7_1_icon_regression_fixture.pptx"
    contact_sheet = output_dir / "magic_layer_v7_1_icon_regression_fixture_contact_sheet.png"
    roles = sorted(curated_v7_1_manifest.get("roles", []), key=lambda row: (row.get("priority", ""), row.get("role_id", "")))
    variants_by_role = curated_v7_1_manifest.get("variants_by_role", {})
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    cells: list[dict[str, Any]] = []
    for slide_index, (size_px, size_in) in enumerate(SIZE_ROWS, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cells.extend(_add_fixture_slide(slide, roles, variants_by_role, slide_index, size_px, size_in))
    prs.save(pptx_path)

    render_report: dict[str, Any] = {"render_status": "not_requested", "slides": []}
    if render:
        render_report = _render_fixture(pptx_path, output_dir / "rendered")
        _copy_named_renders(render_report, output_dir, cells)
    _build_contact_sheet(render_report, contact_sheet, roles)
    cells_path = output_dir / "magic_layer_v7_1_icon_regression_fixture_cells.json"
    cells_path.write_text(json.dumps({"cells": cells}, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    rendered_slide_count = int(render_report.get("rendered_slide_count") or len(render_report.get("slides", [])))
    return {
        "schema_name": "icon_regression_fixture_v7_1_report",
        "status": "passed" if pptx_path.exists() and (not render or rendered_slide_count == 3) else "blocked",
        "fixture_pptx_path": pptx_path.as_posix(),
        "cell_manifest_path": cells_path.as_posix(),
        "contact_sheet_path": contact_sheet.as_posix(),
        "fixture_slide_count": 3,
        "rendered_slide_count": rendered_slide_count,
        "fixture_rendered": (not render) or rendered_slide_count == 3,
        "cell_count": len(cells),
        "svg_media_icon_count": len(cells),
        "native_vector_conversion_count": 0,
        "semantic_raster_icon_count": 0,
        "render_report": render_report,
        "cells": cells,
    }


def _add_fixture_slide(
    slide: Any,
    roles: list[dict[str, Any]],
    variants_by_role: dict[str, dict[str, str]],
    slide_index: int,
    size_px: int,
    icon_size_in: float,
) -> list[dict[str, Any]]:
    _add_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, RGBColor(255, 255, 255))
    dark_x = 6.74
    _add_rect(slide, dark_x, 0.46, 6.33, 6.86, RGBColor(7, 16, 24))
    title = slide.shapes.add_textbox(Inches(0.25), Inches(0.12), Inches(7), Inches(0.25))
    title.text_frame.text = f"Magic Layer v7.1 PowerPoint SVG Fixture - {size_px}px"
    title.text_frame.paragraphs[0].font.size = Pt(13)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)
    cols = 6
    x0 = 0.28
    y0 = 0.58
    cell_w = 1.04
    cell_h = 0.58
    cells: list[dict[str, Any]] = []
    for idx, role in enumerate(roles):
        col = idx % cols
        row = idx // cols
        for background, offset, label_color, variant_name in (
            ("light", 0.0, RGBColor(15, 23, 42), "light"),
            ("dark", 6.7, RGBColor(248, 250, 252), "dark"),
        ):
            x = x0 + offset + col * cell_w
            y = y0 + row * cell_h
            icon_x = x + 0.04
            icon_y = y + 0.06
            svg_path = Path(variants_by_role[role["role_id"]][variant_name])
            _add_svg(slide, svg_path, icon_x, icon_y, icon_size_in, icon_size_in, role["role_id"], background, size_px)
            _add_label(slide, role["role_id"], x + 0.30, y + 0.03, label_color)
            cells.append(
                {
                    "role_id": role["role_id"],
                    "priority": role.get("priority"),
                    "background": background,
                    "size_px": size_px,
                    "slide_index": slide_index,
                    "bbox_px": _bbox_px(icon_x, icon_y, icon_size_in, icon_size_in, margin_px=6),
                    "icon_svg_path": svg_path.as_posix(),
                    "insertion_route": "true_svg_media_insertion",
                }
            )
    return cells


def _add_rect(slide: Any, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def _add_label(slide: Any, role_id: str, x: float, y: float, color: RGBColor) -> None:
    label = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.70), Inches(0.18))
    label.text_frame.text = role_id[:18]
    p = label.text_frame.paragraphs[0]
    p.font.size = Pt(5.6)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def _add_svg(slide: Any, svg_path: Path, x: float, y: float, w: float, h: float, role_id: str, background: str, size_px: int) -> None:
    image_part = ImagePart(
        slide.part.package.next_image_partname("svg"),
        SVG_CONTENT_TYPE,
        slide.part.package,
        svg_path.read_bytes(),
        svg_path.name,
    )
    r_id = slide.part.relate_to(image_part, RT.IMAGE)
    shape_id = slide.shapes._next_shape_id
    safe_role = role_id.replace("&", "and").replace("<", "").replace(">", "")
    slide.shapes._grpSp.add_pic(
        shape_id,
        f"SVG Icon {safe_role} {background} {size_px}px",
        f"svg-icon:{safe_role}:{background}:{size_px}",
        r_id,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    slide.shapes._recalculate_extents()


def _bbox_px(x: float, y: float, w: float, h: float, *, margin_px: int) -> list[int]:
    return [
        max(0, round(x * RENDER_DPI) - margin_px),
        max(0, round(y * RENDER_DPI) - margin_px),
        min(round(SLIDE_W_IN * RENDER_DPI), round((x + w) * RENDER_DPI) + margin_px),
        min(round(SLIDE_H_IN * RENDER_DPI), round((y + h) * RENDER_DPI) + margin_px),
    ]


def _render_fixture(pptx_path: Path, render_dir: Path) -> dict[str, Any]:
    try:
        from src.presentation_agent.qa.render_pptx_preview import render_pptx_preview

        return render_pptx_preview(
            pptx_path=pptx_path,
            output_dir=render_dir,
            manifest_path=render_dir / "render_manifest.json",
            backend="auto",
            dpi=RENDER_DPI,
        )
    except Exception as exc:  # pragma: no cover - renderer environment dependent
        return {"render_status": "skipped", "errors": [{"error": str(exc)}], "slides": []}


def _copy_named_renders(render_report: dict[str, Any], output_dir: Path, cells: list[dict[str, Any]]) -> None:
    names_by_slide = {
        1: "magic_layer_v7_1_icon_regression_fixture_slide_16px.png",
        2: "magic_layer_v7_1_icon_regression_fixture_slide_24px.png",
        3: "magic_layer_v7_1_icon_regression_fixture_slide_32px.png",
    }
    paths_by_slide: dict[int, str] = {}
    for row in render_report.get("slides", []):
        slide_index = int(row.get("slide_index", 0))
        source = Path(row.get("rendered_image_path") or "")
        if source.exists() and slide_index in names_by_slide:
            dest = output_dir / names_by_slide[slide_index]
            shutil.copy2(source, dest)
            paths_by_slide[slide_index] = dest.as_posix()
    for cell in cells:
        cell["render_path"] = paths_by_slide.get(cell["slide_index"], "")


def _build_contact_sheet(render_report: dict[str, Any], output: Path, roles: list[dict[str, Any]]) -> None:
    slides = [Path(row.get("rendered_image_path") or "") for row in render_report.get("slides", [])]
    images = [Image.open(path).convert("RGB") for path in slides if path.exists()]
    if not images:
        image = Image.new("RGB", (1000, 320), "#071018")
        draw = ImageDraw.Draw(image)
        draw.text((24, 24), f"Fixture PPTX created for {len(roles)} roles; render unavailable", fill="#F8FAFC", font=ImageFont.load_default())
        image.save(output)
        return
    thumbs = []
    for image in images:
        image.thumbnail((480, 270), Image.Resampling.LANCZOS)
        thumbs.append(image.copy())
    sheet = Image.new("RGB", (500 * len(thumbs), 300), "#071018")
    for idx, image in enumerate(thumbs):
        sheet.paste(image, (idx * 500 + 10, 15))
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)
