"""Visual clutter and shape budget guard for editable PPTX decks."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_GOLDEN_PPTX = Path("outputs/golden_template_masters.pptx")
DEFAULT_FINAL_PPTX = Path("outputs/final_deck_large_premium.pptx")
DEFAULT_CONTRACTS_DIR = Path("outputs/template_contracts")
DEFAULT_GOLDEN_REPORT = Path("outputs/golden_template_masters_report.json")
DEFAULT_TEMPLATE_USABILITY_REPORT = Path("outputs/template_usability_report.json")
DEFAULT_JSON_REPORT = Path("outputs/visual_clutter_report.json")
DEFAULT_MD_REPORT = Path("outputs/visual_clutter_report.md")

EMU_PER_INCH = 914400
SLIDE_AREA_16_9 = 13.333 * 7.5
FOOTER_Y = 6.45

DENSITY_LIMITS = {
    "low": {"decoration_to_text_ratio": 7.0, "line_count": 18, "chrome_area_ratio": 0.42},
    "medium": {"decoration_to_text_ratio": 9.0, "line_count": 38, "chrome_area_ratio": 0.52},
    "high": {"decoration_to_text_ratio": 12.0, "line_count": 78, "chrome_area_ratio": 0.68},
}

ARCHETYPE_ALIASES = {
    "closing": "creative_cover",
    "concept_relationship": "concept_relationship_venn",
    "cover": "creative_cover",
    "data_dashboard": "kpi_donut_chart",
    "sequence": "work_support_sequence",
    "visual_toc": "visual_table_of_contents",
}


def build_visual_clutter_report_from_files(
    *,
    golden_pptx_path: str | Path = DEFAULT_GOLDEN_PPTX,
    final_pptx_path: str | Path = DEFAULT_FINAL_PPTX,
    contracts_dir: str | Path = DEFAULT_CONTRACTS_DIR,
    golden_report_path: str | Path = DEFAULT_GOLDEN_REPORT,
    template_usability_report_path: str | Path = DEFAULT_TEMPLATE_USABILITY_REPORT,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
) -> dict[str, Any]:
    report = build_visual_clutter_report(
        golden_pptx_path=golden_pptx_path,
        final_pptx_path=final_pptx_path,
        contracts_dir=contracts_dir,
        golden_report_path=golden_report_path,
        template_usability_report_path=template_usability_report_path,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return report


def build_visual_clutter_report(
    *,
    golden_pptx_path: str | Path,
    final_pptx_path: str | Path,
    contracts_dir: str | Path,
    golden_report_path: str | Path,
    template_usability_report_path: str | Path,
) -> dict[str, Any]:
    contracts = _load_contracts(Path(contracts_dir))
    golden_map = _golden_archetype_map(Path(golden_report_path))
    usability_map = _usability_archetype_map(Path(template_usability_report_path))
    decks = []
    findings: list[dict[str, Any]] = []
    for scope, pptx_path, archetype_map in (
        ("golden_template_masters", Path(golden_pptx_path), golden_map),
        ("final_deck_large_premium", Path(final_pptx_path), usability_map),
    ):
        deck_report = _inspect_deck(scope, pptx_path, archetype_map, contracts)
        decks.append(deck_report)
        findings.extend(deck_report["findings"])
    severe = sum(1 for finding in findings if finding["severity"] == "severe")
    warning = sum(1 for finding in findings if finding["severity"] == "warning")
    return {
        "schema_name": "visual_clutter_report",
        "schema_version": "1.0",
        "status": "failed" if severe else "issues_reported" if warning else "passed",
        "findings_summary": {"total": len(findings), "severe": severe, "warning": warning},
        "contracts_dir": _display_path(Path(contracts_dir)),
        "golden_report_path": _display_path(Path(golden_report_path)),
        "template_usability_report_path": _display_path(Path(template_usability_report_path)),
        "decks": decks,
        "findings": findings,
        "guard_policy": {
            "shape_count_exceeds_contract_budget": "severe",
            "decorative_objects_exceed_max_ornament_density": "severe",
            "ornament_overlaps_primary_text_zone": "severe",
            "background_complexity_reduces_readability": "severe",
            "footer_citation_strip_dominates_content": "severe",
            "card_chrome_consumes_more_area_than_card_text": "severe",
            "inspection_like_rendering": "severe",
            "potential_editability_performance_risk": "warning",
        },
    }


def _inspect_deck(scope: str, pptx_path: Path, archetype_map: dict[int, str], contracts: dict[str, dict[str, Any]]) -> dict[str, Any]:
    if not pptx_path.exists():
        return {
            "scope": scope,
            "pptx_path": _display_path(pptx_path),
            "status": "skipped",
            "slide_count": 0,
            "slides": [],
            "findings": [],
            "findings_summary": {"total": 0, "severe": 0, "warning": 0},
        }
    deck = Presentation(pptx_path)
    slide_w = deck.slide_width / EMU_PER_INCH
    slide_h = deck.slide_height / EMU_PER_INCH
    slide_area = max(0.01, slide_w * slide_h)
    slides: list[dict[str, Any]] = []
    findings: list[dict[str, Any]] = []
    for index, slide in enumerate(deck.slides, start=1):
        archetype_id = _canonical_archetype_id(archetype_map.get(index, "unknown"))
        contract = contracts.get(archetype_id, {})
        stats = _slide_stats(slide, slide_area)
        slide_findings = _slide_findings(scope, index, archetype_id, stats, contract)
        findings.extend(slide_findings)
        slides.append(
            {
                "slide_number": index,
                "archetype_id": archetype_id,
                "contract_used": contract.get("_contract_path"),
                "metrics": stats,
                "finding_count": len(slide_findings),
                "finding_codes": [finding["code"] for finding in slide_findings],
            }
        )
    severe = sum(1 for finding in findings if finding["severity"] == "severe")
    warning = sum(1 for finding in findings if finding["severity"] == "warning")
    return {
        "scope": scope,
        "pptx_path": _display_path(pptx_path),
        "status": "failed" if severe else "issues_reported" if warning else "passed",
        "slide_count": len(slides),
        "findings_summary": {"total": len(findings), "severe": severe, "warning": warning},
        "slides": slides,
        "findings": findings,
    }


def _slide_stats(slide: Any, slide_area: float) -> dict[str, Any]:
    shape_count = 0
    line_count = 0
    picture_count = 0
    table_count = 0
    chart_count = 0
    text_shapes: list[dict[str, float]] = []
    chrome_shapes: list[dict[str, float]] = []
    ornament_shapes: list[dict[str, float]] = []
    footer_area = 0.0
    total_chrome_area = 0.0
    card_chrome_area = 0.0
    card_text_area = 0.0
    inspection_text_hits = 0
    for shape in slide.shapes:
        shape_count += 1
        bounds = _shape_bounds(shape)
        area = bounds["w"] * bounds["h"]
        area_ratio = area / slide_area
        has_text = bool(getattr(shape, "has_text_frame", False) and str(shape.text or "").strip())
        if not has_text and shape.shape_type != MSO_SHAPE_TYPE.PICTURE and area_ratio >= 0.82:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            line_count += 1
            ornament_shapes.append(bounds)
            total_chrome_area += min(area, 0.04)
        elif getattr(shape, "has_table", False):
            table_count += 1
            chrome_shapes.append(bounds)
            total_chrome_area += area * 0.1
        elif getattr(shape, "has_chart", False):
            chart_count += 1
            chrome_shapes.append(bounds)
            total_chrome_area += area * 0.1
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
        elif has_text:
            text_shapes.append(bounds)
            text = str(shape.text or "").lower()
            if any(token in text for token in ("safe margins", "archetype:", "density:", "layout-board-", "slot_id", "component_id")):
                inspection_text_hits += 1
            if _looks_like_card_text(bounds):
                card_text_area += area
        else:
            chrome_shapes.append(bounds)
            if _looks_like_small_ornament(bounds):
                total_chrome_area += area
            elif _looks_like_card_panel(bounds):
                total_chrome_area += area * 0.18
            else:
                total_chrome_area += min(area * 0.08, 0.25)
            if _looks_like_small_ornament(bounds):
                ornament_shapes.append(bounds)
            if _looks_like_card_panel(bounds):
                card_chrome_area += area
        footer_area += _footer_intersection_area(bounds)
    decorative_count = max(0, shape_count - len(text_shapes) - table_count - chart_count - picture_count)
    primary_text_shapes = [bounds for bounds in text_shapes if bounds["y"] < FOOTER_Y]
    ornament_intrusions = _ornament_intrusions(ornament_shapes, primary_text_shapes)
    return {
        "shape_count": shape_count,
        "line_ornament_count": line_count,
        "ornament_shape_count": len(ornament_shapes),
        "text_box_count": len(text_shapes),
        "card_count_estimate": _estimated_card_count(text_shapes),
        "background_ornament_density": round(len(ornament_shapes) / max(1.0, slide_area), 3),
        "footer_citation_footprint_ratio": round(min(1.0, footer_area / slide_area), 4),
        "visual_chrome_area_ratio": round(min(1.0, total_chrome_area / slide_area), 4),
        "text_zone_intrusion_count": len(ornament_intrusions),
        "text_zone_intrusions": ornament_intrusions[:5],
        "shape_to_content_ratio": round(decorative_count / max(1, len(text_shapes) + table_count + chart_count), 3),
        "decorative_shape_count": decorative_count,
        "picture_count": picture_count,
        "table_count": table_count,
        "chart_count": chart_count,
        "card_chrome_area_ratio": round(card_chrome_area / max(0.01, card_text_area), 3) if card_text_area else 0.0,
        "inspection_text_hits": inspection_text_hits,
        "potential_editability_performance_risk": shape_count > 350 or decorative_count > 260,
    }


def _slide_findings(scope: str, slide_number: int, archetype_id: str, stats: dict[str, Any], contract: dict[str, Any]) -> list[dict[str, Any]]:
    findings: list[dict[str, Any]] = []
    budget = contract.get("decorative_budget") if isinstance(contract.get("decorative_budget"), dict) else {}
    max_shapes = int(budget.get("max_shape_count_target") or 0)
    density = str(budget.get("max_ornament_density") or "medium").lower()
    density_limits = DENSITY_LIMITS.get(density, DENSITY_LIMITS["medium"])
    if max_shapes and stats["shape_count"] > max_shapes:
        findings.append(_finding("SHAPE_COUNT_EXCEEDS_CONTRACT_BUDGET", "severe", scope, slide_number, archetype_id, "Shape count exceeds the template contract budget.", {"shape_count": stats["shape_count"], "limit": max_shapes}))
    if stats["shape_to_content_ratio"] > density_limits["decoration_to_text_ratio"]:
        findings.append(_finding("DECORATIVE_OBJECTS_EXCEED_ORNAMENT_DENSITY", "severe", scope, slide_number, archetype_id, "Decorative object volume exceeds the max ornament density.", {"shape_to_content_ratio": stats["shape_to_content_ratio"], "limit": density_limits["decoration_to_text_ratio"], "max_ornament_density": density}))
    if stats["line_ornament_count"] > density_limits["line_count"]:
        findings.append(_finding("BACKGROUND_COMPLEXITY_REDUCES_READABILITY", "severe", scope, slide_number, archetype_id, "Line and ornament count may reduce readability.", {"line_ornament_count": stats["line_ornament_count"], "limit": density_limits["line_count"]}))
    if stats["visual_chrome_area_ratio"] > density_limits["chrome_area_ratio"]:
        findings.append(_finding("BACKGROUND_COMPLEXITY_REDUCES_READABILITY", "severe", scope, slide_number, archetype_id, "Visual chrome consumes too much slide area.", {"visual_chrome_area_ratio": stats["visual_chrome_area_ratio"], "limit": density_limits["chrome_area_ratio"]}))
    if stats["text_zone_intrusion_count"] > 0:
        findings.append(_finding("ORNAMENT_OVERLAPS_PRIMARY_TEXT_ZONE", "severe", scope, slide_number, archetype_id, "Ornament geometry overlaps primary text zones.", {"intrusion_count": stats["text_zone_intrusion_count"], "examples": stats["text_zone_intrusions"]}))
    if stats["footer_citation_footprint_ratio"] > 0.16:
        findings.append(_finding("FOOTER_CITATION_STRIP_DOMINATES_CONTENT", "severe", scope, slide_number, archetype_id, "Footer/citation footprint dominates the slide content area.", {"footer_citation_footprint_ratio": stats["footer_citation_footprint_ratio"], "limit": 0.16}))
    if stats["card_chrome_area_ratio"] > 1.25:
        findings.append(_finding("CARD_CHROME_EXCEEDS_CARD_TEXT_AREA", "severe", scope, slide_number, archetype_id, "Card chrome consumes more area than card text.", {"card_chrome_area_ratio": stats["card_chrome_area_ratio"], "limit": 1.25}))
    if stats["inspection_text_hits"] > 0:
        findings.append(_finding("INSPECTION_LIKE_RENDERING", "severe", scope, slide_number, archetype_id, "Slide contains inspection/debug text.", {"inspection_text_hits": stats["inspection_text_hits"]}))
    if stats["potential_editability_performance_risk"]:
        findings.append(_finding("POTENTIAL_EDITABILITY_PERFORMANCE_RISK", "warning", scope, slide_number, archetype_id, "High object count may make the slide harder to edit.", {"shape_count": stats["shape_count"], "decorative_shape_count": stats["decorative_shape_count"]}))
    return findings


def _load_contracts(directory: Path) -> dict[str, dict[str, Any]]:
    contracts: dict[str, dict[str, Any]] = {}
    if not directory.exists():
        return contracts
    for path in sorted(directory.glob("*.contract.json")):
        payload = json.loads(path.read_text(encoding="utf-8"))
        if isinstance(payload, dict) and payload.get("archetype_id"):
            payload = dict(payload)
            payload["_contract_path"] = _display_path(path)
            contracts[str(payload["archetype_id"])] = payload
    return contracts


def _canonical_archetype_id(archetype_id: str) -> str:
    return ARCHETYPE_ALIASES.get(str(archetype_id), str(archetype_id))


def _golden_archetype_map(path: Path) -> dict[int, str]:
    payload = _load_json(path)
    result: dict[int, str] = {}
    for item in payload.get("compiled_layouts") or []:
        if isinstance(item, dict) and item.get("slide_number"):
            result[int(item["slide_number"])] = str(item.get("archetype_id") or "unknown")
    return result


def _usability_archetype_map(path: Path) -> dict[int, str]:
    payload = _load_json(path)
    result: dict[int, str] = {}
    final = payload.get("final_deck_large_premium") if isinstance(payload.get("final_deck_large_premium"), dict) else {}
    for item in final.get("slides") or []:
        if isinstance(item, dict) and item.get("slide_number"):
            result[int(item["slide_number"])] = str(item.get("archetype_id") or "unknown")
    return result


def _shape_bounds(shape: Any) -> dict[str, float]:
    return {
        "x": float(shape.left) / EMU_PER_INCH,
        "y": float(shape.top) / EMU_PER_INCH,
        "w": max(0.001, float(shape.width) / EMU_PER_INCH),
        "h": max(0.001, float(shape.height) / EMU_PER_INCH),
    }


def _looks_like_small_ornament(bounds: dict[str, float]) -> bool:
    area = bounds["w"] * bounds["h"]
    return area <= 0.18 or bounds["w"] <= 0.12 or bounds["h"] <= 0.12


def _looks_like_card_panel(bounds: dict[str, float]) -> bool:
    area = bounds["w"] * bounds["h"]
    return 0.45 <= area <= 4.5 and bounds["w"] >= 0.9 and bounds["h"] >= 0.45


def _looks_like_card_text(bounds: dict[str, float]) -> bool:
    return 0.45 <= bounds["h"] <= 1.8 and 0.7 <= bounds["w"] <= 4.2


def _estimated_card_count(text_shapes: list[dict[str, float]]) -> int:
    return sum(1 for bounds in text_shapes if _looks_like_card_text(bounds))


def _ornament_intrusions(ornaments: list[dict[str, float]], text_shapes: list[dict[str, float]]) -> list[dict[str, Any]]:
    intrusions: list[dict[str, Any]] = []
    for ornament_index, ornament in enumerate(ornaments, start=1):
        for text_index, text in enumerate(text_shapes, start=1):
            overlap = _intersection_area(ornament, text)
            if overlap > 0.01 and overlap / max(0.001, ornament["w"] * ornament["h"]) > 0.25:
                intrusions.append({"ornament_index": ornament_index, "text_index": text_index, "overlap_area": round(overlap, 4)})
                break
    return intrusions


def _intersection_area(a: dict[str, float], b: dict[str, float]) -> float:
    x0 = max(a["x"], b["x"])
    y0 = max(a["y"], b["y"])
    x1 = min(a["x"] + a["w"], b["x"] + b["w"])
    y1 = min(a["y"] + a["h"], b["y"] + b["h"])
    if x1 <= x0 or y1 <= y0:
        return 0.0
    return (x1 - x0) * (y1 - y0)


def _footer_intersection_area(bounds: dict[str, float]) -> float:
    footer = {"x": 0.0, "y": FOOTER_Y, "w": 13.333, "h": 7.5 - FOOTER_Y}
    return _intersection_area(bounds, footer)


def _finding(code: str, severity: str, scope: str, slide_number: int, archetype_id: str, message: str, details: dict[str, Any]) -> dict[str, Any]:
    return {
        "code": code,
        "severity": severity,
        "scope": scope,
        "slide_number": slide_number,
        "archetype_id": archetype_id,
        "message": message,
        "details": details,
    }


def _markdown_report(report: dict[str, Any]) -> str:
    summary = report["findings_summary"]
    lines = [
        "# Visual Clutter Guard Report",
        "",
        f"Status: `{report['status']}`",
        f"Findings: `{summary['total']}` total, `{summary['severe']}` severe, `{summary['warning']}` warnings",
        "",
        "## Decks",
        "",
    ]
    for deck in report.get("decks") or []:
        deck_summary = deck.get("findings_summary") or {}
        lines.append(f"- `{deck.get('scope')}`: `{deck.get('status')}`, slides `{deck.get('slide_count')}`, severe `{deck_summary.get('severe')}`, warnings `{deck_summary.get('warning')}`")
    lines.extend(["", "## Findings", ""])
    for finding in report.get("findings") or []:
        lines.append(f"- `{finding['severity']}` `{finding['code']}` slide `{finding['slide_number']}` `{finding['scope']}` `{finding['archetype_id']}`: {finding['message']}")
    return "\n".join(lines) + "\n"


def _load_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    payload = json.loads(path.read_text(encoding="utf-8"))
    return payload if isinstance(payload, dict) else {}


def _display_path(path: Path) -> str:
    try:
        return str(path.relative_to(Path.cwd())).replace("\\", "/")
    except ValueError:
        return str(path).replace("\\", "/")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Guard editable decks against visual clutter and shape-budget regressions.")
    parser.add_argument("--golden-pptx", type=Path, default=DEFAULT_GOLDEN_PPTX)
    parser.add_argument("--final-pptx", type=Path, default=DEFAULT_FINAL_PPTX)
    parser.add_argument("--contracts-dir", type=Path, default=DEFAULT_CONTRACTS_DIR)
    parser.add_argument("--golden-report", type=Path, default=DEFAULT_GOLDEN_REPORT)
    parser.add_argument("--template-usability-report", type=Path, default=DEFAULT_TEMPLATE_USABILITY_REPORT)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        report = build_visual_clutter_report_from_files(
            golden_pptx_path=args.golden_pptx,
            final_pptx_path=args.final_pptx,
            contracts_dir=args.contracts_dir,
            golden_report_path=args.golden_report,
            template_usability_report_path=args.template_usability_report,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"VISUAL_CLUTTER_GUARD_FAILED {type(exc).__name__}: {exc}")
        return 1
    print(f"WROTE {args.json_report}")
    print(f"WROTE {args.md_report}")
    print(f"VISUAL_CLUTTER_GUARD {report['status']}")
    return 1 if report["status"] == "failed" else 0


if __name__ == "__main__":
    raise SystemExit(main())
