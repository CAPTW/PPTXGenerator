"""Object geometry extraction for D07.1 layout reflow audits."""

from __future__ import annotations

import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_object_geometry_ledger(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    objects: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for z_order, shape in enumerate(slide.shapes, start=1):
            text = shape.text_frame.text if getattr(shape, "has_text_frame", False) and shape.text_frame is not None else ""
            record = {
                "slide_index": slide_index,
                "slide_id": f"d07_slide_{slide_index:02d}",
                "object_index": z_order,
                "z_order": z_order,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "object_type": _object_type(shape),
                "role": infer_geometry_role(shape.name, text),
                "x_in": round(float(shape.left) / 914400, 4),
                "y_in": round(float(shape.top) / 914400, 4),
                "w_in": round(float(shape.width) / 914400, 4),
                "h_in": round(float(shape.height) / 914400, 4),
                "bbox_norm": [
                    round(float(shape.left) / slide_w, 6),
                    round(float(shape.top) / slide_h, 6),
                    round(float(shape.width) / slide_w, 6),
                    round(float(shape.height) / slide_h, 6),
                ],
                "has_text": bool(text),
                "text": text,
                "text_length": len(text),
                "content_bearing": _content_bearing(shape.name, text),
                "is_picture": shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
            }
            objects.append(record)
    return {
        "schema_name": "object_geometry_ledger",
        "status": "passed",
        "pptx_path": pptx_path.as_posix(),
        "slide_count": len(prs.slides),
        "object_count": len(objects),
        "objects": objects,
        "media_count": _pptx_media_count(pptx_path),
    }


def split_geometry_ledgers(object_ledger: dict[str, Any]) -> dict[str, dict[str, Any]]:
    objects = object_ledger.get("objects") or []
    text_objects = [obj for obj in objects if obj.get("has_text")]
    source_footer = [obj for obj in objects if obj.get("role") in {"source_footer_text", "source_footer_strip"}]
    visual_objects = [obj for obj in objects if not obj.get("has_text")]
    return {
        "text_box_geometry_ledger": {
            "schema_name": "text_box_geometry_ledger",
            "status": "passed",
            "text_box_count": len(text_objects),
            "objects": text_objects,
        },
        "visual_object_geometry_ledger": {
            "schema_name": "visual_object_geometry_ledger",
            "status": "passed",
            "visual_object_count": len(visual_objects),
            "objects": visual_objects,
        },
        "source_footer_geometry_ledger": {
            "schema_name": "source_footer_geometry_ledger",
            "status": "passed" if source_footer else "failed",
            "source_footer_object_count": len(source_footer),
            "objects": source_footer,
        },
        "z_order_geometry_ledger": {
            "schema_name": "z_order_geometry_ledger",
            "status": "passed",
            "slide_count": object_ledger.get("slide_count"),
            "slides": _z_order_by_slide(objects),
        },
    }


def infer_geometry_role(name: str, text: str = "") -> str:
    lowered = name.lower()
    has_text = bool(text)
    if lowered.endswith("_bg") or "background" in lowered or lowered.endswith("_background_shape"):
        return "background"
    if "citation_footer" in lowered or "footer" in lowered and text:
        return "source_footer_text"
    if "footer_source_strip" in lowered or "footer_strip" in lowered:
        return "source_footer_strip"
    if has_text and "title" in lowered:
        return "title_text"
    if has_text and "subtitle" in lowered:
        return "subtitle_text"
    if has_text and ("body" in lowered or lowered.endswith("_text") or "_text_" in lowered):
        return "body_text"
    if "frame" in lowered or "panel" in lowered or "card" in lowered or "strip" in lowered or "label" in lowered:
        return "panel"
    if "table" in lowered or "_r" in lowered and "_c" in lowered:
        return "table"
    if "chart" in lowered or "bar_" in lowered:
        return "chart"
    if "icon" in lowered:
        return "icon"
    if "connector" in lowered or "line" in lowered:
        return "connector"
    return "visual_object"


def _object_type(shape: Any) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if "CONNECTOR" in str(shape.shape_type) or "LINE" in str(shape.shape_type):
        return "connector"
    if getattr(shape, "has_text_frame", False) and shape.text_frame is not None and shape.text_frame.text:
        return "text_box"
    return "shape"


def _content_bearing(name: str, text: str) -> bool:
    role = infer_geometry_role(name, text)
    if role in {"panel", "visual_object", "connector"} and not text:
        return False
    return bool(text) or role in {"chart", "table", "icon", "title_text", "body_text", "source_footer_text"}


def _z_order_by_slide(objects: list[dict[str, Any]]) -> list[dict[str, Any]]:
    slides: dict[int, list[dict[str, Any]]] = {}
    for obj in objects:
        slides.setdefault(int(obj["slide_index"]), []).append(
            {
                "object_index": obj.get("object_index", obj.get("z_order", 0)),
                "name": obj["name"],
                "role": obj["role"],
                "z_order": obj.get("z_order", obj.get("object_index", 0)),
                "content_bearing": obj.get("content_bearing", False),
            }
        )
    return [{"slide_index": slide_index, "objects": items} for slide_index, items in sorted(slides.items())]


def _pptx_media_count(path: Path) -> int:
    with zipfile.ZipFile(path) as archive:
        return len([name for name in archive.namelist() if name.startswith("ppt/media/")])
