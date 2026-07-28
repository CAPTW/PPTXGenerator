from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..slide_scene import NativeTable, TableCell, TextRun
from .common import append_warning, stable_scene_shape_name
from .style import SceneStyleContext, resolve_fill_style, resolve_text_style


_EMU_PER_INCH = 914400
_CELL_ALIGNMENTS: dict[str, PP_ALIGN] = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def render_native_table(
    slide: Any,
    native_table: NativeTable,
    *,
    style_context: SceneStyleContext,
    used_shape_names: set[str],
    slide_number: int,
    slide_id: str,
    warnings: list[Any],
) -> str:
    row_count = 1 + len(native_table.rows)
    column_count = len(native_table.headers)
    shape = slide.shapes.add_table(
        row_count,
        column_count,
        Inches(native_table.bounds.x),
        Inches(native_table.bounds.y),
        Inches(native_table.bounds.width),
        Inches(native_table.bounds.height),
    )
    shape.name = stable_scene_shape_name("native_table", native_table.object_id, used_shape_names)
    table = shape.table

    column_widths = _column_widths(native_table, slide_number=slide_number, slide_id=slide_id, warnings=warnings)
    for index, width_emu in enumerate(_dimension_emu_values(column_widths, native_table.bounds.width)):
        table.columns[index].width = width_emu

    for index, height_emu in enumerate(_dimension_emu_values(_row_heights(native_table.bounds.height, row_count), native_table.bounds.height)):
        table.rows[index].height = height_emu

    if native_table.fit.mode not in {"none", "wrap", "fail"}:
        append_warning(
            warnings,
            code="fit_policy_not_enforced",
            severity="warning",
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=native_table.object_id,
            message=f"NativeTable fit policy {native_table.fit.mode!r} is not enforced exactly in the scene compile path.",
            details={"fit_mode": native_table.fit.mode},
        )

    wrap = native_table.fit.mode != "none"
    for column_index, header_cell in enumerate(native_table.headers):
        _render_table_cell(
            table.cell(0, column_index),
            header_cell,
            is_header=True,
            style_context=style_context,
            wrap=wrap,
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=native_table.object_id,
        )
    for row_index, row in enumerate(native_table.rows, start=1):
        for column_index, cell in enumerate(row):
            _render_table_cell(
                table.cell(row_index, column_index),
                cell,
                is_header=False,
                style_context=style_context,
                wrap=wrap,
                slide_number=slide_number,
                slide_id=slide_id,
                object_id=native_table.object_id,
            )
    return shape.name


def _column_widths(
    native_table: NativeTable,
    *,
    slide_number: int,
    slide_id: str,
    warnings: list[Any],
) -> list[float]:
    column_count = len(native_table.headers)
    if native_table.column_widths:
        widths = [round(width, 6) for width in native_table.column_widths]
        total = round(sum(widths), 6)
        if abs(total - native_table.bounds.width) <= 0.000001:
            widths[-1] = round(native_table.bounds.width - sum(widths[:-1]), 6)
            return widths
        scale = native_table.bounds.width / total
        normalized = [round(width * scale, 6) for width in widths]
        normalized[-1] = round(native_table.bounds.width - sum(normalized[:-1]), 6)
        append_warning(
            warnings,
            code="table_column_widths_normalized",
            severity="warning",
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=native_table.object_id,
            message="NativeTable column widths were normalized to match table bounds.",
            details={"declared_total": total, "table_width": round(native_table.bounds.width, 6)},
        )
        return normalized
    base_width = round(native_table.bounds.width / column_count, 6)
    widths = [base_width for _ in range(max(column_count - 1, 0))]
    widths.append(round(native_table.bounds.width - sum(widths), 6))
    return widths


def _row_heights(total_height: float, row_count: int) -> list[float]:
    base_height = round(total_height / row_count, 6)
    heights = [base_height for _ in range(max(row_count - 1, 0))]
    heights.append(round(total_height - sum(heights), 6))
    return heights


def _dimension_emu_values(values_in_inches: list[float], total_in_inches: float) -> list[int]:
    if len(values_in_inches) == 1:
        return [Inches(total_in_inches)]
    emu_values = [Inches(value) for value in values_in_inches[:-1]]
    emu_values.append(int(round(total_in_inches * _EMU_PER_INCH)) - sum(emu_values))
    return emu_values


def _render_table_cell(
    cell: Any,
    cell_data: TableCell,
    *,
    is_header: bool,
    style_context: SceneStyleContext,
    wrap: bool,
    slide_number: int,
    slide_id: str,
    object_id: str,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = resolve_fill_style(
        style_context,
        cell_data.fill,
        "#E2E8F0" if is_header else "#FFFFFF",
        slide_number=slide_number,
        slide_id=slide_id,
        object_id=object_id,
    ).color_rgb

    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = wrap
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraphs = _paragraph_groups(cell_data.runs)
    for paragraph_index, group in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if paragraph_index == 0 else text_frame.add_paragraph()
        paragraph.alignment = _CELL_ALIGNMENTS[cell_data.align]
        for run_index, (source_run, text) in enumerate(group):
            run = paragraph.add_run() if run_index > 0 else paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            _apply_cell_run_style(
                run,
                source_run.model_copy(update={"text": text}),
                style_context=style_context,
                is_header=is_header,
                slide_number=slide_number,
                slide_id=slide_id,
                object_id=object_id,
            )


def _paragraph_groups(runs: list[TextRun]) -> list[list[tuple[TextRun, str]]]:
    paragraphs: list[list[tuple[TextRun, str]]] = [[]]
    for source_run in runs:
        parts = source_run.text.split("\n")
        for part_index, part in enumerate(parts):
            paragraphs[-1].append((source_run, part))
            if part_index < len(parts) - 1:
                paragraphs.append([])
    non_empty = [group for group in paragraphs if group]
    return non_empty or [[(TextRun(text=""), "")]]


def _apply_cell_run_style(
    run: Any,
    source_run: TextRun,
    *,
    style_context: SceneStyleContext,
    is_header: bool,
    slide_number: int,
    slide_id: str,
    object_id: str,
) -> None:
    role = "claim" if is_header else "body"
    text_style = resolve_text_style(
        style_context,
        source_run,
        role,
        slide_number=slide_number,
        slide_id=slide_id,
        object_id=object_id,
    )
    run.text = source_run.text
    run.font.name = text_style.font_name
    run.font.size = Pt(text_style.size_pt)
    run.font.bold = text_style.bold or is_header
    run.font.italic = text_style.italic
    run.font.color.rgb = text_style.color_rgb
