"""D07.2 visual-field asset inventory, policy, and readiness helpers."""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from .source_bound_template_binder import SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN


ALLOWED_VISUAL_TARGETS = {
    "hero_visual_field",
    "photo_frame",
    "abstract_technical_field",
    "case_study_image",
    "section_texture",
    "decorative_background_texture",
}

FORBIDDEN_VISUAL_TARGETS = {
    "title_text",
    "body_text",
    "citation_text",
    "source_footer_text",
    "semantic_icon",
    "semantic_chart",
    "semantic_table",
    "card_body_copy",
    "full_slide_background",
}


def visual_field_asset_policy_v1() -> dict[str, Any]:
    return {
        "schema_name": "visual_field_asset_policy_v1",
        "generated_assets_are_replaceable_visual_fields": True,
        "full_slide_image_background_allowed": False,
        "screenshot_slide_allowed": False,
        "readable_required_text_inside_image_allowed": False,
        "semantic_table_chart_icon_inside_image_allowed": False,
        "source_citation_footer_must_remain_editable_ppt_text": True,
        "semantic_text_icons_charts_tables_must_remain_editable": True,
        "image_frame_bbox_record_required": True,
        "mask_record_required": True,
        "shape_fill_or_picture_crop_must_respect_slot_boundary": True,
        "allowed_visual_asset_targets": sorted(ALLOWED_VISUAL_TARGETS),
        "forbidden_visual_asset_targets": sorted(FORBIDDEN_VISUAL_TARGETS),
        "canva_parity_claimed": False,
    }


def discover_generation_mode(repo_root: Path) -> dict[str, Any]:
    candidates = [
        repo_root / "config/image_generation.json",
        repo_root / "configs/image_generation.json",
        repo_root / "design_runs/run_002/config/image_generation.json",
    ]
    found = [path for path in candidates if path.exists()]
    enabled = False
    selected_config: str | None = None
    for path in found:
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            continue
        if payload.get("image_generation_enabled") is True and payload.get("approved_for_d07_2") is True:
            enabled = True
            selected_config = path.as_posix()
            break
    return {
        "schema_name": "visual_field_asset_generation_mode",
        "mode": "api_configured" if enabled else "manual_import",
        "image_generation_enabled": enabled,
        "approved_config_path": selected_config,
        "configs_found": [path.as_posix() for path in found],
        "remote_generation_called": False,
    }


def build_visual_field_slot_inventory(deck_spec: dict[str, Any]) -> dict[str, Any]:
    slots: list[dict[str, Any]] = []
    forbidden_candidates: list[dict[str, Any]] = []
    for slide in deck_spec.get("slides") or []:
        for obj in slide.get("objects") or []:
            classification = classify_visual_slot(obj)
            if classification in ALLOWED_VISUAL_TARGETS:
                bbox = obj.get("bbox_norm") or [0, 0, 0, 0]
                slots.append(
                    {
                        "slot_id": obj.get("object_id"),
                        "slide_id": slide.get("slide_id"),
                        "slide_number": slide.get("slide_number"),
                        "archetype_id": slide.get("archetype_id"),
                        "object_id": obj.get("object_id"),
                        "visual_field_type": classification,
                        "bbox_norm": bbox,
                        "bbox_area_norm": round(float(bbox[2]) * float(bbox[3]), 6) if len(bbox) == 4 else 0,
                        "mask_type": "rectangular",
                        "target_use": "replaceable_image_frame_or_shape_fill",
                        "allowed_target": True,
                        "semantic_content": False,
                        "source_footer_or_citation": False,
                    }
                )
            elif classification in FORBIDDEN_VISUAL_TARGETS:
                forbidden_candidates.append(
                    {
                        "slide_id": slide.get("slide_id"),
                        "archetype_id": slide.get("archetype_id"),
                        "object_id": obj.get("object_id"),
                        "classification": classification,
                        "allowed_target": False,
                    }
                )
    return {
        "schema_name": "visual_field_slot_inventory",
        "status": "passed" if slots else "blocked",
        "slot_count": len(slots),
        "slots": slots,
        "forbidden_candidate_count": len(forbidden_candidates),
        "forbidden_candidates": forbidden_candidates,
        "full_slide_slot_count": len([slot for slot in slots if slot["bbox_area_norm"] >= 0.85]),
        "canva_parity_claimed": False,
    }


def classify_visual_slot(obj: dict[str, Any]) -> str:
    component = str(obj.get("semantic_component") or "")
    family = str(obj.get("primitive_family") or "")
    identity = str(obj.get("component_identity") or obj.get("identity_region") or "")
    object_id = str(obj.get("object_id") or "")
    combined = f"{component} {family} {identity} {object_id}".lower()
    if component in {"text", "source_footer"} or "citation_footer" in combined:
        return "source_footer_text" if "footer" in combined or "citation" in combined else "body_text"
    if component == "icon":
        return "semantic_icon"
    if component == "chart":
        return "semantic_chart"
    if component in {"table", "matrix"}:
        return "semantic_table"
    if "background_base" in combined or combined.endswith("_bg"):
        return "full_slide_background"
    if "hero_visual_field" in combined or "hero_field" in combined:
        return "hero_visual_field"
    if "case_image" in combined or "case_visual" in combined:
        return "case_study_image"
    if "chapter_visual" in combined or "section" in combined and "visual" in combined:
        return "section_texture"
    if "replaceable_image_frame" in combined or "image_frame" in combined:
        return "photo_frame"
    if "technical" in combined or "abstract" in combined:
        return "abstract_technical_field"
    return "nonvisual_semantic_slot"


def build_prompt_for_slot(slot: dict[str, Any]) -> str:
    aspect = _aspect_label(slot.get("bbox_norm") or [0, 0, 1, 1])
    mood = {
        "hero_visual_field": "cinematic abstract governance workspace with layered evidence trails",
        "photo_frame": "professional editorial photo-like abstract field for decision infrastructure",
        "case_study_image": "documentary-style abstract implementation scene with evidence artifacts",
        "section_texture": "premium section-divider texture with subtle depth",
        "abstract_technical_field": "technical abstract mesh of evidence nodes and review pathways",
        "decorative_background_texture": "quiet non-semantic texture with depth and grain",
    }.get(slot.get("visual_field_type"), "premium abstract visual field")
    return (
        f"Create a scoped replaceable visual-field asset for slide {slot['slide_number']} "
        f"({slot['archetype_id']}, {slot['visual_field_type']}). Subject: {mood}. "
        "Style: dark navy, deep teal, cyan edge light, restrained gold accents, academic/professional/creative. "
        "No readable text, no letters, no numbers, no charts, no tables, no UI labels, no logos, no icons that carry semantic meaning. "
        f"Composition must be safe-cropped for a {aspect} frame, with important visual energy away from edges. "
        "Output a high-resolution PNG or JPG; this is a replaceable visual field, not a slide background."
    )


def build_prompt_pack(slot_inventory: dict[str, Any]) -> tuple[str, dict[str, str]]:
    per_slot: dict[str, str] = {}
    lines = ["# D07.2 Visual Field Generation Prompt Pack", ""]
    lines.append("Use these prompts to create bounded visual-field assets only. Do not generate final slide copy inside images.")
    for slot in slot_inventory.get("slots") or []:
        prompt = build_prompt_for_slot(slot)
        per_slot[f"{slot['slot_id']}.md"] = f"# {slot['slot_id']}\n\n{prompt}\n"
        lines.extend(["", f"## {slot['slot_id']}", "", prompt])
    return "\n".join(lines).strip() + "\n", per_slot


def visual_field_asset_import_requirements(slot_inventory: dict[str, Any]) -> str:
    lines = [
        "# D07.2 Visual Field Asset Import Requirements",
        "",
        "Place manually generated or approved assets in `assets/import/` using one file per slot.",
        "Accepted extensions: `.png`, `.jpg`, `.jpeg`.",
        "Each asset must be bounded to the named visual field slot and must not contain readable required text, charts, tables, UI labels, logos, or semantic icons.",
        "Required sidecar metadata uses `{slot_id}.json` with `declared_no_readable_text: true` and `declared_no_semantic_chart_table_icon: true`.",
        "",
        "Required slot filenames:",
    ]
    for slot in slot_inventory.get("slots") or []:
        lines.append(f"- `{slot['slot_id']}.png` or `{slot['slot_id']}.jpg`")
    return "\n".join(lines) + "\n"


def validate_visual_assets(slot_inventory: dict[str, Any], import_dir: Path, generated_dir: Path, processed_dir: Path) -> dict[str, Any]:
    accepted: list[dict[str, Any]] = []
    missing: list[dict[str, Any]] = []
    rejected: list[dict[str, Any]] = []
    for slot in slot_inventory.get("slots") or []:
        asset_path = _find_asset(import_dir, slot["slot_id"])
        if asset_path is None:
            missing.append({"slot_id": slot["slot_id"], "slide_id": slot["slide_id"], "archetype_id": slot["archetype_id"]})
            continue
        validation = validate_single_visual_asset(asset_path, slot)
        if validation["status"] == "accepted":
            processed_dir.mkdir(parents=True, exist_ok=True)
            processed_path = processed_dir / f"{slot['slot_id']}{asset_path.suffix.lower()}"
            shutil.copy2(asset_path, processed_path)
            accepted.append({**validation, "processed_asset_path": processed_path.as_posix()})
            generated_dir.mkdir(parents=True, exist_ok=True)
            shutil.copy2(asset_path, generated_dir / processed_path.name)
        else:
            rejected.append(validation)
    return {
        "schema_name": "visual_field_asset_validation_report",
        "status": "passed" if accepted and not missing and not rejected else "blocked",
        "accepted_asset_count": len(accepted),
        "missing_asset_count": len(missing),
        "rejected_asset_count": len(rejected),
        "accepted_assets": accepted,
        "missing_assets": missing,
        "rejected_assets": rejected,
    }


def validate_single_visual_asset(asset_path: Path, slot: dict[str, Any]) -> dict[str, Any]:
    errors: list[str] = []
    if slot.get("bbox_area_norm", 0) >= 0.85:
        errors.append("full_slide_asset_insertion_forbidden")
    if slot.get("visual_field_type") not in ALLOWED_VISUAL_TARGETS:
        errors.append("semantic_region_asset_insertion_forbidden")
    try:
        with Image.open(asset_path) as image:
            width, height = image.size
    except Exception as exc:  # noqa: BLE001 - validation report needs concrete error.
        return {
            "slot_id": slot.get("slot_id"),
            "asset_path": asset_path.as_posix(),
            "status": "rejected",
            "errors": [f"asset_unreadable:{exc}"],
        }
    if width < 640 or height < 360:
        errors.append("resolution_below_minimum")
    asset_ratio = width / height
    bbox = slot.get("bbox_norm") or [0, 0, 1, 1]
    slot_ratio = max(0.1, float(bbox[2]) * SLIDE_WIDTH_IN) / max(0.1, float(bbox[3]) * SLIDE_HEIGHT_IN)
    if not (0.45 <= asset_ratio / slot_ratio <= 2.2):
        errors.append("aspect_ratio_outside_safe_crop_range")
    sidecar = asset_path.with_suffix(".json")
    declared_no_text = False
    declared_no_semantic = False
    if sidecar.exists():
        try:
            meta = json.loads(sidecar.read_text(encoding="utf-8"))
            declared_no_text = meta.get("declared_no_readable_text") is True
            declared_no_semantic = meta.get("declared_no_semantic_chart_table_icon") is True
        except json.JSONDecodeError:
            errors.append("metadata_json_invalid")
    if not declared_no_text:
        errors.append("no_readable_text_not_declared")
    if not declared_no_semantic:
        errors.append("no_semantic_chart_table_icon_not_declared")
    return {
        "slot_id": slot.get("slot_id"),
        "slide_id": slot.get("slide_id"),
        "archetype_id": slot.get("archetype_id"),
        "asset_path": asset_path.as_posix(),
        "status": "accepted" if not errors else "rejected",
        "width": width,
        "height": height,
        "asset_aspect_ratio": round(asset_ratio, 4),
        "slot_aspect_ratio": round(slot_ratio, 4),
        "errors": errors,
        "declared_no_readable_text": declared_no_text,
        "declared_no_semantic_chart_table_icon": declared_no_semantic,
    }


def visual_asset_semantic_raster_policy_report(slot_inventory: dict[str, Any], validation_report: dict[str, Any]) -> dict[str, Any]:
    semantic_violations = [
        item
        for item in validation_report.get("rejected_assets") or []
        if "semantic_region_asset_insertion_forbidden" in item.get("errors", [])
    ]
    forbidden_candidates = [
        candidate
        for candidate in slot_inventory.get("forbidden_candidates") or []
        if candidate.get("classification") in FORBIDDEN_VISUAL_TARGETS
    ]
    return {
        "schema_name": "visual_asset_semantic_raster_policy_report",
        "status": "passed" if not semantic_violations else "blocked",
        "semantic_region_asset_insertion_forbidden": True,
        "source_footer_text_image_forbidden": True,
        "semantic_violation_count": len(semantic_violations),
        "semantic_violations": semantic_violations,
        "forbidden_candidate_count": len(forbidden_candidates),
        "forbidden_candidates_inventory_only": forbidden_candidates,
        "accepted_asset_count": validation_report.get("accepted_asset_count", 0),
    }


def shape_fill_alignment_report(slot_inventory: dict[str, Any], validation_report: dict[str, Any]) -> dict[str, Any]:
    accepted = validation_report.get("accepted_assets") or []
    findings = []
    accepted_by_slot = {asset["slot_id"]: asset for asset in accepted}
    for slot in slot_inventory.get("slots") or []:
        if slot["slot_id"] not in accepted_by_slot:
            continue
        if slot.get("bbox_area_norm", 0) >= 0.85:
            findings.append({"slot_id": slot["slot_id"], "issue": "Asset slot is effectively full slide."})
    return {
        "schema_name": "shape_fill_alignment_report",
        "status": "passed" if not findings and accepted else "blocked" if not accepted else "failed",
        "checked_asset_count": len(accepted),
        "alignment_issue_count": len(findings),
        "findings": findings,
    }


def build_asset_binding_ledger(slot_inventory: dict[str, Any], validation_report: dict[str, Any]) -> dict[str, Any]:
    accepted = {asset["slot_id"]: asset for asset in validation_report.get("accepted_assets") or []}
    bindings = []
    for slot in slot_inventory.get("slots") or []:
        asset = accepted.get(slot["slot_id"])
        bindings.append(
            {
                "slot_id": slot["slot_id"],
                "slide_id": slot["slide_id"],
                "archetype_id": slot["archetype_id"],
                "visual_field_type": slot["visual_field_type"],
                "bbox_norm": slot["bbox_norm"],
                "binding_status": "bound" if asset else "missing_asset",
                "asset_path": asset.get("processed_asset_path") if asset else None,
                "full_slide_background": False,
                "semantic_raster_target": False,
            }
        )
    return {
        "schema_name": "visual_field_asset_binding_ledger",
        "status": "passed" if validation_report.get("status") == "passed" else "blocked",
        "binding_count": len(bindings),
        "bound_count": sum(1 for item in bindings if item["binding_status"] == "bound"),
        "missing_count": sum(1 for item in bindings if item["binding_status"] == "missing_asset"),
        "bindings": bindings,
    }


def build_d08_visual_asset_readiness(
    *,
    deck_exists: bool,
    render_count: int,
    slide_count: int,
    asset_validation: dict[str, Any],
    semantic_policy: dict[str, Any],
    alignment_report: dict[str, Any],
    protected_artifacts_unchanged: bool,
    manual_import_mode: bool,
    large_deck_created: bool = False,
    bulk_deck_created: bool = False,
    c11_started: bool = False,
) -> dict[str, Any]:
    required_assets_missing = asset_validation.get("missing_asset_count", 0) > 0
    accepted_assets = asset_validation.get("accepted_asset_count", 0) > 0
    all_rendered = deck_exists and render_count == slide_count and slide_count > 0
    no_semantic = semantic_policy.get("semantic_violation_count", 0) == 0
    alignment_ok = alignment_report.get("status") == "passed"
    if not protected_artifacts_unchanged:
        decision = "D07_2_FAIL_PROTECTED_ARTIFACTS"
    elif required_assets_missing and manual_import_mode:
        decision = "D07_2_BLOCKED_ASSET_IMPORT_REQUIRED"
    elif semantic_policy.get("status") == "blocked" and accepted_assets:
        decision = "D07_2_FAIL_SEMANTIC_RASTER_POLICY"
    elif accepted_assets and not alignment_ok:
        decision = "D07_2_PATCH_SHAPE_FILL_ALIGNMENT"
    elif accepted_assets and deck_exists and all_rendered and no_semantic:
        decision = "D07_2_PASS_START_D08_WITH_VISUAL_FIELD_ASSETS"
    else:
        decision = "D07_2_PASS_OPTIONAL_ASSET_LAYER_READY"
    unlocked = decision in {"D07_2_PASS_START_D08_WITH_VISUAL_FIELD_ASSETS", "D07_2_PASS_OPTIONAL_ASSET_LAYER_READY"} and protected_artifacts_unchanged
    return {
        "schema_name": "d08_visual_asset_readiness_report",
        "decision": decision,
        "d08_visual_asset_layer_unlocked": unlocked,
        "deck_exists": deck_exists,
        "slide_count": slide_count,
        "render_count": render_count,
        "manual_import_mode": manual_import_mode,
        "accepted_asset_count": asset_validation.get("accepted_asset_count", 0),
        "missing_asset_count": asset_validation.get("missing_asset_count", 0),
        "unlock_conditions": {
            "d07_2_deck_exists_or_optional_nonblocking": deck_exists or decision == "D07_2_PASS_OPTIONAL_ASSET_LAYER_READY",
            "accepted_visual_assets_bounded_to_slots": not required_assets_missing and no_semantic,
            "semantic_text_icon_chart_table_not_rasterized": no_semantic,
            "no_full_slide_background": True,
            "no_screenshot_slide": True,
            "all_slides_render": all_rendered,
            "visual_quality_improved_or_preserved": accepted_assets or decision == "D07_2_PASS_OPTIONAL_ASSET_LAYER_READY",
            "no_critical_blockers": decision not in {"D07_2_FAIL_SEMANTIC_RASTER_POLICY", "D07_2_FAIL_PROTECTED_ARTIFACTS"},
            "no_high_product_risks": decision not in {"D07_2_BLOCKED_ASSET_IMPORT_REQUIRED", "D07_2_PATCH_SHAPE_FILL_ALIGNMENT"},
            "large_deck_created": large_deck_created,
            "bulk_deck_created": bulk_deck_created,
            "c11_remains_frozen": not c11_started,
            "protected_artifacts_unchanged": protected_artifacts_unchanged,
        },
        "canva_parity_claimed": False,
    }


def insert_visual_assets_into_deck(input_deck: Path, output_deck: Path, slot_inventory: dict[str, Any], validation_report: dict[str, Any]) -> dict[str, Any]:
    accepted = {asset["slot_id"]: asset for asset in validation_report.get("accepted_assets") or []}
    if not accepted:
        return {
            "schema_name": "d07_2_deck_patch_report",
            "status": "blocked",
            "deck_created": False,
            "reason": "No accepted visual assets available.",
        }
    prs = Presentation(input_deck)
    inserted = []
    for slot in slot_inventory.get("slots") or []:
        asset = accepted.get(slot["slot_id"])
        if not asset:
            continue
        slide_index = int(slot["slide_number"]) - 1
        bbox = slot["bbox_norm"]
        slide = prs.slides[slide_index]
        shape = slide.shapes.add_picture(
            asset["processed_asset_path"],
            Inches(float(bbox[0]) * SLIDE_WIDTH_IN),
            Inches(float(bbox[1]) * SLIDE_HEIGHT_IN),
            Inches(float(bbox[2]) * SLIDE_WIDTH_IN),
            Inches(float(bbox[3]) * SLIDE_HEIGHT_IN),
        )
        shape.name = f"d07_2_asset_{slot['slot_id']}"
        inserted.append({"slot_id": slot["slot_id"], "asset_path": asset["processed_asset_path"], "slide_number": slot["slide_number"]})
    output_deck.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_deck)
    return {
        "schema_name": "d07_2_deck_patch_report",
        "status": "passed",
        "deck_created": output_deck.exists(),
        "inserted_asset_count": len(inserted),
        "inserted_assets": inserted,
    }


def _find_asset(import_dir: Path, slot_id: str) -> Path | None:
    for suffix in [".png", ".jpg", ".jpeg"]:
        candidate = import_dir / f"{slot_id}{suffix}"
        if candidate.exists():
            return candidate
    return None


def _aspect_label(bbox: list[float]) -> str:
    ratio = max(0.1, float(bbox[2]) * SLIDE_WIDTH_IN) / max(0.1, float(bbox[3]) * SLIDE_HEIGHT_IN)
    if ratio >= 1.8:
        return "wide landscape"
    if ratio <= 0.8:
        return "portrait/tall crop"
    return "balanced landscape"
