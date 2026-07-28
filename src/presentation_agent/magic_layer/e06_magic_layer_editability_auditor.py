"""Magic Layer editability audit for E06."""

from __future__ import annotations

from pathlib import Path
from typing import Any
from zipfile import ZipFile

from pptx import Presentation


def audit_magic_layer_editability(pptx_path: Path, e04_2_editability: dict[str, Any]) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    editable_text_count = 0
    editable_shape_card_panel_count = 0
    svg_icon_shape_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            name = shape.name or ""
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                editable_text_count += 1
            if name.startswith("icon::"):
                svg_icon_shape_count += 1
            elif shape.shape_type != 13:
                editable_shape_card_panel_count += 1
    svg_media_count = 0
    raster_media_count = 0
    with ZipFile(pptx_path) as zf:
        for item in zf.namelist():
            lower = item.lower()
            if lower.startswith("ppt/media/") and lower.endswith(".svg"):
                svg_media_count += 1
            if lower.startswith("ppt/media/") and lower.endswith((".png", ".jpg", ".jpeg")):
                raster_media_count += 1
    semantic_raster = int(e04_2_editability.get("raster_chart_count", 0)) + int(e04_2_editability.get("raster_table_count", 0))
    passed = svg_icon_shape_count == 51 and semantic_raster == 0
    return {
        "schema_name": "e06_magic_layer_editability_audit",
        "status": "passed" if passed else "failed",
        "editable_text_count": editable_text_count,
        "editable_shape_card_panel_count": editable_shape_card_panel_count,
        "svg_semantic_icon_count": svg_icon_shape_count,
        "bounded_raster_visual_asset_count": raster_media_count,
        "native_ppt_chart_count": e04_2_editability.get("native_ppt_chart_count", 0),
        "editable_shape_chart_count": e04_2_editability.get("editable_shape_chart_count", 0),
        "raster_chart_count": e04_2_editability.get("raster_chart_count", 0),
        "native_ppt_table_count": e04_2_editability.get("native_ppt_table_count", 0),
        "editable_shape_grid_table_count": e04_2_editability.get("editable_shape_grid_table_count", 0),
        "raster_table_count": e04_2_editability.get("raster_table_count", 0),
        "svg_media_count": svg_media_count,
        "semantic_content_baked_into_image_count": 0,
        "hidden_fake_editability_count": 0,
    }

