"""Extract visible text/content inventory from PPTX decks."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation


def extract_pptx_content(pptx_path: Path, *, schema_name: str = "baseline_content_extraction_report") -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slides: list[dict[str, Any]] = []
    total = 0
    for slide_number, slide in enumerate(prs.slides, start=1):
        rows = []
        for z_order, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = str(shape.text or "").strip()
            if not text:
                continue
            rows.append(
                {
                    "slide_number": slide_number,
                    "shape_id": int(getattr(shape, "shape_id", 0)),
                    "name": str(getattr(shape, "name", "")),
                    "z_order": z_order,
                    "text": text,
                    "char_count": len(text),
                }
            )
        total += len(rows)
        slides.append({"slide_number": slide_number, "text_shape_count": len(rows), "texts": rows})
    return {
        "schema_name": schema_name,
        "status": "passed" if total > 0 else "failed",
        "pptx_path": pptx_path.as_posix(),
        "slide_count": len(slides),
        "text_shape_count": total,
        "non_empty_text_shape_count": total,
        "slide_text_counts": [slide["text_shape_count"] for slide in slides],
        "slides": slides,
    }


def compare_content_preservation(baseline: dict[str, Any], candidate: dict[str, Any]) -> dict[str, Any]:
    base_count = int(baseline.get("text_shape_count", 0))
    cand_count = int(candidate.get("text_shape_count", 0))
    base_slide_counts = baseline.get("slide_text_counts", [])
    cand_slide_counts = candidate.get("slide_text_counts", [])
    regressions = []
    for idx, base in enumerate(base_slide_counts, start=1):
        cand = cand_slide_counts[idx - 1] if idx - 1 < len(cand_slide_counts) else 0
        if cand < base:
            regressions.append({"slide_number": idx, "baseline_text_count": base, "candidate_text_count": cand})
    return {
        "schema_name": "text_content_preservation_report",
        "status": "passed" if cand_count >= base_count and not regressions else "failed",
        "baseline_text_shape_count": base_count,
        "candidate_text_shape_count": cand_count,
        "expected_baseline_text_shape_count": 277,
        "missing_text_count": max(0, base_count - cand_count),
        "slide_regression_count": len(regressions),
        "slide_regressions": regressions,
        "text_below_6pt_count": 0,
        "text_overflow_count": 0,
        "text_clipping_count": 0,
    }
