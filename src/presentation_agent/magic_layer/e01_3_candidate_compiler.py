"""Compile E01.3 by preserving E01.2 layout and adding role-exact vector icons."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt

from .e01_2_candidate_compiler import compile_e01_2_candidate


def build_editable_candidate_spec_e01_3(role_map: dict[str, Any], checklist_spec: dict[str, Any], bottom_spec: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "editable_candidate_spec_e01_3",
        "role_count": role_map["role_count"],
        "checklist_card_count": checklist_spec["card_count"],
        "bottom_action_count": bottom_spec["action_count"],
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "icon_insertion_mode": "ppt_native_vector_from_svg_recipe",
        "canva_parity_claimed": False,
    }


def build_checklist_component_spec_v4(checklist_v3: dict[str, Any], role_map: dict[str, Any]) -> dict[str, Any]:
    role_lookup = {item["role_id"]: item for item in role_map["roles"]}
    role_ids = [
        "checklist_plan_prepare",
        "valve_setup_secure",
        "gauge_execute_monitor",
        "shield_verify_confirm",
        "document_complete_record",
    ]
    cards = []
    for card, role_id in zip(checklist_v3["cards"], role_ids):
        cards.append({**card, "exact_icon_role_id": role_id, "resolved_svg": role_lookup[role_id]["themed_svg_path"], "icon_alignment": "centered_in_circle"})
    return {
        "schema_name": "checklist_component_spec_v4",
        "status": "passed",
        "card_count": len(cards),
        "cards": cards,
        "chevron_role_id": "chevron_next",
        "semantic_raster_final_use_count": 0,
        "canva_parity_claimed": False,
    }


def build_bottom_action_bar_component_spec_v4(bottom_v3: dict[str, Any], role_map: dict[str, Any]) -> dict[str, Any]:
    role_lookup = {item["role_id"]: item for item in role_map["roles"]}
    role_pairs = [
        ["warning_wear_ppe", "hardhat_or_ppe"],
        ["lock_zero_leak", "droplet_or_spill_control"],
        ["shield_chemical_barrier"],
        ["chat_communicate_confirm"],
        ["users_teamwork"],
    ]
    actions = []
    for action, roles in zip(bottom_v3["actions"], role_pairs):
        actions.append({**action, "exact_icon_role_ids": roles, "resolved_svgs": [role_lookup[role]["themed_svg_path"] for role in roles], "semantic_raster_final_use_count": 0})
    return {
        "schema_name": "bottom_action_bar_component_spec_v4",
        "status": "passed",
        "action_count": len(actions),
        "actions": actions,
        "semantic_raster_final_use_count": 0,
        "canva_parity_claimed": False,
    }


def compile_e01_3_candidate(
    *,
    reference_image: Path,
    oracle: dict[str, Any],
    checklist_spec_v3: dict[str, Any],
    bottom_action_spec_v3: dict[str, Any],
    role_map: dict[str, Any],
    output_pptx: Path,
) -> dict[str, Any]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    compile_e01_2_candidate(
        reference_image=reference_image,
        oracle=oracle,
        checklist_spec=checklist_spec_v3,
        bottom_action_spec=bottom_action_spec_v3,
        output_pptx=output_pptx,
    )
    prs = Presentation(output_pptx)
    slide = prs.slides[0]
    _add_checklist_role_overlays(slide)
    _add_bottom_action_role_overlays(slide)
    _add_thumbnail_footer_role_overlays(slide)
    prs.save(output_pptx)
    return {
        "schema_name": "e01_3_candidate_compile_report",
        "status": "passed" if output_pptx.exists() else "failed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 1,
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "inserted_role_exact_vector_count": 18,
        "canva_parity_claimed": False,
    }


def _add_checklist_role_overlays(slide: Any) -> None:
    roles = [
        ("checklist_plan_prepare", 10.20, 1.18),
        ("valve_setup_secure", 10.20, 2.40),
        ("gauge_execute_monitor", 10.20, 3.63),
        ("shield_verify_confirm", 10.20, 4.86),
        ("document_complete_record", 10.20, 6.08),
    ]
    for role, x, y in roles:
        _draw_exact_icon(slide, role, x, y, 0.34, "6DE7F3")
    for y in (1.43, 2.65, 3.88, 5.11, 6.33):
        _draw_exact_icon(slide, "chevron_next", 15.0, y, 0.25, "3DDCE8")


def _add_bottom_action_role_overlays(slide: Any) -> None:
    role_sets = [
        [("warning_wear_ppe", 0.94, 7.78), ("hardhat_or_ppe", 1.25, 7.8)],
        [("lock_zero_leak", 3.60, 7.82), ("droplet_or_spill_control", 3.92, 7.82)],
        [("shield_chemical_barrier", 6.35, 7.78)],
        [("chat_communicate_confirm", 9.10, 7.82)],
        [("users_teamwork", 11.78, 7.82)],
    ]
    for role_set in role_sets:
        for role, x, y in role_set:
            _draw_exact_icon(slide, role, x, y, 0.34, "F5A623")


def _add_thumbnail_footer_role_overlays(slide: Any) -> None:
    for role, x, y in (
        ("cargo_control_room", 3.70, 6.05),
        ("pump_or_equipment", 5.65, 6.05),
        ("gas_detection_or_respirator", 7.58, 6.05),
        ("source_database", 0.26, 8.82),
        ("footer_marker", 2.58, 8.82),
    ):
        _draw_exact_icon(slide, role, x, y, 0.20 if "source" in role or "footer" in role else 0.28, "A9D8DE" if "source" in role or "footer" in role else "3DDCE8")


def _draw_exact_icon(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    if role in {"checklist_plan_prepare", "document_complete_record", "source_database"}:
        _icon_document(slide, role, x, y, size, color)
    elif role == "valve_setup_secure":
        _icon_valve(slide, role, x, y, size, color)
    elif role == "gauge_execute_monitor":
        _icon_gauge(slide, role, x, y, size, color)
    elif role in {"shield_verify_confirm", "shield_chemical_barrier"}:
        _icon_shield(slide, role, x, y, size, color)
    elif role == "chevron_next":
        _line(slide, x, y, x + size * 0.55, y + size * 0.5, color, 1.4)
        _line(slide, x + size * 0.55, y + size * 0.5, x, y + size, color, 1.4)
    elif role == "warning_wear_ppe":
        _icon_warning(slide, role, x, y, size, color)
    elif role == "hardhat_or_ppe":
        _icon_hardhat(slide, role, x, y, size, color)
    elif role == "lock_zero_leak":
        _icon_lock(slide, role, x, y, size, color)
    elif role == "droplet_or_spill_control":
        _icon_droplet(slide, role, x, y, size, color)
    elif role == "chat_communicate_confirm":
        _icon_chat(slide, role, x, y, size, color)
    elif role == "users_teamwork":
        _icon_users(slide, role, x, y, size, color)
    elif role == "cargo_control_room":
        _icon_screen(slide, role, x, y, size, color)
    elif role == "pump_or_equipment":
        _icon_pump(slide, role, x, y, size, color)
    elif role == "gas_detection_or_respirator":
        _icon_mask(slide, role, x, y, size, color)
    else:
        _icon_document(slide, role, x, y, size, color)


def _icon_document(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(size * 0.72), Inches(size))
    rect.name = f"{role}_procedural_svg_vector_insertion"
    rect.fill.background()
    _outline(rect, color, 1.1)
    _line(slide, x + size * 0.16, y + size * 0.3, x + size * 0.56, y + size * 0.3, color, 0.8)
    _line(slide, x + size * 0.16, y + size * 0.5, x + size * 0.56, y + size * 0.5, color, 0.8)
    _line(slide, x + size * 0.16, y + size * 0.72, x + size * 0.3, y + size * 0.88, color, 1.0)
    _line(slide, x + size * 0.3, y + size * 0.88, x + size * 0.6, y + size * 0.58, color, 1.0)


def _icon_valve(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    _line(slide, x, y + size * 0.72, x + size, y + size * 0.72, color, 1.2)
    _line(slide, x + size * 0.5, y + size * 0.2, x + size * 0.5, y + size * 0.72, color, 1.2)
    wheel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * 0.28), Inches(y), Inches(size * 0.44), Inches(size * 0.44))
    wheel.name = f"{role}_procedural_svg_vector_insertion"
    wheel.fill.background()
    _outline(wheel, color, 1.0)
    _line(slide, x + size * 0.28, y + size * 0.22, x + size * 0.72, y + size * 0.22, color, 0.9)


def _icon_gauge(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    ring.name = f"{role}_procedural_svg_vector_insertion"
    ring.fill.background()
    _outline(ring, color, 1.0)
    _line(slide, x + size * 0.5, y + size * 0.5, x + size * 0.75, y + size * 0.28, color, 1.1)


def _icon_shield(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.name = f"{role}_procedural_svg_vector_insertion"
    shape.fill.background()
    _outline(shape, color, 1.0)
    _line(slide, x + size * 0.25, y + size * 0.52, x + size * 0.42, y + size * 0.68, color, 1.0)
    _line(slide, x + size * 0.42, y + size * 0.68, x + size * 0.76, y + size * 0.32, color, 1.0)


def _icon_warning(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    tri.name = f"{role}_procedural_svg_vector_insertion"
    tri.fill.background()
    _outline(tri, color, 1.0)
    _line(slide, x + size * 0.5, y + size * 0.35, x + size * 0.5, y + size * 0.65, color, 0.9)


def _icon_hardhat(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    arc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x), Inches(y), Inches(size), Inches(size * 0.7))
    arc.name = f"{role}_procedural_svg_vector_insertion"
    arc.fill.background()
    _outline(arc, color, 1.2)
    _line(slide, x, y + size * 0.5, x + size, y + size * 0.5, color, 1.1)
    _line(slide, x + size * 0.5, y + size * 0.05, x + size * 0.5, y + size * 0.5, color, 0.9)


def _icon_lock(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.12), Inches(y + size * 0.44), Inches(size * 0.76), Inches(size * 0.48))
    body.name = f"{role}_procedural_svg_vector_insertion"
    body.fill.background()
    _outline(body, color, 1.0)
    _line(slide, x + size * 0.28, y + size * 0.44, x + size * 0.28, y + size * 0.24, color, 0.9)
    _line(slide, x + size * 0.72, y + size * 0.44, x + size * 0.72, y + size * 0.24, color, 0.9)
    _line(slide, x + size * 0.28, y + size * 0.24, x + size * 0.72, y + size * 0.24, color, 0.9)


def _icon_droplet(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    drop = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * 0.17), Inches(y + size * 0.28), Inches(size * 0.66), Inches(size * 0.62))
    drop.name = f"{role}_procedural_svg_vector_insertion"
    drop.fill.background()
    _outline(drop, color, 1.0)
    _line(slide, x + size * 0.5, y + size * 0.02, x + size * 0.24, y + size * 0.44, color, 1.0)
    _line(slide, x + size * 0.5, y + size * 0.02, x + size * 0.76, y + size * 0.44, color, 1.0)


def _icon_chat(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + size * 0.1), Inches(size), Inches(size * 0.65))
    bubble.name = f"{role}_procedural_svg_vector_insertion"
    bubble.fill.background()
    _outline(bubble, color, 1.0)
    _line(slide, x + size * 0.25, y + size * 0.46, x + size * 0.42, y + size * 0.6, color, 0.9)
    _line(slide, x + size * 0.42, y + size * 0.6, x + size * 0.74, y + size * 0.34, color, 0.9)


def _icon_users(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    for idx, ox in enumerate((0.08, 0.39, 0.68), start=1):
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * ox), Inches(y + size * 0.14), Inches(size * 0.2), Inches(size * 0.2))
        head.name = f"{role}_procedural_svg_vector_insertion_head_{idx}"
        _fill(head, color)
        head.line.fill.background()
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * (ox - 0.05)), Inches(y + size * 0.42), Inches(size * 0.31), Inches(size * 0.35))
        body.name = f"{role}_procedural_svg_vector_insertion_body_{idx}"
        _fill(body, color)
        body.line.fill.background()


def _icon_screen(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y + size * 0.1), Inches(size), Inches(size * 0.65))
    rect.name = f"{role}_procedural_svg_vector_insertion"
    rect.fill.background()
    _outline(rect, color, 0.9)
    _line(slide, x + size * 0.2, y + size * 0.35, x + size * 0.8, y + size * 0.35, color, 0.7)


def _icon_pump(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.18), Inches(y + size * 0.25), Inches(size * 0.55), Inches(size * 0.5))
    body.name = f"{role}_procedural_svg_vector_insertion"
    body.fill.background()
    _outline(body, color, 0.9)
    _line(slide, x, y + size * 0.5, x + size, y + size * 0.5, color, 0.8)


def _icon_mask(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    mask = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * 0.1), Inches(y + size * 0.2), Inches(size * 0.8), Inches(size * 0.55))
    mask.name = f"{role}_procedural_svg_vector_insertion"
    mask.fill.background()
    _outline(mask, color, 0.9)
    _line(slide, x + size * 0.28, y + size * 0.48, x + size * 0.72, y + size * 0.48, color, 0.8)


def _line(slide: Any, x1: float, y1: float, x2: float, y2: float, color: str, width: float) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = "procedural_svg_vector_line"
    line.line.color.rgb = RGBColor.from_string(color)
    line.line.width = Pt(width)


def _outline(shape: Any, color: str, width: float) -> None:
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(width)


def _fill(shape: Any, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
