"""Deterministic deck QA against compiled output, build state, and continuity rules."""

from __future__ import annotations

import json
import os
import re
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..compat.legacy_non_pptx import WorkflowGate
from ..pptx_compiler import BuildManifest, SlideBuildLinkage, _forbidden_visible_text_reason, load_pptx_compile_file
from ..qa_compile_policy import summarize_qa_verdict
from ..compat.legacy_non_pptx import (
    AssetKind,
    AssetManifest,
    AssetStatus,
    Blueprint,
    ContentTier,
    ContractModel,
    DeckConstitution,
    DeckMode,
    DesignSystem,
    FindingStatus,
    FrameFit,
    LectureFamily,
    LayoutLibrary,
    QAFindingGovernanceDisposition,
    QAGovernance,
    QARecommendationType,
    QAReport,
    QAFinding,
    QALayer,
    QASlideResult,
    QASummary,
    QASeverity,
    QAStatus,
    SlideArchetype,
    SlideEvidenceClass,
    SlideLedger,
    SlideIntent,
    SlideRole,
    SlideRange,
    StateCapsule,
    StageStatus,
    VisualType,
    VizManifest,
    load_state_file,
    save_state_file,
)
from .qa_governance import refresh_qa_governance
from .state_schemas import normalize_continuity_guidance_and_mirror

TERM_RE = re.compile(r'"([^"]+)"')
PUNCTUATION_ENDINGS = (".", "!", "?", ":", ";")
SUSPICIOUS_OPTIMIZATION_FRAMING = (
    "problem formulation",
    "optimality conditions",
    "method selection",
    "structurally appropriate",
    "graduate optimization",
    "convexity",
    "kkt",
    "line search",
)
CHROME_PHRASES = (
    "roadmap",
    "why now",
    "why it matters",
    "comparison takeaway",
    "phase 1",
    "phase 2",
    "phase 3",
)
MAPPING_ARCHETYPES = {
    SlideArchetype.TWO_COLUMN_MAPPING_TABLE,
    SlideArchetype.CORRESPONDENCE_MATRIX,
}
PROCESS_ARCHETYPES = {
    SlideArchetype.PROCESS_FLOW,
    SlideArchetype.STEP_BY_STEP_MECHANISM,
}
WORKED_EXAMPLE_ARCHETYPES = {
    SlideArchetype.WORKED_EXAMPLE_STATE_TABLE,
}
APPENDIX_CARD_ARCHETYPES = {
    SlideArchetype.APPENDIX_THEMED_EVIDENCE_CLUSTER,
    SlideArchetype.APPENDIX_ANNOTATED_EXCERPT_CLUSTER,
    SlideArchetype.APPENDIX_EVIDENCE_CLUSTER,
}
APPENDIX_TABLE_ARCHETYPES = {
    SlideArchetype.APPENDIX_SOURCE_LOCATION_MATRIX,
    SlideArchetype.APPENDIX_COMPARISON_EVIDENCE_CLUSTER,
    SlideArchetype.APPENDIX_SOURCE_MAP,
}
APPENDIX_ARCHETYPES = APPENDIX_CARD_ARCHETYPES | APPENDIX_TABLE_ARCHETYPES
TEXT_CARD_ARCHETYPES = {
    "anchor-concept-card",
    "comparison-callout-cluster",
    "generic-text-card",
}


@dataclass(slots=True)
class ShapeBox:
    kind: str
    left: float
    top: float
    width: float
    height: float
    text: str
    is_title: bool = False
    row_count: int = 0
    column_count: int = 0


@dataclass(slots=True)
class SlideSnapshot:
    pptx_index: int
    title_text: str
    all_text: str
    text_char_count: int
    paragraph_count: int
    picture_count: int
    table_count: int
    text_blocks: list[str]
    table_blocks: list[str]
    text_shape_count: int
    placeholder_count: int
    auto_shape_count: int
    group_count: int
    connector_count: int
    shape_type_counts: dict[str, int]
    shape_boxes: list[ShapeBox]


class DeckQAOutputs(ContractModel):
    qa_report: QAReport
    qa_governance: QAGovernance | None = None
    slide_ledger: SlideLedger
    slide_build_linkage: SlideBuildLinkage
    state_capsule: StateCapsule | None = None
    compiled_deck_text: dict[str, Any] | None = None
    compiled_deck_shape_census: dict[str, Any] | None = None
    compiled_deck_authoring_audit: dict[str, Any] | None = None
    compiled_deck_thumbnail_strip_path: str | None = None


def _norm(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _normalized_block(text: str) -> str:
    normalized = _norm(text).lower()
    normalized = re.sub(r"^[\-\u2022]\s*", "", normalized)
    return normalized


def _dedupe(items: list[str]) -> list[str]:
    seen: set[str] = set()
    ordered: list[str] = []
    for item in items:
        if item and item not in seen:
            seen.add(item)
            ordered.append(item)
    return ordered


def _resolve_continuity_policy_inputs(state_capsule: StateCapsule | None) -> tuple[list[Any] | None, list[str] | None]:
    if state_capsule is None:
        return None, None
    continuity_alerts = list(state_capsule.continuity_alerts)
    if continuity_alerts:
        return continuity_alerts, None
    continuity_guidance, _continuity_warning_mirror = normalize_continuity_guidance_and_mirror(
        continuity_guidance=getattr(state_capsule, "continuity_guidance", []),
        continuity_warnings=getattr(state_capsule, "continuity_warnings", []),
    )
    return None, continuity_guidance or None


def _title_stem(text: str) -> str:
    stem = _norm(text).lower()
    stem = stem.replace("evidence cluster:", "").strip()
    stem = re.sub(r"\s+as a bridge$", "", stem)
    stem = re.sub(r"\s+in the cycle$", "", stem)
    stem = re.sub(r"\s+tradeoffs$", "", stem)
    stem = re.sub(r"\s+pitfalls$", "", stem)
    return stem


def _formula_prefix(text: str, *, words: int = 3) -> str:
    tokens = [token for token in re.split(r"[^a-z0-9]+", _norm(text).lower()) if token]
    return " ".join(tokens[:words])


def _duplicate_text_flags(slide, entry) -> list[str]:
    title_norm = _norm(slide.title).lower()
    title_stem = _title_stem(slide.title)
    message_norm = _norm(entry.main_message).lower()
    takeaway_norm = _norm(entry.one_line_takeaway).lower()
    flags = list(slide.duplicate_text_flags)
    if title_norm and (title_norm == message_norm or title_norm == takeaway_norm):
        flags.append("title-repeated-as-main-claim")
    for item in slide.core_content:
        item_norm = _norm(item).lower()
        if not item_norm:
            continue
        if item_norm in {message_norm, takeaway_norm}:
            flags.append("body-repeats-main-claim")
        if title_stem and title_stem in item_norm and len(title_stem) >= 18:
            flags.append("body-repeats-title")
    return _dedupe(flags)


def _expected_archetypes(slide) -> set[SlideArchetype]:
    title = _norm(slide.title).lower()
    expected: set[SlideArchetype] = set()
    if slide.slide_archetype is not None:
        return {slide.slide_archetype}
    if slide.deck_mode == DeckMode.APPENDIX:
        if title.startswith("source map for") or title.startswith("evidence summary for") or title.startswith("evidence cluster:"):
            expected.add(SlideArchetype.APPENDIX_EVIDENCE_CLUSTER)
        return expected
    if slide.slide_intent == SlideIntent.MAPPING_BRIDGE or any(token in title for token in {" map", "encoding", "fitness is external"}):
        expected.update(MAPPING_ARCHETYPES)
    if slide.slide_intent == SlideIntent.MECHANISM_WALKTHROUGH or "cycle" in title or title.startswith("how one"):
        expected.update(PROCESS_ARCHETYPES)
    if slide.slide_intent == SlideIntent.WORKED_EXAMPLE or "worked example" in title or "toy population" in title:
        expected.update(WORKED_EXAMPLE_ARCHETYPES)
    if slide.slide_intent == SlideIntent.COMPARISON_TRADEOFF:
        expected.update({SlideArchetype.COMPARISON_MATRIX, SlideArchetype.LIMITATION_PITFALL_CALLOUT})
    if slide.slide_intent == SlideIntent.MISCONCEPTION_PITFALL or "limitation" in title or "metaphor" in title:
        expected.add(SlideArchetype.LIMITATION_PITFALL_CALLOUT)
    if slide.slide_intent == SlideIntent.APPLICATION_VIGNETTE or "use case" in title or "useful when" in title:
        expected.add(SlideArchetype.APPLICATION_VIGNETTE)
    if slide.slide_intent == SlideIntent.SUMMARY_INTEGRATION or title.startswith("synthesis:"):
        expected.add(SlideArchetype.SYNTHESIS_INTEGRATION)
    return expected


def _expected_visual_is_missing(slide, link, snapshot: SlideSnapshot | None) -> bool:
    if snapshot is None:
        return False
    archetype = slide.slide_archetype
    if archetype in MAPPING_ARCHETYPES or archetype == SlideArchetype.WORKED_EXAMPLE_STATE_TABLE or slide.slide_intent in {SlideIntent.MAPPING_BRIDGE, SlideIntent.WORKED_EXAMPLE}:
        return snapshot.table_count < 1
    if archetype in PROCESS_ARCHETYPES or slide.slide_intent == SlideIntent.MECHANISM_WALKTHROUGH:
        return link.layout_family != "process-flow"
    return False


def _deck_range(slide_ledger: SlideLedger) -> SlideRange:
    end = slide_ledger.entries[-1].slide_number if slide_ledger.entries else 1
    return SlideRange(start=1, end=end)


def _is_lecture_blueprint(blueprint: Blueprint) -> bool:
    return blueprint.chosen_workflow == "graduate-lecture-clustered" or blueprint.main_story_slide_budget is not None


def _allows_native_visual_without_external_asset(layout_family: str | None, visual_type: VisualType) -> bool:
    if layout_family in {"process-flow", "concept-explainer", "definition-theorem"}:
        return visual_type in {
            VisualType.PROCESS,
            VisualType.TIMELINE,
            VisualType.DECISION_PATH,
            VisualType.FRAMEWORK,
            VisualType.HIERARCHY,
            VisualType.INFOGRAPHIC,
            VisualType.METRIC_SUMMARY,
        }
    if layout_family in {"summary", "comparison", "appendix-reference"}:
        return visual_type in {VisualType.PROCESS, VisualType.COMPARISON, VisualType.TEXT, VisualType.QUOTE}
    return False


def _path(text: str, root: Path | None) -> Path:
    path = Path(text)
    if path.is_absolute():
        return path
    candidates = []
    if root is not None:
        candidates.append((root / path).resolve())
    candidates.append(path.resolve())
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[0] if candidates else path


def _numbering_label(slide_number: int, appendix: bool) -> str:
    return f"Appendix {slide_number}" if appendix else str(slide_number)


def _normalized_return_content_tier(entry, slide) -> ContentTier | None:
    if entry.deck_mode == DeckMode.APPENDIX:
        return ContentTier.APPENDIX_ONLY
    if entry.content_tier == ContentTier.APPENDIX_ONLY:
        if slide is not None and slide.deck_mode != DeckMode.APPENDIX and slide.content_tier != ContentTier.APPENDIX_ONLY:
            return slide.content_tier
        return ContentTier.SUPPORTING_EXAMPLE
    if entry.content_tier is not None:
        return entry.content_tier
    if slide is not None:
        return slide.content_tier
    return ContentTier.LECTURE_CORE


def _qa_status(findings: list[QAFinding]) -> QAStatus:
    if any(f.severity in {QASeverity.CRITICAL, QASeverity.MAJOR} for f in findings):
        return QAStatus.FAIL
    return QAStatus.CONDITIONAL_PASS if findings else QAStatus.PASS


def _slide_status(findings: list[QAFinding]) -> tuple[QAStatus, int, int]:
    blocking = sum(1 for f in findings if f.blocking)
    warnings = sum(1 for f in findings if not f.blocking)
    return _qa_status(findings), warnings, blocking


def _shape_type_name(shape) -> str:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "unknown"
    return getattr(shape_type, "name", str(shape_type)).lower()


def _iter_nested_shapes(container) -> Any:
    shapes = getattr(container, "shapes", container)
    for shape in shapes:
        yield shape
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            yield from _iter_nested_shapes(child_shapes)


def _snapshots(pptx_path: Path) -> tuple[list[SlideSnapshot], str | None]:
    try:
        presentation = Presentation(str(pptx_path))
    except Exception as exc:  # pragma: no cover
        return [], str(exc)
    slide_width = float(presentation.slide_width or 1)
    slide_height = float(presentation.slide_height or 1)
    snapshots: list[SlideSnapshot] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_text = _norm(title_shape.text if title_shape is not None else "")
        text_char_count = 0
        paragraph_count = 0
        picture_count = 0
        table_count = 0
        text_blocks: list[str] = []
        table_blocks: list[str] = []
        text_shape_count = 0
        placeholder_count = 0
        auto_shape_count = 0
        group_count = 0
        connector_count = 0
        shape_type_counts: Counter[str] = Counter()
        shape_boxes: list[ShapeBox] = []
        for shape in _iter_nested_shapes(slide.shapes):
            shape_type_name = _shape_type_name(shape)
            shape_type_counts[shape_type_name] += 1
            left = max(0.0, float(getattr(shape, "left", 0) or 0) / slide_width)
            top = max(0.0, float(getattr(shape, "top", 0) or 0) / slide_height)
            width = max(0.0, float(getattr(shape, "width", 0) or 0) / slide_width)
            height = max(0.0, float(getattr(shape, "height", 0) or 0) / slide_height)
            if getattr(shape, "is_placeholder", False):
                placeholder_count += 1
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                group_count += 1
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE:
                connector_count += 1
            if getattr(shape, "shape_type", None) in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.TEXT_BOX}:
                auto_shape_count += 1
            if getattr(shape, "has_table", False):
                table_count += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = _norm(cell.text)
                        if cell_text:
                            table_blocks.append(cell_text)
            if getattr(shape, "has_text_frame", False):
                text_shape_count += 1
                text = _norm(shape.text_frame.text)
                text_char_count += len(text)
                paragraph_count += sum(1 for p in shape.text_frame.paragraphs if _norm(p.text))
                if text:
                    text_blocks.append(text)
            else:
                text = ""
            if getattr(shape, "has_table", False):
                row_count = len(shape.table.rows)
                column_count = len(shape.table.columns)
                kind = "table"
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                row_count = 0
                column_count = 0
                kind = "picture"
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE:
                row_count = 0
                column_count = 0
                kind = "line"
            elif getattr(shape, "shape_type", None) in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM}:
                row_count = 0
                column_count = 0
                kind = "auto_shape"
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TEXT_BOX:
                row_count = 0
                column_count = 0
                kind = "text_box"
            else:
                row_count = 0
                column_count = 0
                kind = shape_type_name
            if width > 0 and height > 0:
                shape_boxes.append(
                    ShapeBox(
                        kind=kind,
                        left=round(left, 4),
                        top=round(top, 4),
                        width=round(width, 4),
                        height=round(height, 4),
                        text=text,
                        is_title=shape is title_shape,
                        row_count=row_count,
                        column_count=column_count,
                    )
                )
        if not title_text:
            for block in text_blocks:
                if re.fullmatch(r"(appendix\s+)?\d+", block, flags=re.IGNORECASE):
                    continue
                title_text = block
                break
        snapshots.append(
            SlideSnapshot(
                index,
                title_text,
                " ".join(text_blocks + table_blocks),
                text_char_count,
                paragraph_count,
                picture_count,
                table_count,
                text_blocks,
                table_blocks,
                text_shape_count,
                placeholder_count,
                auto_shape_count,
                group_count,
                connector_count,
                dict(sorted(shape_type_counts.items())),
                shape_boxes,
            )
        )
    return snapshots, None


def _write_json_artifact(path: Path, payload: dict[str, Any]) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def _load_optional_json(path: Path) -> dict[str, Any] | None:
    if not path.is_file():
        return None
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None
    return payload if isinstance(payload, dict) else None


def _compiled_duplicate_blocks(snapshot: SlideSnapshot) -> list[str]:
    title_norm = _normalized_block(snapshot.title_text)
    normalized_text_blocks = [_normalized_block(block) for block in snapshot.text_blocks if _normalized_block(block)]
    normalized_table_blocks = [_normalized_block(block) for block in snapshot.table_blocks if _normalized_block(block)]
    duplicates: list[str] = []
    repeated_counts = Counter(normalized_text_blocks + normalized_table_blocks)
    for block, count in repeated_counts.items():
        if len(block) >= 24 and count > 1:
            duplicates.append(block)
    if title_norm:
        for block in normalized_text_blocks[1:] + normalized_table_blocks:
            if len(block) < max(18, min(28, len(title_norm))):
                continue
            if block == title_norm or title_norm in block or block in title_norm:
                duplicates.append(block)
    return _dedupe(duplicates)


def _compiled_chrome_hits(snapshot: SlideSnapshot) -> list[dict[str, str]]:
    hits: list[dict[str, str]] = []
    seen: set[tuple[str, str]] = set()
    for block in snapshot.text_blocks + snapshot.table_blocks:
        block_norm = _normalized_block(block)
        if not block_norm:
            continue
        for phrase in CHROME_PHRASES:
            if phrase in block_norm:
                key = (phrase, block)
                if key in seen:
                    continue
                seen.add(key)
                hits.append({"phrase": phrase, "block": block})
    return hits


def _compiled_opening_formula(snapshot: SlideSnapshot) -> str:
    title_norm = _normalized_block(snapshot.title_text)
    generic_openings = {
        "core definition",
        "main limitation",
        "evidence summary",
        "source lane 1",
        "source lane 2",
        "source lane 3",
    }
    for block in snapshot.text_blocks:
        block_norm = _normalized_block(block)
        if not block_norm or block_norm == title_norm:
            continue
        if re.fullmatch(r"(step\s+)?\d+", block_norm):
            continue
        if snapshot.connector_count > 0 and len(block_norm.split()) <= 4:
            continue
        if block_norm in generic_openings or block_norm.startswith("use the "):
            continue
        return _formula_prefix(block_norm, words=4)
    return ""


def _cluster_overflow(
    indices: list[int],
    *,
    allowed: int,
    max_gap: int = 1,
) -> tuple[int, list[int]]:
    if not indices:
        return 0, []
    overflow = 0
    overflow_slides: list[int] = []
    cluster: list[int] = [indices[0]]
    for index in indices[1:]:
        if index - cluster[-1] <= max_gap + 1:
            cluster.append(index)
            continue
        if len(cluster) > allowed:
            overflow += len(cluster) - allowed
            overflow_slides.extend(cluster[allowed:])
        cluster = [index]
    if len(cluster) > allowed:
        overflow += len(cluster) - allowed
        overflow_slides.extend(cluster[allowed:])
    return overflow, overflow_slides


def _sequence_overflow(
    sequence: list[tuple[int, str]],
    *,
    allowed: int,
) -> tuple[int, list[int], int]:
    if not sequence:
        return 0, [], 0
    overflow = 0
    overflow_slides: list[int] = []
    max_run = 0
    previous_value: str | None = None
    run_numbers: list[int] = []
    for slide_number, value in sequence:
        if value and value == previous_value:
            run_numbers.append(slide_number)
        else:
            if len(run_numbers) > allowed:
                overflow += len(run_numbers) - allowed
                overflow_slides.extend(run_numbers[allowed:])
            previous_value = value
            run_numbers = [slide_number]
        max_run = max(max_run, len(run_numbers))
    if len(run_numbers) > allowed:
        overflow += len(run_numbers) - allowed
        overflow_slides.extend(run_numbers[allowed:])
    return overflow, overflow_slides, max_run


def _is_bridge_shell(slide, realized_archetype: str) -> bool:
    return slide.slide_intent == SlideIntent.MAPPING_BRIDGE or realized_archetype in {"mapping-table", "correspondence-matrix"}


def _is_cycle_shell(slide, realized_archetype: str, title_text: str) -> bool:
    title = _norm(title_text).lower()
    cycle_keywords = (
        "selection",
        "crossover",
        "mutation",
        "operator",
        "cycle",
        "pressure",
        "variation",
        "linkage",
        "building blocks",
    )
    if not any(keyword in title for keyword in cycle_keywords):
        return False
    return (
        slide.slide_intent in {SlideIntent.MECHANISM_WALKTHROUGH, SlideIntent.COMPARISON_TRADEOFF, SlideIntent.MISCONCEPTION_PITFALL}
        or realized_archetype in {"process-flow", "comparison-matrix", "comparison-callout-cluster", "worked-example-state-table"}
    )


def _compiled_title_template(title_text: str) -> str:
    normalized = _normalized_block(title_text)
    if not normalized:
        return ""
    if (
        normalized.startswith("source map for ")
        or normalized.startswith("evidence summary for ")
        or normalized.startswith("evidence cluster: ")
        or normalized.endswith(" evidence cluster")
        or normalized.endswith(" source cluster")
    ):
        return "appendix-reference-card"
    return _formula_prefix(normalized, words=3)


def _compiled_source_map_title_pattern(title_text: str) -> bool:
    normalized = _normalized_block(title_text)
    return normalized.startswith(("source map for ", "evidence summary for ", "evidence cluster: "))


def _appendix_geometry_signature(realized_archetype: str, snapshot: SlideSnapshot) -> str:
    _ = snapshot
    return realized_archetype


def _shape_area(box: ShapeBox) -> float:
    return max(0.0, box.width) * max(0.0, box.height)


def _non_title_boxes(snapshot: SlideSnapshot) -> list[ShapeBox]:
    return [
        box
        for box in snapshot.shape_boxes
        if not box.is_title and box.kind != "line" and _shape_area(box) >= 0.003
    ]


def _quantize_fraction(value: float, *, step: float = 0.08) -> str:
    if value <= 0:
        return "0.00"
    return f"{round(value / step) * step:.2f}"


def _visual_signature(snapshot: SlideSnapshot, realized_archetype: str) -> str:
    major_boxes = sorted(_non_title_boxes(snapshot), key=_shape_area, reverse=True)
    parts = [realized_archetype]
    for box in major_boxes[:4]:
        part = ":".join(
            [
                box.kind,
                _quantize_fraction(box.left),
                _quantize_fraction(box.top),
                _quantize_fraction(box.width),
                _quantize_fraction(box.height),
            ]
        )
        if box.kind == "table":
            part = f"{part}:r{min(box.row_count, 8)}:c{min(box.column_count, 6)}"
        parts.append(part)
    if snapshot.connector_count:
        parts.append(f"ln{min(snapshot.connector_count, 4)}")
    if snapshot.picture_count:
        parts.append(f"pic{min(snapshot.picture_count, 4)}")
    return "|".join(parts)


def _is_text_card_like(snapshot: SlideSnapshot, realized_archetype: str) -> bool:
    if snapshot.table_count > 0 or snapshot.picture_count > 0 or snapshot.connector_count > 0:
        return False
    if realized_archetype not in TEXT_CARD_ARCHETYPES and not realized_archetype.startswith("appendix-"):
        return False
    content_boxes = _non_title_boxes(snapshot)
    return len(content_boxes) >= 3


def _is_chrome_dominant(snapshot: SlideSnapshot) -> bool:
    chrome_boxes = [
        box
        for box in _non_title_boxes(snapshot)
        if any(phrase in _normalized_block(box.text) for phrase in CHROME_PHRASES)
    ]
    if not chrome_boxes:
        return False
    content_boxes = _non_title_boxes(snapshot)
    total_chars = sum(len(_normalized_block(box.text)) for box in content_boxes if box.text)
    chrome_chars = sum(len(_normalized_block(box.text)) for box in chrome_boxes if box.text)
    total_area = sum(_shape_area(box) for box in content_boxes) or 1.0
    chrome_area = sum(_shape_area(box) for box in chrome_boxes)
    return (
        len(chrome_boxes) >= 2
        or chrome_chars / max(total_chars, 1) >= 0.22
        or chrome_area / total_area >= 0.28
    )


def _missing_visual_center(snapshot: SlideSnapshot, realized_archetype: str) -> bool:
    if snapshot.table_count > 0 or snapshot.picture_count > 0:
        return False
    if realized_archetype == "process-flow" and snapshot.connector_count >= 1:
        return False
    content_boxes = _non_title_boxes(snapshot)
    if not content_boxes:
        return True
    if any(
        _shape_area(box) >= 0.12
        and 0.18 <= box.left + box.width / 2 <= 0.82
        and 0.18 <= box.top + box.height / 2 <= 0.82
        for box in content_boxes
    ):
        return False
    largest_area = max(_shape_area(box) for box in content_boxes)
    return snapshot.text_shape_count >= 6 and len(content_boxes) >= 5 and largest_area < 0.12


def _text_overflow_risk(snapshot: SlideSnapshot) -> bool:
    for box in _non_title_boxes(snapshot):
        if box.kind not in {"text_box", "auto_shape"} or not box.text:
            continue
        normalized = _normalized_block(box.text)
        if len(normalized) < 90:
            continue
        estimated_capacity = max(90, int(_shape_area(box) * 2400))
        if len(normalized) > estimated_capacity:
            return True
    return False


def _box_overlap_fraction(left: ShapeBox, right: ShapeBox) -> float:
    overlap_left = max(left.left, right.left)
    overlap_top = max(left.top, right.top)
    overlap_right = min(left.left + left.width, right.left + right.width)
    overlap_bottom = min(left.top + left.height, right.top + right.height)
    overlap_width = max(0.0, overlap_right - overlap_left)
    overlap_height = max(0.0, overlap_bottom - overlap_top)
    overlap_area = overlap_width * overlap_height
    return overlap_area / max(min(_shape_area(left), _shape_area(right)), 0.0001)


def _overlap_collision_risk(snapshot: SlideSnapshot) -> bool:
    content_boxes = [
        box
        for box in _non_title_boxes(snapshot)
        if box.text or box.kind in {"table", "picture"}
    ]
    for index, left in enumerate(content_boxes):
        if _shape_area(left) < 0.01:
            continue
        for right in content_boxes[index + 1 :]:
            if _shape_area(right) < 0.01:
                continue
            if _box_overlap_fraction(left, right) >= 0.18:
                return True
    return False


def _weak_title_body_balance(entry, snapshot: SlideSnapshot) -> bool:
    if entry.slide_role in {SlideRole.TITLE, SlideRole.SECTION_DIVIDER}:
        return False
    title_chars = len(_normalized_block(snapshot.title_text))
    body_chars = max(snapshot.text_char_count - title_chars, 0)
    if title_chars >= 56 and body_chars < 32:
        return True
    return title_chars >= max(64, body_chars)


def _missing_support_marker(slide, snapshot: SlideSnapshot) -> bool:
    if slide.evidence_class not in {SlideEvidenceClass.SOURCE_BACKED, SlideEvidenceClass.APPENDIX_SUPPORT}:
        return False
    labels = [ref.label.lower() for ref in slide.production_bridge.source_material_refs[:2]]
    supporting = [_normalized_block(item) for item in slide.supporting_evidence[:2]]
    all_text = snapshot.all_text.lower()
    if "source:" in all_text or "evidence:" in all_text:
        return False
    if labels and any(label in all_text for label in labels):
        return False
    if supporting and any(item in all_text for item in supporting if item):
        return False
    return bool(labels or supporting)


def _weak_visual_anchor(slide, snapshot: SlideSnapshot, realized_archetype: str) -> bool:
    if slide.visual_type in {VisualType.TEXT, VisualType.QUOTE}:
        return False
    if snapshot.picture_count > 0 or snapshot.table_count > 0:
        return False
    if realized_archetype == "process-flow" and snapshot.connector_count >= 1:
        return False
    content_boxes = _non_title_boxes(snapshot)
    if not content_boxes:
        return True
    largest_area = max(_shape_area(box) for box in content_boxes)
    return snapshot.text_shape_count >= 4 and largest_area < 0.1


def _realized_archetype(entry, slide, link, snapshot: SlideSnapshot | None) -> str:
    if snapshot is None:
        return "unavailable"
    title = _norm(snapshot.title_text or entry.title).lower()
    if entry.deck_mode == DeckMode.APPENDIX:
        if slide.slide_archetype == SlideArchetype.APPENDIX_SOURCE_LOCATION_MATRIX and snapshot.table_count > 0:
            return "appendix-source-location-matrix"
        if slide.slide_archetype == SlideArchetype.APPENDIX_COMPARISON_EVIDENCE_CLUSTER and snapshot.table_count > 0:
            return "appendix-comparison-evidence-cluster"
        if slide.slide_archetype == SlideArchetype.APPENDIX_SOURCE_MAP and snapshot.table_count > 0:
            return "appendix-source-map"
        if slide.slide_archetype == SlideArchetype.APPENDIX_THEMED_EVIDENCE_CLUSTER and snapshot.table_count == 0 and snapshot.auto_shape_count >= 2:
            return "appendix-themed-evidence-cluster"
        if slide.slide_archetype == SlideArchetype.APPENDIX_ANNOTATED_EXCERPT_CLUSTER and snapshot.table_count == 0 and snapshot.auto_shape_count >= 2:
            return "appendix-annotated-excerpt-cluster"
        if title.endswith((" source-location matrix", " source matrix")) and snapshot.table_count > 0:
            return "appendix-source-location-matrix"
        if title.endswith(" evidence comparison") and snapshot.table_count > 0:
            return "appendix-comparison-evidence-cluster"
        if title.endswith(" cross-reference map") and snapshot.table_count > 0:
            return "appendix-source-map"
        if title.endswith((" excerpt cluster", " source excerpts")) and snapshot.table_count == 0 and snapshot.auto_shape_count >= 2:
            return "appendix-annotated-excerpt-cluster"
        if title.startswith(("source map for", "evidence summary for", "evidence cluster:", "selection analogy evidence cluster", "algorithm loop evidence cluster")):
            return "appendix-evidence-cluster"
        if title.endswith(" evidence cluster") and snapshot.table_count == 0 and snapshot.auto_shape_count >= 2:
            return "appendix-themed-evidence-cluster"
        if title.endswith(" evidence cluster") and snapshot.table_count > 0:
            return "appendix-evidence-cluster"
        return "appendix-reference"
    if snapshot.table_count > 0:
        if slide.slide_archetype == SlideArchetype.WORKED_EXAMPLE_STATE_TABLE:
            return "worked-example-state-table"
        if slide.slide_archetype == SlideArchetype.CORRESPONDENCE_MATRIX:
            return "correspondence-matrix"
        if slide.slide_archetype == SlideArchetype.COMPARISON_MATRIX or slide.slide_intent == SlideIntent.COMPARISON_TRADEOFF:
            return "comparison-matrix"
        if slide.slide_archetype == SlideArchetype.TWO_COLUMN_MAPPING_TABLE:
            return "mapping-table"
        if slide.slide_intent == SlideIntent.WORKED_EXAMPLE or "worked example" in title or "toy population" in title:
            return "worked-example-state-table"
        if "map" in title or "correspondence" in title:
            return "correspondence-matrix"
        if "encoding" in title or "operator" in title or "fitness" in title:
            return "mapping-table"
        return "table-structured"
    if link.layout_family == "process-flow" and snapshot.connector_count >= 1 and (snapshot.text_shape_count >= 5 or "cycle" in title or title.startswith("how one")):
        return "process-flow"
    if link.layout_family == "appendix-reference":
        return "appendix-reference"
    if link.layout_family == "cover":
        return "title-orientation"
    if link.layout_family == "comparison" and snapshot.text_shape_count >= 4:
        return "comparison-callout-cluster"
    if link.layout_family == "concept-explainer" and snapshot.text_shape_count >= 4:
        return "anchor-concept-card"
    return "generic-text-card"


def _compiled_expected_visual_missing(slide, link, snapshot: SlideSnapshot | None, realized_archetype: str) -> bool:
    if snapshot is None:
        return False
    expected = _expected_archetypes(slide)
    if not expected:
        return False
    if slide.slide_archetype in MAPPING_ARCHETYPES or expected & MAPPING_ARCHETYPES:
        return snapshot.table_count < 1
    if slide.slide_archetype == SlideArchetype.COMPARISON_MATRIX:
        return snapshot.table_count < 1
    if slide.slide_archetype in APPENDIX_TABLE_ARCHETYPES:
        return snapshot.table_count < 1
    if slide.slide_archetype in {
        SlideArchetype.APPENDIX_THEMED_EVIDENCE_CLUSTER,
        SlideArchetype.APPENDIX_ANNOTATED_EXCERPT_CLUSTER,
    }:
        return snapshot.auto_shape_count < 2
    if slide.slide_archetype == SlideArchetype.APPENDIX_EVIDENCE_CLUSTER:
        return False
    if slide.slide_archetype == SlideArchetype.WORKED_EXAMPLE_STATE_TABLE or expected & WORKED_EXAMPLE_ARCHETYPES:
        stage_labels = [row[0].strip().lower() for row in slide.authoring_payload.get("rows", []) if isinstance(row, list) and row]
        state_progression_missing = bool(stage_labels) and any(label not in snapshot.all_text.lower() for label in stage_labels[: min(len(stage_labels), 3)])
        return snapshot.table_count < 1 or state_progression_missing
    if slide.slide_archetype in PROCESS_ARCHETYPES or expected & PROCESS_ARCHETYPES:
        return realized_archetype != "process-flow" or snapshot.connector_count < 1
    return False


def _compiled_realization_mismatch(slide, realized_archetype: str, expected_visual_missing: bool) -> bool:
    if expected_visual_missing:
        return True
    expected = _expected_archetypes(slide)
    if not expected:
        return False
    if slide.slide_archetype in MAPPING_ARCHETYPES:
        return realized_archetype not in {"mapping-table", "correspondence-matrix", "table-structured", "worked-example-state-table"}
    if slide.slide_archetype == SlideArchetype.COMPARISON_MATRIX:
        return realized_archetype != "comparison-matrix"
    if slide.slide_archetype in PROCESS_ARCHETYPES:
        return realized_archetype != "process-flow"
    if slide.slide_archetype == SlideArchetype.APPENDIX_THEMED_EVIDENCE_CLUSTER:
        return realized_archetype != "appendix-themed-evidence-cluster"
    if slide.slide_archetype == SlideArchetype.APPENDIX_SOURCE_LOCATION_MATRIX:
        return realized_archetype != "appendix-source-location-matrix"
    if slide.slide_archetype == SlideArchetype.APPENDIX_ANNOTATED_EXCERPT_CLUSTER:
        return realized_archetype != "appendix-annotated-excerpt-cluster"
    if slide.slide_archetype == SlideArchetype.APPENDIX_COMPARISON_EVIDENCE_CLUSTER:
        return realized_archetype != "appendix-comparison-evidence-cluster"
    if slide.slide_archetype == SlideArchetype.APPENDIX_SOURCE_MAP:
        return realized_archetype != "appendix-source-map"
    if slide.slide_archetype == SlideArchetype.APPENDIX_EVIDENCE_CLUSTER:
        return realized_archetype not in {"appendix-evidence-cluster", "appendix-themed-evidence-cluster", "appendix-reference"}
    return False


def _export_slide_pngs_with_powerpoint(pptx_path: Path, export_dir: Path) -> list[Path]:
    try:  # pragma: no cover - depends on local PowerPoint availability
        import pythoncom
        import win32com.client
    except Exception:  # pragma: no cover
        return []
    application = None
    deck = None
    pythoncom.CoInitialize()
    try:  # pragma: no cover - exercised only on Windows with PowerPoint installed
        export_dir.mkdir(parents=True, exist_ok=True)
        application = win32com.client.DispatchEx("PowerPoint.Application")
        application.Visible = 1
        deck = application.Presentations.Open(str(pptx_path), WithWindow=False)
        deck.Export(str(export_dir), "PNG", 320, 180)
        return sorted(export_dir.glob("*.png")) + sorted(export_dir.glob("*.PNG"))
    except Exception:
        return []
    finally:  # pragma: no cover - cleanup path
        if deck is not None:
            try:
                deck.Close()
            except Exception:
                pass
        if application is not None:
            try:
                application.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def _render_text_thumbnail_strip(output_path: Path, snapshots: list[SlideSnapshot]) -> Path:
    from PIL import Image, ImageDraw

    thumb_width = 320
    thumb_height = 180
    gap = 12
    columns = 3
    rows = max(1, (len(snapshots) + columns - 1) // columns)
    canvas = Image.new("RGB", (columns * thumb_width + (columns + 1) * gap, rows * thumb_height + (rows + 1) * gap), "#F4F4F0")
    draw = ImageDraw.Draw(canvas)

    def _wrap(text: str, width: int = 34) -> list[str]:
        words = _norm(text).split()
        lines: list[str] = []
        current: list[str] = []
        for word in words:
            trial = " ".join(current + [word])
            if len(trial) <= width:
                current.append(word)
            else:
                if current:
                    lines.append(" ".join(current))
                current = [word]
        if current:
            lines.append(" ".join(current))
        return lines[:6]

    for position, snapshot in enumerate(snapshots):
        row = position // columns
        column = position % columns
        left = gap + column * (thumb_width + gap)
        top = gap + row * (thumb_height + gap)
        right = left + thumb_width
        bottom = top + thumb_height
        draw.rounded_rectangle((left, top, right, bottom), radius=16, fill="#FFFFFF", outline="#D8D6D0", width=2)
        draw.text((left + 12, top + 10), f"{snapshot.pptx_index}", fill="#6E6A61")
        title = snapshot.title_text or f"Slide {snapshot.pptx_index}"
        y = top + 34
        for line in _wrap(title, width=30)[:2]:
            draw.text((left + 12, y), line, fill="#1E1E1E")
            y += 18
        body_blocks = []
        for block in snapshot.text_blocks:
            if _normalized_block(block) == _normalized_block(title):
                continue
            body_blocks.append(block)
            if len(body_blocks) == 2:
                break
        meta_line = f"text {snapshot.text_shape_count} | table {snapshot.table_count} | pic {snapshot.picture_count}"
        draw.text((left + 12, bottom - 26), meta_line, fill="#58544D")
        y = max(y + 6, top + 76)
        for block in body_blocks:
            for line in _wrap(block, width=34)[:2]:
                if y > bottom - 48:
                    break
                draw.text((left + 12, y), line, fill="#4C4A44")
                y += 16
    output_path.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(output_path)
    return output_path


def _write_compiled_deck_thumbnail_strip(
    *,
    pptx_path: Path,
    output_path: Path,
    snapshots: list[SlideSnapshot],
) -> Path:
    export_dir = output_path.parent / "compiled-deck-thumbnails"
    exported: list[Path] = []
    if os.getenv("PPTXLOCAL_ENABLE_POWERPOINT_EXPORT") == "1":
        exported = _export_slide_pngs_with_powerpoint(pptx_path, export_dir)
    if not exported:
        return _render_text_thumbnail_strip(output_path, snapshots)

    from PIL import Image, ImageDraw

    thumb_width = 320
    thumb_height = 180
    gap = 12
    columns = 3
    rows = max(1, (len(exported) + columns - 1) // columns)
    canvas = Image.new("RGB", (columns * thumb_width + (columns + 1) * gap, rows * thumb_height + (rows + 1) * gap), "#F4F4F0")
    draw = ImageDraw.Draw(canvas)
    for position, image_path in enumerate(exported):
        row = position // columns
        column = position % columns
        left = gap + column * (thumb_width + gap)
        top = gap + row * (thumb_height + gap)
        image = Image.open(image_path).convert("RGB")
        image.thumbnail((thumb_width, thumb_height))
        paste_left = left + (thumb_width - image.width) // 2
        paste_top = top + (thumb_height - image.height) // 2
        canvas.paste(image, (paste_left, paste_top))
        draw.rectangle((left, top, left + thumb_width, top + thumb_height), outline="#D8D6D0", width=2)
        draw.text((left + 10, top + 8), f"{position + 1}", fill="#FFFFFF")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(output_path)
    return output_path


def _shape_box_from_payload(payload: dict[str, Any]) -> ShapeBox:
    return ShapeBox(
        kind=str(payload.get("kind") or "unknown"),
        left=float(payload.get("left") or 0.0),
        top=float(payload.get("top") or 0.0),
        width=float(payload.get("width") or 0.0),
        height=float(payload.get("height") or 0.0),
        text=str(payload.get("text") or ""),
        is_title=bool(payload.get("is_title")),
        row_count=int(payload.get("row_count") or 0),
        column_count=int(payload.get("column_count") or 0),
    )


def _schematic_fill(box: ShapeBox, deck_mode: str) -> str:
    if box.is_title:
        return "#212121"
    if box.kind == "table":
        return "#D9E7F5"
    if box.kind == "picture":
        return "#D7EAD8"
    if deck_mode == "appendix":
        return "#EFE9DD" if box.kind == "auto_shape" else "#F7F2EA"
    return "#E8E3D8" if box.kind == "auto_shape" else "#F4F1EA"


def _draw_text_density(draw, box: ShapeBox, *, left: int, top: int, width: int, height: int) -> None:
    if not box.text or width < 18 or height < 10:
        return
    usable_width = max(8, width - 10)
    lines = min(4, max(1, len(_normalized_block(box.text)) // 40 + 1))
    start_y = top + 5
    spacing = max(6, min(10, (height - 10) // max(lines, 1)))
    for offset in range(lines):
        y = start_y + offset * spacing
        if y >= top + height - 4:
            break
        draw.line((left + 5, y, left + 5 + usable_width, y), fill="#5C5750", width=2)


def _render_schematic_slide_thumbnail(
    *,
    output_path: Path,
    slide_payload: dict[str, Any],
    audit_payload: dict[str, Any],
    thumb_width: int = 320,
    thumb_height: int = 180,
) -> Path:
    from PIL import Image, ImageDraw

    image = Image.new("RGB", (thumb_width, thumb_height), "#FBFAF7")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((1, 1, thumb_width - 2, thumb_height - 2), radius=12, outline="#D8D3C9", width=2, fill="#FFFFFF")
    deck_mode = str(slide_payload.get("deck_mode") or "main-story")
    if deck_mode == "appendix":
        draw.rectangle((0, 0, thumb_width, 10), fill="#D9C6A8")
    else:
        draw.rectangle((0, 0, thumb_width, 10), fill="#B8C7D9")
    for box_payload in slide_payload.get("shape_boxes", []):
        if not isinstance(box_payload, dict):
            continue
        box = _shape_box_from_payload(box_payload)
        left = int(box.left * thumb_width)
        top = int(box.top * thumb_height)
        width = max(2, int(box.width * thumb_width))
        height = max(2, int(box.height * thumb_height))
        if box.kind == "line":
            draw.line((left, top, left + width, top + height), fill="#7D8FA3", width=3)
            continue
        fill = _schematic_fill(box, deck_mode)
        outline = "#B3AEA6" if not box.is_title else "#212121"
        draw.rounded_rectangle((left, top, left + width, top + height), radius=6, fill=fill, outline=outline, width=2)
        if box.kind == "table" and box.row_count > 0 and box.column_count > 0:
            for row in range(1, min(box.row_count, 6)):
                y = top + int(height * row / box.row_count)
                draw.line((left, y, left + width, y), fill="#A5B5C5", width=1)
            for column in range(1, min(box.column_count, 5)):
                x = left + int(width * column / box.column_count)
                draw.line((x, top, x, top + height), fill="#A5B5C5", width=1)
        _draw_text_density(draw, box, left=left, top=top, width=width, height=height)
    slide_number = int(slide_payload.get("pptx_index") or slide_payload.get("slide_number") or 0)
    draw.rounded_rectangle((10, 12, 44, 34), radius=8, fill="#1E1E1E")
    draw.text((19, 17), str(slide_number), fill="#FFFFFF")
    archetype = str(audit_payload.get("realized_archetype") or "")
    if archetype:
        draw.text((12, thumb_height - 18), archetype[:28], fill="#5B574F")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path)
    return output_path


def _write_visual_thumbnail_assets(
    *,
    output_dir: Path,
    deck_title: str,
    pptx_path: str,
    shape_slides: list[dict[str, Any]],
    audit_slides: list[dict[str, Any]],
) -> tuple[Path, Path, Path]:
    audit_by_number = {
        int(payload.get("slide_number") or 0): payload
        for payload in audit_slides
        if isinstance(payload, dict) and isinstance(payload.get("slide_number"), int)
    }
    output_dir.mkdir(parents=True, exist_ok=True)
    index_payload = {
        "schema_name": "compiled_deck_thumbnail_index",
        "schema_version": "1.0",
        "deck_title": deck_title,
        "pptx_path": pptx_path,
        "render_mode": "schematic-geometry",
        "slides": [],
    }
    thumb_width = 320
    thumb_height = 180
    gap = 12
    columns = 4
    rows = max(1, (len(shape_slides) + columns - 1) // columns)
    from PIL import Image, ImageDraw

    strip = Image.new(
        "RGB",
        (columns * thumb_width + (columns + 1) * gap, rows * thumb_height + (rows + 1) * gap),
        "#F4F4F0",
    )
    strip_draw = ImageDraw.Draw(strip)
    for position, slide_payload in enumerate(shape_slides):
        slide_number = int(slide_payload.get("slide_number") or 0)
        audit_payload = audit_by_number.get(slide_number, {})
        thumb_path = output_dir / f"slide-{slide_number:03d}.png"
        _render_schematic_slide_thumbnail(
            output_path=thumb_path,
            slide_payload=slide_payload,
            audit_payload=audit_payload,
            thumb_width=thumb_width,
            thumb_height=thumb_height,
        )
        thumb_image = Image.open(thumb_path).convert("RGB")
        row = position // columns
        column = position % columns
        left = gap + column * (thumb_width + gap)
        top = gap + row * (thumb_height + gap)
        strip.paste(thumb_image, (left, top))
        strip_draw.rectangle((left, top, left + thumb_width, top + thumb_height), outline="#D8D6D0", width=2)
        index_payload["slides"].append(
            {
                "slide_number": slide_number,
                "slide_id": slide_payload.get("slide_id"),
                "pptx_index": slide_payload.get("pptx_index"),
                "deck_mode": slide_payload.get("deck_mode"),
                "title_text": audit_payload.get("title_text") or slide_payload.get("title_text"),
                "realized_archetype": audit_payload.get("realized_archetype"),
                "visual_signature": audit_payload.get("visual_signature"),
                "thumbnail_path": str(thumb_path),
            }
        )
    strip_path = output_dir.parent / "compiled-deck-thumbnail-strip.png"
    strip_path.parent.mkdir(parents=True, exist_ok=True)
    strip.save(strip_path)
    index_path = output_dir.parent / "compiled-deck-thumbnail-index.json"
    index_path.write_text(json.dumps(index_payload, ensure_ascii=False, indent=2), encoding="utf-8")
    summary_path = output_dir.parent / "compiled-deck-visual-review-summary.json"
    return strip_path, index_path, summary_path


def _finding(
    *,
    finding_id: str,
    severity: QASeverity,
    qa_layer: QALayer,
    category: str,
    summary: str,
    remediation_skill: str,
    recommendation_type: QARecommendationType,
    recommendation: str,
    slide_number: int | None = None,
    slide_id: str | None = None,
    slide_range: SlideRange | None = None,
    build_link_index: int | None = None,
    blocking: bool = False,
    tags: list[str] | None = None,
) -> QAFinding:
    payload: dict[str, object] = {
        "finding_id": finding_id,
        "severity": severity,
        "status": FindingStatus.OPEN,
        "qa_layer": qa_layer,
        "category": category,
        "summary": summary,
        "remediation_skill": remediation_skill,
        "recommendation_type": recommendation_type,
        "recommendation": recommendation,
        "blocking": blocking,
        "tags": tags or [],
    }
    if slide_number is not None:
        payload["slide_number"] = slide_number
    if slide_id is not None:
        payload["slide_id"] = slide_id
    if slide_range is not None:
        payload["slide_range"] = slide_range
    if build_link_index is not None:
        payload["build_link_index"] = build_link_index
    return QAFinding.model_validate(payload)


def run_deck_qa(
    blueprint: Blueprint,
    design_system: DesignSystem,
    deck_constitution: DeckConstitution,
    layout_library: LayoutLibrary,
    slide_ledger: SlideLedger,
    asset_manifest: AssetManifest,
    viz_manifest: VizManifest,
    build_manifest: BuildManifest,
    slide_build_linkage: SlideBuildLinkage,
    *,
    state_capsule: StateCapsule | None = None,
    prior_report: QAReport | None = None,
    qa_governance: QAGovernance | None = None,
    checked_artifacts: list[str] | None = None,
    artifact_root: str | Path | None = None,
) -> DeckQAOutputs:
    root = Path(artifact_root).resolve() if artifact_root is not None else None
    deck_range = _deck_range(slide_ledger)
    findings: list[QAFinding] = []
    blueprint_by_number = {slide.slide_number: slide for slide in blueprint.slides}
    linkage_by_number = {slide.slide_number: slide for slide in slide_build_linkage.slides}
    asset_by_id = {asset.asset_id: asset for asset in asset_manifest.assets}
    approved_request_ids = {asset.request_id for asset in asset_manifest.assets if asset.status in {AssetStatus.APPROVED, AssetStatus.READY}}
    viz_by_spec_id = {record.spec.spec_id: record for record in viz_manifest.visuals}
    color_tokens = {token.token for token in design_system.color_tokens}
    typography_tokens = {token.token for token in design_system.typography_tokens}
    layout_index = {pattern.pattern_id: pattern for pattern in layout_library.patterns}
    pptx_path = _path(build_manifest.pptx_path, root)
    snapshots, pptx_error = _snapshots(pptx_path) if pptx_path.is_file() else ([], None)
    snapshot_by_index = {item.pptx_index: item for item in snapshots}
    lecture_mode = _is_lecture_blueprint(blueprint)
    authoring_preview = _load_optional_json(root / "state" / "authoring-preview.json") if root is not None else None
    authoring_preview_slides = {
        int(slide_payload.get("slide_number")): slide_payload
        for slide_payload in (authoring_preview.get("slides", []) if isinstance(authoring_preview, dict) else [])
        if isinstance(slide_payload, dict) and isinstance(slide_payload.get("slide_number"), int)
    }

    if len(blueprint.slides) != len(slide_ledger.entries) or build_manifest.slide_count != len(slide_ledger.entries):
        findings.append(_finding(finding_id="qa-build-count", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="numbering", summary="Blueprint, ledger, and build output disagree on slide count.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Rebuild the deck from a reconciled blueprint and ledger set.", slide_range=deck_range, blocking=True, tags=["build", "numbering"]))
    if not pptx_path.is_file():
        findings.append(_finding(finding_id="qa-pptx-missing", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="build", summary="The compiled PPTX file is missing at the recorded build path.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Re-run compile-pptx and persist the deck artifact before QA.", slide_range=deck_range, blocking=True, tags=["build", "missing-file"]))
    elif pptx_error is not None:
        findings.append(_finding(finding_id="qa-pptx-open", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="build", summary=f"The compiled PPTX could not be opened for QA: {pptx_error}", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Rebuild the PPTX and confirm it opens cleanly before ship.", slide_range=deck_range, blocking=True, tags=["build", "corrupt-output"]))
    elif len(snapshots) != build_manifest.slide_count:
        findings.append(_finding(finding_id="qa-pptx-count", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="numbering", summary="The compiled PPTX slide count does not match the build manifest.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Rebuild the deck and regenerate the build manifest from the same approved inputs.", slide_range=deck_range, blocking=True, tags=["build", "numbering"]))

    if blueprint.main_story_slide_budget is not None:
        actual = blueprint.main_story_actual_slide_count or sum(1 for slide in blueprint.slides if slide.deck_mode == DeckMode.MAIN_STORY)
        budget = blueprint.main_story_slide_budget
        if actual < budget.start or actual > budget.end:
            findings.append(
                _finding(
                    finding_id="qa-main-budget",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="lecture-budget",
                    summary=f"Main-story slide count {actual} falls outside the approved lecture budget {budget.start}-{budget.end}.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Re-cluster or compress the lecture so the main story stays within the approved slide budget.",
                    slide_range=deck_range,
                    blocking=lecture_mode,
                    tags=["lecture", "budget", "main-story"],
                )
            )
    if blueprint.appendix_slide_budget is not None:
        actual = blueprint.appendix_actual_slide_count or sum(1 for slide in blueprint.slides if slide.deck_mode == DeckMode.APPENDIX)
        budget = blueprint.appendix_slide_budget
        if actual < budget.start or actual > budget.end:
            findings.append(
                _finding(
                    finding_id="qa-appendix-budget",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="lecture-budget",
                    summary=f"Appendix slide count {actual} falls outside the approved support budget {budget.start}-{budget.end}.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Trim or consolidate appendix-only support so it stays inside the approved budget band.",
                    slide_range=deck_range,
                    blocking=lecture_mode,
                    tags=["lecture", "budget", "appendix"],
                )
            )
    if any(entry.deck_mode == DeckMode.APPENDIX for entry in slide_ledger.entries) and blueprint.appendix_start is None:
        findings.append(
            _finding(
                finding_id="qa-appendix-boundary-missing",
                severity=QASeverity.CRITICAL,
                qa_layer=QALayer.DECK,
                category="appendix-boundary",
                summary="Appendix slides exist but the blueprint does not record an explicit appendix boundary.",
                remediation_skill="deck-orchestrator",
                recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP,
                recommendation="Rebuild the blueprint with an explicit appendix boundary before ship.",
                slide_range=deck_range,
                blocking=True,
                tags=["appendix", "boundary"],
            )
        )

    if lecture_mode and blueprint.lecture_family is not None:
        main_story_slides = [slide for slide in blueprint.slides if slide.deck_mode == DeckMode.MAIN_STORY]
        main_story_text = " ".join(
            " ".join(
                [
                    slide.title,
                    slide.one_line_takeaway,
                    slide.main_message,
                    " ".join(slide.concept_ids),
                ]
            ).lower()
            for slide in main_story_slides
        )
        first_five = main_story_slides[:5]
        if blueprint.lecture_family != LectureFamily.OPTIMIZATION_METHOD:
            for slide in first_five:
                text = " ".join([slide.title, slide.one_line_takeaway, slide.main_message]).lower()
                suspicious = [term for term in SUSPICIOUS_OPTIMIZATION_FRAMING if term in text]
                if suspicious:
                    findings.append(
                        _finding(
                            finding_id=f"qa-topic-drift-{slide.slide_number:03d}",
                            severity=QASeverity.MAJOR,
                            qa_layer=QALayer.DECK,
                            category="topic-drift",
                            summary=f"Slide {slide.slide_number} uses optimization-family framing that does not fit the selected lecture family: {', '.join(suspicious)}.",
                            remediation_skill="deck-orchestrator",
                            recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                            recommendation="Rebuild the early lecture arc from the concept graph so it opens on the selected teaching family rather than on generic optimization rhetoric.",
                            slide_number=slide.slide_number,
                            slide_id=f"s{slide.slide_number:03d}",
                            blocking=True,
                            tags=["lecture", "topic-drift", "topic_drift_critic"],
                        )
                    )

        pattern_sequence = [
            slide.slide_archetype.value
            if blueprint.lecture_family == LectureFamily.CONCEPT_TO_ALGORITHM_MAPPING and slide.slide_archetype is not None
            else slide.slide_intent.value
            for slide in main_story_slides
            if (
                blueprint.lecture_family == LectureFamily.CONCEPT_TO_ALGORITHM_MAPPING
                and slide.slide_archetype is not None
            )
            or slide.slide_intent is not None
        ]
        main_story_intents = [slide.slide_intent for slide in main_story_slides if slide.slide_intent is not None]
        intent_counts = Counter(intent.value for intent in main_story_intents)
        max_run = 0
        current_run = 0
        previous_pattern: str | None = None
        for pattern in pattern_sequence:
            if pattern == previous_pattern:
                current_run += 1
            else:
                previous_pattern = pattern
                current_run = 1
            max_run = max(max_run, current_run)
        pattern_counts = Counter(pattern_sequence)
        if max_run > 2 or any(count > max(4, len(pattern_sequence) // 2) for count in pattern_counts.values()):
            findings.append(
                _finding(
                    finding_id="qa-template-repetition",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="template-repetition",
                    summary="The main story overuses the same slide intent or scaffold pattern, which recreates the old repetitive lecture failure mode.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Re-synthesize the teaching plan so identical intent clusters are bounded and the lecture keeps changing pedagogical moves.",
                    slide_range=deck_range,
                    blocking=False,
                    tags=["lecture", "repetition", "template_repetition_critic"],
                )
            )

        if blueprint.lecture_family == LectureFamily.CONCEPT_TO_ALGORITHM_MAPPING:
            missing_moves: list[str] = []
            if intent_counts.get(SlideIntent.MAPPING_BRIDGE.value, 0) < 2:
                missing_moves.append("two mapping-bridge slides")
            if intent_counts.get(SlideIntent.MECHANISM_WALKTHROUGH.value, 0) < 1:
                missing_moves.append("one mechanism walkthrough")
            if intent_counts.get(SlideIntent.WORKED_EXAMPLE.value, 0) < 1:
                missing_moves.append("one worked example")
            if intent_counts.get(SlideIntent.MISCONCEPTION_PITFALL.value, 0) < 1:
                missing_moves.append("one misconception or limitation slide")
            if (
                intent_counts.get(SlideIntent.COMPARISON_TRADEOFF.value, 0)
                + intent_counts.get(SlideIntent.APPLICATION_VIGNETTE.value, 0)
            ) < 1:
                missing_moves.append("one comparison/tradeoff or application slide")
            if missing_moves:
                findings.append(
                    _finding(
                        finding_id="qa-teaching-utility",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.DECK,
                        category="teaching-utility",
                        summary=f"The lecture plan is missing required pedagogical moves: {', '.join(missing_moves)}.",
                        remediation_skill="deck-orchestrator",
                        recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                        recommendation="Rebuild the teaching plan from the concept graph so the main story teaches with bridges, mechanisms, examples, and limits instead of summarizing sections.",
                        slide_range=deck_range,
                        blocking=True,
                        tags=["lecture", "teaching_utility_critic"],
                    )
                )

        if blueprint.lecture_family != LectureFamily.OPTIMIZATION_METHOD:
            missing_concepts = [
                concept
                for concept in blueprint.central_concepts[:6]
                if _norm(concept).lower() not in main_story_text
            ]
            if missing_concepts:
                findings.append(
                    _finding(
                        finding_id="qa-concept-coverage",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.DECK,
                        category="concept-coverage",
                        summary=f"Central concepts from the concept graph are missing or weakly surfaced in the main story: {', '.join(missing_concepts[:4])}.",
                        remediation_skill="deck-orchestrator",
                        recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                        recommendation="Re-synthesize the lecture so central concept-graph nodes appear explicitly in the main-story titles, takeaways, or mapped concept ids.",
                        slide_range=deck_range,
                        blocking=False,
                        tags=["lecture", "concept_coverage_critic"],
                    )
                )

    for index, warning in enumerate(build_manifest.warnings, start=1):
        findings.append(_finding(finding_id=f"qa-build-warning-{index:03d}", severity=QASeverity.MINOR, qa_layer=QALayer.OBJECT, category="build-warning", summary=warning, remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.SAFE_TO_DEFER, recommendation="Resolve the compiler warning if it affects export safety or output completeness.", slide_range=deck_range, tags=["build", "warning"]))
    if build_manifest.slide_ratio != blueprint.slide_ratio:
        findings.append(_finding(finding_id="qa-ratio", severity=QASeverity.MINOR, qa_layer=QALayer.OBJECT, category="build", summary=f"Build ratio {build_manifest.slide_ratio} differs from blueprint ratio {blueprint.slide_ratio}.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT, recommendation="Align compile-time ratio settings with the approved blueprint.", slide_range=deck_range, tags=["build", "ratio"]))

    term_pairs = [([_norm(term).lower() for term in TERM_RE.findall(rule)][0], [_norm(term).lower() for term in TERM_RE.findall(rule)][1:]) for rule in deck_constitution.terminology_rules if len(TERM_RE.findall(rule)) >= 2]
    mapping_family = lecture_mode and blueprint.lecture_family == LectureFamily.CONCEPT_TO_ALGORITHM_MAPPING
    main_story_archetype_sequence: list[tuple[int, str]] = []
    main_story_title_stems: Counter[str] = Counter()
    main_story_title_stem_slides: dict[str, list[int]] = {}
    main_story_opening_formulas: Counter[str] = Counter()
    main_story_opening_slides: dict[str, list[int]] = {}
    main_story_bridge_shell_indices: list[int] = []
    main_story_cycle_shell_indices: list[int] = []
    main_story_visual_signature_sequence: list[tuple[int, str]] = []
    appendix_clone_signatures: list[tuple[int, str]] = []
    appendix_geometry_signatures: list[tuple[int, str]] = []
    appendix_visual_signature_sequence: list[tuple[int, str]] = []
    appendix_formulaic_title_slides: list[int] = []
    compiled_chrome_slides: list[int] = []
    compiled_duplicate_slides: list[int] = []
    compiled_realization_mismatch_slides: list[int] = []
    compiled_expected_visual_missing_slides: list[int] = []
    compiled_chrome_dominant_slides: list[int] = []
    compiled_missing_visual_center_slides: list[int] = []
    compiled_text_overflow_risk_slides: list[int] = []
    compiled_overlap_risk_slides: list[int] = []
    compiled_missing_support_marker_slides: list[int] = []
    compiled_weak_title_body_balance_slides: list[int] = []
    compiled_weak_visual_anchor_slides: list[int] = []
    compiled_text_card_like_slides: list[int] = []
    compiled_repeated_geometry_slides: list[int] = []
    compiled_appendix_visual_clone_slides: list[int] = []
    compiled_truth_mismatches: list[str] = []
    compiled_audit_slides: list[dict[str, Any]] = []
    compiled_text_slides: list[dict[str, Any]] = []
    compiled_shape_slides: list[dict[str, Any]] = []

    for entry in slide_ledger.entries:
        slide = blueprint_by_number.get(entry.slide_number)
        link = linkage_by_number.get(entry.slide_number)
        snapshot = snapshot_by_index.get(link.pptx_index) if link is not None else None
        if slide is None:
            findings.append(_finding(finding_id=f"qa-ledger-{entry.slide_number:03d}", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="ledger", summary=f"Ledger slide {entry.slide_number} does not exist in the blueprint.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Reconcile slide numbering and lineage between the blueprint and ledger.", slide_number=entry.slide_number, slide_id=entry.slide_id, blocking=True, tags=["ledger", "numbering"]))
            continue
        if link is None:
            findings.append(_finding(finding_id=f"qa-linkage-{entry.slide_number:03d}", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="build", summary=f"Slide {entry.slide_number} is missing from slide-build-linkage.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Recompile the deck and regenerate slide linkage artifacts.", slide_number=entry.slide_number, slide_id=entry.slide_id, blocking=True, tags=["build", "linkage"]))
            continue
        if entry.compile_status != StageStatus.COMPLETE or link.compile_status != StageStatus.COMPLETE or link.missing_dependencies:
            findings.append(_finding(finding_id=f"qa-compile-{entry.slide_number:03d}", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="build", summary=f"Slide {entry.slide_number} still has unresolved compile state.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Resolve compile blockers and regenerate slide-build-linkage before ship.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=True, tags=["build", "compile-status"]))
        if slide.layout_pattern_id != link.layout_pattern_id:
            findings.append(_finding(finding_id=f"qa-layout-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.OBJECT, category="design-drift", summary=f"Slide {entry.slide_number} compiled with a layout pattern different from the approved blueprint.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT, recommendation="Compile the slide using the approved layout pattern or revise the blueprint explicitly.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["layout", "design-system"]))
        if link.numbering_label != _numbering_label(entry.slide_number, entry.deck_mode.value == "appendix"):
            findings.append(_finding(finding_id=f"qa-numbering-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="numbering", summary=f"Slide {entry.slide_number} uses inconsistent numbering label `{link.numbering_label}`.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT, recommendation="Restore numbering labels so appendix and main-story slides stay visually distinct.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["numbering", "appendix"]))
        pattern = layout_index.get(slide.layout_pattern_id)
        if pattern is None:
            findings.append(_finding(finding_id=f"qa-layout-missing-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.OBJECT, category="layout-compatibility", summary=f"Slide {entry.slide_number} references layout pattern `{slide.layout_pattern_id}` but the active layout library does not define it.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Use only layout ids from the approved layout library or update the library deterministically.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=lecture_mode, tags=["layout", "compatibility"]))
        else:
            if slide.slide_role not in pattern.slide_roles:
                findings.append(_finding(finding_id=f"qa-layout-role-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="layout-compatibility", summary=f"Slide {entry.slide_number} uses layout `{pattern.pattern_id}` with unsupported role `{slide.slide_role.value}`.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Remap the slide to a layout that explicitly supports its role.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=lecture_mode, tags=["layout", "compatibility"]))
            if slide.visual_type not in pattern.supported_visual_types:
                findings.append(_finding(finding_id=f"qa-layout-visual-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="layout-compatibility", summary=f"Slide {entry.slide_number} uses layout `{pattern.pattern_id}` with unsupported visual `{slide.visual_type.value}`.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Remap the slide to a layout that explicitly supports the chosen content type.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=lecture_mode, tags=["layout", "compatibility"]))
        if entry.content_tier != slide.content_tier:
            findings.append(_finding(finding_id=f"qa-tier-drift-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="content-tier", summary=f"Slide {entry.slide_number} content tier drifted between blueprint and ledger.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Keep content tier routing stable across blueprint and ledger artifacts.", slide_number=entry.slide_number, slide_id=entry.slide_id, tags=["lecture", "appendix", "routing"]))
        if entry.deck_mode != DeckMode.APPENDIX and slide.content_tier == ContentTier.APPENDIX_ONLY:
            findings.append(_finding(finding_id=f"qa-tier-main-{entry.slide_number:03d}", severity=QASeverity.CRITICAL, qa_layer=QALayer.DECK, category="content-tier", summary=f"Slide {entry.slide_number} is appendix-only content but still appears in the main story.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Move appendix-only support behind the appendix boundary before ship.", slide_number=entry.slide_number, slide_id=entry.slide_id, blocking=True, tags=["lecture", "routing", "appendix"]))
        if entry.deck_mode == DeckMode.APPENDIX and slide.content_tier != ContentTier.APPENDIX_ONLY:
            findings.append(_finding(finding_id=f"qa-tier-appendix-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="content-tier", summary=f"Slide {entry.slide_number} is in appendix mode but not marked appendix-only.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Mark appendix slides explicitly as appendix-only support.", slide_number=entry.slide_number, slide_id=entry.slide_id, tags=["lecture", "routing", "appendix"]))
        if slide.presenter_notes and not link.notes_present:
            findings.append(_finding(finding_id=f"qa-notes-{entry.slide_number:03d}", severity=QASeverity.MINOR, qa_layer=QALayer.OBJECT, category="build", summary=f"Slide {entry.slide_number} lost approved presenter notes during compile.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT, recommendation="Restore presenter notes when re-running compile.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["notes", "build"]))
        for asset_id in link.asset_ids:
            asset = asset_by_id.get(asset_id)
            if asset is None or not _path(asset.local_path, root).is_file():
                findings.append(_finding(finding_id=f"qa-asset-{entry.slide_number:03d}-{asset_id}", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="asset", summary=f"Missing asset file for slide {entry.slide_number}: {asset.local_path if asset is not None else asset_id}.", remediation_skill="document-asset-crop" if asset is None or asset.asset_kind in {AssetKind.DOCUMENT_CROP, AssetKind.IMAGE} else "structured-visuals", recommendation_type=QARecommendationType.NEEDS_ASSET_REGENERATION, recommendation="Regenerate the missing compile-ready asset or switch to the recorded fallback route.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=True, tags=["asset", "missing-file"]))
            elif asset.asset_kind == AssetKind.DOCUMENT_CROP and asset.limitations:
                findings.append(_finding(finding_id=f"qa-crop-{entry.slide_number:03d}", severity=QASeverity.MINOR, qa_layer=QALayer.SLIDE, category="crop-adequacy", summary=f"Slide {entry.slide_number} uses a crop with recorded limitations: {', '.join(asset.limitations)}.", remediation_skill="document-asset-crop", recommendation_type=QARecommendationType.SAFE_TO_DEFER, recommendation="Review whether the crop limitations are acceptable in final context.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["crop", "readability"]))
        for spec_id in link.viz_spec_ids:
            record = viz_by_spec_id.get(spec_id)
            if record is None:
                findings.append(_finding(finding_id=f"qa-viz-{entry.slide_number:03d}-{spec_id}", severity=QASeverity.CRITICAL, qa_layer=QALayer.OBJECT, category="visual", summary=f"Slide {entry.slide_number} links structured visual `{spec_id}` but no viz-manifest record exists.", remediation_skill="structured-visuals", recommendation_type=QARecommendationType.NEEDS_ASSET_REGENERATION, recommendation="Restore the missing viz-manifest record or re-render the structured visual.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=True, tags=["visual", "manifest"]))
                continue
            if set(record.applied_color_tokens) - color_tokens or set(record.applied_typography_tokens) - typography_tokens:
                findings.append(_finding(finding_id=f"qa-viz-tokens-{entry.slide_number:03d}-{spec_id}", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="design-drift", summary=f"Slide {entry.slide_number} uses structured-visual tokens outside the approved design system.", remediation_skill="structured-visuals", recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT, recommendation="Re-render the structured visual using only approved color and typography tokens.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["design-system", "visual"]))
            if record.spec.readability.frame_fit == FrameFit.SPLIT_RECOMMENDED and record.fallback_output_path is None:
                findings.append(_finding(finding_id=f"qa-viz-fit-{entry.slide_number:03d}-{spec_id}", severity=QASeverity.MAJOR, qa_layer=QALayer.SLIDE, category="frame-fit", summary=f"Slide {entry.slide_number} uses a structured visual that recommends splitting the frame but has no fallback output.", remediation_skill="structured-visuals", recommendation_type=QARecommendationType.NEEDS_FALLBACK_ROUTE, recommendation="Generate the simpler fallback variant or move the dense variant out of the one-slide frame.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["visual", "density"]))
        if entry.asset_request_ids and any(kind.value in {"document-crop", "image"} for kind in entry.asset_dependency_kinds) and not any(request_id in approved_request_ids for request_id in entry.asset_request_ids):
            findings.append(_finding(finding_id=f"qa-asset-link-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.OBJECT, category="asset", summary=f"Slide {entry.slide_number} has source-asset requests but no approved or ready asset record.", remediation_skill="document-asset-crop", recommendation_type=QARecommendationType.NEEDS_ASSET_REGENERATION, recommendation="Resolve the asset request or accept the structured fallback before ship.", slide_number=entry.slide_number, slide_id=entry.slide_id, tags=["asset", "crop"]))
        if snapshot is None:
            continue
        forbidden_reason = _forbidden_visible_text_reason(snapshot.all_text)
        if forbidden_reason is not None:
            findings.append(_finding(finding_id=f"qa-visible-forbidden-{entry.slide_number:03d}", severity=QASeverity.CRITICAL, qa_layer=QALayer.SLIDE, category="forbidden-visible-text", summary=f"Slide {entry.slide_number} exposes forbidden internal or placeholder text ({forbidden_reason}).", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Sanitize the slide before ship and block any path that emits internal helper or fallback text.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=True, tags=["lecture", "sanitizer", "visible-text"]))
        compiled_title = snapshot.title_text or entry.title
        realized_archetype = _realized_archetype(entry, slide, link, snapshot)
        compiled_duplicate_blocks = _compiled_duplicate_blocks(snapshot)
        compiled_chrome = _compiled_chrome_hits(snapshot)
        expected_visual_missing = _compiled_expected_visual_missing(slide, link, snapshot, realized_archetype)
        realization_mismatch = _compiled_realization_mismatch(slide, realized_archetype, expected_visual_missing)
        stage_labels = [
            row[0].strip().lower()
            for row in slide.authoring_payload.get("rows", [])
            if isinstance(row, list) and row and isinstance(row[0], str)
        ]
        worked_example_progression_missing = (
            slide.slide_archetype == SlideArchetype.WORKED_EXAMPLE_STATE_TABLE
            and bool(stage_labels)
            and any(label not in snapshot.all_text.lower() for label in stage_labels[: min(len(stage_labels), 3)])
        )
        process_order_missing = slide.slide_archetype in PROCESS_ARCHETYPES and snapshot.connector_count < 1
        mapping_structure_missing = slide.slide_archetype in MAPPING_ARCHETYPES and snapshot.table_count < 1
        generic_text_collapse = (
            slide.slide_archetype in MAPPING_ARCHETYPES
            or slide.slide_archetype in PROCESS_ARCHETYPES
            or slide.slide_archetype == SlideArchetype.WORKED_EXAMPLE_STATE_TABLE
            or slide.slide_archetype == SlideArchetype.COMPARISON_MATRIX
        ) and realized_archetype == "generic-text-card"
        opening_formula = _compiled_opening_formula(snapshot)
        title_stem = _title_stem(compiled_title)
        title_template = _compiled_title_template(compiled_title)
        appendix_geometry_signature = _appendix_geometry_signature(realized_archetype, snapshot) if entry.deck_mode == DeckMode.APPENDIX else ""
        visual_signature = _visual_signature(snapshot, realized_archetype)
        text_card_like = _is_text_card_like(snapshot, realized_archetype)
        chrome_dominant = _is_chrome_dominant(snapshot)
        missing_visual_center = _missing_visual_center(snapshot, realized_archetype)
        text_overflow_risk = _text_overflow_risk(snapshot)
        overlap_risk = _overlap_collision_risk(snapshot)
        missing_support_marker = _missing_support_marker(slide, snapshot)
        weak_title_body_balance = _weak_title_body_balance(entry, snapshot)
        weak_visual_anchor = _weak_visual_anchor(slide, snapshot, realized_archetype)

        compiled_text_slides.append(
            {
                "slide_number": entry.slide_number,
                "slide_id": entry.slide_id,
                "pptx_index": link.pptx_index,
                "deck_mode": entry.deck_mode.value,
                "title_text": compiled_title,
                "text_blocks": list(snapshot.text_blocks),
                "table_blocks": list(snapshot.table_blocks),
                "all_text": snapshot.all_text,
            }
        )
        compiled_shape_slides.append(
            {
                "slide_number": entry.slide_number,
                "slide_id": entry.slide_id,
                "pptx_index": link.pptx_index,
                "deck_mode": entry.deck_mode.value,
                "layout_family": link.layout_family,
                "picture_count": snapshot.picture_count,
                "table_count": snapshot.table_count,
                "text_shape_count": snapshot.text_shape_count,
                "placeholder_count": snapshot.placeholder_count,
                "auto_shape_count": snapshot.auto_shape_count,
                "group_count": snapshot.group_count,
                "connector_count": snapshot.connector_count,
                "shape_type_counts": snapshot.shape_type_counts,
                "shape_boxes": [
                    {
                        "kind": box.kind,
                        "left": box.left,
                        "top": box.top,
                        "width": box.width,
                        "height": box.height,
                        "text": box.text,
                        "is_title": box.is_title,
                        "row_count": box.row_count,
                        "column_count": box.column_count,
                    }
                    for box in snapshot.shape_boxes
                ],
            }
        )
        compiled_audit_slides.append(
            {
                "slide_number": entry.slide_number,
                "slide_id": entry.slide_id,
                "pptx_index": link.pptx_index,
                "deck_mode": entry.deck_mode.value,
                "title_text": compiled_title,
                "title_template": title_template,
                "planned_archetype": slide.slide_archetype.value if slide.slide_archetype is not None else None,
                "realized_archetype": realized_archetype,
                "chosen_layout_family": link.layout_family,
                "chrome_hits": compiled_chrome,
                "duplicate_blocks": compiled_duplicate_blocks,
                "expected_visual_missing": expected_visual_missing,
                "archetype_realization_mismatch": realization_mismatch,
                "opening_formula": opening_formula,
                "appendix_geometry_signature": appendix_geometry_signature or None,
                "visual_signature": visual_signature,
                "text_card_like": text_card_like,
                "chrome_dominant": chrome_dominant,
                "missing_visual_center": missing_visual_center,
                "text_overflow_risk": text_overflow_risk,
                "overlap_risk": overlap_risk,
                "missing_support_marker": missing_support_marker,
                "weak_title_body_balance": weak_title_body_balance,
                "weak_visual_anchor": weak_visual_anchor,
                "picture_count": snapshot.picture_count,
                "table_count": snapshot.table_count,
                "text_shape_count": snapshot.text_shape_count,
            }
        )
        if mapping_family:
            preview_slide = authoring_preview_slides.get(entry.slide_number)
            if entry.deck_mode == DeckMode.MAIN_STORY:
                main_story_archetype_sequence.append((entry.slide_number, realized_archetype))
                main_story_visual_signature_sequence.append((entry.slide_number, visual_signature))
                if title_stem:
                    main_story_title_stems[title_stem] += 1
                    main_story_title_stem_slides.setdefault(title_stem, []).append(entry.slide_number)
                if opening_formula:
                    main_story_opening_formulas[opening_formula] += 1
                    main_story_opening_slides.setdefault(opening_formula, []).append(entry.slide_number)
                if _is_bridge_shell(slide, realized_archetype):
                    main_story_bridge_shell_indices.append(entry.slide_number)
                if _is_cycle_shell(slide, realized_archetype, compiled_title):
                    main_story_cycle_shell_indices.append(entry.slide_number)
            else:
                appendix_clone_signature = "|".join(
                    [
                        title_template or "appendix",
                        appendix_geometry_signature or realized_archetype,
                    ]
                )
                appendix_clone_signatures.append((entry.slide_number, appendix_clone_signature))
                appendix_geometry_signatures.append((entry.slide_number, appendix_geometry_signature or realized_archetype))
                appendix_visual_signature_sequence.append((entry.slide_number, visual_signature))
                if _compiled_source_map_title_pattern(compiled_title):
                    appendix_formulaic_title_slides.append(entry.slide_number)
            if text_card_like:
                compiled_text_card_like_slides.append(entry.slide_number)
            if chrome_dominant:
                compiled_chrome_dominant_slides.append(entry.slide_number)
            if missing_visual_center and realized_archetype != "title-orientation":
                compiled_missing_visual_center_slides.append(entry.slide_number)
            if text_overflow_risk:
                compiled_text_overflow_risk_slides.append(entry.slide_number)
            if overlap_risk:
                compiled_overlap_risk_slides.append(entry.slide_number)
            if missing_support_marker:
                compiled_missing_support_marker_slides.append(entry.slide_number)
            if weak_title_body_balance:
                compiled_weak_title_body_balance_slides.append(entry.slide_number)
            if weak_visual_anchor:
                compiled_weak_visual_anchor_slides.append(entry.slide_number)

            if compiled_chrome:
                compiled_chrome_slides.append(entry.slide_number)
                findings.append(
                    _finding(
                        finding_id=f"COMPILED_DECK_CHROME_PRESENT-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="repeated-chrome",
                        summary=f"Slide {entry.slide_number} exposes generic chrome in the compiled deck: {', '.join(hit['phrase'] for hit in compiled_chrome)}.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Remove roadmap, phase, and generic bridge chrome unless the compiled layout explicitly requires it.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "repeated_chrome_critic", "compiled_deck_truth"],
                    )
                )
            if chrome_dominant:
                findings.append(
                    _finding(
                        finding_id=f"qa-chrome-dominance-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="chrome-dominance",
                        summary=f"Slide {entry.slide_number} lets chrome or meta scaffolding dominate the visible composition.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Reduce chrome area so the main teaching object, not helper scaffolding, is the visual center of gravity.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "chrome_dominance_critic", "compiled_deck_truth"],
                    )
                )
            if compiled_duplicate_blocks:
                compiled_duplicate_slides.append(entry.slide_number)
                findings.append(
                    _finding(
                        finding_id=f"COMPILED_DECK_TITLE_BODY_DUPLICATION-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="title-body-duplication",
                        summary=f"Slide {entry.slide_number} repeats visible title/body language in the compiled deck: {', '.join(compiled_duplicate_blocks[:3])}.",
                        remediation_skill="deck-orchestrator",
                        recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                        recommendation="Rewrite the compiled slide so the title, body, and any support blocks do not restate the same sentence.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "title_body_duplication_critic", "compiled_deck_truth"],
                    )
                )
            if expected_visual_missing:
                compiled_expected_visual_missing_slides.append(entry.slide_number)
                findings.append(
                    _finding(
                        finding_id=f"COMPILED_DECK_VISUAL_NEEDED_BUT_MISSING-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="visual-needed-but-missing",
                        summary=f"Slide {entry.slide_number} expects mapping/process/example geometry, but the compiled slide still lacks the required table or flow structure.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Make the compiled slide realize the expected table, matrix, flow, or worked-example geometry.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "visual_needed_but_missing_critic", "compiled_deck_truth"],
                    )
                )
            if missing_visual_center and entry.deck_mode == DeckMode.MAIN_STORY and realized_archetype != "title-orientation":
                findings.append(
                    _finding(
                        finding_id=f"qa-missing-visual-center-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="missing-visual-center",
                        summary=f"Slide {entry.slide_number} has no clear visual center of gravity in the compiled layout.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Consolidate the slide around one dominant object or panel instead of scattering equal-weight text fragments.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "missing_visual_center_critic", "compiled_deck_truth"],
                    )
                )
            if text_overflow_risk:
                findings.append(
                    _finding(
                        finding_id=f"qa-text-overflow-risk-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="text-overflow-risk",
                        summary=f"Slide {entry.slide_number} packs more text into a visible box than the current frame is likely to hold cleanly.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Increase the text area or trim the visible copy before the slide reaches render validation.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["render", "overflow-risk"],
                    )
                )
            if overlap_risk:
                findings.append(
                    _finding(
                        finding_id=f"qa-overlap-risk-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="overlap-risk",
                        summary=f"Slide {entry.slide_number} contains overlapping visible content boxes that are likely to collide in the rendered deck.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Separate the colliding text, picture, or table boxes so the slide reads without visual collision.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["render", "collision-risk"],
                    )
                )
            if missing_support_marker:
                findings.append(
                    _finding(
                        finding_id=f"qa-support-marker-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="missing-support-marker",
                        summary=f"Slide {entry.slide_number} is source-backed support but the compiled deck no longer shows an explicit source or evidence marker.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Restore an explicit source or evidence marker on the visible slide before ship.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["evidence", "marker"],
                    )
                )
            if weak_title_body_balance:
                findings.append(
                    _finding(
                        finding_id=f"qa-title-body-balance-{entry.slide_number:03d}",
                        severity=QASeverity.MINOR,
                        qa_layer=QALayer.SLIDE,
                        category="weak-title-body-balance",
                        summary=f"Slide {entry.slide_number} is title-heavy relative to the visible body, which weakens hierarchy and follow-through.",
                        remediation_skill="deck-orchestrator",
                        recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                        recommendation="Tighten the title or add a stronger supporting body so the title/body hierarchy stays balanced.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["hierarchy", "title-body"],
                    )
                )
            if weak_visual_anchor and entry.deck_mode == DeckMode.MAIN_STORY:
                findings.append(
                    _finding(
                        finding_id=f"qa-weak-visual-anchor-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="weak-visual-anchor",
                        summary=f"Slide {entry.slide_number} claims a visual-led teaching move but the compiled deck still reads as loose text without one strong anchor.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Promote one dominant visual or structured object so the slide stops reading like a text scatter.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["hierarchy", "visual-anchor"],
                    )
                )
            if mapping_structure_missing:
                findings.append(
                    _finding(
                        finding_id=f"qa-mapping-structure-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="mapping-slide-missing-correspondence-structure",
                        summary=f"Slide {entry.slide_number} should realize a correspondence or mapping table, but the compiled deck still lacks table structure.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Render the mapping slide as a table or matrix with visible row/column correspondence.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "mapping_slide_missing_correspondence_structure", "compiled_deck_truth"],
                    )
                )
            if process_order_missing:
                findings.append(
                    _finding(
                        finding_id=f"qa-process-ordered-steps-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="process-flow-missing-ordered-steps",
                        summary=f"Slide {entry.slide_number} is meant to teach a process, but the compiled deck lacks explicit ordered-step connectors.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Render the process slide as ordered steps with directional flow geometry.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "process_flow_missing_ordered_steps", "compiled_deck_truth"],
                    )
                )
            if worked_example_progression_missing:
                findings.append(
                    _finding(
                        finding_id=f"qa-worked-example-state-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="worked-example-missing-state-progression",
                        summary=f"Slide {entry.slide_number} should show an initial, evaluation, and next-state progression, but the compiled deck omits that state sequence.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Render the worked example with explicit state progression before or alongside the state table.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "worked_example_missing_state_progression", "compiled_deck_truth"],
                    )
                )
            if realization_mismatch:
                compiled_realization_mismatch_slides.append(entry.slide_number)
                findings.append(
                    _finding(
                        finding_id=f"COMPILED_DECK_ARCHETYPE_REALIZATION_MISMATCH-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="archetype-mismatch",
                        summary=f"Slide {entry.slide_number} compiles as `{realized_archetype}` even though its teaching move expects a stronger realized structure.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Make the compiled slide realize its mapping, process, or worked-example archetype instead of collapsing into a generic card form.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "archetype_mismatch_critic", "compiled_deck_truth"],
                    )
                )
            if generic_text_collapse:
                findings.append(
                    _finding(
                        finding_id=f"qa-generic-text-card-{entry.slide_number:03d}",
                        severity=QASeverity.MAJOR,
                        qa_layer=QALayer.SLIDE,
                        category="archetype-realized-as-generic-text-card",
                        summary=f"Slide {entry.slide_number} collapsed into a generic text card even though the approved archetype requires a stronger visual structure.",
                        remediation_skill="pptx-compiler",
                        recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                        recommendation="Realize the approved archetype with actual table, matrix, flow, or worked-example objects.",
                        slide_number=entry.slide_number,
                        slide_id=entry.slide_id,
                        build_link_index=link.pptx_index,
                        tags=["lecture", "archetype_realized_as_generic_text_card", "compiled_deck_truth"],
                    )
                )
            if isinstance(preview_slide, dict):
                preview_chrome = preview_slide.get("chrome_blocks_used")
                if (not preview_chrome) and compiled_chrome:
                    compiled_truth_mismatches.append(
                        f"Slide {entry.slide_number}: authoring-preview recorded no chrome blocks, but the compiled deck still shows {', '.join(hit['phrase'] for hit in compiled_chrome)}."
                    )
                preview_duplicate_flags = preview_slide.get("duplicate_text_flags")
                if (not preview_duplicate_flags) and compiled_duplicate_blocks:
                    compiled_truth_mismatches.append(
                        f"Slide {entry.slide_number}: authoring-preview recorded no duplicate text flags, but compiled duplicate blocks remain."
                    )
                preview_archetype = str(preview_slide.get('slide_archetype') or '').strip().lower()
                if realization_mismatch and preview_archetype and preview_archetype != "generic-text-card":
                    compiled_truth_mismatches.append(
                        f"Slide {entry.slide_number}: authoring-preview targeted `{preview_archetype}`, but the compiled deck realizes `{realized_archetype}`."
                    )
        if lecture_mode:
            if len(slide.core_content) > 3 or any(len(_norm(item)) > 110 for item in slide.core_content):
                findings.append(_finding(finding_id=f"qa-core-copy-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.SLIDE, category="teaching-copy", summary=f"Slide {entry.slide_number} carries more on-slide teaching copy than the lecture policy allows.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Compress the slide to a title, takeaway, and at most three short teaching bullets.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=False, tags=["lecture", "copy", "density"]))
        if len(entry.title) > 72 or len(entry.one_line_takeaway) > 120:
            findings.append(_finding(finding_id=f"qa-clarity-{entry.slide_number:03d}", severity=QASeverity.MINOR, qa_layer=QALayer.SLIDE, category="clarity", summary=f"Slide {entry.slide_number} may violate three-second comprehension because the title or takeaway is too long.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Tighten the title or takeaway so the slide claim reads faster.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["clarity", "three-second"]))
        if _norm(link.title) not in snapshot.all_text:
            findings.append(_finding(finding_id=f"qa-title-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.SLIDE, category="content-drift", summary=f"Slide {entry.slide_number} title in the compiled deck does not match the approved title.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT, recommendation="Recompile the slide so the visible title matches the approved ledger and build linkage.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["title", "content"]))
        warn_limit = 700 if entry.deck_mode.value == "appendix" else 560
        fail_limit = 900 if entry.deck_mode.value == "appendix" else 760
        warn_paragraph_limit = 10
        fail_paragraph_limit = 14
        if lecture_mode and entry.deck_mode == DeckMode.MAIN_STORY:
            warn_limit -= 60
            fail_limit -= 80
            warn_paragraph_limit = 8
            fail_paragraph_limit = 12
        if entry.visual_type.value in {"text", "quote"}:
            warn_limit += 100
            fail_limit += 120
        if entry.visual_type.value == "table" or snapshot.table_count:
            warn_limit += 180
            fail_limit += 220
            warn_paragraph_limit = 14
            fail_paragraph_limit = 18
        layout_family = link.layout_family if link is not None else None
        if layout_family == "process-flow":
            warn_limit += 80
            fail_limit += 100
            warn_paragraph_limit += 6
            fail_paragraph_limit += 6
        if mapping_family and slide.slide_archetype in {
            SlideArchetype.ANCHOR_CONCEPT_CARD,
            SlideArchetype.LIMITATION_PITFALL_CALLOUT,
            SlideArchetype.APPLICATION_VIGNETTE,
            SlideArchetype.SYNTHESIS_INTEGRATION,
        }:
            warn_limit += 180
            fail_limit += 220
            warn_paragraph_limit += 8
            fail_paragraph_limit += 10
        if mapping_family and slide.slide_archetype in PROCESS_ARCHETYPES:
            warn_limit += 60
            fail_limit += 80
            warn_paragraph_limit += 2
            fail_paragraph_limit += 2
        evaluate_density = mapping_family or entry.visual_type.value in {"text", "quote"} or snapshot.table_count > 0 or layout_family == "appendix-reference"
        if evaluate_density:
            if snapshot.text_char_count > fail_limit or snapshot.paragraph_count > fail_paragraph_limit:
                findings.append(_finding(finding_id=f"qa-density-major-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.SLIDE, category="density", summary=f"Slide {entry.slide_number} is likely too dense for slide-native reading at presentation distance.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT, recommendation="Reduce on-slide text or move supporting detail into notes or appendix.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["density", "frame-fit"]))
            elif snapshot.text_char_count > warn_limit or snapshot.paragraph_count > warn_paragraph_limit:
                findings.append(_finding(finding_id=f"qa-density-warning-{entry.slide_number:03d}", severity=QASeverity.MINOR, qa_layer=QALayer.SLIDE, category="density", summary=f"Slide {entry.slide_number} is approaching a text-density threshold that can weaken quick comprehension.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.SAFE_TO_DEFER, recommendation="Consider trimming the on-slide copy if a later pass needs more visual breathing room.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, tags=["density", "three-second"]))
        if entry.visual_type.value not in {"text", "quote"} and not _allows_native_visual_without_external_asset(layout_family, entry.visual_type) and not (link.linked_paths or snapshot.picture_count or snapshot.table_count):
            findings.append(_finding(finding_id=f"qa-visual-missing-{entry.slide_number:03d}", severity=QASeverity.CRITICAL, qa_layer=QALayer.SLIDE, category="hierarchy", summary=f"Slide {entry.slide_number} expects a visual asset but none appears in the compiled deck.", remediation_skill="pptx-compiler", recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP, recommendation="Restore the missing visual placement or block ship until the fallback route is compiled cleanly.", slide_number=entry.slide_number, slide_id=entry.slide_id, build_link_index=link.pptx_index, blocking=True, tags=["visual", "compile"]))

    if blueprint.appendix_start is not None:
        for entry in slide_ledger.entries:
            if (entry.slide_number >= blueprint.appendix_start and entry.deck_mode.value != "appendix") or (entry.slide_number < blueprint.appendix_start and entry.deck_mode.value == "appendix"):
                findings.append(_finding(finding_id=f"qa-appendix-{entry.slide_number:03d}", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="section-drift", summary=f"Slide {entry.slide_number} crosses the approved appendix boundary.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Restore the appendix boundary in the blueprint and ledger before ship.", slide_number=entry.slide_number, slide_id=entry.slide_id, tags=["appendix", "section"]))
    planned_authoring_metrics = authoring_preview.get("repetition_metrics", {}) if isinstance(authoring_preview, dict) else {}
    if not isinstance(planned_authoring_metrics, dict):
        planned_authoring_metrics = {}
    compiled_chrome_block_count = 0
    compiled_repeated_title_stem_count = 0
    compiled_repeated_rhetorical_opening_count = 0
    compiled_repeated_archetype_count = 0
    compiled_repeated_bridge_shell_count = 0
    compiled_repeated_cycle_cluster_count = 0
    compiled_appendix_clone_count = 0
    compiled_repeated_source_map_title_pattern_count = 0
    compiled_repeated_appendix_geometry_pattern_count = 0
    compiled_repeated_geometry_count = 0
    compiled_repeated_geometry_rate = 0.0
    compiled_text_card_overuse_rate = 0.0
    compiled_chrome_dominance_rate = 0.0
    compiled_missing_visual_center_count = 0
    compiled_appendix_visual_clone_run_length = 0
    appendix_clone_rate = 0.0
    compiled_repetitive_motion_count = 0
    appendix_clone_run_slides: list[int] = []
    appendix_source_map_title_pattern_slides: list[int] = []
    appendix_geometry_pattern_slides: list[int] = []
    repeated_geometry_slides: list[int] = []
    repeated_motion_slides: list[int] = []
    repeated_title_stem_slides: list[int] = []
    repeated_bridge_shell_slides: list[int] = []
    repeated_cycle_cluster_slides: list[int] = []
    if mapping_family:
        compiled_chrome_block_count = sum(len(slide_payload["chrome_hits"]) for slide_payload in compiled_audit_slides)
        repeated_stems = {stem: count for stem, count in main_story_title_stems.items() if stem and count > 1}
        compiled_repeated_title_stem_count = len(repeated_stems)
        repeated_title_stem_slides = sorted(
            {
                slide_number
                for stem, slide_numbers in main_story_title_stem_slides.items()
                if repeated_stems.get(stem)
                for slide_number in slide_numbers
            }
        )
        max_run = 0
        current_run = 0
        previous_archetype: str | None = None
        run_start_slide: int | None = None
        repeated_archetype_run_slides: list[int] = []
        for slide_number, archetype in main_story_archetype_sequence:
            if archetype == previous_archetype:
                current_run += 1
            else:
                if previous_archetype is not None and current_run > 2 and run_start_slide is not None:
                    overflow = current_run - 2
                    compiled_repeated_archetype_count += overflow
                    repeated_archetype_run_slides.extend(
                        range(run_start_slide + 2, run_start_slide + current_run)
                    )
                current_run = 1
                previous_archetype = archetype
                run_start_slide = slide_number
            max_run = max(max_run, current_run)
        if previous_archetype is not None and current_run > 2 and run_start_slide is not None:
            overflow = current_run - 2
            compiled_repeated_archetype_count += overflow
            repeated_archetype_run_slides.extend(range(run_start_slide + 2, run_start_slide + current_run))
        repeated_formulas = {
            formula: count for formula, count in main_story_opening_formulas.items() if formula and count > 2
        }
        compiled_repeated_rhetorical_opening_count = sum(count - 2 for count in repeated_formulas.values())
        repeated_opening_slides = sorted(
            {
                slide_number
                for formula, slide_numbers in main_story_opening_slides.items()
                if repeated_formulas.get(formula)
                for slide_number in slide_numbers[2:]
            }
        )
        compiled_repeated_bridge_shell_count, repeated_bridge_shell_slides = _cluster_overflow(
            sorted(main_story_bridge_shell_indices),
            allowed=4,
        )
        compiled_repeated_cycle_cluster_count, repeated_cycle_cluster_slides = _cluster_overflow(
            sorted(main_story_cycle_shell_indices),
            allowed=4,
        )
        compiled_repeated_geometry_count, repeated_geometry_slides, geometry_max_run = _sequence_overflow(
            main_story_visual_signature_sequence,
            allowed=2,
        )
        main_story_slide_numbers = {slide_number for slide_number, _ in main_story_visual_signature_sequence}
        compiled_repeated_geometry_slides = sorted(set(repeated_geometry_slides))
        compiled_repeated_geometry_rate = round(
            compiled_repeated_geometry_count / len(main_story_visual_signature_sequence),
            3,
        ) if main_story_visual_signature_sequence else 0.0
        compiled_text_card_overuse_rate = round(
            len({slide_number for slide_number in compiled_text_card_like_slides if slide_number in main_story_slide_numbers})
            / max(len(main_story_visual_signature_sequence), 1),
            3,
        )
        compiled_chrome_dominance_rate = round(
            len(set(compiled_chrome_dominant_slides)) / max(len(compiled_audit_slides), 1),
            3,
        )
        compiled_missing_visual_center_count = len(
            {
                slide_number
                for slide_number in compiled_missing_visual_center_slides
                if slide_number in main_story_slide_numbers
            }
        )
        repeated_motion_slides = sorted(
            set(
                repeated_archetype_run_slides
                + repeated_title_stem_slides
                + repeated_opening_slides
                + repeated_bridge_shell_slides
                + repeated_cycle_cluster_slides
                + repeated_geometry_slides
            )
        )
        compiled_repetitive_motion_count = (
            compiled_repeated_archetype_count
            + compiled_repeated_title_stem_count
            + compiled_repeated_rhetorical_opening_count
            + compiled_repeated_bridge_shell_count
            + compiled_repeated_cycle_cluster_count
            + compiled_repeated_geometry_count
        )
        previous_signature: str | None = None
        appendix_run_length = 0
        appendix_run_numbers: list[int] = []
        for slide_number, signature in appendix_clone_signatures:
            if signature == previous_signature:
                appendix_run_length += 1
                appendix_run_numbers.append(slide_number)
            else:
                if appendix_run_length > 2:
                    compiled_appendix_clone_count += appendix_run_length - 2
                    appendix_clone_run_slides.extend(appendix_run_numbers)
                previous_signature = signature
                appendix_run_length = 1
                appendix_run_numbers = [slide_number]
        if appendix_run_length > 2:
            compiled_appendix_clone_count += appendix_run_length - 2
            appendix_clone_run_slides.extend(appendix_run_numbers)
        appendix_visual_overflow, compiled_appendix_visual_clone_slides, compiled_appendix_visual_clone_run_length = _sequence_overflow(
            appendix_visual_signature_sequence,
            allowed=2,
        )
        if appendix_visual_overflow > 0:
            compiled_appendix_clone_count += appendix_visual_overflow
            appendix_clone_run_slides.extend(compiled_appendix_visual_clone_slides)
        if len(appendix_formulaic_title_slides) > 2:
            compiled_repeated_source_map_title_pattern_count = len(appendix_formulaic_title_slides) - 2
            appendix_source_map_title_pattern_slides = appendix_formulaic_title_slides[2:]
        geometry_counts = Counter(signature for _, signature in appendix_geometry_signatures)
        for signature, count in geometry_counts.items():
            if count <= 5:
                continue
            compiled_repeated_appendix_geometry_pattern_count += count - 5
            appendix_geometry_pattern_slides.extend(
                slide_number for slide_number, candidate_signature in appendix_geometry_signatures if candidate_signature == signature
            )
        appendix_slide_total = len(appendix_geometry_signatures)
        appendix_clone_rate = round(
            (
                len(
                    set(
                        appendix_clone_run_slides
                        + appendix_source_map_title_pattern_slides
                        + appendix_geometry_pattern_slides
                        + compiled_appendix_visual_clone_slides
                    )
                )
                / appendix_slide_total
            ),
            3,
        ) if appendix_slide_total else 0.0
        if compiled_repeated_title_stem_count > 0:
            findings.append(
                _finding(
                    finding_id="REPEATED_TITLE_STEM",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="repeated-title-stem",
                    summary=f"The compiled main story still repeats title stems across slides {', '.join(str(number) for number in repeated_title_stem_slides[:6])}.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Rewrite or merge the repeated titled shells so later slides add a different teaching move instead of a renamed restatement.",
                    slide_range=deck_range,
                    tags=["lecture", "repeated_title_stem_critic", "compiled_deck_truth"],
                )
            )
        if compiled_repeated_bridge_shell_count > 0:
            findings.append(
                _finding(
                    finding_id="REPEATED_BRIDGE_SHELL",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="bridge-shell-repetition",
                    summary=f"The compiled main story still contains an overlong bridge-shell cluster on slides {', '.join(str(number) for number in repeated_bridge_shell_slides[:6])}.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Merge or remove the weaker bridge cards so the lecture stops revisiting the same mapping shell with minor noun substitutions.",
                    slide_range=deck_range,
                    tags=["lecture", "repeated_bridge_shell_critic", "compiled_deck_truth"],
                )
            )
        if compiled_repeated_cycle_cluster_count > 0:
            findings.append(
                _finding(
                    finding_id="REPEATED_CYCLE_CLUSTER",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="cycle-cluster-repetition",
                    summary=f"The compiled main story still carries too many cycle/operator shells in one cluster on slides {', '.join(str(number) for number in repeated_cycle_cluster_slides[:6])}.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Compress repeated cycle/operator material into a smaller set of overview, matrix, example, and tradeoff slides.",
                    slide_range=deck_range,
                    tags=["lecture", "repeated_cycle_cluster_critic", "compiled_deck_truth"],
                )
            )
        if compiled_repetitive_motion_count > 0:
            findings.append(
                _finding(
                    finding_id="COMPILED_DECK_REPETITIVE_MOTION",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="deck-motion-repetition",
                    summary=(
                        f"The compiled main story repeats the same motion too often (max same-archetype run {max_run}; "
                        f"repeated title stems: {', '.join(sorted(repeated_stems)[:4]) or 'none'}; "
                        f"repeated openings: {', '.join(sorted(repeated_formulas)[:3]) or 'none'}; "
                        f"bridge-shell overflow: {compiled_repeated_bridge_shell_count}; cycle-cluster overflow: {compiled_repeated_cycle_cluster_count})."
                    ),
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Rewrite or merge repeated bridge and cycle shells so the compiled lecture changes teaching motion instead of repeating one card rhythm.",
                    slide_range=deck_range,
                    tags=["lecture", "deck_motion_repetition_critic", "compiled_deck_truth"],
                )
            )
        if compiled_repeated_geometry_count > 0:
            findings.append(
                _finding(
                    finding_id="VISUAL_MONOTONY",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="visual-monotony",
                    summary=(
                        f"The compiled main story still repeats the same realized geometry too often "
                        f"(geometry overflow {compiled_repeated_geometry_count}, max visual-signature run {geometry_max_run})."
                    ),
                    remediation_skill="pptx-compiler",
                    recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                    recommendation="Break long runs of the same realized geometry so mapping, process, example, and synthesis slides read as distinct objects in the compiled deck.",
                    slide_range=deck_range,
                    tags=["lecture", "visual_monotony_critic", "compiled_deck_truth"],
                )
            )
        if compiled_text_card_overuse_rate > 0.4:
            findings.append(
                _finding(
                    finding_id="TEXT_CARD_OVERUSE",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="text-card-overuse",
                    summary=f"The compiled main story still relies on text-card composition too often (rate {compiled_text_card_overuse_rate:.3f}).",
                    remediation_skill="pptx-compiler",
                    recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                    recommendation="Shift more main-story slides into tables, matrices, flows, or worked examples so the lecture is not carried by repeated text-card scaffolding.",
                    slide_range=deck_range,
                    tags=["lecture", "text_card_overuse_critic", "compiled_deck_truth"],
                )
            )
        if compiled_chrome_dominant_slides:
            findings.append(
                _finding(
                    finding_id="CHROME_DOMINANCE",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="chrome-dominance",
                    summary=f"Chrome or meta scaffolding dominates the compiled layout on slides {', '.join(str(number) for number in sorted(set(compiled_chrome_dominant_slides))[:6])}.",
                    remediation_skill="pptx-compiler",
                    recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                    recommendation="Suppress or shrink chrome blocks until the main content object regains dominance.",
                    slide_range=deck_range,
                    tags=["lecture", "chrome_dominance_critic", "compiled_deck_truth"],
                )
            )
        if compiled_appendix_clone_count > 0:
            findings.append(
                _finding(
                    finding_id="COMPILED_DECK_APPENDIX_CLONE_RUN",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="appendix-clone",
                    summary=f"The compiled appendix still contains a clone run of near-identical evidence cards on slides {', '.join(str(number) for number in appendix_clone_run_slides[:6])}.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Cluster appendix evidence by theme so source traceability survives without repeating one near-identical card per topic.",
                    slide_range=deck_range,
                    tags=["lecture", "appendix_clone_critic", "compiled_deck_truth"],
                )
            )
        if compiled_appendix_visual_clone_run_length > 2:
            findings.append(
                _finding(
                    finding_id="APPENDIX_VISUAL_CLONE",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="appendix-visual-clone",
                    summary=f"The compiled appendix still runs the same realized visual geometry for {compiled_appendix_visual_clone_run_length} slides in sequence.",
                    remediation_skill="pptx-compiler",
                    recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                    recommendation="Alternate themed clusters, matrices, excerpt layouts, and comparison evidence so appendix support does not read as a clone factory.",
                    slide_range=deck_range,
                    tags=["lecture", "appendix_visual_clone_critic", "compiled_deck_truth"],
                )
            )
        if compiled_repeated_source_map_title_pattern_count > 0:
            findings.append(
                _finding(
                    finding_id="REPEATED_SOURCE_MAP_TITLE_PATTERN",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="repeated-source-map-title-pattern",
                    summary=f"The compiled appendix still repeats source-map title shells on slides {', '.join(str(number) for number in appendix_source_map_title_pattern_slides[:6])}.",
                    remediation_skill="deck-orchestrator",
                    recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE,
                    recommendation="Use theme-led appendix titles and reserve explicit source-map wording for the few slides that truly need cross-reference treatment.",
                    slide_range=deck_range,
                    tags=["lecture", "repeated_source_map_title_pattern", "compiled_deck_truth"],
                )
            )
        if compiled_repeated_appendix_geometry_pattern_count > 0:
            findings.append(
                _finding(
                    finding_id="REPEATED_APPENDIX_GEOMETRY_PATTERN",
                    severity=QASeverity.MAJOR,
                    qa_layer=QALayer.DECK,
                    category="repeated-appendix-geometry-pattern",
                    summary=f"The compiled appendix still overuses one appendix geometry pattern on slides {', '.join(str(number) for number in sorted(set(appendix_geometry_pattern_slides))[:6])}.",
                    remediation_skill="pptx-compiler",
                    recommendation_type=QARecommendationType.NEEDS_LAYOUT_ADJUSTMENT,
                    recommendation="Distribute appendix support across themed clusters, source matrices, excerpt clusters, and comparison layouts instead of one repeated evidence-card geometry.",
                    slide_range=deck_range,
                    tags=["lecture", "repeated_appendix_geometry_pattern", "compiled_deck_truth"],
                )
            )
        if isinstance(planned_authoring_metrics.get("chrome_block_count"), int) and planned_authoring_metrics["chrome_block_count"] == 0 and compiled_chrome_block_count > 0:
            compiled_truth_mismatches.append(
                f"authoring-preview chrome_block_count=0, but compiled deck chrome_block_count={compiled_chrome_block_count}."
            )
        if isinstance(planned_authoring_metrics.get("main_story_repeated_title_stem_count"), int) and planned_authoring_metrics["main_story_repeated_title_stem_count"] == 0 and compiled_repeated_title_stem_count > 0:
            compiled_truth_mismatches.append(
                f"authoring-preview main_story_repeated_title_stem_count=0, but compiled deck repeated_title_stem_count={compiled_repeated_title_stem_count}."
            )
        if isinstance(planned_authoring_metrics.get("repeated_title_stem_count"), int) and planned_authoring_metrics["repeated_title_stem_count"] == 0 and compiled_repeated_title_stem_count > 0:
            compiled_truth_mismatches.append(
                f"authoring-preview repeated_title_stem_count=0, but compiled deck repeated_title_stem_count={compiled_repeated_title_stem_count}."
            )
        if isinstance(planned_authoring_metrics.get("main_story_bridge_shell_count"), int) and planned_authoring_metrics["main_story_bridge_shell_count"] == 0 and compiled_repeated_bridge_shell_count > 0:
            compiled_truth_mismatches.append(
                f"authoring-preview main_story_bridge_shell_count=0, but compiled deck repeated_bridge_shell_count={compiled_repeated_bridge_shell_count}."
            )
        if isinstance(planned_authoring_metrics.get("main_story_cycle_cluster_count"), int) and planned_authoring_metrics["main_story_cycle_cluster_count"] == 0 and compiled_repeated_cycle_cluster_count > 0:
            compiled_truth_mismatches.append(
                f"authoring-preview main_story_cycle_cluster_count=0, but compiled deck repeated_cycle_cluster_count={compiled_repeated_cycle_cluster_count}."
            )
        if compiled_truth_mismatches:
            findings.append(
                _finding(
                    finding_id="COMPILED_DECK_TRUTH_MISMATCH",
                    severity=QASeverity.CRITICAL,
                    qa_layer=QALayer.DECK,
                    category="compiled-deck-truth-mismatch",
                    summary=f"Planned authoring claims disagree with the compiled deck on {len(compiled_truth_mismatches)} item(s).",
                    remediation_skill="deck-qa",
                    recommendation_type=QARecommendationType.FIX_NOW_BEFORE_SHIP,
                    recommendation="Recompute reporting from the compiled PPTX so a clean report cannot coexist with a visibly repetitive deck.",
                    slide_range=deck_range,
                    blocking=True,
                    tags=["lecture", "compiled_deck_truth"],
                )
            )
    if len({entry.title.endswith(PUNCTUATION_ENDINGS) for entry in slide_ledger.entries}) > 1:
        findings.append(_finding(finding_id="qa-title-style", severity=QASeverity.MINOR, qa_layer=QALayer.DECK, category="title-style", summary="Compiled slide titles mix terminal punctuation styles, which weakens title continuity.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.SAFE_TO_DEFER, recommendation="Normalize title punctuation so title style stays consistent across the deck.", slide_range=deck_range, tags=["titles", "continuity"]))
    for preferred, discouraged_terms in term_pairs:
        for entry in slide_ledger.entries:
            text = " ".join([entry.title, entry.one_line_takeaway, entry.main_message]).lower()
            used = [term for term in discouraged_terms if _norm(term) and _norm(term) in text]
            if used and preferred not in text:
                findings.append(_finding(finding_id=f"qa-term-{entry.slide_number:03d}-{preferred.replace(' ', '-')}", severity=QASeverity.MAJOR if entry.deck_mode.value != "appendix" else QASeverity.MINOR, qa_layer=QALayer.DECK, category="terminology", summary=f"Slide {entry.slide_number} uses discouraged term(s) {', '.join(used)} instead of the approved term `{preferred}`.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Normalize deck terminology so the approved term stays stable across titles, body copy, and notes.", slide_number=entry.slide_number, slide_id=entry.slide_id, tags=["continuity", "terminology"]))

    divider_numbers = [entry.slide_number for entry in slide_ledger.entries if entry.slide_role == entry.slide_role.SECTION_DIVIDER and entry.deck_mode == DeckMode.MAIN_STORY]
    if lecture_mode:
        allowed_dividers = max(2, len([section for section in blueprint.story_architecture if section.deck_mode == DeckMode.MAIN_STORY]) // 2 + 1)
        if len(divider_numbers) > allowed_dividers:
            findings.append(_finding(finding_id="qa-divider-rhythm", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="lecture-rhythm", summary=f"The main story uses {len(divider_numbers)} divider slides, which is too many for the current lecture structure.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Reduce repeated opener rhythm slides and keep dividers only for major lecture boundaries.", slide_range=deck_range, blocking=False, tags=["lecture", "rhythm", "divider"]))
        if any((current - previous) <= 2 for previous, current in zip(divider_numbers, divider_numbers[1:])):
            findings.append(_finding(finding_id="qa-divider-spacing", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="lecture-rhythm", summary="Section divider slides appear too close together to function as meaningful lecture pacing resets.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Remove or merge closely spaced divider slides so the lecture pacing remains intentional.", slide_range=deck_range, blocking=False, tags=["lecture", "rhythm", "divider"]))

        main_story_entries = [entry for entry in slide_ledger.entries if entry.deck_mode == DeckMode.MAIN_STORY]
        supporting_examples = [entry for entry in main_story_entries if entry.content_tier == ContentTier.SUPPORTING_EXAMPLE]
        if main_story_entries and len(supporting_examples) > max(6, len(main_story_entries) // 3):
            findings.append(_finding(finding_id="qa-support-ratio", severity=QASeverity.MAJOR, qa_layer=QALayer.DECK, category="appendix-routing", summary="Too much supporting-example content remains in the main lecture sequence.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Move backup examples and derivation-heavy support behind the appendix boundary.", slide_range=deck_range, blocking=False, tags=["lecture", "routing", "appendix"]))

        section_counts = Counter(entry.section for entry in main_story_entries if entry.slide_role != entry.slide_role.SECTION_DIVIDER)
        for section_name, count in section_counts.items():
            if count > 8 and section_name not in {"Orientation", "Closing"}:
                findings.append(_finding(finding_id=f"qa-section-pacing-{_norm(section_name).lower().replace(' ', '-')}", severity=QASeverity.MINOR, qa_layer=QALayer.DECK, category="section-pacing", summary=f"Section `{section_name}` runs for {count} main-story slides, which is likely too long for one lecture module.", remediation_skill="deck-orchestrator", recommendation_type=QARecommendationType.NEEDS_UPSTREAM_CONTENT_CHANGE, recommendation="Re-cluster the section into smaller concept modules or move the overflow to appendix support.", slide_range=deck_range, tags=["lecture", "pacing", "section"]))

    shape_type_totals: Counter[str] = Counter()
    total_text_shapes = 0
    total_placeholders = 0
    total_auto_shapes = 0
    total_tables = 0
    total_pictures = 0
    total_groups = 0
    total_connectors = 0
    for slide_payload in compiled_shape_slides:
        shape_type_totals.update(slide_payload.get("shape_type_counts", {}))
        total_text_shapes += int(slide_payload.get("text_shape_count", 0) or 0)
        total_placeholders += int(slide_payload.get("placeholder_count", 0) or 0)
        total_auto_shapes += int(slide_payload.get("auto_shape_count", 0) or 0)
        total_tables += int(slide_payload.get("table_count", 0) or 0)
        total_pictures += int(slide_payload.get("picture_count", 0) or 0)
        total_groups += int(slide_payload.get("group_count", 0) or 0)
        total_connectors += int(slide_payload.get("connector_count", 0) or 0)

    compiled_issue_slides = {
        "chrome_slides": sorted(set(compiled_chrome_slides)),
        "chrome_dominance_slides": sorted(set(compiled_chrome_dominant_slides)),
        "title_body_duplication_slides": sorted(set(compiled_duplicate_slides)),
        "archetype_realization_mismatch_slides": sorted(set(compiled_realization_mismatch_slides)),
        "expected_visual_missing_slides": sorted(set(compiled_expected_visual_missing_slides)),
        "missing_visual_center_slides": sorted(set(compiled_missing_visual_center_slides)),
        "text_overflow_risk_slides": sorted(set(compiled_text_overflow_risk_slides)),
        "overlap_risk_slides": sorted(set(compiled_overlap_risk_slides)),
        "missing_support_marker_slides": sorted(set(compiled_missing_support_marker_slides)),
        "weak_title_body_balance_slides": sorted(set(compiled_weak_title_body_balance_slides)),
        "weak_visual_anchor_slides": sorted(set(compiled_weak_visual_anchor_slides)),
        "text_card_like_slides": sorted(set(compiled_text_card_like_slides)),
        "appendix_clone_slides": sorted(set(appendix_clone_run_slides)),
        "appendix_visual_clone_slides": sorted(set(compiled_appendix_visual_clone_slides)),
        "source_map_title_pattern_slides": sorted(set(appendix_source_map_title_pattern_slides)),
        "appendix_geometry_pattern_slides": sorted(set(appendix_geometry_pattern_slides)),
        "repeated_title_stem_slides": sorted(set(repeated_title_stem_slides)),
        "bridge_shell_slides": sorted(set(repeated_bridge_shell_slides)),
        "cycle_cluster_slides": sorted(set(repeated_cycle_cluster_slides)),
        "repeated_geometry_slides": sorted(set(compiled_repeated_geometry_slides)),
        "repetitive_motion_slides": sorted(set(repeated_motion_slides)),
    }
    compiled_metrics = {
        "slide_count": len(compiled_audit_slides),
        "chrome_block_count": compiled_chrome_block_count,
        "chrome_dominance_rate": compiled_chrome_dominance_rate,
        "repeated_title_stem_count": compiled_repeated_title_stem_count,
        "repeated_rhetorical_opening_count": compiled_repeated_rhetorical_opening_count,
        "title_body_duplication_count": len(set(compiled_duplicate_slides)),
        "repeated_archetype_count": compiled_repeated_archetype_count,
        "repeated_geometry_count": compiled_repeated_geometry_count,
        "repeated_geometry_rate": compiled_repeated_geometry_rate,
        "repeated_bridge_shell_count": compiled_repeated_bridge_shell_count,
        "repeated_cycle_cluster_count": compiled_repeated_cycle_cluster_count,
        "text_card_overuse_rate": compiled_text_card_overuse_rate,
        "appendix_clone_count": compiled_appendix_clone_count,
        "appendix_clone_rate": appendix_clone_rate,
        "appendix_visual_clone_run_length": compiled_appendix_visual_clone_run_length,
        "repeated_source_map_title_pattern_count": compiled_repeated_source_map_title_pattern_count,
        "repeated_appendix_geometry_pattern_count": compiled_repeated_appendix_geometry_pattern_count,
        "expected_visual_missing_count": len(set(compiled_expected_visual_missing_slides)),
        "mapping_process_example_missing_structure_count": len(set(compiled_expected_visual_missing_slides)),
        "archetype_realization_mismatch_count": len(set(compiled_realization_mismatch_slides)),
        "missing_visual_center_count": compiled_missing_visual_center_count,
        "text_overflow_risk_count": len(set(compiled_text_overflow_risk_slides)),
        "overlap_risk_count": len(set(compiled_overlap_risk_slides)),
        "missing_support_marker_count": len(set(compiled_missing_support_marker_slides)),
        "weak_title_body_balance_count": len(set(compiled_weak_title_body_balance_slides)),
        "weak_visual_anchor_count": len(set(compiled_weak_visual_anchor_slides)),
        "repetitive_motion_count": compiled_repetitive_motion_count,
        "truth_mismatch_count": len(compiled_truth_mismatches),
    }
    compiled_deck_text = {
        "schema_name": "compiled_deck_text",
        "schema_version": "1.0",
        "deck_title": build_manifest.deck_title,
        "pptx_path": str(pptx_path),
        "slides": compiled_text_slides,
    }
    compiled_deck_shape_census = {
        "schema_name": "compiled_deck_shape_census",
        "schema_version": "1.0",
        "deck_title": build_manifest.deck_title,
        "pptx_path": str(pptx_path),
        "aggregate": {
            "slide_count": len(compiled_shape_slides),
            "text_shape_count": total_text_shapes,
            "placeholder_count": total_placeholders,
            "auto_shape_count": total_auto_shapes,
            "table_count": total_tables,
            "picture_count": total_pictures,
            "group_count": total_groups,
            "connector_count": total_connectors,
            "shape_type_counts": dict(sorted(shape_type_totals.items())),
        },
        "slides": compiled_shape_slides,
    }
    compiled_deck_authoring_audit = {
        "schema_name": "compiled_deck_authoring_audit",
        "schema_version": "1.0",
        "deck_title": build_manifest.deck_title,
        "pptx_path": str(pptx_path),
        "planned_metrics": dict(sorted(planned_authoring_metrics.items())),
        "compiled_metrics": compiled_metrics,
        "issue_slides": compiled_issue_slides,
        "truth_mismatches": compiled_truth_mismatches,
        "slides": compiled_audit_slides,
    }

    findings_by_slide: dict[int, list[QAFinding]] = {}
    for finding in findings:
        if finding.slide_number is not None:
            findings_by_slide.setdefault(finding.slide_number, []).append(finding)
    updated_entries = []
    slide_results: list[QASlideResult] = []
    for entry in slide_ledger.entries:
        entry_findings = findings_by_slide.get(entry.slide_number, [])
        status, warning_count, blocking_count = _slide_status(entry_findings)
        note = f"QA {status.value}: {len(entry_findings)} findings, {blocking_count} blocking."
        approved_slide = blueprint_by_number.get(entry.slide_number)
        updated_entry_payload = entry.model_dump(mode="python")
        updated_entry_payload.update(
            {
                "content_tier": _normalized_return_content_tier(entry, approved_slide),
                "qa_status": status,
                "change_note": _norm(f"{entry.change_note or ''} {note}"),
                "unresolved_blockers": _dedupe((entry.unresolved_blockers or []) + [finding.summary for finding in entry_findings if finding.blocking]) or None,
            }
        )
        updated_entries.append(type(entry).model_validate(updated_entry_payload))
        link = linkage_by_number.get(entry.slide_number)
        slide_results.append(QASlideResult(slide_number=entry.slide_number, slide_id=entry.slide_id, qa_status=status, layout_pattern_id=entry.layout_pattern_id, compile_status=entry.compile_status, build_link_index=link.pptx_index if link is not None else None, finding_ids=[finding.finding_id for finding in entry_findings], warning_count=warning_count, blocking_count=blocking_count, tags=_dedupe([tag for finding in entry_findings for tag in finding.tags])))
    updated_links = []
    for link in slide_build_linkage.slides:
        link_findings = findings_by_slide.get(link.slide_number, [])
        status, warning_count, blocking_count = _slide_status(link_findings)
        updated_links.append(link.model_copy(update={"qa_status": status, "qa_warning_count": warning_count, "qa_blocking_count": blocking_count, "qa_finding_ids": [finding.finding_id for finding in link_findings], "qa_notes": [finding.summary for finding in link_findings[:3]]}))

    round_base = state_capsule.qa_round if state_capsule is not None else prior_report.bounded_round if prior_report is not None else 0
    max_rounds = state_capsule.max_qa_rounds if state_capsule is not None else prior_report.max_rounds if prior_report is not None else 2
    bounded_round = min(round_base + 1, max_rounds)
    report_status = _qa_status(findings)
    summary = QASummary(
        slide_count=build_manifest.slide_count,
        finding_count=len(findings),
        blocking_count=sum(1 for finding in findings if finding.blocking),
        severity_counts=dict(Counter(finding.severity.value for finding in findings)),
        layer_counts=dict(Counter(finding.qa_layer.value for finding in findings)),
        recommendation_counts=dict(Counter(finding.recommendation_type.value for finding in findings)),
        pass_slide_count=sum(1 for result in slide_results if result.qa_status == QAStatus.PASS),
        conditional_slide_count=sum(1 for result in slide_results if result.qa_status == QAStatus.CONDITIONAL_PASS),
        fail_slide_count=sum(1 for result in slide_results if result.qa_status == QAStatus.FAIL),
    )
    qa_report = QAReport(
        report_id=f"qa-{build_manifest.deck_title.lower().replace(' ', '-')}",
        deck_title=build_manifest.deck_title,
        qa_status=report_status,
        audited_scope=f"Compiled deck audit for {build_manifest.slide_count} slides",
        findings=findings,
        summary=summary,
        slide_results=slide_results,
        recommended_actions=_dedupe([finding.recommendation for finding in findings]),
        bounded_round=bounded_round,
        max_rounds=max_rounds,
        drift_checks=[
            "Object-level build and asset integrity",
            "Slide-level readability and frame-fit",
            "Terminology, title-style, and appendix continuity",
            "Layout-library and design-token adherence",
            "Lecture budget, appendix routing, and divider pacing",
            "topic_drift_critic for lecture-family mismatch in main-story framing",
            "template_repetition_critic for repeated pedagogical scaffold patterns",
            "teaching_utility_critic for missing bridges, examples, mechanisms, and limits",
            "concept_coverage_critic for missing central concept-graph nodes",
            "title_body_duplication_critic for redundant title/body/takeaway shells",
            "repeated_chrome_critic for roadmap, why-now, why-it-matters, and phase chrome",
            "archetype_mismatch_critic for generic cards where mapping/flow/example geometry is required",
            "visual_needed_but_missing_critic for mapping, process, and worked-example slides without the expected structure",
            "archetype_realized_as_generic_text_card for archetypes that collapse into plain text panels",
            "process_flow_missing_ordered_steps for process slides without directional step geometry",
            "mapping_slide_missing_correspondence_structure for mapping slides missing row/column correspondence",
            "worked_example_missing_state_progression for examples that omit state transitions",
            "visual_monotony_critic for repeated realized geometry in the compiled deck",
            "chrome_dominance_critic for slides where helper chrome outweighs the main teaching object",
            "missing_visual_center_critic for layouts without one clear dominant object",
            "text_card_overuse_critic for main-story decks that lean too heavily on text-card composition",
            "appendix_visual_clone_critic for appendix runs that still look like near-identical support cards",
            "appendix_clone_critic for formulaic appendix evidence runs",
            "repeated_source_map_title_pattern for appendix title-shell repetition",
            "repeated_appendix_geometry_pattern for appendix geometry overuse in the compiled deck",
            "deck_motion_repetition_critic for long same-archetype runs and repeated title stems",
            "repeated_bridge_shell_critic for overlong clusters of mapping/bridge shells in the compiled deck",
            "repeated_cycle_cluster_critic for overlong clusters of cycle/operator shells in the compiled deck",
            "repeated_title_stem_critic for repeated compiled title stems across the main story",
            "compiled deck truth audit for realized text, shape geometry, and report-vs-PPTX mismatches",
            "Forbidden visible helper/debug text sanitation",
        ],
        checked_artifacts=checked_artifacts
        or [
            "state/blueprint.json",
            "state/design-system.json",
            "state/deck-constitution.json",
            "state/layout-library.json",
            "state/slide-ledger.json",
            "state/asset-manifest.json",
            "state/viz-manifest.json",
            "state/qa-governance.json",
            build_manifest.pptx_path,
            build_manifest.linkage_path,
        ],
        stop_condition_reached=report_status == QAStatus.FAIL and bounded_round >= max_rounds,
    )
    governance_report = refresh_qa_governance(
        qa_report,
        prior_report=prior_report,
        persisted_governance=qa_governance,
    )
    governance_by_finding_id = {
        finding_status.finding_id: finding_status
        for finding_status in governance_report.finding_statuses
        if finding_status.current_report_present
    }
    governed_findings: list[QAFinding] = []
    for finding in qa_report.findings:
        governed_status = governance_by_finding_id.get(finding.finding_id)
        finding_status = FindingStatus.OPEN
        if governed_status is not None:
            if governed_status.disposition == QAFindingGovernanceDisposition.WAIVED:
                finding_status = FindingStatus.WAIVED
            elif governed_status.disposition == QAFindingGovernanceDisposition.ACCEPTED_RISK:
                finding_status = FindingStatus.ACCEPTED
        governed_findings.append(finding.model_copy(update={"status": finding_status}))
    render_check_failure_codes: list[str] = []
    if not pptx_path.is_file():
        render_check_failure_codes.append("compiled-pptx-missing")
    elif pptx_error is not None:
        render_check_failure_codes.append("compiled-pptx-open-failed")
    elif len(snapshots) != build_manifest.slide_count:
        render_check_failure_codes.append("compiled-pptx-slide-count-mismatch")
    render_checks_present = not render_check_failure_codes and bool(snapshots) and len(snapshots) == build_manifest.slide_count
    continuity_alerts, continuity_guidance_lines = _resolve_continuity_policy_inputs(state_capsule)
    verdict_summary = summarize_qa_verdict(
        qa_report=qa_report,
        build_manifest=build_manifest,
        render_checks_present=render_checks_present,
        render_check_failure_codes=render_check_failure_codes,
        continuity_alerts=continuity_alerts,
        continuity_guidance_lines=continuity_guidance_lines,
    )
    qa_report = qa_report.model_copy(
        update={
            "findings": governed_findings,
            "verdict_summary": verdict_summary,
            "governance_report_id": governance_report.governance_id,
            "governance_summary": governance_report.summary,
        }
    )
    updated_capsule = None
    if state_capsule is not None:
        pending = [finding.recommendation for finding in governed_findings if finding.status == FindingStatus.OPEN]
        if report_status == QAStatus.PASS:
            pending = ["QA passed for the current compiled scope."]
        continuity_guidance = _dedupe(
            list(getattr(state_capsule, "continuity_guidance", []) or [])
            + [finding.summary for finding in governed_findings if finding.qa_layer == QALayer.DECK]
        )
        updated_capsule = state_capsule.model_copy(
            update={
                "active_gate": WorkflowGate.PRODUCTION_AND_QA,
                "qa_round": bounded_round,
                "pending_actions": _dedupe(pending),
                "open_issues": _dedupe([finding.summary for finding in governed_findings if finding.status == FindingStatus.OPEN]),
                "continuity_guidance": continuity_guidance,
                # `continuity_warnings` remains a deprecated compatibility mirror even
                # after structured continuity alerts become the primary policy surface,
                # so it mirrors `continuity_guidance` instead of drifting separately.
                "continuity_warnings": list(continuity_guidance),
            }
        )
    return DeckQAOutputs(
        qa_report=qa_report,
        qa_governance=governance_report,
        slide_ledger=SlideLedger(deck_title=slide_ledger.deck_title, entries=updated_entries, continuity_notes=_dedupe(slide_ledger.continuity_notes + [finding.summary for finding in findings if finding.qa_layer == QALayer.DECK and finding.severity in {QASeverity.MAJOR, QASeverity.CRITICAL}])),
        slide_build_linkage=SlideBuildLinkage(deck_title=slide_build_linkage.deck_title, pptx_path=slide_build_linkage.pptx_path, slides=updated_links),
        state_capsule=updated_capsule,
        compiled_deck_text=compiled_deck_text,
        compiled_deck_shape_census=compiled_deck_shape_census,
        compiled_deck_authoring_audit=compiled_deck_authoring_audit,
    )


def run_deck_qa_from_files(
    blueprint_path: str | Path,
    design_system_path: str | Path,
    deck_constitution_path: str | Path,
    layout_library_path: str | Path,
    slide_ledger_path: str | Path,
    asset_manifest_path: str | Path,
    viz_manifest_path: str | Path,
    build_manifest_path: str | Path,
    slide_build_linkage_path: str | Path,
    *,
    state_capsule_path: str | Path | None = None,
    prior_report_path: str | Path | None = None,
    qa_governance_path: str | Path | None = None,
    artifact_root: str | Path | None = None,
) -> DeckQAOutputs:
    blueprint = load_state_file(blueprint_path)
    design_system = load_state_file(design_system_path)
    deck_constitution = load_state_file(deck_constitution_path)
    layout_library = load_state_file(layout_library_path)
    slide_ledger = load_state_file(slide_ledger_path)
    asset_manifest = load_state_file(asset_manifest_path)
    viz_manifest = load_state_file(viz_manifest_path)
    build_manifest = load_pptx_compile_file(build_manifest_path)
    slide_build_linkage = load_pptx_compile_file(slide_build_linkage_path)
    state_capsule = load_state_file(state_capsule_path) if state_capsule_path is not None and Path(state_capsule_path).is_file() else None
    prior_report = load_state_file(prior_report_path) if prior_report_path is not None and Path(prior_report_path).is_file() else None
    qa_governance = load_state_file(qa_governance_path) if qa_governance_path is not None and Path(qa_governance_path).is_file() else None
    if blueprint.schema_name != "blueprint" or design_system.schema_name != "design_system" or deck_constitution.schema_name != "deck_constitution" or layout_library.schema_name != "layout_library" or slide_ledger.schema_name != "slide_ledger" or asset_manifest.schema_name != "asset_manifest" or viz_manifest.schema_name != "viz_manifest" or build_manifest.schema_name != "build_manifest" or slide_build_linkage.schema_name != "slide_build_linkage":
        raise TypeError("deck QA received incompatible state artifacts")
    if state_capsule is not None and state_capsule.schema_name != "state_capsule":
        raise TypeError(f"expected state_capsule, found {state_capsule.schema_name}")
    if prior_report is not None and prior_report.schema_name != "qa_report":
        raise TypeError(f"expected qa_report, found {prior_report.schema_name}")
    if qa_governance is not None and qa_governance.schema_name != "qa_governance":
        raise TypeError(f"expected qa_governance, found {qa_governance.schema_name}")
    return run_deck_qa(
        blueprint=blueprint,
        design_system=design_system,
        deck_constitution=deck_constitution,
        layout_library=layout_library,
        slide_ledger=slide_ledger,
        asset_manifest=asset_manifest,
        viz_manifest=viz_manifest,
        build_manifest=build_manifest,
        slide_build_linkage=slide_build_linkage,
        state_capsule=state_capsule,
        prior_report=prior_report,
        qa_governance=qa_governance,
        artifact_root=artifact_root if artifact_root is not None else Path.cwd(),
    )


def write_deck_qa_outputs(outputs: DeckQAOutputs, output_dir: str | Path) -> dict[str, Path]:
    root = Path(output_dir)
    root.mkdir(parents=True, exist_ok=True)
    written = {
        "qa_report": save_state_file(outputs.qa_report, root / "qa-report.json"),
        "slide_ledger": save_state_file(outputs.slide_ledger, root / "slide-ledger.json"),
        "slide_build_linkage": save_state_file(outputs.slide_build_linkage, root / "slide-build-linkage.json"),
    }
    if outputs.qa_governance is not None:
        written["qa_governance"] = save_state_file(outputs.qa_governance, root / "qa-governance.json")
    pptx_path = Path(outputs.slide_build_linkage.pptx_path)
    if not pptx_path.is_absolute():
        candidates = [
            (root.parent / pptx_path).resolve(),
            (root / pptx_path).resolve(),
            pptx_path.resolve(),
        ]
        for candidate in candidates:
            if candidate.exists():
                pptx_path = candidate
                break
    if outputs.compiled_deck_text is not None:
        written["compiled_deck_text"] = _write_json_artifact(root / "compiled-deck-text.json", outputs.compiled_deck_text)
    if outputs.compiled_deck_shape_census is not None:
        written["compiled_deck_shape_census"] = _write_json_artifact(root / "compiled-deck-shape-census.json", outputs.compiled_deck_shape_census)
    if outputs.compiled_deck_shape_census is not None and outputs.compiled_deck_authoring_audit is not None:
        try:
            strip_path, index_path, summary_path = _write_visual_thumbnail_assets(
                output_dir=root / "compiled-deck-thumbnails",
                deck_title=str(outputs.compiled_deck_authoring_audit.get("deck_title") or outputs.qa_report.deck_title),
                pptx_path=str(pptx_path),
                shape_slides=[
                    payload
                    for payload in outputs.compiled_deck_shape_census.get("slides", [])
                    if isinstance(payload, dict)
                ],
                audit_slides=[
                    payload
                    for payload in outputs.compiled_deck_authoring_audit.get("slides", [])
                    if isinstance(payload, dict)
                ],
            )
            written["compiled_deck_thumbnail_strip"] = strip_path
            written["compiled_deck_thumbnail_index"] = index_path
            visual_summary = {
                "schema_name": "compiled_deck_visual_review_summary",
                "schema_version": "1.0",
                "deck_title": str(outputs.compiled_deck_authoring_audit.get("deck_title") or outputs.qa_report.deck_title),
                "pptx_path": str(pptx_path),
                "render_mode": "schematic-geometry",
                "thumbnail_strip_path": str(strip_path),
                "thumbnail_index_path": str(index_path),
                "metrics": {
                    key: value
                    for key, value in outputs.compiled_deck_authoring_audit.get("compiled_metrics", {}).items()
                    if key in {
                        "repeated_geometry_count",
                        "repeated_geometry_rate",
                        "text_card_overuse_rate",
                        "chrome_dominance_rate",
                        "missing_visual_center_count",
                        "appendix_visual_clone_run_length",
                        "mapping_process_example_missing_structure_count",
                    }
                },
                "issue_slides": {
                    key: value
                    for key, value in outputs.compiled_deck_authoring_audit.get("issue_slides", {}).items()
                    if key in {
                        "repeated_geometry_slides",
                        "text_card_like_slides",
                        "chrome_dominance_slides",
                        "missing_visual_center_slides",
                        "appendix_visual_clone_slides",
                        "expected_visual_missing_slides",
                    }
                },
                "slides": [
                    {
                        "slide_number": payload.get("slide_number"),
                        "pptx_index": payload.get("pptx_index"),
                        "deck_mode": payload.get("deck_mode"),
                        "title_text": payload.get("title_text"),
                        "realized_archetype": payload.get("realized_archetype"),
                        "visual_signature": payload.get("visual_signature"),
                        "text_card_like": payload.get("text_card_like"),
                        "chrome_dominant": payload.get("chrome_dominant"),
                        "missing_visual_center": payload.get("missing_visual_center"),
                        "expected_visual_missing": payload.get("expected_visual_missing"),
                    }
                    for payload in outputs.compiled_deck_authoring_audit.get("slides", [])
                    if isinstance(payload, dict)
                ],
            }
            written["compiled_deck_visual_review_summary"] = _write_json_artifact(summary_path, visual_summary)
        except Exception:  # pragma: no cover - artifact best effort, gate will catch missing path
            pass
    if outputs.compiled_deck_authoring_audit is not None:
        audit_payload = dict(outputs.compiled_deck_authoring_audit)
        if "compiled_deck_thumbnail_strip" in written:
            audit_payload["thumbnail_strip_path"] = str(written["compiled_deck_thumbnail_strip"])
        if "compiled_deck_thumbnail_index" in written:
            audit_payload["thumbnail_index_path"] = str(written["compiled_deck_thumbnail_index"])
        if "compiled_deck_visual_review_summary" in written:
            audit_payload["visual_review_summary_path"] = str(written["compiled_deck_visual_review_summary"])
        written["compiled_deck_authoring_audit"] = _write_json_artifact(root / "compiled-deck-authoring-audit.json", audit_payload)
    if outputs.state_capsule is not None:
        written["state_capsule"] = save_state_file(outputs.state_capsule, root / "state-capsule.json")
    return written


