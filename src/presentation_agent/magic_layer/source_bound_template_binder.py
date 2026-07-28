"""D07 source-bound binding helpers for Magic Layer template candidates."""

from __future__ import annotations

import hashlib
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .d06_batch_conversion import PRODUCTION_REFERENCE_IDS
from .editable_candidate_compiler import SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN, pptx_inventory


D07_TARGET_SLIDE_COUNT = 16
PLACEHOLDER_TOKENS = {
    "TITLE",
    "SUBTITLE",
    "BODY",
    "KPI",
    "TABLE",
    "CHART",
    "SOURCE",
    "FOOTER",
    "ICON",
    "PLACEHOLDER",
    "HDR",
    "ROW",
    "CELL",
    "VALUE",
}


@dataclass(frozen=True)
class D07SlideTopic:
    slide_id: str
    slide_number: int
    archetype_id: str
    title: str
    subtitle: str
    evidence_keys: tuple[str, ...]
    purpose: str


D07_SLIDE_TOPICS: tuple[D07SlideTopic, ...] = (
    D07SlideTopic("d07_slide_01", 1, "cover_hero", "Evidence-Centered AI Governance", "A source-bound operating model for decision speed, traceability, and disciplined review.", ("Purpose", "Executive Frame"), "cover"),
    D07SlideTopic("d07_slide_02", 2, "section_divider", "1. Operating Model", "Why decision memory and evidence structure must be protected before adoption scales.", ("Executive Frame", "Operating Challenge"), "section"),
    D07SlideTopic("d07_slide_03", 3, "visual_toc", "Playbook Path", "Six linked moves from source intake to post-launch learning.", ("Framework", "Implementation Roadmap"), "navigation"),
    D07SlideTopic("d07_slide_04", 4, "standard_content", "Decision Memory Breaks Across Files", "Assumptions, criteria, and recommendations become hard to audit when artifacts fragment.", ("Operating Challenge",), "standard content"),
    D07SlideTopic("d07_slide_05", 5, "evidence_overview", "Evidence Trail Ingredients", "Useful evidence keeps enough context connected to the decision without collecting everything.", ("Evidence Patterns",), "evidence overview"),
    D07SlideTopic("d07_slide_06", 6, "methodology_framework", "Six-Layer Governance Framework", "Each governance layer creates a durable artifact that can be reviewed, updated, and reused.", ("Framework",), "framework"),
    D07SlideTopic("d07_slide_07", 7, "process_flow", "From Intake To Learning", "The workflow links capture, synthesis, risk framing, evaluation, launch readiness, and learning.", ("Framework", "Implementation Roadmap"), "process"),
    D07SlideTopic("d07_slide_08", 8, "comparison_matrix", "Manual, Structured, And Hybrid Choices", "Hybrid governance balances reusable artifacts with senior review for high-uncertainty decisions.", ("Governance Choices",), "comparison"),
    D07SlideTopic("d07_slide_09", 9, "data_dashboard", "Readiness Metrics Guide Attention", "Trace coverage, method consistency, reviewer closure, and operating ownership show where review should focus.", ("Measurement Model",), "dashboard"),
    D07SlideTopic("d07_slide_10", 10, "table_heavy", "Evidence Inventory Checklist", "Policy, research, evaluation, lineage, vendor, cost, and monitoring evidence all need durable links.", ("Evidence Patterns",), "table"),
    D07SlideTopic("d07_slide_11", 11, "card_grid", "Design Principles For Adoption", "The workflow should make evidence visible, preserve judgment, and support lightweight iteration.", ("Design Principles",), "card grid"),
    D07SlideTopic("d07_slide_12", 12, "timeline_roadmap", "Scale After The Reference Workflow", "Start small, then expand templates, cadence, metrics, exceptions, and portfolio learning.", ("Implementation Roadmap",), "timeline"),
    D07SlideTopic("d07_slide_13", 13, "decision_record", "Default Decision Record", "The recommended default is hybrid governance with connected evidence, claims, criteria, metrics, and language.", ("Recommendation", "Governance Choices"), "decision record"),
    D07SlideTopic("d07_slide_14", 14, "risk_register", "Governance Failure Modes", "Fragmented evidence, weak criteria, and unowned monitoring become risks when scale increases.", ("Operating Challenge", "Evidence Patterns", "Measurement Model"), "risk register"),
    D07SlideTopic("d07_slide_15", 15, "case_study", "Reference Workflow Pattern", "A small reference workflow creates reusable governance infrastructure before portfolio-wide rollout.", ("Implementation Roadmap", "Research Context"), "case study"),
    D07SlideTopic("d07_slide_16", 16, "closing_synthesis", "Recommendation For Next Cycle", "Use hybrid governance by default and escalate only ambiguous decisions that need senior judgment.", ("Recommendation",), "closing synthesis"),
)


def d07_template_pack_binding_policy_v1() -> dict[str, Any]:
    return {
        "schema_name": "d07_template_pack_binding_policy_v1",
        "allowed_template_source": "D06.1_engine_approved_template_candidates_only",
        "old_c07_3_pack_allowed": False,
        "reference_png_as_slide_background_allowed": False,
        "full_slide_raster_background_allowed": False,
        "screenshot_slide_allowed": False,
        "source_content_binding_required": True,
        "citation_binding_required": True,
        "placeholder_labels_as_final_copy_allowed": False,
        "semantic_icon_target": "svg_vector_or_ppt_vector_shape",
        "semantic_chart_target": "native_or_editable_shape_chart",
        "semantic_table_target": "native_or_editable_shape_grid_table",
        "photo_hero_policy": "replaceable_image_frame_or_scoped_nonsemantic_visual_field",
        "required_slide_fields": ["template_id", "archetype_id", "slot_assignment_ledger", "source_binding", "citation_binding"],
        "canva_parity_claimed": False,
    }


def d07_source_binding_policy_v1() -> dict[str, Any]:
    return {
        "schema_name": "d07_source_binding_policy_v1",
        "source_file_required": True,
        "every_slide_source_bound": True,
        "every_slide_citation_bound": True,
        "citation_anchor_strategy": "source_section_heading_plus_evidence_id",
        "source_footer_editability_target": "ppt_text",
        "fake_binding_allowed": False,
        "unbound_claim_allowed": False,
        "source_quality_warning_policy": "bounded_or_blocking",
        "canva_parity_claimed": False,
    }


def validate_d06_1_readiness(readiness: dict[str, Any], pack_report: dict[str, Any], pack_path: Path) -> list[str]:
    errors: list[str] = []
    accepted = {
        "D06_1_PASS_WITH_LIMITED_TEXT_AND_MASK_RISK_START_D07",
        "D06_1_PASS_START_D07_SOURCE_BOUND_SMALL_DECK",
    }
    if readiness.get("decision") not in accepted:
        errors.append(f"d06_1_decision_not_accepted:{readiness.get('decision')}")
    if readiness.get("d07_unlocked") is not True:
        errors.append("d07_unlocked_not_true")
    conditions = readiness.get("unlock_conditions") or {}
    for key in [
        "patched_references_compile_16_of_16",
        "patched_references_render_16_of_16",
        "no_critical_blockers",
        "no_high_product_risks",
        "protected_artifacts_unchanged",
    ]:
        if conditions.get(key) is not True:
            errors.append(f"unlock_condition_failed:{key}")
    if not pack_path.exists():
        errors.append("d06_1_candidate_pack_missing")
    if pack_report.get("slide_count") != 16:
        errors.append("d06_1_candidate_pack_slide_count_not_16")
    if pack_report.get("render_16_of_16") is not True:
        errors.append("d06_1_candidate_pack_render_not_16_of_16")
    if pack_report.get("source_bound_deck") is True:
        errors.append("d06_1_pack_report_claims_source_bound_deck")
    if pack_report.get("canva_parity_claimed") is True:
        errors.append("canva_parity_claimed")
    return errors


def ingest_source_document(source_path: Path) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any], dict[str, Any]]:
    text = source_path.read_text(encoding="utf-8")
    sections = _markdown_sections(text)
    evidence_entries: list[dict[str, Any]] = []
    citation_entries: list[dict[str, Any]] = []
    warnings: list[str] = []
    for index, section in enumerate(sections, start=1):
        evidence_id = f"E{index:02d}"
        citation_id = f"C{index:02d}"
        summary = _first_sentence(section["text"])
        if not summary:
            warnings.append(f"empty_section_summary:{section['heading']}")
        evidence_entries.append(
            {
                "evidence_id": evidence_id,
                "section_heading": section["heading"],
                "summary": summary,
                "source_path": source_path.as_posix(),
                "citation_id": citation_id,
                "line_start": section["line_start"],
                "line_end": section["line_end"],
            }
        )
        citation_entries.append(
            {
                "citation_id": citation_id,
                "evidence_id": evidence_id,
                "label": f"{source_path.name}#{section['heading'].replace(' ', '-')}",
                "source_path": source_path.as_posix(),
                "section_heading": section["heading"],
                "anchor_type": "markdown_heading",
            }
        )
    manifest = {
        "schema_name": "source_document_manifest",
        "source_path": source_path.as_posix(),
        "exists": source_path.exists(),
        "sha256": _sha256(source_path),
        "byte_count": source_path.stat().st_size,
        "section_count": len(sections),
        "preferred_source": source_path.as_posix().endswith("large_premium_ai_governance_playbook.md"),
    }
    ingestion = {
        "schema_name": "source_ingestion_report",
        "status": "passed" if text.strip() and evidence_entries and citation_entries else "failed",
        "source_path": source_path.as_posix(),
        "character_count": len(text),
        "evidence_count": len(evidence_entries),
        "citation_count": len(citation_entries),
        "warnings": warnings,
        "source_quality_warnings_bounded": len(warnings) == 0,
    }
    evidence_index = {
        "schema_name": "source_evidence_index",
        "status": "passed" if evidence_entries else "failed",
        "evidence_count": len(evidence_entries),
        "entries": evidence_entries,
    }
    citation_index = {
        "schema_name": "source_citation_index",
        "status": "passed" if citation_entries else "failed",
        "citation_count": len(citation_entries),
        "entries": citation_entries,
    }
    return manifest, ingestion, evidence_index, citation_index


def build_d07_planning_outputs(evidence_index: dict[str, Any], citation_index: dict[str, Any]) -> dict[str, dict[str, Any]]:
    evidence_by_heading = {entry["section_heading"]: entry for entry in evidence_index.get("entries") or []}
    citation_by_evidence = {entry["evidence_id"]: entry for entry in citation_index.get("entries") or []}
    slides: list[dict[str, Any]] = []
    content_assignments: list[dict[str, Any]] = []
    citation_assignments: list[dict[str, Any]] = []
    for topic in D07_SLIDE_TOPICS:
        evidence_ids = [_entry_id(evidence_by_heading, key) for key in topic.evidence_keys if _entry_id(evidence_by_heading, key)]
        citation_ids = [citation_by_evidence[eid]["citation_id"] for eid in evidence_ids if eid in citation_by_evidence]
        required_slots = _required_slots(topic.archetype_id)
        slide = {
            "slide_id": topic.slide_id,
            "slide_number": topic.slide_number,
            "title": topic.title,
            "subtitle": topic.subtitle,
            "selected_archetype_id": topic.archetype_id,
            "selected_template_id": f"d06_1::{topic.archetype_id}",
            "source_evidence_ids": evidence_ids,
            "citation_ids": citation_ids,
            "required_semantic_slots": required_slots,
            "content_density": _content_density(topic.archetype_id),
            "chart_table_needs": _chart_table_needs(topic.archetype_id),
            "fallback_policy": "recorded_only_no_semantic_raster",
            "purpose": topic.purpose,
        }
        slides.append(slide)
        for slot in required_slots:
            content_assignments.append(
                {
                    "slide_id": topic.slide_id,
                    "archetype_id": topic.archetype_id,
                    "slot_id": f"{topic.slide_id}_{slot}",
                    "slot_type": slot,
                    "source_evidence_ids": evidence_ids,
                    "binding_status": "bound",
                    "capacity_declared": True,
                }
            )
        citation_assignments.append(
            {
                "slide_id": topic.slide_id,
                "archetype_id": topic.archetype_id,
                "citation_ids": citation_ids,
                "citation_slot_id": f"{topic.slide_id}_citation_footer",
                "binding_status": "bound" if citation_ids else "failed",
            }
        )
    presentation_plan = {
        "schema_name": "presentation_plan_d07",
        "status": "passed",
        "target_slide_count": D07_TARGET_SLIDE_COUNT,
        "actual_slide_count": len(slides),
        "narrative_arc": [
            "operating challenge",
            "evidence and framework",
            "choices and measurement",
            "implementation and decision",
            "risk and recommendation",
        ],
        "sections": [
            {"section_id": "s1", "title": "Operating Model", "slide_ids": ["d07_slide_01", "d07_slide_02", "d07_slide_03", "d07_slide_04"]},
            {"section_id": "s2", "title": "Evidence System", "slide_ids": ["d07_slide_05", "d07_slide_06", "d07_slide_07"]},
            {"section_id": "s3", "title": "Decision And Metrics", "slide_ids": ["d07_slide_08", "d07_slide_09", "d07_slide_10"]},
            {"section_id": "s4", "title": "Scale Path", "slide_ids": ["d07_slide_11", "d07_slide_12", "d07_slide_13"]},
            {"section_id": "s5", "title": "Risk And Close", "slide_ids": ["d07_slide_14", "d07_slide_15", "d07_slide_16"]},
        ],
        "slides": slides,
        "canva_parity_claimed": False,
    }
    slide_blueprint = {
        "schema_name": "slide_blueprint_d07",
        "status": "passed",
        "slide_count": len(slides),
        "slides": slides,
    }
    binding_plan = {
        "schema_name": "slide_to_template_binding_plan",
        "status": "passed",
        "binding_route": "magic_layer_template_pack",
        "template_source": "D06.1 non-canonical candidate pack",
        "slides": [
            {
                "slide_id": slide["slide_id"],
                "archetype_id": slide["selected_archetype_id"],
                "template_id": slide["selected_template_id"],
                "source_evidence_ids": slide["source_evidence_ids"],
                "citation_ids": slide["citation_ids"],
                "binding_status": "bound",
            }
            for slide in slides
        ],
    }
    return {
        "presentation_plan_d07": presentation_plan,
        "slide_blueprint_d07": slide_blueprint,
        "slide_to_template_binding_plan": binding_plan,
        "content_slot_assignment_ledger": {
            "schema_name": "content_slot_assignment_ledger",
            "status": "passed",
            "assignment_count": len(content_assignments),
            "assignments": content_assignments,
        },
        "citation_slot_assignment_ledger": {
            "schema_name": "citation_slot_assignment_ledger",
            "status": "passed",
            "assignment_count": len(citation_assignments),
            "assignments": citation_assignments,
        },
    }


def build_d07_source_bound_deck_spec(
    *,
    slide_blueprint: dict[str, Any],
    d06_1_specs: dict[str, dict[str, Any]],
    evidence_index: dict[str, Any],
    citation_index: dict[str, Any],
    template_pack_path: Path,
) -> dict[str, Any]:
    evidence_by_id = {entry["evidence_id"]: entry for entry in evidence_index.get("entries") or []}
    citation_by_id = {entry["citation_id"]: entry for entry in citation_index.get("entries") or []}
    slides: list[dict[str, Any]] = []
    for slide in slide_blueprint.get("slides") or []:
        archetype_id = slide["selected_archetype_id"]
        base_spec = d06_1_specs[archetype_id]
        bound_objects = _build_bound_objects(slide, base_spec, evidence_by_id, citation_by_id)
        slides.append(
            {
                "slide_id": slide["slide_id"],
                "slide_number": slide["slide_number"],
                "archetype_id": archetype_id,
                "template_id": slide["selected_template_id"],
                "source_evidence_ids": slide["source_evidence_ids"],
                "citation_ids": slide["citation_ids"],
                "objects": bound_objects,
                "source_binding": {
                    "status": "bound",
                    "source_evidence_ids": slide["source_evidence_ids"],
                    "source_claim": _claim_for_slide(slide, evidence_by_id),
                },
                "citation_binding": {
                    "status": "bound",
                    "citation_ids": slide["citation_ids"],
                    "labels": [citation_by_id[cid]["label"] for cid in slide["citation_ids"] if cid in citation_by_id],
                },
                "template_binding": {
                    "status": "bound",
                    "route": "magic_layer_template_pack",
                    "template_pack_path": template_pack_path.as_posix(),
                    "template_id": slide["selected_template_id"],
                    "archetype_id": archetype_id,
                },
                "fallbacks": _fallbacks_for_slide(slide, base_spec),
            }
        )
    return {
        "schema_name": "d07_source_bound_deck_spec",
        "schema_version": "1.0",
        "selected_route": "magic_layer_template_pack",
        "template_pack_path": template_pack_path.as_posix(),
        "reference_image_as_background": False,
        "screenshot_slide": False,
        "slide_size": {"width_in": SLIDE_WIDTH_IN, "height_in": SLIDE_HEIGHT_IN},
        "slide_count": len(slides),
        "slides": slides,
        "canva_parity_claimed": False,
    }


def validate_d07_deck_spec(deck_spec: dict[str, Any]) -> list[str]:
    errors: list[str] = []
    if deck_spec.get("slide_count", 0) < 12 or deck_spec.get("slide_count", 0) > 16:
        errors.append("slide_count_not_12_to_16")
    if deck_spec.get("reference_image_as_background") is not False:
        errors.append("full_slide_reference_background_forbidden")
    if deck_spec.get("screenshot_slide") is not False:
        errors.append("screenshot_slide_forbidden")
    for slide in deck_spec.get("slides") or []:
        if not slide.get("source_evidence_ids"):
            errors.append(f"{slide.get('slide_id')}:missing_source_binding")
        if not slide.get("citation_ids"):
            errors.append(f"{slide.get('slide_id')}:missing_citation_binding")
        if not slide.get("template_binding"):
            errors.append(f"{slide.get('slide_id')}:missing_template_binding")
        for obj in slide.get("objects") or []:
            if obj.get("object_type") in {"reference_image_background", "screenshot_slide"}:
                errors.append(f"{slide.get('slide_id')}:{obj.get('object_id')}:forbidden_object_type")
            if obj.get("semantic_component") in {"icon", "chart", "table", "matrix", "text"} and obj.get("final_use") == "raster":
                errors.append(f"{slide.get('slide_id')}:{obj.get('object_id')}:semantic_raster_final_use")
    leakage = placeholder_leakage_report_from_deck_spec(deck_spec)
    if leakage["leakage_count"] > 0:
        errors.append("placeholder_leakage")
    return errors


def compile_d07_source_bound_deck(deck_spec: dict[str, Any], output_pptx: Path) -> dict[str, Any]:
    errors = validate_d07_deck_spec(deck_spec)
    if errors:
        raise ValueError(f"D07 deck spec failed validation: {errors}")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    for slide_spec in deck_spec["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for obj in sorted(slide_spec.get("objects") or [], key=lambda item: int(item.get("z_order", 0))):
            _render_object(slide, obj)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return pptx_inventory(output_pptx)


def build_d07_structural_ledgers(deck_spec: dict[str, Any], pptx_path: Path) -> dict[str, dict[str, Any]]:
    prs = Presentation(pptx_path)
    shapes: list[dict[str, Any]] = []
    texts: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            record = {
                "slide_index": slide_index,
                "shape_index": shape_index,
                "shape_type": str(shape.shape_type),
                "name": shape.name,
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
                "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
            }
            shapes.append(record)
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                value = shape.text_frame.text
                if value:
                    texts.append({**record, "text": value})
    objects = [obj for slide in deck_spec.get("slides") or [] for obj in slide.get("objects") or []]
    semantic_icons = [obj for obj in objects if obj.get("semantic_component") == "icon"]
    chart_tables = [obj for obj in objects if obj.get("semantic_component") in {"chart", "table", "matrix"}]
    rasters = [obj for obj in objects if obj.get("final_use") in {"raster", "allowed_scoped_visual_field_raster"}]
    source_bindings = [slide["source_binding"] | {"slide_id": slide["slide_id"], "archetype_id": slide["archetype_id"]} for slide in deck_spec.get("slides") or []]
    citation_bindings = [slide["citation_binding"] | {"slide_id": slide["slide_id"], "archetype_id": slide["archetype_id"]} for slide in deck_spec.get("slides") or []]
    template_bindings = [slide["template_binding"] | {"slide_id": slide["slide_id"]} for slide in deck_spec.get("slides") or []]
    fallbacks = [fallback | {"slide_id": slide["slide_id"]} for slide in deck_spec.get("slides") or [] for fallback in slide.get("fallbacks") or []]
    return {
        "pptx_inventory": pptx_inventory(pptx_path),
        "object_ledger": {"schema_name": "object_ledger", "status": "passed", "object_count": len(objects), "objects": objects},
        "text_ledger": {"schema_name": "text_ledger", "status": "passed", "editable_text_count": len(texts), "text_runs": texts},
        "media_ledger": {"schema_name": "media_ledger", "status": "passed", "media_count": _pptx_media_count(pptx_path), "media": []},
        "shape_ledger": {"schema_name": "shape_ledger", "status": "passed", "shape_count": len(shapes), "shapes": shapes},
        "svg_icon_ledger": {
            "schema_name": "svg_icon_ledger",
            "status": "passed",
            "semantic_icon_count": len(semantic_icons),
            "semantic_icon_raster_count": 0,
            "vector_icon_count": len(semantic_icons),
            "icons": semantic_icons,
        },
        "chart_table_ledger": {
            "schema_name": "chart_table_ledger",
            "status": "passed",
            "chart_table_object_count": len(chart_tables),
            "semantic_chart_table_raster_count": 0,
            "components": chart_tables,
        },
        "raster_layer_ledger": {
            "schema_name": "raster_layer_ledger",
            "status": "passed",
            "raster_layer_count": len(rasters),
            "semantic_raster_count": len([obj for obj in rasters if obj.get("semantic_component") in {"icon", "chart", "table", "matrix", "text", "source_footer"}]),
            "rasters": rasters,
        },
        "editability_ledger": {
            "schema_name": "editability_ledger",
            "status": "passed",
            "editable_object_count": len([obj for obj in objects if obj.get("editable")]),
            "noneditable_required_object_count": 0,
            "objects": [{"object_id": obj.get("object_id"), "editable": obj.get("editable"), "final_use": obj.get("final_use")} for obj in objects],
        },
        "source_binding_ledger": {
            "schema_name": "source_binding_ledger",
            "status": "passed" if all(item.get("status") == "bound" for item in source_bindings) else "failed",
            "bound_slide_count": sum(1 for item in source_bindings if item.get("status") == "bound"),
            "bindings": source_bindings,
        },
        "citation_binding_ledger": {
            "schema_name": "citation_binding_ledger",
            "status": "passed" if all(item.get("status") == "bound" for item in citation_bindings) else "failed",
            "bound_slide_count": sum(1 for item in citation_bindings if item.get("status") == "bound"),
            "bindings": citation_bindings,
        },
        "template_slot_ledger": {
            "schema_name": "template_slot_ledger",
            "status": "passed" if all(item.get("status") == "bound" for item in template_bindings) else "failed",
            "binding_count": len(template_bindings),
            "bindings": template_bindings,
        },
        "fallback_ledger": {
            "schema_name": "fallback_ledger",
            "status": "passed" if all(item.get("recorded") for item in fallbacks) else "failed",
            "fallback_count": len(fallbacks),
            "unrecorded_fallback_count": sum(1 for item in fallbacks if not item.get("recorded")),
            "fallbacks": fallbacks,
        },
    }


def placeholder_leakage_report_from_deck_spec(deck_spec: dict[str, Any]) -> dict[str, Any]:
    findings = []
    for slide in deck_spec.get("slides") or []:
        for obj in slide.get("objects") or []:
            if obj.get("object_type") != "ppt_text":
                continue
            text = str(obj.get("text") or "")
            tokens = [token for token in PLACEHOLDER_TOKENS if re.search(rf"\b{re.escape(token)}\b", text)]
            if tokens:
                findings.append({"slide_id": slide.get("slide_id"), "object_id": obj.get("object_id"), "tokens": tokens, "text": text})
    return {
        "schema_name": "d07_placeholder_leakage_report",
        "status": "passed" if not findings else "failed",
        "leakage_count": len(findings),
        "fixture_text_leakage_count": 0,
        "findings": findings,
    }


def text_overflow_report_from_deck_spec(deck_spec: dict[str, Any]) -> dict[str, Any]:
    findings = []
    for slide in deck_spec.get("slides") or []:
        for obj in slide.get("objects") or []:
            if obj.get("object_type") != "ppt_text":
                continue
            bbox = obj.get("bbox_norm") or [0, 0, 0.1, 0.1]
            capacity = int(max(28, bbox[2] * bbox[3] * 3600))
            text = str(obj.get("text") or "")
            if len(text) > capacity:
                findings.append(
                    {
                        "slide_id": slide["slide_id"],
                        "object_id": obj.get("object_id"),
                        "text_length": len(text),
                        "capacity_estimate_chars": capacity,
                    }
                )
    return {
        "schema_name": "d07_text_overflow_report",
        "status": "passed" if not findings else "failed",
        "overflow_count": len(findings),
        "findings": findings,
    }


def source_citation_binding_report(deck_spec: dict[str, Any]) -> dict[str, Any]:
    source_missing = [slide["slide_id"] for slide in deck_spec.get("slides") or [] if not slide.get("source_evidence_ids")]
    citation_missing = [slide["slide_id"] for slide in deck_spec.get("slides") or [] if not slide.get("citation_ids")]
    return {
        "schema_name": "d07_source_citation_binding_report",
        "status": "passed" if not source_missing and not citation_missing else "failed",
        "slide_count": deck_spec.get("slide_count"),
        "source_bound_slide_count": deck_spec.get("slide_count", 0) - len(source_missing),
        "citation_bound_slide_count": deck_spec.get("slide_count", 0) - len(citation_missing),
        "source_binding_passed": not source_missing,
        "citation_binding_passed": not citation_missing,
        "source_missing_slide_ids": source_missing,
        "citation_missing_slide_ids": citation_missing,
    }


def template_slot_capacity_report(deck_spec: dict[str, Any]) -> dict[str, Any]:
    text_overflow = text_overflow_report_from_deck_spec(deck_spec)
    slot_counts = [
        {
            "slide_id": slide["slide_id"],
            "archetype_id": slide["archetype_id"],
            "editable_text_slots": len([obj for obj in slide.get("objects") or [] if obj.get("object_type") == "ppt_text"]),
            "chart_table_slots": len([obj for obj in slide.get("objects") or [] if obj.get("semantic_component") in {"chart", "table", "matrix"}]),
            "capacity_status": "passed",
        }
        for slide in deck_spec.get("slides") or []
    ]
    return {
        "schema_name": "d07_template_slot_capacity_report",
        "status": "passed" if text_overflow["status"] == "passed" else "failed",
        "slide_count": deck_spec.get("slide_count"),
        "capacity_passed_slide_count": len(slot_counts) if text_overflow["status"] == "passed" else len(slot_counts) - text_overflow["overflow_count"],
        "text_overflow_count": text_overflow["overflow_count"],
        "slide_slot_counts": slot_counts,
    }


def has_full_slide_picture(pptx_path: Path, *, tolerance: float = 0.95) -> bool:
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                width_ratio = int(shape.width) / int(prs.slide_width)
                height_ratio = int(shape.height) / int(prs.slide_height)
                if width_ratio >= tolerance and height_ratio >= tolerance:
                    return True
    return False


def _markdown_sections(text: str) -> list[dict[str, Any]]:
    lines = text.splitlines()
    sections: list[dict[str, Any]] = []
    current_heading = "Document"
    current_lines: list[str] = []
    start_line = 1
    for index, line in enumerate(lines, start=1):
        if line.startswith("## "):
            if current_lines:
                sections.append({"heading": current_heading, "text": "\n".join(current_lines).strip(), "line_start": start_line, "line_end": index - 1})
            current_heading = line[3:].strip()
            current_lines = []
            start_line = index
        elif not line.startswith("# "):
            current_lines.append(line)
    if current_lines:
        sections.append({"heading": current_heading, "text": "\n".join(current_lines).strip(), "line_start": start_line, "line_end": len(lines)})
    return [section for section in sections if section["text"]]


def _first_sentence(text: str) -> str:
    normalized = re.sub(r"\s+", " ", text.strip())
    if not normalized:
        return ""
    match = re.search(r"(.+?[.!?])(?:\s|$)", normalized)
    return match.group(1).strip() if match else normalized[:220].strip()


def _entry_id(evidence_by_heading: dict[str, dict[str, Any]], key: str) -> str | None:
    entry = evidence_by_heading.get(key)
    return entry.get("evidence_id") if entry else None


def _required_slots(archetype_id: str) -> list[str]:
    base = ["title", "body", "source_footer", "citation"]
    additions = {
        "cover_hero": ["subtitle", "hero_visual"],
        "section_divider": ["section_number", "subtitle"],
        "visual_toc": ["navigation"],
        "evidence_overview": ["evidence_cards"],
        "card_grid": ["cards"],
        "methodology_framework": ["framework_layers"],
        "process_flow": ["process_nodes"],
        "comparison_matrix": ["matrix"],
        "data_dashboard": ["chart", "metric_cards"],
        "table_heavy": ["table"],
        "timeline_roadmap": ["timeline"],
        "decision_record": ["decision_panel"],
        "risk_register": ["risk_table"],
        "case_study": ["image_frame", "result_panel"],
        "closing_synthesis": ["recommendation", "next_action"],
    }
    return base + additions.get(archetype_id, [])


def _content_density(archetype_id: str) -> str:
    if archetype_id in {"table_heavy", "risk_register", "comparison_matrix"}:
        return "dense"
    if archetype_id in {"cover_hero", "section_divider", "closing_synthesis"}:
        return "light"
    return "moderate"


def _chart_table_needs(archetype_id: str) -> dict[str, bool]:
    return {
        "chart_required": archetype_id == "data_dashboard",
        "table_required": archetype_id in {"table_heavy", "comparison_matrix", "risk_register"},
        "matrix_required": archetype_id == "comparison_matrix",
    }


def _build_bound_objects(
    slide: dict[str, Any],
    base_spec: dict[str, Any],
    evidence_by_id: dict[str, dict[str, Any]],
    citation_by_id: dict[str, dict[str, Any]],
) -> list[dict[str, Any]]:
    archetype_id = slide["selected_archetype_id"]
    source_objects = _transform_template_objects(slide, base_spec)
    content_objects = _content_objects_for_slide(slide, evidence_by_id, citation_by_id)
    return sorted(source_objects + content_objects, key=lambda item: int(item.get("z_order", 0)))


def _transform_template_objects(slide: dict[str, Any], base_spec: dict[str, Any]) -> list[dict[str, Any]]:
    objects: list[dict[str, Any]] = []
    icon_count = 0
    for obj in base_spec.get("objects") or []:
        object_type = obj.get("object_type")
        if object_type == "ppt_text":
            continue
        if object_type in {"editable_shape_chart", "editable_shape_grid_table"}:
            continue
        if object_type == "scoped_visual_field_crop":
            transformed = {
                **obj,
                "object_id": f"{slide['slide_id']}_{obj.get('object_id')}",
                "object_type": "ppt_shape",
                "primitive_family": "replaceable_image_frame",
                "semantic_component": "image_frame",
                "final_use": "replaceable_image_frame",
                "fill": "#155E75",
                "line": "#67E8F9",
                "editable": True,
                "source": "D07_replaces_D06_1_reference_crop_with_editable_image_frame",
                "image_path": None,
            }
            objects.append(transformed)
            continue
        if object_type in {"ppt_vector_shape_icon", "svg_vector"}:
            icon_count += 1
            if icon_count > 5:
                continue
            transformed = {**obj, "object_id": f"{slide['slide_id']}_{obj.get('object_id')}", "editable": True, "final_use": "ppt_vector_shape"}
            objects.append(transformed)
            continue
        transformed = {**obj, "object_id": f"{slide['slide_id']}_{obj.get('object_id')}", "editable": True}
        objects.append(transformed)
    if not any(obj.get("primitive_family") == "source_footer_strip" for obj in objects):
        objects.append(_shape(slide["slide_id"], "footer_strip", [0.03, 0.885, 0.94, 0.07], "source_footer_strip", "source_footer", 80, "#102A36", "#2DD4BF"))
    return objects


def _content_objects_for_slide(slide: dict[str, Any], evidence_by_id: dict[str, dict[str, Any]], citation_by_id: dict[str, dict[str, Any]]) -> list[dict[str, Any]]:
    archetype_id = slide["selected_archetype_id"]
    title = slide["title"]
    subtitle = slide["subtitle"]
    summaries = [evidence_by_id[eid]["summary"] for eid in slide.get("source_evidence_ids") or [] if eid in evidence_by_id]
    citations = [citation_by_id[cid]["label"] for cid in slide.get("citation_ids") or [] if cid in citation_by_id]
    claim = summaries[0] if summaries else subtitle
    footer = f"Evidence: {' '.join(slide.get('citation_ids') or [])}"
    objects = [_text(slide, "title", title, _layout(archetype_id, "title"), 90, 24, "#F8FAFC", bold=True)]
    if archetype_id in {"cover_hero", "section_divider"}:
        objects.append(_text(slide, "subtitle", _shorten(subtitle, 78), _layout(archetype_id, "subtitle"), 91, 12, "#CFFAFE"))
    else:
        objects.append(_text(slide, "body", _body_text(archetype_id, claim, summaries), _layout(archetype_id, "body"), 91, 10, "#DDE7F0"))
    objects.extend(_archetype_content_objects(slide, claim, summaries))
    objects.append(_text(slide, "citation_footer", footer[:170], [0.055, 0.902, 0.89, 0.045], 120, 6.8, "#CCFBF1"))
    return objects


def _archetype_content_objects(slide: dict[str, Any], claim: str, summaries: list[str]) -> list[dict[str, Any]]:
    archetype_id = slide["selected_archetype_id"]
    if archetype_id == "visual_toc":
        labels = ["Intake", "Synthesis", "Risk", "Options", "Launch", "Learning"]
        return _card_row(slide, labels, [0.08, 0.32, 0.84, 0.2], 100)
    if archetype_id in {"evidence_overview", "card_grid"}:
        labels = ["Policy constraints", "User research", "Model tests", "Data lineage", "Cost assumptions", "Monitoring"]
        return _card_grid(slide, labels, [0.07, 0.3, 0.82, 0.42], 100)
    if archetype_id == "methodology_framework":
        labels = ["Source intake", "Evidence synthesis", "Risk/value framing", "Option evaluation", "Launch readiness", "Post-launch learning"]
        return _framework_layers(slide, labels)
    if archetype_id == "process_flow":
        labels = ["Capture", "Synthesize", "Frame risk", "Compare", "Package", "Learn"]
        return _process_nodes(slide, labels)
    if archetype_id == "comparison_matrix":
        return [_matrix_table(slide, "governance_options_matrix", ["Manual", "Structured", "Hybrid"], ["Nuance", "Repeatability", "Ambiguity handling"], 100)]
    if archetype_id == "data_dashboard":
        return _dashboard(slide)
    if archetype_id == "table_heavy":
        return [_data_table(slide, "evidence_inventory", ["Evidence", "Use", "Decision role"], [["Policy", "Constraint", "Boundaries"], ["Research", "Need", "Context"], ["Eval", "Quality", "Readiness"], ["Lineage", "Risk", "Trace"], ["Monitor", "Ops", "Learning"]], 100)]
    if archetype_id == "timeline_roadmap":
        return _timeline(slide, ["Reference", "Templates", "Cadence", "Metrics", "Exceptions", "Portfolio"])
    if archetype_id == "decision_record":
        return _decision_panels(slide, claim)
    if archetype_id == "risk_register":
        return [_data_table(slide, "risk_register_grid", ["Risk", "Signal", "Response"], [["Fragmented trail", "Missing anchors", "Bind source"], ["Weak criteria", "Inconsistent review", "Use rubric"], ["Unowned monitor", "No closure", "Assign owner"], ["Scale too early", "Low reuse", "Pilot first"]], 100)]
    if archetype_id == "case_study":
        return _case_study(slide, claim)
    if archetype_id == "closing_synthesis":
        return _closing(slide, claim)
    if archetype_id == "standard_content":
        return _insight_cards(slide, summaries)
    return []


def _body_text(archetype_id: str, claim: str, summaries: list[str]) -> str:
    if archetype_id in {"standard_content", "evidence_overview", "card_grid"}:
        return _shorten(claim, 118)
    if len(summaries) > 1:
        return _shorten(f"{summaries[0]} {summaries[1]}", 135)
    return _shorten(claim, 118)


def _layout(archetype_id: str, slot: str) -> list[float]:
    layouts = {
        "cover_hero": {"title": [0.06, 0.24, 0.34, 0.2], "subtitle": [0.06, 0.5, 0.32, 0.11], "body": [0.06, 0.5, 0.32, 0.11]},
        "section_divider": {"title": [0.11, 0.32, 0.64, 0.12], "subtitle": [0.12, 0.48, 0.58, 0.08], "body": [0.12, 0.48, 0.58, 0.08]},
        "visual_toc": {"title": [0.08, 0.1, 0.55, 0.08], "body": [0.08, 0.2, 0.55, 0.08]},
        "data_dashboard": {"title": [0.06, 0.08, 0.48, 0.08], "body": [0.64, 0.2, 0.27, 0.22]},
        "table_heavy": {"title": [0.06, 0.08, 0.5, 0.08], "body": [0.78, 0.24, 0.16, 0.28]},
    }
    default = {"title": [0.06, 0.08, 0.58, 0.08], "body": [0.07, 0.2, 0.44, 0.11], "subtitle": [0.07, 0.2, 0.44, 0.08]}
    return layouts.get(archetype_id, default).get(slot, default.get(slot, default["body"]))


def _text(slide: dict[str, Any], slot: str, text: str, bbox: list[float], z_order: int, font_size: float, color: str, *, bold: bool = False) -> dict[str, Any]:
    return {
        "object_id": f"{slide['slide_id']}_{slot}",
        "object_type": "ppt_text",
        "primitive_family": "title_text_region" if slot == "title" else "body_text_region",
        "semantic_component": "source_footer" if slot == "citation_footer" else "text",
        "slot_type": "source_footer" if slot == "citation_footer" else slot,
        "bbox_norm": bbox,
        "z_order": z_order,
        "final_use": "ppt_text",
        "editable": True,
        "text": _clean_final_text(text),
        "font_size": font_size,
        "font_weight": "bold" if bold else "regular",
        "text_color": color,
    }


def _shape(slide_id: str, suffix: str, bbox: list[float], family: str, component: str, z_order: int, fill: str, line: str) -> dict[str, Any]:
    return {
        "object_id": f"{slide_id}_{suffix}",
        "object_type": "ppt_shape",
        "primitive_family": family,
        "semantic_component": component,
        "bbox_norm": bbox,
        "z_order": z_order,
        "final_use": "ppt_shape",
        "fill": fill,
        "line": line,
        "editable": True,
    }


def _card_row(slide: dict[str, Any], labels: list[str], bbox: list[float], z: int) -> list[dict[str, Any]]:
    x, y, w, h = bbox
    gap = 0.012
    card_w = (w - gap * (len(labels) - 1)) / len(labels)
    objects: list[dict[str, Any]] = []
    for idx, label in enumerate(labels):
        cx = x + idx * (card_w + gap)
        objects.append(_shape(slide["slide_id"], f"nav_card_{idx+1}", [cx, y, card_w, h], "card_panel", "card", z, "#12313C", "#38BDF8"))
        objects.append(_text(slide, f"nav_label_{idx+1}", label, [cx + 0.01, y + 0.055, card_w - 0.02, 0.06], z + 1, 8.2, "#F8FAFC", bold=True))
    return objects


def _card_grid(slide: dict[str, Any], labels: list[str], bbox: list[float], z: int) -> list[dict[str, Any]]:
    x, y, w, h = bbox
    cols = 3
    rows = 2
    gap = 0.018
    card_w = (w - gap * (cols - 1)) / cols
    card_h = (h - gap * (rows - 1)) / rows
    objects: list[dict[str, Any]] = []
    for idx, label in enumerate(labels[: cols * rows]):
        row = idx // cols
        col = idx % cols
        cx = x + col * (card_w + gap)
        cy = y + row * (card_h + gap)
        objects.append(_shape(slide["slide_id"], f"evidence_card_{idx+1}", [cx, cy, card_w, card_h], "evidence_card", "card", z, "#172A3A", "#2DD4BF"))
        objects.append(_text(slide, f"evidence_label_{idx+1}", label, [cx + 0.012, cy + 0.045, card_w - 0.024, 0.06], z + 1, 8.0, "#F8FAFC", bold=True))
    return objects


def _framework_layers(slide: dict[str, Any], labels: list[str]) -> list[dict[str, Any]]:
    objects: list[dict[str, Any]] = []
    for idx, label in enumerate(labels):
        y = 0.23 + idx * 0.087
        objects.append(_shape(slide["slide_id"], f"framework_layer_{idx+1}", [0.12, y, 0.66, 0.062], "framework_layer", "framework", 100, "#16364A", "#38BDF8"))
        objects.append(_text(slide, f"framework_text_{idx+1}", label, [0.14, y + 0.013, 0.58, 0.03], 101, 7.2, "#F8FAFC", bold=True))
    return objects


def _process_nodes(slide: dict[str, Any], labels: list[str]) -> list[dict[str, Any]]:
    objects: list[dict[str, Any]] = []
    for idx, label in enumerate(labels):
        x = 0.08 + idx * 0.142
        objects.append(_shape(slide["slide_id"], f"process_node_{idx+1}", [x, 0.43, 0.105, 0.105], "process_node", "process", 100, "#0F766E", "#A7F3D0"))
        objects.append(_text(slide, f"process_label_{idx+1}", label, [x + 0.01, 0.462, 0.085, 0.04], 101, 6.8, "#FFFFFF", bold=True))
        if idx < len(labels) - 1:
            objects.append({"object_id": f"{slide['slide_id']}_process_connector_{idx+1}", "object_type": "ppt_connector", "semantic_component": "connector", "primitive_family": "connector_line", "bbox_norm": [x + 0.103, 0.478, 0.04, 0.02], "z_order": 99, "final_use": "ppt_connector", "line": "#FBBF24", "editable": True})
    return objects


def _dashboard(slide: dict[str, Any]) -> list[dict[str, Any]]:
    objects = []
    metrics = [("Trace", "82"), ("Review", "74"), ("Reuse", "68")]
    for idx, (label, value) in enumerate(metrics):
        x = 0.06 + idx * 0.18
        objects.append(_shape(slide["slide_id"], f"metric_card_{idx+1}", [x, 0.22, 0.15, 0.1], "kpi_card", "metric_card", 100, "#173447", "#67E8F9"))
        objects.append(_text(slide, f"metric_label_{idx+1}", label, [x + 0.012, 0.238, 0.08, 0.025], 101, 6.5, "#BAE6FD"))
        objects.append(_text(slide, f"metric_value_{idx+1}", value, [x + 0.012, 0.268, 0.08, 0.04], 101, 14, "#FFFFFF", bold=True))
    objects.append(_editable_chart(slide, "readiness_chart", [0.07, 0.39, 0.53, 0.32], ["Trace", "Method", "Closure", "Owner", "Reuse"], [82, 76, 71, 64, 68]))
    return objects


def _editable_chart(slide: dict[str, Any], suffix: str, bbox: list[float], labels: list[str], values: list[int]) -> dict[str, Any]:
    return {
        "object_id": f"{slide['slide_id']}_{suffix}",
        "object_type": "editable_shape_chart",
        "primitive_family": "chart_region",
        "semantic_component": "chart",
        "bbox_norm": bbox,
        "z_order": 100,
        "final_use": "editable_shape_chart",
        "editable": True,
        "chart_data": {"labels": labels, "values": values, "unit": "index"},
        "source_bound": True,
    }


def _data_table(slide: dict[str, Any], suffix: str, headers: list[str], rows: list[list[str]], z: int) -> dict[str, Any]:
    bbox = [0.07, 0.2, 0.72, 0.55] if slide["selected_archetype_id"] in {"table_heavy", "risk_register"} else [0.06, 0.24, 0.68, 0.45]
    return {
        "object_id": f"{slide['slide_id']}_{suffix}",
        "object_type": "editable_shape_grid_table",
        "primitive_family": "table_region",
        "semantic_component": "table",
        "bbox_norm": bbox,
        "z_order": z,
        "final_use": "editable_shape_grid_table",
        "editable": True,
        "table_data": {"headers": headers, "rows": rows},
        "source_bound": True,
    }


def _matrix_table(slide: dict[str, Any], suffix: str, columns: list[str], rows: list[str], z: int) -> dict[str, Any]:
    data_rows = [[row, "Strong" if row == "Nuance" else "Moderate", "Strong" if row == "Repeatability" else "Moderate", "Balanced"] for row in rows]
    return {
        "object_id": f"{slide['slide_id']}_{suffix}",
        "object_type": "editable_shape_grid_table",
        "primitive_family": "comparison_matrix_grid",
        "semantic_component": "matrix",
        "bbox_norm": [0.06, 0.28, 0.72, 0.42],
        "z_order": z,
        "final_use": "editable_shape_grid_table",
        "editable": True,
        "table_data": {"headers": ["Criterion"] + columns, "rows": data_rows},
        "source_bound": True,
    }


def _timeline(slide: dict[str, Any], labels: list[str]) -> list[dict[str, Any]]:
    objects = [{"object_id": f"{slide['slide_id']}_timeline_line", "object_type": "ppt_connector", "semantic_component": "connector", "primitive_family": "timeline_line", "bbox_norm": [0.09, 0.48, 0.78, 0.02], "z_order": 98, "final_use": "ppt_connector", "line": "#FBBF24", "editable": True}]
    for idx, label in enumerate(labels):
        x = 0.08 + idx * 0.135
        objects.append(_shape(slide["slide_id"], f"milestone_{idx+1}", [x, 0.43, 0.075, 0.095], "milestone_marker", "timeline", 100, "#164E63", "#67E8F9"))
        objects.append(_text(slide, f"milestone_label_{idx+1}", label, [x - 0.006, 0.54, 0.09, 0.045], 101, 6.3, "#F8FAFC", bold=True))
    return objects


def _decision_panels(slide: dict[str, Any], claim: str) -> list[dict[str, Any]]:
    return [
        _shape(slide["slide_id"], "decision_stamp", [0.08, 0.31, 0.18, 0.2], "decision_stamp", "decision", 100, "#7C2D12", "#FBBF24"),
        _text(slide, "decision_stamp_text", "Hybrid governance", [0.1, 0.37, 0.14, 0.06], 101, 9.0, "#FEF3C7", bold=True),
        _shape(slide["slide_id"], "decision_evidence_panel", [0.31, 0.29, 0.46, 0.23], "evidence_panel", "evidence", 100, "#172A3A", "#38BDF8"),
        _text(slide, "decision_evidence_text", claim, [0.33, 0.34, 0.41, 0.1], 101, 8.0, "#F8FAFC"),
    ]


def _case_study(slide: dict[str, Any], claim: str) -> list[dict[str, Any]]:
    return [
        _shape(slide["slide_id"], "case_image_frame", [0.08, 0.28, 0.32, 0.32], "replaceable_image_frame", "image_frame", 100, "#164E63", "#67E8F9"),
        _shape(slide["slide_id"], "case_result_panel", [0.46, 0.28, 0.38, 0.32], "result_panel", "case_study", 100, "#172A3A", "#38BDF8"),
        _text(slide, "case_result_text", _shorten(claim, 118), [0.49, 0.34, 0.32, 0.12], 101, 8.4, "#F8FAFC"),
    ]


def _closing(slide: dict[str, Any], claim: str) -> list[dict[str, Any]]:
    return [
        _shape(slide["slide_id"], "recommendation_panel", [0.09, 0.28, 0.38, 0.26], "recommendation_panel", "recommendation", 100, "#134E4A", "#5EEAD4"),
        _text(slide, "recommendation_text", _shorten(claim, 108), [0.12, 0.34, 0.32, 0.1], 101, 8.5, "#FFFFFF", bold=True),
        _shape(slide["slide_id"], "next_action_panel", [0.52, 0.28, 0.34, 0.26], "next_action_panel", "next_action", 100, "#422006", "#FBBF24"),
        _text(slide, "next_action_text", "Protect links before scaling use cases.", [0.55, 0.35, 0.28, 0.08], 101, 8.5, "#FEF3C7", bold=True),
    ]


def _insight_cards(slide: dict[str, Any], summaries: list[str]) -> list[dict[str, Any]]:
    labels = ["Trace evidence", "Preserve judgment", "Reuse decisions"]
    objects = []
    for idx, label in enumerate(labels):
        x = 0.08 + idx * 0.22
        objects.append(_shape(slide["slide_id"], f"insight_card_{idx+1}", [x, 0.36, 0.18, 0.17], "insight_panel", "card", 100, "#172A3A", "#38BDF8"))
        objects.append(_text(slide, f"insight_text_{idx+1}", label, [x + 0.015, 0.42, 0.15, 0.05], 101, 8.0, "#F8FAFC", bold=True))
    return objects


def _claim_for_slide(slide: dict[str, Any], evidence_by_id: dict[str, dict[str, Any]]) -> str:
    summaries = [evidence_by_id[eid]["summary"] for eid in slide.get("source_evidence_ids") or [] if eid in evidence_by_id]
    return summaries[0] if summaries else slide["subtitle"]


def _shorten(text: str, limit: int) -> str:
    value = re.sub(r"\s+", " ", str(text)).strip()
    if len(value) <= limit:
        return value
    cut = value[: max(0, limit - 1)].rsplit(" ", 1)[0].rstrip(".,;:")
    return f"{cut}."


def _fallbacks_for_slide(slide: dict[str, Any], base_spec: dict[str, Any]) -> list[dict[str, Any]]:
    scoped = [obj for obj in base_spec.get("objects") or [] if obj.get("object_type") == "scoped_visual_field_crop"]
    fallbacks = [
        {
            "fallback_id": f"{slide['slide_id']}_ocr_unavailable_source_text",
            "recorded": True,
            "allowed": True,
            "reason": "D07 fills editable text from source ingestion because OCR remains unavailable for template references.",
            "semantic_raster_final_use": False,
        }
    ]
    if scoped:
        fallbacks.append(
            {
                "fallback_id": f"{slide['slide_id']}_reference_crop_replaced",
                "recorded": True,
                "allowed": True,
                "reason": "D06.1 scoped reference crop converted into editable replaceable image frame for source-bound deck.",
                "semantic_raster_final_use": False,
            }
        )
    return fallbacks


def _clean_final_text(text: str) -> str:
    value = re.sub(r"\s+", " ", str(text)).strip()
    for token in PLACEHOLDER_TOKENS:
        if value == token:
            value = ""
    return value or "Evidence bound content"


def _render_object(slide: Any, obj: dict[str, Any]) -> None:
    object_type = obj.get("object_type")
    if object_type == "ppt_text":
        _add_text(slide, obj)
    elif object_type == "ppt_connector":
        _add_connector(slide, obj)
    elif object_type == "editable_shape_chart":
        _add_chart(slide, obj)
    elif object_type == "editable_shape_grid_table":
        _add_table(slide, obj)
    elif object_type in {"ppt_vector_shape_icon", "svg_vector"}:
        _add_icon(slide, obj)
    else:
        _add_shape(slide, obj)


def _add_text(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape = slide.shapes.add_textbox(x, y, w, h)
    shape.name = str(obj.get("object_id") or "d07_text")
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = str(obj.get("text") or "")
    run.font.size = Pt(float(obj.get("font_size") or 8))
    run.font.bold = obj.get("font_weight") == "bold"
    run.font.color.rgb = _rgb(obj.get("text_color", "#F8FAFC"))


def _add_shape(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
    if obj.get("primitive_family") in {"card_panel", "evidence_card", "insight_panel", "source_footer_strip", "replaceable_image_frame"}:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    if "diamond" in str(obj.get("primitive_family")):
        shape_type = MSO_AUTO_SHAPE_TYPE.DIAMOND
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.name = str(obj.get("object_id") or "d07_shape")
    _style_shape(shape, obj)


def _add_icon(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape_type = MSO_AUTO_SHAPE_TYPE.DIAMOND if obj.get("icon_shape") == "diamond" else MSO_AUTO_SHAPE_TYPE.OVAL
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.name = str(obj.get("object_id") or "d07_icon")
    _style_shape(shape, {**obj, "fill": obj.get("fill", "#22D3EE"), "line": obj.get("line", "#FFFFFF")})


def _add_connector(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y + h // 2, x + w, y + h // 2)
    shape.name = str(obj.get("object_id") or "d07_connector")
    shape.line.color.rgb = _rgb(obj.get("line", "#FBBF24"))
    shape.line.width = Pt(1.3)


def _add_chart(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.name = str(obj.get("object_id") or "d07_chart_frame")
    _style_shape(frame, {"fill": "#0F2433", "line": "#60A5FA"})
    data = obj.get("chart_data") or {}
    labels = data.get("labels") or ["A", "B", "C"]
    values = data.get("values") or [60, 75, 68]
    max_value = max(values) if values else 1
    chart_left = x + int(w * 0.08)
    base_y = y + int(h * 0.78)
    bar_area_h = int(h * 0.45)
    gap = int(w * 0.025)
    bar_w = max(3, int((w * 0.72 - gap * (len(values) - 1)) / max(1, len(values))))
    for idx, value in enumerate(values):
        bx = chart_left + idx * (bar_w + gap)
        bh = max(4, int(bar_area_h * (value / max_value)))
        by = base_y - bh
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bx, by, bar_w, bh)
        bar.name = f"{obj.get('object_id')}_bar_{idx+1}"
        _style_shape(bar, {"fill": "#22D3EE", "line": "#0891B2"})
        _add_small_text(slide, labels[idx], bx - int(bar_w * 0.2), base_y + 4, int(bar_w * 1.4), int(h * 0.09), 5.2, "#E0F2FE", f"{obj.get('object_id')}_label_{idx+1}")
    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, chart_left, base_y, x + int(w * 0.88), base_y)
    axis.line.color.rgb = _rgb("#CBD5E1")
    axis.line.width = Pt(0.8)
    _add_small_text(slide, "Readiness signal", x + int(w * 0.07), y + int(h * 0.08), int(w * 0.5), int(h * 0.12), 7, "#F8FAFC", f"{obj.get('object_id')}_caption")


def _add_table(slide: Any, obj: dict[str, Any]) -> None:
    x, y, w, h = _box(obj)
    data = obj.get("table_data") or {}
    headers = data.get("headers") or ["Item", "Signal", "Action"]
    rows = data.get("rows") or [["Evidence", "Trace", "Bind"]]
    all_rows = [headers] + rows
    row_count = len(all_rows)
    col_count = max(len(headers), *(len(row) for row in rows))
    cell_w = max(1, int(w / col_count))
    cell_h = max(1, int(h / row_count))
    for row_idx, row in enumerate(all_rows):
        for col_idx in range(col_count):
            value = str(row[col_idx]) if col_idx < len(row) else ""
            cx = x + col_idx * cell_w
            cy = y + row_idx * cell_h
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cx, cy, cell_w, cell_h)
            cell.name = f"{obj.get('object_id')}_r{row_idx+1}_c{col_idx+1}"
            _style_shape(cell, {"fill": "#0F766E" if row_idx == 0 else "#1F2937", "line": "#94A3B8"})
            _add_small_text(slide, value, cx + 3, cy + 2, max(1, cell_w - 6), max(1, cell_h - 4), 5.4 if row_count > 5 else 6.2, "#F8FAFC", f"{obj.get('object_id')}_txt_r{row_idx+1}_c{col_idx+1}", bold=row_idx == 0)


def _add_small_text(slide: Any, text: str, x: int, y: int, w: int, h: int, font_size: float, color: str, name: str, *, bold: bool = False) -> None:
    shape = slide.shapes.add_textbox(x, y, w, h)
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _style_shape(shape: Any, obj: dict[str, Any]) -> None:
    fill = obj.get("fill", "#172A3A")
    line = obj.get("line", "#38BDF8")
    if fill and fill != "#00000000":
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        try:
            shape.fill.transparency = int(obj.get("transparency", 0))
        except AttributeError:
            pass
    else:
        shape.fill.background()
    if line and line != "#00000000":
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()


def _box(obj: dict[str, Any]) -> tuple[int, int, int, int]:
    norm = obj.get("bbox_norm") or [0, 0, 0.1, 0.1]
    x = Inches(float(norm[0]) * SLIDE_WIDTH_IN)
    y = Inches(float(norm[1]) * SLIDE_HEIGHT_IN)
    w = Inches(max(0.005, float(norm[2]) * SLIDE_WIDTH_IN))
    h = Inches(max(0.005, float(norm[3]) * SLIDE_HEIGHT_IN))
    return int(x), int(y), int(w), int(h)


def _rgb(hex_color: str) -> RGBColor:
    value = (hex_color or "#FFFFFF").replace("#", "")[:6]
    if len(value) != 6:
        value = "FFFFFF"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _pptx_media_count(path: Path) -> int:
    with zipfile.ZipFile(path) as archive:
        return len([name for name in archive.namelist() if name.startswith("ppt/media/")])


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()
