"""Compile an accepted E03 archetype set into one editable template pack PPTX."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt


SLIDE_W_IN = 16.0
SLIDE_H_IN = 9.0
SLIDE_W_PX = 1672
SLIDE_H_PX = 941


def compile_editable_template_pack(accepted_archetypes: list[dict[str, Any]], output_dir: Path) -> tuple[dict[str, Any], dict[str, Any]]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    for item in accepted_archetypes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_pack_slide(slide, item["archetype_id"], item["object_graph"].get("nodes", []))
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "editable_template_pack.pptx"
    prs.save(pptx_path)
    spec = {
        "schema_name": "editable_template_pack_spec",
        "slide_count": len(accepted_archetypes),
        "accepted_archetypes": [item["archetype_id"] for item in accepted_archetypes],
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "canva_parity_claimed": False,
    }
    report = {"schema_name": "e03_template_pack_compile_report", "status": "passed" if pptx_path.exists() else "failed", "pptx_path": pptx_path.as_posix(), "slide_count": len(accepted_archetypes), "canva_parity_claimed": False}
    return spec, report


def render_template_pack_contact_sheet(accepted_archetypes: list[dict[str, Any]], output_path: Path) -> Path:
    thumbs = []
    for item in accepted_archetypes:
        image = Image.new("RGB", (420, 236), "#061526")
        draw = ImageDraw.Draw(image, "RGBA")
        draw.text((16, 14), item["archetype_id"], fill=(248, 250, 252, 255))
        for node in item["object_graph"].get("nodes", []):
            bbox = node.get("bbox_norm") or {}
            x = round(float(bbox.get("x", 0)) * 420)
            y = round(float(bbox.get("y", 0)) * 236)
            w = round(float(bbox.get("w", 0)) * 420)
            h = round(float(bbox.get("h", 0)) * 236)
            role = node.get("semantic_role")
            if role in {"primary_chart", "table_region", "comparison_matrix"}:
                draw.rectangle((x, y, x + w, y + h), outline=(244, 180, 63, 180), width=2)
            elif role != "background_base":
                draw.rectangle((x, y, x + w, y + h), outline=(45, 212, 255, 120), width=1)
        thumbs.append(image)
    cols = 4
    rows = max(1, (len(thumbs) + cols - 1) // cols)
    sheet = Image.new("RGB", (cols * 420, rows * 236), "#061526")
    for index, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((index % cols) * 420, (index // cols) * 236))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    return output_path


def _add_pack_slide(slide: Any, archetype_id: str, nodes: list[dict[str, Any]]) -> None:
    _background(slide)
    for node in sorted(nodes, key=lambda item: item.get("z_order", 0)):
        role = node.get("semantic_role")
        if role == "background_base":
            continue
        x, y, w, h = _inches(node.get("bbox_norm") or {"x": 0.06, "y": 0.08, "w": 0.4, "h": 0.1})
        if role == "primary_chart":
            _chart(slide, node["object_id"], x, y, w, h)
        elif role in {"table_region", "comparison_matrix"}:
            _table(slide, node["object_id"], x, y, w, h)
        elif str(role).endswith("text_region") or role == "source_footer_text":
            _textbox(slide, node["object_id"], x, y, w, h, _placeholder(role), size=30 if role == "title_text_region" else 12)
        elif role == "hero_visual_field":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.name = node["object_id"]
            _shape_style(shape, "0C313D", "2DD4FF")
        elif role in {"connector_line", "timeline_axis", "phase_rail", "progress_indicator"}:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y + h / 2), Inches(x + w), Inches(y + h / 2))
            line.name = node["object_id"]
            line.line.color.rgb = RGBColor.from_string("2DD4FF")
            line.line.width = Pt(1.1)
        else:
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.name = node.get("object_id", role)
            _shape_style(shape, "0C313D", "2DD4FF")
    _textbox(slide, f"{archetype_id}_pack_label", 0.04, 0.02, 0.30, 0.04, archetype_id, size=8)


def _background(slide: Any) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    shape.name = "background_base"
    _shape_style(shape, "061526", "061526")
    shape.line.fill.background()


def _textbox(slide: Any, name: str, x: float, y: float, w: float, h: float, text: str, *, size: int) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string("F8FAFC")


def _chart(slide: Any, name: str, x: float, y: float, w: float, h: float) -> None:
    data = ChartData()
    data.categories = ["A", "B", "C", "D"]
    data.add_series("Series", (32, 46, 38, 58))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    slide.shapes[-1].name = name


def _table(slide: Any, name: str, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_table(4, 4, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    for row in range(4):
        for col in range(4):
            shape.table.cell(row, col).text = "Slot"


def _shape_style(shape: Any, fill: str, line: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(1.0)


def _inches(bbox: dict[str, float]) -> tuple[float, float, float, float]:
    return bbox["x"] * SLIDE_W_IN, bbox["y"] * SLIDE_H_IN, bbox["w"] * SLIDE_W_IN, bbox["h"] * SLIDE_H_IN


def _placeholder(role: str) -> str:
    return "TITLE PLACEHOLDER" if role == "title_text_region" else "Editable slot"
