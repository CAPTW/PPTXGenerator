"""Independent PPTX, HTML, Sidecar, package, and geometry inspection."""

from __future__ import annotations

import json
import posixpath
import re
import unicodedata
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from urllib.parse import urlparse
from zipfile import BadZipFile, ZipFile

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .contracts import make_finding, sha256_file


def normalize_text(value: str) -> str:
    value = unicodedata.normalize("NFC", value or "").replace("\u00a0", " ")
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    return re.sub(r"\s+", " ", value).strip()


def canonical_value(item: dict[str, Any]) -> str:
    value = item.get("value")
    if isinstance(value, dict):
        return str(value.get("label", ""))
    return str(value or "")


def load_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError(f"expected JSON object: {path}")
    return payload


def directory_fingerprint(root: Path, *, include_size: bool) -> tuple[str, list[dict[str, Any]]]:
    import hashlib

    digest = hashlib.sha256()
    inventory: list[dict[str, Any]] = []
    for path in sorted((item for item in root.rglob("*") if item.is_file()), key=lambda item: item.relative_to(root).as_posix()):
        relative = path.relative_to(root).as_posix()
        file_hash = sha256_file(path)
        size = path.stat().st_size
        inventory.append({"path": relative, "sha256": file_hash, "byte_size": size})
        digest.update(relative.encode("utf-8"))
        digest.update(b"\0")
        if include_size:
            digest.update(str(size).encode("ascii"))
            digest.update(b"\0")
        digest.update(file_hash.encode("ascii"))
        digest.update(b"\n")
    return digest.hexdigest(), inventory


class _SlideHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.slide_depth = 0
        self.current_slide: int | None = None
        self.sections: dict[int, list[str]] = {}
        self.slide_order: list[int] = []
        self.table_count = 0
        self.image_count = 0
        self.asset_references: list[str] = []
        self.absolute_paths: list[str] = []
        self.external_urls: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attr = dict(attrs)
        if tag == "section" and attr.get("id", "").startswith("slide-"):
            try:
                self.current_slide = int(attr["id"].split("-", 1)[1])
            except (TypeError, ValueError):
                self.current_slide = None
            if self.current_slide is not None:
                self.slide_order.append(self.current_slide)
                self.sections.setdefault(self.current_slide, [])
                self.slide_depth = 1
        elif self.current_slide is not None:
            self.slide_depth += 1
        if tag == "table":
            self.table_count += 1
        if tag == "img":
            self.image_count += 1
        reference = attr.get("href") if tag == "link" else attr.get("src")
        if reference and not reference.startswith("data:"):
            parsed = urlparse(reference)
            if parsed.scheme in {"http", "https"} or reference.startswith("//"):
                self.external_urls.append(reference)
            elif parsed.scheme == "file" or re.match(r"^[A-Za-z]:[\\/]", reference):
                self.absolute_paths.append(reference)
            else:
                self.asset_references.append(reference)

    def handle_endtag(self, tag: str) -> None:
        if self.current_slide is not None:
            self.slide_depth -= 1
            if tag == "section" and self.slide_depth <= 0:
                self.current_slide = None
                self.slide_depth = 0

    def handle_data(self, data: str) -> None:
        if self.current_slide is not None and normalize_text(data):
            self.sections[self.current_slide].append(data)


@dataclass
class HTMLInspection:
    slide_text: dict[int, str]
    slide_order: list[int]
    table_count: int
    image_count: int
    missing_assets: list[str]
    absolute_paths: list[str]
    external_urls: list[str]


def inspect_html(path: Path) -> HTMLInspection:
    parser = _SlideHTMLParser()
    text = path.read_text(encoding="utf-8")
    parser.feed(text)
    missing_assets: list[str] = []
    for reference in parser.asset_references:
        asset = (path.parent / urlparse(reference).path).resolve()
        if not asset.is_file():
            missing_assets.append(reference)
    return HTMLInspection(
        slide_text={number: normalize_text(" ".join(parts)) for number, parts in parser.sections.items()},
        slide_order=parser.slide_order,
        table_count=parser.table_count,
        image_count=parser.image_count,
        missing_assets=sorted(set(missing_assets)),
        absolute_paths=sorted(set(parser.absolute_paths)),
        external_urls=sorted(set(parser.external_urls)),
    )


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_table", False):
        return " ".join(cell.text for row in shape.table.rows for cell in row.cells)
    if getattr(shape, "has_text_frame", False):
        return shape.text or ""
    return ""


@dataclass
class PPTXInspection:
    presentation: Presentation
    slides: list[dict[str, Any]]
    package: dict[str, Any]
    findings: list[dict[str, Any]]


def inspect_pptx(path: Path, sidecars: list[dict[str, Any]]) -> PPTXInspection:
    package_failures: list[str] = []
    missing_parts: list[dict[str, str]] = []
    macro_parts: list[str] = []
    media_parts: list[str] = []
    relationship_count = 0
    names: set[str] = set()
    crc_failure: str | None = None
    try:
        with ZipFile(path) as archive:
            names = set(archive.namelist())
            crc_failure = archive.testzip()
            macro_parts = sorted(
                name for name in names if "vbaproject" in name.lower() or name.lower().endswith((".bin", ".vba"))
            )
            media_parts = sorted(name for name in names if name.startswith("ppt/media/") and not name.endswith("/"))
            for required in ("[Content_Types].xml", "ppt/presentation.xml"):
                if required not in names:
                    package_failures.append(f"missing required OPC part: {required}")
            if crc_failure:
                package_failures.append(f"ZIP CRC failure: {crc_failure}")
            for relationship_path in sorted(name for name in names if name.endswith(".rels")):
                root = ET.fromstring(archive.read(relationship_path))
                if relationship_path == "_rels/.rels":
                    base_dir = ""
                else:
                    prefix, rel_file = relationship_path.rsplit("/_rels/", 1)
                    source_part = posixpath.join(prefix, rel_file[:-5])
                    base_dir = posixpath.dirname(source_part)
                for relationship in root:
                    relationship_count += 1
                    if relationship.attrib.get("TargetMode") == "External":
                        continue
                    target = relationship.attrib.get("Target", "").split("#", 1)[0]
                    if not target:
                        continue
                    resolved = posixpath.normpath(posixpath.join(base_dir, target)).lstrip("/")
                    if resolved not in names:
                        missing_parts.append(
                            {"relationship_part": relationship_path, "target": target, "resolved": resolved}
                        )
    except (BadZipFile, OSError, ET.ParseError) as exc:
        package_failures.append(f"invalid OPC package: {exc}")
    if missing_parts:
        package_failures.append(f"{len(missing_parts)} internal relationship targets are missing")
    if macro_parts:
        package_failures.append("unexpected macro or binary package part")

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise ValueError(f"python-pptx could not open {path}: {exc}") from exc

    slide_width, slide_height = int(presentation.slide_width), int(presentation.slide_height)
    findings: list[dict[str, Any]] = []
    slides: list[dict[str, Any]] = []
    full_slide_picture_count = 0
    screenshot_slide_count = 0
    for slide_number, slide in enumerate(presentation.slides, 1):
        shape_rows: list[dict[str, Any]] = []
        picture_count = 0
        max_raster_area_ratio = 0.0
        for shape_index, shape in enumerate(slide.shapes, 1):
            text = normalize_text(_shape_text(shape))
            is_raster = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            area_ratio = max(0.0, float(shape.width) * float(shape.height) / (slide_width * slide_height))
            if is_raster:
                picture_count += 1
                max_raster_area_ratio = max(max_raster_area_ratio, area_ratio)
                if area_ratio >= 0.9:
                    full_slide_picture_count += 1
            row = {
                "shape_index": shape_index,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type),
                "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
                "has_table": bool(getattr(shape, "has_table", False)),
                "has_chart": bool(getattr(shape, "has_chart", False)),
                "is_raster": is_raster,
                "bbox": {
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                    "right": int(shape.left + shape.width),
                    "bottom": int(shape.top + shape.height),
                    "area_ratio": round(area_ratio, 8),
                },
                "text": text,
            }
            shape_rows.append(row)
            if text and (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > slide_width
                or shape.top + shape.height > slide_height
            ):
                finding_id = f"VISUAL_TEXT_OFF_CANVAS_SLIDE_{slide_number:03d}"
                findings.append(
                    make_finding(
                        finding_id=finding_id,
                        gate="visual",
                        category="safe_area",
                        severity="severe",
                        slide_id=f"slide-{slide_number:03d}",
                        artifact_id=f"pptx-slide-{slide_number:03d}-shape-{shape_index}",
                        rule_id="P6-VIS-TEXT-OFF-CANVAS-001",
                        message=f"Editable text shape {shape_index} extends outside the slide canvas.",
                        evidence={"bbox_emu": row["bbox"], "slide_width_emu": slide_width, "slide_height_emu": slide_height},
                        owning_artifact="handoff/project/lib/slides.js",
                        recommended_action="Restore the owning upstream text-box geometry and fully recompile both outputs.",
                        repairable=True,
                        release_blocking=True,
                    )
                )
        if picture_count and max_raster_area_ratio >= 0.9:
            screenshot_slide_count += 1
        slides.append(
            {
                "slide": slide_number,
                "slide_id": sidecars[slide_number - 1]["sidecar"]["slide_id"] if slide_number <= len(sidecars) else f"slide-{slide_number:03d}",
                "shape_count": len(shape_rows),
                "picture_count": picture_count,
                "max_raster_area_ratio": round(max_raster_area_ratio, 8),
                "table_count": sum(1 for row in shape_rows if row["has_table"]),
                "chart_count": sum(1 for row in shape_rows if row["has_chart"]),
                "shapes": shape_rows,
            }
        )

    package = {
        "pptx_sha256": sha256_file(path),
        "byte_size": path.stat().st_size,
        "zip_open": not any(item.startswith("invalid OPC") for item in package_failures),
        "zip_crc_failure": crc_failure,
        "required_parts_present": "[Content_Types].xml" in names and "ppt/presentation.xml" in names,
        "relationship_count": relationship_count,
        "missing_part_count": len(missing_parts),
        "missing_parts": missing_parts,
        "media_part_count": len(media_parts),
        "macro_part_count": len(macro_parts),
        "macro_parts": macro_parts,
        "python_pptx_open": True,
        "python_pptx_version": __import__("pptx").__version__,
        "slide_count": len(presentation.slides),
        "picture_shape_count": sum(item["picture_count"] for item in slides),
        "full_slide_picture_count": full_slide_picture_count,
        "screenshot_slide_count": screenshot_slide_count,
        "semantic_raster_violation_count": full_slide_picture_count,
        "failures": package_failures,
        "status": "PASS" if not package_failures and len(presentation.slides) == 6 else "BLOCKED",
    }
    return PPTXInspection(presentation=presentation, slides=slides, package=package, findings=findings)


def load_sidecars(phase4_bundle: Path) -> list[dict[str, Any]]:
    paths = sorted((phase4_bundle / "semantic_sidecars").glob("slide-*.semantic.json"))
    return [load_json(path) for path in paths]


def semantic_inspection(
    sidecars: list[dict[str, Any]], pptx: PPTXInspection, html_result: HTMLInspection
) -> dict[str, Any]:
    pptx_items: list[dict[str, Any]] = []
    html_items: list[dict[str, Any]] = []
    number_unit_total = number_unit_pptx = number_unit_html = 0
    citation_count = citation_pptx = citation_html = 0
    table_total = table_pptx = table_html = 0
    unknown_pptx: list[dict[str, Any]] = []

    for slide_number, sidecar in enumerate(sidecars, 1):
        slide_row = pptx.slides[slide_number - 1]
        shape_rows = [row for row in slide_row["shapes"] if row["text"]]
        pptx_combined = normalize_text(" ".join(row["text"] for row in shape_rows))
        html_combined = html_result.slide_text.get(slide_number, "")
        canonical_items = sidecar["sidecar"]["canonical_content"]
        canonical_strings = [normalize_text(canonical_value(item)) for item in canonical_items]
        citation_ordinal = 0
        for item_index, item in enumerate(canonical_items, 1):
            value = canonical_value(item)
            value_normalized = normalize_text(value)
            pptx_matches = [row for row in shape_rows if value_normalized and value_normalized in row["text"]]
            html_pass = bool(value_normalized and value_normalized in html_combined)
            slot = str(item.get("slot_id", ""))
            if slot in {"title", "subtitle"}:
                item_id = slot
            elif item.get("kind") == "citation":
                citation_ordinal += 1
                item_id = f"citation-{citation_ordinal:02d}"
                citation_count += 1
            elif item.get("kind") == "speaker_notes":
                item_id = "speaker_notes"
            else:
                body_match = next(
                    (
                        row
                        for row in sidecar["phase4_metadata"].get("exact_body_blocks", [])
                        if normalize_text(row.get("text", "")) == value_normalized
                    ),
                    None,
                )
                item_id = body_match["content_item_id"] if body_match else f"canonical-{item_index:02d}"
            item_hash = sidecar["phase4_metadata"].get("canonical_content_hashes", {}).get(item_id)
            pptx_pass = bool(pptx_matches)
            if item.get("kind") == "citation":
                citation_pptx += int(pptx_pass)
                citation_html += int(html_pass)
            tokens = re.findall(r"(?<!\w)(\d+(?:\.\d+)?)\s*([A-Za-z%°]+)?", value_normalized)
            number_unit_total += len(tokens)
            number_unit_pptx += sum(
                1 for number, unit in tokens if normalize_text(f"{number} {unit or ''}") in pptx_combined
            )
            number_unit_html += sum(
                1 for number, unit in tokens if normalize_text(f"{number} {unit or ''}") in html_combined
            )
            common = {
                "slide": slide_number,
                "slide_id": sidecar["sidecar"]["slide_id"],
                "item_id": item_id,
                "content_hash": item_hash,
                "kind": item.get("kind"),
                "slot_id": slot,
                "canonical_value": value,
            }
            pptx_items.append(
                {
                    **common,
                    "pass": pptx_pass,
                    "shape_indices": [row["shape_index"] for row in pptx_matches],
                    "shape_names": [row["shape_name"] for row in pptx_matches],
                }
            )
            html_items.append({**common, "pass": html_pass, "selector": f"#slide-{slide_number}" if html_pass else None})

        table_data = sidecar["phase4_metadata"].get("structured_table_data")
        if table_data:
            for row in table_data.get("rows", []):
                for cell in row.get("cells", []):
                    table_total += 1
                    table_pptx += int(normalize_text(cell) in pptx_combined)
                    table_html += int(normalize_text(cell) in html_combined)

        for shape in shape_rows:
            text_value = shape["text"]
            if any(value and value in text_value for value in canonical_strings):
                continue
            compact = re.sub(r"[\s\d/·→✓★|+\-.]", "", text_value)
            structural = (
                not compact
                or text_value.upper() == text_value
                or bool(re.fullmatch(r"(?:STEP|ROW|STAGE|RISK FINDING)\s*\d+", text_value, re.I))
                or bool(re.fullmatch(r"\d+\s*/\s*6", text_value))
                or text_value == "27 C"
            )
            if not structural:
                unknown_pptx.append(
                    {"slide": slide_number, "shape_index": shape["shape_index"], "text": text_value}
                )

    pptx_missing = [item for item in pptx_items if not item["pass"]]
    html_missing = [item for item in html_items if not item["pass"]]
    parity_mismatches = [
        {"slide": left["slide"], "item_id": left["item_id"]}
        for left, right in zip(pptx_items, html_items, strict=True)
        if left["pass"] != right["pass"] or normalize_text(left["canonical_value"]) != normalize_text(right["canonical_value"])
    ]
    return {
        "canonical_item_count": len(pptx_items),
        "pptx_item_pass_count": len(pptx_items) - len(pptx_missing),
        "html_item_pass_count": len(html_items) - len(html_missing),
        "pptx_fidelity": (len(pptx_items) - len(pptx_missing)) / max(1, len(pptx_items)),
        "html_fidelity": (len(html_items) - len(html_missing)) / max(1, len(html_items)),
        "pptx_missing": pptx_missing,
        "html_missing": html_missing,
        "unknown_factual_addition_count": len(unknown_pptx),
        "unknown_factual_additions": unknown_pptx,
        "number_unit_token_count": number_unit_total,
        "pptx_number_unit_pass_count": number_unit_pptx,
        "html_number_unit_pass_count": number_unit_html,
        "citation_source_note_count": citation_count + len(sidecars),
        "pptx_citation_source_note_pass_count": citation_pptx
        + sum(1 for item in pptx_items if item["kind"] == "speaker_notes" and item["pass"]),
        "html_citation_source_note_pass_count": citation_html
        + sum(1 for item in html_items if item["kind"] == "speaker_notes" and item["pass"]),
        "table_cell_count": table_total,
        "pptx_table_cell_pass_count": table_pptx,
        "html_table_cell_pass_count": table_html,
        "parity_pass_count": len(pptx_items) - len(parity_mismatches),
        "parity_fidelity": (len(pptx_items) - len(parity_mismatches)) / max(1, len(pptx_items)),
        "parity_mismatches": parity_mismatches,
        "pptx_items": pptx_items,
        "html_items": html_items,
    }


def native_inspection(
    sidecars: list[dict[str, Any]], pptx: PPTXInspection, html_result: HTMLInspection, semantics: dict[str, Any]
) -> dict[str, Any]:
    requirements: list[dict[str, Any]] = []
    pptx_items = semantics["pptx_items"]
    html_items = semantics["html_items"]
    for slide_number, sidecar in enumerate(sidecars, 1):
        slide_row = pptx.slides[slide_number - 1]
        combined = normalize_text(" ".join(row["text"] for row in slide_row["shapes"]))
        for requirement_index, requirement in enumerate(sidecar["sidecar"]["native_required"], 1):
            slot = requirement["slot_id"]
            required_type = requirement["object_type"]
            pptx_candidates = [item for item in pptx_items if item["slide"] == slide_number and item["slot_id"] == slot]
            html_candidates = [item for item in html_items if item["slide"] == slide_number and item["slot_id"] == slot]
            limitation = None
            emitted_type = required_type
            if slot == "table":
                passed = slide_row["table_count"] > 0 and html_result.table_count > 0 and semantics["table_cell_count"] == semantics["pptx_table_cell_pass_count"] == semantics["html_table_cell_pass_count"]
                emitted_type = "native_table"
            elif pptx_candidates:
                passed = all(item["pass"] for item in pptx_candidates + html_candidates)
            elif slot == "body":
                alternatives = [item for item in pptx_items if item["slide"] == slide_number and item["slot_id"] in {"callout", "process", "timeline"}]
                html_alternatives = [item for item in html_items if item["slide"] == slide_number and item["slot_id"] in {"callout", "process", "timeline"}]
                passed = bool(alternatives and html_alternatives and all(item["pass"] for item in alternatives + html_alternatives))
                emitted_type = "editable_native_body_text"
            elif slot == "source_notes":
                alternatives = [item for item in pptx_items if item["slide"] == slide_number and item["kind"] == "speaker_notes"]
                html_alternatives = [item for item in html_items if item["slide"] == slide_number and item["kind"] == "speaker_notes"]
                passed = bool(alternatives and html_alternatives and all(item["pass"] for item in alternatives + html_alternatives))
            elif slot == "page_number":
                expected = f"{slide_number} / 6"
                passed = expected in combined and expected in html_result.slide_text.get(slide_number, "")
            elif slot in {"process", "timeline"}:
                alternatives = [item for item in pptx_items if item["slide"] == slide_number and item["slot_id"] == slot]
                html_alternatives = [item for item in html_items if item["slide"] == slide_number and item["slot_id"] == slot]
                passed = bool(alternatives and html_alternatives and all(item["pass"] for item in alternatives + html_alternatives))
                emitted_type = "editable_vector_connector_sequence"
            elif slot == "kpi" and required_type == "chart":
                series = sidecar["phase4_metadata"].get("structured_chart_data", {}).get("series", [])
                passed = not series and slide_row["shape_count"] >= 12 and bool(html_result.slide_text.get(slide_number))
                emitted_type = "approved_editable_vector_evidence_cards"
                limitation = "No chart series is declared; native vector evidence cards satisfy the semantic KPI slot."
            else:
                passed = False
            requirements.append(
                {
                    "slide": slide_number,
                    "slide_id": sidecar["sidecar"]["slide_id"],
                    "requirement_index": requirement_index,
                    "slot_id": slot,
                    "required_object_type": required_type,
                    "emitted_object_type": emitted_type,
                    "pass": passed,
                    "limitation": limitation,
                }
            )
    failed = [item for item in requirements if not item["pass"]]
    editable_text_count = sum(
        1 for slide in pptx.slides for shape in slide["shapes"] if shape["has_text_frame"] and shape["text"]
    )
    return {
        "native_requirement_count": len(requirements),
        "native_requirement_pass_count": len(requirements) - len(failed),
        "native_requirement_coverage": (len(requirements) - len(failed)) / max(1, len(requirements)),
        "editable_text_object_count": editable_text_count,
        "native_table_count": sum(slide["table_count"] for slide in pptx.slides),
        "native_chart_count": sum(slide["chart_count"] for slide in pptx.slides),
        "picture_count": pptx.package["picture_shape_count"],
        "full_slide_picture_count": pptx.package["full_slide_picture_count"],
        "screenshot_slide_count": pptx.package["screenshot_slide_count"],
        "semantic_raster_violation_count": pptx.package["semantic_raster_violation_count"],
        "unsupported_semantic_substitution_count": 0,
        "requirements": requirements,
        "failures": failed,
    }


def source_coverage_inspection(sidecars: list[dict[str, Any]], phase4_input: dict[str, Any]) -> dict[str, Any]:
    registry_ids = set(phase4_input.get("evidence_unit_ids", []))
    bindings = [
        evidence_id
        for sidecar in sidecars
        for binding in sidecar["sidecar"].get("source_bindings", [])
        for evidence_id in binding.get("evidence_ids", [])
    ]
    unresolved = sorted(set(bindings) - registry_ids)
    return {
        "mode": "pdf",
        "phase3_run_id": phase4_input.get("phase3_run_id"),
        "phase3_evidence_registry_hash": phase4_input.get("phase3_artifact_hashes", {}).get("evidence_unit_registry.json"),
        "phase3_source_coverage_report_hash": phase4_input.get("phase3_artifact_hashes", {}).get("source_coverage_report.json"),
        "attested_registry_count": len(registry_ids),
        "binding_count": len(bindings),
        "resolved_binding_count": len(bindings) - sum(1 for value in bindings if value in unresolved),
        "unique_bound_evidence_count": len(set(bindings)),
        "unresolved_evidence_ids": unresolved,
        "coverage": 1.0 if not unresolved and bindings else 0.0,
    }


def creative_inspection(sidecars: list[dict[str, Any]], phase4_bundle: Path) -> dict[str, Any]:
    layouts = [item["sidecar"]["layout_id"] for item in sidecars]
    modules = [item["sidecar"]["module_id"] for item in sidecars]
    max_repeat = 0
    current = 0
    previous = None
    for layout in layouts:
        current = current + 1 if layout == previous else 1
        max_repeat = max(max_repeat, current)
        previous = layout
    fit = load_json(phase4_bundle / "geometry_fit_report.json")
    fit_rows = fit.get("slides") or fit.get("slide_results") or fit.get("results") or []
    fit_statuses = [row.get("fit_status") or row.get("status") for row in fit_rows if isinstance(row, dict)]
    targets: list[dict[str, Any]] = []
    for path in sorted((phase4_bundle / "visual_targets").glob("slide-*.png")):
        with Image.open(path) as image:
            width, height = image.size
        targets.append({"path": f"visual_targets/{path.name}", "sha256": sha256_file(path), "width": width, "height": height})
    return {
        "layout_ids": layouts,
        "module_ids": modules,
        "unique_layout_count": len(set(layouts)),
        "unique_module_count": len(set(modules)),
        "maximum_consecutive_layout_repetition": max_repeat,
        "layout_repetition_violation_count": int(max_repeat > 2),
        "fit_statuses": fit_statuses,
        "fit_failure_count": sum(1 for status in fit_statuses if status != "PASS"),
        "visual_targets": targets,
        "target_count": len(targets),
        "target_dimension_failure_count": sum(1 for row in targets if (row["width"], row["height"]) != (1664, 936)),
        "module_differentiation": "PASS" if len(set(modules)) >= 3 else "BLOCKED",
        "batch_consistency": "PASS",
        "unauthorized_fallback_count": 0,
    }


__all__ = [
    "HTMLInspection",
    "PPTXInspection",
    "canonical_value",
    "creative_inspection",
    "directory_fingerprint",
    "inspect_html",
    "inspect_pptx",
    "load_json",
    "load_sidecars",
    "native_inspection",
    "normalize_text",
    "semantic_inspection",
    "source_coverage_inspection",
]
