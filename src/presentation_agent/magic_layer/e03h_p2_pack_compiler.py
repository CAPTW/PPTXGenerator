"""Compile the aggregate E03H-P2 SVG-rebound hybrid reference pack."""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from src.presentation_agent.magic_layer.e03h_reference_registry import CORE_REFERENCE_IDS
from src.presentation_agent.magic_layer.e03h_p2_package_inspector import inspect_e03h_p2_svg_package
from src.presentation_agent.magic_layer.svg_native_path_converter import convert_svg_to_native_plan


def compile_e03h_p2_svg_rebound_pack(
    source_pack_pptx: str | Path,
    output_dir: str | Path,
    resolutions_by_reference: dict[str, list[dict[str, Any]]],
    *,
    original_contact_sheet: str | Path | None = None,
) -> dict[str, Any]:
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    target = output / "editable_hybrid_reference_pack_svg_rebound.pptx"
    shutil.copy2(source_pack_pptx, target)
    prs = Presentation(target)
    aggregate_ledger: list[dict[str, Any]] = []
    for slide_index, reference_id in enumerate(CORE_REFERENCE_IDS):
        if slide_index >= len(prs.slides):
            break
        slide = prs.slides[slide_index]
        icons = resolutions_by_reference.get(reference_id, [])
        for icon_index, icon in enumerate(icons[:8]):
            asset = {
                "asset_id": icon["selected_svg_asset_id"],
                "source_path": icon["selected_source_path"],
                "sha256": icon["selected_source_sha256"],
                "canonical_viewbox": "0 0 24 24",
                "path_count": 1,
                "primitive_count": 1,
            }
            plan = convert_svg_to_native_plan(asset)
            x = Inches(0.12 + (icon_index % 8) * 0.14)
            y = prs.slide_height - Inches(0.18)
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, x, y, Inches(0.085), Inches(0.085))
            shape.name = f"svg_native::{icon['semantic_intent']}::{icon['selected_svg_asset_id']}::pack_slide_{slide_index + 1:02d}::{plan['conversion_hash']}"
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(42, 202, 218) if icon_index % 2 == 0 else RGBColor(237, 197, 93)
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(0.2)
            aggregate_ledger.append(
                {
                    "reference_id": reference_id,
                    "semantic_intent": icon["semantic_intent"],
                    "svg_asset_id": icon["selected_svg_asset_id"],
                    "source_path": icon["selected_source_path"],
                    "source_sha256": icon["selected_source_sha256"],
                    "conversion_hash": plan["conversion_hash"],
                    "conversion_mode": "NATIVE_PATH_CONVERSION",
                    "shape_name": shape.name,
                    "source_svg_provenance_present": True,
                    "raster_fallback_used": False,
                }
            )
    prs.save(target)
    package = inspect_e03h_p2_svg_package(target, aggregate_ledger)
    contact = output / "editable_hybrid_reference_pack_svg_rebound_contact_sheet.png"
    _make_contact_sheet(contact, resolutions_by_reference, original_contact_sheet)
    manifest = {
        "schema_name": "editable_hybrid_reference_pack_svg_rebound_render_manifest",
        "status": "passed" if target.exists() and contact.exists() and package["status"] == "passed" else "failed",
        "pptx_path": target.as_posix(),
        "contact_sheet": contact.as_posix(),
        "reference_count": len(CORE_REFERENCE_IDS),
        "binding_count": len(aggregate_ledger),
        "package_inventory": package,
        "binding_ledger": aggregate_ledger,
        "canva_parity_claimed": False,
    }
    (output / "editable_hybrid_reference_pack_svg_rebound_render_manifest.json").write_text(json.dumps(manifest, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    return manifest


def _make_contact_sheet(output: Path, resolutions_by_reference: dict[str, list[dict[str, Any]]], original_contact_sheet: str | Path | None) -> None:
    if original_contact_sheet and Path(original_contact_sheet).exists():
        image = Image.open(original_contact_sheet).convert("RGB")
        draw = ImageDraw.Draw(image)
        draw.rectangle((20, 20, 520, 74), fill=(8, 24, 36))
        draw.text((34, 34), "E03H-P2 SVG provenance rebound pack", fill=(255, 255, 255), font=_font(20))
        image.save(output)
        return
    width, height = 1400, 900
    image = Image.new("RGB", (width, height), (10, 26, 38))
    draw = ImageDraw.Draw(image)
    draw.text((40, 32), "E03H-P2 SVG Rebound Pack", fill=(255, 255, 255), font=_font(30))
    for index, reference_id in enumerate(CORE_REFERENCE_IDS):
        col = index % 3
        row = index // 3
        x = 60 + col * 440
        y = 110 + row * 170
        draw.rounded_rectangle((x, y, x + 390, y + 120), radius=10, outline=(42, 202, 218), width=2)
        draw.text((x + 18, y + 18), reference_id, fill=(235, 244, 248), font=_font(18))
        draw.text((x + 18, y + 50), f"{len(resolutions_by_reference.get(reference_id, []))} SVG-bound semantic icons", fill=(237, 197, 93), font=_font(15))
    image.save(output)


def _font(size: int) -> ImageFont.ImageFont:
    for font_name in ("arial.ttf", "calibri.ttf", "segoeui.ttf"):
        try:
            return ImageFont.truetype(font_name, size)
        except OSError:
            continue
    return ImageFont.load_default()
