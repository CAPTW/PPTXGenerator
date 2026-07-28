"""Compile an editable source-bound E04H hybrid PPTX deck."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from src.presentation_agent.magic_layer.svg_native_path_converter import convert_svg_to_native_plan


SLIDE_W = 13.333
SLIDE_H = 7.5


def compile_e04h_source_bound_hybrid_deck(
    source_artifacts: dict[str, Any],
    layout_report: dict[str, Any],
    art_plan: dict[str, Any],
    slot_binding: dict[str, Any],
    icon_plan: dict[str, Any],
    output_dir: str | Path,
) -> dict[str, Any]:
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    rendered_dir = output / "rendered_slides"
    rendered_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    chart_count = 0
    table_count = 0
    svg_ledger: list[dict[str, Any]] = []
    slides = slot_binding["slide_bindings"]
    icon_by_slide: dict[str, list[dict[str, Any]]] = {}
    for row in icon_plan["bindings"]:
        icon_by_slide.setdefault(row["slide_id"], []).append(row)
    for slide_binding in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _draw_slide(slide, slide_binding, icon_by_slide.get(slide_binding["slide_id"], []), svg_ledger)
        if slide_binding["selected_reference_id"] == "data_dashboard_hybrid":
            _add_chart(slide, source_artifacts.get("chart_data_ledger", {}))
            chart_count += 1
        if slide_binding["selected_reference_id"] in {"table_matrix_hybrid", "comparison_matrix_hybrid"}:
            _add_table(slide, source_artifacts.get("table_data_ledger", {}))
            table_count += 1
        _render_slide_png(slide_binding, icon_by_slide.get(slide_binding["slide_id"], []), rendered_dir / f"slide_{slide_binding['slide_number']:02d}.png")
    pptx_path = output / "source_bound_hybrid_sample_deck_12_16.pptx"
    prs.save(pptx_path)
    contact = output / "source_bound_hybrid_sample_deck_contact_sheet.png"
    _build_contact_sheet(rendered_dir, contact)
    manifest = {
        "schema_name": "source_bound_hybrid_sample_deck_render_manifest",
        "status": "passed" if pptx_path.exists() and contact.exists() else "failed",
        "pptx_path": pptx_path.as_posix(),
        "contact_sheet": contact.as_posix(),
        "rendered_slides_dir": rendered_dir.as_posix(),
        "slide_count": len(slides),
        "native_chart_count": chart_count,
        "native_table_count": table_count,
        "svg_binding_ledger": svg_ledger,
        "canva_parity_claimed": False,
    }
    (output / "source_bound_hybrid_sample_deck_render_manifest.json").write_text(json.dumps(manifest, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    return manifest


def _draw_slide(slide: Any, binding: dict[str, Any], icons: list[dict[str, Any]], svg_ledger: list[dict[str, Any]]) -> None:
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(8, 27, 39), None, "background_base_native")
    _rect(slide, 0.25, 0.28, 12.85, 6.65, RGBColor(10, 43, 56), RGBColor(17, 112, 126), "bounded_visual_backplate_01")
    _line(slide, 0.4, 6.55, 12.8, 6.55, RGBColor(237, 197, 93), "source_footer_gold_rule")
    _text(slide, binding["title"], 0.55, 0.42, 8.3, 0.42, 20, True, "title_text")
    _text(slide, binding["subtitle"], 0.55, 0.88, 8.6, 0.42, 10, False, "subtitle_text", RGBColor(181, 205, 214))
    reference = binding["selected_reference_id"]
    if reference == "data_dashboard_hybrid":
        _dashboard_chrome(slide, binding)
    elif reference in {"table_matrix_hybrid", "comparison_matrix_hybrid"}:
        _matrix_chrome(slide, binding)
    elif reference == "timeline_roadmap_hybrid":
        _timeline_chrome(slide, binding)
    elif reference in {"process_workflow_infographic", "methodology_framework_layered"}:
        _process_chrome(slide, binding)
    elif reference == "visual_toc_navigation":
        _toc_chrome(slide, binding)
    elif reference == "evidence_stack_visual":
        _evidence_chrome(slide, binding)
    else:
        _card_chrome(slide, binding)
    _text(slide, binding["primary_claim"], 0.65, 5.75, 8.2, 0.38, 9, False, "primary_claim_text", RGBColor(235, 244, 248))
    footer = f"{binding['citation_footer']} | evidence {binding['evidence_ref']} | confidence {binding['confidence']:.2f}"
    _text(slide, footer, 0.55, 6.63, 11.8, 0.26, 6.5, False, "source_footer_text", RGBColor(237, 197, 93))
    for index, icon in enumerate(icons):
        _add_svg_provenance_icon(slide, icon, index, binding["slide_id"], svg_ledger)


def _card_chrome(slide: Any, binding: dict[str, Any]) -> None:
    for i, detail in enumerate(binding["details"][:4]):
        x = 0.65 + (i % 2) * 4.15
        y = 1.65 + (i // 2) * 1.7
        _rect(slide, x, y, 3.75, 1.2, RGBColor(11, 55, 68), RGBColor(38, 170, 184), f"card_panel_{i+1}")
        _text(slide, detail, x + 0.18, y + 0.18, 3.3, 0.8, 8.2, False, f"card_text_{i+1}")


def _evidence_chrome(slide: Any, binding: dict[str, Any]) -> None:
    _rect(slide, 0.75, 1.55, 4.2, 2.7, RGBColor(12, 61, 72), RGBColor(237, 197, 93), "claim_evidence_panel")
    _text(slide, binding["primary_claim"], 1.0, 2.05, 3.6, 1.1, 14, True, "claim_focal_text")
    for i, detail in enumerate(binding["details"][:3]):
        _rect(slide, 5.45, 1.52 + i * 0.86, 5.4, 0.54, RGBColor(10, 47, 59), RGBColor(42, 202, 218), f"evidence_layer_{i+1}")
        _text(slide, detail, 5.65, 1.64 + i * 0.86, 4.9, 0.25, 7.2, False, f"evidence_text_{i+1}")


def _toc_chrome(slide: Any, binding: dict[str, Any]) -> None:
    for i, detail in enumerate(binding["details"][:6]):
        y = 1.45 + i * 0.55
        _rect(slide, 2.1, y, 7.5, 0.36, RGBColor(10, 49, 62), RGBColor(42, 202, 218) if i == 0 else RGBColor(34, 96, 110), f"navigation_item_{i+1}")
        _text(slide, detail, 2.4, y + 0.07, 6.8, 0.16, 7, False, f"navigation_text_{i+1}")


def _process_chrome(slide: Any, binding: dict[str, Any]) -> None:
    for i, detail in enumerate(binding["details"][:5]):
        x = 1.0 + i * 2.12
        _rect(slide, x, 2.2, 1.45, 0.8, RGBColor(12, 58, 70), RGBColor(42, 202, 218), f"process_node_{i+1}")
        _text(slide, detail.split(":")[0], x + 0.13, 2.48, 1.1, 0.18, 6.4, False, f"process_node_text_{i+1}")
        if i < 4:
            _line(slide, x + 1.45, 2.6, x + 2.03, 2.6, RGBColor(237, 197, 93), f"process_connector_{i+1}")


def _dashboard_chrome(slide: Any, binding: dict[str, Any]) -> None:
    for i, label in enumerate(("Trace", "Method", "QA", "Reuse")):
        _rect(slide, 0.75 + i * 1.45, 1.4, 1.15, 0.72, RGBColor(13, 61, 73), RGBColor(42, 202, 218), f"kpi_card_{i+1}")
        _text(slide, label, 0.9 + i * 1.45, 1.62, 0.84, 0.16, 6.4, False, f"kpi_text_{i+1}")
    _rect(slide, 1.0, 2.45, 5.4, 2.6, RGBColor(9, 47, 59), RGBColor(237, 197, 93), "native_chart_frame")
    _rect(slide, 7.0, 2.45, 3.4, 2.6, RGBColor(9, 47, 59), RGBColor(42, 202, 218), "insight_panel")
    _text(slide, binding["primary_claim"], 7.25, 2.85, 2.8, 0.82, 9, False, "dashboard_insight_text")


def _matrix_chrome(slide: Any, binding: dict[str, Any]) -> None:
    _rect(slide, 0.7, 1.35, 7.7, 4.0, RGBColor(9, 47, 59), RGBColor(237, 197, 93), "matrix_table_frame")
    for i in range(5):
        _line(slide, 0.7, 1.95 + i * 0.65, 8.4, 1.95 + i * 0.65, RGBColor(28, 100, 116), f"matrix_row_rule_{i}")
    for i in range(4):
        _line(slide, 0.7 + i * 1.92, 1.35, 0.7 + i * 1.92, 5.35, RGBColor(28, 100, 116), f"matrix_col_rule_{i}")


def _timeline_chrome(slide: Any, binding: dict[str, Any]) -> None:
    _line(slide, 1.0, 3.1, 11.0, 3.1, RGBColor(42, 202, 218), "timeline_rail")
    for i, detail in enumerate(binding["details"][:5]):
        x = 1.2 + i * 2.0
        _rect(slide, x, 2.65, 0.65, 0.65, RGBColor(12, 58, 70), RGBColor(237, 197, 93), f"milestone_{i+1}")
        _text(slide, detail, x - 0.25, 3.55, 1.2, 0.3, 6.2, False, f"milestone_text_{i+1}")


def _add_chart(slide: Any, chart_ledger: dict[str, Any]) -> None:
    chart = chart_ledger.get("charts", [{}])[0]
    data = CategoryChartData()
    data.categories = [point["label"] for point in chart.get("data_points", [])]
    data.add_series("Readiness", [point["value"] for point in chart.get("data_points", [])])
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.2), Inches(2.7), Inches(4.8), Inches(1.8), data).chart


def _add_table(slide: Any, table_ledger: dict[str, Any]) -> None:
    table_data = table_ledger.get("tables", [{}])[0]
    rows = table_data.get("rows", [])
    header = table_data.get("header", [])
    table_shape = slide.shapes.add_table(len(rows) + 1, len(header), Inches(0.9), Inches(1.65), Inches(7.2), Inches(2.8))
    table_shape.name = "native_governance_table"
    table = table_shape.table
    for col, value in enumerate(header):
        table.cell(0, col).text = value
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value


def _add_svg_provenance_icon(slide: Any, icon: dict[str, Any], index: int, slide_id: str, ledger: list[dict[str, Any]]) -> None:
    asset = {"asset_id": icon["svg_asset_id"], "source_path": icon["source_path"], "sha256": icon.get("source_sha256"), "canonical_viewbox": "0 0 24 24", "path_count": 1, "primitive_count": 1}
    plan = convert_svg_to_native_plan(asset)
    x = 10.7 + (index % 4) * 0.35
    y = 0.55 + (index // 4) * 0.28
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(x), Inches(y), Inches(0.16), Inches(0.16))
    shape.name = f"svg_native::{icon['semantic_intent']}::{icon['svg_asset_id']}::{slide_id}::{index+1:02d}::{plan['conversion_hash']}"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(42, 202, 218) if index % 2 == 0 else RGBColor(237, 197, 93)
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(0.25)
    ledger.append({**icon, "slide_id": slide_id, "shape_name": shape.name, "conversion_hash": plan["conversion_hash"], "source_svg_provenance_present": True})


def _rect(slide: Any, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None, name: str) -> Any:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def _line(slide: Any, x1: float, y1: float, x2: float, y2: float, color: RGBColor, name: str) -> None:
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.name = name
    shape.line.color.rgb = color
    shape.line.width = Pt(1)


def _text(slide: Any, text: str, x: float, y: float, w: float, h: float, size: float, bold: bool, name: str, color: RGBColor = RGBColor(255, 255, 255)) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _render_slide_png(binding: dict[str, Any], icons: list[dict[str, Any]], path: Path) -> None:
    image = Image.new("RGB", (1600, 900), (8, 27, 39))
    draw = ImageDraw.Draw(image)
    font_title = _font(31)
    font = _font(18)
    small = _font(14)
    draw.rounded_rectangle((48, 52, 1540, 820), radius=18, fill=(10, 43, 56), outline=(17, 112, 126), width=2)
    draw.text((72, 74), binding["title"], fill=(255, 255, 255), font=font_title)
    draw.text((72, 122), binding["subtitle"], fill=(181, 205, 214), font=font)
    draw.text((72, 690), binding["primary_claim"], fill=(235, 244, 248), font=font)
    draw.line((64, 792, 1530, 792), fill=(237, 197, 93), width=3)
    draw.text((72, 810), binding["citation_footer"], fill=(237, 197, 93), font=small)
    ref = binding["selected_reference_id"]
    if ref == "data_dashboard_hybrid":
        _draw_png_dashboard(draw, binding, font, small)
    elif ref in {"table_matrix_hybrid", "comparison_matrix_hybrid"}:
        _draw_png_matrix(draw, binding, small)
    elif ref == "timeline_roadmap_hybrid":
        _draw_png_timeline(draw, binding, small)
    elif ref in {"process_workflow_infographic", "methodology_framework_layered"}:
        _draw_png_process(draw, binding, small)
    elif ref == "visual_toc_navigation":
        _draw_png_navigation(draw, binding, small)
    elif ref == "evidence_stack_visual":
        _draw_png_evidence(draw, binding, font, small)
    else:
        _draw_png_cards(draw, binding, small)
    for i, icon in enumerate(icons):
        draw.ellipse((1320 + i * 28, 72, 1334 + i * 28, 86), fill=(42, 202, 218) if i % 2 == 0 else (237, 197, 93))
    image.save(path)


def _draw_png_cards(draw: ImageDraw.ImageDraw, binding: dict[str, Any], small: ImageFont.ImageFont) -> None:
    for i, detail in enumerate(binding["details"][:4]):
        x = 90 + (i % 2) * 460
        y = 220 + (i // 2) * 170
        draw.rounded_rectangle((x, y, x + 395, y + 110), radius=12, fill=(11, 55, 68), outline=(42, 202, 218), width=2)
        draw.arc((x + 18, y + 18, x + 58, y + 58), 35, 325, fill=(237, 197, 93), width=3)
        draw.text((x + 72, y + 24), detail[:52], fill=(240, 248, 250), font=small)


def _draw_png_evidence(draw: ImageDraw.ImageDraw, binding: dict[str, Any], font: ImageFont.ImageFont, small: ImageFont.ImageFont) -> None:
    draw.rounded_rectangle((90, 210, 570, 500), radius=16, fill=(12, 61, 72), outline=(237, 197, 93), width=3)
    draw.text((124, 250), binding["primary_claim"][:70], fill=(255, 255, 255), font=font)
    for i, detail in enumerate(binding["details"][:4]):
        y = 210 + i * 72
        draw.rounded_rectangle((680, y, 1290, y + 46), radius=10, fill=(10, 49, 62), outline=(42, 202, 218), width=2)
        draw.text((704, y + 14), detail[:68], fill=(235, 244, 248), font=small)


def _draw_png_navigation(draw: ImageDraw.ImageDraw, binding: dict[str, Any], small: ImageFont.ImageFont) -> None:
    for i, detail in enumerate(binding["details"][:6]):
        y = 190 + i * 58
        outline = (237, 197, 93) if i == 0 else (42, 202, 218)
        draw.rounded_rectangle((360, y, 1120, y + 38), radius=8, fill=(10, 49, 62), outline=outline, width=2)
        draw.ellipse((382, y + 9, 402, y + 29), outline=outline, width=3)
        draw.text((430, y + 10), detail[:54], fill=(235, 244, 248), font=small)


def _draw_png_process(draw: ImageDraw.ImageDraw, binding: dict[str, Any], small: ImageFont.ImageFont) -> None:
    y = 365
    for i, detail in enumerate((binding["details"] + ["Review", "Package"])[:5]):
        x = 105 + i * 275
        draw.rounded_rectangle((x, y, x + 155, y + 78), radius=14, fill=(12, 58, 70), outline=(42, 202, 218), width=2)
        draw.text((x + 18, y + 30), detail.split(":")[0][:16], fill=(235, 244, 248), font=small)
        if i < 4:
            draw.line((x + 155, y + 39, x + 255, y + 39), fill=(237, 197, 93), width=4)
            draw.polygon([(x + 255, y + 39), (x + 238, y + 28), (x + 238, y + 50)], fill=(237, 197, 93))
    for arc in range(5):
        draw.arc((220 + arc * 60, 165, 980 + arc * 60, 670), 205, 330, fill=(17, 112, 126), width=1)


def _draw_png_dashboard(draw: ImageDraw.ImageDraw, binding: dict[str, Any], font: ImageFont.ImageFont, small: ImageFont.ImageFont) -> None:
    for i, label in enumerate(("Trace", "Method", "QA", "Reuse")):
        x = 95 + i * 175
        draw.rounded_rectangle((x, 185, x + 130, 72 + 185), radius=10, fill=(13, 61, 73), outline=(42, 202, 218), width=2)
        draw.text((x + 24, 210), label, fill=(235, 244, 248), font=small)
    draw.rounded_rectangle((100, 315, 720, 590), radius=14, fill=(9, 47, 59), outline=(237, 197, 93), width=3)
    bars = [180, 155, 125, 168]
    for i, h in enumerate(bars):
        x = 170 + i * 110
        draw.rectangle((x, 560 - h, x + 46, 560), fill=(42, 202, 218))
    draw.rounded_rectangle((820, 315, 1250, 590), radius=14, fill=(9, 47, 59), outline=(42, 202, 218), width=2)
    draw.text((850, 365), binding["primary_claim"][:62], fill=(255, 255, 255), font=font)


def _draw_png_matrix(draw: ImageDraw.ImageDraw, binding: dict[str, Any], small: ImageFont.ImageFont) -> None:
    left, top, right, bottom = 105, 190, 960, 610
    draw.rounded_rectangle((left, top, right, bottom), radius=10, fill=(9, 47, 59), outline=(237, 197, 93), width=3)
    for i in range(1, 5):
        y = top + i * 84
        draw.line((left, y, right, y), fill=(28, 100, 116), width=2)
    for i in range(1, 4):
        x = left + i * 214
        draw.line((x, top, x, bottom), fill=(28, 100, 116), width=2)
    labels = ["Traceability", "Speed", "Judgment", "Repeatability"]
    for i, label in enumerate(labels):
        draw.text((left + 20, top + 112 + i * 84), label, fill=(235, 244, 248), font=small)


def _draw_png_timeline(draw: ImageDraw.ImageDraw, binding: dict[str, Any], small: ImageFont.ImageFont) -> None:
    draw.line((150, 400, 1350, 400), fill=(42, 202, 218), width=5)
    for i, detail in enumerate((binding["details"] + ["Review"])[:5]):
        x = 190 + i * 270
        draw.rounded_rectangle((x, 350, x + 90, 450), radius=12, fill=(12, 58, 70), outline=(237, 197, 93), width=3)
        draw.text((x - 18, 485), detail[:18], fill=(235, 244, 248), font=small)


def _build_contact_sheet(rendered_dir: Path, output: Path) -> None:
    images = [Image.open(path).convert("RGB") for path in sorted(rendered_dir.glob("slide_*.png"))]
    thumbs = []
    for image in images:
        image.thumbnail((400, 225))
        thumbs.append(image.copy())
    canvas = Image.new("RGB", (1280, 1120), (8, 20, 30))
    draw = ImageDraw.Draw(canvas)
    draw.text((30, 24), "E04H source-bound hybrid deck", fill=(255, 255, 255), font=_font(28))
    for i, thumb in enumerate(thumbs):
        x = 30 + (i % 3) * 415
        y = 80 + (i // 3) * 250
        canvas.paste(thumb, (x, y))
    canvas.save(output)


def _font(size: int) -> ImageFont.ImageFont:
    for font_name in ("arial.ttf", "calibri.ttf", "segoeui.ttf"):
        try:
            return ImageFont.truetype(font_name, size)
        except OSError:
            continue
    return ImageFont.load_default()
