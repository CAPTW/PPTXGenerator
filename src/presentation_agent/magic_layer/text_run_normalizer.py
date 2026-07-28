"""Text run normalization checks for E01.2."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


def style_text_frame(text_frame: Any, *, font_family: str, font_size_pt: int, color: str, bold: bool, alignment: int | None = None) -> None:
    if alignment is not None:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = alignment
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_family
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)


def inspect_text_run_formatting(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    runs: list[dict[str, Any]] = []
    defects: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False) or not shape.text_frame.text.strip():
                continue
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                if not paragraph.runs and paragraph.text:
                    defects.append({"shape_name": shape.name, "defect": "paragraph_text_without_explicit_run"})
                for run_index, run in enumerate(paragraph.runs, start=1):
                    record = {
                        "slide_index": slide_index,
                        "shape_index": shape_index,
                        "shape_name": shape.name,
                        "paragraph_index": paragraph_index,
                        "run_index": run_index,
                        "text": run.text,
                        "font_name": run.font.name,
                        "font_size_pt": round(run.font.size.pt, 2) if run.font.size is not None else None,
                        "bold": run.font.bold,
                        "color_rgb": str(run.font.color.rgb) if run.font.color.rgb is not None else None,
                        "paragraph_alignment": str(paragraph.alignment),
                    }
                    runs.append(record)
                    if not record["font_name"]:
                        defects.append({**record, "defect": "missing_font_family"})
                    if record["font_size_pt"] is None:
                        defects.append({**record, "defect": "missing_font_size"})
                    if record["color_rgb"] in {None, "000000"}:
                        defects.append({**record, "defect": "missing_or_black_font_color"})
    return {
        "schema_name": "text_run_normalization_report",
        "status": "passed" if not defects else "failed",
        "text_run_count": len(runs),
        "defect_count": len(defects),
        "runs": runs,
        "defects": defects,
        "canva_parity_claimed": False,
    }

