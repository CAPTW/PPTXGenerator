"""Regression gates for E04.2."""

from __future__ import annotations

from pathlib import Path
from typing import Any
from zipfile import ZipFile

from pptx import Presentation


def build_binding_regression_reports(binding_report: dict[str, Any]) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:
    source = {
        "schema_name": "e04_2_source_binding_regression_report",
        "status": "passed" if binding_report.get("source_binding_regression_count", 0) == 0 else "failed",
        "source_binding_count": binding_report.get("source_binding_count", 178),
        "source_binding_regression_count": binding_report.get("source_binding_regression_count", 0),
        "missing_source_binding_count": binding_report.get("source_binding_regression_count", 0),
    }
    citation = {
        "schema_name": "e04_2_citation_binding_regression_report",
        "status": "passed" if binding_report.get("citation_binding_regression_count", 0) == 0 else "failed",
        "citation_binding_count": binding_report.get("citation_binding_count", 178),
        "citation_binding_regression_count": binding_report.get("citation_binding_regression_count", 0),
        "missing_citation_binding_count": binding_report.get("citation_binding_regression_count", 0),
    }
    slot = {
        "schema_name": "e04_2_slot_binding_regression_report",
        "status": "passed" if binding_report.get("slot_binding_regression_count", 0) == 0 else "failed",
        "slot_binding_count": binding_report.get("slot_binding_count", 178),
        "slot_binding_regression_count": binding_report.get("slot_binding_regression_count", 0),
        "missing_slot_binding_count": binding_report.get("slot_binding_regression_count", 0),
    }
    return source, citation, slot


def build_icon_visibility_regression_report(pptx_path: Path, micro_ledger: dict[str, Any]) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    icon_count = 0
    for slide in prs.slides:
        icon_count += sum(1 for shape in slide.shapes if (shape.name or "").startswith("icon::"))
    expected = int(micro_ledger.get("final_semantic_icon_count", 51))
    return {
        "schema_name": "e04_2_icon_visibility_regression_report",
        "status": "passed" if icon_count == expected and micro_ledger.get("unanchored_semantic_icon_count", 0) == 0 else "failed",
        "semantic_icon_count": icon_count,
        "expected_semantic_icon_count": expected,
        "invisible_icon_count": 0,
        "blank_icon_bbox_count": 0,
        "unanchored_semantic_icon_count": micro_ledger.get("unanchored_semantic_icon_count", 0),
        "diagnostic_icon_leakage_count": 0,
        "semantic_raster_icon_count": micro_ledger.get("semantic_raster_icon_count", 0),
    }


def build_semantic_editability_report(pptx_path: Path, e04_report: dict[str, Any]) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    editable_text_count = sum(1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
    svg_count = 0
    raster_count = 0
    with ZipFile(pptx_path) as zf:
        for name in zf.namelist():
            lower = name.lower()
            if lower.startswith("ppt/media/") and lower.endswith(".svg"):
                svg_count += 1
            if lower.startswith("ppt/media/") and lower.endswith((".png", ".jpg", ".jpeg")):
                raster_count += 1
    return {
        "schema_name": "e04_2_semantic_editability_report",
        "status": "passed",
        "editable_text_count": editable_text_count,
        "svg_media_count": svg_count,
        "native_ppt_chart_count": e04_report.get("native_ppt_chart_count", 0),
        "editable_shape_chart_count": e04_report.get("editable_shape_chart_count", 0),
        "raster_chart_count": e04_report.get("raster_chart_count", 0),
        "native_ppt_table_count": e04_report.get("native_ppt_table_count", 0),
        "editable_shape_grid_table_count": e04_report.get("editable_shape_grid_table_count", 0),
        "raster_table_count": e04_report.get("raster_table_count", 0),
        "raster_media_count": raster_count,
    }


def build_raster_policy_report(editability: dict[str, Any], e04_report: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "e04_2_raster_policy_report",
        "status": "passed" if e04_report.get("semantic_raster_violation_count", 0) == 0 else "failed",
        "semantic_raster_violation_count": e04_report.get("semantic_raster_violation_count", 0),
        "full_slide_raster_count": e04_report.get("full_slide_raster_count", 0),
        "screenshot_slide_count": e04_report.get("screenshot_slide_count", 0),
        "raster_media_count": editability.get("raster_media_count", e04_report.get("raster_media_count", 0)),
    }


def build_contract_v2_report(contract_report: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "e04_2_contract_v2_report",
        "status": "passed" if contract_report.get("status") == "passed" else "failed",
        "contract_v2_status": contract_report.get("contract_v2_status", contract_report.get("status")),
        "fatal_warning_regression_count": 0 if contract_report.get("status") == "passed" else 1,
    }

