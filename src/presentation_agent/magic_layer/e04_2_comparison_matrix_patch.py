"""Slide 9 comparison matrix patch for E04.2."""

from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400


def patch_comparison_matrix_slide(slide) -> dict[str, Any]:
    return polish_dense_text_slide(slide, slide_number=9, archetype_id="comparison_matrix")


def polish_dense_text_slide(slide, *, slide_number: int, archetype_id: str) -> dict[str, Any]:
    updated = 0
    below_before = 0
    below_after = 0
    min_before: float | None = None
    min_after: float | None = None
    rows = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
            continue
        y_in = shape.top / EMU_PER_INCH
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.font.size:
                    continue
                before = float(run.font.size.pt)
                min_before = before if min_before is None else min(min_before, before)
                if before < 6.0:
                    below_before += 1
                    target = 7.0 if y_in <= 1.9 else 6.2
                    run.font.size = Pt(target)
                    updated += 1
                    rows.append({"shape_id": shape.shape_id, "shape_name": shape.name, "text": run.text, "before_pt": before, "after_pt": target})
                after = float(run.font.size.pt)
                min_after = after if min_after is None else min(min_after, after)
                if after < 6.0:
                    below_after += 1
        shape.text_frame.margin_left = Inches(0.03)
        shape.text_frame.margin_right = Inches(0.03)
        shape.text_frame.margin_top = Inches(0.01)
        shape.text_frame.margin_bottom = Inches(0.01)
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return {
        "schema_name": f"e04_2_slide_{slide_number:02d}_{archetype_id}_patch_report",
        "status": "passed" if below_after == 0 else "failed",
        "slide_number": slide_number,
        "archetype_id": archetype_id,
        "updated_text_run_count": updated,
        "text_below_6pt_before_count": below_before,
        "text_below_6pt_after_count": below_after,
        "minimum_font_before_pt": min_before,
        "minimum_font_after_pt": min_after,
        "patch_actions": ["font_minimum_raised", "text_box_margins_tightened", "vertical_anchor_middle"],
        "rows": rows,
    }

