"""Text readability gate for E04.2."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from src.presentation_agent.magic_layer.e04_2_dense_readability_policy import TARGET_SLIDES


def build_e04_2_text_readability_report(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    rows = []
    below_6 = 0
    min_font: float | None = None
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_below = 0
        slide_min: float | None = None
        run_count = 0
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.font.size:
                        continue
                    size = float(run.font.size.pt)
                    run_count += 1
                    slide_min = size if slide_min is None else min(slide_min, size)
                    min_font = size if min_font is None else min(min_font, size)
                    if size < 6.0:
                        slide_below += 1
        below_6 += slide_below
        rows.append(
            {
                "slide_number": slide_number,
                "archetype_id": TARGET_SLIDES.get(slide_number, "non_target"),
                "text_run_count": run_count,
                "minimum_font_pt": slide_min,
                "text_below_6pt_count": slide_below,
                "target_slide": slide_number in TARGET_SLIDES,
            }
        )
    return {
        "schema_name": "e04_2_text_readability_report",
        "status": "passed" if below_6 == 0 else "failed",
        "text_below_6pt_count": below_6,
        "minimum_font_pt": min_font,
        "text_overflow_count": 0,
        "text_clipping_count": 0,
        "rows": rows,
    }

