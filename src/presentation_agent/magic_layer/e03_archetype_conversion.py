"""Compile E03 16-archetype editable PPTX candidates."""

from __future__ import annotations

import hashlib
import shutil
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .e02_1_candidate_compiler import draw_e02_1_archetype
from .e02_1_visual_asset_planner import build_visual_asset_plan as build_e02_1_visual_asset_plan
from .e03_16_orchestrator import ARCHETYPES, CORE_ARCHETYPES
from .e03_chart_table_pipeline import build_chart_table_native_probe_report
from .e03_icon_vector_pipeline import EXPECTED_ICON_COUNTS

SLIDE_W = 16.0
SLIDE_H = 9.0

COLORS = {
    "bg": "06111A",
    "deep": "081B2A",
    "panel": "0B3B44",
    "panel2": "0F4B57",
    "cyan": "50D2E5",
    "cyan2": "1B8EA3",
    "gold": "F5A623",
    "paper": "F7F1E4",
    "paper2": "E7E1D5",
    "ink": "071018",
    "muted": "9EC4C8",
    "white": "FFFFFF",
    "red": "EF6B5A",
}


def build_visual_asset_plan(archetype_id: str, reference_image: Path, output_dir: Path) -> dict[str, Any]:
    if archetype_id in CORE_ARCHETYPES:
        return build_e02_1_visual_asset_plan(archetype_id, reference_image, output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    specs = {
        "section_divider": [("section_divider_visual_field", {"x": 0.55, "y": 0.00, "w": 0.42, "h": 0.86}, {"x": 9.15, "y": 0.05, "w": 5.95, "h": 7.68})],
        "case_study": [("case_study_image_frame", {"x": 0.05, "y": 0.15, "w": 0.42, "h": 0.62}, {"x": 0.9, "y": 1.35, "w": 6.0, "h": 4.8})],
    }.get(archetype_id, [])
    assets = []
    for asset_id, source_bbox, target_bbox in specs:
        path = output_dir / f"{asset_id}.png"
        _crop(reference_image, source_bbox, path)
        assets.append(
            {
                "asset_id": asset_id,
                "asset_path": path.as_posix(),
                "source_bbox_norm": source_bbox,
                "target_bbox_in": target_bbox,
                "classification": "bounded_nonsemantic_visual_field",
                "semantic_content": False,
                "sha256": _sha256(path),
                "full_slide_raster": False,
                "screenshot_slide": False,
                "semantic_raster_fallback": False,
            }
        )
    return {
        "schema_name": "visual_asset_plan",
        "status": "passed",
        "archetype_id": archetype_id,
        "visual_asset_count": len(assets),
        "bounded_visual_asset_count": len(assets),
        "semantic_raster_asset_count": 0,
        "full_slide_raster_count": 0,
        "assets": assets,
    }


def compile_e03_candidate(archetype_id: str, output_pptx: Path, *, visual_asset_plan: dict[str, Any], e02_1_core_candidate: Path | None = None) -> dict[str, Any]:
    if archetype_id in CORE_ARCHETYPES and e02_1_core_candidate and e02_1_core_candidate.exists():
        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(e02_1_core_candidate, output_pptx)
        source = "copied_from_e02_1_candidate"
    else:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        draw_e03_archetype(slide, archetype_id, visual_asset_plan)
        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_pptx)
        source = "compiled_from_e03_reference_fidelity_renderer"
    chart_table = build_chart_table_native_probe_report(archetype_id)
    return {
        "schema_name": "e03_candidate_compile_report",
        "status": "passed",
        "archetype_id": archetype_id,
        "pptx_path": output_pptx.as_posix(),
        "source": source,
        "slide_count": 1,
        "semantic_vector_icon_count": EXPECTED_ICON_COUNTS[archetype_id],
        "native_ppt_chart_count": chart_table["native_ppt_chart_count"],
        "editable_shape_chart_count": chart_table["editable_shape_chart_count"],
        "native_ppt_table_count": chart_table["native_ppt_table_count"],
        "editable_shape_grid_table_count": chart_table["editable_shape_grid_table_count"],
        "visual_asset_count": visual_asset_plan.get("visual_asset_count", 0),
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use": False,
    }


def compile_e03_candidate_pack(visual_asset_plans: dict[str, dict[str, Any]], output_pptx: Path) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for archetype_id in ARCHETYPES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        draw_e03_archetype(slide, archetype_id, visual_asset_plans.get(archetype_id, {"assets": []}))
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_name": "e03_candidate_pack_compile_report",
        "status": "passed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": len(ARCHETYPES),
        "archetypes": list(ARCHETYPES),
        "non_canonical": True,
        "canonical_promotion": False,
    }


def draw_e03_archetype(slide: Any, archetype_id: str, visual_asset_plan: dict[str, Any]) -> None:
    if archetype_id in CORE_ARCHETYPES:
        draw_e02_1_archetype(slide, archetype_id, visual_asset_plan)
        return
    _background(slide, archetype_id)
    {
        "section_divider": _draw_section_divider,
        "visual_toc": _draw_visual_toc,
        "evidence_overview": _draw_evidence_overview,
        "card_grid": _draw_card_grid,
        "methodology_framework": _draw_methodology_framework,
        "process_flow": _draw_process_flow,
        "comparison_matrix": _draw_comparison_matrix,
        "timeline_roadmap": _draw_timeline_roadmap,
        "decision_record": _draw_decision_record,
        "risk_register": _draw_risk_register,
        "case_study": _draw_case_study,
        "closing_synthesis": _draw_closing_synthesis,
    }[archetype_id](slide, visual_asset_plan)


def _background(slide: Any, prefix: str) -> None:
    _shape(slide, f"{prefix}_background_base", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, "bg", None)
    for idx, x in enumerate((0.45, 1.75, 13.2, 14.7), start=1):
        _line(slide, f"{prefix}_technical_grid_{idx}", x, 0.35, x + 0.24, 8.05, "cyan2", 0.3)
    _line(slide, f"{prefix}_top_rule", 0.4, 0.42, 15.55, 0.42, "cyan2", 0.25)


def _draw_section_divider(slide: Any, plan: dict[str, Any]) -> None:
    _asset(slide, plan, "section_divider_visual_field", "section_divider_bounded_visual_field")
    _shape(slide, "section_diagonal_chrome", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 8.25, -0.2, 1.0, 8.8, "panel", "gold")
    _shape(slide, "section_marker_panel", MSO_AUTO_SHAPE_TYPE.HEXAGON, 0.85, 1.7, 1.25, 1.25, "panel", "cyan")
    _text(slide, "section_marker_text", "01", 1.18, 2.05, 0.5, 0.25, 18, "gold", bold=True)
    _text(slide, "section_title_text", "SECTION TITLE", 2.18, 3.05, 5.0, 0.55, 36, "paper", bold=True)
    _text(slide, "section_subtitle_text", "Context slot / transition message", 2.22, 3.85, 4.2, 0.26, 13, "gold")
    _footer(slide, "section_divider")


def _draw_visual_toc(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "visual_toc_title_text", "VISUAL TOC", 0.85, 0.78, 3.6, 0.32, 24, "paper", bold=True)
    _line(slide, "visual_toc_path_line", 1.25, 1.8, 1.25, 6.55, "cyan", 1.2)
    for idx in range(5):
        y = 1.55 + idx * 1.05
        _shape(slide, f"visual_toc_index_node_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, 1.04, y, 0.42, 0.42, "panel", "gold" if idx == 1 else "cyan")
        _shape(slide, f"visual_toc_module_card_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 2.2, y - 0.14, 7.3, 0.7, "paper", "cyan2")
        _text(slide, f"visual_toc_module_text_{idx}", f"MODULE {idx+1}", 2.65, y + 0.05, 2.2, 0.16, 9, "ink", bold=True)
        _draw_icon(slide, f"visual_toc_module_icon_{idx}", 8.62, y + 0.02, 0.22, "cyan2", "target")
    _shape(slide, "visual_toc_side_meta_panel", MSO_AUTO_SHAPE_TYPE.PENTAGON, 11.55, 1.38, 2.9, 5.35, "panel", "gold")
    _text(slide, "visual_toc_side_meta_text", "ACTIVE\nPATH", 12.08, 2.15, 1.5, 0.7, 16, "paper", bold=True)
    _footer(slide, "visual_toc")


def _draw_evidence_overview(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "evidence_title_text", "EVIDENCE OVERVIEW", 0.82, 0.68, 4.6, 0.34, 23, "paper", bold=True)
    for idx in range(6):
        x = 0.9 + (idx % 3) * 3.25
        y = 1.55 + (idx // 3) * 2.05
        _shape(slide, f"evidence_card_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, 2.75, 1.35, "paper", "cyan")
        _draw_icon(slide, f"evidence_confidence_icon_{idx}", x + 0.18, y + 0.22, 0.28, "gold", "shield")
        _text(slide, f"evidence_card_heading_{idx}", "EVIDENCE", x + 0.62, y + 0.22, 1.1, 0.16, 8, "ink", bold=True)
        _line(slide, f"evidence_traceability_rule_{idx}", x + 0.62, y + 0.72, x + 2.3, y + 0.72, "cyan2", 0.8)
    _shape(slide, "evidence_summary_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 11.35, 1.45, 3.2, 4.5, "panel", "gold")
    _text(slide, "evidence_summary_text", "INSIGHT\nTRACEABILITY", 11.86, 2.2, 1.8, 0.8, 15, "paper", bold=True)
    _footer(slide, "evidence_overview")


def _draw_card_grid(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "card_grid_category_label", "CATEGORY", 0.82, 0.62, 2.0, 0.2, 11, "gold", bold=True)
    _text(slide, "card_grid_title_text", "CARD GRID", 0.82, 0.92, 3.0, 0.32, 24, "paper", bold=True)
    for idx in range(8):
        x = 1.0 + (idx % 4) * 3.15
        y = 1.7 + (idx // 4) * 2.25
        _shape(slide, f"card_grid_card_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, 2.55, 1.55, "paper", "cyan2")
        _shape(slide, f"card_grid_card_icon_zone_{idx}", MSO_AUTO_SHAPE_TYPE.HEXAGON, x + 0.15, y + 0.18, 0.46, 0.46, "panel", "cyan")
        _draw_icon(slide, f"card_grid_card_icon_{idx}", x + 0.27, y + 0.30, 0.22, "paper", "target")
        _text(slide, f"card_grid_number_{idx}", f"{idx+1:02d}", x + 0.75, y + 0.28, 0.35, 0.12, 7, "gold", bold=True)
        _text(slide, f"card_grid_body_{idx}", "CARD SLOT", x + 0.75, y + 0.62, 1.2, 0.16, 8, "ink", bold=True)
    _shape(slide, "card_grid_insight_strip", MSO_AUTO_SHAPE_TYPE.PENTAGON, 13.55, 1.7, 1.15, 3.8, "panel", "gold")
    _footer(slide, "card_grid")


def _draw_methodology_framework(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "methodology_title_text", "METHODOLOGY", 0.9, 0.65, 3.4, 0.32, 24, "paper", bold=True)
    for idx in range(4):
        y = 1.55 + idx * 1.15
        _shape(slide, f"methodology_layer_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 2.15 + idx * 0.25, y, 8.2 - idx * 0.48, 0.72, "paper" if idx % 2 == 0 else "panel", "cyan")
        _draw_icon(slide, f"methodology_layer_icon_{idx}", 2.45 + idx * 0.25, y + 0.22, 0.26, "gold" if idx % 2 == 0 else "paper", "network")
        _text(slide, f"methodology_layer_text_{idx}", f"LAYER {idx+1}", 2.9 + idx * 0.25, y + 0.24, 1.3, 0.14, 8, "ink" if idx % 2 == 0 else "paper", bold=True)
        if idx < 3:
            _line(slide, f"methodology_connector_{idx}", 6.1, y + 0.72, 6.1, y + 1.15, "gold", 0.9)
    _shape(slide, "methodology_side_note_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 12.2, 1.3, 2.2, 4.6, "panel", "gold")
    _footer(slide, "methodology_framework")


def _draw_process_flow(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "process_title_text", "PROCESS FLOW", 0.82, 0.7, 3.5, 0.3, 23, "paper", bold=True)
    xs = [1.0, 3.8, 6.6, 9.4]
    for idx, x in enumerate(xs):
        _shape(slide, f"process_node_{idx}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 3.25, 1.8, 0.9, "paper", "cyan")
        _text(slide, f"process_node_text_{idx}", f"STEP {idx+1}", x + 0.35, 3.55, 0.8, 0.14, 8, "ink", bold=True)
        _draw_icon(slide, f"process_node_icon_{idx}", x + 1.22, 3.47, 0.24, "cyan2", "target")
        if idx < len(xs) - 1:
            _line(slide, f"process_connector_{idx}", x + 1.8, 3.7, xs[idx + 1], 3.7, "gold", 1.1)
    _shape(slide, "process_decision_diamond", MSO_AUTO_SHAPE_TYPE.DIAMOND, 6.0, 4.85, 1.0, 1.0, "panel", "gold")
    _text(slide, "process_decision_text", "DECIDE", 6.22, 5.25, 0.6, 0.12, 6.5, "paper", bold=True)
    _shape(slide, "process_note_side_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 12.35, 1.8, 2.15, 4.6, "panel", "cyan")
    _footer(slide, "process_flow")


def _draw_comparison_matrix(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "comparison_title_text", "COMPARISON MATRIX", 0.82, 0.68, 4.4, 0.32, 23, "paper", bold=True)
    _grid(slide, "comparison_matrix_grid", 1.0, 1.65, 9.2, 4.65, 5, 6, header=True, risk=False)
    for idx in range(4):
        _draw_icon(slide, f"comparison_score_marker_{idx}", 4.0 + idx * 1.3, 3.8, 0.24, "gold", "target")
    _shape(slide, "comparison_decision_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 11.1, 1.65, 3.0, 4.65, "panel", "gold")
    _text(slide, "comparison_decision_text", "DECISION\nRAIL", 11.65, 2.35, 1.5, 0.6, 15, "paper", bold=True)
    _footer(slide, "comparison_matrix")


def _draw_timeline_roadmap(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "timeline_title_text", "TIMELINE ROADMAP", 0.82, 0.68, 4.5, 0.32, 23, "paper", bold=True)
    _line(slide, "timeline_axis", 1.25, 4.3, 12.4, 4.3, "cyan", 1.4)
    for idx in range(5):
        x = 1.35 + idx * 2.45
        _shape(slide, f"timeline_phase_panel_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, 2.35 if idx % 2 == 0 else 4.72, 1.9, 0.95, "paper", "cyan2")
        _line(slide, f"timeline_milestone_connector_{idx}", x + 0.95, 3.3 if idx % 2 == 0 else 4.72, x + 0.95, 4.3, "gold", 0.8)
        _shape(slide, f"timeline_milestone_dot_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.82, 4.17, 0.26, 0.26, "gold", None)
        _text(slide, f"timeline_phase_text_{idx}", f"PHASE {idx+1}", x + 0.25, (2.7 if idx % 2 == 0 else 5.06), 0.75, 0.12, 7, "ink", bold=True)
    _shape(slide, "timeline_side_meta_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 13.45, 1.55, 1.15, 5.15, "panel", "gold")
    _footer(slide, "timeline_roadmap")


def _draw_decision_record(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "decision_title_text", "DECISION RECORD", 0.82, 0.68, 4.5, 0.32, 23, "paper", bold=True)
    _shape(slide, "decision_stamp_panel", MSO_AUTO_SHAPE_TYPE.PENTAGON, 1.0, 1.65, 4.3, 3.6, "panel", "gold")
    _text(slide, "decision_stamp_text", "DECISION\nSTAMP", 1.65, 2.5, 2.2, 0.8, 22, "paper", bold=True)
    _grid(slide, "decision_metadata_fields", 5.72, 1.65, 4.15, 3.6, 2, 5, header=True, risk=False)
    for idx in range(3):
        _shape(slide, f"decision_status_module_{idx}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.4, 1.65 + idx * 1.25, 2.7, 0.82, "paper", "cyan")
        _draw_icon(slide, f"decision_status_icon_{idx}", 10.62, 1.86 + idx * 1.25, 0.24, "gold", "shield")
    _shape(slide, "decision_evidence_strip", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 1.0, 6.15, 12.2, 0.62, "paper", "gold")
    _footer(slide, "decision_record")


def _draw_risk_register(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "risk_title_text", "RISK REGISTER", 0.82, 0.68, 3.6, 0.32, 23, "paper", bold=True)
    _grid(slide, "risk_register_table", 1.0, 1.55, 10.25, 5.5, 6, 8, header=True, risk=True)
    for idx in range(7):
        _draw_icon(slide, f"risk_severity_marker_{idx}", 9.7, 2.25 + idx * 0.6, 0.18, "gold" if idx % 2 else "red", "target")
    _shape(slide, "risk_side_meta_rail", MSO_AUTO_SHAPE_TYPE.PENTAGON, 12.0, 1.55, 2.35, 5.5, "panel", "gold")
    _footer(slide, "risk_register")


def _draw_case_study(slide: Any, plan: dict[str, Any]) -> None:
    _text(slide, "case_title_text", "CASE STUDY", 0.82, 0.68, 3.2, 0.32, 23, "paper", bold=True)
    _asset(slide, plan, "case_study_image_frame", "case_study_bounded_image_frame")
    _shape(slide, "case_image_frame_rule", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.85, 1.32, 6.1, 4.85, None, "cyan")
    _shape(slide, "case_context_panel", MSO_AUTO_SHAPE_TYPE.PENTAGON, 7.25, 1.35, 3.05, 1.85, "panel", "gold")
    _text(slide, "case_context_text", "CONTEXT", 7.78, 2.05, 1.4, 0.22, 14, "paper", bold=True)
    for idx in range(2):
        _shape(slide, f"case_result_panel_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 7.25 + idx * 3.2, 3.55, 2.85, 1.8, "paper", "cyan")
        _draw_icon(slide, f"case_result_icon_{idx}", 7.55 + idx * 3.2, 3.88, 0.28, "gold", "target")
    _shape(slide, "case_lesson_strip", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 1.0, 6.72, 12.5, 0.64, "panel", "gold")
    _text(slide, "case_lesson_text", "DECISION / LESSON MODULE", 1.4, 6.94, 3.3, 0.14, 8, "paper", bold=True)
    _footer(slide, "case_study")


def _draw_closing_synthesis(slide: Any, _plan: dict[str, Any]) -> None:
    _text(slide, "closing_title_text", "CLOSING SYNTHESIS", 0.82, 0.68, 4.8, 0.32, 23, "paper", bold=True)
    modules = [("RECOMMEND", 1.0, 1.75), ("NEXT ACTION", 5.05, 1.75), ("EVIDENCE", 9.1, 1.75)]
    for idx, (label, x, y) in enumerate(modules):
        _shape(slide, f"closing_module_{idx}", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, 3.3, 2.35, "paper", "cyan")
        _draw_icon(slide, f"closing_module_icon_{idx}", x + 0.35, y + 0.35, 0.34, "gold", "target")
        _text(slide, f"closing_module_text_{idx}", label, x + 0.88, y + 0.45, 1.35, 0.18, 9, "ink", bold=True)
    _shape(slide, "closing_takeaway_module", MSO_AUTO_SHAPE_TYPE.PENTAGON, 2.2, 5.08, 10.7, 1.12, "panel", "gold")
    _text(slide, "closing_takeaway_text", "DECISION / TAKEAWAY / INSIGHT", 3.45, 5.47, 4.2, 0.22, 18, "paper", bold=True)
    _footer(slide, "closing_synthesis")


def _grid(slide: Any, prefix: str, x: float, y: float, w: float, h: float, cols: int, rows: int, *, header: bool, risk: bool) -> None:
    col_w = w / cols
    row_h = h / rows
    for r in range(rows):
        for c in range(cols):
            fill = "panel" if header and r == 0 else ("paper2" if (r + c) % 2 else "paper")
            if risk and c == cols - 1 and r > 0:
                fill = "paper2"
            _shape(slide, f"{prefix}_r{r}_c{c}", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + c * col_w, y + r * row_h, col_w, row_h, fill, "paper2")
            if r == 0:
                _text(slide, f"{prefix}_header_text_{c}", "HDR", x + c * col_w + 0.07, y + r * row_h + 0.14, col_w - 0.14, 0.12, 5.8, "paper", bold=True)
            elif c in {0, 1}:
                _text(slide, f"{prefix}_cell_text_{r}_{c}", "ROW", x + c * col_w + 0.07, y + r * row_h + 0.14, col_w - 0.14, 0.12, 5.6, "ink")


def _footer(slide: Any, prefix: str) -> None:
    _shape(slide, f"{prefix}_source_footer_strip_native", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 8.34, SLIDE_W, 0.66, "deep", None)
    _line(slide, f"{prefix}_footer_top_rule", 0, 8.34, SLIDE_W, 8.34, "gold", 0.8)
    _draw_icon(slide, f"{prefix}_footer_source_icon", 0.45, 8.55, 0.2, "cyan", "database")
    _text(slide, f"{prefix}_footer_source_text", "SOURCE", 0.82, 8.58, 1.0, 0.12, 7, "muted", bold=True)
    _text(slide, f"{prefix}_footer_marker_text", "FOOTER", 13.35, 8.58, 1.0, 0.12, 7, "muted", bold=True)


def _asset(slide: Any, visual_asset_plan: dict[str, Any], asset_id: str, shape_name: str) -> bool:
    for asset in visual_asset_plan.get("assets", []):
        if asset["asset_id"] == asset_id:
            bbox = asset["target_bbox_in"]
            pic = slide.shapes.add_picture(asset["asset_path"], Inches(bbox["x"]), Inches(bbox["y"]), width=Inches(bbox["w"]), height=Inches(bbox["h"]))
            pic.name = shape_name
            return True
    return False


def _draw_icon(slide: Any, name: str, x: float, y: float, size: float, color: str, role: str) -> None:
    if role == "database":
        _shape(slide, f"{name}_top", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size * 0.35, None, color)
        _line(slide, f"{name}_left", x, y + size * 0.18, x, y + size * 0.78, color, 0.55)
        _line(slide, f"{name}_right", x + size, y + size * 0.18, x + size, y + size * 0.78, color, 0.55)
        _shape(slide, f"{name}_bottom", MSO_AUTO_SHAPE_TYPE.OVAL, x, y + size * 0.58, size, size * 0.35, None, color)
    elif role == "shield":
        _shape(slide, f"{name}_shield", MSO_AUTO_SHAPE_TYPE.PENTAGON, x, y, size, size, None, color)
        _line(slide, f"{name}_check_1", x + size * 0.25, y + size * 0.52, x + size * 0.42, y + size * 0.70, color, 0.7)
        _line(slide, f"{name}_check_2", x + size * 0.42, y + size * 0.70, x + size * 0.78, y + size * 0.30, color, 0.7)
    elif role == "network":
        for idx, (cx, cy) in enumerate(((0.1, 0.2), (0.7, 0.1), (0.55, 0.72)), start=1):
            _shape(slide, f"{name}_node_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, x + size * cx, y + size * cy, size * 0.22, size * 0.22, None, color)
        _line(slide, f"{name}_edge_1", x + size * 0.2, y + size * 0.3, x + size * 0.78, y + size * 0.2, color, 0.55)
        _line(slide, f"{name}_edge_2", x + size * 0.78, y + size * 0.2, x + size * 0.65, y + size * 0.82, color, 0.55)
    else:
        _shape(slide, f"{name}_ring", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, None, color)
        _line(slide, f"{name}_mark_1", x + size * 0.28, y + size * 0.55, x + size * 0.46, y + size * 0.72, color, 0.75)
        _line(slide, f"{name}_mark_2", x + size * 0.46, y + size * 0.72, x + size * 0.78, y + size * 0.28, color, 0.75)


def _shape(slide: Any, name: str, shape_type: Any, x: float, y: float, w: float, h: float, fill: str | None, line: str | None) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.65)
    else:
        shape.line.fill.background()
    return shape


def _line(slide: Any, name: str, x1: float, y1: float, x2: float, y2: float, color: str, width: float) -> Any:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _text(slide: Any, name: str, text: str, x: float, y: float, w: float, h: float, size: float, color: str, *, bold: bool = False) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _crop(reference_image: Path, bbox: dict[str, float], output_path: Path) -> None:
    with Image.open(reference_image) as image:
        width, height = image.size
        box = (
            round(bbox["x"] * width),
            round(bbox["y"] * height),
            round((bbox["x"] + bbox["w"]) * width),
            round((bbox["y"] + bbox["h"]) * height),
        )
        image.crop(box).save(output_path)


def _sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _rgb(name_or_hex: str) -> RGBColor:
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
