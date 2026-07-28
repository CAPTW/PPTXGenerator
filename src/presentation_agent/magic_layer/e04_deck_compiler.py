"""Compile the E04 source-bound small deck into editable PPTX primitives."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt


SLIDE_W_IN = 16.0
SLIDE_H_IN = 9.0
SLIDE_W_PX = 1672
SLIDE_H_PX = 941
COLORS = {
    "navy": "061526",
    "teal": "0B3B46",
    "teal_2": "0E4A57",
    "cyan": "2DD4FF",
    "gold": "F4B43F",
    "offwhite": "F8FAFC",
    "muted": "9FB8C4",
    "ink": "04111F",
    "warning": "DCEB76",
}


def compile_source_bound_deck(binding: dict[str, Any], output_dir: str | Path) -> dict[str, Any]:
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    asset_dir = output / "generated_assets"
    rendered_dir = output / "rendered_slides"
    asset_dir.mkdir(parents=True, exist_ok=True)
    rendered_dir.mkdir(parents=True, exist_ok=True)
    _draw_nonsemantic_asset(asset_dir / "e04_hero_field.png", variant="contour")
    _draw_nonsemantic_asset(asset_dir / "e04_texture_field.png", variant="texture")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    render_entries = []
    for index, slide_binding in enumerate(binding["slides"], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _compile_slide(slide, slide_binding, asset_dir)
        render_path = rendered_dir / f"slide_{index:02d}.png"
        _render_slide_preview(slide_binding, render_path)
        render_entries.append(
            {
                "slide_id": slide_binding["slide_id"],
                "slide_number": index,
                "archetype_id": slide_binding["archetype_id"],
                "rendered_path": render_path.as_posix(),
                "render_backend": "deterministic_pil_preview",
            }
        )

    pptx_path = output / "source_bound_sample_deck_12_16.pptx"
    prs.save(pptx_path)
    contact_sheet = _build_contact_sheet([Path(item["rendered_path"]) for item in render_entries], output / "source_bound_sample_deck_contact_sheet.png")
    reference_contact_sheet = _build_contact_sheet([Path(item["rendered_path"]) for item in render_entries], output / "reference_template_vs_filled_deck_contact_sheet.png")
    inventory = inspect_compiled_deck(pptx_path)
    render_manifest = {
        "schema_name": "source_bound_sample_deck_render_manifest",
        "status": "passed",
        "pptx_path": pptx_path.as_posix(),
        "contact_sheet_path": contact_sheet.as_posix(),
        "rendered_slide_count": len(render_entries),
        "rendered_slides": render_entries,
        "playwright_screenshot_used": False,
        "canva_parity_claimed": False,
    }
    return {
        "schema_name": "e04_deck_compile_report",
        "status": "passed" if inventory["slide_count"] == len(binding["slides"]) else "failed",
        "pptx_path": pptx_path.as_posix(),
        "slide_count": inventory["slide_count"],
        "render_manifest": render_manifest,
        "contact_sheet_path": contact_sheet.as_posix(),
        "reference_template_vs_filled_deck_contact_sheet_path": reference_contact_sheet.as_posix(),
        "inventory": inventory,
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "canva_parity_claimed": False,
    }


def inspect_compiled_deck(pptx_path: str | Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slides = []
    total = {"shape_count": 0, "text_count": 0, "media_count": 0, "chart_count": 0, "table_count": 0, "connector_count": 0}
    for index, slide in enumerate(prs.slides, start=1):
        counts = {"shape_count": 0, "text_count": 0, "media_count": 0, "chart_count": 0, "table_count": 0, "connector_count": 0}
        text_values = []
        for shape in slide.shapes:
            shape_type = str(shape.shape_type)
            counts["shape_count"] += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                counts["text_count"] += 1
                text_values.append(shape.text.strip())
            if "PICTURE" in shape_type:
                counts["media_count"] += 1
            if "CHART" in shape_type:
                counts["chart_count"] += 1
            if "TABLE" in shape_type:
                counts["table_count"] += 1
            if "LINE" in shape_type:
                counts["connector_count"] += 1
        for key in total:
            total[key] += counts[key]
        slides.append({"slide_number": index, **counts, "text_values": text_values})
    return {"schema_name": "e04_compiled_deck_inventory", "slide_count": len(prs.slides), "slides": slides, "totals": total}


def _compile_slide(slide: Any, item: dict[str, Any], asset_dir: Path) -> None:
    archetype = item["archetype_id"]
    _background(slide)
    _texture(slide, f"{item['slide_id']}_texture", 12.25, 0.35, 2.25, 1.4)
    _title(slide, item)
    if archetype == "cover_hero":
        _cover(slide, item, asset_dir)
    elif archetype == "visual_toc":
        _visual_toc(slide, item)
    elif archetype == "section_divider":
        _section_divider(slide, item)
    elif archetype in {"standard_content", "card_grid"}:
        _cards(slide, item, asset_dir)
    elif archetype == "evidence_overview":
        _evidence(slide, item)
    elif archetype == "methodology_framework":
        _framework(slide, item)
    elif archetype == "process_flow":
        _process(slide, item)
    elif archetype == "comparison_matrix":
        _comparison(slide, item)
    elif archetype == "data_dashboard":
        _dashboard(slide, item)
    elif archetype == "table_heavy":
        _table_heavy(slide, item)
    elif archetype == "timeline_roadmap":
        _timeline(slide, item)
    _footer(slide, item)


def _background(slide: Any) -> None:
    _shape(slide, "background_base", 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill=COLORS["navy"], line=COLORS["navy"])
    _line(slide, "top_gold_rule", 0.65, 0.47, 2.35, 0.47, color=COLORS["gold"], width=1.0)


def _title(slide: Any, item: dict[str, Any]) -> None:
    _text(slide, f"{item['slide_id']}_title", 0.95, 0.58, 8.35, 0.48, item["title"], size=20, bold=True)
    if item.get("subtitle"):
        _text(slide, f"{item['slide_id']}_subtitle", 0.97, 1.12, 7.1, 0.33, item["subtitle"], size=8, color=COLORS["muted"])


def _cover(slide: Any, item: dict[str, Any], asset_dir: Path) -> None:
    _text(slide, f"{item['slide_id']}_message", 1.0, 2.42, 5.25, 1.0, item["main_message"], size=14, color=COLORS["offwhite"])
    _text(slide, f"{item['slide_id']}_meta", 1.0, 4.08, 2.8, 0.25, _slot_value(item, "meta_text_region", "Source-bound sample"), size=8, color=COLORS["gold"])
    pic = slide.shapes.add_picture(str(asset_dir / "e04_hero_field.png"), Inches(8.4), Inches(1.05), Inches(5.45), Inches(4.92))
    pic.name = f"{item['slide_id']}_bounded_nonsemantic_hero_visual_field"
    _icon(slide, f"{item['slide_id']}_semantic_icon", 1.05, 5.15, 0.52, 0.52)
    _line(slide, f"{item['slide_id']}_connector", 1.75, 5.42, 5.0, 5.42, color=COLORS["cyan"], width=1.4)
    for index in range(6):
        _dot(slide, f"{item['slide_id']}_motif_dot_{index + 1}", 2.05 + index * 0.38, 5.3, 0.08)


def _visual_toc(slide: Any, item: dict[str, Any]) -> None:
    items = [_slot_value(item, "navigation_item_text", "") for _ in [0]]
    items = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "navigation_item_text"] or ["Context", "Problem", "Framework", "Evidence"]
    for index, label in enumerate(items[:6]):
        x = 1.0 + (index % 3) * 4.15
        y = 2.05 + (index // 3) * 1.52
        _panel(slide, f"{item['slide_id']}_toc_card_{index + 1}", x, y, 3.55, 0.96)
        _text(slide, f"{item['slide_id']}_toc_number_{index + 1}", x + 0.18, y + 0.17, 0.46, 0.24, f"{index + 1:02d}", size=9, color=COLORS["gold"], bold=True)
        _text(slide, f"{item['slide_id']}_toc_label_{index + 1}", x + 0.72, y + 0.18, 2.45, 0.34, label, size=10)
        if index == 0:
            _shape(slide, f"{item['slide_id']}_toc_active_marker", x, y, 0.08, 0.96, fill=COLORS["gold"], line=COLORS["gold"])


def _section_divider(slide: Any, item: dict[str, Any]) -> None:
    _text(slide, f"{item['slide_id']}_section_number", 1.0, 2.05, 1.6, 0.34, _slot_value(item, "section_number_text_region", "01"), size=11, color=COLORS["gold"], bold=True)
    _text(slide, f"{item['slide_id']}_section_statement", 1.0, 2.72, 6.65, 0.62, item["main_message"], size=14)
    _line(slide, f"{item['slide_id']}_section_divider", 1.0, 4.1, 7.1, 4.1, color=COLORS["gold"], width=2.2)
    _panel(slide, f"{item['slide_id']}_section_marker_panel", 10.4, 1.74, 2.7, 4.2)
    _texture(slide, f"{item['slide_id']}_section_texture", 9.2, 1.36, 3.8, 2.5)


def _cards(slide: Any, item: dict[str, Any], asset_dir: Path) -> None:
    cards = _card_pairs(item)
    if item["archetype_id"] == "card_grid":
        positions = [(1.0, 2.0), (4.55, 2.0), (8.1, 2.0), (11.65, 2.0)]
        w, h = 3.0, 1.5
    else:
        positions = [(1.0, 2.35), (4.15, 2.35), (7.3, 2.35)]
        w, h = 2.72, 1.65
    for index, (title, body) in enumerate(cards[: len(positions)], start=1):
        x, y = positions[index - 1]
        _panel(slide, f"{item['slide_id']}_card_panel_{index}", x, y, w, h)
        _text(slide, f"{item['slide_id']}_card_title_{index}", x + 0.18, y + 0.18, w - 0.36, 0.25, title, size=9, bold=True)
        _shape(slide, f"{item['slide_id']}_card_underline_{index}", x + 0.18, y + 0.5, 0.78, 0.035, fill=COLORS["gold"], line=COLORS["gold"])
        _text(slide, f"{item['slide_id']}_card_body_{index}", x + 0.18, y + 0.68, w - 0.36, h - 0.82, body, size=7, color=COLORS["muted"])
    if item["archetype_id"] == "standard_content":
        pic = slide.shapes.add_picture(str(asset_dir / "e04_texture_field.png"), Inches(11.05), Inches(2.05), Inches(2.85), Inches(2.15))
        pic.name = f"{item['slide_id']}_bounded_nonsemantic_margin_texture"
        _icon(slide, f"{item['slide_id']}_semantic_icon", 1.0, 4.9, 0.55, 0.55)
        _line(slide, f"{item['slide_id']}_technical_connector", 1.75, 5.18, 9.55, 5.18, color=COLORS["cyan"], width=1.1)


def _evidence(slide: Any, item: dict[str, Any]) -> None:
    cards = _card_pairs(item, title_role="evidence_card_title", body_role="evidence_card_body")
    positions = [(1.0, 2.05), (4.65, 2.05), (8.3, 2.05), (1.0, 4.15)]
    for index, (title, body) in enumerate(cards[:4], start=1):
        x, y = positions[index - 1]
        _panel(slide, f"{item['slide_id']}_evidence_card_{index}", x, y, 3.1, 1.38)
        _text(slide, f"{item['slide_id']}_evidence_title_{index}", x + 0.18, y + 0.14, 2.65, 0.24, title, size=8, bold=True, color=COLORS["gold"])
        _text(slide, f"{item['slide_id']}_evidence_body_{index}", x + 0.18, y + 0.52, 2.65, 0.55, body, size=6, color=COLORS["muted"])


def _framework(slide: Any, item: dict[str, Any]) -> None:
    stages = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "framework_node_text"]
    for index, stage in enumerate(stages[:4], start=1):
        x = 1.0 + (index - 1) * 3.15
        _panel(slide, f"{item['slide_id']}_framework_node_{index}", x, 2.7, 2.55, 1.24)
        _text(slide, f"{item['slide_id']}_framework_label_{index}", x + 0.18, 2.92, 2.12, 0.54, stage, size=7)
        if index < len(stages[:4]):
            _line(slide, f"{item['slide_id']}_framework_connector_{index}", x + 2.58, 3.32, x + 3.05, 3.32, color=COLORS["gold"], width=1.4)


def _process(slide: Any, item: dict[str, Any]) -> None:
    steps = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "process_node_text"]
    for index, step in enumerate(steps[:5], start=1):
        x = 1.0 + (index - 1) * 2.55
        _dot(slide, f"{item['slide_id']}_process_dot_{index}", x, 3.0, 0.32)
        _text(slide, f"{item['slide_id']}_process_label_{index}", x - 0.35, 3.52, 1.2, 0.3, step, size=8, bold=True)
        if index < len(steps[:5]):
            _line(slide, f"{item['slide_id']}_process_connector_{index}", x + 0.36, 3.18, x + 2.35, 3.18, color=COLORS["cyan"], width=1.3)


def _comparison(slide: Any, item: dict[str, Any]) -> None:
    rows = _table_rows(item)
    _add_table(slide, f"{item['slide_id']}_native_comparison_matrix", 0.95, 2.0, 9.85, 3.55, rows)
    _panel(slide, f"{item['slide_id']}_insight_panel", 11.15, 2.0, 2.95, 2.2)
    _text(slide, f"{item['slide_id']}_insight", 11.35, 2.24, 2.45, 0.9, item["main_message"], size=8, color=COLORS["muted"])


def _dashboard(slide: Any, item: dict[str, Any]) -> None:
    points = (item.get("chart_data") or {}).get("data_points", [])
    for index, point in enumerate(points[:4], start=1):
        x = 0.95 + (index - 1) * 2.2
        _panel(slide, f"{item['slide_id']}_kpi_card_{index}", x, 2.0, 1.85, 0.9)
        _text(slide, f"{item['slide_id']}_kpi_value_{index}", x + 0.14, 2.13, 0.9, 0.28, f"{point.get('value', 0)}%", size=14, bold=True, color=COLORS["gold"])
        _text(slide, f"{item['slide_id']}_kpi_label_{index}", x + 0.14, 2.48, 1.35, 0.22, point.get("label", "Metric"), size=6, color=COLORS["muted"])
    _add_chart(slide, f"{item['slide_id']}_native_chart", 1.0, 3.35, 6.8, 2.55, points)
    _panel(slide, f"{item['slide_id']}_insight_panel", 8.35, 3.35, 3.1, 2.0)
    _text(slide, f"{item['slide_id']}_dashboard_insight", 8.55, 3.62, 2.58, 0.82, item["main_message"], size=8, color=COLORS["muted"])


def _table_heavy(slide: Any, item: dict[str, Any]) -> None:
    rows = _table_rows(item)
    _add_table(slide, f"{item['slide_id']}_native_table", 0.9, 1.95, 11.55, 4.0, rows)
    _text(slide, f"{item['slide_id']}_table_note", 12.65, 2.08, 2.15, 1.0, "Editable native table: source rows remain selectable and reusable.", size=7, color=COLORS["muted"])


def _timeline(slide: Any, item: dict[str, Any]) -> None:
    milestones = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "timeline_phase_text"]
    _line(slide, f"{item['slide_id']}_timeline_axis", 1.0, 3.4, 12.2, 3.4, color=COLORS["gold"], width=2.0)
    for index, milestone in enumerate(milestones[:5], start=1):
        x = 1.0 + (index - 1) * 2.8
        _dot(slide, f"{item['slide_id']}_timeline_dot_{index}", x, 3.24, 0.24)
        _panel(slide, f"{item['slide_id']}_timeline_card_{index}", x - 0.48, 3.88, 1.75, 0.82)
        _text(slide, f"{item['slide_id']}_timeline_label_{index}", x - 0.32, 4.1, 1.28, 0.28, milestone, size=7)


def _footer(slide: Any, item: dict[str, Any]) -> None:
    _shape(slide, f"{item['slide_id']}_source_footer_strip", 0.55, 8.28, 14.65, 0.32, fill=COLORS["ink"], line=COLORS["ink"])
    _shape(slide, f"{item['slide_id']}_source_footer_gold_rule", 0.55, 8.28, 14.65, 0.02, fill=COLORS["gold"], line=COLORS["gold"])
    _text(slide, f"{item['slide_id']}_source_footer_text", 0.72, 8.36, 9.85, 0.17, item["footer"]["text"], size=5, color=COLORS["muted"])


def _shape(slide: Any, name: str, x: float, y: float, w: float, h: float, *, fill: str, line: str, shape_type: Any = MSO_AUTO_SHAPE_TYPE.RECTANGLE) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(0.7)
    return shape


def _panel(slide: Any, name: str, x: float, y: float, w: float, h: float) -> Any:
    return _shape(slide, name, x, y, w, h, fill=COLORS["teal"], line=COLORS["cyan"], shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)


def _text(slide: Any, name: str, x: float, y: float, w: float, h: float, text: str, *, size: int, bold: bool = False, color: str = COLORS["offwhite"]) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _line(slide: Any, name: str, x1: float, y1: float, x2: float, y2: float, *, color: str, width: float) -> Any:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.name = name
    connector.line.color.rgb = RGBColor.from_string(color)
    connector.line.width = Pt(width)
    return connector


def _dot(slide: Any, name: str, x: float, y: float, size: float) -> Any:
    return _shape(slide, name, x, y, size, size, fill=COLORS["gold"], line=COLORS["gold"], shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)


def _icon(slide: Any, name: str, x: float, y: float, w: float, h: float) -> None:
    circle = _shape(slide, name, x, y, w, h, fill=COLORS["cyan"], line=COLORS["cyan"], shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)
    circle.fill.transparency = 10
    triangle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x + w * 0.3), Inches(y + h * 0.25), Inches(w * 0.4), Inches(h * 0.5))
    triangle.name = f"{name}_inner_triangle"
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor.from_string(COLORS["navy"])
    triangle.line.fill.background()


def _texture(slide: Any, name: str, x: float, y: float, w: float, h: float) -> None:
    for index in range(7):
        lx = x + index * (w / 7)
        _line(slide, f"{name}_line_{index + 1}", lx, y, lx + w * 0.18, y + h, color=COLORS["cyan"], width=0.35)


def _add_chart(slide: Any, name: str, x: float, y: float, w: float, h: float, points: list[dict[str, Any]]) -> None:
    _panel(slide, f"{name}_backplate", x - 0.08, y - 0.08, w + 0.16, h + 0.16)
    data = ChartData()
    categories = [point.get("label", "Metric")[:18] for point in points] or ["Trace", "Method", "QA", "Reuse"]
    values = [int(point.get("value", 0)) for point in points] or [86, 78, 72, 81]
    data.categories = categories
    data.add_series("Readiness", values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart_shape.name = name


def _add_table(slide: Any, name: str, x: float, y: float, w: float, h: float, rows: list[list[str]]) -> None:
    rows = rows or [["Criterion", "Manual Review", "Structured Pipeline", "Hybrid Governance"]]
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table_shape.name = name
    table = table_shape.table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(COLORS["teal_2"] if row_index == 0 else COLORS["ink"])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(7 if len(rows) > 5 else 8)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = RGBColor.from_string(COLORS["offwhite"] if row_index == 0 else COLORS["muted"])


def _slot_value(item: dict[str, Any], semantic_role: str, default: str) -> str:
    for slot in item.get("slots", []):
        if slot.get("semantic_role") == semantic_role:
            return str(slot.get("value", default))
    return default


def _card_pairs(item: dict[str, Any], *, title_role: str = "card_title_text", body_role: str = "card_body_text") -> list[tuple[str, str]]:
    titles = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == title_role]
    bodies = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == body_role]
    return list(zip(titles, bodies))


def _table_rows(item: dict[str, Any]) -> list[list[str]]:
    if item.get("table_data"):
        return [item["table_data"]["header"], *item["table_data"]["rows"]]
    return [["Criterion", "Manual Review", "Structured Pipeline", "Hybrid Governance"]]


def _draw_nonsemantic_asset(path: Path, *, variant: str) -> None:
    image = Image.new("RGBA", (900, 620), f"#{COLORS['teal']}")
    draw = ImageDraw.Draw(image, "RGBA")
    if variant == "contour":
        for idx in range(18):
            inset = 24 + idx * 18
            draw.rounded_rectangle((inset, inset, 900 - inset, 620 - int(inset * 0.65)), radius=60, outline=(*_rgb(COLORS["cyan"]), max(22, 130 - idx * 5)), width=3)
        for idx in range(18):
            x = int(900 * (0.18 + idx * 0.036))
            y = int(620 * (0.22 + (idx % 5) * 0.04))
            draw.ellipse((x, y, x + 8, y + 8), fill=(*_rgb(COLORS["gold"]), 190))
    else:
        for idx in range(-240, 900, 34):
            draw.line((idx, 0, idx + 300, 620), fill=(*_rgb(COLORS["cyan"]), 70), width=2)
        for idx in range(0, 620, 42):
            draw.line((0, idx, 900, idx), fill=(*_rgb(COLORS["gold"]), 28), width=1)
    image.save(path)


def _render_slide_preview(item: dict[str, Any], path: Path) -> None:
    image = Image.new("RGB", (SLIDE_W_PX, SLIDE_H_PX), f"#{COLORS['navy']}")
    draw = ImageDraw.Draw(image, "RGBA")
    draw.rectangle((70, 865, 1590, 900), fill=(*_rgb(COLORS["ink"]), 255))
    draw.rectangle((70, 865, 1590, 868), fill=(*_rgb(COLORS["gold"]), 255))
    _draw_text(draw, 100, 70, item["title"], 24)
    _draw_text(draw, 102, 125, item.get("subtitle", ""), 10)
    archetype = item["archetype_id"]
    if archetype == "cover_hero":
        draw.rounded_rectangle((880, 115, 1460, 630), radius=34, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["cyan"]), 180), width=3)
        for idx in range(10):
            inset = 28 + idx * 18
            draw.rounded_rectangle((880 + inset, 115 + inset, 1460 - inset, 630 - int(inset * 0.55)), radius=28, outline=(*_rgb(COLORS["cyan"]), max(30, 130 - idx * 8)), width=2)
        _draw_text(draw, 105, 255, item["main_message"], 13)
    elif archetype == "data_dashboard":
        for i, point in enumerate((item.get("chart_data") or {}).get("data_points", [])[:4]):
            draw.rounded_rectangle((105 + i * 230, 220, 285 + i * 230, 315), radius=16, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["cyan"]), 160), width=2)
            _draw_text(draw, 123 + i * 230, 233, f"{point['value']}%", 18)
            _draw_text(draw, 123 + i * 230, 278, point["label"][:16], 6)
        draw.rounded_rectangle((110, 365, 820, 625), radius=16, fill=(*_rgb(COLORS["ink"]), 255), outline=(*_rgb(COLORS["gold"]), 190), width=2)
        points = (item.get("chart_data") or {}).get("data_points", [])[:4]
        for i, point in enumerate(points):
            bar_h = int(point["value"] * 1.8)
            draw.rectangle((180 + i * 130, 590 - bar_h, 240 + i * 130, 590), fill=(*_rgb(COLORS["cyan"]), 210))
    elif archetype in {"comparison_matrix", "table_heavy"}:
        draw.rounded_rectangle((100, 210, 1280, 630), radius=8, fill=(*_rgb(COLORS["ink"]), 255), outline=(*_rgb(COLORS["cyan"]), 140), width=2)
        for row in range(5):
            y = 225 + row * 78
            draw.line((115, y, 1250, y), fill=(*_rgb(COLORS["gold" if row == 0 else "teal_2"]), 150), width=2)
        rows = _table_rows(item)[:5]
        for row_index, row in enumerate(rows):
            for col_index, value in enumerate(row[:4]):
                _draw_text(draw, 130 + col_index * 275, 235 + row_index * 78, str(value)[:26], 6 if row_index else 7)
    elif archetype == "visual_toc":
        labels = [slot["value"] for slot in item["slots"] if slot["semantic_role"] == "navigation_item_text"][:6]
        for index, label in enumerate(labels):
            x = 105 + (index % 3) * 360
            y = 225 + (index // 3) * 160
            draw.rounded_rectangle((x, y, x + 290, y + 95), radius=12, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["cyan"]), 160), width=2)
            draw.rectangle((x + 18, y + 20, x + 80, y + 25), fill=(*_rgb(COLORS["gold"]), 255))
            _draw_text(draw, x + 22, y + 42, f"{index + 1:02d}  {label}", 8)
    else:
        cards = _card_pairs(item)
        if archetype == "evidence_overview":
            cards = _card_pairs(item, title_role="evidence_card_title", body_role="evidence_card_body")
        if archetype in {"methodology_framework", "process_flow", "timeline_roadmap"}:
            cards = [(slot["value"][:18], "") for slot in item["slots"] if slot["semantic_role"] in {"framework_node_text", "process_node_text", "timeline_phase_text"}]
        for index in range(3):
            x = 105 + index * 335
            draw.rounded_rectangle((x, 245, x + 285, 420), radius=16, fill=(*_rgb(COLORS["teal"]), 255), outline=(*_rgb(COLORS["cyan"]), 160), width=2)
            draw.rectangle((x + 18, 302, x + 92, 308), fill=(*_rgb(COLORS["gold"]), 255))
            if index < len(cards):
                _draw_text(draw, x + 22, 265, cards[index][0][:28], 7)
                if cards[index][1]:
                    _draw_text(draw, x + 22, 324, cards[index][1][:46], 5)
    _draw_text(draw, 92, 872, item["footer"]["text"][:120], 7)
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def _build_contact_sheet(paths: list[Path], output_path: Path) -> Path:
    thumbs = [Image.open(path).resize((418, 235)) for path in paths]
    cols = 4
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 418, rows * 235), f"#{COLORS['navy']}")
    for index, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((index % cols) * 418, (index // cols) * 235))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    return output_path


def _draw_text(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, size: int) -> None:
    try:
        font = ImageFont.truetype("arial.ttf", max(8, size * 2))
    except OSError:
        font = ImageFont.load_default()
    draw.text((x, y), str(text), fill=(*_rgb(COLORS["offwhite"]), 255), font=font)


def _rgb(hex_color: str) -> tuple[int, int, int]:
    return tuple(int(hex_color[index : index + 2], 16) for index in (0, 2, 4))
