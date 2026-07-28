"""Compile isolated E02 editable candidates for the four core archetypes."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = 16.0
SLIDE_H = 9.0

COLORS = {
    "background": "06111A",
    "panel": "0B3B44",
    "panel2": "103F4A",
    "cyan": "50D2E5",
    "cyan_dim": "1B8EA3",
    "gold": "F5A623",
    "paper": "F7F1E4",
    "ink": "071018",
    "muted": "9EC4C8",
    "coral": "EF6B5A",
    "white": "FFFFFF",
}


def compile_e02_candidate(archetype_id: str, output_pptx: Path) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_archetype_slide(slide, archetype_id)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_name": "e02_candidate_compile_report",
        "status": "passed",
        "archetype_id": archetype_id,
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 1,
        **candidate_expected_counts(archetype_id),
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use": False,
    }


def compile_e02_candidate_pack(archetype_ids: list[str], output_pptx: Path) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for archetype_id in archetype_ids:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        draw_archetype_slide(slide, archetype_id)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_name": "e02_candidate_pack_compile_report",
        "status": "passed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": len(archetype_ids),
        "archetypes": archetype_ids,
        "non_canonical": True,
        "canonical_promotion": False,
    }


def candidate_expected_counts(archetype_id: str) -> dict[str, int]:
    return {
        "cover_hero": {"semantic_vector_icon_count": 1, "native_or_editable_chart_count": 0, "native_or_editable_table_count": 0},
        "standard_content": {"semantic_vector_icon_count": 4, "native_or_editable_chart_count": 0, "native_or_editable_table_count": 0},
        "data_dashboard": {"semantic_vector_icon_count": 4, "native_or_editable_chart_count": 1, "native_or_editable_table_count": 0},
        "table_heavy": {"semantic_vector_icon_count": 2, "native_or_editable_chart_count": 0, "native_or_editable_table_count": 1},
    }[archetype_id]


def required_semantic_slots(archetype_id: str) -> list[str]:
    return {
        "cover_hero": ["title", "subtitle", "hero_visual_field", "meta_bar", "source_footer_strip"],
        "standard_content": ["title", "content_card_group", "body_text_regions", "insight_or_takeaway", "source_footer_strip"],
        "data_dashboard": ["title", "kpi_cards", "primary_chart", "secondary_chart_or_insight_panel", "source_footer_strip"],
        "table_heavy": ["title", "table_region", "header_band", "row_groups", "source_footer_strip", "optional_kpi_chips"],
    }[archetype_id]


def draw_archetype_slide(slide: Any, archetype_id: str) -> None:
    _background(slide)
    if archetype_id == "cover_hero":
        _draw_cover_hero(slide)
    elif archetype_id == "standard_content":
        _draw_standard_content(slide)
    elif archetype_id == "data_dashboard":
        _draw_data_dashboard(slide)
    elif archetype_id == "table_heavy":
        _draw_table_heavy(slide)
    else:
        raise ValueError(archetype_id)


def _background(slide: Any) -> None:
    bg = _rect(slide, "background_base_native", 0, 0, SLIDE_W, SLIDE_H, "background", None)
    bg.line.fill.background()
    for idx, x in enumerate((0.6, 2.0, 13.7, 15.0), start=1):
        line = _line(slide, f"technical_overlay_grid_line_{idx}", x, 0.35, x + 0.35, 7.6, "cyan_dim", 0.4)
        line.line.transparency = 55


def _draw_cover_hero(slide: Any) -> None:
    _text(slide, "cover_title_text", "TITLE", 0.72, 1.35, 5.6, 0.85, 42, "paper", bold=True)
    _text(slide, "cover_subtitle_text", "SUBTITLE / VALUE PROMISE", 0.76, 3.15, 5.4, 0.6, 20, "gold")
    _rect(slide, "cover_meta_bar_native", 0.72, 5.68, 4.5, 0.44, "panel", "cyan")
    _text(slide, "cover_meta_text", "META / DATE / OWNER", 0.95, 5.78, 4.0, 0.2, 10, "muted", bold=True)
    _rect(slide, "cover_hero_visual_field_frame", 8.05, 0.78, 6.8, 6.7, "panel2", "cyan")
    _rect(slide, "cover_hero_replaceable_image_frame", 8.45, 1.18, 6.0, 5.9, "background", "gold")
    _text(slide, "cover_hero_visual_slot_label", "REPLACEABLE HERO VISUAL FIELD", 9.0, 3.9, 5.0, 0.3, 11, "muted", bold=True)
    _draw_icon(slide, "cover_meta_vector_icon", 5.36, 5.72, 0.36, "gold", "target")
    _footer(slide, "SOURCE / FOOTER")


def _draw_standard_content(slide: Any) -> None:
    _text(slide, "standard_title_text", "TITLE", 0.72, 0.52, 6.2, 0.45, 26, "paper", bold=True)
    card_data = [
        ("01", "CARD HEADING", "Editable body text region with capacity for a concise content module."),
        ("02", "CARD HEADING", "Card panels are PPT-native shapes with editable labels and vector icons."),
        ("03", "CARD HEADING", "Observed content groups remain separate semantic layer objects."),
        ("04", "CARD HEADING", "Footer and source strips stay editable and protected."),
    ]
    for idx, (num, heading, body) in enumerate(card_data):
        x = 0.82 + (idx % 2) * 5.05
        y = 1.55 + (idx // 2) * 2.45
        _rect(slide, f"standard_card_{idx+1}_panel_native", x, y, 4.65, 1.9, "paper", "cyan_dim")
        _draw_icon(slide, f"standard_card_{idx+1}_vector_icon", x + 0.24, y + 0.28, 0.42, "gold", "document")
        _text(slide, f"standard_card_{idx+1}_number_text", num, x + 0.78, y + 0.2, 0.5, 0.3, 10, "cyan_dim", bold=True)
        _text(slide, f"standard_card_{idx+1}_heading_text", heading, x + 1.18, y + 0.22, 2.9, 0.28, 12, "ink", bold=True)
        _text(slide, f"standard_card_{idx+1}_body_text", body, x + 0.78, y + 0.75, 3.45, 0.75, 10.5, "ink")
    _rect(slide, "standard_insight_panel_native", 11.25, 1.55, 3.7, 4.36, "panel", "gold")
    _text(slide, "standard_insight_text", "INSIGHT / TAKEAWAY", 11.65, 2.08, 2.8, 0.8, 16, "paper", bold=True)
    _text(slide, "standard_insight_body_text", "Editable takeaway region with protected capacity and no raster fallback.", 11.65, 3.12, 2.75, 1.4, 12, "muted")
    _footer(slide, "SOURCE / FOOTER")


def _draw_data_dashboard(slide: Any) -> None:
    _text(slide, "dashboard_title_text", "TITLE", 0.72, 0.45, 6.2, 0.45, 26, "paper", bold=True)
    kpis = [("KPI", "42"), ("KPI", "16"), ("KPI", "11"), ("KPI", "07")]
    for idx, (label, value) in enumerate(kpis):
        x = 0.82 + idx * 3.55
        _rect(slide, f"dashboard_kpi_{idx+1}_card_native", x, 1.28, 3.0, 1.1, "panel", "cyan_dim")
        _draw_icon(slide, f"dashboard_kpi_{idx+1}_vector_icon", x + 0.24, 1.62, 0.32, "gold", "gauge")
        _text(slide, f"dashboard_kpi_{idx+1}_label_text", label, x + 0.75, 1.48, 1.0, 0.25, 9, "muted", bold=True)
        _text(slide, f"dashboard_kpi_{idx+1}_value_text", value, x + 0.75, 1.74, 1.1, 0.45, 22, "paper", bold=True)
    _rect(slide, "dashboard_primary_chart_panel_native", 0.82, 2.9, 9.2, 3.85, "paper", "cyan_dim")
    _text(slide, "dashboard_chart_title_text", "PRIMARY CHART", 1.15, 3.18, 2.4, 0.25, 11, "ink", bold=True)
    _bar_chart(slide, 1.3, 3.78, 7.6, 2.25)
    _rect(slide, "dashboard_insight_panel_native", 10.65, 2.9, 4.4, 3.85, "panel", "gold")
    _text(slide, "dashboard_insight_heading_text", "INSIGHT PANEL", 11.05, 3.25, 3.3, 0.35, 14, "paper", bold=True)
    _text(slide, "dashboard_insight_body_text", "Editable commentary and secondary signal area. Template-stage data is explicit and replaceable.", 11.05, 3.92, 3.35, 1.35, 12, "muted")
    _footer(slide, "SOURCE / FOOTER")


def _draw_table_heavy(slide: Any) -> None:
    _text(slide, "table_title_text", "TITLE", 0.72, 0.45, 6.2, 0.45, 26, "paper", bold=True)
    _rect(slide, "table_region_outer_panel_native", 0.82, 1.45, 14.3, 5.85, "paper", "cyan_dim")
    columns = ["TABLE", "OWNER", "SIGNAL", "CADENCE", "DECISION", "SOURCE"]
    rows = [["ROW", "OWNER", "SIGNAL", "CADENCE", "DECISION", "CITE"] for _ in range(5)]
    _shape_grid_table(slide, 1.04, 1.75, 13.85, 4.95, columns, rows)
    _rect(slide, "table_optional_kpi_chips_native", 1.0, 7.0, 13.9, 0.38, "panel", "gold")
    _text(slide, "table_optional_kpi_chips_text", "OPTIONAL KPI CHIPS / NOTES", 1.25, 7.08, 5.0, 0.18, 9, "paper", bold=True)
    _draw_icon(slide, "table_source_vector_icon", 14.25, 7.05, 0.22, "gold", "database")
    _footer(slide, "SOURCE / FOOTER")


def _footer(slide: Any, text: str) -> None:
    _rect(slide, "source_footer_strip_native", 0, 8.36, SLIDE_W, 0.64, "panel", None)
    _line(slide, "source_footer_top_rule", 0, 8.36, SLIDE_W, 8.36, "gold", 1.1)
    _text(slide, "source_footer_text_editable", text, 0.55, 8.58, 5.5, 0.2, 8, "muted")


def _bar_chart(slide: Any, x: float, y: float, w: float, h: float) -> None:
    values = [0.42, 0.54, 0.63, 0.72, 0.86]
    labels = ["A", "B", "C", "D", "E"]
    _line(slide, "dashboard_chart_axis_y", x, y, x, y + h, "ink", 0.7)
    _line(slide, "dashboard_chart_axis_x", x, y + h, x + w, y + h, "ink", 0.7)
    bar_w = w / 7
    for idx, value in enumerate(values):
        bx = x + 0.55 + idx * (bar_w + 0.42)
        bh = h * value
        _rect(slide, f"dashboard_primary_chart_bar_{idx+1}_editable_shape", bx, y + h - bh, bar_w, bh, "cyan_dim" if idx < 4 else "gold", None)
        _text(slide, f"dashboard_primary_chart_label_{idx+1}", labels[idx], bx + 0.12, y + h + 0.12, 0.35, 0.18, 8, "ink")


def _shape_grid_table(slide: Any, x: float, y: float, w: float, h: float, columns: list[str], rows: list[list[str]]) -> None:
    col_w = w / len(columns)
    row_h = h / (len(rows) + 1)
    for c_idx, col in enumerate(columns):
        _rect(slide, f"table_header_cell_{c_idx+1}_native", x + c_idx * col_w, y, col_w, row_h, "panel", "paper")
        _text(slide, f"table_header_text_{c_idx+1}", col, x + c_idx * col_w + 0.08, y + 0.13, col_w - 0.16, 0.18, 7.5, "paper", bold=True)
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            fill = "white" if r_idx % 2 == 0 else "paper"
            _rect(slide, f"table_row_{r_idx+1}_cell_{c_idx+1}_native", x + c_idx * col_w, y + (r_idx + 1) * row_h, col_w, row_h, fill, "muted")
            _text(slide, f"table_row_{r_idx+1}_cell_{c_idx+1}_text", value, x + c_idx * col_w + 0.08, y + (r_idx + 1) * row_h + 0.13, col_w - 0.16, 0.18, 7, "ink")


def _draw_icon(slide: Any, name: str, x: float, y: float, size: float, color: str, role: str) -> None:
    if role in {"document", "database"}:
        _rect(slide, f"{name}_vector_body", x, y, size * 0.74, size, None, color)
        _line(slide, f"{name}_vector_line_1", x + size * 0.18, y + size * 0.34, x + size * 0.58, y + size * 0.34, color, 0.8)
        _line(slide, f"{name}_vector_line_2", x + size * 0.18, y + size * 0.54, x + size * 0.58, y + size * 0.54, color, 0.8)
    elif role == "gauge":
        _oval(slide, f"{name}_vector_ring", x, y, size, size, None, color)
        _line(slide, f"{name}_vector_needle", x + size * 0.5, y + size * 0.55, x + size * 0.78, y + size * 0.32, color, 1.0)
    else:
        _oval(slide, f"{name}_vector_ring", x, y, size, size, None, color)
        _line(slide, f"{name}_vector_mark_1", x + size * 0.25, y + size * 0.52, x + size * 0.45, y + size * 0.70, color, 1.0)
        _line(slide, f"{name}_vector_mark_2", x + size * 0.45, y + size * 0.70, x + size * 0.78, y + size * 0.28, color, 1.0)


def _rect(slide: Any, name: str, x: float, y: float, w: float, h: float, fill: str | None, line: str | None) -> Any:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def _oval(slide: Any, name: str, x: float, y: float, w: float, h: float, fill: str | None, line: str | None) -> Any:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.9)
    else:
        shape.line.fill.background()
    return shape


def _line(slide: Any, name: str, x1: float, y1: float, x2: float, y2: float, color: str, width: float) -> Any:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _text(slide: Any, name: str, text: str, x: float, y: float, w: float, h: float, size: float, color: str, *, bold: bool = False) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _rgb(name_or_hex: str) -> RGBColor:
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
