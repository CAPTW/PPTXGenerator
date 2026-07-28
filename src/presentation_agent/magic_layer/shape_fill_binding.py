"""Bind accepted D07.2 visual assets into PPT image-frame slots."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from .source_bound_template_binder import SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN


def build_shape_fill_patch_plan(slot_file_map: dict[str, Any], validation_report: dict[str, Any]) -> dict[str, Any]:
    accepted = {asset["slot_id"]: asset for asset in validation_report.get("accepted_assets") or []}
    entries: list[dict[str, Any]] = []
    for slot in slot_file_map.get("entries") or []:
        asset = accepted.get(slot["slot_id"])
        entries.append(
            {
                "slot_id": slot["slot_id"],
                "slide_id": slot["slide_id"],
                "slide_number": slot.get("slide_number"),
                "archetype_id": slot["archetype_id"],
                "role": slot["role"],
                "bbox_norm": slot["bbox_norm"],
                "patch_status": "ready" if asset else "blocked_missing_or_rejected_asset",
                "asset_path": asset.get("processed_asset_path") if asset else None,
                "insert_mode": "picture_in_existing_visual_field_slot",
                "full_slide_background": False,
                "screenshot_slide": False,
                "semantic_raster_target": False,
                "preserve_native_overlays_and_text": True,
            }
        )
    ready_count = sum(1 for entry in entries if entry["patch_status"] == "ready")
    return {
        "schema_name": "shape_fill_patch_plan",
        "status": "passed" if ready_count == len(entries) and entries else "blocked",
        "slot_count": len(entries),
        "ready_patch_count": ready_count,
        "blocked_patch_count": len(entries) - ready_count,
        "entries": entries,
        "canva_parity_claimed": False,
    }


def patch_deck_with_shape_fill_assets(input_deck: Path, output_deck: Path, patch_plan: dict[str, Any]) -> dict[str, Any]:
    ready = [entry for entry in patch_plan.get("entries") or [] if entry.get("patch_status") == "ready"]
    if patch_plan.get("status") != "passed" or not ready:
        return {
            "schema_name": "shape_fill_deck_patch_report",
            "status": "blocked",
            "deck_created": False,
            "reason": "Not all required visual assets are accepted.",
            "inserted_asset_count": 0,
        }
    prs = Presentation(input_deck)
    inserted: list[dict[str, Any]] = []
    for entry in ready:
        slide_index = int(entry["slide_number"]) - 1
        slide = prs.slides[slide_index]
        bbox = entry["bbox_norm"]
        picture = slide.shapes.add_picture(
            entry["asset_path"],
            Inches(float(bbox[0]) * SLIDE_WIDTH_IN),
            Inches(float(bbox[1]) * SLIDE_HEIGHT_IN),
            Inches(float(bbox[2]) * SLIDE_WIDTH_IN),
            Inches(float(bbox[3]) * SLIDE_HEIGHT_IN),
        )
        picture.name = f"d07_2_visual_asset_{entry['slot_id']}"
        _move_shape_after_slot(slide, picture, entry["slot_id"])
        inserted.append(
            {
                "slot_id": entry["slot_id"],
                "slide_number": entry["slide_number"],
                "asset_path": entry["asset_path"],
                "bbox_norm": bbox,
                "full_slide_background": False,
                "semantic_raster_target": False,
            }
        )
    output_deck.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_deck)
    return {
        "schema_name": "shape_fill_deck_patch_report",
        "status": "passed",
        "deck_created": output_deck.exists(),
        "input_deck": input_deck.as_posix(),
        "output_deck": output_deck.as_posix(),
        "inserted_asset_count": len(inserted),
        "inserted_assets": inserted,
    }


def shape_fill_alignment_report_from_plan(patch_plan: dict[str, Any], validation_report: dict[str, Any]) -> dict[str, Any]:
    findings: list[dict[str, Any]] = []
    for entry in patch_plan.get("entries") or []:
        bbox = entry.get("bbox_norm") or [0, 0, 0, 0]
        area = float(bbox[2]) * float(bbox[3]) if len(bbox) == 4 else 1
        if area >= 0.85:
            findings.append({"slot_id": entry["slot_id"], "issue": "visual asset slot is effectively full slide"})
        if entry.get("semantic_raster_target"):
            findings.append({"slot_id": entry["slot_id"], "issue": "visual asset targets semantic raster region"})
    accepted = validation_report.get("accepted_asset_count", 0)
    return {
        "schema_name": "shape_fill_alignment_report",
        "status": "passed" if accepted and not findings and patch_plan.get("status") == "passed" else "blocked" if not accepted else "failed",
        "checked_asset_count": accepted,
        "alignment_issue_count": len(findings),
        "findings": findings,
    }


def _move_shape_after_slot(slide: Any, shape: Any, slot_id: str) -> None:
    """Move inserted picture immediately after its slot placeholder in z-order."""

    sp_tree = slide.shapes._spTree  # noqa: SLF001 - python-pptx has no public z-order API.
    picture_element = shape._element  # noqa: SLF001
    target_index: int | None = None
    for index, candidate in enumerate(sp_tree):
        name = candidate.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
        if name is not None and name.get("name") == slot_id:
            target_index = index
            break
    if target_index is None:
        return
    sp_tree.remove(picture_element)
    sp_tree.insert(target_index + 1, picture_element)
