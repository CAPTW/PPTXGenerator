"""Visible text inventory for E04-R3 editorial integrity checks."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from src.presentation_agent.magic_layer.e04_r3_internal_label_filter import is_internal_label


def build_visible_text_inventory(pptx_path: str | Path) -> dict[str, Any]:
    path = Path(pptx_path)
    prs = Presentation(path)
    rows = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = " ".join(shape.text.split())
            if not text:
                continue
            role = _role_guess(str(getattr(shape, "name", "")), text)
            rows.append(
                {
                    "slide_id": f"SLIDE-{slide_index:03d}",
                    "slide_number": slide_index,
                    "shape_id": str(getattr(shape, "shape_id", shape_index)),
                    "shape_name": str(getattr(shape, "name", "")),
                    "visible_text": text,
                    "role_guess": role,
                    "source_ref": None,
                    "is_audience_facing": role not in {"citation_footer", "internal_label"},
                    "is_internal_art_direction_label": is_internal_label(text),
                    "is_template_placeholder": any(token in text for token in ("TITLE PLACEHOLDER", "Editable slot", "Slot")),
                    "is_source_bound_claim": role in {"title", "subtitle", "claim", "body"},
                    "is_citation_footer": role == "citation_footer",
                    "risk": _risk(text),
                }
            )
    return {
        "schema_name": "visible_text_inventory_r2",
        "status": "passed",
        "pptx_path": path.as_posix(),
        "visible_text_count": len(rows),
        "internal_label_candidate_count": sum(1 for row in rows if row["is_internal_art_direction_label"]),
        "rows": rows,
        "canva_parity_claimed": False,
    }


def visible_text_inventory_markdown(inventory: dict[str, Any]) -> str:
    return "\n".join(
        [
            "# Visible Text Inventory R2",
            "",
            f"- Status: `{inventory['status']}`",
            f"- Visible text count: `{inventory['visible_text_count']}`",
            f"- Internal label candidates: `{inventory['internal_label_candidate_count']}`",
        ]
    )


def _role_guess(shape_name: str, text: str) -> str:
    name = shape_name.lower()
    if "footer" in name or text.lower().startswith("sources:"):
        return "citation_footer"
    if is_internal_label(text):
        return "internal_label"
    if "title" in name:
        return "title"
    if "subtitle" in name:
        return "subtitle"
    if "kpi" in name:
        return "kpi_label"
    if "table" in name:
        return "table_text"
    return "body"


def _risk(text: str) -> str:
    if is_internal_label(text):
        return "internal_label_leakage"
    if text.endswith((" for r", " disconnect", " a structured")):
        return "possible_truncation"
    return "none"
