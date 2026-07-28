"""Read-only E04 source-bound deck design quality gate."""

from __future__ import annotations

import json
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation

from scripts.run_e01x_self_describing_ps_layer_integration import protected_report, protected_snapshot, run_protect_check
from src.presentation_agent.magic_layer.e04_focal_object_gate import build_focal_object_report
from src.presentation_agent.magic_layer.e04_r2_art_direction import load_e04_r2_art_direction_plan
from src.presentation_agent.magic_layer.e04_skeleton_similarity import build_skeleton_similarity_report
from src.presentation_agent.magic_layer.e04_slide_rhythm import build_slide_rhythm_report
from src.presentation_agent.magic_layer.e04_visual_hierarchy_gate import build_visual_hierarchy_report


REPO_ROOT = Path(__file__).resolve().parents[3]
E04_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e04_source_bound_small_deck_with_e03_r2_pack"
E04_DQ_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e04_dq_source_bound_design_quality_gate"
EXPECTED_E04_DECISION = "E04_PASS_SOURCE_BOUND_SMALL_DECK_WITH_E03_R2_PACK"


def run_e04_design_quality_gate(e04_root: str | Path = E04_ROOT, output_dir: str | Path = E04_DQ_ROOT) -> dict[str, Any]:
    input_root = Path(e04_root)
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    protected_before = protected_snapshot()
    protect_pre = run_protect_check()
    if not protect_pre:
        final = _override("E04_DQ_FAIL_PROTECTED_ARTIFACTS", False, False, "protected artifact precheck failed", {})
        _write_json(output / "e04_design_quality_override.json", final)
        _write_md(output / "e04_design_quality_override.md", _simple_md("E04 Design Quality Override", final))
        return final

    structural = _structural_status(input_root)
    skeleton = build_skeleton_similarity_report(input_root)
    rhythm = build_slide_rhythm_report(input_root, skeleton)
    focal = build_focal_object_report(input_root)
    hierarchy = build_visual_hierarchy_report(input_root)
    complexity = build_object_complexity_vs_design_quality_report(input_root, hierarchy)
    interpretation = build_source_content_visual_interpretation_report(input_root)
    premium_pass = all(
        report["status"] == "passed"
        for report in (skeleton, rhythm, focal, hierarchy, complexity, interpretation)
    )
    decision = "E04_DQ_PASS_READY_FOR_E04_R2_OR_E05" if premium_pass else _patch_decision(skeleton, hierarchy, interpretation)
    override = _override(
        decision,
        structural["source_bound_structural_pass"],
        premium_pass,
        "source-bound structural pass is preserved, but premium deck design quality is blocked" if not premium_pass else "premium deck design quality passed",
        structural,
    )
    e05 = {
        "schema_name": "e05_readiness_override",
        "status": "blocked" if not premium_pass else "passed",
        "e05_unlocked": premium_pass,
        "reason": "premium deck design quality must pass before E05" if not premium_pass else "E04-DQ passed",
        "structural_decision": structural["structural_decision"],
        "design_quality_decision": decision,
        "canva_parity_claimed": False,
    }
    design_report = {
        "schema_name": "source_bound_deck_design_quality_report",
        "status": "passed" if premium_pass else "failed",
        "decision": decision,
        "source_bound_structural_pass": structural["source_bound_structural_pass"],
        "semantic_editability_pass": structural["semantic_editability_pass"],
        "premium_deck_design_quality_pass": premium_pass,
        "e05_unlocked": premium_pass,
        "skeleton_similarity_status": skeleton["status"],
        "slide_rhythm_status": rhythm["status"],
        "focal_object_status": focal["status"],
        "visual_hierarchy_status": hierarchy["status"],
        "object_complexity_vs_design_quality_status": complexity["status"],
        "source_content_visual_interpretation_status": interpretation["status"],
        "canva_parity_claimed": False,
    }
    manifest = {
        "schema_name": "e04_dq_manifest",
        "generated_at": _now(),
        "input_folder": _rel(input_root),
        "output_folder": _rel(output),
        "original_e04_report_modified": False,
        "new_deck_generated": False,
        "large_deck_generated": False,
        "e05_started": False,
        "canonical_promotion": False,
        "canva_parity_claimed": False,
        "decision": decision,
    }
    reports = {
        "e04_dq_manifest.json": manifest,
        "e04_design_quality_override.json": override,
        "source_bound_deck_design_quality_report.json": design_report,
        "skeleton_similarity_report.json": skeleton,
        "slide_rhythm_report.json": rhythm,
        "focal_object_report.json": focal,
        "visual_hierarchy_report.json": hierarchy,
        "object_complexity_vs_design_quality_report.json": complexity,
        "source_content_visual_interpretation_report.json": interpretation,
        "e05_readiness_override.json": e05,
    }
    for filename, payload in reports.items():
        _write_json(output / filename, payload)
    for filename in [
        "e04_design_quality_override",
        "source_bound_deck_design_quality_report",
        "skeleton_similarity_report",
        "slide_rhythm_report",
        "focal_object_report",
        "visual_hierarchy_report",
        "object_complexity_vs_design_quality_report",
        "source_content_visual_interpretation_report",
    ]:
        _write_md(output / f"{filename}.md", _simple_md(filename.replace("_", " ").title(), reports[f"{filename}.json"]))

    protected_after = protected_snapshot()
    protected_md, protected_ok = protected_report(protected_before, protected_after)
    protect_post = run_protect_check()
    protected_md += f"\n\n- npm protect precheck: `passed`\n- npm protect postcheck: `{'passed' if protect_post else 'failed'}`\n"
    if not protected_ok or not protect_post:
        override = _override("E04_DQ_FAIL_PROTECTED_ARTIFACTS", structural["source_bound_structural_pass"], False, "protected artifact postcheck failed", structural)
        _write_json(output / "e04_design_quality_override.json", override)
        _write_md(output / "e04_design_quality_override.md", _simple_md("E04 Design Quality Override", override))
    _write_md(output / "protected_artifact_check_report.md", protected_md)
    return override


def build_object_complexity_vs_design_quality_report(e04_root: str | Path, hierarchy_report: dict[str, Any]) -> dict[str, Any]:
    root = Path(e04_root)
    pptx_path = root / "source_bound_sample_deck_12_16.pptx"
    if not pptx_path.exists():
        pptx_path = root / "source_bound_sample_deck_r2_12_16.pptx"
    if not pptx_path.exists():
        pptx_path = root / "source_bound_sample_deck_r3_12_16.pptx"
    prs = Presentation(pptx_path)
    rows = []
    for index, slide in enumerate(prs.slides, start=1):
        shape_types = Counter(str(shape.shape_type) for shape in slide.shapes)
        shape_count = sum(shape_types.values())
        connector_count = sum(count for key, count in shape_types.items() if "LINE" in key)
        rows.append(
            {
                "slide_number": index,
                "object_count": shape_count,
                "connector_or_line_count": connector_count,
                "design_quality_warning": connector_count > 6 or shape_count > 22,
            }
        )
    warning_count = sum(1 for row in rows if row["design_quality_warning"])
    hierarchy_failed = hierarchy_report["status"] != "passed"
    return {
        "schema_name": "object_complexity_vs_design_quality_report",
        "status": "failed" if hierarchy_failed else "passed",
        "deck_object_count": sum(row["object_count"] for row in rows),
        "warning_slide_count": warning_count,
        "principle": "object count is reported but not rewarded as design quality",
        "finding": "structural complexity increased without enough hierarchy improvement" if hierarchy_failed else "object complexity supports hierarchy",
        "slides": rows,
        "canva_parity_claimed": False,
    }


def build_source_content_visual_interpretation_report(e04_root: str | Path) -> dict[str, Any]:
    root = Path(e04_root)
    art_plan = load_e04_r2_art_direction_plan(root)
    if art_plan:
        rows = [
            {
                "slide_id": slide["slide_id"],
                "slide_number": slide["slide_number"],
                "archetype_id": slide["archetype_id"],
                "interpretation_status": slide["source_visual_interpretation_status"],
                "finding": slide["source_visual_interpretation"],
            }
            for slide in art_plan.get("slides", [])
        ]
        failures = [row for row in rows if row["interpretation_status"] != "passed"]
        return {
            "schema_name": "source_content_visual_interpretation_report",
            "status": "failed" if failures else "passed",
            "slide_count": len(rows),
            "failure_count": len(failures),
            "failures": failures,
            "slides": rows,
            "art_direction_plan_used": True,
            "canva_parity_claimed": False,
        }
    blueprints = _read_json(root / "slide_blueprint_v1.json")
    rows = []
    weak_archetypes = {"standard_content", "evidence_overview", "card_grid", "methodology_framework", "process_flow", "timeline_roadmap"}
    for slide in blueprints.get("slides", []):
        archetype = slide["archetype_id"]
        passed = archetype not in weak_archetypes
        rows.append(
            {
                "slide_id": slide["slide_id"],
                "slide_number": slide["slide_number"],
                "archetype_id": archetype,
                "interpretation_status": "passed" if passed else "failed",
                "finding": _interpretation_finding(archetype),
            }
        )
    failures = [row for row in rows if row["interpretation_status"] != "passed"]
    return {
        "schema_name": "source_content_visual_interpretation_report",
        "status": "failed" if failures else "passed",
        "slide_count": len(rows),
        "failure_count": len(failures),
        "failures": failures,
        "slides": rows,
        "canva_parity_claimed": False,
    }


def _structural_status(e04_root: Path) -> dict[str, Any]:
    if (e04_root / "e04_final_decision.json").exists():
        final = _read_json(e04_root / "e04_final_decision.json")
    else:
        final = _read_json(e04_root / "e04_r2_final_decision.json")
    editability = _read_json(e04_root / "semantic_editability_ledger.json")
    raster = _read_json(e04_root / "semantic_raster_violation_report.json")
    citations = _read_json(e04_root / "citation_coverage_report.json")
    return {
        "structural_decision": final.get("decision"),
        "source_bound_structural_pass": final.get("decision") in {EXPECTED_E04_DECISION, "E04_R2_PASS_READY_FOR_E05_SOURCE_BOUND_SMALL_DECK"} and final.get("status") == "passed",
        "semantic_editability_pass": editability.get("status") == "passed",
        "semantic_raster_violation_count": raster.get("semantic_raster_violation_count", 0),
        "citation_coverage_status": citations.get("status"),
    }


def _override(
    decision: str,
    structural_pass: bool,
    premium_pass: bool,
    reason: str,
    structural: dict[str, Any],
) -> dict[str, Any]:
    return {
        "schema_name": "e04_design_quality_override",
        "status": "passed" if premium_pass else "failed",
        "decision": decision,
        "structural_decision": structural.get("structural_decision", EXPECTED_E04_DECISION),
        "source_bound_structural_pass": structural_pass,
        "semantic_editability_pass": structural.get("semantic_editability_pass", False),
        "premium_deck_design_quality_pass": premium_pass,
        "visual_quality_decision": "E04_PREMIUM_DECK_DESIGN_QUALITY_PASSED" if premium_pass else "E04_PATCH_PREMIUM_DECK_DESIGN_REQUIRED",
        "e05_unlocked": premium_pass,
        "reason": reason,
        "canva_parity_claimed": False,
    }


def _patch_decision(
    skeleton: dict[str, Any],
    hierarchy: dict[str, Any],
    interpretation: dict[str, Any],
) -> str:
    if skeleton["status"] == "failed" and hierarchy["status"] == "failed":
        return "E04_DQ_PATCH_DECK_ART_DIRECTION_REQUIRED"
    if hierarchy["status"] == "failed":
        return "E04_DQ_PATCH_VISUAL_HIERARCHY_REQUIRED"
    if interpretation["status"] == "failed":
        return "E04_DQ_PATCH_COMPONENT_RECIPES_REQUIRED"
    return "E04_DQ_PATCH_LAYOUT_SELECTOR_REQUIRED"


def _interpretation_finding(archetype_id: str) -> str:
    return {
        "cover_hero": "hero field gives the opening a visual anchor",
        "visual_toc": "navigation labels form a sequence system",
        "comparison_matrix": "source comparison is bound into a readable table/matrix",
        "data_dashboard": "metric values are plotted visually",
        "table_heavy": "source rows are represented as a native table",
        "standard_content": "problem content still reads as equal cards rather than a risk/emphasis composition",
        "evidence_overview": "evidence is carded but lacks a strong claim/proof/source hierarchy",
        "card_grid": "source artifacts are arranged as generic cards",
        "methodology_framework": "framework stages are not diagrammatic enough",
        "process_flow": "sequence is present but visual flow is too weak",
        "timeline_roadmap": "roadmap sequence is present but the timeline rail is under-emphasized",
    }.get(archetype_id, "generic source interpretation")


def _read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def _write_md(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.rstrip() + "\n", encoding="utf-8")


def _simple_md(title: str, payload: dict[str, Any]) -> str:
    lines = [f"# {title}", "", f"- Status: `{payload.get('status', 'n/a')}`"]
    for key in (
        "decision",
        "structural_decision",
        "source_bound_structural_pass",
        "semantic_editability_pass",
        "premium_deck_design_quality_pass",
        "e05_unlocked",
        "near_identical_body_composition_ratio",
        "equal_weight_card_rhythm_count",
        "weak_focal_object_count",
        "average_visual_hierarchy_score",
        "failure_count",
        "reason",
    ):
        if key in payload:
            lines.append(f"- {key}: `{payload[key]}`")
    if payload.get("failures"):
        lines.append("")
        lines.append("## Findings")
        for item in payload["failures"]:
            lines.append(f"- {item if isinstance(item, str) else item.get('finding', item)}")
    return "\n".join(lines)


def _rel(path: Path) -> str:
    try:
        return path.resolve().relative_to(REPO_ROOT).as_posix()
    except ValueError:
        return path.resolve().as_posix()


def _now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()
