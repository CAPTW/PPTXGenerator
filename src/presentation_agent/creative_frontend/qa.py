"""Semantic, creative, and editability QA for creative-front-end PPTX runs."""

from __future__ import annotations

import argparse
import json
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..generator_contracts import (
    validateCreativeTemplateArchitecture,
    validatePresentationArchitecture,
    validateSlideSemanticSidecar,
)
from .semantic_validation import (
    validate_creative_template_architecture_semantics,
    validate_presentation_architecture_semantics,
    validate_sidecar_semantics,
)


EMU_PER_INCH = 914400
FULL_SLIDE_AREA_RATIO = 0.88


def build_creative_frontend_qa_report(
    *,
    presentation_architecture: dict[str, Any],
    creative_template_architecture: dict[str, Any],
    editable_template_spec: dict[str, Any],
    sidecars: list[dict[str, Any]],
    pptx_path: str | Path,
    render_manifest: dict[str, Any] | None = None,
    image_policy_report: dict[str, Any] | None = None,
) -> dict[str, Any]:
    validatePresentationArchitecture(presentation_architecture)
    validateCreativeTemplateArchitecture(creative_template_architecture)
    for sidecar in sidecars:
        validateSlideSemanticSidecar(sidecar)
    validate_presentation_architecture_semantics(presentation_architecture)
    validate_creative_template_architecture_semantics(
        creative_template_architecture,
        presentation_architecture,
        editable_template_spec,
    )
    validate_sidecar_semantics(sidecars, presentation_architecture, creative_template_architecture)

    pptx_file = Path(pptx_path)
    deck = Presentation(pptx_file)
    ordered_slide_ids = [slide["slide_id"] for slide in presentation_architecture["slides"]]
    sidecar_by_id = {sidecar["slide_id"]: sidecar for sidecar in sidecars}
    findings: list[dict[str, Any]] = []
    slide_reports: list[dict[str, Any]] = []
    object_totals: Counter[str] = Counter()
    semantic_checks = 0
    semantic_passes = 0
    native_checks = 0
    native_passes = 0
    total_raster_area_ratio = 0.0
    full_slide_picture_count = 0

    if len(deck.slides) != len(ordered_slide_ids):
        findings.append(
            _finding(
                "SLIDE_COUNT_MISMATCH",
                "severe",
                f"PPTX has {len(deck.slides)} slides; semantic architecture has {len(ordered_slide_ids)}.",
            )
        )

    slide_width = deck.slide_width / EMU_PER_INCH
    slide_height = deck.slide_height / EMU_PER_INCH
    slide_area = max(0.01, slide_width * slide_height)
    for index, slide_id in enumerate(ordered_slide_ids, start=1):
        sidecar = sidecar_by_id[slide_id]
        if index > len(deck.slides):
            slide_reports.append({"slide_id": slide_id, "slide_index": index, "status": "missing"})
            continue
        slide = deck.slides[index - 1]
        inventory = _slide_inventory(slide, slide_area, slide_width, slide_height)
        object_totals.update(inventory["counts"])
        total_raster_area_ratio += inventory["picture_area_ratio"]
        full_slide_picture_count += inventory["full_slide_picture_count"]
        visible_text = inventory["visible_text"]
        notes_text = inventory["notes_text"]
        semantic_results: list[dict[str, Any]] = []

        for item in sidecar["canonical_content"]:
            result = _check_content_item(item, visible_text, notes_text, inventory)
            semantic_results.append(result)
            semantic_checks += result["check_count"]
            semantic_passes += result["pass_count"]
            if result["status"] != "passed":
                findings.append(
                    _finding(
                        "SEMANTIC_CONTENT_MISSING",
                        "severe",
                        f"{slide_id}: canonical {item['kind']} content for slot {item['slot_id']} was not preserved.",
                        slide_id=slide_id,
                        slide_index=index,
                        details={"slot_id": item["slot_id"], "kind": item["kind"], "missing": result["missing"]},
                    )
                )

        native_results: list[dict[str, Any]] = []
        for requirement in sidecar["native_required"]:
            passed, limitation = _native_requirement_passes(requirement, inventory, semantic_results)
            native_checks += 1
            native_passes += int(passed)
            native_results.append({**requirement, "status": "passed" if passed else "failed", "limitation": limitation})
            if not passed:
                findings.append(
                    _finding(
                        "NATIVE_REQUIREMENT_MISSING",
                        "severe",
                        f"{slide_id}: native {requirement['object_type']} requirement for {requirement['slot_id']} was not met.",
                        slide_id=slide_id,
                        slide_index=index,
                        details=requirement,
                    )
                )

        slide_reports.append(
            {
                "slide_id": slide_id,
                "slide_index": index,
                "layout_id": sidecar["layout_id"],
                "semantic_results": semantic_results,
                "native_results": native_results,
                "object_counts": inventory["counts"],
                "picture_area_ratio": round(inventory["picture_area_ratio"], 6),
                "full_slide_picture_count": inventory["full_slide_picture_count"],
            }
        )

    evidence_registry_ids = {item["evidence_id"] for item in presentation_architecture["evidence_registry"]}
    bound_evidence_ids = {
        evidence_id
        for sidecar in sidecars
        for binding in sidecar["source_bindings"]
        for evidence_id in binding["evidence_ids"]
    }
    source_coverage = len(bound_evidence_ids & evidence_registry_ids) / max(1, len(evidence_registry_ids))
    semantic_fidelity = semantic_passes / max(1, semantic_checks)
    native_editability = native_passes / max(1, native_checks)
    raster_coverage = total_raster_area_ratio / max(1, len(deck.slides))

    if full_slide_picture_count:
        findings.append(_finding("FULL_SLIDE_RASTER_FORBIDDEN", "severe", "One or more full-slide picture shapes were found."))
    if image_policy_report and image_policy_report.get("status") != "passed":
        findings.append(_finding("IMAGE_POLICY_FAILED", "severe", "The final deck image-policy report did not pass."))

    creative = _creative_metrics(creative_template_architecture)
    render = _render_metrics(render_manifest, len(deck.slides))
    if render["status"] == "failed":
        findings.append(_finding("RENDER_QA_FAILED", "severe", "Local PPTX render did not complete for every slide."))

    semantic_status = "passed" if semantic_fidelity == 1.0 and source_coverage == 1.0 else "failed"
    editability_status = "passed" if native_editability == 1.0 and full_slide_picture_count == 0 else "failed"
    creative_status = "passed" if creative["non_pass_fit_count"] == 0 and not creative["repetition_violations"] else "needs_revision"
    severe_count = sum(1 for finding in findings if finding["severity"] == "severe")
    overall_status = "failed" if severe_count else "needs_revision" if creative_status == "needs_revision" else "passed"
    limitations = [
        "Shape requirements without trace metadata are verified at slide-object level, not by exact semantic shape identity.",
        "Raster coverage sums picture bounding-box areas and may over-count overlapping pictures.",
    ]
    if not image_policy_report:
        limitations.append("Template-reference image hash reuse is INCONCLUSIVE because no image-policy report was supplied.")

    return {
        "schema_name": "creative_frontend_qa_report",
        "schema_version": "1.0",
        "status": overall_status,
        "pptx_path": pptx_file.as_posix(),
        "semantic_qa": {
            "status": semantic_status,
            "semantic_fidelity": round(semantic_fidelity, 4),
            "source_coverage": round(source_coverage, 4),
            "semantic_check_count": semantic_checks,
            "semantic_pass_count": semantic_passes,
            "evidence_registry_count": len(evidence_registry_ids),
            "bound_evidence_count": len(bound_evidence_ids & evidence_registry_ids),
        },
        "creative_qa": {"status": creative_status, **creative},
        "editability_qa": {
            "status": editability_status,
            "native_editability": round(native_editability, 4),
            "native_requirement_count": native_checks,
            "native_requirement_pass_count": native_passes,
            "object_totals": dict(sorted(object_totals.items())),
            "raster_coverage": round(raster_coverage, 4),
            "full_slide_picture_count": full_slide_picture_count,
            "image_policy_status": image_policy_report.get("status") if image_policy_report else "inconclusive",
        },
        "render_qa": render,
        "slide_count": len(deck.slides),
        "slides": slide_reports,
        "findings_summary": {
            "total": len(findings),
            "severe": severe_count,
            "warning": sum(1 for finding in findings if finding["severity"] == "warning"),
        },
        "findings": findings,
        "limitations": limitations,
    }


def build_creative_frontend_qa_report_from_files(
    *,
    presentation_architecture_path: str | Path,
    creative_template_architecture_path: str | Path,
    editable_template_spec_path: str | Path,
    semantic_sidecar_dir: str | Path,
    pptx_path: str | Path,
    json_report_path: str | Path,
    md_report_path: str | Path,
    render_manifest_path: str | Path | None = None,
    image_policy_report_path: str | Path | None = None,
) -> Path:
    presentation = _load_json(presentation_architecture_path)
    creative = _load_json(creative_template_architecture_path)
    spec = _load_json(editable_template_spec_path)
    sidecar_dir = Path(semantic_sidecar_dir)
    sidecar_by_id = {
        payload["slide_id"]: payload
        for path in sidecar_dir.glob("*.semantic.json")
        for payload in [_load_json(path)]
    }
    sidecars = [sidecar_by_id[slide["slide_id"]] for slide in presentation["slides"]]
    render_manifest = _load_json(render_manifest_path) if render_manifest_path and Path(render_manifest_path).exists() else None
    image_policy = _load_json(image_policy_report_path) if image_policy_report_path and Path(image_policy_report_path).exists() else None
    report = build_creative_frontend_qa_report(
        presentation_architecture=presentation,
        creative_template_architecture=creative,
        editable_template_spec=spec,
        sidecars=sidecars,
        pptx_path=pptx_path,
        render_manifest=render_manifest,
        image_policy_report=image_policy,
    )
    output = Path(json_report_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=False) + "\n", encoding="utf-8")
    markdown = Path(md_report_path)
    markdown.parent.mkdir(parents=True, exist_ok=True)
    markdown.write_text(_markdown_report(report), encoding="utf-8")
    return output


def _slide_inventory(slide: Any, slide_area: float, slide_width: float, slide_height: float) -> dict[str, Any]:
    counts: Counter[str] = Counter()
    visible_parts: list[str] = []
    tables: list[list[list[str]]] = []
    charts: list[dict[str, Any]] = []
    picture_area_ratio = 0.0
    full_slide_picture_count = 0
    for shape in slide.shapes:
        counts["shapes"] += 1
        if getattr(shape, "has_text_frame", False):
            counts["text_shapes"] += 1
            if str(shape.text or "").strip():
                visible_parts.append(str(shape.text))
        if getattr(shape, "has_table", False):
            counts["tables"] += 1
            tables.append([[str(cell.text) for cell in row.cells] for row in shape.table.rows])
        if getattr(shape, "has_chart", False):
            counts["charts"] += 1
            charts.append(_chart_payload(shape.chart))
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            counts["pictures"] += 1
            x = shape.left / EMU_PER_INCH
            y = shape.top / EMU_PER_INCH
            w = shape.width / EMU_PER_INCH
            h = shape.height / EMU_PER_INCH
            area_ratio = w * h / slide_area
            picture_area_ratio += area_ratio
            if area_ratio >= FULL_SLIDE_AREA_RATIO and x <= 0.2 and y <= 0.2 and w >= slide_width * 0.9 and h >= slide_height * 0.9:
                full_slide_picture_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            counts["connectors"] += 1
        if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.GROUP}:
            counts["editable_shapes"] += 1
    try:
        notes_text = str(slide.notes_slide.notes_text_frame.text or "")
    except Exception:
        notes_text = ""
    return {
        "counts": dict(counts),
        "visible_text": "\n".join(visible_parts),
        "notes_text": notes_text,
        "tables": tables,
        "charts": charts,
        "picture_area_ratio": picture_area_ratio,
        "full_slide_picture_count": full_slide_picture_count,
    }


def _check_content_item(item: dict[str, Any], visible_text: str, notes_text: str, inventory: dict[str, Any]) -> dict[str, Any]:
    kind = item["kind"]
    slot_id = item["slot_id"]
    if kind == "image_need":
        return {"slot_id": slot_id, "kind": kind, "status": "passed", "check_count": 0, "pass_count": 0, "missing": []}
    if kind == "table":
        passed = _table_matches(item["value"], inventory["tables"])
        return _single_result(item, passed, [] if passed else ["native table data"])
    if kind == "chart":
        passed = _chart_matches(item["value"], inventory["charts"])
        return _single_result(item, passed, [] if passed else ["native chart data"])
    if kind == "citation":
        expected = [str(citation.get("label") or citation.get("source") or "") for citation in item["value"] if isinstance(citation, dict)]
        return _text_result(item, expected, visible_text)
    if kind == "speaker_notes":
        return _text_result(item, [str(item["value"])], notes_text)
    expected = _atomic_text_values(item["value"])
    search_text = visible_text if slot_id == "title" else f"{visible_text}\n{notes_text}"
    return _text_result(item, expected, search_text)


def _single_result(item: dict[str, Any], passed: bool, missing: list[str]) -> dict[str, Any]:
    return {
        "slot_id": item["slot_id"],
        "kind": item["kind"],
        "status": "passed" if passed else "failed",
        "check_count": 1,
        "pass_count": int(passed),
        "missing": missing,
    }


def _text_result(item: dict[str, Any], expected: list[str], search_text: str) -> dict[str, Any]:
    normalized_search = _normalize_text(search_text)
    missing = [value for value in expected if _normalize_text(value) not in normalized_search]
    return {
        "slot_id": item["slot_id"],
        "kind": item["kind"],
        "status": "passed" if not missing else "failed",
        "check_count": len(expected),
        "pass_count": len(expected) - len(missing),
        "missing": missing,
    }


def _atomic_text_values(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value if str(item).strip()]
    if isinstance(value, dict):
        return [f"{key}: {item}" for key, item in value.items()]
    text = str(value or "").strip()
    return [text] if text else []


def _table_matches(expected: Any, actual_tables: list[list[list[str]]]) -> bool:
    if not isinstance(expected, dict):
        return False
    expected_rows = []
    if expected.get("headers"):
        expected_rows.append([str(value) for value in expected["headers"]])
    expected_rows.extend([[str(value) for value in row] for row in expected.get("rows") or []])
    normalized_expected = [[_normalize_text(cell) for cell in row] for row in expected_rows]
    return any([[ _normalize_text(cell) for cell in row] for row in table] == normalized_expected for table in actual_tables)


def _chart_payload(chart: Any) -> dict[str, Any]:
    try:
        categories = [str(category.label) for category in chart.plots[0].categories]
    except Exception:
        categories = []
    series = []
    for item in chart.series:
        try:
            values = [float(value) if value is not None else None for value in item.values]
        except Exception:
            values = []
        series.append({"name": str(item.name), "values": values})
    return {"categories": categories, "series": series}


def _chart_matches(expected: Any, actual_charts: list[dict[str, Any]]) -> bool:
    if not isinstance(expected, dict):
        return False
    expected_categories = [str(value) for value in expected.get("categories") or []]
    expected_series = [
        {"name": str(item.get("name") or ""), "values": [float(value) for value in item.get("values") or []]}
        for item in expected.get("series") or []
        if isinstance(item, dict)
    ]
    return any(chart["categories"] == expected_categories and chart["series"] == expected_series for chart in actual_charts)


def _native_requirement_passes(
    requirement: dict[str, Any],
    inventory: dict[str, Any],
    semantic_results: list[dict[str, Any]],
) -> tuple[bool, str | None]:
    object_type = requirement["object_type"]
    slot_id = requirement["slot_id"]
    counts = inventory["counts"]
    semantic_match = any(result["slot_id"] == slot_id and result["status"] == "passed" for result in semantic_results)
    if object_type == "table":
        return counts.get("tables", 0) > 0, None
    if object_type == "chart":
        return counts.get("charts", 0) > 0, None
    if object_type == "connector":
        return counts.get("connectors", 0) > 0, None
    if object_type == "speaker_notes":
        return bool(inventory["notes_text"].strip()) and semantic_match, None
    if object_type == "footer":
        return semantic_match, None
    if object_type == "text_box":
        return semantic_match and counts.get("text_shapes", 0) > 0, None
    if object_type == "svg_icon":
        return counts.get("pictures", 0) > 0 or counts.get("editable_shapes", 0) > 0, "SVG identity is not trace-verified."
    if object_type == "shape":
        return counts.get("editable_shapes", 0) > 0 and (semantic_match or slot_id not in {result["slot_id"] for result in semantic_results}), "Shape identity is not trace-verified."
    return False, None


def _creative_metrics(creative: dict[str, Any]) -> dict[str, Any]:
    decisions = creative["slide_fit_decisions"]
    status_counts = Counter(decision["status"] for decision in decisions)
    layout_ids = [decision["layout_id"] for decision in decisions]
    family_ids = [decision["template_family_id"] for decision in decisions]
    max_layout_run = _max_consecutive(layout_ids)
    max_family_run = _max_consecutive(family_ids)
    policy = creative["fit_policy"]
    violations = []
    if max_layout_run > int(policy["max_consecutive_layout"]):
        violations.append({"kind": "layout", "actual": max_layout_run, "maximum": policy["max_consecutive_layout"]})
    if max_family_run > int(policy["max_consecutive_family"]):
        violations.append({"kind": "family", "actual": max_family_run, "maximum": policy["max_consecutive_family"]})
    signatures = [tuple(module["differentiation_signature"]) for module in creative["modules"]]
    return {
        "fit_status_counts": dict(sorted(status_counts.items())),
        "non_pass_fit_count": sum(count for status, count in status_counts.items() if status != "pass"),
        "distinct_layout_count": len(set(layout_ids)),
        "distinct_family_count": len(set(family_ids)),
        "max_consecutive_layout": max_layout_run,
        "max_consecutive_family": max_family_run,
        "repetition_violations": violations,
        "module_differentiation_signature_count": len(set(signatures)),
        "module_count": len(signatures),
    }


def _render_metrics(render_manifest: dict[str, Any] | None, slide_count: int) -> dict[str, Any]:
    if not render_manifest:
        return {"status": "inconclusive", "render_status": "not_supplied", "rendered_slide_count": 0, "slide_count": slide_count}
    rendered = int(render_manifest.get("rendered_slide_count") or 0)
    expected = int(render_manifest.get("slide_count") or slide_count)
    passed = render_manifest.get("render_status") == "rendered" and rendered == expected == slide_count and not render_manifest.get("errors")
    return {
        "status": "passed" if passed else "failed",
        "render_status": render_manifest.get("render_status"),
        "backend": render_manifest.get("backend"),
        "rendered_slide_count": rendered,
        "slide_count": expected,
        "error_count": len(render_manifest.get("errors") or []),
    }


def _max_consecutive(values: list[str]) -> int:
    best = current = 0
    previous = None
    for value in values:
        current = current + 1 if value == previous else 1
        best = max(best, current)
        previous = value
    return best


def _normalize_text(value: Any) -> str:
    return " ".join(str(value or "").split()).casefold()


def _finding(
    code: str,
    severity: str,
    message: str,
    *,
    slide_id: str | None = None,
    slide_index: int | None = None,
    details: dict[str, Any] | None = None,
) -> dict[str, Any]:
    return {
        "code": code,
        "severity": severity,
        "message": message,
        "slide_id": slide_id,
        "slide_index": slide_index,
        "details": details or {},
    }


def _markdown_report(report: dict[str, Any]) -> str:
    semantic = report["semantic_qa"]
    creative = report["creative_qa"]
    editable = report["editability_qa"]
    lines = [
        "# Creative Front-End QA",
        "",
        f"Overall status: `{report['status']}`",
        f"Semantic fidelity: `{semantic['semantic_fidelity']:.1%}`",
        f"Source coverage: `{semantic['source_coverage']:.1%}`",
        f"Native editability: `{editable['native_editability']:.1%}`",
        f"Raster coverage: `{editable['raster_coverage']:.1%}`",
        f"Full-slide pictures: `{editable['full_slide_picture_count']}`",
        f"Creative fit: `{creative['status']}` ({creative['non_pass_fit_count']} non-pass decisions)",
        f"Local render: `{report['render_qa']['status']}`",
        "",
        "## Findings",
        "",
    ]
    if report["findings"]:
        lines.extend(f"- [{item['severity']}] {item['code']}: {item['message']}" for item in report["findings"])
    else:
        lines.append("- None")
    lines.extend(["", "## Limitations", ""])
    lines.extend(f"- {item}" for item in report["limitations"])
    return "\n".join(lines) + "\n"


def _load_json(path: str | Path) -> dict[str, Any]:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Run semantic, creative, and editability QA on a creative-front-end PPTX run.")
    parser.add_argument("--presentation-architecture", type=Path, required=True)
    parser.add_argument("--creative-template-architecture", type=Path, required=True)
    parser.add_argument("--template-spec", type=Path, required=True)
    parser.add_argument("--semantic-sidecars", type=Path, required=True)
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--render-manifest", type=Path, default=None)
    parser.add_argument("--image-policy-report", type=Path, default=None)
    parser.add_argument("--json-report", type=Path, required=True)
    parser.add_argument("--md-report", type=Path, required=True)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        output = build_creative_frontend_qa_report_from_files(
            presentation_architecture_path=args.presentation_architecture,
            creative_template_architecture_path=args.creative_template_architecture,
            editable_template_spec_path=args.template_spec,
            semantic_sidecar_dir=args.semantic_sidecars,
            pptx_path=args.pptx,
            render_manifest_path=args.render_manifest,
            image_policy_report_path=args.image_policy_report,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
        report = _load_json(output)
    except Exception as exc:
        print(f"CREATIVE_FRONTEND_QA_FAILED {exc}")
        return 1
    print(f"WROTE {output}")
    return 0 if report["status"] == "passed" else 2 if report["status"] == "needs_revision" else 1


if __name__ == "__main__":
    raise SystemExit(main())
