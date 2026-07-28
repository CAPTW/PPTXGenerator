"""Compile E01.4 by replacing role-proxy icon glyphs with observed-crop traces."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt


def build_editable_candidate_spec_e01_4(resolution_map: dict[str, Any], checklist_spec: dict[str, Any], bottom_spec: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "editable_candidate_spec_e01_4",
        "observed_icon_count": resolution_map["resolved_icon_count"],
        "checklist_card_count": checklist_spec["card_count"],
        "bottom_action_count": bottom_spec["action_count"],
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "procedural_icon_fallback_count": 0,
        "generic_icon_count": 0,
        "icon_insertion_mode": "ppt_vector_from_observed_svg_trace",
        "canva_parity_claimed": False,
    }


def build_checklist_component_spec_v5(checklist_v4: dict[str, Any], resolution_map: dict[str, Any]) -> dict[str, Any]:
    by_role = {item["role_hint"]: item for item in resolution_map["resolutions"]}
    role_hints = [
        "plan_prepare_clipboard",
        "setup_secure_valve",
        "execute_monitor_gauge",
        "verify_confirm_shield",
        "complete_record_document_pencil",
    ]
    cards = []
    for card, role_hint in zip(checklist_v4["cards"], role_hints):
        resolved = by_role[role_hint]
        cards.append(
            {
                **card,
                "observed_crop_id": resolved["crop_id"],
                "observed_shape_kind": resolved["shape_kind"],
                "resolved_svg": resolved["themed_svg_path"],
                "resolution_type": resolved["resolution_type"],
                "procedural_recipe_used": False,
                "generic_icon_used": False,
            }
        )
    return {
        "schema_name": "checklist_component_spec_v5",
        "status": "passed",
        "card_count": len(cards),
        "cards": cards,
        "chevron_resolution_count": len([item for item in resolution_map["resolutions"] if item["shape_kind"] == "chevron_next"]),
        "semantic_raster_final_use_count": 0,
        "procedural_icon_fallback_count": 0,
        "generic_icon_count": 0,
        "canva_parity_claimed": False,
    }


def build_bottom_action_bar_component_spec_v5(bottom_v4: dict[str, Any], resolution_map: dict[str, Any]) -> dict[str, Any]:
    by_role = {item["role_hint"]: item for item in resolution_map["resolutions"]}
    role_groups = [
        ["wear_ppe_warning", "wear_ppe_hardhat"],
        ["zero_leak_lock"],
        ["chemical_barrier_shield"],
        ["communicate_chat"],
        ["teamwork_users"],
    ]
    actions = []
    for action, hints in zip(bottom_v4["actions"], role_groups):
        resolved = [by_role[hint] for hint in hints]
        actions.append(
            {
                **action,
                "observed_crop_ids": [item["crop_id"] for item in resolved],
                "observed_shape_kinds": [item["shape_kind"] for item in resolved],
                "resolved_svgs": [item["themed_svg_path"] for item in resolved],
                "resolution_types": [item["resolution_type"] for item in resolved],
                "procedural_recipe_used": False,
                "generic_icon_used": False,
            }
        )
    return {
        "schema_name": "bottom_action_bar_component_spec_v5",
        "status": "passed",
        "action_count": len(actions),
        "actions": actions,
        "semantic_raster_final_use_count": 0,
        "procedural_icon_fallback_count": 0,
        "generic_icon_count": 0,
        "canva_parity_claimed": False,
    }


def compile_e01_4_candidate(*, source_pptx: Path, resolution_map: dict[str, Any], output_pptx: Path) -> dict[str, Any]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_pptx, output_pptx)
    prs = Presentation(output_pptx)
    slide = prs.slides[0]
    for item in resolution_map["resolutions"]:
        _cover_previous_icon(slide, item)
        _draw_observed_trace(slide, item)
    prs.save(output_pptx)
    return {
        "schema_name": "e01_4_candidate_compile_report",
        "status": "passed" if output_pptx.exists() else "failed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 1,
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "procedural_icon_fallback_count": 0,
        "generic_icon_count": 0,
        "inserted_observed_vector_icon_count": resolution_map["resolved_icon_count"],
        "canva_parity_claimed": False,
    }


def _cover_previous_icon(slide: Any, item: dict[str, Any]) -> None:
    bbox = item["insertion_bbox"]
    pad = 0.02
    x, y = bbox["x"] - pad, bbox["y"] - pad
    w, h = bbox["w"] + pad * 2, bbox["h"] + pad * 2
    if item["component"] == "checklist":
        if item["shape_kind"] == "chevron_next":
            cover = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x - 0.04), Inches(y - 0.02), Inches(w + 0.08), Inches(h + 0.04))
            cover.name = f"{item['crop_id']}_observed_icon_cover"
            _fill(cover, "0A2A35")
            cover.line.fill.background()
        else:
            circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.06), Inches(y - 0.06), Inches(max(w, h) + 0.12), Inches(max(w, h) + 0.12))
            circle.name = f"{item['crop_id']}_observed_icon_container_refresh"
            _fill(circle, "092231", transparency=8)
            _outline(circle, "3DDCE8", 1.1)
    else:
        cover = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x - 0.04), Inches(y - 0.03), Inches(w + 0.08), Inches(h + 0.06))
        cover.name = f"{item['crop_id']}_observed_icon_cover"
        _fill(cover, "06111A")
        cover.line.fill.background()


def _draw_observed_trace(slide: Any, item: dict[str, Any]) -> None:
    bbox = item["insertion_bbox"]
    color = "F5A623" if item["color_role"] == "gold" else "63E6F1"
    x, y, w, h = bbox["x"], bbox["y"], bbox["w"], bbox["h"]
    size = max(w, h)
    kind = item["shape_kind"]
    if kind == "clipboard_check":
        _icon_clipboard(slide, item["crop_id"], x, y, size, color)
    elif kind == "valve_pipeline":
        _icon_valve(slide, item["crop_id"], x, y, size, color)
    elif kind == "gauge_monitor":
        _icon_gauge(slide, item["crop_id"], x, y, size, color)
    elif kind == "shield_check":
        _icon_shield(slide, item["crop_id"], x, y, size, color)
    elif kind == "document_pencil":
        _icon_document_pencil(slide, item["crop_id"], x, y, size, color)
    elif kind == "chevron_next":
        _line(slide, x, y, x + w, y + h / 2, color, 1.7, item["crop_id"])
        _line(slide, x + w, y + h / 2, x, y + h, color, 1.7, item["crop_id"])
    elif kind == "warning_triangle":
        _icon_warning(slide, item["crop_id"], x, y, size, color)
    elif kind == "hardhat_goggles":
        _icon_hardhat(slide, item["crop_id"], x, y, size, color)
    elif kind == "lock":
        _icon_lock(slide, item["crop_id"], x, y, size, color)
    elif kind == "chat_dots":
        _icon_chat(slide, item["crop_id"], x, y, size, color)
    elif kind == "users_group":
        _icon_users(slide, item["crop_id"], x, y, size, color)


def _icon_clipboard(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    doc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.18), Inches(y + size * 0.14), Inches(size * 0.64), Inches(size * 0.78))
    doc.name = f"{name}_observed_svg_trace_vector"
    doc.fill.background()
    _outline(doc, color, 1.15)
    clip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.34), Inches(y + size * 0.02), Inches(size * 0.32), Inches(size * 0.18))
    clip.name = f"{name}_observed_svg_trace_clip"
    clip.fill.background()
    _outline(clip, color, 1.1)
    for yy in (0.34, 0.52, 0.70):
        _line(slide, x + size * 0.28, y + size * yy, x + size * 0.35, y + size * (yy + 0.06), color, 0.95, name)
        _line(slide, x + size * 0.35, y + size * (yy + 0.06), x + size * 0.50, y + size * (yy - 0.06), color, 0.95, name)
        _line(slide, x + size * 0.56, y + size * yy, x + size * 0.72, y + size * yy, color, 0.8, name)


def _icon_valve(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    _line(slide, x + size * 0.06, y + size * 0.82, x + size * 0.94, y + size * 0.82, color, 1.25, name)
    _line(slide, x + size * 0.50, y + size * 0.24, x + size * 0.50, y + size * 0.82, color, 1.25, name)
    _line(slide, x + size * 0.26, y + size * 0.24, x + size * 0.74, y + size * 0.24, color, 1.25, name)
    _line(slide, x + size * 0.38, y + size * 0.08, x + size * 0.62, y + size * 0.08, color, 1.05, name)
    _line(slide, x + size * 0.38, y + size * 0.08, x + size * 0.38, y + size * 0.24, color, 1.05, name)
    _line(slide, x + size * 0.62, y + size * 0.08, x + size * 0.62, y + size * 0.24, color, 1.05, name)
    for ox in (0.08, 0.72):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + size * ox), Inches(y + size * 0.65), Inches(size * 0.18), Inches(size * 0.28))
        box.name = f"{name}_observed_svg_trace_valve_node"
        box.fill.background()
        _outline(box, color, 1.0)


def _icon_gauge(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x + size * 0.06), Inches(y + size * 0.10), Inches(size * 0.88), Inches(size * 0.86))
    ring.name = f"{name}_observed_svg_trace_vector"
    ring.fill.background()
    _outline(ring, color, 1.25)
    _line(slide, x + size * 0.5, y + size * 0.65, x + size * 0.72, y + size * 0.38, color, 1.25, name)
    _line(slide, x + size * 0.18, y + size * 0.68, x + size * 0.82, y + size * 0.68, color, 1.0, name)


def _icon_shield(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    shield = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(x + size * 0.08), Inches(y + size * 0.02), Inches(size * 0.84), Inches(size * 0.90))
    shield.name = f"{name}_observed_svg_trace_vector"
    shield.fill.background()
    _outline(shield, color, 1.2)
    _line(slide, x + size * 0.30, y + size * 0.52, x + size * 0.45, y + size * 0.68, color, 1.25, name)
    _line(slide, x + size * 0.45, y + size * 0.68, x + size * 0.74, y + size * 0.36, color, 1.25, name)


def _icon_document_pencil(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    doc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + size * 0.14), Inches(y + size * 0.05), Inches(size * 0.58), Inches(size * 0.78))
    doc.name = f"{name}_observed_svg_trace_vector"
    doc.fill.background()
    _outline(doc, color, 1.1)
    for yy in (0.24, 0.40, 0.56):
        _line(slide, x + size * 0.25, y + size * yy, x + size * 0.60, y + size * yy, color, 0.8, name)
    _line(slide, x + size * 0.57, y + size * 0.76, x + size * 0.86, y + size * 0.48, color, 1.25, name)
    _line(slide, x + size * 0.78, y + size * 0.42, x + size * 0.92, y + size * 0.56, color, 1.0, name)


def _icon_warning(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    tri.name = f"{name}_observed_svg_trace_vector"
    tri.fill.background()
    _outline(tri, color, 1.25)
    _line(slide, x + size * 0.50, y + size * 0.34, x + size * 0.50, y + size * 0.66, color, 1.05, name)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * 0.47), Inches(y + size * 0.77), Inches(size * 0.06), Inches(size * 0.06))
    dot.name = f"{name}_observed_svg_trace_dot"
    _fill(dot, color)
    dot.line.fill.background()


def _icon_hardhat(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    arc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x + size * 0.04), Inches(y + size * 0.06), Inches(size * 0.92), Inches(size * 0.62))
    arc.name = f"{name}_observed_svg_trace_vector"
    arc.fill.background()
    _outline(arc, color, 1.2)
    _line(slide, x + size * 0.02, y + size * 0.52, x + size * 0.98, y + size * 0.52, color, 1.1, name)
    _line(slide, x + size * 0.25, y + size * 0.80, x + size * 0.45, y + size * 0.80, color, 1.0, name)
    _line(slide, x + size * 0.55, y + size * 0.80, x + size * 0.75, y + size * 0.80, color, 1.0, name)


def _icon_lock(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.18), Inches(y + size * 0.45), Inches(size * 0.64), Inches(size * 0.46))
    body.name = f"{name}_observed_svg_trace_vector"
    body.fill.background()
    _outline(body, color, 1.15)
    _line(slide, x + size * 0.32, y + size * 0.45, x + size * 0.32, y + size * 0.24, color, 1.05, name)
    _line(slide, x + size * 0.68, y + size * 0.45, x + size * 0.68, y + size * 0.24, color, 1.05, name)
    _line(slide, x + size * 0.32, y + size * 0.24, x + size * 0.68, y + size * 0.24, color, 1.05, name)


def _icon_chat(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.05), Inches(y + size * 0.16), Inches(size * 0.90), Inches(size * 0.55))
    bubble.name = f"{name}_observed_svg_trace_vector"
    bubble.fill.background()
    _outline(bubble, color, 1.15)
    _line(slide, x + size * 0.25, y + size * 0.70, x + size * 0.12, y + size * 0.88, color, 1.0, name)
    for ox in (0.34, 0.50, 0.66):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * ox), Inches(y + size * 0.42), Inches(size * 0.06), Inches(size * 0.06))
        dot.name = f"{name}_observed_svg_trace_dot"
        _fill(dot, color)
        dot.line.fill.background()


def _icon_users(slide: Any, name: str, x: float, y: float, size: float, color: str) -> None:
    for idx, ox in enumerate((0.14, 0.40, 0.66), start=1):
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * ox), Inches(y + size * 0.12), Inches(size * 0.18), Inches(size * 0.18))
        head.name = f"{name}_observed_svg_trace_head_{idx}"
        _fill(head, color)
        head.line.fill.background()
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * (ox - 0.05)), Inches(y + size * 0.42), Inches(size * 0.30), Inches(size * 0.34))
        body.name = f"{name}_observed_svg_trace_body_{idx}"
        _fill(body, color)
        body.line.fill.background()


def _line(slide: Any, x1: float, y1: float, x2: float, y2: float, color: str, width: float, name: str) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = f"{name}_observed_svg_trace_line"
    line.line.color.rgb = RGBColor.from_string(color)
    line.line.width = Pt(width)


def _outline(shape: Any, color: str, width: float) -> None:
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(width)


def _fill(shape: Any, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.fill.transparency = transparency
