"""Contract-first PPTX compiler for E06.2."""

from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation

from src.presentation_agent.magic_layer.e06_2_contract_chart_table_factory import add_contract_chart_table
from src.presentation_agent.magic_layer.e06_2_contract_object_factory import build_contract_compile_object_plan
from src.presentation_agent.magic_layer.e06_2_contract_shape_factory import add_contract_placeholder, add_contract_shape
from src.presentation_agent.magic_layer.e06_2_contract_svg_icon_factory import build_svg_icon_instructions, inject_svg_icons
from src.presentation_agent.magic_layer.e06_2_contract_text_factory import add_contract_text


def compile_contract_pptx(
    contract: dict[str, Any],
    output_pptx: Path,
    *,
    baseline_pptx: Path,
    icon_root: Path,
) -> dict[str, Any]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        media_dir = Path(tmp) / "media"
        media_dir.mkdir()
        media_map = _extract_baseline_media(baseline_pptx, media_dir)
        prs = Presentation()
        prs.slide_width = int(contract["slide_size"]["width_emu"])
        prs.slide_height = int(contract["slide_size"]["height_emu"])
        blank_layout = prs.slide_layouts[6]
        compiled_count = 0
        fallback_count = 0
        for slide_contract in contract.get("slides", []):
            slide = prs.slides.add_slide(blank_layout)
            for obj in sorted(slide_contract.get("objects", []), key=lambda row: row.get("z_order", 0)):
                compiled_count += 1
                if obj.get("object_type") == "semantic_icon":
                    add_contract_placeholder(slide, obj)
                elif obj.get("object_type") in {"text", "source_footer"} and str(obj.get("text_excerpt") or "").strip():
                    add_contract_text(slide, obj)
                elif obj.get("object_type") == "image_field" and _add_image_if_available(slide, obj, media_map):
                    continue
                elif obj.get("object_type") in {"chart_region", "table_region"}:
                    add_contract_chart_table(slide, obj)
                else:
                    add_contract_shape(slide, obj)
                    if obj.get("object_type") == "image_field":
                        fallback_count += 1
        prs.save(output_pptx)
    instructions, icon_resolution = build_svg_icon_instructions(contract, icon_root)
    injection = inject_svg_icons(output_pptx, instructions)
    plan = build_contract_compile_object_plan(contract)
    return {
        "schema_name": "contract_first_compile_report",
        "status": "passed" if output_pptx.exists() and icon_resolution["status"] == "passed" and injection["status"] == "passed" else "failed",
        "compiled_pptx_path": output_pptx.as_posix(),
        "slides_compiled": len(contract.get("slides", [])),
        "objects_compiled_from_contract": compiled_count,
        "image_field_shape_fallback_count": fallback_count,
        "svg_icon_resolution": icon_resolution,
        "svg_icon_injection": injection,
        "object_plan_summary": {
            "object_count": plan.get("object_count", 0),
            "object_counts_by_type": plan.get("object_counts_by_type", {}),
        },
        "semantic_raster_icon_count": 0,
        "copied_existing_slide_xml_wholesale": False,
    }


def _extract_baseline_media(baseline_pptx: Path, media_dir: Path) -> dict[str, Path]:
    media: dict[str, Path] = {}
    if not baseline_pptx.exists():
        return media
    with zipfile.ZipFile(baseline_pptx, "r") as zf:
        for name in zf.namelist():
            if not name.startswith("ppt/media/"):
                continue
            target = media_dir / Path(name).name
            target.write_bytes(zf.read(name))
            media[f"/{name}"] = target
            media[name] = target
    return media


def _add_image_if_available(slide: Any, obj: dict[str, Any], media_map: dict[str, Path]) -> bool:
    media = obj.get("media", {})
    content_type = media.get("content_type", "")
    if content_type == "image/svg+xml":
        return False
    partname = media.get("partname")
    image_path = media_map.get(str(partname or ""))
    if not image_path or not image_path.exists():
        return False
    bbox = obj["bbox_emu"]
    try:
        picture = slide.shapes.add_picture(str(image_path), int(bbox["x"]), int(bbox["y"]), max(1, int(bbox["w"])), max(1, int(bbox["h"])))
        from src.presentation_agent.magic_layer.e06_2_contract_object_factory import contract_shape_name

        picture.name = contract_shape_name(obj)
        return True
    except Exception:
        return False


def render_contract_deck(pptx_path: Path, output_root: Path, *, prefix: str = "slide") -> dict[str, Any]:
    from src.presentation_agent.qa.render_pptx_preview import render_pptx_preview

    raw_dir = output_root / "renders" / f"_{prefix}_raw"
    raw_dir.mkdir(parents=True, exist_ok=True)
    report = render_pptx_preview(pptx_path=pptx_path, output_dir=raw_dir, manifest_path=output_root / "renders" / f"{prefix}_render_manifest.json", backend="auto", dpi=144)
    rendered = []
    for idx, row in enumerate(report.get("slides", []), start=1):
        source = Path(row.get("rendered_image_path") or "")
        if source.exists():
            target = output_root / "renders" / f"{prefix}-{idx:03d}.png"
            shutil.copy2(source, target)
            row["rendered_image_path"] = target.as_posix()
            rendered.append(target)
    report["rendered_slide_count"] = len(rendered)
    report["expected_slide_count"] = 16
    report["rendered_paths"] = [path.as_posix() for path in rendered]
    report["status"] = "passed" if len(rendered) == 16 else "failed"
    return report
