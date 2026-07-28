"""Build a non-canonical SVG/vector icon regression fixture PPTX."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.image import ImagePart
from pptx.util import Inches, Pt

try:
    import cairosvg
except Exception:  # pragma: no cover - optional dependency in some environments
    cairosvg = None


SVG_CONTENT_TYPE = "image/svg+xml"


def build_icon_regression_fixture(curated_manifest: dict[str, Any], output_dir: Path, *, render: bool = True) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "magic_layer_v7_icon_regression_fixture.pptx"
    render_path = output_dir / "magic_layer_v7_icon_regression_fixture_render.png"
    preview_path = output_dir / "magic_layer_v7_icon_regression_fixture_preview.png"
    roles = sorted(curated_manifest.get("roles", []), key=lambda row: (row.get("priority", ""), row.get("role_id", "")))
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for size_label, icon_size in (("16px", 0.17), ("24px", 0.25), ("32px", 0.33)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_fixture_slide(slide, roles, size_label, icon_size)
    prs.save(pptx_path)
    _build_preview_sheet(roles, preview_path)
    fixture_render_exists = False
    render_report: dict[str, Any] = {"render_status": "not_requested"}
    if render:
        render_report = _render_fixture(pptx_path, output_dir / "rendered")
        first_slide = None
        for slide_row in render_report.get("slides", []):
            candidate = Path(slide_row.get("rendered_image_path") or slide_row.get("preview_path") or slide_row.get("path") or "")
            if candidate.exists():
                first_slide = candidate
                break
        if first_slide:
            shutil.copy2(first_slide, render_path)
            fixture_render_exists = True
    return {
        "schema_name": "icon_regression_fixture_report",
        "status": "passed" if pptx_path.exists() and (fixture_render_exists or not render) else "blocked",
        "fixture_pptx_path": pptx_path.as_posix(),
        "fixture_render_path": render_path.as_posix(),
        "fixture_preview_path": preview_path.as_posix(),
        "fixture_render_exists": fixture_render_exists or (not render and pptx_path.exists()),
        "render_report": render_report,
        "role_count": len(roles),
        "svg_vector_icon_count": len(roles),
        "raster_icon_count": 0,
    }


def _add_fixture_slide(slide: Any, roles: list[dict[str, Any]], size_label: str, icon_size: float) -> None:
    _add_background(slide, RGBColor(248, 250, 252))
    title = slide.shapes.add_textbox(Inches(0.25), Inches(0.12), Inches(6), Inches(0.25))
    title.text_frame.text = f"Magic Layer v7 Icon Regression Fixture - {size_label}"
    title.text_frame.paragraphs[0].font.size = Pt(13)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)
    dark_panel = slide.shapes.add_shape(1, Inches(6.75), Inches(0.5), Inches(6.25), Inches(6.75))
    dark_panel.fill.solid()
    dark_panel.fill.fore_color.rgb = RGBColor(7, 16, 24)
    dark_panel.line.color.rgb = RGBColor(7, 16, 24)
    cols = 6
    x0 = 0.28
    y0 = 0.58
    cell_w = 1.04
    cell_h = 0.58
    for idx, role in enumerate(roles[:66]):
        col = idx % cols
        row = idx // cols
        for side, offset, color in (("light", 0.0, RGBColor(15, 23, 42)), ("dark", 6.7, RGBColor(248, 250, 252))):
            x = x0 + offset + col * cell_w
            y = y0 + row * cell_h
            _add_svg(slide, Path(role["svg_path"]), x, y + 0.03, icon_size, icon_size, role["role_id"])
            label = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.02), Inches(0.72), Inches(0.18))
            label.text_frame.text = role["role_id"][:18]
            p = label.text_frame.paragraphs[0]
            p.font.size = Pt(5.8)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.LEFT


def _add_background(slide: Any, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, 0, 0, Inches(13.333333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def _add_svg(slide: Any, svg_path: Path, x: float, y: float, w: float, h: float, role_id: str) -> None:
    image_part = ImagePart(
        slide.part.package.next_image_partname("svg"),
        SVG_CONTENT_TYPE,
        slide.part.package,
        svg_path.read_bytes(),
        svg_path.name,
    )
    r_id = slide.part.relate_to(image_part, RT.IMAGE)
    shape_id = slide.shapes._next_shape_id
    safe_role = role_id.replace("&", "and").replace("<", "").replace(">", "")
    slide.shapes._grpSp.add_pic(
        shape_id,
        f"SVG Icon {safe_role}",
        f"svg-icon:{safe_role}",
        r_id,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    slide.shapes._recalculate_extents()


def _render_fixture(pptx_path: Path, render_dir: Path) -> dict[str, Any]:
    try:
        from src.presentation_agent.qa.render_pptx_preview import render_pptx_preview

        return render_pptx_preview(
            pptx_path=pptx_path,
            output_dir=render_dir,
            manifest_path=render_dir / "render_manifest.json",
            backend="auto",
            dpi=144,
        )
    except Exception as exc:  # pragma: no cover - renderer environment dependent
        return {"render_status": "skipped", "errors": [{"error": str(exc)}], "slides": []}


def _build_preview_sheet(roles: list[dict[str, Any]], output: Path) -> None:
    cols = 6
    cell_w, cell_h = 170, 96
    rows = max(1, (len(roles) + cols - 1) // cols)
    sheet = Image.new("RGB", (cols * cell_w, rows * cell_h), "#F8FAFC")
    draw = ImageDraw.Draw(sheet)
    font = ImageFont.load_default()
    for idx, role in enumerate(roles):
        x = (idx % cols) * cell_w
        y = (idx // cols) * cell_h
        draw.rectangle((x, y, x + cell_w, y + cell_h), outline="#CBD5E1", fill="#FFFFFF")
        draw.text((x + 8, y + 8), role["role_id"][:24], fill="#0F172A", font=font)
        _paste_svg(sheet, Path(role["svg_path"]), x + 58, y + 34)
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)


def _paste_svg(sheet: Image.Image, svg_path: Path, x: int, y: int) -> None:
    if cairosvg is None or not svg_path.exists():
        return
    temp = svg_path.with_suffix(".e03_4_preview.png")
    try:
        cairosvg.svg2png(url=svg_path.as_posix(), write_to=temp.as_posix(), output_width=42, output_height=42)
        image = Image.open(temp).convert("RGBA")
        sheet.paste(image, (x, y), image)
    finally:
        temp.unlink(missing_ok=True)
