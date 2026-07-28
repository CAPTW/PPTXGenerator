"""Assemble the non-canonical E03.5 16-archetype pack with v7.1 SVG icons."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from .e03_5_icon_v7_1_inserter import _add_badge, _add_svg, _bbox_px, _icon_position


def assemble_e03_5_pack(
    baseline_pack: Path,
    output_pptx: Path,
    archetypes: list[str],
    resolved_rows_by_archetype: dict[str, list[dict[str, Any]]],
) -> dict[str, Any]:
    """Copy the E03.3 pack and insert v7.1 semantic SVG media per slide."""

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    if not baseline_pack.exists():
        return {
            "schema_name": "e03_5_pack_assembly_report",
            "status": "blocked",
            "reason": "baseline_pack_missing",
            "pptx_path": output_pptx.as_posix(),
            "slide_count": 0,
            "icon_v7_1_usage_count": 0,
            "true_svg_media_insertion_count": 0,
            "native_vector_conversion_count": 0,
            "raster_semantic_icon_count": 0,
            "non_canonical": True,
            "canonical_promotion": False,
            "source_bound_deck_created": False,
        }

    shutil.copy2(baseline_pack, output_pptx)
    prs = Presentation(output_pptx)
    usage_rows: list[dict[str, Any]] = []
    for slide_index, archetype in enumerate(archetypes):
        if slide_index >= len(prs.slides):
            break
        slide = prs.slides[slide_index]
        resolved_rows = resolved_rows_by_archetype.get(archetype, [])
        for idx, row in enumerate(resolved_rows):
            x, y, size = _icon_position(idx, len(resolved_rows))
            _add_badge(slide, x, y, size)
            shape_id = _add_svg(slide, Path(row["themed_svg_path"]), x + 0.035, y + 0.035, size - 0.07, size - 0.07, row["semantic_role"])
            usage_rows.append(
                {
                    **row,
                    "archetype_id": archetype,
                    "slide_index": slide_index + 1,
                    "slide_id": f"{archetype}_pack_slide",
                    "shape_id": shape_id,
                    "object_name": f"SVG Icon {row['semantic_role']}",
                    "bbox": {"x": x + 0.035, "y": y + 0.035, "w": size - 0.07, "h": size - 0.07},
                    "bbox_px": _bbox_px(x + 0.035, y + 0.035, size - 0.07, size - 0.07, margin_px=7),
                    "visible_in_render": None,
                    "raster_fallback": False,
                }
            )
    prs.save(output_pptx)
    return {
        "schema_name": "e03_5_pack_assembly_report",
        "status": "passed" if output_pptx.exists() and len(prs.slides) >= len(archetypes) else "blocked",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": len(prs.slides),
        "archetypes": archetypes,
        "icon_v7_1_usage_count": len(usage_rows),
        "true_svg_media_insertion_count": len(usage_rows),
        "native_vector_conversion_count": 0,
        "raster_semantic_icon_count": 0,
        "rows": usage_rows,
        "non_canonical": True,
        "canonical_promotion": False,
        "source_bound_deck_created": False,
    }


def render_e03_5_pack(pack_pptx: Path, output_dir: Path, contact_sheet: Path, *, expected_count: int = 16) -> dict[str, Any]:
    """Render the pack through the local PPTX renderer and build a contact sheet."""

    if not pack_pptx.exists():
        return {"schema_name": "e03_5_pack_render_report", "status": "blocked", "rendered_slide_count": 0, "reason": "pack_missing"}

    from src.presentation_agent.qa.render_pptx_preview import render_pptx_preview

    output_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = output_dir / "render_manifest.json"
    render_report = render_pptx_preview(pptx_path=pack_pptx, output_dir=output_dir, manifest_path=manifest_path, backend="auto", dpi=144)
    slides = [Path(row.get("rendered_image_path") or "") for row in render_report.get("slides", []) if Path(row.get("rendered_image_path") or "").exists()]
    _render_contact_sheet(slides[:expected_count], contact_sheet)
    return {
        "schema_name": "e03_5_pack_render_report",
        "status": "passed" if len(slides) >= expected_count and contact_sheet.exists() else "failed",
        "rendered_slide_count": len(slides),
        "expected_slide_count": expected_count,
        "contact_sheet": contact_sheet.as_posix(),
        "render_manifest": manifest_path.as_posix(),
    }


def _render_contact_sheet(slides: list[Path], output: Path) -> None:
    cell_w, cell_h = 320, 205
    cols = 4
    rows = max(1, (len(slides) + cols - 1) // cols)
    sheet = Image.new("RGB", (cols * cell_w, rows * cell_h + 40), "#071018")
    draw = ImageDraw.Draw(sheet)
    font = ImageFont.load_default()
    draw.text((18, 14), "E03.5 16-template candidate pack", fill="#F8FAFC", font=font)
    for idx, path in enumerate(slides):
        x = (idx % cols) * cell_w
        y = (idx // cols) * cell_h + 40
        draw.rectangle((x, y, x + cell_w, y + cell_h), fill="#111827")
        image = Image.open(path).convert("RGB")
        image.thumbnail((cell_w - 24, cell_h - 30), Image.Resampling.LANCZOS)
        sheet.paste(image, (x + (cell_w - image.width) // 2, y + 22))
        draw.text((x + 8, y + 8), f"slide {idx + 1}", fill="#F2A900", font=font)
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)
