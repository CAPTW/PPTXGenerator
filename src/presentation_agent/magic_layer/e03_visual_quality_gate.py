"""Visual quality regression gate for the E03 editable template pack."""

from __future__ import annotations

import json
import shutil
import subprocess
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation

from scripts.run_e01x_self_describing_ps_layer_integration import protected_report, protected_snapshot
from src.presentation_agent.magic_layer.e03_archetype_registry import CORE_12_ARCHETYPE_IDS
from src.presentation_agent.magic_layer.e03_archetype_visual_identity import build_archetype_visual_identity_report
from src.presentation_agent.magic_layer.e03_placeholder_overdominance import evaluate_placeholder_overdominance, is_placeholder_text
from src.presentation_agent.magic_layer.e03_premium_gate import evaluate_premium_template_pack_gate
from src.presentation_agent.magic_layer.e03_visual_richness_metrics import build_visual_richness_score_report


REPO_ROOT = Path(__file__).resolve().parents[3]
E03_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e03_12_16_archetype_ps_layer_template_pack"
E01XP_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e01x_p_visual_slot_fidelity_patch"
E02_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e02_4core_ps_layer_archetype_conversion"
E03_VQ_ROOT = REPO_ROOT / "design_runs/run_002/outputs/magic_layer_engine_e03_vq_visual_quality_gate"


def run_e03_visual_quality_gate(input_dir: Path = E03_ROOT, output_dir: Path = E03_VQ_ROOT) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    protected_before = protected_snapshot()
    if not _run_protect_check():
        report = _protected_failure(output_dir, "protected_artifact_precheck_failed")
        return report

    input_report = validate_vq_inputs(input_dir)
    structural_decision = _read_json(input_dir / "e03_final_decision.json").get("decision") if (input_dir / "e03_final_decision.json").exists() else None
    pack_report = _read_json(input_dir / "e03_template_pack_report.json") if (input_dir / "e03_template_pack_report.json").exists() else {}
    inventory = inspect_template_pack(input_dir / "editable_template_pack.pptx")
    _render_contact_sheet(input_dir, output_dir / "e03_rendered_contact_sheet.png")

    slide_records = build_visual_slide_records(input_dir, inventory)
    identity_report = build_archetype_visual_identity_report(slide_records)
    richness_report = build_visual_richness_score_report(slide_records)
    placeholder_report = evaluate_placeholder_overdominance(slide_records)
    smart_object_report = build_smart_object_visual_field_usage_report(slide_records)
    motif_report = build_decorative_motif_usage_report(slide_records)
    e01x_regression = build_e03_vs_e01x_p_regression_report(slide_records)
    e02_regression = build_e03_vs_e02_regression_report(input_dir)
    design_regression = build_e03_visual_design_regression_report(
        identity_report=identity_report,
        richness_report=richness_report,
        placeholder_report=placeholder_report,
        smart_object_report=smart_object_report,
        motif_report=motif_report,
        e01x_regression=e01x_regression,
    )
    structural_pass = (
        structural_decision == "E03_PASS_READY_FOR_E04_SOURCE_BOUND_SMALL_DECK"
        and pack_report.get("status") == "passed"
        and input_report["status"] == "passed"
    )
    premium_gate = evaluate_premium_template_pack_gate(
        structural_native_editability_pass=structural_pass,
        semantic_raster_violations=int(_sum_archetype_metric(pack_report, "semantic_raster_violation_count")),
        unknown_content_bearing_layers=int(_sum_archetype_metric(pack_report, "unknown_content_bearing_count")),
        duplicate_bbox_collisions=int(_sum_archetype_metric(pack_report, "duplicate_bbox_collision_count")),
        visual_richness_report=richness_report,
        identity_report=identity_report,
        placeholder_report=placeholder_report,
        regression_report=e01x_regression,
        protected_artifacts_unchanged=True,
    )
    quality_override = build_e03_quality_override(str(structural_decision), premium_gate)
    e04_override = build_e04_readiness_override(quality_override, premium_gate)

    protected_after = protected_snapshot()
    protected_md, protected_ok = protected_report(protected_before, protected_after)
    protected_post_ok = _run_protect_check()
    protected_md += f"\n\n- npm protect precheck: `passed`\n- npm protect postcheck: `{'passed' if protected_post_ok else 'failed'}`\n"
    if not (protected_ok and protected_post_ok):
        premium_gate = {**premium_gate, "status": "failed", "decision": "E03_VQ_FAIL_PROTECTED_ARTIFACTS", "e04_readiness": False}
        quality_override = build_e03_quality_override(str(structural_decision), premium_gate)
        e04_override = build_e04_readiness_override(quality_override, premium_gate)

    artifacts = {
        "e03_vq_manifest.json": build_vq_manifest(input_dir, output_dir, premium_gate, quality_override),
        "e03_quality_override.json": quality_override,
        "e03_visual_design_regression_report.json": design_regression,
        "e03_vs_e01x_p_regression_report.json": e01x_regression,
        "e03_vs_e02_regression_report.json": e02_regression,
        "archetype_visual_identity_report.json": identity_report,
        "visual_richness_score_report.json": richness_report,
        "placeholder_overdominance_report.json": placeholder_report,
        "smart_object_visual_field_usage_report.json": smart_object_report,
        "decorative_motif_usage_report.json": motif_report,
        "premium_template_pack_gate_report.json": premium_gate,
        "e04_readiness_override.json": e04_override,
    }
    for filename, payload in artifacts.items():
        _write_json(output_dir / filename, payload)
    _write_md(output_dir / "e03_quality_override.md", _simple_md("E03 Quality Override", quality_override))
    _write_md(output_dir / "e03_visual_design_regression_report.md", _regression_md(design_regression))
    _write_md(output_dir / "archetype_visual_identity_report.md", _identity_md(identity_report))
    _write_md(output_dir / "visual_richness_score_report.md", _richness_md(richness_report))
    _write_md(output_dir / "placeholder_overdominance_report.md", _simple_md("Placeholder Overdominance Report", placeholder_report))
    _write_md(output_dir / "premium_template_pack_gate_report.md", _simple_md("Premium Template Pack Gate Report", premium_gate))
    _write_md(output_dir / "protected_artifact_check_report.md", protected_md)
    return premium_gate


def validate_vq_inputs(input_dir: Path) -> dict[str, Any]:
    required = [
        "editable_template_pack.pptx",
        "template_pack_registry.json",
        "component_library.json",
        "layout_selector_contract.json",
        "local_model_slot_filling_contract.json",
        "e03_final_decision.json",
        "e03_template_pack_report.json",
    ]
    missing = [name for name in required if not (input_dir / name).exists()]
    for archetype_id in CORE_12_ARCHETYPE_IDS:
        for name in ("rendered_candidate.png", "final_reference.png", "object_graph_v1.json"):
            if not (input_dir / "archetypes" / archetype_id / name).exists():
                missing.append(f"archetypes/{archetype_id}/{name}")
    if not (E01XP_ROOT / "patched_rendered_candidate.png").exists():
        missing.append(rel(E01XP_ROOT / "patched_rendered_candidate.png"))
    if not (E02_ROOT / "e02_4core_conversion_report.json").exists():
        missing.append(rel(E02_ROOT / "e02_4core_conversion_report.json"))
    return {
        "schema_name": "e03_vq_input_validation",
        "status": "passed" if not missing else "failed",
        "missing": missing,
        "canva_parity_claimed": False,
    }


def inspect_template_pack(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slides: list[dict[str, Any]] = []
    all_text: list[str] = []
    for slide_index, slide in enumerate(prs.slides):
        text_values: list[str] = []
        chart_count = 0
        table_count = 0
        media_count = 0
        connector_count = 0
        for shape in slide.shapes:
            shape_type = str(shape.shape_type)
            if bool(getattr(shape, "has_chart", False)):
                chart_count += 1
            if bool(getattr(shape, "has_table", False)):
                table_count += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_values.append(cell.text.strip())
            if "PICTURE" in shape_type or int(getattr(shape.shape_type, "value", -999)) == 13:
                media_count += 1
            if "CONNECTOR" in shape_type or "LINE" in shape_type:
                connector_count += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_values.append(shape.text.strip())
        all_text.extend(text_values)
        placeholder_count = sum(1 for value in text_values if is_placeholder_text(value))
        text_count = len(text_values)
        slides.append(
            {
                "slide_index": slide_index,
                "shape_count": len(slide.shapes),
                "text_count": text_count,
                "media_count": media_count,
                "chart_count": chart_count,
                "table_count": table_count,
                "connector_vector_count": connector_count,
                "placeholder_text_count": placeholder_count,
                "placeholder_text_ratio": round(placeholder_count / text_count, 4) if text_count else 0.0,
                "text_values": text_values,
            }
        )
    repeated = Counter(value for value in all_text if is_placeholder_text(value))
    return {
        "schema_name": "e03_template_pack_visual_inventory",
        "status": "passed",
        "pptx_path": rel(pptx_path),
        "slide_count": len(slides),
        "slides": slides,
        "shape_count": sum(slide["shape_count"] for slide in slides),
        "text_count": sum(slide["text_count"] for slide in slides),
        "media_count": sum(slide["media_count"] for slide in slides),
        "chart_count": sum(slide["chart_count"] for slide in slides),
        "table_count": sum(slide["table_count"] for slide in slides),
        "connector_vector_count": sum(slide["connector_vector_count"] for slide in slides),
        "placeholder_text_ratio": round(sum(slide["placeholder_text_count"] for slide in slides) / max(1, sum(slide["text_count"] for slide in slides)), 4),
        "repeated_placeholder_string_count": sum(count for count in repeated.values() if count > 1),
        "repeated_placeholder_strings": dict(sorted(repeated.items())),
        "canva_parity_claimed": False,
    }


def build_visual_slide_records(input_dir: Path, inventory: dict[str, Any]) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    for index, archetype_id in enumerate(CORE_12_ARCHETYPE_IDS):
        object_graph = _read_json(input_dir / "archetypes" / archetype_id / "object_graph_v1.json")
        nodes = object_graph.get("nodes", [])
        roles = [str(node.get("semantic_role")) for node in nodes]
        role_counts = Counter(roles)
        slide = dict(inventory["slides"][index]) if index < len(inventory.get("slides", [])) else {}
        decorative_count = sum(role_counts[role] for role in ("decorative_texture", "technical_overlay", "card_underline", "semantic_icon", "active_marker", "progress_indicator"))
        accent_count = max(0, int(slide.get("shape_count", 0)) - int(slide.get("text_count", 0)) - int(slide.get("chart_count", 0)) - int(slide.get("table_count", 0)) - int(slide.get("media_count", 0)) - 1)
        component_count = len(
            {
                role
                for role in roles
                if role
                and role
                not in {
                    "background_base",
                    "title_text_region",
                    "subtitle_text_region",
                    "source_footer_text",
                }
            }
        )
        record = {
            **slide,
            "archetype_id": archetype_id,
            "roles": roles,
            "role_counts": dict(sorted(role_counts.items())),
            "connector_count": int(slide.get("connector_vector_count", 0)),
            "decorative_motif_count": decorative_count,
            "accent_shape_count": accent_count,
            "archetype_component_count": component_count,
            "card_panel_count": role_counts["card_panel"],
            "hero_visual_field_count": role_counts["hero_visual_field"],
            "has_footer_system": role_counts["source_footer_strip"] >= 1 or role_counts["source_footer_text"] >= 1,
            "requires_visual_field": archetype_id in {"cover_hero", "standard_content"},
            "placeholder_ratio": float(slide.get("placeholder_text_ratio", 0.0)),
        }
        record["metrics"] = {key: value for key, value in record.items() if key not in {"roles", "role_counts", "text_values", "metrics"}}
        records.append(record)
    return records


def build_e03_quality_override(structural_decision: str, premium_gate: dict[str, Any]) -> dict[str, Any]:
    passed = premium_gate.get("status") == "passed"
    return {
        "schema_name": "e03_quality_override",
        "structural_decision": structural_decision,
        "visual_quality_decision": "E03_VISUAL_DESIGN_PASS" if passed else "E03_PATCH_VISUAL_DESIGN_REQUIRED",
        "e04_unlock": bool(passed),
        "reason": "structural/native editability passed but premium visual design failed" if not passed else "structural/native editability and premium visual design passed",
        "premium_gate_decision": premium_gate.get("decision"),
        "failures": premium_gate.get("failures", []),
        "canva_parity_claimed": False,
    }


def build_e04_readiness_override(quality_override: dict[str, Any], premium_gate: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "e04_readiness_override",
        "status": "ready" if quality_override["e04_unlock"] else "blocked",
        "e04_unlock": quality_override["e04_unlock"],
        "structural_decision_preserved": quality_override["structural_decision"],
        "visual_quality_decision": quality_override["visual_quality_decision"],
        "premium_gate_decision": premium_gate.get("decision"),
        "e04_started": False,
        "source_bound_deck_generated": False,
        "large_deck_generated": False,
        "canva_parity_claimed": False,
    }


def build_e03_visual_design_regression_report(**reports: dict[str, Any]) -> dict[str, Any]:
    failures: list[str] = []
    for report in reports.values():
        if report.get("status") != "passed":
            failures.extend(report.get("failures", [report.get("schema_name", "unknown_report_failed")]))
    return {
        "schema_name": "e03_visual_design_regression_report",
        "status": "passed" if not failures else "failed",
        "structural_native_editability_separated": True,
        "visual_quality_regression_detected": bool(failures),
        "reports": reports,
        "failures": sorted(set(str(failure) for failure in failures)),
        "canva_parity_claimed": False,
    }


def build_e03_vs_e01x_p_regression_report(slide_records: list[dict[str, Any]]) -> dict[str, Any]:
    standard = next((record for record in slide_records if record["archetype_id"] == "standard_content"), {})
    fixture_exists = (E01XP_ROOT / "patched_rendered_candidate.png").exists()
    failures: list[str] = []
    if not fixture_exists:
        failures.append("e01x_p_patched_standard_content_fixture_missing")
    if float(standard.get("placeholder_text_ratio", 0.0)) >= 0.70:
        failures.append("standard_content_visual_language_reads_as_wireframe")
    if int(standard.get("media_count", 0)) == 0:
        failures.append("standard_content_loses_hero_texture_or_visual_field_materiality")
    if int(standard.get("decorative_motif_count", 0)) < 3 or int(standard.get("connector_vector_count", 0)) < 1:
        failures.append("card_underlines_icon_connector_motifs_degraded_or_generic")
    if failures:
        failures.append("standard_content_visually_simpler_than_e01x_p")
    return {
        "schema_name": "e03_vs_e01x_p_regression_report",
        "status": "passed" if not failures else "failed",
        "e01x_p_fixture": rel(E01XP_ROOT / "patched_rendered_candidate.png"),
        "e03_standard_content_metrics": {
            "placeholder_text_ratio": standard.get("placeholder_text_ratio"),
            "media_count": standard.get("media_count"),
            "connector_vector_count": standard.get("connector_vector_count"),
            "decorative_motif_count": standard.get("decorative_motif_count"),
        },
        "failures": sorted(set(failures)),
        "canva_parity_claimed": False,
    }


def build_e03_vs_e02_regression_report(input_dir: Path) -> dict[str, Any]:
    e02 = _read_json(E02_ROOT / "e02_4core_conversion_report.json") if (E02_ROOT / "e02_4core_conversion_report.json").exists() else {"archetypes": {}}
    e03 = _read_json(input_dir / "e03_template_pack_report.json") if (input_dir / "e03_template_pack_report.json").exists() else {"archetypes": {}}
    rows = []
    failures: list[str] = []
    for archetype_id in ("cover_hero", "standard_content", "data_dashboard", "table_heavy"):
        e02_row = e02.get("archetypes", {}).get(archetype_id, {})
        e03_row = e03.get("archetypes", {}).get(archetype_id, {})
        row = {
            "archetype_id": archetype_id,
            "e02_status": e02_row.get("status"),
            "e03_status": e03_row.get("status"),
            "e02_native_chart_table_decision": e02_row.get("native_chart_table_decision"),
            "e03_native_chart_table_decision": e03_row.get("native_chart_table_decision"),
            "e03_semantic_raster_violation_count": e03_row.get("semantic_raster_violation_count"),
            "e03_unknown_content_bearing_count": e03_row.get("unknown_content_bearing_count"),
            "e03_duplicate_bbox_collision_count": e03_row.get("duplicate_bbox_collision_count"),
        }
        rows.append(row)
        if e03_row.get("status") != "passed":
            failures.append(f"{archetype_id}_e03_structural_regression")
    return {
        "schema_name": "e03_vs_e02_regression_report",
        "status": "passed" if not failures else "failed",
        "rows": rows,
        "failures": failures,
        "canva_parity_claimed": False,
    }


def build_smart_object_visual_field_usage_report(slide_records: list[dict[str, Any]]) -> dict[str, Any]:
    rows = []
    failures = []
    for record in slide_records:
        requires = bool(record.get("requires_visual_field"))
        row = {
            "archetype_id": record["archetype_id"],
            "requires_visual_field": requires,
            "hero_visual_field_count": record.get("hero_visual_field_count", 0),
            "media_count": record.get("media_count", 0),
            "status": "passed",
        }
        if requires and int(record.get("hero_visual_field_count", 0)) < 1:
            row["status"] = "failed"
            failures.append(f"{record['archetype_id']}_hero_visual_field_missing")
        elif requires and int(record.get("media_count", 0)) == 0:
            row["status"] = "warning"
            failures.append(f"{record['archetype_id']}_visual_field_rendered_as_generic_shape")
        rows.append(row)
    return {
        "schema_name": "smart_object_visual_field_usage_report",
        "status": "passed" if not failures else "failed",
        "rows": rows,
        "failures": failures,
        "canva_parity_claimed": False,
    }


def build_decorative_motif_usage_report(slide_records: list[dict[str, Any]]) -> dict[str, Any]:
    rows = [
        {
            "archetype_id": record["archetype_id"],
            "decorative_motif_count": record.get("decorative_motif_count", 0),
            "connector_vector_count": record.get("connector_vector_count", 0),
        }
        for record in slide_records
    ]
    slides_with_motifs = sum(1 for row in rows if int(row["decorative_motif_count"]) + int(row["connector_vector_count"]) > 0)
    failures = []
    if slides_with_motifs < 8:
        failures.append("decorative_motif_usage_absent_across_most_slides")
    return {
        "schema_name": "decorative_motif_usage_report",
        "status": "passed" if not failures else "failed",
        "slides_with_motifs": slides_with_motifs,
        "rows": rows,
        "failures": failures,
        "canva_parity_claimed": False,
    }


def build_vq_manifest(input_dir: Path, output_dir: Path, premium_gate: dict[str, Any], quality_override: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema_name": "e03_vq_manifest",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "input_dir": rel(input_dir),
        "output_dir": rel(output_dir),
        "final_decision": premium_gate.get("decision"),
        "structural_decision": quality_override.get("structural_decision"),
        "visual_quality_decision": quality_override.get("visual_quality_decision"),
        "e04_unlocked": quality_override.get("e04_unlock"),
        "e04_started": False,
        "source_bound_deck_generated": False,
        "large_deck_generated": False,
        "canonical_promotion": False,
        "canva_parity_claimed": False,
    }


def _render_contact_sheet(input_dir: Path, output_path: Path) -> Path:
    thumbs = []
    for archetype_id in CORE_12_ARCHETYPE_IDS:
        path = input_dir / "archetypes" / archetype_id / "rendered_candidate.png"
        if path.exists():
            thumb = Image.open(path).convert("RGB").resize((420, 236))
        else:
            thumb = Image.new("RGB", (420, 236), "#061526")
            ImageDraw.Draw(thumb).text((18, 18), archetype_id, fill=(248, 250, 252))
        thumbs.append(thumb)
    sheet = Image.new("RGB", (4 * 420, 3 * 236), "#061526")
    for index, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((index % 4) * 420, (index // 4) * 236))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    return output_path


def _sum_archetype_metric(report: dict[str, Any], metric: str) -> int:
    return sum(int(value.get(metric, 0) or 0) for value in report.get("archetypes", {}).values())


def _protected_failure(output_dir: Path, reason: str) -> dict[str, Any]:
    payload = {
        "schema_name": "premium_template_pack_gate_report",
        "status": "failed",
        "decision": "E03_VQ_FAIL_PROTECTED_ARTIFACTS",
        "reason": reason,
        "e04_readiness": False,
        "canva_parity_claimed": False,
    }
    _write_json(output_dir / "premium_template_pack_gate_report.json", payload)
    _write_json(output_dir / "e04_readiness_override.json", {"schema_name": "e04_readiness_override", "status": "blocked", "e04_unlock": False, "reason": reason, "canva_parity_claimed": False})
    return payload


def _run_protect_check() -> bool:
    npm = shutil.which("npm.cmd") or shutil.which("npm")
    if not npm:
        return False
    return subprocess.run([npm, "run", "protect:check"], cwd=REPO_ROOT, capture_output=True, text=True, check=False).returncode == 0


def _simple_md(title: str, payload: dict[str, Any]) -> str:
    lines = [f"# {title}", ""]
    for key, value in payload.items():
        if not isinstance(value, (dict, list)):
            lines.append(f"- {key}: `{value}`")
    if payload.get("failures"):
        lines.extend(["", "## Failures"])
        lines.extend(f"- `{failure}`" for failure in payload["failures"])
    return "\n".join(lines)


def _regression_md(report: dict[str, Any]) -> str:
    lines = ["# E03 Visual Design Regression Report", "", f"- status: `{report['status']}`", f"- visual_quality_regression_detected: `{report['visual_quality_regression_detected']}`", ""]
    for failure in report.get("failures", []):
        lines.append(f"- `{failure}`")
    return "\n".join(lines)


def _identity_md(report: dict[str, Any]) -> str:
    lines = ["# Archetype Visual Identity Report", "", f"- status: `{report['status']}`", ""]
    for archetype_id, row in report.get("archetypes", {}).items():
        lines.append(f"- {archetype_id}: `{row['status']}`")
    return "\n".join(lines)


def _richness_md(report: dict[str, Any]) -> str:
    lines = ["# Visual Richness Score Report", "", f"- status: `{report['status']}`", f"- average_visual_richness_score: `{report['average_visual_richness_score']}`", ""]
    for score in report.get("slide_scores", []):
        lines.append(f"- {score['archetype_id']}: `{score['score']}`")
    return "\n".join(lines)


def _read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")


def _write_md(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.strip() + "\n", encoding="utf-8")


def rel(path: Path) -> str:
    try:
        return path.resolve().relative_to(REPO_ROOT.resolve()).as_posix()
    except ValueError:
        return path.resolve().as_posix()
