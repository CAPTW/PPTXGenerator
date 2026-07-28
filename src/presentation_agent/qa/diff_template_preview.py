"""Approximate render/template QA report for template preview decks."""

from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_PPTX_PATH = Path("outputs/template_preview.pptx")
DEFAULT_RENDER_DIR = Path("outputs/template_preview_png")
DEFAULT_TEMPLATE_SPEC = Path("outputs/editable_template_spec.json")
DEFAULT_PREVIEW_MANIFEST = Path("outputs/template_preview_manifest.json")
DEFAULT_TEMPLATE_IMAGE_MANIFEST = Path("outputs/template_images/template_image_manifest.json")
DEFAULT_JSON_REPORT = Path("outputs/template_diff_report.json")
DEFAULT_MD_REPORT = Path("outputs/template_diff_report.md")


def build_template_diff_report(
    *,
    pptx_path: str | Path = DEFAULT_PPTX_PATH,
    render_dir: str | Path = DEFAULT_RENDER_DIR,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    preview_manifest_path: str | Path = DEFAULT_PREVIEW_MANIFEST,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    render_report_path: str | Path | None = None,
) -> dict[str, Any]:
    pptx_file = Path(pptx_path)
    render_root = Path(render_dir)
    template_spec = _load_json(template_spec_path)
    preview_manifest = _load_json(preview_manifest_path) if Path(preview_manifest_path).exists() else {}
    template_image_manifest = _load_json(template_image_manifest_path) if Path(template_image_manifest_path).exists() else {}
    render_report = _load_json(render_report_path) if render_report_path and Path(render_report_path).exists() else None

    deck = Presentation(pptx_file)
    layouts = {layout["layout_id"]: layout for layout in template_spec.get("layouts", [])}
    reference_by_archetype = _reference_images_by_archetype(template_image_manifest)
    compiled_layouts = preview_manifest.get("compiled_layouts") or _fallback_compiled_layouts(template_spec, len(deck.slides))
    slides: list[dict[str, Any]] = []
    all_findings: list[dict[str, Any]] = []

    for index, pptx_slide in enumerate(deck.slides, start=1):
        compiled = compiled_layouts[index - 1] if index - 1 < len(compiled_layouts) else {}
        layout_id = compiled.get("layout_id") or ""
        archetype_id = compiled.get("archetype_id") or layouts.get(layout_id, {}).get("archetype_id")
        rendered_path = render_root / f"slide-{index:03d}.png"
        reference_path = Path(reference_by_archetype[archetype_id]) if archetype_id in reference_by_archetype else None
        slide_report = _inspect_slide(
            slide_index=index,
            pptx_slide=pptx_slide,
            layout=layouts.get(layout_id),
            layout_id=layout_id,
            archetype_id=archetype_id,
            rendered_path=rendered_path,
            reference_path=reference_path,
        )
        slides.append(slide_report)
        all_findings.extend(slide_report["warnings"])

    if render_report and render_report.get("status") == "skipped":
        all_findings.append(
            {
                "code": "RENDER_SKIPPED",
                "severity": "warning",
                "message": "PNG render was skipped because no local renderer was available.",
            }
        )

    severe_findings = [finding for finding in all_findings if finding.get("severity") == "severe"]
    warning_findings = [finding for finding in all_findings if finding.get("severity") == "warning"]
    report = {
        "schema_name": "template_diff_report",
        "schema_version": "1.0",
        "pptx_path": _display_path(pptx_file),
        "render_dir": _display_path(render_root),
        "template_spec_path": _display_path(Path(template_spec_path)),
        "slide_count": len(slides),
        "status": "failed" if severe_findings else "issues_reported" if warning_findings else "passed",
        "qa_blocks_deck_generation": bool(severe_findings),
        "findings_summary": {
            "total": len(all_findings),
            "severe": len(severe_findings),
            "warning": len(warning_findings),
            "info": sum(1 for finding in all_findings if finding.get("severity") == "info"),
        },
        "render_status": render_report.get("status") if render_report else "not_provided",
        "slides": slides,
        "findings": all_findings,
    }
    return report


def build_template_diff_report_from_files(
    *,
    pptx_path: str | Path = DEFAULT_PPTX_PATH,
    render_dir: str | Path = DEFAULT_RENDER_DIR,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    preview_manifest_path: str | Path = DEFAULT_PREVIEW_MANIFEST,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    render_report_path: str | Path | None = None,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
) -> Path:
    report = build_template_diff_report(
        pptx_path=pptx_path,
        render_dir=render_dir,
        template_spec_path=template_spec_path,
        preview_manifest_path=preview_manifest_path,
        template_image_manifest_path=template_image_manifest_path,
        render_report_path=render_report_path,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return json_path


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Build an approximate template preview diff QA report.")
    parser.add_argument("--pptx", type=Path, default=DEFAULT_PPTX_PATH)
    parser.add_argument("--render-dir", type=Path, default=DEFAULT_RENDER_DIR)
    parser.add_argument("--template-spec", type=Path, default=DEFAULT_TEMPLATE_SPEC)
    parser.add_argument("--preview-manifest", type=Path, default=DEFAULT_PREVIEW_MANIFEST)
    parser.add_argument("--template-image-manifest", type=Path, default=DEFAULT_TEMPLATE_IMAGE_MANIFEST)
    parser.add_argument("--render-report", type=Path, default=DEFAULT_RENDER_DIR / "render_preview_report.json")
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        output = build_template_diff_report_from_files(
            pptx_path=args.pptx,
            render_dir=args.render_dir,
            template_spec_path=args.template_spec,
            preview_manifest_path=args.preview_manifest,
            template_image_manifest_path=args.template_image_manifest,
            render_report_path=args.render_report,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"DIFF_TEMPLATE_PREVIEW_FAILED {exc}")
        return 1
    print(f"WROTE {output}")
    return 0


def _inspect_slide(
    *,
    slide_index: int,
    pptx_slide: Any,
    layout: dict[str, Any] | None,
    layout_id: str,
    archetype_id: str | None,
    rendered_path: Path,
    reference_path: Path | None,
) -> dict[str, Any]:
    warnings: list[dict[str, Any]] = []
    rendered_metrics = _image_metrics(rendered_path) if rendered_path.exists() else None
    reference_metrics = _image_metrics(reference_path) if reference_path and reference_path.exists() else None
    structural = _slide_structural_metrics(pptx_slide)

    if rendered_metrics is None:
        warnings.append(_warning("RENDERED_IMAGE_MISSING", "warning", "Rendered slide image is missing.", slide_index))
    else:
        if abs((rendered_metrics["width_px"] / rendered_metrics["height_px"]) - (16 / 9)) > 0.08:
            warnings.append(_warning("CANVAS_SIZE_MISMATCH", "warning", "Rendered slide is outside expected 16:9 tolerance.", slide_index))
        if rendered_metrics["blank_area_ratio"] > 0.92:
            warnings.append(_warning("BLANK_AREA_WARNING", "warning", "Rendered slide appears mostly blank.", slide_index))

    palette_difference = None
    canvas_size_mismatch = None
    if rendered_metrics and reference_metrics:
        palette_difference = _palette_distance(rendered_metrics["mean_rgb"], reference_metrics["mean_rgb"])
        canvas_size_mismatch = (rendered_metrics["width_px"], rendered_metrics["height_px"]) != (
            reference_metrics["width_px"],
            reference_metrics["height_px"],
        )
    elif rendered_metrics:
        canvas_size_mismatch = False

    if structural["full_slide_picture_count"]:
        warnings.append(_warning("FULL_SLIDE_IMAGE_VIOLATION", "severe", "A picture shape covers most of the slide.", slide_index))
    if structural["editable_object_count"] <= 0:
        warnings.append(_warning("NO_EDITABLE_CONTENT_OBJECTS", "severe", "No editable text/table/chart/shape content objects were detected.", slide_index))
    if layout is not None:
        expected_slots = {slot["slot_id"] for slot in layout.get("slots", [])}
        if "title" in expected_slots and not structural["has_title_text"]:
            warnings.append(_warning("MISSING_TITLE", "severe", "Expected title slot was not detectable in PPTX text.", slide_index))
        if archetype_id != "cover_hero" and "footer" in expected_slots and not structural["has_footer_text"]:
            warnings.append(_warning("MISSING_FOOTER", "severe", "Expected footer system was not detectable in PPTX text.", slide_index))
    if structural["text_overflow_risk_count"]:
        warnings.append(_warning("TEXT_OVERFLOW_RISK", "warning", "One or more text boxes may overflow by rough character-area heuristic.", slide_index))
    if structural["object_count"] < 4:
        warnings.append(_warning("OBJECT_COUNT_LOW", "warning", "Slide has a low editable object count for a template preview.", slide_index))

    return {
        "slide_index": slide_index,
        "layout_id": layout_id,
        "archetype_id": archetype_id,
        "rendered_image_path": _display_path(rendered_path),
        "reference_image_path": _display_path(reference_path) if reference_path else None,
        "canvas_size_mismatch": canvas_size_mismatch,
        "rough_color_palette_difference": palette_difference,
        "blank_area_ratio": rendered_metrics["blank_area_ratio"] if rendered_metrics else None,
        "text_overflow_warning": structural["text_overflow_risk_count"] > 0,
        "object_count": structural["object_count"],
        "object_count_warning": structural["object_count"] < 4,
        "missing_footer_warning": any(warning["code"] == "MISSING_FOOTER" for warning in warnings),
        "missing_title_warning": any(warning["code"] == "MISSING_TITLE" for warning in warnings),
        "full_slide_image_violation_warning": structural["full_slide_picture_count"] > 0,
        "warnings": warnings,
    }


def _slide_structural_metrics(slide: Any) -> dict[str, Any]:
    object_count = 0
    editable_count = 0
    full_slide_picture_count = 0
    has_title_text = False
    has_footer_text = False
    text_overflow_risk_count = 0
    for shape in slide.shapes:
        object_count += 1
        if getattr(shape, "has_text_frame", False) or shape.has_table or getattr(shape, "has_chart", False):
            editable_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and _is_full_slide_shape(shape):
            full_slide_picture_count += 1
        text = str(getattr(shape, "text", "") or "")
        if "title" in text.lower() or len(text.strip()) >= 8 and shape.top < 1_300_000:
            has_title_text = True
        if "footer" in text.lower() or "source" in text.lower() or shape.top > 5_900_000:
            has_footer_text = True
        if getattr(shape, "has_text_frame", False) and _text_overflow_risk(shape, text):
            text_overflow_risk_count += 1
    return {
        "object_count": object_count,
        "editable_object_count": editable_count,
        "full_slide_picture_count": full_slide_picture_count,
        "has_title_text": has_title_text,
        "has_footer_text": has_footer_text,
        "text_overflow_risk_count": text_overflow_risk_count,
    }


def _image_metrics(path: Path | None) -> dict[str, Any] | None:
    if path is None or not path.exists():
        return None
    with Image.open(path) as image:
        rgb = image.convert("RGB")
        stat = ImageStat.Stat(rgb)
        grayscale = rgb.convert("L")
        lum = ImageStat.Stat(grayscale)
        width, height = rgb.size
        histogram = grayscale.histogram()
        whiteish = sum(histogram[242:])
        blackish = sum(histogram[:12])
        total = width * height
        blank_area_ratio = max(whiteish, blackish) / total if total else 1.0
        return {
            "width_px": width,
            "height_px": height,
            "mean_rgb": tuple(float(value) for value in stat.mean),
            "stddev_luminance": float(lum.stddev[0]),
            "blank_area_ratio": round(blank_area_ratio, 6),
        }


def _reference_images_by_archetype(template_image_manifest: dict[str, Any]) -> dict[str, str]:
    result: dict[str, str] = {}
    for record in template_image_manifest.get("images") or []:
        if isinstance(record, dict) and record.get("archetype_id") and record.get("image_output_path"):
            result[str(record["archetype_id"])] = str(record["image_output_path"])
    return result


def _fallback_compiled_layouts(template_spec: dict[str, Any], slide_count: int) -> list[dict[str, Any]]:
    layouts = template_spec.get("layouts") or []
    return [
        {
            "layout_id": layout.get("layout_id"),
            "archetype_id": layout.get("archetype_id"),
            "slide_number": index + 1,
        }
        for index, layout in enumerate(layouts[:slide_count])
    ]


def _text_overflow_risk(shape: Any, text: str) -> bool:
    if not text.strip():
        return False
    width_in = shape.width / 914400
    height_in = shape.height / 914400
    capacity = max(20, width_in * height_in * 85)
    return len(text) > capacity


def _is_full_slide_shape(shape: Any) -> bool:
    width_in = shape.width / 914400
    height_in = shape.height / 914400
    left_in = shape.left / 914400
    top_in = shape.top / 914400
    return left_in <= 0.15 and top_in <= 0.15 and width_in >= 12.0 and height_in >= 6.7


def _palette_distance(first: tuple[float, float, float], second: tuple[float, float, float]) -> float:
    return round(math.sqrt(sum((a - b) ** 2 for a, b in zip(first, second))) / 441.67295593, 6)


def _warning(code: str, severity: str, message: str, slide_index: int | None = None) -> dict[str, Any]:
    payload = {"code": code, "severity": severity, "message": message}
    if slide_index is not None:
        payload["slide_index"] = slide_index
    return payload


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Template Diff QA Report",
        "",
        f"Status: `{report['status']}`",
        f"Slides: `{report['slide_count']}`",
        f"Findings: `{report['findings_summary']['total']}` total, `{report['findings_summary']['severe']}` severe",
        f"QA blocks deck generation: `{str(report['qa_blocks_deck_generation']).lower()}`",
        "",
        "| Slide | Layout | Render | Reference | Findings |",
        "|---:|---|---|---|---|",
    ]
    for slide in report["slides"]:
        findings = ", ".join(warning["code"] for warning in slide["warnings"]) or "none"
        lines.append(
            f"| {slide['slide_index']} | `{slide['layout_id']}` | `{slide['rendered_image_path']}` | `{slide.get('reference_image_path') or ''}` | {findings} |"
        )
    lines.append("")
    if report["findings"]:
        lines.append("## Findings")
        lines.append("")
        for finding in report["findings"]:
            slide = f" slide {finding['slide_index']}" if "slide_index" in finding else ""
            lines.append(f"- `{finding['severity']}` `{finding['code']}`{slide}: {finding['message']}")
    return "\n".join(lines) + "\n"


def _load_json(path: str | Path | None) -> Any:
    if path is None:
        return None
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _display_path(path: Path | None) -> str | None:
    if path is None:
        return None
    return str(path.as_posix())


if __name__ == "__main__":
    raise SystemExit(main())
