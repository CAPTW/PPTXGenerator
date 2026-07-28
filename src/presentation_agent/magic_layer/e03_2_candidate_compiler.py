"""Compile the E03.2 visual_toc golden placement candidate."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 16.0
SLIDE_H = 9.0
COLORS = {
    "bg": "06111A",
    "deep": "071522",
    "panel": "07344A",
    "panel2": "0A5260",
    "cyan": "4AD6E8",
    "cyan2": "1B8EA3",
    "gold": "D89400",
    "gold2": "F5A623",
    "paper": "F7F3EA",
    "paper2": "EDE8DE",
    "ink": "071018",
    "muted": "9EC4C8",
    "white": "FFFFFF",
}


def compile_e03_2_candidate(output_pptx: Path, graph: dict[str, Any], reconstruction_plan: dict[str, Any]) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_visual_toc_golden(slide)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_name": "e03_2_candidate_compile_report",
        "status": "passed",
        "target_archetype": graph["target_archetype"],
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 1,
        "object_graph_nodes": len(graph["nodes"]),
        "reconstruction_plan_status": reconstruction_plan["status"],
        "semantic_text_as_ppt_text": True,
        "semantic_icons_as_native_vector": True,
        "native_ppt_chart_count": 0,
        "native_ppt_table_count": 0,
        "editable_shape_chart_count": 0,
        "editable_shape_grid_table_count": 0,
        "raster_media_count": 0,
        "full_slide_raster_count": 0,
        "screenshot_slide_count": 0,
        "semantic_raster_violation_count": 0,
    }


def draw_visual_toc_golden(slide: Any) -> None:
    _rect(slide, "background_base", 0, 0, SLIDE_W, SLIDE_H, "bg", None)
    _header(slide)
    _main_stage(slide)
    _progress_path(slide)
    _module_cards(slide)
    _reading_path(slide)
    _right_meta_panel(slide)
    _footer(slide)


def _header(slide: Any) -> None:
    _rect(slide, "dark_header_region", 0, 0, SLIDE_W, 2.32, "deep", None)
    for idx, size in enumerate((2.4, 1.8, 1.2), start=1):
        _shape(slide, f"header_left_tech_ring_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, -0.95 - idx * 0.04, -0.55 - idx * 0.04, size, size, None, "cyan2", 0.45)
    for idx in range(10):
        _shape(slide, f"header_left_dot_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, 0.52 + idx * 0.04, 1.24 + idx * 0.035, 0.035, 0.035, "cyan2", None)
    _text(slide, "title_index_label", "INDEX", 1.45, 0.38, 0.72, 0.18, 10.5, "gold2", bold=True)
    _line(slide, "title_index_rule", 2.22, 0.49, 4.98, 0.49, "gold2", 0.95)
    _shape(slide, "title_index_rule_dot", MSO_AUTO_SHAPE_TYPE.OVAL, 4.96, 0.455, 0.085, 0.085, "gold2", None)
    _text(slide, "title_region", "TITLE", 1.40, 0.82, 2.70, 0.62, 36, "white", bold=True)
    _line(slide, "title_gold_underline", 1.43, 1.82, 2.16, 1.82, "gold2", 1.5)
    _line(slide, "header_meta_divider", 4.95, 0.86, 4.95, 1.77, "gold2", 0.9)
    _text(slide, "header_meta_region_index", "INDEX", 5.48, 1.13, 0.64, 0.14, 9.2, "gold2", bold=True)
    _text(slide, "header_meta_region_meta", "META", 5.48, 1.41, 0.55, 0.12, 8.0, "cyan", bold=True)
    for idx, (x, y) in enumerate(((11.36, 0.28), (12.42, 0.42), (13.10, 0.78), (14.78, 0.62), (15.05, 0.38)), start=1):
        _shape(slide, f"technical_overlay_node_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, 0.09, 0.09, "cyan", "cyan", 0.35)
    traces = [
        (11.85, 1.18, 13.65, 1.18),
        (12.1, 0.88, 14.2, 0.88),
        (13.0, 0.55, 14.7, 0.55),
        (13.8, 0.35, 15.4, 0.35),
        (11.9, 1.48, 15.85, 1.48),
    ]
    for idx, (x1, y1, x2, y2) in enumerate(traces, start=1):
        _line(slide, f"technical_overlay_trace_{idx}", x1, y1, x2, y2, "cyan2", 0.8)
    for r in range(6):
        for c in range(14):
            _shape(slide, f"technical_overlay_dot_{r}_{c}", MSO_AUTO_SHAPE_TYPE.OVAL, 11.3 + c * 0.12, 0.06 + r * 0.12, 0.025, 0.025, "muted", None)


def _main_stage(slide: Any) -> None:
    _rect(slide, "main_stage_region", 0, 2.25, SLIDE_W, 5.55, "paper", "paper2", 0.25)
    _shape(slide, "main_stage_header_notch_dark", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 11.35, 1.92, 4.75, 0.42, "deep", None)
    _line(slide, "main_stage_top_gold_rule_left", 0.0, 2.25, 11.35, 2.25, "gold2", 1.2)
    _line(slide, "main_stage_top_gold_rule_notch", 11.35, 2.25, 11.78, 1.82, "gold2", 1.2)
    _line(slide, "main_stage_top_gold_rule_right", 11.78, 1.82, 16.0, 1.82, "gold2", 1.2)


def _progress_path(slide: Any) -> None:
    centers = [1.34, 3.47, 5.78, 7.95, 9.98, 11.78]
    y = 3.05
    _line(slide, "progress_path_region", centers[0] + 0.2, y, centers[-1] - 0.2, y, "cyan2", 0.95)
    for idx in range(5):
        _line(slide, f"progress_dotted_connector_{idx}", centers[idx] + 0.36, y, centers[idx + 1] - 0.36, y, "gold2" if idx == 1 else "cyan2", 0.8)
        _shape(slide, f"progress_arrow_{idx}", MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, centers[idx + 1] - 0.55, y - 0.065, 0.12, 0.13, "deep", "deep", 0.2)
    for idx, cx in enumerate(centers, start=1):
        active = idx == 2
        size = 0.78 if active else 0.54
        _shape(slide, f"progress_circle_{idx:02d}", MSO_AUTO_SHAPE_TYPE.OVAL, cx - size / 2, y - size / 2, size, size, "gold" if active else "panel", "cyan", 1.0)
        if active:
            _shape(slide, f"progress_circle_{idx:02d}_halo", MSO_AUTO_SHAPE_TYPE.OVAL, cx - 0.43, y - 0.43, 0.86, 0.86, None, "gold2", 1.2)
        _text(slide, f"progress_number_{idx:02d}", f"{idx:02d}", cx - 0.17, y - 0.12, 0.34, 0.16, 11, "white", bold=True)


def _module_cards(slide: Any) -> None:
    cards = [
        ("module_card_01", 0.40, "target", False),
        ("module_card_02_active", 2.42, "database", True),
        ("module_card_03", 4.78, "network", False),
        ("module_card_04", 6.92, "shield", False),
        ("module_card_05", 9.08, "chart", False),
        ("module_card_06", 11.02, "document", False),
    ]
    y, h = 3.48, 3.38
    for idx, (name, x, icon, active) in enumerate(cards, start=1):
        w = 2.12 if active else 1.86
        fill = "paper" if active else "panel"
        line = "gold2" if active else "cyan2"
        _shape(slide, name, MSO_AUTO_SHAPE_TYPE.SNIP_2_DIAG_RECTANGLE, x, y, w, h, fill, line, 1.1)
        _line(slide, f"{name}_title_rule", x + 0.31, y + 0.68, x + w - 0.31, y + 0.68, "gold" if active else "cyan", 0.45)
        _text(slide, f"{name}_index_label", "INDEX", x + 0.68, y + 0.38, 0.64, 0.11, 8.0, "gold" if active else "cyan", bold=True)
        _focus_corners(slide, f"{name}_focus", x + w / 2 - 0.34, y + 1.08, 0.68, "gold" if active else "cyan")
        _icon(slide, f"{name}_icon", icon, x + w / 2 - 0.28, y + 1.18, 0.56, "gold" if active else "cyan")
        _text(slide, f"{name}_icon_text", "ICON", x + w / 2 - 0.26, y + 2.06, 0.52, 0.12, 7.5, "gold" if active else "cyan", bold=True)
        _rect(slide, f"{name}_bottom_band", x, y + 2.55, w, 0.70, "gold" if active else "deep", None)
        _text(slide, f"{name}_module_label", "MODULE", x + 0.55, y + 2.76, 0.76, 0.16, 8.5, "white", bold=True)
        for dot in range(7 if active else 6):
            dx = x + 0.43 + dot * 0.18
            _shape(slide, f"{name}_status_dot_{dot}", MSO_AUTO_SHAPE_TYPE.OVAL, dx, y + 3.08, 0.09, 0.09, None if dot != idx % 6 else "cyan", "white" if active else "cyan", 0.55)


def _reading_path(slide: Any) -> None:
    _line(slide, "reading_path_region", 0.95, 7.28, 5.45, 7.28, "gold2", 0.8)
    _shape(slide, "reading_path_start_dot", MSO_AUTO_SHAPE_TYPE.OVAL, 0.89, 7.22, 0.12, 0.12, "paper", "gold2", 0.8)
    _text(slide, "reading_path_text", "READING PATH", 5.68, 7.21, 1.20, 0.12, 8.2, "gold", bold=True)
    _line(slide, "reading_path_region_right", 7.20, 7.28, 11.98, 7.28, "gold2", 0.8)
    _line(slide, "reading_path_upturn", 11.98, 7.28, 11.98, 6.92, "gold2", 0.8)
    _shape(slide, "reading_path_arrow", MSO_AUTO_SHAPE_TYPE.UP_ARROW, 11.90, 6.84, 0.16, 0.20, "gold2", None)


def _right_meta_panel(slide: Any) -> None:
    _shape(slide, "right_meta_panel", MSO_AUTO_SHAPE_TYPE.SNIP_2_DIAG_RECTANGLE, 12.76, 2.18, 2.92, 5.15, "panel", "cyan", 0.85)
    _rect(slide, "right_meta_panel_header", 12.77, 2.18, 2.85, 0.42, "panel2", None)
    _text(slide, "right_meta_panel_title", "META", 13.02, 2.34, 0.56, 0.14, 8.8, "white", bold=True)
    _shape(slide, "right_meta_icon_dotted_ring", MSO_AUTO_SHAPE_TYPE.OVAL, 13.62, 3.02, 1.42, 1.42, None, "cyan", 0.65)
    _icon(slide, "right_meta_book_icon", "book", 13.88, 3.38, 0.88, "cyan")
    _text(slide, "right_meta_icon_text", "ICON", 13.96, 4.35, 0.52, 0.12, 7.0, "cyan", bold=True)
    _text(slide, "right_meta_label", "META", 13.05, 5.25, 0.58, 0.12, 8.0, "cyan", bold=True)
    _line(slide, "right_meta_label_rule", 13.05, 5.50, 15.18, 5.50, "cyan2", 0.55)
    markers = [(MSO_AUTO_SHAPE_TYPE.OVAL, "circle"), (MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, "triangle"), (MSO_AUTO_SHAPE_TYPE.RECTANGLE, "square")]
    for idx, (shape_type, label) in enumerate(markers):
        y = 5.82 + idx * 0.45
        _shape(slide, f"right_meta_marker_{label}", shape_type, 13.08, y, 0.18, 0.18, None, "cyan", 0.65)
        _text(slide, f"right_meta_item_{idx}", "INDEX", 13.48, y + 0.02, 0.55, 0.08, 5.9, "paper")
        _line(slide, f"right_meta_item_rule_{idx}", 14.04, y + 0.09, 15.12, y + 0.09, "muted", 0.35)


def _footer(slide: Any) -> None:
    _rect(slide, "source_footer_strip", 0, 7.78, SLIDE_W, 1.22, "deep", None)
    _line(slide, "source_footer_top_gold_rule", 0, 7.78, SLIDE_W, 7.78, "gold2", 1.1)
    _shape(slide, "footer_source_cluster_hex", MSO_AUTO_SHAPE_TYPE.HEXAGON, 0.38, 8.16, 0.60, 0.60, "panel", "cyan", 0.65)
    _icon(slide, "footer_source_cluster_icon", "network", 0.50, 8.28, 0.35, "gold2")
    _text(slide, "footer_source_title", "SOURCE", 1.22, 8.25, 0.82, 0.14, 8.4, "white", bold=True)
    _text(slide, "footer_source_meta", "META", 1.22, 8.52, 0.55, 0.11, 7.6, "cyan", bold=True)
    _line(slide, "footer_separator_1", 3.28, 8.06, 3.28, 8.60, "gold2", 0.55)
    _icon(slide, "footer_meta_calendar_icon", "calendar", 3.70, 8.18, 0.33, "gold2")
    _text(slide, "footer_meta_calendar_text", "META", 4.26, 8.30, 0.54, 0.10, 6.8, "cyan", bold=True)
    _line(slide, "footer_meta_calendar_rule", 4.73, 8.40, 5.90, 8.40, "muted", 0.45)
    _line(slide, "footer_separator_2", 6.36, 8.06, 6.36, 8.60, "gold2", 0.55)
    _icon(slide, "footer_meta_user_icon", "user", 6.86, 8.16, 0.36, "gold2")
    _text(slide, "footer_meta_user_text", "META", 7.26, 8.30, 0.50, 0.10, 6.8, "cyan", bold=True)
    _line(slide, "footer_meta_user_rule", 7.74, 8.40, 8.90, 8.40, "muted", 0.45)
    _line(slide, "footer_separator_3", 11.35, 8.06, 11.35, 8.60, "gold2", 0.55)
    _text(slide, "footer_label_title", "FOOTER", 11.68, 8.27, 0.82, 0.12, 8.2, "white", bold=True)
    _text(slide, "footer_label_meta", "META", 11.68, 8.52, 0.50, 0.10, 7.0, "cyan", bold=True)
    _shape(slide, "footer_gold_wedge", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 14.45, 7.78, 1.78, 1.22, "gold2", None)
    _icon(slide, "footer_gold_wedge_icon", "target", 14.78, 8.24, 0.42, "ink")
    _text(slide, "footer_gold_wedge_text", "ICON", 15.25, 8.36, 0.38, 0.10, 6.8, "ink", bold=True)


def _focus_corners(slide: Any, prefix: str, x: float, y: float, size: float, color: str) -> None:
    l = 0.14
    coords = [
        (x, y, x + l, y),
        (x, y, x, y + l),
        (x + size - l, y, x + size, y),
        (x + size, y, x + size, y + l),
        (x, y + size, x + l, y + size),
        (x, y + size - l, x, y + size),
        (x + size - l, y + size, x + size, y + size),
        (x + size, y + size - l, x + size, y + size),
    ]
    for idx, (x1, y1, x2, y2) in enumerate(coords):
        _line(slide, f"{prefix}_{idx}", x1, y1, x2, y2, color, 0.7)


def _icon(slide: Any, name: str, role: str, x: float, y: float, size: float, color: str) -> None:
    if role in {"target", "chart"}:
        _shape(slide, f"{name}_ring", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, None, color, 0.8)
        _line(slide, f"{name}_h", x + size * 0.15, y + size * 0.50, x + size * 0.85, y + size * 0.50, color, 0.7)
        _line(slide, f"{name}_v", x + size * 0.50, y + size * 0.15, x + size * 0.50, y + size * 0.85, color, 0.7)
        if role == "chart":
            for idx, h in enumerate((0.18, 0.30, 0.43)):
                _line(slide, f"{name}_bar_{idx}", x + size * (0.25 + idx * 0.18), y + size * 0.75, x + size * (0.25 + idx * 0.18), y + size * (0.75 - h), color, 0.8)
            _line(slide, f"{name}_trend", x + size * 0.23, y + size * 0.44, x + size * 0.76, y + size * 0.22, color, 0.8)
    elif role == "database":
        for off in (0.0, 0.18, 0.36):
            _shape(slide, f"{name}_db_{off}", MSO_AUTO_SHAPE_TYPE.OVAL, x + size * 0.08, y + size * (0.10 + off), size * 0.84, size * 0.26, None, color, 0.75)
        _line(slide, f"{name}_db_l", x + size * 0.08, y + size * 0.22, x + size * 0.08, y + size * 0.64, color, 0.75)
        _line(slide, f"{name}_db_r", x + size * 0.92, y + size * 0.22, x + size * 0.92, y + size * 0.64, color, 0.75)
    elif role == "network":
        pts = [(0.5, 0.1), (0.15, 0.48), (0.5, 0.88), (0.84, 0.48)]
        for idx, (px, py) in enumerate(pts):
            _shape(slide, f"{name}_node_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, x + size * px - 0.035, y + size * py - 0.035, 0.07, 0.07, color, color, 0.25)
        for idx, ((x1, y1), (x2, y2)) in enumerate(zip(pts, pts[1:] + pts[:1])):
            _line(slide, f"{name}_edge_{idx}", x + size * x1, y + size * y1, x + size * x2, y + size * y2, color, 0.7)
    elif role == "shield":
        _shape(slide, f"{name}_shield", MSO_AUTO_SHAPE_TYPE.PENTAGON, x + size * 0.14, y + size * 0.06, size * 0.72, size * 0.82, None, color, 0.9)
        _line(slide, f"{name}_check_1", x + size * 0.32, y + size * 0.50, x + size * 0.44, y + size * 0.63, color, 0.9)
        _line(slide, f"{name}_check_2", x + size * 0.44, y + size * 0.63, x + size * 0.70, y + size * 0.34, color, 0.9)
    elif role == "document":
        _shape(slide, f"{name}_doc", MSO_AUTO_SHAPE_TYPE.SNIP_1_RECTANGLE, x + size * 0.18, y + size * 0.08, size * 0.62, size * 0.78, None, color, 0.85)
        for idx in range(3):
            _line(slide, f"{name}_doc_line_{idx}", x + size * 0.32, y + size * (0.34 + idx * 0.16), x + size * 0.65, y + size * (0.34 + idx * 0.16), color, 0.65)
    elif role == "book":
        _shape(slide, f"{name}_left_page", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y + size * 0.10, size * 0.45, size * 0.58, None, color, 0.85)
        _shape(slide, f"{name}_right_page", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x + size * 0.45, y + size * 0.10, size * 0.45, size * 0.58, None, color, 0.85)
        _line(slide, f"{name}_spine", x + size * 0.45, y + size * 0.14, x + size * 0.45, y + size * 0.78, color, 0.75)
    elif role == "calendar":
        _shape(slide, f"{name}_cal", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, size, size * 0.78, None, color, 0.75)
        _line(slide, f"{name}_cal_rule", x, y + size * 0.22, x + size, y + size * 0.22, color, 0.75)
        for idx in range(2):
            _line(slide, f"{name}_cal_tick_{idx}", x + size * (0.28 + idx * 0.34), y - size * 0.08, x + size * (0.28 + idx * 0.34), y + size * 0.10, color, 0.75)
    elif role == "user":
        _shape(slide, f"{name}_head", MSO_AUTO_SHAPE_TYPE.OVAL, x + size * 0.28, y, size * 0.44, size * 0.44, None, color, 0.75)
        _shape(slide, f"{name}_body", MSO_AUTO_SHAPE_TYPE.ARC, x + size * 0.12, y + size * 0.45, size * 0.76, size * 0.45, None, color, 0.75)
    else:
        _shape(slide, f"{name}_icon", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, None, color, 0.75)


def _rect(slide: Any, name: str, x: float, y: float, w: float, h: float, fill: str | None, line: str | None, line_width: float = 0.6) -> Any:
    return _shape(slide, name, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill, line, line_width)


def _shape(slide: Any, name: str, shape_type: Any, x: float, y: float, w: float, h: float, fill: str | None, line: str | None, line_width: float = 0.6) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _line(slide: Any, name: str, x1: float, y1: float, x2: float, y2: float, color: str, width: float) -> Any:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _text(slide: Any, name: str, text: str, x: float, y: float, w: float, h: float, size: float, color: str, *, bold: bool = False) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos Display" if size >= 20 else "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _rgb(name_or_hex: str) -> RGBColor:
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
