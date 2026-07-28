"""Insert curated v7.1 SVG icons into PPTX candidates."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.image import ImagePart
from pptx.util import Inches


SVG_CONTENT_TYPE = "image/svg+xml"
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
RENDER_DPI = 144


def insert_icon_v7_1_svg_media(source_pptx: Path, output_pptx: Path, archetype: str, resolved_rows: list[dict[str, Any]], *, slide_index: int = 0) -> dict[str, Any]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(source_pptx)
    slide = prs.slides[slide_index]
    usage_rows: list[dict[str, Any]] = []
    for idx, row in enumerate(resolved_rows):
        x, y, size = _icon_position(idx, len(resolved_rows))
        _add_badge(slide, x, y, size)
        object_id = _add_svg(slide, Path(row["themed_svg_path"]), x + 0.035, y + 0.035, size - 0.07, size - 0.07, row["semantic_role"])
        bbox_px = _bbox_px(x + 0.035, y + 0.035, size - 0.07, size - 0.07, margin_px=7)
        usage_rows.append(
            {
                **row,
                "slide_id": f"{archetype}_candidate",
                "shape_id": object_id,
                "object_name": f"SVG Icon {row['semantic_role']}",
                "bbox": {"x": x + 0.035, "y": y + 0.035, "w": size - 0.07, "h": size - 0.07},
                "bbox_px": bbox_px,
                "visible_in_render": None,
                "raster_fallback": False,
            }
        )
    prs.save(output_pptx)
    return {
        "schema_name": "icon_v7_1_usage_ledger",
        "status": "passed" if output_pptx.exists() else "failed",
        "archetype_id": archetype,
        "pptx_path": output_pptx.as_posix(),
        "icon_v7_1_usage_count": len(usage_rows),
        "true_svg_media_insertion_count": len(usage_rows),
        "native_vector_conversion_count": 0,
        "raster_semantic_icon_count": 0,
        "rows": usage_rows,
    }


def _icon_position(idx: int, count: int) -> tuple[float, float, float]:
    size = 0.24
    gap = 0.08
    start_x = max(7.9, 12.75 - count * (size + gap))
    return start_x + idx * (size + gap), 0.20, size


def _add_badge(slide: Any, x: float, y: float, size: float) -> None:
    badge = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(size), Inches(size))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
    badge.line.color.rgb = RGBColor(203, 213, 225)


def _add_svg(slide: Any, svg_path: Path, x: float, y: float, w: float, h: float, role: str) -> int:
    image_part = ImagePart(slide.part.package.next_image_partname("svg"), SVG_CONTENT_TYPE, slide.part.package, svg_path.read_bytes(), svg_path.name)
    r_id = slide.part.relate_to(image_part, RT.IMAGE)
    shape_id = slide.shapes._next_shape_id
    safe_role = role.replace("&", "and").replace("<", "").replace(">", "")
    slide.shapes._grpSp.add_pic(shape_id, f"SVG Icon {safe_role}", f"svg-icon:{safe_role}:e03_5", r_id, Inches(x), Inches(y), Inches(w), Inches(h))
    slide.shapes._recalculate_extents()
    return shape_id


def _bbox_px(x: float, y: float, w: float, h: float, *, margin_px: int) -> list[int]:
    return [
        max(0, round(x * RENDER_DPI) - margin_px),
        max(0, round(y * RENDER_DPI) - margin_px),
        min(round(SLIDE_W_IN * RENDER_DPI), round((x + w) * RENDER_DPI) + margin_px),
        min(round(SLIDE_H_IN * RENDER_DPI), round((y + h) * RENDER_DPI) + margin_px),
    ]
