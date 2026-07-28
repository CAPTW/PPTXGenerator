"""Template usability contract generation and QA."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_GOLDEN_PPTX = Path("outputs/golden_template_masters.pptx")
DEFAULT_TEMPLATE_SPEC = Path("outputs/editable_template_spec.final.json")
DEFAULT_CONTRACTS_DIR = Path("outputs/template_contracts")
DEFAULT_FINAL_DECK = Path("outputs/final_deck_large_premium.pptx")
DEFAULT_FINAL_DECK_ASSEMBLY_PLAN = Path("outputs/deck_assembly_plan_large_premium.json")
DEFAULT_JSON_REPORT = Path("outputs/template_usability_report.json")
DEFAULT_MD_REPORT = Path("outputs/template_usability_report.md")

EMU_PER_INCH = 914400

RAW_SLOT_LABELS = (
    "layout-board-",
    "slot_id",
    "component_id",
    "title | text | title_block",
    "cards",
    "timeline_steps",
    "source_notes",
)

DEBUG_TEXT = (
    "SAFE MARGINS",
    "ARCHETYPE:",
    "DENSITY:",
    "FOOTER SYSTEM",
    "IMAGE FRAME ONLY",
    "PREVIEW WARNINGS",
)

FIXTURE_PHRASES = (
    "Hierarchy cue",
    "Editable motif",
    "Decision-ready placeholder",
    "Master Fixture",
    "Template Spec",
    "A canonical master",
    "Creative Academic Systems",
    "Primary cue",
    "Supporting cue",
    "Action cue",
)

REPORTED_CATEGORIES = (
    "title capacity violations",
    "body capacity violations",
    "card capacity violations",
    "table/chart capacity violations",
    "reading order risks",
    "hierarchy risks",
    "fixture phrase hits",
    "raw slot label hits",
    "excessive decoration risk",
    "shape budget risk",
    "slide role mismatch",
)

ROLE_PROFILES: dict[str, dict[str, Any]] = {
    "creative_cover": {
        "presentation_role": "deck opener",
        "capacity": (54, 96, 0, 0, 1, 24, 90, 0, 0, 0, 110),
        "budget": (260, "high", 0.82),
        "required": ["ppt_text", "ppt_shape", "svg_ornament", "photo_frame_image"],
        "overflow": {"compress": True, "split_slide": False, "convert_to_cards": False, "move_to_appendix": False, "move_to_speaker_notes": True},
    },
    "visual_table_of_contents": {
        "presentation_role": "visual navigation overview",
        "capacity": (58, 80, 0, 0, 6, 32, 78, 0, 0, 0, 90),
        "budget": (150, "medium", 0.48),
        "required": ["ppt_text", "ppt_shape", "svg_ornament"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": False, "move_to_speaker_notes": False},
    },
    "section_divider": {
        "presentation_role": "section divider",
        "capacity": (42, 84, 1, 72, 0, 0, 0, 0, 0, 0, 88),
        "budget": (180, "high", 0.62),
        "required": ["ppt_text", "ppt_shape", "svg_ornament", "photo_frame_image"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": False, "move_to_appendix": False, "move_to_speaker_notes": True},
    },
    "research_overview": {
        "presentation_role": "evidence overview",
        "capacity": (70, 100, 4, 92, 5, 30, 86, 0, 0, 0, 120),
        "budget": (135, "medium", 0.44),
        "required": ["ppt_text", "ppt_shape"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "problem_statement": {
        "presentation_role": "problem framing",
        "capacity": (66, 96, 3, 86, 3, 28, 76, 0, 0, 0, 120),
        "budget": (120, "medium", 0.42),
        "required": ["ppt_text", "ppt_shape"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": False, "move_to_speaker_notes": True},
    },
    "research_gap": {
        "presentation_role": "gap synthesis",
        "capacity": (66, 96, 4, 84, 4, 28, 74, 0, 0, 0, 110),
        "budget": (130, "medium", 0.44),
        "required": ["ppt_text", "ppt_shape", "svg_ornament"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "literature_map": {
        "presentation_role": "literature relationship map",
        "capacity": (68, 96, 4, 78, 4, 28, 72, 0, 0, 0, 110),
        "budget": (150, "medium", 0.46),
        "required": ["ppt_text", "ppt_shape", "svg_ornament"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "methodology_framework": {
        "presentation_role": "method framework",
        "capacity": (68, 100, 4, 82, 4, 30, 74, 0, 0, 0, 120),
        "budget": (150, "medium", 0.46),
        "required": ["ppt_text", "ppt_shape", "svg_ornament"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "technical_flow_chart": {
        "presentation_role": "technical flow explanation",
        "capacity": (66, 90, 4, 72, 5, 24, 60, 0, 0, 30, 100),
        "budget": (155, "medium", 0.44),
        "required": ["ppt_text", "ppt_shape", "svg_ornament"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "work_support_sequence": {
        "presentation_role": "step sequence",
        "capacity": (66, 90, 5, 72, 5, 24, 64, 0, 0, 28, 96),
        "budget": (140, "medium", 0.42),
        "required": ["ppt_text", "ppt_shape"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "photo_caption_grid": {
        "presentation_role": "photo evidence grid",
        "capacity": (64, 88, 3, 72, 4, 24, 68, 0, 0, 0, 96),
        "budget": (130, "medium", 0.52),
        "required": ["ppt_text", "ppt_shape", "photo_frame_image"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": False, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "comparison_matrix": {
        "presentation_role": "comparison matrix",
        "capacity": (64, 90, 3, 68, 0, 0, 0, 6, 5, 26, 96),
        "budget": (125, "low", 0.38),
        "required": ["ppt_text", "ppt_shape", "ppt_table"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": False, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "concept_relationship_venn": {
        "presentation_role": "concept relationship diagram",
        "capacity": (66, 90, 3, 72, 3, 24, 68, 0, 0, 0, 96),
        "budget": (140, "medium", 0.44),
        "required": ["ppt_text", "ppt_shape", "svg_ornament"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "three_level_explanation": {
        "presentation_role": "three-level explanation",
        "capacity": (66, 90, 3, 72, 3, 26, 72, 0, 0, 0, 104),
        "budget": (130, "medium", 0.42),
        "required": ["ppt_text", "ppt_shape"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "circular_process": {
        "presentation_role": "cyclical process",
        "capacity": (64, 90, 4, 68, 5, 22, 58, 0, 0, 24, 96),
        "budget": (140, "medium", 0.42),
        "required": ["ppt_text", "ppt_shape", "svg_ornament"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "kpi_donut_chart": {
        "presentation_role": "KPI dashboard",
        "capacity": (62, 86, 3, 64, 4, 22, 56, 0, 0, 22, 86),
        "budget": (130, "low", 0.40),
        "required": ["ppt_text", "ppt_shape", "ppt_chart"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": False, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "timeline_roadmap": {
        "presentation_role": "timeline roadmap",
        "capacity": (64, 88, 5, 68, 5, 22, 60, 0, 0, 24, 96),
        "budget": (135, "medium", 0.42),
        "required": ["ppt_text", "ppt_shape"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
    "data_table_appendix": {
        "presentation_role": "appendix data table",
        "capacity": (62, 84, 2, 62, 0, 0, 0, 8, 6, 20, 86),
        "budget": (115, "low", 0.34),
        "required": ["ppt_text", "ppt_shape", "ppt_table"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": False, "move_to_appendix": True, "move_to_speaker_notes": True},
    },
}


def build_template_usability_report_from_files(
    *,
    golden_pptx_path: str | Path = DEFAULT_GOLDEN_PPTX,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    contracts_dir: str | Path = DEFAULT_CONTRACTS_DIR,
    final_deck_path: str | Path = DEFAULT_FINAL_DECK,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
) -> dict[str, Any]:
    report = build_template_usability_report(
        golden_pptx_path=golden_pptx_path,
        template_spec_path=template_spec_path,
        contracts_dir=contracts_dir,
        final_deck_path=final_deck_path,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return report


def build_template_usability_report(
    *,
    golden_pptx_path: str | Path,
    template_spec_path: str | Path,
    contracts_dir: str | Path,
    final_deck_path: str | Path,
) -> dict[str, Any]:
    spec = _load_json(template_spec_path)
    contracts = generate_template_contracts(spec, Path(contracts_dir))
    golden_report = _inspect_deck(Path(golden_pptx_path), spec, contracts, scope="golden_template_masters", map_by_spec=True)
    final_slide_archetypes = _final_deck_archetypes(spec, DEFAULT_FINAL_DECK_ASSEMBLY_PLAN)
    final_report = (
        _inspect_deck(Path(final_deck_path), spec, contracts, scope="final_deck_large_premium", map_by_spec=False, slide_archetypes=final_slide_archetypes)
        if Path(final_deck_path).exists()
        else {"status": "skipped", "pptx_path": _display_path(Path(final_deck_path)), "slides": [], "findings": []}
    )
    findings = golden_report["findings"] + final_report["findings"]
    severe = sum(1 for finding in findings if finding["severity"] == "severe")
    warning = sum(1 for finding in findings if finding["severity"] == "warning")
    category_summary: dict[str, int] = {category: 0 for category in REPORTED_CATEGORIES}
    for finding in findings:
        category = str(finding.get("category") or "unknown")
        category_summary[category] = category_summary.get(category, 0) + 1
    status = "failed" if severe else "issues_reported" if warning else "passed"
    return {
        "schema_name": "template_usability_report",
        "schema_version": "1.0",
        "status": status,
        "golden_pptx_path": _display_path(Path(golden_pptx_path)),
        "template_contracts_dir": _display_path(Path(contracts_dir)),
        "template_spec_path": _display_path(Path(template_spec_path)),
        "final_deck_path": _display_path(Path(final_deck_path)) if Path(final_deck_path).exists() else None,
        "contract_count": len(contracts),
        "findings_summary": {"total": len(findings), "severe": severe, "warning": warning},
        "category_summary": category_summary,
        "golden_template_masters": golden_report,
        "final_deck_large_premium": final_report,
        "slides": golden_report["slides"],
        "findings": findings,
    }


def generate_template_contracts(spec: dict[str, Any], contracts_dir: Path) -> dict[str, dict[str, Any]]:
    contracts_dir.mkdir(parents=True, exist_ok=True)
    contracts: dict[str, dict[str, Any]] = {}
    for layout in spec.get("layouts") or []:
        if not isinstance(layout, dict):
            continue
        archetype_id = str(layout.get("archetype_id") or "")
        if not archetype_id:
            continue
        contract = _contract_for_layout(layout)
        _validate_contract(contract)
        path = contracts_dir / f"{archetype_id}.contract.json"
        path.write_text(json.dumps(contract, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
        contracts[archetype_id] = contract
    return contracts


def _contract_for_layout(layout: dict[str, Any]) -> dict[str, Any]:
    archetype_id = str(layout.get("archetype_id") or "")
    profile = ROLE_PROFILES.get(archetype_id, _generic_profile())
    capacity = profile["capacity"]
    budget = profile["budget"]
    slot_order = _slot_reading_order(layout)
    motifs = _motifs_for_layout(layout)
    return {
        "schema_name": "template_usability_contract",
        "schema_version": "1.0",
        "archetype_id": archetype_id,
        "layout_id": layout.get("layout_id"),
        "layout_family_id": layout.get("layout_family_id"),
        "presentation_role": profile["presentation_role"],
        "reading_order": slot_order,
        "content_capacity": {
            "title_max_chars": capacity[0],
            "subtitle_max_chars": capacity[1],
            "body_bullet_count": capacity[2],
            "body_bullet_max_chars": capacity[3],
            "card_count": capacity[4],
            "card_title_max_chars": capacity[5],
            "card_body_max_chars": capacity[6],
            "table_max_rows": capacity[7],
            "table_max_columns": capacity[8],
            "chart_label_max_chars": capacity[9],
            "takeaway_max_chars": capacity[10],
        },
        "required_visual_motifs": motifs["required"],
        "optional_visual_motifs": motifs["optional"],
        "decorative_budget": {
            "max_shape_count_target": budget[0],
            "max_ornament_density": budget[1],
            "max_background_coverage": budget[2],
        },
        "required_editable_objects": profile["required"],
        "forbidden_patterns": {
            "raw_slot_labels": list(RAW_SLOT_LABELS),
            "debug_text": list(DEBUG_TEXT),
            "long_bullet_dump": True,
            "fixture_phrases": list(FIXTURE_PHRASES),
            "visual_clutter_over_text": True,
        },
        "overflow_policy": profile["overflow"],
    }


def _inspect_deck(
    pptx_path: Path,
    spec: dict[str, Any],
    contracts: dict[str, dict[str, Any]],
    *,
    scope: str,
    map_by_spec: bool,
    slide_archetypes: list[str] | None = None,
) -> dict[str, Any]:
    if not pptx_path.exists():
        return {"status": "skipped", "pptx_path": _display_path(pptx_path), "slides": [], "findings": []}
    deck = Presentation(pptx_path)
    layouts = [layout for layout in spec.get("layouts") or [] if isinstance(layout, dict)]
    slides: list[dict[str, Any]] = []
    findings: list[dict[str, Any]] = []
    for index, slide in enumerate(deck.slides, start=1):
        if map_by_spec and index - 1 < len(layouts):
            archetype_id = str(layouts[index - 1].get("archetype_id") or "")
        elif slide_archetypes and index - 1 < len(slide_archetypes):
            archetype_id = slide_archetypes[index - 1]
        else:
            archetype_id = _infer_archetype(slide)
        contract = contracts.get(archetype_id) or _contract_for_layout({"archetype_id": archetype_id or "unknown", "slots": []})
        stats = _slide_stats(slide, index, archetype_id, scope)
        slide_findings = _slide_findings(stats, contract, scope)
        findings.extend(slide_findings)
        slides.append(
            {
                "slide_number": index,
                "scope": scope,
                "archetype_id": archetype_id,
                "presentation_role": contract.get("presentation_role"),
                "usable_as_presentation_template": not any(f["severity"] == "severe" for f in slide_findings),
                "finding_count": len(slide_findings),
                "metrics": stats["metrics"],
                "finding_codes": [finding["code"] for finding in slide_findings],
            }
        )
    severe = sum(1 for finding in findings if finding["severity"] == "severe")
    warning = sum(1 for finding in findings if finding["severity"] == "warning")
    return {
        "status": "failed" if severe else "issues_reported" if warning else "passed",
        "pptx_path": _display_path(pptx_path),
        "slide_count": len(slides),
        "findings_summary": {"total": len(findings), "severe": severe, "warning": warning},
        "slides": slides,
        "findings": findings,
    }


def _slide_findings(stats: dict[str, Any], contract: dict[str, Any], scope: str) -> list[dict[str, Any]]:
    findings: list[dict[str, Any]] = []
    cap = contract["content_capacity"]
    budget = contract["decorative_budget"]
    text = stats["text_joined"]
    metrics = stats["metrics"]
    slide_number = stats["slide_number"]
    archetype_id = contract["archetype_id"]

    if metrics["title_chars"] > cap["title_max_chars"]:
        findings.append(_finding("TITLE_CAPACITY_VIOLATION", "warning", "Title text exceeds the contract capacity.", "title capacity violations", slide_number, archetype_id, scope, {"title_chars": metrics["title_chars"], "limit": cap["title_max_chars"]}))
    if metrics["body_bullet_count"] > cap["body_bullet_count"] and cap["body_bullet_count"] > 0:
        findings.append(_finding("BODY_BULLET_COUNT_VIOLATION", "warning", "Body bullet count exceeds the contract capacity.", "body capacity violations", slide_number, archetype_id, scope, {"body_bullet_count": metrics["body_bullet_count"], "limit": cap["body_bullet_count"]}))
    if metrics["max_body_line_chars"] > cap["body_bullet_max_chars"] and cap["body_bullet_max_chars"] > 0:
        findings.append(_finding("BODY_LINE_LENGTH_VIOLATION", "warning", "A body line is too long for the layout contract.", "body capacity violations", slide_number, archetype_id, scope, {"max_body_line_chars": metrics["max_body_line_chars"], "limit": cap["body_bullet_max_chars"]}))
    if metrics["estimated_card_count"] > cap["card_count"] and cap["card_count"] > 0:
        findings.append(_finding("CARD_CAPACITY_VIOLATION", "warning", "Card count exceeds the contract capacity.", "card capacity violations", slide_number, archetype_id, scope, {"estimated_card_count": metrics["estimated_card_count"], "limit": cap["card_count"]}))
    if metrics["max_table_rows"] > cap["table_max_rows"] and cap["table_max_rows"] > 0:
        findings.append(_finding("TABLE_ROW_CAPACITY_VIOLATION", "warning", "Table row count exceeds the contract capacity.", "table/chart capacity violations", slide_number, archetype_id, scope, {"max_table_rows": metrics["max_table_rows"], "limit": cap["table_max_rows"]}))
    if metrics["max_table_columns"] > cap["table_max_columns"] and cap["table_max_columns"] > 0:
        findings.append(_finding("TABLE_COLUMN_CAPACITY_VIOLATION", "warning", "Table column count exceeds the contract capacity.", "table/chart capacity violations", slide_number, archetype_id, scope, {"max_table_columns": metrics["max_table_columns"], "limit": cap["table_max_columns"]}))
    if metrics["reading_order_risk"]:
        findings.append(_finding("READING_ORDER_RISK", "warning", "Geometric reading order does not clearly lead with title/section context before dense content.", "reading order risks", slide_number, archetype_id, scope, metrics["reading_order_details"]))
    if metrics["hierarchy_risk"]:
        findings.append(_finding("HIERARCHY_RISK", "warning", "Text hierarchy may be too weak or ambiguous for source-backed content.", "hierarchy risks", slide_number, archetype_id, scope, metrics["hierarchy_details"]))
    for phrase in _hits(text, contract["forbidden_patterns"]["fixture_phrases"]):
        findings.append(_finding("FIXTURE_PHRASE_HIT", "warning", "Fixture phrase is visible; master may read as a demo instead of a reusable template.", "fixture phrase hits", slide_number, archetype_id, scope, {"phrase": phrase}))
    for phrase in _hits(text, contract["forbidden_patterns"]["raw_slot_labels"]):
        findings.append(_finding("RAW_SLOT_LABEL_HIT", "severe", "Raw slot or semantic key text is visible.", "raw slot label hits", slide_number, archetype_id, scope, {"phrase": phrase}))
    for phrase in _hits(text, contract["forbidden_patterns"]["debug_text"]):
        findings.append(_finding("DEBUG_TEXT_HIT", "severe", "Debug or inspection preview text is visible.", "raw slot label hits", slide_number, archetype_id, scope, {"phrase": phrase}))
    if metrics["shape_count"] > budget["max_shape_count_target"]:
        findings.append(_finding("SHAPE_BUDGET_RISK", "warning", "Shape count exceeds the usability contract budget.", "shape budget risk", slide_number, archetype_id, scope, {"shape_count": metrics["shape_count"], "limit": budget["max_shape_count_target"]}))
    if metrics["decoration_to_text_ratio"] > 9.0 and metrics["text_shape_count"] > 0:
        findings.append(_finding("EXCESSIVE_DECORATION_RISK", "warning", "Decorative object volume may compete with readable text zones.", "excessive decoration risk", slide_number, archetype_id, scope, {"decoration_to_text_ratio": metrics["decoration_to_text_ratio"]}))
    if _role_mismatch(metrics, contract):
        findings.append(_finding("SLIDE_ROLE_MISMATCH", "warning", "Rendered objects do not fully support the declared presentation role.", "slide role mismatch", slide_number, archetype_id, scope, {"presentation_role": contract.get("presentation_role"), "required_editable_objects": contract.get("required_editable_objects")}))
    return findings


def _slide_stats(slide: Any, slide_number: int, archetype_id: str, scope: str) -> dict[str, Any]:
    text_items: list[dict[str, Any]] = []
    shape_count = 0
    text_shape_count = 0
    line_count = 0
    table_count = 0
    chart_count = 0
    picture_count = 0
    max_table_rows = 0
    max_table_columns = 0
    for shape in slide.shapes:
        shape_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            line_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
        if getattr(shape, "has_chart", False):
            chart_count += 1
        if getattr(shape, "has_table", False):
            table_count += 1
            max_table_rows = max(max_table_rows, len(shape.table.rows))
            max_table_columns = max(max_table_columns, len(shape.table.columns))
        if getattr(shape, "has_text_frame", False):
            text = str(shape.text or "").strip()
            if text:
                text_shape_count += 1
                text_items.append(
                    {
                        "text": text,
                        "x": shape.left / EMU_PER_INCH,
                        "y": shape.top / EMU_PER_INCH,
                        "w": shape.width / EMU_PER_INCH,
                        "h": shape.height / EMU_PER_INCH,
                        "font_pt": _max_font_size(shape),
                    }
                )
    text_joined = "\n".join(item["text"] for item in text_items)
    title = _title_candidate(text_items)
    body_lines = _body_lines(text_items, title)
    ordered = sorted(text_items, key=lambda item: (item["y"], item["x"]))
    reading_order_risk, reading_details = _reading_order_risk(ordered, title)
    hierarchy_risk, hierarchy_details = _hierarchy_risk(text_items, title)
    estimated_card_count = sum(1 for item in text_items if 0.7 < item["h"] < 1.7 and len(item["text"]) <= 220)
    decorative_count = max(0, shape_count - text_shape_count - table_count - chart_count - picture_count)
    metrics = {
        "shape_count": shape_count,
        "text_shape_count": text_shape_count,
        "line_count": line_count,
        "table_count": table_count,
        "chart_count": chart_count,
        "picture_count": picture_count,
        "decorative_shape_count": decorative_count,
        "decoration_to_text_ratio": round(decorative_count / max(1, text_shape_count), 3),
        "title_chars": len(title["text"]) if title else 0,
        "body_bullet_count": len(body_lines),
        "max_body_line_chars": max((len(line) for line in body_lines), default=0),
        "estimated_card_count": estimated_card_count,
        "max_table_rows": max_table_rows,
        "max_table_columns": max_table_columns,
        "reading_order_risk": reading_order_risk,
        "reading_order_details": reading_details,
        "hierarchy_risk": hierarchy_risk,
        "hierarchy_details": hierarchy_details,
    }
    return {
        "scope": scope,
        "slide_number": slide_number,
        "archetype_id": archetype_id,
        "text_joined": text_joined,
        "text_items": text_items,
        "metrics": metrics,
    }


def _role_mismatch(metrics: dict[str, Any], contract: dict[str, Any]) -> bool:
    role = str(contract.get("presentation_role") or "")
    required = set(contract.get("required_editable_objects") or [])
    if "ppt_table" in required and metrics["table_count"] == 0:
        return True
    if "ppt_chart" in required and metrics["chart_count"] == 0 and "dashboard" in role.lower():
        return True
    if "photo_frame_image" in required and metrics["picture_count"] == 0 and "photo" in role.lower():
        return True
    if "section divider" in role and metrics["body_bullet_count"] > 3:
        return True
    return False


def _final_deck_archetypes(spec: dict[str, Any], assembly_plan_path: Path) -> list[str]:
    if not assembly_plan_path.exists():
        return []
    try:
        assembly_plan = _load_json(assembly_plan_path)
    except (OSError, json.JSONDecodeError):
        return []
    by_layout_id = {
        str(layout.get("layout_id") or ""): str(layout.get("archetype_id") or "")
        for layout in spec.get("layouts") or []
        if isinstance(layout, dict)
    }
    result: list[str] = []
    for binding in assembly_plan.get("slide_layout_bindings") or []:
        if not isinstance(binding, dict):
            continue
        layout_id = str(binding.get("selected_layout_id") or binding.get("layout_id") or "")
        archetype_id = str(binding.get("archetype_id") or binding.get("slide_type") or by_layout_id.get(layout_id) or "")
        result.append(archetype_id)
    return result


def _title_candidate(text_items: list[dict[str, Any]]) -> dict[str, Any] | None:
    candidates = [item for item in text_items if not _is_footer_text(item["text"]) and len(item["text"]) > 1]
    if not candidates:
        return None
    return sorted(candidates, key=lambda item: (-float(item.get("font_pt") or 0), item["y"], item["x"]))[0]


def _body_lines(text_items: list[dict[str, Any]], title: dict[str, Any] | None) -> list[str]:
    lines: list[str] = []
    title_text = title["text"] if title else None
    for item in text_items:
        if item["text"] == title_text or _is_footer_text(item["text"]):
            continue
        for line in re.split(r"[\n\r]+", item["text"]):
            line = line.strip(" \t-•")
            if len(line) >= 8:
                lines.append(line)
    return lines


def _reading_order_risk(ordered: list[dict[str, Any]], title: dict[str, Any] | None) -> tuple[bool, dict[str, Any]]:
    if not title:
        return True, {"reason": "missing_title_candidate"}
    meaningful = [item for item in ordered if not _is_footer_text(item["text"])]
    title_index = next((index for index, item in enumerate(meaningful) if item is title), len(meaningful))
    if title_index > 2:
        return True, {"reason": "title_not_early_in_geometric_order", "title_index": title_index}
    above_title = [item["text"][:80] for item in meaningful if item["y"] + 0.05 < title["y"] and item is not title]
    if len(above_title) > 3:
        return True, {"reason": "too_many_text_items_above_title", "items": above_title[:5]}
    return False, {"title_index": title_index}


def _hierarchy_risk(text_items: list[dict[str, Any]], title: dict[str, Any] | None) -> tuple[bool, dict[str, Any]]:
    if not title:
        return True, {"reason": "missing_title_candidate"}
    title_size = float(title.get("font_pt") or 0)
    body_sizes = [float(item.get("font_pt") or 0) for item in text_items if item is not title and not _is_footer_text(item["text"])]
    max_body = max(body_sizes, default=0)
    if title_size and max_body and title_size < max_body + 3:
        return True, {"reason": "title_not_visually_dominant", "title_font_pt": title_size, "max_body_font_pt": max_body}
    if len(text_items) > 45 and title_size < 20:
        return True, {"reason": "small_title_in_dense_slide", "text_shape_count": len(text_items), "title_font_pt": title_size}
    return False, {"title_font_pt": title_size, "max_body_font_pt": max_body}


def _max_font_size(shape: Any) -> float:
    sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size:
            sizes.append(paragraph.font.size.pt)
        for run in paragraph.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
    return round(max(sizes, default=0.0), 2)


def _hits(text: str, needles: list[str]) -> list[str]:
    lower = text.lower()
    return [needle for needle in needles if needle.lower() in lower]


def _is_footer_text(text: str) -> bool:
    lower = text.lower()
    return "source" in lower or "footer" in lower or "fixture" in lower or len(text) <= 2


def _infer_archetype(slide: Any) -> str:
    text = "\n".join(str(shape.text or "") for shape in slide.shapes if getattr(shape, "has_text_frame", False)).lower()
    if "appendix" in text or "table" in text:
        return "data_table_appendix"
    if "method" in text:
        return "methodology_framework"
    if "chart" in text or "kpi" in text:
        return "kpi_donut_chart"
    if "timeline" in text:
        return "timeline_roadmap"
    return "research_overview"


def _slot_reading_order(layout: dict[str, Any]) -> list[str]:
    slots = [slot for slot in layout.get("slots") or [] if isinstance(slot, dict)]
    ordered = sorted(slots, key=lambda slot: (float((slot.get("bounds") or {}).get("y") or 0), float((slot.get("bounds") or {}).get("x") or 0)))
    ids = [str(slot.get("slot_id") or "") for slot in ordered if slot.get("slot_id")]
    if "title" in ids:
        ids.insert(0, ids.pop(ids.index("title")))
    if "footer" in ids:
        ids.append(ids.pop(ids.index("footer")))
    return ids or ["title", "content", "footer"]


def _motifs_for_layout(layout: dict[str, Any]) -> dict[str, list[str]]:
    component_ids = {str(slot.get("component_id") or "") for slot in layout.get("slots") or [] if isinstance(slot, dict)}
    required = ["clear title zone", "editable content zone", "persistent footer/source system"]
    optional = sorted(component_ids - {"title_block", "footer_standard"})
    if any("image" in item or "photo" in item for item in component_ids):
        required.append("declared photo frame")
    if any("table" in item for item in component_ids):
        required.append("editable table grid")
    if any("chart" in item for item in component_ids):
        required.append("editable chart module")
    return {"required": required, "optional": optional}


def _validate_contract(contract: dict[str, Any]) -> None:
    required = {
        "archetype_id",
        "presentation_role",
        "reading_order",
        "content_capacity",
        "required_visual_motifs",
        "optional_visual_motifs",
        "decorative_budget",
        "required_editable_objects",
        "forbidden_patterns",
        "overflow_policy",
    }
    missing = sorted(required - set(contract))
    if missing:
        raise ValueError(f"{contract.get('archetype_id', '<unknown>')} contract missing {', '.join(missing)}")


def _generic_profile() -> dict[str, Any]:
    return {
        "presentation_role": "content slide",
        "capacity": (66, 96, 4, 78, 4, 28, 74, 0, 0, 24, 110),
        "budget": (130, "medium", 0.42),
        "required": ["ppt_text", "ppt_shape"],
        "overflow": {"compress": True, "split_slide": True, "convert_to_cards": True, "move_to_appendix": True, "move_to_speaker_notes": True},
    }


def _finding(
    code: str,
    severity: str,
    message: str,
    category: str,
    slide_number: int,
    archetype_id: str,
    scope: str,
    details: dict[str, Any] | None = None,
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "code": code,
        "severity": severity,
        "message": message,
        "category": category,
        "slide_number": slide_number,
        "archetype_id": archetype_id,
        "scope": scope,
    }
    if details is not None:
        payload["details"] = details
    return payload


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Template Usability Report",
        "",
        f"Status: `{report['status']}`",
        f"Contracts: `{report['contract_count']}`",
        f"Findings: `{report['findings_summary']['total']}` total, `{report['findings_summary']['severe']}` severe, `{report['findings_summary']['warning']}` warning",
        "",
        "This QA evaluates whether golden masters can hold real presentation content with readable hierarchy, not just whether they are editable visual objects.",
        "",
        "## Category Summary",
        "",
    ]
    for category, count in sorted(report.get("category_summary", {}).items()):
        lines.append(f"- `{category}`: {count}")
    lines.extend(["", "## Golden Master Slides", "", "| Slide | Archetype | Role | Usable | Findings | Shapes | Text boxes |", "|---:|---|---|---|---:|---:|---:|"])
    for slide in report.get("slides") or []:
        metrics = slide.get("metrics") or {}
        lines.append(
            f"| {slide.get('slide_number')} | `{slide.get('archetype_id')}` | {slide.get('presentation_role')} | `{slide.get('usable_as_presentation_template')}` | {slide.get('finding_count')} | {metrics.get('shape_count')} | {metrics.get('text_shape_count')} |"
        )
    lines.extend(["", "## Findings", ""])
    if report.get("findings"):
        for finding in report["findings"][:80]:
            lines.append(
                f"- Slide {finding.get('slide_number')} `{finding.get('archetype_id')}` `{finding.get('code')}` ({finding.get('severity')}): {finding.get('message')}"
            )
        if len(report["findings"]) > 80:
            lines.append(f"- Additional findings omitted from Markdown: {len(report['findings']) - 80}")
    else:
        lines.append("- None")
    return "\n".join(lines) + "\n"


def _load_json(path: str | Path) -> Any:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _display_path(path: Path) -> str:
    return path.as_posix()


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Evaluate golden master templates against usability contracts.")
    parser.add_argument("--golden-pptx", type=Path, default=DEFAULT_GOLDEN_PPTX)
    parser.add_argument("--template-spec", type=Path, default=DEFAULT_TEMPLATE_SPEC)
    parser.add_argument("--contracts-dir", type=Path, default=DEFAULT_CONTRACTS_DIR)
    parser.add_argument("--final-deck", type=Path, default=DEFAULT_FINAL_DECK)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        report = build_template_usability_report_from_files(
            golden_pptx_path=args.golden_pptx,
            template_spec_path=args.template_spec,
            contracts_dir=args.contracts_dir,
            final_deck_path=args.final_deck,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        print(f"TEMPLATE_USABILITY_QA_FAILED {exc}")
        return 1
    print(f"WROTE {args.json_report}")
    print(f"TEMPLATE_USABILITY_QA {report['status']}")
    return 1 if report["findings_summary"]["severe"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
