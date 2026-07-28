from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.util import Inches

from ..slide_scene import Shape
from .common import append_warning, stable_scene_shape_name
from .style import SceneStyleContext, resolve_fill_style, resolve_stroke_style


_SHAPE_TYPE_MAP: dict[str, MSO_AUTO_SHAPE_TYPE] = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rounded_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "rounded-rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
}


def render_shape(
    slide: Any,
    scene_shape: Shape,
    *,
    style_context: SceneStyleContext,
    used_shape_names: set[str],
    slide_number: int,
    slide_id: str,
    warnings: list[Any],
    trace_prefix: str = "shape",
) -> str | None:
    if scene_shape.shape_type == "line":
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(scene_shape.bounds.x),
            Inches(scene_shape.bounds.y),
            Inches(scene_shape.bounds.right),
            Inches(scene_shape.bounds.bottom),
        )
        shape.name = stable_scene_shape_name(trace_prefix, scene_shape.object_id, used_shape_names)
        _apply_line_style(
            shape.line,
            style_context=style_context,
            stroke=scene_shape.stroke,
            default_hex="#94A3B8",
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=scene_shape.object_id,
        )
        return shape.name

    pptx_shape_type = _SHAPE_TYPE_MAP.get(scene_shape.shape_type)
    if pptx_shape_type is None:
        append_warning(
            warnings,
            code="unsupported_shape_type",
            severity="warning",
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=scene_shape.object_id,
            message=f"Scene Shape type {scene_shape.shape_type!r} is not supported by the scene compile path yet.",
            details={"shape_type": scene_shape.shape_type, "supported_shape_types": sorted(_SHAPE_TYPE_MAP)},
        )
        return None

    shape = slide.shapes.add_shape(
        pptx_shape_type,
        Inches(scene_shape.bounds.x),
        Inches(scene_shape.bounds.y),
        Inches(scene_shape.bounds.width),
        Inches(scene_shape.bounds.height),
    )
    shape.name = stable_scene_shape_name(trace_prefix, scene_shape.object_id, used_shape_names)
    if scene_shape.fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = resolve_fill_style(
            style_context,
            scene_shape.fill,
            "#FFFFFF",
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=scene_shape.object_id,
        ).color_rgb
    else:
        shape.fill.background()
    _apply_line_style(
        shape.line,
        style_context=style_context,
        stroke=scene_shape.stroke,
        default_hex="#94A3B8",
        slide_number=slide_number,
        slide_id=slide_id,
        object_id=scene_shape.object_id,
    )
    return shape.name


def _apply_line_style(
    line: Any,
    *,
    style_context: SceneStyleContext,
    stroke: Any,
    default_hex: str,
    slide_number: int,
    slide_id: str,
    object_id: str,
) -> None:
    if stroke is None:
        line.fill.background()
        return
    line.fill.solid()
    line.fill.fore_color.rgb = resolve_stroke_style(
        style_context,
        stroke,
        default_hex=default_hex,
        slide_number=slide_number,
        slide_id=slide_id,
        object_id=object_id,
    ).color_rgb
