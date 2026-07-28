"""Semantic editability and raster policy review for E05."""

from __future__ import annotations

from pathlib import Path
from typing import Any
from zipfile import ZipFile

from pptx import Presentation


def review_semantic_editability(pptx_path: Path, e04_report: dict[str, Any], micro_ledger: dict[str, Any]) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    editable_text_count = 0
    vector_icon_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                editable_text_count += 1
            name = shape.name or ""
            if name.startswith("icon::"):
                vector_icon_count += 1

    svg_media_count = 0
    raster_media_count = 0
    if pptx_path.exists():
        with ZipFile(pptx_path) as zf:
            for name in zf.namelist():
                lower = name.lower()
                if lower.startswith("ppt/media/") and lower.endswith(".svg"):
                    svg_media_count += 1
                if lower.startswith("ppt/media/") and lower.endswith((".png", ".jpg", ".jpeg")):
                    raster_media_count += 1

    semantic_raster = int(e04_report.get("semantic_raster_violation_count", 0)) + int(micro_ledger.get("semantic_raster_icon_count", 0))
    return {
        "schema_name": "e05_semantic_editability_review",
        "status": "passed" if semantic_raster == 0 and vector_icon_count == micro_ledger.get("final_semantic_icon_count", vector_icon_count) else "failed",
        "editable_text_count": editable_text_count,
        "semantic_icon_count": micro_ledger.get("final_semantic_icon_count", vector_icon_count),
        "vector_icon_shape_count": vector_icon_count,
        "svg_media_count": svg_media_count,
        "native_ppt_chart_count": e04_report.get("native_ppt_chart_count", 0),
        "editable_shape_chart_count": e04_report.get("editable_shape_chart_count", 0),
        "raster_chart_count": e04_report.get("raster_chart_count", 0),
        "native_ppt_table_count": e04_report.get("native_ppt_table_count", 0),
        "editable_shape_grid_table_count": e04_report.get("editable_shape_grid_table_count", 0),
        "raster_table_count": e04_report.get("raster_table_count", 0),
        "raster_media_count": raster_media_count,
        "semantic_raster_violation_count": semantic_raster,
        "full_slide_raster_count": e04_report.get("full_slide_raster_count", 0),
        "screenshot_slide_count": e04_report.get("screenshot_slide_count", 0),
        "unknown_content_bearing_count": e04_report.get("unknown_content_bearing_count", 0),
    }


def review_raster_policy(e04_report: dict[str, Any], editability: dict[str, Any]) -> dict[str, Any]:
    semantic_raster = int(editability.get("semantic_raster_violation_count", 0))
    full_slide = int(editability.get("full_slide_raster_count", 0))
    screenshot = int(editability.get("screenshot_slide_count", 0))
    return {
        "schema_name": "e05_raster_policy_review",
        "status": "passed" if semantic_raster == 0 and full_slide == 0 and screenshot == 0 else "failed",
        "semantic_raster_violation_count": semantic_raster,
        "full_slide_raster_count": full_slide,
        "screenshot_slide_count": screenshot,
        "raster_media_count": editability.get("raster_media_count", e04_report.get("raster_media_count", 0)),
        "allowed_raster_policy": "bounded nonsemantic visual/photo/texture fields only",
    }

