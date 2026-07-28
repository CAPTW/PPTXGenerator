"""Extract object coordinates from the E06 baseline PPTX."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation


EMU_PER_INCH = 914400


def extract_pptx_coordinates(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    slides: list[dict[str, Any]] = []
    media_counts: dict[str, int] = {}
    for slide_number, slide in enumerate(prs.slides, start=1):
        objects: list[dict[str, Any]] = []
        for z_order, shape in enumerate(slide.shapes):
            obj = _shape_to_object(shape, slide_number, z_order, slide_w_emu, slide_h_emu)
            objects.append(obj)
            content_type = obj.get("media", {}).get("content_type")
            if content_type:
                media_counts[content_type] = media_counts.get(content_type, 0) + 1
        slides.append(
            {
                "slide_number": slide_number,
                "slide_id": f"slide-{slide_number:03d}",
                "archetype_id": _archetype_for_slide(slide_number),
                "object_count": len(objects),
                "objects": objects,
            }
        )
    object_count = sum(slide["object_count"] for slide in slides)
    semantic_icon_count = sum(1 for slide in slides for obj in slide["objects"] if obj["object_type"] == "semantic_icon")
    return {
        "schema_name": "pptx_coordinate_extraction_report",
        "status": "passed",
        "pptx_path": pptx_path.as_posix(),
        "slide_count": len(slides),
        "slide_size": {
            "width_emu": slide_w_emu,
            "height_emu": slide_h_emu,
            "width_in": round(slide_w_emu / EMU_PER_INCH, 4),
            "height_in": round(slide_h_emu / EMU_PER_INCH, 4),
        },
        "object_count": object_count,
        "semantic_icon_count": semantic_icon_count,
        "svg_media_count": media_counts.get("image/svg+xml", 0),
        "media_counts_by_content_type": media_counts,
        "slides": slides,
    }


def parse_semantic_icon_name(name: str) -> dict[str, str] | None:
    if not name.startswith("icon::"):
        return None
    parts = name.split("::")
    if len(parts) < 5:
        return None
    return {
        "source_slide_id": parts[1],
        "role": parts[2],
        "slot_type": parts[3],
        "component_id": parts[4],
    }


def _shape_to_object(shape: Any, slide_number: int, z_order: int, slide_w_emu: int, slide_h_emu: int) -> dict[str, Any]:
    name = str(getattr(shape, "name", "") or f"shape_{z_order}")
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    icon = parse_semantic_icon_name(name)
    object_type = _classify_object(shape, name)
    semantic_role = icon["role"] if icon else _semantic_role_from_name(name, object_type)
    component_id = icon["component_id"] if icon else _component_id_from_name(name, slide_number)
    text = _shape_text(shape)
    object_id = f"s{slide_number:03d}_z{z_order:04d}_id{getattr(shape, 'shape_id', z_order)}_{_slug(name)}"
    media = _media_info(shape)
    return {
        "object_id": object_id,
        "slide_id": f"slide-{slide_number:03d}",
        "slide_number": slide_number,
        "archetype_id": _archetype_for_slide(slide_number),
        "shape_id": int(getattr(shape, "shape_id", z_order)),
        "name": name,
        "object_type": object_type,
        "shape_type": str(getattr(shape, "shape_type", "")),
        "semantic_role": semantic_role,
        "component_id": component_id,
        "bbox_norm": _bbox(left, top, width, height, slide_w_emu, slide_h_emu),
        "bbox_in": _bbox(left, top, width, height, EMU_PER_INCH, EMU_PER_INCH),
        "bbox_emu": {"x": left, "y": top, "w": width, "h": height},
        "z_order": z_order,
        "editable": object_type not in {"image_field"} or media.get("content_type") == "image/svg+xml",
        "content_bearing": _content_bearing(object_type, text, name),
        "source_binding_id": _source_binding_id(name, text),
        "citation_binding_id": _citation_binding_id(name, text),
        "constraints": _constraints_for(object_type),
        "text_excerpt": text[:120],
        "media": media,
        "icon_metadata": icon or {},
    }


def _bbox(left: int, top: int, width: int, height: int, denom_w: int, denom_h: int) -> dict[str, float]:
    return {
        "x": round(left / denom_w, 6),
        "y": round(top / denom_h, 6),
        "w": round(width / denom_w, 6),
        "h": round(height / denom_h, 6),
    }


def _classify_object(shape: Any, name: str) -> str:
    lower = name.lower()
    if name.startswith("icon::"):
        return "semantic_icon"
    if name.startswith("icon_bg::"):
        return "icon_background"
    if "source" in lower or "footer" in lower or "citation" in lower:
        return "source_footer"
    if "table" in lower or "matrix" in lower or "register" in lower:
        return "table_region"
    if "chart" in lower or "dashboard" in lower or "kpi" in lower:
        return "chart_region"
    if "card" in lower or "panel" in lower or "module" in lower:
        return "card_region"
    if str(getattr(shape, "shape_type", "")).startswith("PICTURE"):
        return "image_field"
    if str(getattr(shape, "shape_type", "")).startswith("LINE"):
        return "line"
    if _shape_text(shape).strip():
        return "text"
    return "shape"


def _shape_text(shape: Any) -> str:
    if not bool(getattr(shape, "has_text_frame", False)):
        return ""
    try:
        return str(shape.text or "").strip()
    except Exception:
        return ""


def _media_info(shape: Any) -> dict[str, Any]:
    rid = getattr(getattr(shape, "_element", None), "blip_rId", None)
    if not rid:
        return {}
    try:
        part = shape.part.related_part(rid)
        return {
            "relationship_id": rid,
            "content_type": getattr(part, "content_type", None),
            "partname": str(getattr(part, "partname", "")),
        }
    except Exception as exc:  # pragma: no cover - defensive for malformed media relationships
        return {"relationship_id": rid, "error": str(exc)}


def _semantic_role_from_name(name: str, object_type: str) -> str:
    lower = name.lower()
    for token in ("source", "citation", "footer", "table", "chart", "card", "kpi", "title", "body"):
        if token in lower:
            return token
    return object_type


def _component_id_from_name(name: str, slide_number: int) -> str:
    clean = re.sub(r"[^a-zA-Z0-9_]+", "_", name).strip("_").lower()
    return clean or f"slide_{slide_number:03d}_component"


def _source_binding_id(name: str, text: str) -> str | None:
    if "source" in name.lower() or "source" in text.lower():
        return f"source_binding::{_slug(name)}"
    return None


def _citation_binding_id(name: str, text: str) -> str | None:
    if "citation" in name.lower() or re.search(r"\bc\d+\b|\[\d+\]", text.lower()):
        return f"citation_binding::{_slug(name)}"
    return None


def _content_bearing(object_type: str, text: str, name: str) -> bool:
    return object_type in {"semantic_icon", "source_footer", "chart_region", "table_region"} or bool(text.strip()) or name.startswith("icon::")


def _constraints_for(object_type: str) -> dict[str, Any]:
    if object_type == "semantic_icon":
        return {"bbox_diff_norm_max": 0.003, "anchor_required": True, "size_token_required": True}
    if object_type == "source_footer":
        return {"bbox_diff_norm_max": 0.005, "source_footer_required": True}
    return {"bbox_diff_norm_max": 0.005}


def _slug(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9_]+", "_", value).strip("_").lower()[:80]


def _archetype_for_slide(slide_number: int) -> str:
    archetypes = [
        "cover_hero",
        "visual_toc",
        "section_divider",
        "standard_content",
        "evidence_overview",
        "card_grid",
        "methodology_framework",
        "process_flow",
        "comparison_matrix",
        "data_dashboard",
        "table_heavy",
        "timeline_roadmap",
        "decision_record",
        "risk_register",
        "case_study",
        "closing_synthesis",
    ]
    return archetypes[slide_number - 1] if 1 <= slide_number <= len(archetypes) else f"slide_{slide_number:03d}"
