"""Parallel, opt-in SceneDeck to PPTX compiler for low-risk object classes.

This module intentionally supports only ``TextBox``, ``ImageObject``,
``NativeTable``, bounded ``NativeChart``, and low-risk structural support
objects on the scene path. The default Blueprint -> SlideIR -> PPTX path
remains unchanged. Unsupported scene objects are reported deterministically and
skipped rather than silently dropped.
"""

from __future__ import annotations

import hashlib
import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field

from .pptx_renderers.chart import render_native_chart
from .pptx_renderers.callout import render_callout
from .pptx_renderers.common import (
    SceneCompileWarning,
    append_warning,
    warning_sort_key,
)
from .pptx_renderers.divider import render_divider_line
from .pptx_renderers.image import render_image_object
from .pptx_renderers.shape import render_shape
from .pptx_renderers.style import build_scene_style_context, resolve_fill_style, scene_style_warning_counts
from .pptx_renderers.table import render_native_table
from .pptx_renderers.text import render_text_box
from .slide_scene import Callout, DividerLine, ImageObject, NativeChart, NativeTable, SceneDeck, Shape, TextBox


SCENE_PPTX_COMPILE_REPORT_VERSION = "0.1"


class SceneCompileModel(BaseModel):
    model_config = ConfigDict(extra="forbid", populate_by_name=True)


class SceneRenderedObjectSummary(SceneCompileModel):
    object_id: str
    kind: str
    shape_name: str


class ScenePptxSlideSummary(SceneCompileModel):
    slide_number: int
    slide_id: str
    rendered_objects: list[SceneRenderedObjectSummary] = Field(default_factory=list)
    skipped_object_ids: list[str] = Field(default_factory=list)
    warning_count: int = 0


class ScenePptxCompileReport(SceneCompileModel):
    report_version: str = SCENE_PPTX_COMPILE_REPORT_VERSION
    pptx_path: str
    scene_deck_path: str | None = None
    structural_hash: str
    slide_count: int
    rendered_text_object_count: int
    rendered_image_object_count: int
    rendered_native_table_count: int
    rendered_native_chart_count: int
    rendered_shape_object_count: int
    rendered_divider_object_count: int
    rendered_callout_object_count: int
    rendered_background_motif_count: int
    rendered_background_shape_count: int
    rendered_background_divider_count: int
    unsupported_background_motif_count: int
    style_warning_count: int = 0
    unresolved_theme_token_count: int = 0
    unresolved_font_token_count: int = 0
    unresolved_spacing_token_count: int = 0
    fallback_style_count: int = 0
    style_alias_count: int = 0
    deprecated_style_alias_count: int = 0
    ambiguous_style_alias_count: int = 0
    skipped_object_count: int
    supported_object_kinds: list[str] = Field(
        default_factory=lambda: ["text_box", "image", "native_table", "native_chart", "shape", "divider", "callout"]
    )
    warnings: list[SceneCompileWarning] = Field(default_factory=list)
    slides: list[ScenePptxSlideSummary] = Field(default_factory=list)

    def to_stable_payload(self, *, include_paths: bool = True) -> dict[str, Any]:
        return scene_pptx_compile_report_to_stable_payload(self, include_paths=include_paths)

    def to_stable_json(self) -> str:
        return scene_pptx_compile_report_to_stable_json(self)


@dataclass
class ScenePptxCompileOutputs:
    scene_deck: SceneDeck
    report: ScenePptxCompileReport
    pptx_path: Path


def compile_pptx_from_scene_deck(
    scene_deck: SceneDeck,
    output_dir: str | Path,
    *,
    root: str | Path | None = None,
    pptx_name: str = "deck.pptx",
    scene_deck_path: str | Path | None = None,
) -> ScenePptxCompileOutputs:
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    asset_output_dir = output_path / "scene-assets"
    root_path = Path(root).resolve() if root is not None else None
    presentation = Presentation()
    presentation.slide_width = int(scene_deck.slide_width * 914400)
    presentation.slide_height = int(scene_deck.slide_height * 914400)
    blank_layout = presentation.slide_layouts[6]

    warnings: list[SceneCompileWarning] = []
    style_context = build_scene_style_context(scene_deck, warnings)
    slide_summaries: list[ScenePptxSlideSummary] = []
    rendered_text_object_count = 0
    rendered_image_object_count = 0
    rendered_native_table_count = 0
    rendered_native_chart_count = 0
    rendered_shape_object_count = 0
    rendered_divider_object_count = 0
    rendered_callout_object_count = 0
    rendered_background_motif_count = 0
    rendered_background_shape_count = 0
    rendered_background_divider_count = 0
    unsupported_background_motif_count = 0
    skipped_object_count = 0

    for scene_slide in sorted(scene_deck.slides, key=lambda item: item.slide_number):
        slide = presentation.slides.add_slide(blank_layout)
        _apply_background_fill(slide, scene_slide, style_context, warnings)
        rendered_objects: list[SceneRenderedObjectSummary] = []
        skipped_object_ids: list[str] = []
        slide_warning_count_before = len(warnings)
        used_shape_names: set[str] = set()
        background_result = _render_background_motifs(
            slide,
            scene_slide=scene_slide,
            style_context=style_context,
            used_shape_names=used_shape_names,
            warnings=warnings,
        )
        rendered_objects.extend(background_result["rendered_objects"])
        skipped_object_ids.extend(background_result["skipped_object_ids"])
        rendered_background_motif_count += background_result["rendered_background_motif_count"]
        rendered_background_shape_count += background_result["rendered_background_shape_count"]
        rendered_background_divider_count += background_result["rendered_background_divider_count"]
        unsupported_background_motif_count += background_result["unsupported_background_motif_count"]
        skipped_object_count += len(background_result["skipped_object_ids"])
        for scene_object in sorted(
            scene_slide.objects,
            key=lambda item: (getattr(item, "z_order", 0), getattr(item, "reading_order", 10_000), item.object_id),
        ):
            if isinstance(scene_object, TextBox):
                shape_name = render_text_box(
                    slide,
                    scene_object,
                    style_context=style_context,
                    used_shape_names=used_shape_names,
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                    warnings=warnings,
                )
                rendered_objects.append(
                    SceneRenderedObjectSummary(object_id=scene_object.object_id, kind=scene_object.kind, shape_name=shape_name)
                )
                rendered_text_object_count += 1
                continue
            if isinstance(scene_object, ImageObject):
                shape_name = render_image_object(
                    slide,
                    scene_object,
                    root=root_path,
                    asset_output_dir=asset_output_dir / f"slide-{scene_slide.slide_number:03d}",
                    used_shape_names=used_shape_names,
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                    warnings=warnings,
                )
                if shape_name is not None:
                    rendered_objects.append(
                        SceneRenderedObjectSummary(object_id=scene_object.object_id, kind=scene_object.kind, shape_name=shape_name)
                    )
                    rendered_image_object_count += 1
                else:
                    skipped_object_ids.append(scene_object.object_id)
                    skipped_object_count += 1
                continue
            if isinstance(scene_object, NativeTable):
                shape_name = render_native_table(
                    slide,
                    scene_object,
                    style_context=style_context,
                    used_shape_names=used_shape_names,
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                    warnings=warnings,
                )
                rendered_objects.append(
                    SceneRenderedObjectSummary(object_id=scene_object.object_id, kind=scene_object.kind, shape_name=shape_name)
                )
                rendered_native_table_count += 1
                continue
            if isinstance(scene_object, NativeChart):
                shape_name = render_native_chart(
                    slide,
                    scene_object,
                    style_context=style_context,
                    used_shape_names=used_shape_names,
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                    warnings=warnings,
                )
                if shape_name is not None:
                    rendered_objects.append(
                        SceneRenderedObjectSummary(object_id=scene_object.object_id, kind=scene_object.kind, shape_name=shape_name)
                    )
                    rendered_native_chart_count += 1
                else:
                    skipped_object_ids.append(scene_object.object_id)
                    skipped_object_count += 1
                continue
            if isinstance(scene_object, Shape):
                shape_name = render_shape(
                    slide,
                    scene_object,
                    style_context=style_context,
                    used_shape_names=used_shape_names,
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                    warnings=warnings,
                )
                if shape_name is not None:
                    rendered_objects.append(
                        SceneRenderedObjectSummary(object_id=scene_object.object_id, kind=scene_object.kind, shape_name=shape_name)
                    )
                    rendered_shape_object_count += 1
                else:
                    skipped_object_ids.append(scene_object.object_id)
                    skipped_object_count += 1
                continue
            if isinstance(scene_object, DividerLine):
                shape_name = render_divider_line(
                    slide,
                    scene_object,
                    style_context=style_context,
                    used_shape_names=used_shape_names,
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                )
                rendered_objects.append(
                    SceneRenderedObjectSummary(object_id=scene_object.object_id, kind=scene_object.kind, shape_name=shape_name)
                )
                rendered_divider_object_count += 1
                continue
            if isinstance(scene_object, Callout):
                shape_name = render_callout(
                    slide,
                    scene_object,
                    style_context=style_context,
                    used_shape_names=used_shape_names,
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                    warnings=warnings,
                )
                rendered_objects.append(
                    SceneRenderedObjectSummary(object_id=scene_object.object_id, kind=scene_object.kind, shape_name=shape_name)
                )
                rendered_callout_object_count += 1
                continue
            append_warning(
                warnings,
                code="unsupported_scene_object",
                severity="warning",
                slide_number=scene_slide.slide_number,
                slide_id=scene_slide.slide_id,
                object_id=scene_object.object_id,
                message=f"Scene object kind {scene_object.kind!r} is not supported by the scene compile path yet.",
                details={"kind": scene_object.kind},
            )
            skipped_object_ids.append(scene_object.object_id)
            skipped_object_count += 1
        slide_summaries.append(
            ScenePptxSlideSummary(
                slide_number=scene_slide.slide_number,
                slide_id=scene_slide.slide_id,
                rendered_objects=rendered_objects,
                skipped_object_ids=skipped_object_ids,
                warning_count=len(warnings) - slide_warning_count_before,
            )
        )

    pptx_path = (output_path / pptx_name).resolve()
    presentation.save(pptx_path)
    sorted_warnings = sorted(warnings, key=warning_sort_key)
    style_counts = scene_style_warning_counts(sorted_warnings)
    report = ScenePptxCompileReport(
        pptx_path=str(pptx_path),
        scene_deck_path=str(Path(scene_deck_path).resolve()) if scene_deck_path is not None else None,
        structural_hash="",
        slide_count=len(slide_summaries),
        rendered_text_object_count=rendered_text_object_count,
        rendered_image_object_count=rendered_image_object_count,
        rendered_native_table_count=rendered_native_table_count,
        rendered_native_chart_count=rendered_native_chart_count,
        rendered_shape_object_count=rendered_shape_object_count,
        rendered_divider_object_count=rendered_divider_object_count,
        rendered_callout_object_count=rendered_callout_object_count,
        rendered_background_motif_count=rendered_background_motif_count,
        rendered_background_shape_count=rendered_background_shape_count,
        rendered_background_divider_count=rendered_background_divider_count,
        unsupported_background_motif_count=unsupported_background_motif_count,
        style_warning_count=style_counts["style_warning_count"],
        unresolved_theme_token_count=style_counts["unresolved_theme_token_count"],
        unresolved_font_token_count=style_counts["unresolved_font_token_count"],
        unresolved_spacing_token_count=style_counts["unresolved_spacing_token_count"],
        fallback_style_count=style_counts["fallback_style_count"],
        style_alias_count=style_counts["style_alias_count"],
        deprecated_style_alias_count=style_counts["deprecated_style_alias_count"],
        ambiguous_style_alias_count=style_counts["ambiguous_style_alias_count"],
        skipped_object_count=skipped_object_count,
        warnings=sorted_warnings,
        slides=slide_summaries,
    )
    return ScenePptxCompileOutputs(
        scene_deck=scene_deck,
        report=report.model_copy(update={"structural_hash": scene_pptx_compile_report_structural_hash(report)}),
        pptx_path=pptx_path,
    )


def compile_pptx_from_scene_deck_file(
    scene_deck_path: str | Path,
    output_dir: str | Path,
    *,
    root: str | Path | None = None,
    pptx_name: str = "deck.pptx",
) -> ScenePptxCompileOutputs:
    path = Path(scene_deck_path)
    scene_deck = SceneDeck.model_validate_json(path.read_text(encoding="utf-8"))
    resolved_root = Path(root).resolve() if root is not None else path.resolve().parent
    return compile_pptx_from_scene_deck(
        scene_deck,
        output_dir,
        root=resolved_root,
        pptx_name=pptx_name,
        scene_deck_path=path,
    )


def write_scene_pptx_compile_report(report: ScenePptxCompileReport, output_path: str | Path) -> Path:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(scene_pptx_compile_report_to_stable_json(report) + "\n", encoding="utf-8")
    return output


def summarize_scene_pptx_compile(report: ScenePptxCompileReport) -> list[str]:
    lines = [
        "SCENE_PPTX_COMPILE "
        f"slides={report.slide_count} "
        f"text={report.rendered_text_object_count} "
        f"images={report.rendered_image_object_count} "
        f"tables={report.rendered_native_table_count} "
        f"charts={report.rendered_native_chart_count} "
        f"motifs={report.rendered_background_motif_count} "
        f"shapes={report.rendered_shape_object_count} "
        f"dividers={report.rendered_divider_object_count} "
        f"callouts={report.rendered_callout_object_count} "
        f"skipped={report.skipped_object_count} "
        f"style_warnings={report.style_warning_count} "
        f"style_aliases={report.style_alias_count} "
        f"warnings={len(report.warnings)}"
    ]
    for warning in report.warnings:
        lines.append(
            f"WARNING slide={warning.slide_number} code={warning.code} severity={warning.severity} {warning.message}"
        )
    return lines


def scene_pptx_compile_report_to_stable_payload(
    report: ScenePptxCompileReport,
    *,
    include_paths: bool = True,
) -> dict[str, Any]:
    payload = report.model_dump(mode="json", exclude_none=True)
    if not include_paths:
        payload.pop("pptx_path", None)
        payload.pop("scene_deck_path", None)
    return _normalize_for_stable_json(payload)


def scene_pptx_compile_report_to_stable_json(report: ScenePptxCompileReport) -> str:
    return json.dumps(
        scene_pptx_compile_report_to_stable_payload(report),
        sort_keys=True,
        separators=(",", ":"),
        ensure_ascii=True,
    )


def scene_pptx_compile_report_structural_hash(report: ScenePptxCompileReport) -> str:
    payload = scene_pptx_compile_report_to_stable_payload(report, include_paths=False)
    payload.pop("structural_hash", None)
    stable_json = json.dumps(
        payload,
        sort_keys=True,
        separators=(",", ":"),
        ensure_ascii=True,
    )
    return hashlib.sha256(stable_json.encode("utf-8")).hexdigest()


def _normalize_for_stable_json(value: Any) -> Any:
    if isinstance(value, dict):
        return {key: _normalize_for_stable_json(value[key]) for key in sorted(value)}
    if isinstance(value, list):
        return [_normalize_for_stable_json(item) for item in value]
    if isinstance(value, float):
        if not math.isfinite(value):
            raise ValueError("Scene compile report cannot serialize non-finite floats")
        normalized = round(value, 6)
        return 0.0 if normalized == 0 else normalized
    return value


def _apply_background_fill(slide: Any, scene_slide: Any, style_context: Any, warnings: list[SceneCompileWarning]) -> None:
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = resolve_fill_style(
            style_context,
            scene_slide.background.fill,
            "#FFFFFF",
            slide_number=scene_slide.slide_number,
            slide_id=scene_slide.slide_id,
            object_id=None,
        ).color_rgb
    except Exception as exc:
        append_warning(
            warnings,
            code="unsupported_slide_background",
            severity="warning",
            slide_number=scene_slide.slide_number,
            slide_id=scene_slide.slide_id,
            object_id=None,
            message="Slide background fill could not be applied in the scene compile path.",
            details={"error": str(exc)},
        )


def _render_background_motifs(
    slide: Any,
    *,
    scene_slide: Any,
    style_context: Any,
    used_shape_names: set[str],
    warnings: list[SceneCompileWarning],
) -> dict[str, Any]:
    rendered_objects: list[SceneRenderedObjectSummary] = []
    skipped_object_ids: list[str] = []
    rendered_background_shape_count = 0
    rendered_background_divider_count = 0
    unsupported_background_motif_count = 0

    for motif in sorted(
        scene_slide.background.motifs,
        key=lambda item: (getattr(item, "z_order", 0), getattr(item, "reading_order", -1), item.object_id),
    ):
        if isinstance(motif, Shape):
            shape_name = render_shape(
                slide,
                motif,
                style_context=style_context,
                used_shape_names=used_shape_names,
                slide_number=scene_slide.slide_number,
                slide_id=scene_slide.slide_id,
                warnings=warnings,
                trace_prefix="background_shape",
            )
            if shape_name is None:
                append_warning(
                    warnings,
                    code="unsupported_background_motif",
                    severity="warning",
                    slide_number=scene_slide.slide_number,
                    slide_id=scene_slide.slide_id,
                    object_id=motif.object_id,
                    message="Background motif shape could not be rendered on the scene compile path.",
                    details={"kind": motif.kind, "shape_type": motif.shape_type},
                )
                skipped_object_ids.append(motif.object_id)
                unsupported_background_motif_count += 1
                continue
            rendered_objects.append(
                SceneRenderedObjectSummary(
                    object_id=motif.object_id,
                    kind="background_shape",
                    shape_name=shape_name,
                )
            )
            rendered_background_shape_count += 1
            continue
        if isinstance(motif, DividerLine):
            shape_name = render_divider_line(
                slide,
                motif,
                style_context=style_context,
                used_shape_names=used_shape_names,
                slide_number=scene_slide.slide_number,
                slide_id=scene_slide.slide_id,
                trace_prefix="background_divider",
            )
            rendered_objects.append(
                SceneRenderedObjectSummary(
                    object_id=motif.object_id,
                    kind="background_divider",
                    shape_name=shape_name,
                )
            )
            rendered_background_divider_count += 1
            continue
        append_warning(
            warnings,
            code="unsupported_background_motif",
            severity="warning",
            slide_number=scene_slide.slide_number,
            slide_id=scene_slide.slide_id,
            object_id=getattr(motif, "object_id", None),
            message="Background motif kind is not supported by the scene compile path.",
            details={"kind": getattr(motif, "kind", "unknown")},
        )
        skipped_object_ids.append(getattr(motif, "object_id", "background-motif"))
        unsupported_background_motif_count += 1

    return {
        "rendered_objects": rendered_objects,
        "skipped_object_ids": skipped_object_ids,
        "rendered_background_motif_count": rendered_background_shape_count + rendered_background_divider_count,
        "rendered_background_shape_count": rendered_background_shape_count,
        "rendered_background_divider_count": rendered_background_divider_count,
        "unsupported_background_motif_count": unsupported_background_motif_count,
    }
