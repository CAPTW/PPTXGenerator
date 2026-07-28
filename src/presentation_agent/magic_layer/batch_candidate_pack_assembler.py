"""Assemble non-canonical D06.1 editable candidate packs."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from .d06_batch_conversion import validate_candidate_pack_path
from .editable_candidate_compiler import SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN, pptx_inventory


def assemble_candidate_pack(specs: list[dict[str, Any]], output_pptx: Path) -> dict[str, Any]:
    errors = validate_candidate_pack_path(output_pptx)
    if errors:
        raise ValueError(errors)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    for spec in specs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for obj in spec.get("objects") or []:
            _object(slide, obj)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    inv = pptx_inventory(output_pptx)
    return {
        "schema_name": "harness_v3_d06_1_16_reference_editable_candidate_pack_report",
        "status": "passed" if inv["slide_count"] == len(specs) else "failed",
        "pack_path": output_pptx.as_posix(),
        "slide_count": inv["slide_count"],
        "expected_slide_count": len(specs),
        "non_canonical": True,
        "source_bound_deck": False,
        "canva_parity_claimed": False,
        "inventory": inv,
    }


def _add_slide_title(slide: Any, spec: dict[str, Any]) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.25), Inches(5.5), Inches(0.35))
    box.text_frame.text = str(spec.get("reference_id") or "reference")
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(12)


def _add_reference_marker(slide: Any, spec: dict[str, Any]) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(6.95), Inches(4.2), Inches(0.25))
    box.text_frame.text = "D06.1 non-canonical template candidate routing slide"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(6)


def _object(slide: Any, obj: dict[str, Any]) -> None:
    from .editable_candidate_compiler import (
        _add_chart_skeleton,
        _add_connector,
        _add_icon_vector,
        _add_scoped_visual_crop,
        _add_shape,
        _add_table_skeleton,
        _add_text,
    )

    object_type = obj.get("object_type")
    if object_type == "ppt_text":
        _add_text(slide, obj)
    elif object_type == "ppt_connector":
        _add_connector(slide, obj)
    elif object_type == "editable_shape_chart":
        _add_chart_skeleton(slide, obj)
    elif object_type == "editable_shape_grid_table":
        _add_table_skeleton(slide, obj)
    elif object_type == "scoped_visual_field_crop":
        _add_scoped_visual_crop(slide, obj)
    elif object_type in {"svg_vector", "ppt_vector_shape_icon"}:
        _add_icon_vector(slide, obj)
    else:
        _add_shape(slide, obj)
