"""Deterministic production gate for lecture-deck planner/compiler/QA changes."""

from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import time
import traceback
from collections import Counter
from collections.abc import Mapping
from enum import StrEnum
from pathlib import Path
from typing import Any

import fitz
from docx import Document
from pydantic import BaseModel, ConfigDict, Field, field_validator
from pptx import Presentation

from .non_pptx_modules.asset_derivation import derive_assets_from_blueprint, write_asset_derivation_outputs
from .non_pptx_modules.deck_qa import run_deck_qa, write_deck_qa_outputs
from .non_pptx_modules.document_asset_crop import run_document_asset_crop, run_document_crop_review, write_document_crop_outputs
from .non_pptx_modules.gate2_planner import plan_gate2, write_gate2_outputs
from .non_pptx_modules.provider_runtime import LLMBackendProof
from .non_pptx_modules.runtime_config import ProviderSettings, parse_provider_option_items
from .non_pptx_modules.structured_visuals import run_structured_visuals, write_structured_visual_outputs
from .non_pptx_modules.workflow_planner import WorkflowBriefInput, plan_workflow_with_provider
from .non_pptx_modules.state_schemas import (
    Blueprint,
    BriefMaterialType,
    CompileEligibility,
    DeliveryMode,
    DeckMode,
    ProjectMaterial,
    QAReport,
    QAVerdictSummary,
    SlideLedger,
    StageStatus,
    VisualType,
    load_state_file,
    save_state_file,
)
from .pptx_compiler import (
    BuildManifest,
    SlideBuildLinkage,
    compile_pptx,
    load_pptx_compile_file,
    write_pptx_compile_outputs,
)


ROOT = Path(__file__).resolve().parents[2]
DEFAULT_POLICY_PATH = ROOT / "state" / "lecture-production-gate-policy.json"
DEFAULT_OUTPUT_ROOT = ROOT / "outputs" / "production-gate"
FORBIDDEN_FALLBACK_MARKERS = (
    "blueprint slide was missing",
    "not approved layout",
    "no approved",
    "placeholder",
    "reconcile the blueprint",
)
DEFAULT_NATIVE_VISUAL_TYPES = (
    VisualType.PROCESS.value,
    VisualType.TIMELINE.value,
    VisualType.DECISION_PATH.value,
    VisualType.FRAMEWORK.value,
    VisualType.HIERARCHY.value,
    VisualType.INFOGRAPHIC.value,
    VisualType.METRIC_SUMMARY.value,
    VisualType.COMPARISON.value,
)
DEFAULT_NATIVE_LAYOUT_FAMILIES = (
    "process-flow",
    "concept-explainer",
    "definition-theorem",
    "summary",
    "comparison",
    "appendix-reference",
)
NATIVE_VISUAL_MISSING_FINDING_ID_PREFIX = "qa-visual-missing-"
NATIVE_VISUAL_MISSING_CATEGORY = "hierarchy"
NATIVE_VISUAL_MISSING_REQUIRED_TAGS = frozenset({"visual", "compile"})
WARNING_SLIDE_PATTERNS: tuple[re.Pattern[str], ...] = (
    re.compile(r"\bslide\s*#?\s*(\d+)\b", re.IGNORECASE),
    re.compile(r"(?<!\w)\((\d+)\)(?!\w)"),
)
WARNING_CONTEXT_KEYWORDS = ("fallback", "placeholder", "layout", "blueprint", "compile", "warning")
BUILD_WARNING_STRING_SURFACE_CODE = "build-warning-string-surface"


class GateContract(BaseModel):
    model_config = ConfigDict(extra="forbid", str_strip_whitespace=True, validate_assignment=True)


class GateProfile(StrEnum):
    CANDIDATE = "candidate"
    RELEASE = "release"


class MaterialMode(StrEnum):
    AUTO = "auto"
    SOURCE = "source"
    SYNTHETIC = "synthetic"


class GateRuleSeverity(StrEnum):
    ERROR = "error"
    WARNING = "warning"


class GateThresholds(GateContract):
    main_story_min: int = 55
    main_story_max: int = 70
    appendix_min: int = 15
    appendix_max: int = 25
    compile_warnings_max: int = 0
    forbidden_visible_text_max: int = 0
    qa_blocking_max: int = 0


class CandidateAllowlistRule(GateContract):
    category: str
    severity: str
    appendix_only: bool = True


class RegressionTargetSpec(GateContract):
    target_id: str
    output_subdir: str
    fixture_kind: str | None = None
    deck_title: str
    topic: str
    audience: list[str] = Field(default_factory=list)
    purpose: str
    delivery_mode: DeliveryMode
    expected_duration_minutes: int
    current_materials: list[ProjectMaterial] = Field(default_factory=list)
    constraints: list[str] = Field(default_factory=list)
    notes: list[str] = Field(default_factory=list)


class ProductionGatePolicy(GateContract):
    policy_id: str
    policy_version: str = "1.0"
    default_profile: GateProfile = GateProfile.CANDIDATE
    pytest_paths: list[str] = Field(default_factory=list)
    required_artifacts: list[str] = Field(default_factory=list)
    thresholds: GateThresholds = Field(default_factory=GateThresholds)
    forbidden_visible_text_substrings: list[str] = Field(default_factory=list)
    forbidden_fallback_markers: list[str] = Field(default_factory=list)
    native_visual_types: list[str] = Field(default_factory=lambda: list(DEFAULT_NATIVE_VISUAL_TYPES))
    native_layout_families: list[str] = Field(default_factory=lambda: list(DEFAULT_NATIVE_LAYOUT_FAMILIES))
    candidate_allowlist: list[CandidateAllowlistRule] = Field(default_factory=list)
    candidate_max_allowlisted_findings: int = 1
    regression_targets: list[RegressionTargetSpec] = Field(default_factory=list)


class PytestExecutionResult(GateContract):
    command: list[str]
    returncode: int
    passed: bool
    duration_seconds: float
    stdout_path: str
    stderr_path: str


class ResidualFinding(GateContract):
    finding_id: str
    category: str
    severity: str
    slide_number: int | None = None
    deck_mode: str | None = None
    blocking: bool = False


class GateMetrics(GateContract):
    target_id: str
    provider_requested: str | None = None
    provider_used: str | None = None
    model_requested: str | None = None
    model_used: str | None = None
    endpoint_requested: str | None = None
    endpoint_used: str | None = None
    transport_used: str | None = None
    strict_structured_output: bool | None = None
    llm_request_count: int = 0
    llm_request_targets: list[str] = Field(default_factory=list)
    llm_backend_proof_path: str | None = None
    qa_status: str = "unknown"
    qa_compile_eligibility: str | None = None
    qa_warning_reason_codes: list[str] = Field(default_factory=list)
    qa_blocking_reason_codes: list[str] = Field(default_factory=list)
    qa_compatibility_warning_codes: list[str] = Field(default_factory=list)
    main_story_slide_count: int = 0
    appendix_slide_count: int = 0
    blueprint_total_slide_count: int | None = None
    blueprint_main_story_slide_count: int | None = None
    blueprint_appendix_slide_count: int | None = None
    slide_ledger_total_slide_count: int | None = None
    slide_ledger_main_story_slide_count: int | None = None
    slide_ledger_appendix_slide_count: int | None = None
    slide_build_linkage_total_slide_count: int | None = None
    slide_build_linkage_main_story_slide_count: int | None = None
    slide_build_linkage_appendix_slide_count: int | None = None
    compiled_pptx_total_slide_count: int | None = None
    compiled_pptx_main_story_slide_count: int | None = None
    compiled_pptx_appendix_slide_count: int | None = None
    compiled_pptx_deck_mode_derivation: str = "not-attempted"
    compiled_pptx_deck_mode_derivation_reason: str | None = None
    compile_warning_count: int = 0
    forbidden_visible_text_count: int = 0
    pptx_scan_forbidden_visible_text_count: int = 0
    validation_summary_forbidden_visible_text_count: int | None = None
    forbidden_visible_text_artifact_match: bool | None = None
    planned_authoring_metrics: dict[str, Any] = Field(default_factory=dict)
    compiled_deck_metrics: dict[str, Any] = Field(default_factory=dict)
    compiled_deck_issue_slides: dict[str, list[int]] = Field(default_factory=dict)
    compiled_deck_truth_mismatches: list[str] = Field(default_factory=list)
    compiled_deck_chrome_block_count: int = 0
    compiled_deck_repeated_title_stem_count: int = 0
    compiled_deck_title_body_duplication_count: int = 0
    compiled_deck_repeated_archetype_count: int = 0
    compiled_deck_repeated_geometry_count: int = 0
    compiled_deck_repeated_geometry_rate: float = 0.0
    compiled_deck_text_card_overuse_rate: float = 0.0
    compiled_deck_chrome_dominance_rate: float = 0.0
    compiled_deck_appendix_clone_count: int = 0
    compiled_deck_appendix_visual_clone_run_length: int = 0
    compiled_deck_expected_visual_missing_count: int = 0
    compiled_deck_archetype_realization_mismatch_count: int = 0
    compiled_deck_missing_visual_center_count: int = 0
    compiled_deck_repetitive_motion_count: int = 0
    compiled_deck_truth_mismatch_count: int = 0
    qa_blocking_count: int = 0
    appendix_boundary_violation_count: int = 0
    layout_compatibility_failure_count: int = 0
    deterministic_fallback_violation_count: int = 0
    native_visual_misclassification_count: int = 0
    native_visual_classifier_mode: str = "not-triggered"
    compiled_pptx_path_source: str | None = None
    validation_summary_pptx_path_status: str | None = None
    validation_summary_pptx_path_hint: str | None = None
    slide_identity: dict[str, Any] = Field(default_factory=dict)
    residual_findings: list[ResidualFinding] = Field(default_factory=list)
    artifact_errors: list[str] = Field(default_factory=list)
    artifact_paths: dict[str, str] = Field(default_factory=dict)


class NativeVisualClassificationOutcome(GateContract):
    is_misclassification: bool = False
    mode: str | None = None
    schema_error: str | None = None


class GateRuleFinding(GateContract):
    rule_id: str
    passed: bool
    severity: GateRuleSeverity
    message: str
    target_id: str | None = None
    profile_impact: list[GateProfile] = Field(default_factory=list)
    affected_slides: list[int] = Field(default_factory=list)
    metric_values: dict[str, Any] = Field(default_factory=dict)

    @field_validator("profile_impact", mode="before")
    @classmethod
    def _coerce_profiles(cls, value: object) -> object:
        if value is None:
            return []
        return value


class TargetBuildExecution(GateContract):
    target_id: str
    output_dir: str
    artifact_paths: dict[str, str] = Field(default_factory=dict)
    validation_summary_path: str | None = None
    build_error: str | None = None
    provider_requested: str | None = None
    model_requested: str | None = None
    endpoint_requested: str | None = None


class TargetGateResult(GateContract):
    target_id: str
    output_dir: str
    passed: bool
    metrics: GateMetrics
    rule_results: list[GateRuleFinding] = Field(default_factory=list)
    artifact_paths: dict[str, str] = Field(default_factory=dict)


class ProductionGateResult(GateContract):
    profile: GateProfile
    passed: bool
    policy_id: str
    policy_version: str
    material_mode: MaterialMode
    pytest_result: PytestExecutionResult
    targets: list[TargetGateResult] = Field(default_factory=list)
    failed_rules: list[GateRuleFinding] = Field(default_factory=list)
    warnings: list[GateRuleFinding] = Field(default_factory=list)
    artifact_paths: dict[str, str] = Field(default_factory=dict)


def _write_json(path: Path, payload: dict[str, Any]) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def _load_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError(f"expected a top-level JSON object in {path}")
    return payload


def _coerce_qa_verdict_summary(
    value: Any,
    *,
    source_label: str,
    artifact_errors: list[str],
) -> QAVerdictSummary | None:
    if value is None:
        return None
    try:
        if isinstance(value, QAVerdictSummary):
            return value
        return QAVerdictSummary.model_validate(value)
    except Exception as exc:
        artifact_errors.append(f"{source_label} is malformed: {type(exc).__name__}: {exc}")
        return None


def _coerce_compile_eligibility(
    value: Any,
    *,
    source_label: str,
    artifact_errors: list[str],
) -> CompileEligibility | None:
    if value is None:
        return None
    raw = value.value if hasattr(value, "value") else value
    normalized = str(raw).strip()
    try:
        return CompileEligibility(normalized)
    except ValueError:
        artifact_errors.append(f"{source_label} is malformed: {value!r}.")
        return None


def _resolve_qa_policy_summary(
    *,
    qa_report: QAReport | None,
    validation_summary: dict[str, Any] | None,
    artifact_errors: list[str],
) -> QAVerdictSummary | None:
    qa_report_summary = _coerce_qa_verdict_summary(
        getattr(qa_report, "verdict_summary", None),
        source_label="qa-report verdict_summary",
        artifact_errors=artifact_errors,
    )
    validation_summary_payload = validation_summary.get("qa_verdict_summary") if isinstance(validation_summary, Mapping) else None
    validation_summary_summary = _coerce_qa_verdict_summary(
        validation_summary_payload,
        source_label="validation-summary qa_verdict_summary",
        artifact_errors=artifact_errors,
    )
    if qa_report is not None and qa_report_summary is not None and qa_report.qa_status != qa_report_summary.qa_status:
        artifact_errors.append("qa-report qa_status does not match verdict_summary.qa_status.")
    if qa_report_summary is not None and validation_summary_summary is not None:
        if qa_report_summary.model_dump(mode="json") != validation_summary_summary.model_dump(mode="json"):
            artifact_errors.append("validation-summary qa_verdict_summary does not match qa-report verdict_summary.")
    return qa_report_summary or validation_summary_summary


def load_production_gate_policy(path: str | Path) -> ProductionGatePolicy:
    return ProductionGatePolicy.model_validate(_load_json(Path(path)))


def _prepare_output_dir(path: Path) -> None:
    if path.exists():
        try:
            def _retry_remove(function: Any, target: str, _excinfo: object) -> None:
                os.chmod(target, 0o700)
                function(target)

            shutil.rmtree(path, onerror=_retry_remove)
        except Exception as exc:  # pragma: no cover
            raise RuntimeError(f"Could not clean production-gate output directory `{path}`: {exc}") from exc
    path.mkdir(parents=True, exist_ok=True)


def _run_pytest_suite(root: Path, pytest_paths: list[str], output_dir: Path) -> PytestExecutionResult:
    command = [sys.executable, "-m", "pytest", *pytest_paths, "-q"]
    start = time.perf_counter()
    result = subprocess.run(command, cwd=root, capture_output=True, text=True, check=False)
    duration = time.perf_counter() - start
    stdout_path = output_dir / "pytest-stdout.txt"
    stderr_path = output_dir / "pytest-stderr.txt"
    stdout_path.write_text(result.stdout, encoding="utf-8")
    stderr_path.write_text(result.stderr, encoding="utf-8")
    execution = PytestExecutionResult(
        command=command,
        returncode=result.returncode,
        passed=result.returncode == 0,
        duration_seconds=round(duration, 3),
        stdout_path=str(stdout_path),
        stderr_path=str(stderr_path),
    )
    _write_json(output_dir / "pytest-result.json", execution.model_dump(mode="json", exclude_none=True))
    return execution


def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run the deterministic lecture production gate.")
    parser.add_argument("--profile", choices=[profile.value for profile in GateProfile], default=GateProfile.CANDIDATE.value)
    parser.add_argument("--policy", type=Path, default=DEFAULT_POLICY_PATH)
    parser.add_argument("--output-dir", type=Path, default=None)
    parser.add_argument("--material-mode", choices=[mode.value for mode in MaterialMode], default=MaterialMode.AUTO.value)
    parser.add_argument("--repo-root", type=Path, default=ROOT)
    parser.add_argument("--provider", default="local-none")
    parser.add_argument("--model")
    parser.add_argument("--endpoint")
    parser.add_argument("--provider-profile")
    parser.add_argument("--provider-option", action="append")
    return parser.parse_args(argv)


def _provider_settings_from_args(args: argparse.Namespace) -> ProviderSettings:
    return ProviderSettings(
        provider=args.provider,
        model=args.model,
        endpoint=args.endpoint,
        profile=args.provider_profile,
        options=parse_provider_option_items(args.provider_option),
    )


def _optimization_fixture_outline() -> list[tuple[str, list[tuple[str, list[str]]]]]:
    return [
        (
            "1 최적화 개요와 해석적 방법",
            [
                ("1.1 최적화 개요", [
                    "1.1.1 최적화 문제의 의미",
                    "1.1.2 설계 변수와 목적 함수",
                    "1.1.3 제약조건과 허용영역",
                    "1.1.4 의사결정 구조",
                    "1.1.5 지역해와 전역해",
                    "1.1.6 공학 설계에서의 역할",
                    "1.1.7 실무 해석 포인트",
                ]),
                ("1.2 최적화 문제의 정형화 예", [
                    "1.2.1 단순 비용 최소화",
                    "1.2.2 성능 최대화 문제",
                    "1.2.3 제약이 있는 설계 예",
                    "1.2.4 다목적 정식화",
                    "1.2.5 모델 단순화와 변수 스케일링",
                ]),
                ("1.3 최적화 문제의 구조", [
                    "1.3.1 볼록성과 비볼록성",
                    "1.3.2 연속/이산 변수",
                    "1.3.3 선형/비선형 구조",
                    "1.3.4 목적 함수의 민감도",
                    "1.3.5 제약의 활성화",
                    "1.3.6 구조에 따른 해법 선택",
                ]),
                ("1.4 단변수 함수의 최적화", [
                    "1.4.1 도함수 기반 조건",
                    "1.4.2 극값의 판정",
                    "1.4.3 이차 미분의 해석",
                    "1.4.4 경계 조건",
                    "1.4.5 탐색 구간의 설정",
                    "1.4.6 예제 계산",
                    "1.4.7 해석적 한계",
                ]),
                ("1.5 다변수 함수의 최적화", [
                    "1.5.1 gradient의 의미",
                    "1.5.2 Hessian의 역할",
                    "1.5.3 stationary point 분류",
                    "1.5.4 등고선 기반 해석",
                    "1.5.5 제약이 없는 경우",
                    "1.5.6 라그랑주 승수",
                    "1.5.7 KKT 조건의 해석",
                ]),
            ],
        ),
        (
            "2 최적화법",
            [
                ("2.1 최적화법의 분류", [
                    "2.1.1 해석적/수치적 방법",
                    "2.1.2 지역/전역 탐색",
                    "2.1.3 문제 구조와 알고리즘 선택",
                ]),
                ("2.2 구배법", [
                    "2.2.1 steepest descent",
                    "2.2.2 Newton method",
                    "2.2.3 quasi-Newton",
                    "2.2.4 line search",
                    "2.2.5 수렴 특성",
                    "2.2.6 적용 시 주의점",
                ]),
                ("2.3 무작위 탐색법", [
                    "2.3.1 확률적 탐색의 개요",
                    "2.3.2 장단점과 적용 맥락",
                ]),
                ("2.4 유전 알고리즘", [
                    "2.4.1 표현과 초기화",
                    "2.4.2 선택/교차/돌연변이",
                    "2.4.3 적합도 설계",
                    "2.4.4 실무 파라미터 설정",
                ]),
                ("2.5 유전 프로그래밍", [
                    "2.5.1 표현식 진화와 구조 탐색",
                ]),
                ("2.6 진화 알고리즘", [
                    "2.6.1 진화전략",
                    "2.6.2 차분진화와 현대적 변형",
                ]),
                ("2.7 지역탐색과 전역탐색", [
                    "2.7.1 exploitation vs exploration",
                    "2.7.2 초기값 민감도",
                    "2.7.3 하이브리드 전략",
                    "2.7.4 중단 조건",
                    "2.7.5 계산비용",
                    "2.7.6 문제 구조별 선택 기준",
                ]),
            ],
        ),
    ]


def _build_optimization_fixture_docx(docx_path: Path) -> None:
    document = Document()
    document.add_heading("최적화 강의 교재", level=0)
    for top_heading, section_groups in _optimization_fixture_outline():
        document.add_paragraph(top_heading, style="Heading 1")
        document.add_paragraph(
            "이 장은 최적화 문제의 구조를 이해하고, 해석적 조건과 수치 알고리즘을 연결하는 강의용 개요이다."
        )
        for section_heading, subsection_headings in section_groups:
            document.add_paragraph(section_heading, style="Heading 2")
            document.add_paragraph(
                "핵심 개념, 선택 기준, 예제 해석을 함께 정리하여 대학원 수준의 강의 흐름을 만든다."
            )
            for subsection_heading in subsection_headings:
                document.add_paragraph(subsection_heading, style="Heading 3")
                document.add_paragraph(
                    f"{subsection_heading}의 정의, 수학적 표현, 설계적 의미, 해법 선택과의 연결을 강의용으로 요약한다."
                )
    document.save(docx_path)


def _build_optimization_fixture_pdf(pdf_path: Path) -> None:
    pdf = fitz.open()
    page_number = 1
    for top_heading, section_groups in _optimization_fixture_outline():
        page = pdf.new_page(width=842, height=595)
        page.insert_text((48, 56), top_heading, fontsize=26)
        page.insert_text((48, 96), "강의용 source map", fontsize=14)
        y = 138
        for section_heading, _subsections in section_groups:
            page.insert_text((60, y), section_heading, fontsize=16)
            y += 28
        page.insert_text((48, 540), f"Page {page_number}", fontsize=10)
        page_number += 1
        for section_heading, subsection_headings in section_groups:
            page = pdf.new_page(width=842, height=595)
            page.insert_text((48, 56), section_heading, fontsize=24)
            y = 110
            for subsection_heading in subsection_headings:
                page.insert_text((60, y), subsection_heading, fontsize=15)
                page.insert_text((84, y + 18), "정의, 조건, 예제, 해법 선택 기준", fontsize=11)
                y += 54
                if y > 500:
                    break
            page.insert_text((48, 540), f"Page {page_number}", fontsize=10)
            page_number += 1
    pdf.save(pdf_path)
    pdf.close()


def _resolve_target_materials(target: RegressionTargetSpec, target_root: Path, material_mode: MaterialMode) -> list[ProjectMaterial]:
    configured = [material for material in target.current_materials if material.path]
    source_ready = bool(configured) and all(Path(material.path).exists() for material in configured)
    if material_mode == MaterialMode.SOURCE:
        if not source_ready:
            raise FileNotFoundError(f"Configured source materials for target `{target.target_id}` are unavailable.")
        return configured
    if material_mode == MaterialMode.AUTO and source_ready:
        return configured
    if target.fixture_kind != "optimization-lecture-synthetic":
        raise FileNotFoundError(f"Target `{target.target_id}` cannot resolve source materials and has no synthetic fixture configured.")

    inputs_dir = target_root / "fixture-inputs"
    inputs_dir.mkdir(parents=True, exist_ok=True)
    docx_path = inputs_dir / "optimization-lecture-fixture.docx"
    pdf_path = inputs_dir / "optimization-lecture-fixture.pdf"
    _build_optimization_fixture_docx(docx_path)
    _build_optimization_fixture_pdf(pdf_path)
    return [
        ProjectMaterial(label=docx_path.name, material_type=BriefMaterialType.DOCUMENT, path=str(docx_path)),
        ProjectMaterial(label=pdf_path.name, material_type=BriefMaterialType.DOCUMENT, path=str(pdf_path)),
    ]


def _pptx_name_for_target(target: RegressionTargetSpec) -> str:
    return f"{target.target_id}.pptx"


def _artifact_paths_for_target(target_root: Path, target: RegressionTargetSpec) -> dict[str, Path]:
    return {
        "workflow_plan": target_root / "state" / "workflow-plan.json",
        "authoring_preview": target_root / "state" / "authoring-preview.json",
        "blueprint": target_root / "state" / "blueprint.json",
        "design_system": target_root / "state" / "design-system.json",
        "deck_constitution": target_root / "state" / "deck-constitution.json",
        "layout_library": target_root / "state" / "layout-library.json",
        "slide_ledger": target_root / "state" / "slide-ledger.json",
        "asset_requests": target_root / "state" / "asset-requests.json",
        "asset_manifest": target_root / "state" / "asset-manifest.json",
        "viz_spec": target_root / "state" / "viz-spec.json",
        "viz_manifest": target_root / "state" / "viz-manifest.json",
        "build_manifest": target_root / "state" / "build-manifest.json",
        "slide_build_linkage": target_root / "state" / "slide-build-linkage.json",
        "qa_report": target_root / "state" / "qa-report.json",
        "compiled_deck_text": target_root / "state" / "compiled-deck-text.json",
        "compiled_deck_shape_census": target_root / "state" / "compiled-deck-shape-census.json",
        "compiled_deck_authoring_audit": target_root / "state" / "compiled-deck-authoring-audit.json",
        "compiled_deck_thumbnail_strip": target_root / "state" / "compiled-deck-thumbnail-strip.png",
        "compiled_deck_thumbnail_index": target_root / "state" / "compiled-deck-thumbnail-index.json",
        "compiled_deck_visual_review_summary": target_root / "state" / "compiled-deck-visual-review-summary.json",
        "validation_summary": target_root / "state" / "validation-summary.json",
        "llm_backend_proof": target_root / "state" / "llm-backend-proof.json",
        "pptx": target_root / "artifacts" / "pptx" / _pptx_name_for_target(target),
    }


def _normalized_value(value: Any) -> str:
    if value is None:
        return ""
    raw = value.value if hasattr(value, "value") else value
    return str(raw).strip().lower()


def _count_blueprint_slides(blueprint: Blueprint) -> tuple[int, int, int]:
    total = len(blueprint.slides)
    main_story = sum(1 for slide in blueprint.slides if slide.deck_mode == DeckMode.MAIN_STORY)
    appendix = sum(1 for slide in blueprint.slides if slide.deck_mode == DeckMode.APPENDIX)
    return total, main_story, appendix


def _count_slide_ledger_entries(slide_ledger: SlideLedger) -> tuple[int, int, int]:
    total = len(slide_ledger.entries)
    main_story = sum(1 for entry in slide_ledger.entries if entry.deck_mode == DeckMode.MAIN_STORY)
    appendix = sum(1 for entry in slide_ledger.entries if entry.deck_mode == DeckMode.APPENDIX)
    return total, main_story, appendix


def _count_slide_build_linkage_entries(slide_build_linkage: SlideBuildLinkage) -> tuple[int, int, int]:
    total = len(slide_build_linkage.slides)
    main_story = sum(1 for entry in slide_build_linkage.slides if entry.deck_mode == DeckMode.MAIN_STORY)
    appendix = sum(1 for entry in slide_build_linkage.slides if entry.deck_mode == DeckMode.APPENDIX)
    return total, main_story, appendix


def _build_slide_registry(entries: list[Any], *, index_attr: str | None = None) -> dict[str, Any]:
    slide_numbers: list[int] = []
    slide_number_duplicates: set[int] = set()
    slide_modes: dict[int, str] = {}
    seen_slide_numbers: set[int] = set()
    index_values: list[int] = []
    index_duplicates: set[int] = set()
    seen_indices: set[int] = set()
    for entry in entries:
        slide_number = getattr(entry, "slide_number", None)
        if isinstance(slide_number, int):
            slide_numbers.append(slide_number)
            if slide_number in seen_slide_numbers:
                slide_number_duplicates.add(slide_number)
            else:
                seen_slide_numbers.add(slide_number)
            slide_modes.setdefault(slide_number, _normalized_value(getattr(entry, "deck_mode", None)))
        if index_attr is not None:
            index_value = getattr(entry, index_attr, None)
            if isinstance(index_value, int):
                index_values.append(index_value)
                if index_value in seen_indices:
                    index_duplicates.add(index_value)
                else:
                    seen_indices.add(index_value)
    registry = {
        "sequence": slide_numbers,
        "slide_number_set": sorted(set(slide_numbers)),
        "slide_number_duplicates": sorted(slide_number_duplicates),
        "deck_modes": slide_modes,
    }
    if index_attr is not None:
        registry["index_sequence"] = index_values
        registry["index_set"] = sorted(set(index_values))
        registry["index_duplicates"] = sorted(index_duplicates)
    return registry


def _compare_slide_identity(left_label: str, left_registry: dict[str, Any], right_label: str, right_registry: dict[str, Any]) -> dict[str, Any]:
    left_sequence = left_registry.get("comparison_sequence", left_registry["sequence"])
    right_sequence = right_registry.get("comparison_sequence", right_registry["sequence"])
    return {
        "left_label": left_label,
        "right_label": right_label,
        "left_only": sorted(set(left_sequence) - set(right_sequence)),
        "right_only": sorted(set(right_sequence) - set(left_sequence)),
        "left_duplicates": list(left_registry["slide_number_duplicates"]),
        "right_duplicates": list(right_registry["slide_number_duplicates"]),
        "sequence_match": left_sequence == right_sequence,
        "left_sequence": list(left_sequence),
        "right_sequence": list(right_sequence),
    }


def _reconcile_slide_identity(
    *,
    blueprint: Blueprint | None,
    slide_ledger: SlideLedger | None,
    slide_build_linkage: SlideBuildLinkage | None,
    compiled_pptx_total_slide_count: int | None,
) -> dict[str, Any]:
    identity: dict[str, Any] = {}
    if blueprint is not None:
        identity["blueprint"] = _build_slide_registry(list(blueprint.slides))
    if slide_ledger is not None:
        identity["slide_ledger"] = _build_slide_registry(list(slide_ledger.entries))
    if slide_build_linkage is not None:
        linkage_registry = _build_slide_registry(list(slide_build_linkage.slides), index_attr="pptx_index")
        ordered_linkage = sorted(slide_build_linkage.slides, key=lambda entry: entry.pptx_index)
        linkage_registry["sequence_by_pptx_index"] = [entry.slide_number for entry in ordered_linkage]
        linkage_registry["comparison_sequence"] = list(linkage_registry["sequence_by_pptx_index"])
        identity["slide_build_linkage"] = linkage_registry

    pairwise: dict[str, Any] = {}
    if "blueprint" in identity and "slide_ledger" in identity:
        pairwise["blueprint_vs_slide_ledger"] = _compare_slide_identity(
            "Blueprint",
            identity["blueprint"],
            "Slide ledger",
            identity["slide_ledger"],
        )
    if "blueprint" in identity and "slide_build_linkage" in identity:
        pairwise["blueprint_vs_slide_build_linkage"] = _compare_slide_identity(
            "Blueprint",
            identity["blueprint"],
            "Slide-build linkage",
            identity["slide_build_linkage"],
        )
    if "slide_ledger" in identity and "slide_build_linkage" in identity:
        pairwise["slide_ledger_vs_slide_build_linkage"] = _compare_slide_identity(
            "Slide ledger",
            identity["slide_ledger"],
            "Slide-build linkage",
            identity["slide_build_linkage"],
        )
    identity["pairwise"] = pairwise

    linkage = identity.get("slide_build_linkage")
    if linkage is not None:
        index_sequence = list(linkage.get("index_sequence", []))
        expected_index_sequence = list(range(1, len(index_sequence) + 1))
        identity["linkage_pptx_index"] = {
            "sequence": index_sequence,
            "duplicates": list(linkage.get("index_duplicates", [])),
            "expected_dense_sequence": expected_index_sequence,
            "dense_sequence_valid": not linkage.get("index_duplicates") and sorted(index_sequence) == expected_index_sequence,
        }
        if compiled_pptx_total_slide_count is not None:
            compiled_positions = list(range(1, compiled_pptx_total_slide_count + 1))
            actual_positions = sorted(linkage.get("index_set", []))
            sequence_by_pptx_index = list(linkage.get("sequence_by_pptx_index", []))
            position_match_sources: dict[str, bool] = {}
            if "blueprint" in identity:
                position_match_sources["blueprint"] = sequence_by_pptx_index == list(identity["blueprint"]["sequence"])
            if "slide_ledger" in identity:
                position_match_sources["slide_ledger"] = sequence_by_pptx_index == list(identity["slide_ledger"]["sequence"])
            identity["compiled_pptx_position"] = {
                "compiled_positions": compiled_positions,
                "linkage_positions": actual_positions,
                "missing_positions": [position for position in compiled_positions if position not in set(actual_positions)],
                "extra_positions": [position for position in actual_positions if position not in set(compiled_positions)],
                "sequence_by_pptx_index": sequence_by_pptx_index,
                "position_match_sources": position_match_sources,
            }

    if {"blueprint", "slide_ledger", "slide_build_linkage"}.issubset(identity):
        all_slide_numbers = sorted(
            set(identity["blueprint"]["slide_number_set"])
            | set(identity["slide_ledger"]["slide_number_set"])
            | set(identity["slide_build_linkage"]["slide_number_set"])
        )
        deck_mode_mismatches: list[dict[str, Any]] = []
        for slide_number in all_slide_numbers:
            modes = {
                "blueprint": identity["blueprint"]["deck_modes"].get(slide_number),
                "slide_ledger": identity["slide_ledger"]["deck_modes"].get(slide_number),
                "slide_build_linkage": identity["slide_build_linkage"]["deck_modes"].get(slide_number),
            }
            if len(set(modes.values())) != 1:
                deck_mode_mismatches.append({"slide_number": slide_number, "modes": modes})
        identity["per_slide_deck_mode_mismatches"] = deck_mode_mismatches
    return identity


def _safe_int_field(
    payload: dict[str, Any],
    field_name: str,
    *,
    default: int,
    artifact_errors: list[str],
) -> int:
    if field_name not in payload:
        return default
    value = payload.get(field_name)
    if isinstance(value, bool):
        artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
        return default
    if isinstance(value, int):
        return value
    if isinstance(value, str):
        stripped = value.strip()
        if re.fullmatch(r"-?\d+", stripped):
            return int(stripped)
    artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
    return default


def _safe_optional_int_field(
    payload: dict[str, Any],
    field_name: str,
    *,
    artifact_errors: list[str],
) -> int | None:
    if field_name not in payload:
        return None
    value = payload.get(field_name)
    if isinstance(value, bool):
        artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
        return None
    if isinstance(value, int):
        return value
    if isinstance(value, str):
        stripped = value.strip()
        if re.fullmatch(r"-?\d+", stripped):
            return int(stripped)
    artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
    return None


def _safe_float_field(
    payload: dict[str, Any],
    field_name: str,
    *,
    default: float,
    artifact_errors: list[str],
) -> float:
    if field_name not in payload:
        return default
    value = payload.get(field_name)
    if isinstance(value, bool):
        artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
        return default
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        stripped = value.strip()
        if re.fullmatch(r"-?\d+(?:\.\d+)?", stripped):
            return float(stripped)
    artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
    return default


def _safe_list_length_field(
    payload: dict[str, Any],
    field_name: str,
    *,
    default: int,
    artifact_errors: list[str],
) -> int:
    if field_name not in payload:
        return default
    value = payload.get(field_name)
    if isinstance(value, list):
        return len(value)
    artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
    return default


def _safe_optional_list_length_field(
    payload: dict[str, Any],
    field_name: str,
    *,
    artifact_errors: list[str],
) -> int | None:
    if field_name not in payload:
        return None
    value = payload.get(field_name)
    if isinstance(value, list):
        return len(value)
    artifact_errors.append(f"validation-summary field `{field_name}` is malformed: {value!r}.")
    return None


def _iter_shape_text(shape: Any) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = (shape.text or "").strip()
        if text:
            texts.append(text)
    table = getattr(shape, "table", None)
    if table is not None:
        for row in table.rows:
            for cell in row.cells:
                text = (cell.text or "").strip()
                if text:
                    texts.append(text)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(_iter_shape_text(child))
    return texts


def _has_structured_compile_warning_signal(qa_policy_summary: QAVerdictSummary | None) -> bool:
    return qa_policy_summary is not None and BUILD_WARNING_STRING_SURFACE_CODE in qa_policy_summary.compatibility_warning_codes


def _resolve_compile_warning_count(
    *,
    qa_policy_summary: QAVerdictSummary | None,
    validation_compile_warning_count: int | None,
    validation_compile_warnings_length: int | None,
    raw_build_warning_count: int,
    raw_build_warnings: list[str],
    artifact_errors: list[str],
) -> int:
    if _has_structured_compile_warning_signal(qa_policy_summary):
        if validation_compile_warning_count is not None:
            if raw_build_warning_count != validation_compile_warning_count:
                artifact_errors.append(
                    "build-manifest warning count does not match structured compile-warning policy surface."
                )
            return validation_compile_warning_count
        if validation_compile_warnings_length is not None:
            if raw_build_warning_count != validation_compile_warnings_length:
                artifact_errors.append(
                    "build-manifest warning count does not match structured compile-warning policy surface."
                )
            return validation_compile_warnings_length
        if raw_build_warning_count > 0:
            return raw_build_warning_count
        if raw_build_warnings:
            return len(raw_build_warnings)
        return 1
    return raw_build_warning_count


def _inspect_compiled_pptx(pptx_path: Path, substrings: list[str]) -> tuple[int, list[dict[str, Any]]]:
    hits: list[dict[str, Any]] = []
    presentation = Presentation(str(pptx_path))
    slide_count = len(presentation.slides)
    lowered_needles = [(needle, needle.lower()) for needle in substrings if needle.strip()]
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            for text in _iter_shape_text(shape):
                lowered_text = text.lower()
                for raw_needle, lowered_needle in lowered_needles:
                    if lowered_needle in lowered_text:
                        hits.append(
                            {
                                "slide_number": slide_index,
                                "matched_substring": raw_needle,
                                "text": text,
                            }
                        )
                        break
    return slide_count, hits


def _scan_forbidden_visible_text(pptx_path: Path, substrings: list[str]) -> list[dict[str, Any]]:
    _slide_count, hits = _inspect_compiled_pptx(pptx_path, substrings)
    return hits


def _appendix_boundary_violation_count(
    blueprint: Blueprint,
    slide_ledger: SlideLedger,
    slide_build_linkage: SlideBuildLinkage,
) -> tuple[int, list[int]]:
    if blueprint.appendix_start is None:
        return 0, []

    ledger_modes = {entry.slide_number: entry.deck_mode for entry in slide_ledger.entries}
    linkage_modes = {entry.slide_number: entry.deck_mode for entry in slide_build_linkage.slides}
    violations: set[int] = set()
    for slide in blueprint.slides:
        expected_appendix = slide.slide_number >= blueprint.appendix_start
        modes = [slide.deck_mode, ledger_modes.get(slide.slide_number), linkage_modes.get(slide.slide_number)]
        if expected_appendix:
            if any(mode is None or mode != DeckMode.APPENDIX for mode in modes):
                violations.add(slide.slide_number)
        else:
            if any(mode == DeckMode.APPENDIX for mode in modes if mode is not None):
                violations.add(slide.slide_number)
    return len(violations), sorted(violations)


def _layout_compatibility_failures(blueprint: Blueprint, layout_library: Any) -> tuple[int, list[int]]:
    pattern_map = {pattern.pattern_id: pattern for pattern in layout_library.patterns}
    failed_slides: list[int] = []
    for slide in blueprint.slides:
        pattern = pattern_map.get(slide.layout_pattern_id)
        if pattern is None:
            failed_slides.append(slide.slide_number)
            continue
        if slide.slide_role not in pattern.slide_roles or slide.visual_type not in pattern.supported_visual_types:
            failed_slides.append(slide.slide_number)
    return len(failed_slides), failed_slides


def _extract_slide_number_from_warning(warning: Any) -> int | None:
    if isinstance(warning, Mapping):
        for key in ("slide_number", "slide", "pptx_index", "build_link_index"):
            value = warning.get(key)
            if isinstance(value, int) and value > 0:
                return value
            if isinstance(value, str) and value.strip().isdigit():
                return int(value.strip())
        text = str(warning.get("summary") or warning.get("message") or warning.get("warning") or "").strip()
    else:
        text = str(warning or "").strip()
    if not text:
        return None

    for pattern in WARNING_SLIDE_PATTERNS:
        match = pattern.search(text)
        if match is None:
            continue
        if pattern.pattern == WARNING_SLIDE_PATTERNS[1].pattern:
            lowered = text.lower()
            if not any(keyword in lowered for keyword in WARNING_CONTEXT_KEYWORDS):
                continue
            numeric_tokens = re.findall(r"\d+", text)
            if len(numeric_tokens) != 1:
                continue
        return int(match.group(1))
    return None


def _deterministic_fallback_violations(
    build_manifest: BuildManifest,
    slide_build_linkage: SlideBuildLinkage,
    layout_library: Any,
    forbidden_fallback_markers: list[str],
) -> tuple[int, list[int]]:
    lowered_markers = [marker.lower() for marker in (forbidden_fallback_markers or FORBIDDEN_FALLBACK_MARKERS) if marker.strip()]
    approved_layouts = {pattern.pattern_id for pattern in layout_library.patterns}
    violation_slides: set[int] = set()

    for warning in build_manifest.warnings:
        lowered_warning = warning.lower()
        if any(marker in lowered_warning for marker in lowered_markers):
            slide_number = _extract_slide_number_from_warning(warning)
            if slide_number is not None:
                violation_slides.add(slide_number)

    for entry in slide_build_linkage.slides:
        family = (entry.layout_family or "").strip().lower()
        pattern_id = (entry.layout_pattern_id or "").strip()
        if not family or "placeholder" in family:
            violation_slides.add(entry.slide_number)
            continue
        if pattern_id not in approved_layouts:
            violation_slides.add(entry.slide_number)
            continue
        if any(marker in family for marker in lowered_markers):
            violation_slides.add(entry.slide_number)

    violation_slides.discard(0)
    return len(violation_slides), sorted(violation_slides)


def _classify_native_visual_misclassification(
    finding: Any,
    slide_number: int | None,
    linkage_by_slide: dict[int, Any],
    *,
    policy: ProductionGatePolicy,
) -> NativeVisualClassificationOutcome:
    required_fields = ("finding_id", "category", "tags", "summary")
    missing_fields = [field for field in required_fields if not hasattr(finding, field)]
    if missing_fields:
        return NativeVisualClassificationOutcome(
            schema_error="qa finding schema missing fields required for native-visual classification: "
            + ", ".join(sorted(missing_fields))
        )
    if slide_number is None:
        return NativeVisualClassificationOutcome()

    normalized_tags = {_normalized_value(tag) for tag in getattr(finding, "tags", [])}
    finding_id = _normalized_value(getattr(finding, "finding_id", ""))
    category = _normalized_value(getattr(finding, "category", ""))
    summary = _normalized_value(getattr(finding, "summary", ""))
    mode: str | None = None

    if (
        finding_id.startswith(NATIVE_VISUAL_MISSING_FINDING_ID_PREFIX)
        and category == NATIVE_VISUAL_MISSING_CATEGORY
        and NATIVE_VISUAL_MISSING_REQUIRED_TAGS.issubset(normalized_tags)
    ):
        mode = "structured"
    elif "expects a visual asset but none appears in the compiled deck" in summary:
        mode = "compatibility-fallback"
    else:
        return NativeVisualClassificationOutcome()

    entry = linkage_by_slide.get(slide_number)
    if entry is None:
        return NativeVisualClassificationOutcome(mode=mode)
    if _normalized_value(entry.visual_type) not in {_normalized_value(value) for value in policy.native_visual_types}:
        return NativeVisualClassificationOutcome(mode=mode)
    if _normalized_value(entry.layout_family) not in {_normalized_value(value) for value in policy.native_layout_families}:
        return NativeVisualClassificationOutcome(mode=mode)
    if entry.missing_dependencies:
        return NativeVisualClassificationOutcome(mode=mode)
    if not (entry.linked_paths or entry.viz_spec_ids or entry.asset_ids):
        return NativeVisualClassificationOutcome(mode=mode)
    return NativeVisualClassificationOutcome(is_misclassification=True, mode=mode)


def _derive_compiled_pptx_deck_mode_counts(
    slide_build_linkage: SlideBuildLinkage | None,
    compiled_pptx_total_slide_count: int | None,
) -> tuple[int | None, int | None, str, str | None]:
    if slide_build_linkage is None:
        return None, None, "unavailable", "slide-build-linkage artifact is missing."
    if compiled_pptx_total_slide_count is None:
        return None, None, "unavailable", "compiled PPTX slide count is unavailable."
    indices = [entry.pptx_index for entry in slide_build_linkage.slides]
    if len(indices) != len(set(indices)):
        return None, None, "unavailable", "slide-build-linkage pptx_index values are not unique."
    if sorted(indices) != list(range(1, len(indices) + 1)):
        return None, None, "unavailable", "slide-build-linkage pptx_index values are not a dense 1..N sequence."
    if len(indices) != compiled_pptx_total_slide_count:
        return None, None, "unavailable", "compiled PPTX count and slide-build-linkage count disagree."
    main_story = sum(1 for entry in slide_build_linkage.slides if entry.deck_mode == DeckMode.MAIN_STORY)
    appendix = sum(1 for entry in slide_build_linkage.slides if entry.deck_mode == DeckMode.APPENDIX)
    return main_story, appendix, "linkage-derived", None


def _resolve_candidate_path(path_text: str | None, *, repo_root: Path, target_root: Path) -> Path | None:
    if not isinstance(path_text, str) or not path_text.strip():
        return None
    candidate = Path(path_text)
    if candidate.is_absolute() or re.match(r"^[A-Za-z]:[\\/]", path_text.strip()):
        return candidate if candidate.exists() else None

    resolved_candidates = [
        (target_root / candidate).resolve(),
        (repo_root / candidate).resolve(),
    ]
    for resolved in resolved_candidates:
        if resolved.exists():
            return resolved
    return None


def _resolve_compiled_pptx_path(
    *,
    repo_root: Path,
    target_root: Path,
    validation_summary: dict[str, Any] | None,
    execution: TargetBuildExecution | None,
) -> tuple[Path | None, dict[str, str | None], list[str]]:
    resolution = {
        "source": "missing",
        "validation_summary_status": "absent",
        "validation_summary_hint": None,
    }
    errors: list[str] = []

    raw_hint = validation_summary.get("pptx_path") if isinstance(validation_summary, dict) else None
    if isinstance(raw_hint, str) and raw_hint.strip():
        resolution["validation_summary_hint"] = raw_hint
        resolved = _resolve_candidate_path(raw_hint, repo_root=repo_root, target_root=target_root)
        if resolved is not None:
            resolution["source"] = "validation-summary"
            resolution["validation_summary_status"] = "accepted"
            return resolved, resolution, errors
        resolution["validation_summary_status"] = "rejected-stale"

    execution_hint = execution.artifact_paths.get("pptx") if execution is not None else None
    if isinstance(execution_hint, str) and execution_hint.strip():
        resolved = _resolve_candidate_path(execution_hint, repo_root=repo_root, target_root=target_root)
        if resolved is not None:
            resolution["source"] = "execution-artifact-path"
            if resolution["validation_summary_status"] == "rejected-stale":
                resolution["validation_summary_status"] = "replaced-by-execution-artifact-path"
            elif resolution["validation_summary_status"] == "absent":
                resolution["validation_summary_status"] = "absent"
            return resolved, resolution, errors

    pptx_candidates = sorted((target_root / "artifacts" / "pptx").glob("*.pptx"))
    if pptx_candidates:
        resolution["source"] = "target-artifact-fallback"
        if resolution["validation_summary_status"] == "rejected-stale":
            resolution["validation_summary_status"] = "replaced-by-target-artifact-fallback"
        return pptx_candidates[0], resolution, errors

    errors.append("Compiled PPTX artifact is missing.")
    if resolution["validation_summary_status"] == "absent":
        resolution["validation_summary_status"] = "missing"
    return None, resolution, errors

def _asset_status_counts(asset_manifest: Any) -> dict[str, int]:
    return dict(sorted(Counter(asset.status.value for asset in asset_manifest.assets).items()))


def _write_validation_summary(
    *,
    path: Path,
    policy: ProductionGatePolicy,
    target: RegressionTargetSpec,
    workflow_plan: Any,
    blueprint: Blueprint,
    layout_library: Any,
    slide_ledger: SlideLedger,
    asset_manifest: Any,
    build_manifest: BuildManifest,
    slide_build_linkage: SlideBuildLinkage,
    qa_report: QAReport,
    llm_backend_proof: LLMBackendProof,
    forbidden_visible_text_findings: list[dict[str, Any]],
    forbidden_fallback_markers: list[str],
    repo_root: Path,
    artifact_paths: dict[str, Path],
) -> Path:
    authoring_preview = _load_json(artifact_paths["authoring_preview"]) if artifact_paths["authoring_preview"].is_file() else None
    compiled_deck_authoring_audit = _load_json(artifact_paths["compiled_deck_authoring_audit"]) if artifact_paths["compiled_deck_authoring_audit"].is_file() else None
    planned_authoring_metrics = authoring_preview.get("repetition_metrics", {}) if isinstance(authoring_preview, dict) else {}
    if not isinstance(planned_authoring_metrics, dict):
        planned_authoring_metrics = {}
    compiled_deck_metrics = compiled_deck_authoring_audit.get("compiled_metrics", {}) if isinstance(compiled_deck_authoring_audit, dict) else {}
    if not isinstance(compiled_deck_metrics, dict):
        compiled_deck_metrics = {}
    compiled_deck_truth_mismatches = compiled_deck_authoring_audit.get("truth_mismatches", []) if isinstance(compiled_deck_authoring_audit, dict) else []
    if not isinstance(compiled_deck_truth_mismatches, list):
        compiled_deck_truth_mismatches = []
    main_story_actual = sum(1 for slide in blueprint.slides if slide.deck_mode == DeckMode.MAIN_STORY)
    appendix_actual = sum(1 for slide in blueprint.slides if slide.deck_mode == DeckMode.APPENDIX)
    appendix_boundary_violation_count, appendix_boundary_slides = _appendix_boundary_violation_count(
        blueprint,
        slide_ledger,
        slide_build_linkage,
    )
    layout_compatibility_failure_count, layout_failure_slides = _layout_compatibility_failures(blueprint, layout_library)
    deterministic_fallback_violation_count, deterministic_fallback_slides = _deterministic_fallback_violations(
        build_manifest,
        slide_build_linkage,
        layout_library,
        forbidden_fallback_markers,
    )
    linkage_by_slide = {entry.slide_number: entry for entry in slide_build_linkage.slides}
    native_visual_misclassification_slides = sorted(
        {
            finding.slide_number
            for finding in qa_report.findings
            if _classify_native_visual_misclassification(
                finding,
                finding.slide_number,
                linkage_by_slide,
                policy=policy,
            ).is_misclassification
        }
    )
    qa_verdict_summary = getattr(qa_report, "verdict_summary", None)
    qa_status_value = qa_report.qa_status.value
    if qa_verdict_summary is not None:
        qa_status_value = qa_verdict_summary.qa_status.value
    payload = {
        "summary_id": f"{target.target_id}-validation-summary",
        "target_id": target.target_id,
        "pptx_path": str(artifact_paths["pptx"]),
        "provider_requested": llm_backend_proof.provider_requested,
        "provider_used": llm_backend_proof.provider_used,
        "model_requested": llm_backend_proof.model_requested,
        "model_used": llm_backend_proof.model_used,
        "endpoint_requested": llm_backend_proof.endpoint_requested,
        "endpoint_used": llm_backend_proof.endpoint_used,
        "transport_used": llm_backend_proof.transport_used,
        "strict_structured_output": llm_backend_proof.strict_structured_output,
        "llm_request_count": llm_backend_proof.llm_request_count,
        "llm_request_targets": llm_backend_proof.llm_request_targets,
        "workflow_option": workflow_plan.workflow_option,
        "main_story_budget": workflow_plan.main_story_slide_count_range.model_dump(mode="json"),
        "appendix_budget": workflow_plan.appendix_candidate_slide_count_range.model_dump(mode="json"),
        "main_story_actual": main_story_actual,
        "appendix_actual": appendix_actual,
        "compile_warning_count": len(build_manifest.warnings),
        "compile_warnings": list(build_manifest.warnings),
        "qa_status": qa_status_value,
        "qa_findings": qa_report.summary.finding_count,
        "qa_blocking": qa_report.summary.blocking_count,
        "qa_severity_counts": dict(sorted(qa_report.summary.severity_counts.items())),
        "qa_category_counts": dict(sorted(Counter(finding.category for finding in qa_report.findings).items())),
        "forbidden_visible_text_findings": forbidden_visible_text_findings,
        "planned_authoring_metrics": dict(sorted(planned_authoring_metrics.items())),
        "compiled_deck_metrics": dict(sorted(compiled_deck_metrics.items())),
        "compiled_deck_truth_mismatches": compiled_deck_truth_mismatches,
        "compiled_deck_truth_mismatch_count": len(compiled_deck_truth_mismatches),
        "appendix_boundary_violation_count": appendix_boundary_violation_count,
        "appendix_boundary_violation_slides": appendix_boundary_slides,
        "layout_compatibility_failure_count": layout_compatibility_failure_count,
        "layout_compatibility_failure_slides": layout_failure_slides,
        "deterministic_fallback_violation_count": deterministic_fallback_violation_count,
        "deterministic_fallback_violation_slides": deterministic_fallback_slides,
        "native_visual_misclassification_count": len(native_visual_misclassification_slides),
        "native_visual_misclassification_slides": native_visual_misclassification_slides,
        "asset_status_counts": _asset_status_counts(asset_manifest),
        "llm_backend_proof": {
            "path": str(artifact_paths["llm_backend_proof"]),
            "request_targets": llm_backend_proof.llm_request_targets,
        },
        "artifact_paths": {
            name: str(path_value.relative_to(repo_root)) if path_value.is_relative_to(repo_root) else str(path_value)
            for name, path_value in artifact_paths.items()
        },
    }
    if qa_verdict_summary is not None:
        payload["qa_verdict_summary"] = qa_verdict_summary.model_dump(mode="json")
        payload["compile_eligibility"] = qa_verdict_summary.compile_eligibility.value
    return _write_json(path, payload)


def _run_regression_target(
    *,
    repo_root: Path,
    target: RegressionTargetSpec,
    gate_output_root: Path,
    material_mode: MaterialMode,
    policy: ProductionGatePolicy,
    provider_settings: ProviderSettings,
) -> TargetBuildExecution:
    target_root = gate_output_root / target.output_subdir
    _prepare_output_dir(target_root)
    state_dir = target_root / "state"
    state_dir.mkdir(parents=True, exist_ok=True)
    artifact_paths = _artifact_paths_for_target(target_root, target)
    error_path = target_root / "state" / "build-error.json"

    try:
        materials = _resolve_target_materials(target, target_root, material_mode)
        brief = WorkflowBriefInput(
            topic=target.topic,
            deck_title=target.deck_title,
            audience=target.audience,
            purpose=target.purpose,
            delivery_mode=target.delivery_mode,
            expected_duration_minutes=target.expected_duration_minutes,
            current_materials=materials,
            constraints=target.constraints,
            notes=target.notes,
        )
        workflow_plan, llm_backend_proof = plan_workflow_with_provider(
            brief,
            provider_settings=provider_settings,
            llm_backend_proof_path=artifact_paths["llm_backend_proof"],
        )
        save_state_file(workflow_plan, artifact_paths["workflow_plan"])

        gate2_outputs = plan_gate2(workflow_plan, brief=brief)
        write_gate2_outputs(gate2_outputs, state_dir)
        approved_blueprint = gate2_outputs.blueprint.model_copy(update={"approval_status": StageStatus.APPROVED})
        save_state_file(approved_blueprint, artifact_paths["blueprint"])

        asset_outputs = derive_assets_from_blueprint(
            blueprint=approved_blueprint,
            design_system=gate2_outputs.design_system,
            deck_constitution=gate2_outputs.deck_constitution,
            layout_library=gate2_outputs.layout_library,
            slide_ledger=gate2_outputs.slide_ledger,
            asset_requests=gate2_outputs.asset_requests,
        )
        write_asset_derivation_outputs(asset_outputs, state_dir)

        structured_visual_outputs = run_structured_visuals(
            viz_spec=asset_outputs.viz_spec,
            design_system=gate2_outputs.design_system,
            slide_ledger=asset_outputs.slide_ledger,
            output_dir=target_root / "artifacts" / "structured-visuals",
            deck_constitution=gate2_outputs.deck_constitution,
            layout_library=gate2_outputs.layout_library,
            asset_requests=asset_outputs.asset_requests,
            blueprint=approved_blueprint,
            root=repo_root,
        )
        write_structured_visual_outputs(structured_visual_outputs, state_dir)

        initial_crop_outputs = run_document_asset_crop(
            asset_requests=asset_outputs.asset_requests,
            slide_ledger=structured_visual_outputs.slide_ledger,
            output_dir=state_dir / "document-crops-initial",
            asset_manifest=structured_visual_outputs.asset_manifest,
            root=repo_root,
        )
        write_document_crop_outputs(initial_crop_outputs, state_dir / "document-crops-initial")

        reviewed_crop_outputs = run_document_crop_review(
            asset_requests=asset_outputs.asset_requests,
            crop_candidates=initial_crop_outputs.crop_candidates,
            asset_manifest=initial_crop_outputs.asset_manifest,
            slide_ledger=initial_crop_outputs.slide_ledger,
            output_dir=state_dir / "document-crops-reviewed",
            root=repo_root,
        )
        write_document_crop_outputs(reviewed_crop_outputs, state_dir / "document-crops-reviewed")
        save_state_file(reviewed_crop_outputs.asset_manifest, artifact_paths["asset_manifest"])
        save_state_file(reviewed_crop_outputs.slide_ledger, artifact_paths["slide_ledger"])

        compile_outputs = compile_pptx(
            blueprint=approved_blueprint,
            design_system=gate2_outputs.design_system,
            deck_constitution=gate2_outputs.deck_constitution,
            layout_library=gate2_outputs.layout_library,
            slide_ledger=reviewed_crop_outputs.slide_ledger,
            asset_manifest=reviewed_crop_outputs.asset_manifest,
            viz_manifest=structured_visual_outputs.viz_manifest,
            output_dir=target_root / "artifacts" / "pptx",
            pptx_name=_pptx_name_for_target(target),
            root=repo_root,
        )
        write_pptx_compile_outputs(compile_outputs, state_dir)

        qa_outputs = run_deck_qa(
            blueprint=approved_blueprint,
            design_system=gate2_outputs.design_system,
            deck_constitution=gate2_outputs.deck_constitution,
            layout_library=gate2_outputs.layout_library,
            slide_ledger=compile_outputs.slide_ledger,
            asset_manifest=reviewed_crop_outputs.asset_manifest,
            viz_manifest=structured_visual_outputs.viz_manifest,
            build_manifest=compile_outputs.build_manifest,
            slide_build_linkage=compile_outputs.slide_build_linkage,
            artifact_root=target_root,
        )
        write_deck_qa_outputs(qa_outputs, state_dir)

        forbidden_visible_text_findings = _scan_forbidden_visible_text(
            artifact_paths["pptx"],
            policy.forbidden_visible_text_substrings,
        )
        validation_summary_path = _write_validation_summary(
            path=artifact_paths["validation_summary"],
            policy=policy,
            target=target,
            workflow_plan=workflow_plan,
            blueprint=approved_blueprint,
            layout_library=gate2_outputs.layout_library,
            slide_ledger=qa_outputs.slide_ledger,
            asset_manifest=reviewed_crop_outputs.asset_manifest,
            build_manifest=compile_outputs.build_manifest,
            slide_build_linkage=qa_outputs.slide_build_linkage,
            qa_report=qa_outputs.qa_report,
            llm_backend_proof=llm_backend_proof,
            forbidden_visible_text_findings=forbidden_visible_text_findings,
            forbidden_fallback_markers=policy.forbidden_fallback_markers,
            repo_root=repo_root,
            artifact_paths=artifact_paths,
        )
        return TargetBuildExecution(
            target_id=target.target_id,
            output_dir=str(target_root),
            artifact_paths={name: str(path_value) for name, path_value in artifact_paths.items()},
            validation_summary_path=str(validation_summary_path),
            provider_requested=provider_settings.provider,
            model_requested=provider_settings.model,
            endpoint_requested=provider_settings.endpoint,
        )
    except Exception as exc:
        error_payload = {
            "target_id": target.target_id,
            "error": f"{type(exc).__name__}: {exc}",
            "traceback": traceback.format_exc(),
        }
        _write_json(error_path, error_payload)
        return TargetBuildExecution(
            target_id=target.target_id,
            output_dir=str(target_root),
            artifact_paths={name: str(path_value) for name, path_value in artifact_paths.items()} | {"build_error": str(error_path)},
            build_error=error_payload["error"],
            provider_requested=provider_settings.provider,
            model_requested=provider_settings.model,
            endpoint_requested=provider_settings.endpoint,
        )


def _load_required_artifacts(
    target_root: Path,
    policy: ProductionGatePolicy,
    *,
    repo_root: Path,
    execution: TargetBuildExecution | None = None,
) -> tuple[dict[str, Any], list[str], dict[str, str], dict[str, str | None]]:
    loaded: dict[str, Any] = {}
    errors: list[str] = []
    artifact_paths: dict[str, str] = {}
    loaders: dict[str, Any] = {
        "workflow_plan": load_state_file,
        "authoring_preview": _load_json,
        "blueprint": load_state_file,
        "design_system": load_state_file,
        "deck_constitution": load_state_file,
        "layout_library": load_state_file,
        "slide_ledger": load_state_file,
        "asset_requests": load_state_file,
        "asset_manifest": load_state_file,
        "viz_spec": load_state_file,
        "viz_manifest": load_state_file,
        "qa_report": load_state_file,
        "build_manifest": load_pptx_compile_file,
        "slide_build_linkage": load_pptx_compile_file,
        "compiled_deck_text": _load_json,
        "compiled_deck_shape_census": _load_json,
        "compiled_deck_authoring_audit": _load_json,
        "compiled_deck_thumbnail_index": _load_json,
        "compiled_deck_visual_review_summary": _load_json,
        "validation_summary": _load_json,
    }
    for relative_path in policy.required_artifacts:
        path = target_root / relative_path
        artifact_paths[relative_path] = str(path)
        if not path.exists():
            errors.append(f"Missing required artifact `{relative_path}`.")
            continue
        key = path.stem.replace("-", "_")
        loader = loaders.get(key)
        if loader is None:
            loaded[key] = path
            continue
        try:
            loaded[key] = loader(path)
        except Exception as exc:
            errors.append(f"Malformed required artifact `{relative_path}`: {type(exc).__name__}: {exc}")

    llm_backend_proof_path = None
    if execution is not None:
        proof_hint = execution.artifact_paths.get("llm_backend_proof")
        if isinstance(proof_hint, str) and proof_hint.strip():
            llm_backend_proof_path = Path(proof_hint)
    if llm_backend_proof_path is None:
        llm_backend_proof_path = target_root / "state" / "llm-backend-proof.json"
    artifact_paths["llm_backend_proof"] = str(llm_backend_proof_path)
    if llm_backend_proof_path.exists():
        try:
            loaded["llm_backend_proof"] = _load_json(llm_backend_proof_path)
        except Exception as exc:
            errors.append(f"Malformed llm backend proof `{llm_backend_proof_path}`: {type(exc).__name__}: {exc}")
    else:
        errors.append(f"Missing llm backend proof `{llm_backend_proof_path}`.")

    truth_artifacts = {
        "authoring_preview": target_root / "state" / "authoring-preview.json",
        "compiled_deck_text": target_root / "state" / "compiled-deck-text.json",
        "compiled_deck_shape_census": target_root / "state" / "compiled-deck-shape-census.json",
        "compiled_deck_authoring_audit": target_root / "state" / "compiled-deck-authoring-audit.json",
        "compiled_deck_thumbnail_strip": target_root / "state" / "compiled-deck-thumbnail-strip.png",
        "compiled_deck_thumbnail_index": target_root / "state" / "compiled-deck-thumbnail-index.json",
        "compiled_deck_visual_review_summary": target_root / "state" / "compiled-deck-visual-review-summary.json",
    }
    for key, path in truth_artifacts.items():
        artifact_paths[key] = str(path)
        if not path.exists():
            if key != "authoring_preview":
                errors.append(f"Missing compiled-deck truth artifact `{path}`.")
            continue
        loader = loaders.get(key)
        if loader is None:
            loaded[key] = path
            continue
        try:
            loaded[key] = loader(path)
        except Exception as exc:
            errors.append(f"Malformed compiled-deck truth artifact `{path}`: {type(exc).__name__}: {exc}")

    validation_summary = loaded.get("validation_summary")
    pptx_path, path_resolution, path_errors = _resolve_compiled_pptx_path(
        repo_root=repo_root,
        target_root=target_root,
        validation_summary=validation_summary if isinstance(validation_summary, dict) else None,
        execution=execution,
    )
    errors.extend(path_errors)
    if pptx_path is not None:
        loaded["pptx_path"] = pptx_path
        artifact_paths["pptx"] = str(pptx_path)
    return loaded, errors, artifact_paths, path_resolution


def collect_target_metrics(
    *,
    target: RegressionTargetSpec,
    execution: TargetBuildExecution,
    policy: ProductionGatePolicy,
    repo_root: str | Path = ROOT,
) -> GateMetrics:
    target_root = Path(execution.output_dir).resolve()
    resolved_repo_root = Path(repo_root).resolve()
    loaded, artifact_errors, artifact_paths, path_resolution = _load_required_artifacts(
        target_root,
        policy,
        repo_root=resolved_repo_root,
        execution=execution,
    )
    metrics = GateMetrics(
        target_id=target.target_id,
        provider_requested=execution.provider_requested,
        model_requested=execution.model_requested,
        endpoint_requested=execution.endpoint_requested,
        artifact_errors=list(artifact_errors),
        artifact_paths=dict(sorted(artifact_paths.items())),
        llm_backend_proof_path=artifact_paths.get("llm_backend_proof"),
        compiled_pptx_path_source=path_resolution.get("source"),
        validation_summary_pptx_path_status=path_resolution.get("validation_summary_status"),
        validation_summary_pptx_path_hint=path_resolution.get("validation_summary_hint"),
    )
    if execution.build_error:
        metrics.artifact_errors.append(execution.build_error)

    blueprint = loaded.get("blueprint")
    slide_ledger = loaded.get("slide_ledger")
    layout_library = loaded.get("layout_library")
    build_manifest = loaded.get("build_manifest")
    slide_build_linkage = loaded.get("slide_build_linkage")
    qa_report = loaded.get("qa_report")
    validation_summary = loaded.get("validation_summary")
    pptx_path = loaded.get("pptx_path")
    llm_backend_proof = loaded.get("llm_backend_proof")
    authoring_preview = loaded.get("authoring_preview")
    compiled_deck_authoring_audit = loaded.get("compiled_deck_authoring_audit")

    if isinstance(authoring_preview, dict):
        repetition_metrics = authoring_preview.get("repetition_metrics")
        if isinstance(repetition_metrics, dict):
            metrics.planned_authoring_metrics = dict(sorted(repetition_metrics.items()))
        elif repetition_metrics is not None:
            metrics.artifact_errors.append(
                f"authoring-preview field `repetition_metrics` is malformed: {repetition_metrics!r}."
            )
    if isinstance(compiled_deck_authoring_audit, dict):
        planned_metrics_payload = compiled_deck_authoring_audit.get("planned_metrics")
        if isinstance(planned_metrics_payload, dict):
            normalized_planned_metrics = dict(sorted(planned_metrics_payload.items()))
            if metrics.planned_authoring_metrics and normalized_planned_metrics != metrics.planned_authoring_metrics:
                metrics.artifact_errors.append(
                    "compiled-deck-authoring-audit planned_metrics do not match authoring-preview repetition_metrics."
                )
            elif not metrics.planned_authoring_metrics:
                metrics.planned_authoring_metrics = normalized_planned_metrics
        elif planned_metrics_payload is not None:
            metrics.artifact_errors.append(
                f"compiled-deck-authoring-audit field `planned_metrics` is malformed: {planned_metrics_payload!r}."
            )
        compiled_metrics_payload = compiled_deck_authoring_audit.get("compiled_metrics")
        if isinstance(compiled_metrics_payload, dict):
            metrics.compiled_deck_metrics = dict(sorted(compiled_metrics_payload.items()))
            metrics.compiled_deck_chrome_block_count = _safe_int_field(
                compiled_metrics_payload,
                "chrome_block_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_repeated_title_stem_count = _safe_int_field(
                compiled_metrics_payload,
                "repeated_title_stem_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_title_body_duplication_count = _safe_int_field(
                compiled_metrics_payload,
                "title_body_duplication_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_repeated_archetype_count = _safe_int_field(
                compiled_metrics_payload,
                "repeated_archetype_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_repeated_geometry_count = _safe_int_field(
                compiled_metrics_payload,
                "repeated_geometry_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_repeated_geometry_rate = _safe_float_field(
                compiled_metrics_payload,
                "repeated_geometry_rate",
                default=0.0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_text_card_overuse_rate = _safe_float_field(
                compiled_metrics_payload,
                "text_card_overuse_rate",
                default=0.0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_chrome_dominance_rate = _safe_float_field(
                compiled_metrics_payload,
                "chrome_dominance_rate",
                default=0.0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_appendix_clone_count = _safe_int_field(
                compiled_metrics_payload,
                "appendix_clone_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_appendix_visual_clone_run_length = _safe_int_field(
                compiled_metrics_payload,
                "appendix_visual_clone_run_length",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_expected_visual_missing_count = _safe_int_field(
                compiled_metrics_payload,
                "expected_visual_missing_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_archetype_realization_mismatch_count = _safe_int_field(
                compiled_metrics_payload,
                "archetype_realization_mismatch_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_missing_visual_center_count = _safe_int_field(
                compiled_metrics_payload,
                "missing_visual_center_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_repetitive_motion_count = _safe_int_field(
                compiled_metrics_payload,
                "repetitive_motion_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
            metrics.compiled_deck_truth_mismatch_count = _safe_int_field(
                compiled_metrics_payload,
                "truth_mismatch_count",
                default=0,
                artifact_errors=metrics.artifact_errors,
            )
        else:
            metrics.artifact_errors.append(
                f"compiled-deck-authoring-audit field `compiled_metrics` is malformed: {compiled_metrics_payload!r}."
            )
        issue_slides_payload = compiled_deck_authoring_audit.get("issue_slides")
        if isinstance(issue_slides_payload, dict):
            normalized_issue_slides: dict[str, list[int]] = {}
            for key, value in issue_slides_payload.items():
                if isinstance(value, list):
                    normalized_issue_slides[str(key)] = [item for item in value if isinstance(item, int)]
                else:
                    metrics.artifact_errors.append(
                        f"compiled-deck-authoring-audit issue_slides entry `{key}` is malformed: {value!r}."
                    )
            metrics.compiled_deck_issue_slides = normalized_issue_slides
        elif issue_slides_payload is not None:
            metrics.artifact_errors.append(
                f"compiled-deck-authoring-audit field `issue_slides` is malformed: {issue_slides_payload!r}."
            )
        truth_mismatches_payload = compiled_deck_authoring_audit.get("truth_mismatches")
        if isinstance(truth_mismatches_payload, list):
            metrics.compiled_deck_truth_mismatches = [str(item) for item in truth_mismatches_payload]
        elif truth_mismatches_payload is not None:
            metrics.artifact_errors.append(
                f"compiled-deck-authoring-audit field `truth_mismatches` is malformed: {truth_mismatches_payload!r}."
            )

    if isinstance(llm_backend_proof, dict):
        metrics.provider_requested = str(llm_backend_proof.get("provider_requested") or metrics.provider_requested or "").strip() or None
        metrics.provider_used = str(llm_backend_proof.get("provider_used") or "").strip() or None
        metrics.model_requested = str(llm_backend_proof.get("model_requested") or metrics.model_requested or "").strip() or None
        metrics.model_used = str(llm_backend_proof.get("model_used") or "").strip() or None
        metrics.endpoint_requested = str(llm_backend_proof.get("endpoint_requested") or metrics.endpoint_requested or "").strip() or None
        metrics.endpoint_used = str(llm_backend_proof.get("endpoint_used") or "").strip() or None
        metrics.transport_used = str(llm_backend_proof.get("transport_used") or "").strip() or None
        strict_value = llm_backend_proof.get("strict_structured_output")
        if isinstance(strict_value, bool):
            metrics.strict_structured_output = strict_value
        elif strict_value is not None:
            metrics.artifact_errors.append(f"llm backend proof field `strict_structured_output` is malformed: {strict_value!r}.")
        metrics.llm_request_count = _safe_int_field(
            llm_backend_proof,
            "llm_request_count",
            default=metrics.llm_request_count,
            artifact_errors=metrics.artifact_errors,
        )
        targets_value = llm_backend_proof.get("llm_request_targets")
        if isinstance(targets_value, list):
            metrics.llm_request_targets = [str(item).strip() for item in targets_value if str(item).strip()]
        elif targets_value is not None:
            metrics.artifact_errors.append(f"llm backend proof field `llm_request_targets` is malformed: {targets_value!r}.")
        if metrics.provider_requested and metrics.provider_requested != "local-none":
            if metrics.provider_used != metrics.provider_requested:
                metrics.artifact_errors.append(
                    f"llm backend proof provider mismatch: requested {metrics.provider_requested!r}, used {metrics.provider_used!r}."
                )
            if metrics.llm_request_count < 1:
                metrics.artifact_errors.append(
                    f"provider `{metrics.provider_requested}` was requested but llm_request_count={metrics.llm_request_count}."
                )

    if blueprint is not None:
        (
            metrics.blueprint_total_slide_count,
            metrics.blueprint_main_story_slide_count,
            metrics.blueprint_appendix_slide_count,
        ) = _count_blueprint_slides(blueprint)
        metrics.main_story_slide_count = metrics.blueprint_main_story_slide_count or 0
        metrics.appendix_slide_count = metrics.blueprint_appendix_slide_count or 0
    if slide_ledger is not None:
        (
            metrics.slide_ledger_total_slide_count,
            metrics.slide_ledger_main_story_slide_count,
            metrics.slide_ledger_appendix_slide_count,
        ) = _count_slide_ledger_entries(slide_ledger)
    if slide_build_linkage is not None:
        (
            metrics.slide_build_linkage_total_slide_count,
            metrics.slide_build_linkage_main_story_slide_count,
            metrics.slide_build_linkage_appendix_slide_count,
        ) = _count_slide_build_linkage_entries(slide_build_linkage)
    raw_build_warnings = list(build_manifest.warnings) if build_manifest is not None else []
    raw_build_warning_count = len(raw_build_warnings)
    qa_policy_summary = _resolve_qa_policy_summary(
        qa_report=qa_report,
        validation_summary=validation_summary if isinstance(validation_summary, dict) else None,
        artifact_errors=metrics.artifact_errors,
    )
    validation_compile_warning_count = (
        _safe_optional_int_field(
            validation_summary,
            "compile_warning_count",
            artifact_errors=metrics.artifact_errors,
        )
        if isinstance(validation_summary, dict)
        else None
    )
    validation_compile_warnings_length = (
        _safe_optional_list_length_field(
            validation_summary,
            "compile_warnings",
            artifact_errors=metrics.artifact_errors,
        )
        if isinstance(validation_summary, dict)
        else None
    )
    metrics.compile_warning_count = _resolve_compile_warning_count(
        qa_policy_summary=qa_policy_summary,
        validation_compile_warning_count=validation_compile_warning_count,
        validation_compile_warnings_length=validation_compile_warnings_length,
        raw_build_warning_count=raw_build_warning_count,
        raw_build_warnings=raw_build_warnings,
        artifact_errors=metrics.artifact_errors,
    )
    if qa_policy_summary is not None:
        metrics.qa_status = qa_policy_summary.qa_status.value
        metrics.qa_compile_eligibility = qa_policy_summary.compile_eligibility.value
        metrics.qa_warning_reason_codes = list(qa_policy_summary.warning_reason_codes)
        metrics.qa_blocking_reason_codes = list(qa_policy_summary.blocking_reason_codes)
        metrics.qa_compatibility_warning_codes = list(qa_policy_summary.compatibility_warning_codes)
    elif qa_report is not None:
        metrics.qa_status = qa_report.qa_status.value
    if qa_report is not None:
        metrics.qa_blocking_count = qa_report.summary.blocking_count
    if pptx_path is not None:
        try:
            metrics.compiled_pptx_total_slide_count, forbidden_hits = _inspect_compiled_pptx(
                pptx_path,
                policy.forbidden_visible_text_substrings,
            )
            metrics.pptx_scan_forbidden_visible_text_count = len(forbidden_hits)
            metrics.forbidden_visible_text_count = len(forbidden_hits)
        except Exception as exc:
            metrics.artifact_errors.append(f"Compiled PPTX could not be scanned directly: {type(exc).__name__}: {exc}")
    metrics.slide_identity = _reconcile_slide_identity(
        blueprint=blueprint,
        slide_ledger=slide_ledger,
        slide_build_linkage=slide_build_linkage,
        compiled_pptx_total_slide_count=metrics.compiled_pptx_total_slide_count,
    )
    for comparison_key, comparison in metrics.slide_identity.get("pairwise", {}).items():
        if (
            comparison.get("left_only")
            or comparison.get("right_only")
            or comparison.get("left_duplicates")
            or comparison.get("right_duplicates")
            or not comparison.get("sequence_match", False)
        ):
            metrics.artifact_errors.append(
                f"slide identity mismatch for `{comparison_key}`: left_only={comparison.get('left_only', [])}, "
                f"right_only={comparison.get('right_only', [])}, left_duplicates={comparison.get('left_duplicates', [])}, "
                f"right_duplicates={comparison.get('right_duplicates', [])}, sequence_match={comparison.get('sequence_match', False)}."
            )
    linkage_index = metrics.slide_identity.get("linkage_pptx_index", {})
    if linkage_index and (
        linkage_index.get("duplicates") or not linkage_index.get("dense_sequence_valid", False)
    ):
        metrics.artifact_errors.append(
            "slide-build-linkage pptx_index sequence is invalid: "
            f"duplicates={linkage_index.get('duplicates', [])}, "
            f"expected_dense_sequence={linkage_index.get('expected_dense_sequence', [])}, "
            f"actual_sequence={linkage_index.get('sequence', [])}."
        )
    compiled_position = metrics.slide_identity.get("compiled_pptx_position", {})
    if compiled_position and (
        compiled_position.get("missing_positions")
        or compiled_position.get("extra_positions")
        or any(match is False for match in compiled_position.get("position_match_sources", {}).values())
    ):
        metrics.artifact_errors.append(
            "compiled PPTX positions do not align with slide-build linkage ordering: "
            f"missing_positions={compiled_position.get('missing_positions', [])}, "
            f"extra_positions={compiled_position.get('extra_positions', [])}, "
            f"position_match_sources={compiled_position.get('position_match_sources', {})}."
        )
    if metrics.slide_identity.get("per_slide_deck_mode_mismatches"):
            metrics.artifact_errors.append(
                "per-slide deck mode mismatches were found across blueprint, ledger, and linkage."
            )
    validation_compile_eligibility = None
    if isinstance(validation_summary, dict):
        validation_compile_eligibility = _coerce_compile_eligibility(
            validation_summary.get("compile_eligibility"),
            source_label="validation-summary compile_eligibility",
            artifact_errors=metrics.artifact_errors,
        )
        if validation_compile_eligibility is not None:
            if metrics.qa_compile_eligibility is None:
                metrics.qa_compile_eligibility = validation_compile_eligibility.value
            elif metrics.qa_compile_eligibility != validation_compile_eligibility.value:
                metrics.artifact_errors.append("validation-summary compile_eligibility does not match structured QA verdict.")
    if isinstance(validation_summary, dict):
        for field_name, metric_value in (
            ("provider_requested", metrics.provider_requested),
            ("provider_used", metrics.provider_used),
            ("model_requested", metrics.model_requested),
            ("model_used", metrics.model_used),
            ("endpoint_requested", metrics.endpoint_requested),
            ("endpoint_used", metrics.endpoint_used),
            ("transport_used", metrics.transport_used),
        ):
            summary_value = validation_summary.get(field_name)
            normalized_summary = None if summary_value is None else str(summary_value).strip() or None
            if metric_value != normalized_summary:
                metrics.artifact_errors.append(f"validation-summary {field_name} does not match llm backend proof.")
        strict_value = validation_summary.get("strict_structured_output")
        if strict_value is not None:
            if not isinstance(strict_value, bool):
                metrics.artifact_errors.append(
                    f"validation-summary field `strict_structured_output` is malformed: {strict_value!r}."
                )
            elif metrics.strict_structured_output is not None and metrics.strict_structured_output != strict_value:
                metrics.artifact_errors.append("validation-summary strict_structured_output does not match llm backend proof.")
        if metrics.llm_request_count != _safe_int_field(
            validation_summary,
            "llm_request_count",
            default=metrics.llm_request_count,
            artifact_errors=metrics.artifact_errors,
        ):
            metrics.artifact_errors.append("validation-summary llm_request_count does not match llm backend proof.")
        llm_targets = validation_summary.get("llm_request_targets")
        if isinstance(llm_targets, list):
            normalized_targets = [str(item).strip() for item in llm_targets if str(item).strip()]
            if metrics.llm_request_targets != normalized_targets:
                metrics.artifact_errors.append("validation-summary llm_request_targets do not match llm backend proof.")
        elif llm_targets is not None:
            metrics.artifact_errors.append(f"validation-summary field `llm_request_targets` is malformed: {llm_targets!r}.")
        planned_metrics_summary = validation_summary.get("planned_authoring_metrics")
        if isinstance(planned_metrics_summary, dict):
            if metrics.planned_authoring_metrics and planned_metrics_summary != metrics.planned_authoring_metrics:
                metrics.artifact_errors.append("validation-summary planned_authoring_metrics do not match authoring-preview.")
        elif planned_metrics_summary is not None:
            metrics.artifact_errors.append(f"validation-summary field `planned_authoring_metrics` is malformed: {planned_metrics_summary!r}.")
        compiled_metrics_summary = validation_summary.get("compiled_deck_metrics")
        if isinstance(compiled_metrics_summary, dict):
            if metrics.compiled_deck_metrics and compiled_metrics_summary != metrics.compiled_deck_metrics:
                metrics.artifact_errors.append("validation-summary compiled_deck_metrics do not match compiled-deck authoring audit.")
        elif compiled_metrics_summary is not None:
            metrics.artifact_errors.append(f"validation-summary field `compiled_deck_metrics` is malformed: {compiled_metrics_summary!r}.")
        compiled_truth_mismatch_count = _safe_int_field(
            validation_summary,
            "compiled_deck_truth_mismatch_count",
            default=metrics.compiled_deck_truth_mismatch_count,
            artifact_errors=metrics.artifact_errors,
        )
        if compiled_truth_mismatch_count != metrics.compiled_deck_truth_mismatch_count:
            metrics.artifact_errors.append(
                "validation-summary compiled_deck_truth_mismatch_count does not match compiled-deck authoring audit."
            )
        metrics.validation_summary_forbidden_visible_text_count = _safe_list_length_field(
            validation_summary,
            "forbidden_visible_text_findings",
            default=0,
            artifact_errors=metrics.artifact_errors,
        )
        if metrics.main_story_slide_count != _safe_int_field(
            validation_summary,
            "main_story_actual",
            default=metrics.main_story_slide_count,
            artifact_errors=metrics.artifact_errors,
        ):
            metrics.artifact_errors.append("validation-summary main_story_actual does not match blueprint count.")
        if metrics.appendix_slide_count != _safe_int_field(
            validation_summary,
            "appendix_actual",
            default=metrics.appendix_slide_count,
            artifact_errors=metrics.artifact_errors,
        ):
            metrics.artifact_errors.append("validation-summary appendix_actual does not match blueprint count.")
        if validation_compile_warning_count is not None and metrics.compile_warning_count != validation_compile_warning_count:
            if _has_structured_compile_warning_signal(qa_policy_summary):
                metrics.artifact_errors.append(
                    "validation-summary compile_warning_count does not match structured compile-warning policy surface."
                )
            else:
                metrics.artifact_errors.append("validation-summary compile_warning_count does not match build-manifest.")
        summary_qa_status = validation_summary.get("qa_status", metrics.qa_status)
        if metrics.qa_status != str(summary_qa_status):
            if qa_policy_summary is not None:
                metrics.artifact_errors.append("validation-summary qa_status does not match structured QA verdict.")
            else:
                metrics.artifact_errors.append("validation-summary qa_status does not match qa-report.")
        if metrics.qa_blocking_count != _safe_int_field(
            validation_summary,
            "qa_blocking",
            default=metrics.qa_blocking_count,
            artifact_errors=metrics.artifact_errors,
        ):
            metrics.artifact_errors.append("validation-summary qa_blocking does not match qa-report.")
        if metrics.compiled_pptx_total_slide_count is not None:
            summary_forbidden = metrics.validation_summary_forbidden_visible_text_count or 0
            metrics.forbidden_visible_text_artifact_match = summary_forbidden == metrics.pptx_scan_forbidden_visible_text_count
            if not metrics.forbidden_visible_text_artifact_match:
                metrics.artifact_errors.append(
                    "validation-summary forbidden visible text count does not match direct PPTX scan."
                )
    elif pptx_path is not None:
        metrics.forbidden_visible_text_artifact_match = True

    if blueprint is not None and slide_ledger is not None and slide_build_linkage is not None:
        appendix_violation_count, _ = _appendix_boundary_violation_count(blueprint, slide_ledger, slide_build_linkage)
        metrics.appendix_boundary_violation_count = appendix_violation_count
        if isinstance(validation_summary, dict) and _safe_int_field(
            validation_summary,
            "appendix_boundary_violation_count",
            default=appendix_violation_count,
            artifact_errors=metrics.artifact_errors,
        ) != appendix_violation_count:
            metrics.artifact_errors.append(
                "validation-summary appendix boundary count does not match blueprint/ledger/linkage evaluation."
            )
    if blueprint is not None and layout_library is not None:
        layout_failure_count, _ = _layout_compatibility_failures(blueprint, layout_library)
        metrics.layout_compatibility_failure_count = layout_failure_count
        if isinstance(validation_summary, dict) and _safe_int_field(
            validation_summary,
            "layout_compatibility_failure_count",
            default=layout_failure_count,
            artifact_errors=metrics.artifact_errors,
        ) != layout_failure_count:
            metrics.artifact_errors.append(
                "validation-summary layout compatibility count does not match blueprint/layout-library evaluation."
            )
    if build_manifest is not None and slide_build_linkage is not None and layout_library is not None:
        fallback_violation_count, _ = _deterministic_fallback_violations(
            build_manifest,
            slide_build_linkage,
            layout_library,
            policy.forbidden_fallback_markers,
        )
        metrics.deterministic_fallback_violation_count = fallback_violation_count
        if isinstance(validation_summary, dict) and _safe_int_field(
            validation_summary,
            "deterministic_fallback_violation_count",
            default=fallback_violation_count,
            artifact_errors=metrics.artifact_errors,
        ) != fallback_violation_count:
            metrics.artifact_errors.append(
                "validation-summary fallback violation count does not match build-manifest/linkage evaluation."
            )

    if qa_report is not None and blueprint is not None and slide_build_linkage is not None:
        deck_mode_by_slide = {slide.slide_number: slide.deck_mode.value for slide in blueprint.slides}
        linkage_by_slide = {entry.slide_number: entry for entry in slide_build_linkage.slides}
        metrics.residual_findings = [
            ResidualFinding(
                finding_id=finding.finding_id,
                category=finding.category,
                severity=finding.severity.value,
                slide_number=finding.slide_number,
                deck_mode=deck_mode_by_slide.get(finding.slide_number) if finding.slide_number is not None else None,
                blocking=finding.blocking,
            )
            for finding in qa_report.findings
        ]
        classifier_modes: set[str] = set()
        native_visual_misclassification_count = 0
        for finding in qa_report.findings:
            outcome = _classify_native_visual_misclassification(
                finding,
                finding.slide_number,
                linkage_by_slide,
                policy=policy,
            )
            if outcome.schema_error:
                metrics.artifact_errors.append(outcome.schema_error)
            if outcome.mode:
                classifier_modes.add(outcome.mode)
            if outcome.is_misclassification:
                native_visual_misclassification_count += 1
        if not classifier_modes:
            metrics.native_visual_classifier_mode = "not-triggered"
        elif len(classifier_modes) == 1:
            metrics.native_visual_classifier_mode = next(iter(classifier_modes))
        else:
            metrics.native_visual_classifier_mode = "mixed"
        metrics.native_visual_misclassification_count = native_visual_misclassification_count
        qa_category_counts = Counter(finding.category for finding in qa_report.findings)
        compiled_truth_expectations = (
            ("repeated-chrome", metrics.compiled_deck_chrome_block_count, "compiled deck chrome"),
            ("title-body-duplication", metrics.compiled_deck_title_body_duplication_count, "compiled deck title/body duplication"),
            ("archetype-mismatch", metrics.compiled_deck_archetype_realization_mismatch_count, "compiled deck archetype realization mismatches"),
            ("visual-needed-but-missing", metrics.compiled_deck_expected_visual_missing_count, "compiled deck expected visuals missing"),
            ("appendix-clone", metrics.compiled_deck_appendix_clone_count, "compiled deck appendix clone runs"),
            ("deck-motion-repetition", metrics.compiled_deck_repetitive_motion_count, "compiled deck repetitive motion"),
            ("compiled-deck-truth-mismatch", metrics.compiled_deck_truth_mismatch_count, "compiled deck truth mismatches"),
        )
        for category, count, label in compiled_truth_expectations:
            if count > 0 and qa_category_counts.get(category, 0) == 0:
                metrics.artifact_errors.append(
                    f"qa-report says no `{category}` findings, but {label}={count} in compiled-deck audit."
                )
        if isinstance(validation_summary, dict) and _safe_int_field(
            validation_summary,
            "native_visual_misclassification_count",
            default=native_visual_misclassification_count,
            artifact_errors=metrics.artifact_errors,
        ) != native_visual_misclassification_count:
            metrics.artifact_errors.append(
                "validation-summary native visual misclassification count does not match qa/linkage evaluation."
            )

    (
        metrics.compiled_pptx_main_story_slide_count,
        metrics.compiled_pptx_appendix_slide_count,
        metrics.compiled_pptx_deck_mode_derivation,
        metrics.compiled_pptx_deck_mode_derivation_reason,
    ) = _derive_compiled_pptx_deck_mode_counts(slide_build_linkage, metrics.compiled_pptx_total_slide_count)

    metrics.artifact_errors = list(dict.fromkeys(metrics.artifact_errors))
    return metrics


def _allowlisted(finding: ResidualFinding, allowlist: list[CandidateAllowlistRule]) -> bool:
    normalized_category = finding.category.strip().lower()
    normalized_severity = finding.severity.strip().lower()
    normalized_mode = (finding.deck_mode or "").strip().lower()
    for rule in allowlist:
        if normalized_category != rule.category.strip().lower():
            continue
        if normalized_severity != rule.severity.strip().lower():
            continue
        if rule.appendix_only and normalized_mode != DeckMode.APPENDIX.value:
            continue
        return True
    return False


def evaluate_gate_metrics(
    *,
    profile: GateProfile,
    target: RegressionTargetSpec,
    metrics: GateMetrics,
    policy: ProductionGatePolicy,
) -> list[GateRuleFinding]:
    thresholds = policy.thresholds
    findings: list[GateRuleFinding] = []

    def add_rule(
        rule_id: str,
        passed: bool,
        message: str,
        *,
        severity: GateRuleSeverity = GateRuleSeverity.ERROR,
        affected_slides: list[int] | None = None,
        metric_values: dict[str, Any] | None = None,
        profile_impact: list[GateProfile] | None = None,
    ) -> None:
        findings.append(
            GateRuleFinding(
                rule_id=rule_id,
                passed=passed,
                severity=severity,
                message=message,
                target_id=target.target_id,
                affected_slides=affected_slides or [],
                metric_values=metric_values or {},
                profile_impact=profile_impact or [GateProfile.CANDIDATE, GateProfile.RELEASE],
            )
        )

    def count_values(prefix: str) -> dict[str, int | None]:
        return {
            "total": getattr(metrics, f"{prefix}_total_slide_count"),
            "main_story": getattr(metrics, f"{prefix}_main_story_slide_count"),
            "appendix": getattr(metrics, f"{prefix}_appendix_slide_count"),
        }

    def add_pair_count_rule(rule_id: str, left_prefix: str, right_prefix: str, left_label: str, right_label: str) -> None:
        left_counts = count_values(left_prefix)
        right_counts = count_values(right_prefix)
        missing_sources = [
            source
            for source, counts in ((left_label, left_counts), (right_label, right_counts))
            if any(value is None for value in counts.values())
        ]
        if missing_sources:
            add_rule(
                rule_id,
                False,
                f"{left_label} and {right_label} slide counts could not be fully compared because metrics are missing for: {', '.join(missing_sources)}.",
                metric_values={left_prefix: left_counts, right_prefix: right_counts},
            )
            return
        passed = left_counts == right_counts
        add_rule(
            rule_id,
            passed,
            f"{left_label} and {right_label} slide counts match."
            if passed
            else f"{left_label} and {right_label} slide counts disagree.",
            metric_values={left_prefix: left_counts, right_prefix: right_counts},
        )

    def add_identity_rule(rule_id: str, comparison_key: str, message_prefix: str) -> None:
        comparison = metrics.slide_identity.get("pairwise", {}).get(comparison_key, {})
        if not comparison:
            add_rule(
                rule_id,
                False,
                f"{message_prefix} could not be evaluated because reconciliation data is unavailable.",
                metric_values={"slide_identity": metrics.slide_identity},
            )
            return
        passed = (
            not comparison.get("left_only")
            and not comparison.get("right_only")
            and not comparison.get("left_duplicates")
            and not comparison.get("right_duplicates")
            and comparison.get("sequence_match", False)
        )
        add_rule(
            rule_id,
            passed,
            f"{message_prefix} matches exactly."
            if passed
            else f"{message_prefix} disagrees on set, duplicates, or order.",
            metric_values=comparison,
        )

    add_rule(
        "REQUIRED_ARTIFACTS",
        not metrics.artifact_errors,
        "All required artifacts are present and well-formed."
        if not metrics.artifact_errors
        else "Required artifacts are missing or malformed: " + "; ".join(metrics.artifact_errors),
        metric_values={"artifact_errors": metrics.artifact_errors},
    )
    add_pair_count_rule(
        "BLUEPRINT_LEDGER_COUNT_MATCH",
        "blueprint",
        "slide_ledger",
        "Blueprint",
        "Slide ledger",
    )
    add_pair_count_rule(
        "BLUEPRINT_LINKAGE_COUNT_MATCH",
        "blueprint",
        "slide_build_linkage",
        "Blueprint",
        "Slide-build linkage",
    )
    add_pair_count_rule(
        "LEDGER_LINKAGE_COUNT_MATCH",
        "slide_ledger",
        "slide_build_linkage",
        "Slide ledger",
        "Slide-build linkage",
    )
    add_identity_rule(
        "BLUEPRINT_LEDGER_SLIDE_IDENTITY_MATCH",
        "blueprint_vs_slide_ledger",
        "Blueprint and slide ledger slide identity",
    )
    add_identity_rule(
        "BLUEPRINT_LINKAGE_SLIDE_IDENTITY_MATCH",
        "blueprint_vs_slide_build_linkage",
        "Blueprint and slide-build linkage slide identity",
    )
    add_identity_rule(
        "LEDGER_LINKAGE_SLIDE_IDENTITY_MATCH",
        "slide_ledger_vs_slide_build_linkage",
        "Slide ledger and slide-build linkage slide identity",
    )
    linkage_index = metrics.slide_identity.get("linkage_pptx_index", {})
    linkage_index_valid = bool(linkage_index) and bool(linkage_index.get("dense_sequence_valid")) and not linkage_index.get("duplicates")
    add_rule(
        "LINKAGE_PPTX_INDEX_SEQUENCE_VALID",
        linkage_index_valid,
        "Slide-build linkage PPTX indices form a dense 1..N sequence with no duplicates."
        if linkage_index_valid
        else "Slide-build linkage PPTX indices are missing, duplicated, or not a dense 1..N sequence.",
        metric_values=linkage_index or {"slide_identity": metrics.slide_identity},
    )
    compiled_position = metrics.slide_identity.get("compiled_pptx_position", {})
    position_sources = compiled_position.get("position_match_sources", {})
    linkage_position_match = bool(compiled_position) and not compiled_position.get("missing_positions") and not compiled_position.get(
        "extra_positions"
    )
    if position_sources:
        linkage_position_match = linkage_position_match and all(position_sources.values())
    else:
        linkage_position_match = False
    add_rule(
        "LINKAGE_PPTX_POSITION_MATCH",
        linkage_position_match,
        "Compiled PPTX positions align with slide-build linkage ordering."
        if linkage_position_match
        else "Compiled PPTX positions do not align with slide-build linkage ordering.",
        metric_values=compiled_position or {"slide_identity": metrics.slide_identity},
    )
    deck_mode_mismatches = metrics.slide_identity.get("per_slide_deck_mode_mismatches", [])
    add_rule(
        "PER_SLIDE_DECK_MODE_MATCH",
        not deck_mode_mismatches,
        "Per-slide deck mode agrees across blueprint, ledger, and linkage."
        if not deck_mode_mismatches
        else f"Per-slide deck mode mismatches found on {len(deck_mode_mismatches)} slide(s).",
        affected_slides=[entry["slide_number"] for entry in deck_mode_mismatches if isinstance(entry.get("slide_number"), int)],
        metric_values={"per_slide_deck_mode_mismatches": deck_mode_mismatches},
    )
    compiled_metric_values = {
        "compiled_pptx": {
            "total": metrics.compiled_pptx_total_slide_count,
            "main_story": metrics.compiled_pptx_main_story_slide_count,
            "appendix": metrics.compiled_pptx_appendix_slide_count,
            "path_source": metrics.compiled_pptx_path_source,
            "deck_mode_derivation": metrics.compiled_pptx_deck_mode_derivation,
            "deck_mode_derivation_reason": metrics.compiled_pptx_deck_mode_derivation_reason,
        },
        "slide_build_linkage": count_values("slide_build_linkage"),
    }
    compiled_total_match = (
        metrics.compiled_pptx_total_slide_count is not None
        and metrics.slide_build_linkage_total_slide_count is not None
        and metrics.compiled_pptx_total_slide_count == metrics.slide_build_linkage_total_slide_count
    )
    compiled_mode_counts_available = (
        metrics.compiled_pptx_main_story_slide_count is not None
        and metrics.compiled_pptx_appendix_slide_count is not None
        and metrics.slide_build_linkage_main_story_slide_count is not None
        and metrics.slide_build_linkage_appendix_slide_count is not None
    )
    compiled_mode_match = (
        compiled_mode_counts_available
        and metrics.compiled_pptx_main_story_slide_count == metrics.slide_build_linkage_main_story_slide_count
        and metrics.compiled_pptx_appendix_slide_count == metrics.slide_build_linkage_appendix_slide_count
    )
    compiled_rule_passed = compiled_total_match and (compiled_mode_match if compiled_mode_counts_available else True)
    if metrics.compiled_pptx_total_slide_count is None or metrics.slide_build_linkage_total_slide_count is None:
        compiled_rule_message = "Compiled PPTX slide count could not be reconciled against slide-build linkage."
        compiled_rule_passed = False
    elif compiled_mode_counts_available:
        compiled_rule_message = (
            "Compiled PPTX total and deck-mode slide counts match slide-build linkage."
            if compiled_rule_passed
            else "Compiled PPTX total or deck-mode slide counts disagree with slide-build linkage."
        )
    else:
        compiled_rule_message = (
            "Compiled PPTX total slide count matches slide-build linkage; deck-mode counts were not safely derivable."
            if compiled_total_match
            else "Compiled PPTX total slide count disagrees with slide-build linkage."
        )
    add_rule(
        "COMPILED_SLIDE_COUNT_MATCH",
        compiled_rule_passed,
        compiled_rule_message,
        metric_values=compiled_metric_values,
    )
    add_rule(
        "COMPILE_WARNINGS_ZERO",
        metrics.compile_warning_count <= thresholds.compile_warnings_max,
        "Compile warnings are zero."
        if metrics.compile_warning_count <= thresholds.compile_warnings_max
        else f"Compile warnings must be 0, found {metrics.compile_warning_count}.",
        metric_values={"compile_warning_count": metrics.compile_warning_count},
    )
    add_rule(
        "FORBIDDEN_VISIBLE_TEXT",
        metrics.forbidden_visible_text_count <= thresholds.forbidden_visible_text_max,
        "No forbidden visible internal/helper text was detected."
        if metrics.forbidden_visible_text_count <= thresholds.forbidden_visible_text_max
        else f"Forbidden visible text findings must be 0, found {metrics.forbidden_visible_text_count}.",
        metric_values={
            "forbidden_visible_text_count": metrics.forbidden_visible_text_count,
            "pptx_scan_forbidden_visible_text_count": metrics.pptx_scan_forbidden_visible_text_count,
            "validation_summary_forbidden_visible_text_count": metrics.validation_summary_forbidden_visible_text_count,
            "compiled_pptx_path_source": metrics.compiled_pptx_path_source,
        },
    )
    add_rule(
        "FORBIDDEN_VISIBLE_TEXT_ARTIFACT_MATCH",
        metrics.forbidden_visible_text_artifact_match is not False,
        "Validation summary and direct PPTX scan agree on forbidden visible text findings."
        if metrics.forbidden_visible_text_artifact_match is not False
        else "Validation summary and direct PPTX scan disagree on forbidden visible text findings.",
        metric_values={
            "pptx_scan_forbidden_visible_text_count": metrics.pptx_scan_forbidden_visible_text_count,
            "validation_summary_forbidden_visible_text_count": metrics.validation_summary_forbidden_visible_text_count,
            "forbidden_visible_text_artifact_match": metrics.forbidden_visible_text_artifact_match,
        },
    )
    add_rule(
        "QA_BLOCKING_ZERO",
        metrics.qa_blocking_count <= thresholds.qa_blocking_max,
        "QA blocking findings are zero."
        if metrics.qa_blocking_count <= thresholds.qa_blocking_max
        else f"QA blocking findings must be 0, found {metrics.qa_blocking_count}.",
        metric_values={"qa_blocking_count": metrics.qa_blocking_count},
    )
    add_rule(
        "LECTURE_MAIN_BUDGET",
        thresholds.main_story_min <= metrics.main_story_slide_count <= thresholds.main_story_max,
        "Main-story slide count is inside the approved lecture band."
        if thresholds.main_story_min <= metrics.main_story_slide_count <= thresholds.main_story_max
        else f"Main-story slide count must stay within {thresholds.main_story_min}-{thresholds.main_story_max}, found {metrics.main_story_slide_count}.",
        metric_values={"main_story_slide_count": metrics.main_story_slide_count},
    )
    add_rule(
        "LECTURE_APPENDIX_BUDGET",
        thresholds.appendix_min <= metrics.appendix_slide_count <= thresholds.appendix_max,
        "Appendix slide count is inside the approved lecture band."
        if thresholds.appendix_min <= metrics.appendix_slide_count <= thresholds.appendix_max
        else f"Appendix slide count must stay within {thresholds.appendix_min}-{thresholds.appendix_max}, found {metrics.appendix_slide_count}.",
        metric_values={"appendix_slide_count": metrics.appendix_slide_count},
    )
    add_rule(
        "APPENDIX_BOUNDARY",
        metrics.appendix_boundary_violation_count == 0,
        "Main-story and appendix boundary is clean."
        if metrics.appendix_boundary_violation_count == 0
        else f"Appendix boundary violations found: {metrics.appendix_boundary_violation_count}.",
        metric_values={"appendix_boundary_violation_count": metrics.appendix_boundary_violation_count},
    )
    add_rule(
        "LAYOUT_COMPATIBILITY",
        metrics.layout_compatibility_failure_count == 0,
        "Layout/content compatibility holds for the compiled lecture."
        if metrics.layout_compatibility_failure_count == 0
        else f"Layout/content compatibility failed on {metrics.layout_compatibility_failure_count} slides.",
        metric_values={"layout_compatibility_failure_count": metrics.layout_compatibility_failure_count},
    )
    add_rule(
        "DETERMINISTIC_FALLBACK",
        metrics.deterministic_fallback_violation_count == 0,
        "Deterministic fallback behavior stayed inside the approved lecture layout set."
        if metrics.deterministic_fallback_violation_count == 0
        else f"Deterministic fallback violations found: {metrics.deterministic_fallback_violation_count}.",
        metric_values={"deterministic_fallback_violation_count": metrics.deterministic_fallback_violation_count},
    )
    add_rule(
        "NATIVE_VISUAL_CLASSIFICATION",
        metrics.native_visual_misclassification_count == 0,
        "Native slide visuals were not misclassified as missing."
        if metrics.native_visual_misclassification_count == 0
        else f"Native visual misclassification findings found: {metrics.native_visual_misclassification_count}.",
        metric_values={
            "native_visual_misclassification_count": metrics.native_visual_misclassification_count,
            "native_visual_classifier_mode": metrics.native_visual_classifier_mode,
        },
    )
    add_rule(
        "COMPILED_DECK_VISUAL_MONOTONY",
        metrics.compiled_deck_repeated_geometry_count == 0,
        "Compiled deck geometry is not monotonous."
        if metrics.compiled_deck_repeated_geometry_count == 0
        else f"Compiled deck still repeats realized geometry {metrics.compiled_deck_repeated_geometry_count} time(s).",
        affected_slides=metrics.compiled_deck_issue_slides.get("repeated_geometry_slides", []),
        metric_values={
            "compiled_deck_repeated_geometry_count": metrics.compiled_deck_repeated_geometry_count,
            "compiled_deck_repeated_geometry_rate": metrics.compiled_deck_repeated_geometry_rate,
        },
    )
    add_rule(
        "COMPILED_DECK_TEXT_CARD_OVERUSE",
        metrics.compiled_deck_text_card_overuse_rate <= 0.4,
        "Compiled deck does not overuse text-card composition."
        if metrics.compiled_deck_text_card_overuse_rate <= 0.4
        else f"Compiled deck text-card overuse rate is {metrics.compiled_deck_text_card_overuse_rate:.3f}.",
        affected_slides=metrics.compiled_deck_issue_slides.get("text_card_like_slides", []),
        metric_values={
            "compiled_deck_text_card_overuse_rate": metrics.compiled_deck_text_card_overuse_rate,
        },
    )
    add_rule(
        "COMPILED_DECK_CHROME_DOMINANCE",
        metrics.compiled_deck_chrome_dominance_rate == 0.0,
        "Compiled deck is not chrome-dominant."
        if metrics.compiled_deck_chrome_dominance_rate == 0.0
        else f"Compiled deck chrome-dominance rate is {metrics.compiled_deck_chrome_dominance_rate:.3f}.",
        affected_slides=metrics.compiled_deck_issue_slides.get("chrome_dominance_slides", []),
        metric_values={
            "compiled_deck_chrome_dominance_rate": metrics.compiled_deck_chrome_dominance_rate,
        },
    )
    add_rule(
        "COMPILED_DECK_MISSING_VISUAL_CENTER",
        metrics.compiled_deck_missing_visual_center_count == 0,
        "Compiled deck keeps a clear visual center on every slide."
        if metrics.compiled_deck_missing_visual_center_count == 0
        else f"Compiled deck has {metrics.compiled_deck_missing_visual_center_count} slide(s) without a clear visual center.",
        affected_slides=metrics.compiled_deck_issue_slides.get("missing_visual_center_slides", []),
        metric_values={
            "compiled_deck_missing_visual_center_count": metrics.compiled_deck_missing_visual_center_count,
        },
    )
    add_rule(
        "COMPILED_DECK_CHROME_PRESENT",
        metrics.compiled_deck_chrome_block_count == 0,
        "Compiled deck has no generic chrome blocks."
        if metrics.compiled_deck_chrome_block_count == 0
        else f"Compiled deck still contains {metrics.compiled_deck_chrome_block_count} chrome block(s).",
        affected_slides=metrics.compiled_deck_issue_slides.get("chrome_slides", []),
        metric_values={"compiled_deck_chrome_block_count": metrics.compiled_deck_chrome_block_count},
    )
    add_rule(
        "COMPILED_DECK_TITLE_BODY_DUPLICATION",
        metrics.compiled_deck_title_body_duplication_count == 0,
        "Compiled deck has no title/body duplication findings."
        if metrics.compiled_deck_title_body_duplication_count == 0
        else f"Compiled deck still has {metrics.compiled_deck_title_body_duplication_count} title/body duplication slide(s).",
        affected_slides=metrics.compiled_deck_issue_slides.get("title_body_duplication_slides", []),
        metric_values={"compiled_deck_title_body_duplication_count": metrics.compiled_deck_title_body_duplication_count},
    )
    add_rule(
        "COMPILED_DECK_ARCHETYPE_REALIZATION_MISMATCH",
        metrics.compiled_deck_archetype_realization_mismatch_count == 0,
        "Compiled deck realizes the expected archetypes."
        if metrics.compiled_deck_archetype_realization_mismatch_count == 0
        else f"Compiled deck still has {metrics.compiled_deck_archetype_realization_mismatch_count} archetype realization mismatch slide(s).",
        affected_slides=metrics.compiled_deck_issue_slides.get("archetype_realization_mismatch_slides", []),
        metric_values={
            "compiled_deck_archetype_realization_mismatch_count": metrics.compiled_deck_archetype_realization_mismatch_count,
            "compiled_deck_expected_visual_missing_count": metrics.compiled_deck_expected_visual_missing_count,
        },
    )
    add_rule(
        "COMPILED_DECK_VISUAL_NEEDED_BUT_MISSING",
        metrics.compiled_deck_expected_visual_missing_count == 0,
        "Compiled deck realizes the required mapping, process, and worked-example structures."
        if metrics.compiled_deck_expected_visual_missing_count == 0
        else f"Compiled deck still omits required visual structure on {metrics.compiled_deck_expected_visual_missing_count} slide(s).",
        affected_slides=metrics.compiled_deck_issue_slides.get("expected_visual_missing_slides", []),
        metric_values={
            "compiled_deck_expected_visual_missing_count": metrics.compiled_deck_expected_visual_missing_count,
        },
    )
    add_rule(
        "COMPILED_DECK_APPENDIX_CLONE_RUN",
        metrics.compiled_deck_appendix_clone_count == 0,
        "Compiled appendix has no clone run."
        if metrics.compiled_deck_appendix_clone_count == 0
        else f"Compiled appendix still has clone-run overflow count {metrics.compiled_deck_appendix_clone_count}.",
        affected_slides=metrics.compiled_deck_issue_slides.get("appendix_clone_slides", []),
        metric_values={"compiled_deck_appendix_clone_count": metrics.compiled_deck_appendix_clone_count},
    )
    add_rule(
        "COMPILED_DECK_APPENDIX_VISUAL_CLONE",
        metrics.compiled_deck_appendix_visual_clone_run_length <= 2,
        "Compiled appendix does not run one visual geometry for too many slides in sequence."
        if metrics.compiled_deck_appendix_visual_clone_run_length <= 2
        else f"Compiled appendix visual clone run length is {metrics.compiled_deck_appendix_visual_clone_run_length}.",
        affected_slides=metrics.compiled_deck_issue_slides.get("appendix_visual_clone_slides", []),
        metric_values={
            "compiled_deck_appendix_visual_clone_run_length": metrics.compiled_deck_appendix_visual_clone_run_length,
        },
    )
    add_rule(
        "COMPILED_DECK_REPETITIVE_MOTION",
        metrics.compiled_deck_repetitive_motion_count == 0,
        "Compiled deck motion is not repetitive."
        if metrics.compiled_deck_repetitive_motion_count == 0
        else f"Compiled deck repetitive-motion count is {metrics.compiled_deck_repetitive_motion_count}.",
        affected_slides=metrics.compiled_deck_issue_slides.get("repetitive_motion_slides", []),
        metric_values={
            "compiled_deck_repetitive_motion_count": metrics.compiled_deck_repetitive_motion_count,
            "compiled_deck_repeated_archetype_count": metrics.compiled_deck_repeated_archetype_count,
            "compiled_deck_repeated_title_stem_count": metrics.compiled_deck_repeated_title_stem_count,
        },
    )
    add_rule(
        "COMPILED_DECK_TRUTH_MATCH",
        metrics.compiled_deck_truth_mismatch_count == 0,
        "Planned metrics and compiled-deck truth agree."
        if metrics.compiled_deck_truth_mismatch_count == 0
        else f"Planned metrics disagree with compiled-deck truth on {metrics.compiled_deck_truth_mismatch_count} item(s).",
        metric_values={"compiled_deck_truth_mismatches": metrics.compiled_deck_truth_mismatches},
    )

    if profile == GateProfile.CANDIDATE:
        allowlisted = [finding for finding in metrics.residual_findings if _allowlisted(finding, policy.candidate_allowlist)]
        non_allowlisted = [finding for finding in metrics.residual_findings if not _allowlisted(finding, policy.candidate_allowlist)]
        main_story_density = [
            finding
            for finding in metrics.residual_findings
            if finding.category.strip().lower() == "density" and (finding.deck_mode or "").strip().lower() == DeckMode.MAIN_STORY.value
        ]
        non_appendix = [
            finding for finding in metrics.residual_findings if (finding.deck_mode or "").strip().lower() != DeckMode.APPENDIX.value
        ]
        qa_status_ok = metrics.qa_status == "pass" or (metrics.qa_status == "conditional-pass" and not non_allowlisted)
        add_rule(
            "QA_CANDIDATE_ALLOWLIST",
            qa_status_ok,
            "Candidate QA status is pass or a narrow allowlisted conditional-pass."
            if qa_status_ok
            else f"Candidate profile only allows pass or allowlisted conditional-pass findings, found qa_status={metrics.qa_status!r}.",
            metric_values={"qa_status": metrics.qa_status, "residual_findings": len(metrics.residual_findings)},
            profile_impact=[GateProfile.CANDIDATE],
        )
        add_rule(
            "QA_CANDIDATE_ALLOWLIST_COUNT",
            len(allowlisted) <= policy.candidate_max_allowlisted_findings,
            "Allowlisted residual findings stay within the candidate cap."
            if len(allowlisted) <= policy.candidate_max_allowlisted_findings
            else f"Candidate profile allows at most {policy.candidate_max_allowlisted_findings} allowlisted residual findings, found {len(allowlisted)}.",
            metric_values={"allowlisted_residual_findings": len(allowlisted)},
            profile_impact=[GateProfile.CANDIDATE],
        )
        add_rule(
            "QA_CANDIDATE_MAIN_STORY_DENSITY_ZERO",
            not main_story_density,
            "Candidate profile has zero main-story density warnings."
            if not main_story_density
            else f"Candidate profile does not allow main-story density warnings, found {len(main_story_density)}.",
            affected_slides=[finding.slide_number for finding in main_story_density if finding.slide_number is not None],
            metric_values={"main_story_density_warnings": len(main_story_density)},
            profile_impact=[GateProfile.CANDIDATE],
        )
        add_rule(
            "QA_CANDIDATE_NON_APPENDIX_ZERO",
            not non_appendix,
            "Candidate profile has zero non-appendix residual findings."
            if not non_appendix
            else f"Candidate profile does not allow residual findings outside appendix, found {len(non_appendix)}.",
            affected_slides=[finding.slide_number for finding in non_appendix if finding.slide_number is not None],
            metric_values={"non_appendix_residual_findings": len(non_appendix)},
            profile_impact=[GateProfile.CANDIDATE],
        )
        if allowlisted:
            add_rule(
                "QA_CANDIDATE_ALLOWLIST_NOTICE",
                True,
                f"Candidate pass uses {len(allowlisted)} allowlisted appendix-only minor density finding(s).",
                severity=GateRuleSeverity.WARNING,
                affected_slides=[finding.slide_number for finding in allowlisted if finding.slide_number is not None],
                metric_values={"allowlisted_residual_findings": len(allowlisted)},
                profile_impact=[GateProfile.CANDIDATE],
            )
    else:
        add_rule(
            "QA_RELEASE_PASS",
            metrics.qa_status == "pass",
            "Release profile requires QA status pass."
            if metrics.qa_status == "pass"
            else f"Release profile requires qa_status=pass, found {metrics.qa_status!r}.",
            metric_values={"qa_status": metrics.qa_status},
            profile_impact=[GateProfile.RELEASE],
        )
        add_rule(
            "QA_RELEASE_FINDINGS_ZERO",
            len(metrics.residual_findings) == 0,
            "Release profile has zero residual findings."
            if len(metrics.residual_findings) == 0
            else f"Release profile requires zero residual findings, found {len(metrics.residual_findings)}.",
            affected_slides=[finding.slide_number for finding in metrics.residual_findings if finding.slide_number is not None],
            metric_values={"residual_findings": len(metrics.residual_findings)},
            profile_impact=[GateProfile.RELEASE],
        )
    return findings


def _target_result_from_execution(
    *,
    profile: GateProfile,
    target: RegressionTargetSpec,
    execution: TargetBuildExecution,
    policy: ProductionGatePolicy,
    repo_root: Path,
) -> TargetGateResult:
    metrics = collect_target_metrics(target=target, execution=execution, policy=policy, repo_root=repo_root)
    rule_results = evaluate_gate_metrics(profile=profile, target=target, metrics=metrics, policy=policy)
    passed = all(
        rule.passed for rule in rule_results if rule.severity == GateRuleSeverity.ERROR and profile in rule.profile_impact
    )
    return TargetGateResult(
        target_id=target.target_id,
        output_dir=execution.output_dir,
        passed=passed,
        metrics=metrics,
        rule_results=rule_results,
        artifact_paths=execution.artifact_paths,
    )


def run_production_gate(
    *,
    profile: GateProfile,
    policy_path: str | Path = DEFAULT_POLICY_PATH,
    output_dir: str | Path | None = None,
    repo_root: str | Path = ROOT,
    material_mode: MaterialMode = MaterialMode.AUTO,
    provider_settings: ProviderSettings | None = None,
) -> ProductionGateResult:
    repo_root_path = Path(repo_root).resolve()
    policy = load_production_gate_policy(policy_path)
    resolved_provider_settings = provider_settings or ProviderSettings()
    resolved_output_dir = Path(output_dir) if output_dir is not None else DEFAULT_OUTPUT_ROOT / profile.value
    resolved_output_dir = resolved_output_dir.resolve()
    _prepare_output_dir(resolved_output_dir)

    pytest_result = _run_pytest_suite(repo_root_path, policy.pytest_paths, resolved_output_dir)
    if not pytest_result.passed:
        pytest_rule = GateRuleFinding(
            rule_id="PYTEST_SUITE",
            passed=False,
            severity=GateRuleSeverity.ERROR,
            message=f"Targeted regression pytest suite failed with return code {pytest_result.returncode}.",
            profile_impact=[GateProfile.CANDIDATE, GateProfile.RELEASE],
            metric_values={
                "returncode": pytest_result.returncode,
                "stdout_path": pytest_result.stdout_path,
                "stderr_path": pytest_result.stderr_path,
            },
        )
        result = ProductionGateResult(
            profile=profile,
            passed=False,
            policy_id=policy.policy_id,
            policy_version=policy.policy_version,
            material_mode=material_mode,
            pytest_result=pytest_result,
            targets=[],
            failed_rules=[pytest_rule],
            warnings=[],
            artifact_paths={
                "output_dir": str(resolved_output_dir),
                "policy_path": str(Path(policy_path).resolve()),
                "pytest_result": str(resolved_output_dir / "pytest-result.json"),
            },
        )
        _write_json(resolved_output_dir / "gate-result.json", result.model_dump(mode="json", exclude_none=True))
        return result

    target_results: list[TargetGateResult] = []
    for target in policy.regression_targets:
        execution = _run_regression_target(
            repo_root=repo_root_path,
            target=target,
            gate_output_root=resolved_output_dir,
            material_mode=material_mode,
            policy=policy,
            provider_settings=resolved_provider_settings,
        )
        target_results.append(
            _target_result_from_execution(
                profile=profile,
                target=target,
                execution=execution,
                policy=policy,
                repo_root=repo_root_path,
            )
        )

    failed_rules = [
        rule
        for target_result in target_results
        for rule in target_result.rule_results
        if rule.severity == GateRuleSeverity.ERROR and not rule.passed and profile in rule.profile_impact
    ]
    warnings = [
        rule
        for target_result in target_results
        for rule in target_result.rule_results
        if rule.severity == GateRuleSeverity.WARNING and profile in rule.profile_impact
    ]
    result = ProductionGateResult(
        profile=profile,
        passed=not failed_rules and all(target_result.passed for target_result in target_results),
        policy_id=policy.policy_id,
        policy_version=policy.policy_version,
        material_mode=material_mode,
        pytest_result=pytest_result,
        targets=target_results,
        failed_rules=failed_rules,
        warnings=warnings,
        artifact_paths={
            "output_dir": str(resolved_output_dir),
            "policy_path": str(Path(policy_path).resolve()),
            "pytest_result": str(resolved_output_dir / "pytest-result.json"),
        },
    )
    _write_json(resolved_output_dir / "gate-result.json", result.model_dump(mode="json", exclude_none=True))
    return result


def _print_gate_result(result: ProductionGateResult) -> None:
    status = "PASS" if result.passed else "FAIL"
    print(f"[production-gate] profile={result.profile.value} status={status}")
    print(
        "[production-gate] pytest="
        f"{'PASS' if result.pytest_result.passed else 'FAIL'} "
        f"duration={result.pytest_result.duration_seconds:.3f}s "
        f"stdout={result.pytest_result.stdout_path}"
    )
    for target in result.targets:
        print(
            "[production-gate] target="
            f"{target.target_id} status={'PASS' if target.passed else 'FAIL'} "
            f"main={target.metrics.main_story_slide_count} "
            f"appendix={target.metrics.appendix_slide_count} "
            f"qa={target.metrics.qa_status} "
            f"compile_warnings={target.metrics.compile_warning_count} "
            f"forbidden_text={target.metrics.forbidden_visible_text_count}"
        )
        for rule in target.rule_results:
            if rule.severity == GateRuleSeverity.WARNING and rule.passed and result.profile in rule.profile_impact:
                print(f"[production-gate] warning {rule.rule_id}: {rule.message}")
    for rule in result.failed_rules:
        print(f"[production-gate] failed {rule.rule_id}: {rule.message}")
    print(f"[production-gate] gate-result={Path(result.artifact_paths['output_dir']) / 'gate-result.json'}")


def main(argv: list[str] | None = None) -> int:
    args = _parse_args(argv)
    result = run_production_gate(
        profile=GateProfile(args.profile),
        policy_path=args.policy,
        output_dir=args.output_dir,
        repo_root=args.repo_root,
        material_mode=MaterialMode(args.material_mode),
        provider_settings=_provider_settings_from_args(args),
    )
    _print_gate_result(result)
    return 0 if result.passed else 1


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
