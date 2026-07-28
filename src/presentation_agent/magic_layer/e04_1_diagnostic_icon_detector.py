"""Detect diagnostic top-right SVG icon clusters in E04 decks."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

EMU_PER_IN = 914400


def is_diagnostic_cluster_bbox(x: float, y: float, w: float, h: float) -> bool:
    return x >= 9.8 and y <= 0.55 and 0.09 <= w <= 0.30 and 0.09 <= h <= 0.30


def detect_diagnostic_icon_leakage(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    rows: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        cluster_count = 0
        for shape in slide.shapes:
            name = getattr(shape, "name", "")
            if "SVG Icon" not in name and not name.startswith("icon::"):
                continue
            x, y, w, h = [value / EMU_PER_IN for value in (shape.left, shape.top, shape.width, shape.height)]
            diagnostic = is_diagnostic_cluster_bbox(x, y, w, h)
            if diagnostic:
                cluster_count += 1
            rows.append(
                {
                    "slide_number": slide_idx,
                    "shape_id": shape.shape_id,
                    "object_name": name,
                    "bbox_in": [round(x, 3), round(y, 3), round(w, 3), round(h, 3)],
                    "diagnostic_leakage": diagnostic,
                    "reason": "top_right_uniform_svg_cluster" if diagnostic else None,
                }
            )
        if cluster_count >= 2:
            for row in rows:
                if row["slide_number"] == slide_idx and row["diagnostic_leakage"]:
                    row["cluster_pattern_detected"] = True
    leakage = [row for row in rows if row.get("diagnostic_leakage")]
    return {
        "schema_name": "diagnostic_icon_leakage_report",
        "status": "passed" if not leakage else "failed",
        "original_svg_icon_count": len(rows),
        "diagnostic_icon_count": len(leakage),
        "diagnostic_leakage_count": len(leakage),
        "rows": rows,
    }
