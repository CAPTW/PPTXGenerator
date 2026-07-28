"""Deterministic PPTX compiler for approved blueprint, asset, and visual state."""

from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from statistics import median
from types import SimpleNamespace
from typing import Any

import cairosvg
import yaml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import Field, field_validator

from .compat.presentation_contracts import WorkflowGate
from .compat.presentation_contracts import (
    AssetKind,
    AssetManifest,
    AssetRecord,
    AssetStatus,
    BatchManifest,
    Blueprint,
    BlueprintSlide,
    ContractModel,
    DeckConstitution,
    DeckMode,
    DesignSystem,
    LayoutLibrary,
    LayoutPattern,
    ProductionMode,
    QAStatus,
    SchemaModel,
    SlideArchetype,
    SlideEvidenceClass,
    SlideLedger,
    SlideLedgerEntry,
    SlideRole,
    StageStatus,
    StateCapsule,
    VisualSourcePreference,
    VisualType,
    VizManifest,
    VizRecord,
)
from .compat.state_io import load_state_file, save_state_file
from .slide_ir import (
    IRRect,
    SlideIRCompileContext,
    SlideIRCompileReport,
    SlideIRBackgroundLayer,
    SlideIRLayoutCandidateScore,
    SlideIRLayoutCriticReport,
    SlideIRComponentStyleReference,
    SlideIRContinuityAnchor,
    SlideIRContinuityFinding,
    SlideIRContinuityReport,
    SlideIRDeckContinuityMetadata,
    SlideIRDeckIdentity,
    SlideIRDensityBudgetReference,
    SlideIRDocument,
    SlideIRGeometryObjectSummary,
    SlideIRGeometrySlideSummary,
    SlideIRObject,
    SlideIRPreviewMetadata,
    SlideIRSafeAreaMask,
    SlideIRSlide,
    SlideIRSpacingRhythmReference,
    SlideIRStylePrior,
    SlideIRSectionStructure,
    SlideIRThemeIdentity,
    SlideIRTypographyScaleReference,
    SlideIRValidationFinding,
    SlideIRValidationReport,
)
from .style_prior import DEFAULT_STYLE_PRIOR_PROVIDER, StylePriorProvider, build_style_prior_context


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
DEFAULT_MARGIN_IN = 0.55
DEFAULT_GAP_IN = 0.24
LAYOUT_CRITIC_REPORT_KEY = "layout_critic_report"
MAX_LAYOUT_CANDIDATES = 3
FORBIDDEN_VISIBLE_TEXT_PATTERNS: tuple[tuple[re.Pattern[str], str], ...] = (
    (re.compile(r"\bprimary visual form\b", re.IGNORECASE), "internal planner text"),
    (re.compile(r"\banchor the slide in local material\b", re.IGNORECASE), "internal planner text"),
    (re.compile(r"\bproduction inputs\b", re.IGNORECASE), "internal planner text"),
    (re.compile(r"\breserve the lower band\b", re.IGNORECASE), "layout helper text"),
    (re.compile(r"\buse the lower-right zone\b", re.IGNORECASE), "layout helper text"),
    (re.compile(r"\bno approved\b", re.IGNORECASE), "fallback debug text"),
    (re.compile(r"\bblueprint slide was missing\b", re.IGNORECASE), "compiler placeholder text"),
    (re.compile(r"\breconcile the blueprint\b", re.IGNORECASE), "compiler placeholder text"),
    (re.compile(r"\bnot approved layout\b", re.IGNORECASE), "layout debug text"),
    (re.compile(r"\bdensity audit\b", re.IGNORECASE), "QA leakage"),
    (re.compile(r"\bvisual safety policy\b", re.IGNORECASE), "placeholder label"),
    (re.compile(r"\bplaceholder helper text\b", re.IGNORECASE), "helper copy leakage"),
)


class CompileContractError(ValueError):
    """Raised when approved Gate 2 artifacts do not satisfy compile-time contract checks."""


class SlideBuildLink(ContractModel):
    slide_number: int
    slide_id: str
    title: str
    pptx_index: int
    deck_mode: DeckMode
    batch_id: str | None = None
    layout_pattern_id: str
    layout_family: str
    visual_type: VisualType
    numbering_label: str
    asset_ids: list[str] = Field(default_factory=list)
    viz_spec_ids: list[str] = Field(default_factory=list)
    linked_paths: list[str] = Field(default_factory=list)
    notes_present: bool = False
    compile_status: StageStatus = StageStatus.DRAFT
    missing_dependencies: list[str] = Field(default_factory=list)
    qa_status: QAStatus | None = None
    qa_warning_count: int = 0
    qa_blocking_count: int = 0
    qa_finding_ids: list[str] = Field(default_factory=list)
    qa_notes: list[str] = Field(default_factory=list)
    remediation_status: StageStatus | None = None
    remediation_finding_ids: list[str] = Field(default_factory=list)
    remediation_batch_ids: list[str] = Field(default_factory=list)
    continuation_notes: list[str] = Field(default_factory=list)


class SlideBuildLinkage(SchemaModel):
    SCHEMA_NAME = "slide_build_linkage"
    SUMMARY = "Per-slide linkage from compiled PPTX slides back to approved assets and structured visuals."

    deck_title: str
    pptx_path: str
    slides: list[SlideBuildLink] = Field(default_factory=list)


class BuildManifest(SchemaModel):
    SCHEMA_NAME = "build_manifest"
    SUMMARY = "PPTX build summary, warnings, and output linkage for one compilation run."

    deck_title: str
    pptx_path: str
    slide_ratio: str
    slide_count: int
    compiled_layout_patterns: list[str] = Field(default_factory=list)
    warnings: list[str] = Field(default_factory=list)
    linkage_path: str
    batch_ids: list[str] = Field(default_factory=list)

    @field_validator("slide_count")
    @classmethod
    def _validate_slide_count(cls, value: int) -> int:
        if value < 1:
            raise ValueError("slide_count must be at least 1")
        return value


class PptxCompileOutputs(ContractModel):
    slide_ir: SlideIRDocument
    compile_report: SlideIRCompileReport
    build_manifest: BuildManifest
    slide_build_linkage: SlideBuildLinkage
    slide_ledger: SlideLedger
    layout_critic_report: SlideIRLayoutCriticReport | None = None
    batch_manifest: BatchManifest | None = None
    state_capsule: StateCapsule | None = None
    pptx_path: Path


class ResolvedVisual(ContractModel):
    record: VizRecord
    output_path: str | None = None
    data_output_path: str | None = None


class ResolvedResources(ContractModel):
    assets: list[AssetRecord] = Field(default_factory=list)
    visuals: list[ResolvedVisual] = Field(default_factory=list)
    missing_dependencies: list[str] = Field(default_factory=list)


@dataclass(slots=True)
class CompilerStyle:
    slide_width: float
    slide_height: float
    margin: float
    gap: float
    colors: dict[str, RGBColor]
    fonts: dict[str, Any]
    root: Path
    raster_dir: Path
    style_prior: SlideIRStylePrior | None = None


@dataclass(slots=True)
class AdaptedSlideSpec:
    slide: BlueprintSlide
    entry: SlideLedgerEntry
    layout_warnings: list[str]
    candidate_families: tuple[str, ...]


@dataclass(slots=True)
class ScoredSlideCandidate:
    slide: SlideIRSlide
    score: SlideIRLayoutCandidateScore


_COMPILE_REVALIDATED_SLIDE_DERIVED_FIELDS: tuple[str, ...] = (
    "layout_slot_map",
    "density_budget",
    "risk_flags",
    "qa_acceptance_hints",
    "visual_intent",
    "supporting_evidence",
)


def _display_path(path: Path, root: Path) -> str:
    try:
        return str(path.relative_to(root))
    except ValueError:
        return str(path)


def _resolve_path(path_text: str | None, root: Path) -> Path | None:
    if not path_text:
        return None
    candidate = Path(path_text)
    if candidate.is_absolute():
        return candidate
    return (root / candidate).resolve()


def _compile_ready_blueprint_slide(slide: BlueprintSlide) -> BlueprintSlide:
    """Rebuild legacy alias slides so compile-time derived fields match the canonical layout id."""
    if slide.layout_pattern_id not in _LEGACY_LAYOUT_PATTERN_ALIASES:
        return slide
    payload = slide.model_dump(mode="json", exclude_none=True)
    for field_name in _COMPILE_REVALIDATED_SLIDE_DERIVED_FIELDS:
        payload.pop(field_name, None)
    return slide.__class__.model_validate(payload)


def _load_notes(path: str | Path | None) -> dict[str, str]:
    if path is None:
        return {}
    notes_path = Path(path)
    if not notes_path.is_file():
        raise FileNotFoundError(notes_path)
    text = notes_path.read_text(encoding="utf-8")
    if notes_path.suffix.lower() in {".yaml", ".yml"}:
        payload = yaml.safe_load(text)
    else:
        payload = json.loads(text)
    if payload is None:
        return {}
    if not isinstance(payload, dict):
        raise ValueError("notes file must contain a top-level object keyed by slide_id or slide number")
    return {str(key): str(value) for key, value in payload.items() if value}


def _normalize_text(text: str | None) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _ir_rect(left: float, top: float, width: float, height: float) -> IRRect:
    return IRRect(left=left, top=top, width=width, height=height)


def _preview_object_id(*, slide_number: int, slot: str, kind: str) -> str:
    safe_slot = slot.replace(" ", "-").lower()
    return f"s{slide_number:03d}:{safe_slot}:{kind}"


def _slide_width_height_from_ratio(slide_ratio: str) -> tuple[float, float]:
    width_value, height_value = _parse_ratio(slide_ratio)
    return width_value, height_value


def _build_preview_objects(slide: BlueprintSlide, entry: SlideLedgerEntry, style: CompilerStyle, family: str) -> list[SlideIRObject]:
    has_visual = slide.visual_type not in {VisualType.TEXT, VisualType.QUOTE}
    has_primary_visual_slot = "primary_visual" in set(slide.layout_slot_map.values())
    use_two_column_layout = has_visual and (
        has_primary_visual_slot or family in {"comparison", "process-flow", "worked-example", "appendix-reference"}
    )

    objects: list[SlideIRObject] = []
    title_box = _ir_rect(
        left=style.margin,
        top=0.46,
        width=style.slide_width - (style.margin * 2),
        height=0.62,
    )
    if slide.title:
        objects.append(
            SlideIRObject(
                object_id=_preview_object_id(slide_number=slide.slide_number, slot="title", kind="text"),
                slot="title",
                kind="text",
                bounds=title_box,
                text=_normalize_text(slide.title),
                source="blueprint.title",
            )
        )

    claim_top = title_box.bottom + style.gap
    claim_box = _ir_rect(
        left=style.margin,
        top=claim_top,
        width=style.slide_width - (style.margin * 2),
        height=0.72,
    )
    if slide.one_line_takeaway:
        objects.append(
            SlideIRObject(
                object_id=_preview_object_id(slide_number=slide.slide_number, slot="claim", kind="text"),
                slot="claim",
                kind="text",
                bounds=claim_box,
                text=_normalize_text(slide.one_line_takeaway),
                source="blueprint.one_line_takeaway",
            )
        )

    body_top = claim_box.bottom + style.gap
    body_width = style.slide_width - (style.margin * 2)
    body_left = style.margin
    visual_left = body_left
    visual_top = body_top
    visual_height = 3.2
    if use_two_column_layout:
        visual_left = (style.slide_width / 2.0) + (style.margin / 2.0)
        visual_width = body_width - (visual_left - style.margin)
        text_width = visual_left - style.margin - style.gap
        body_width = text_width
    else:
        visual_width = 0.0
        text_width = body_width

    main_message_box = _ir_rect(
        left=body_left,
        top=body_top,
        width=body_width,
        height=1.55,
    )
    if slide.main_message:
        objects.append(
            SlideIRObject(
                object_id=_preview_object_id(slide_number=slide.slide_number, slot="main-message", kind="text"),
                slot="main-message",
                kind="text",
                bounds=main_message_box,
                text=_normalize_text(slide.main_message),
                source="blueprint.main_message",
            )
        )

    if has_visual or has_primary_visual_slot:
        visual_top = main_message_box.bottom + style.gap
        visual_box = _ir_rect(
            left=visual_left,
            top=visual_top,
            width=visual_width or (style.slide_width - (style.margin * 2)),
            height=2.95,
        )
        objects.append(
            SlideIRObject(
                object_id=_preview_object_id(slide_number=slide.slide_number, slot="primary-visual", kind="visual"),
                slot="primary-visual",
                kind="visual",
                bounds=visual_box,
                text=None,
                source="slide.visual_type",
                payload={"visual_type": slide.visual_type.value, "layout_family": family},
            )
        )

    for index, evidence_text in enumerate(slide.supporting_evidence[:2]):
        evidence_top = (main_message_box.bottom + 0.2) + (index * 0.88)
        evidence_box = _ir_rect(
            left=body_left,
            top=evidence_top,
            width=body_width,
            height=0.58,
        )
        objects.append(
            SlideIRObject(
                object_id=_preview_object_id(slide_number=slide.slide_number, slot=f"evidence-{index + 1}", kind="evidence"),
                slot=f"evidence-{index + 1}",
                kind="evidence",
                bounds=evidence_box,
                text=_normalize_text(evidence_text),
                source="blueprint.supporting_evidence",
            )
        )

    footer_top = style.slide_height - 0.4
    footer_box = _ir_rect(left=style.margin, top=footer_top, width=0.95, height=0.28)
    objects.append(
        SlideIRObject(
            object_id=_preview_object_id(slide_number=slide.slide_number, slot="footer", kind="annotation"),
            slot="footer",
            kind="annotation",
            bounds=footer_box,
            text=f"p{slide.slide_number}",
            source="compiler.footer",
            payload={"deck_mode": entry.deck_mode.value},
        )
    )

    return objects


def _preview_style_for_slide_ir(
    *,
    slide_width_in: float,
    slide_height_in: float,
    design_system: DesignSystem | None,
) -> CompilerStyle:
    spacing_scale = design_system.spacing_scale if design_system is not None else []
    return CompilerStyle(
        slide_width=slide_width_in,
        slide_height=slide_height_in,
        margin=max(DEFAULT_MARGIN_IN, _spacing(spacing_scale, 4, DEFAULT_MARGIN_IN)),
        gap=max(DEFAULT_GAP_IN, _spacing(spacing_scale, 3, DEFAULT_GAP_IN)),
        colors={},
        fonts={},
        root=Path.cwd(),
        raster_dir=Path.cwd(),
        style_prior=None,
    )


def _slide_ir_structural_fingerprint(slide_ir: SlideIRDocument) -> list[tuple[Any, ...]]:
    fingerprint: list[tuple[Any, ...]] = []
    for slide in slide_ir.slides:
        fingerprint.append(
            (
                slide.slide_number,
                slide.slide_id,
                slide.layout_family,
                tuple(
                    (
                        obj.object_id,
                        obj.slot,
                        obj.kind,
                        round(obj.bounds.left, 4),
                        round(obj.bounds.top, 4),
                        round(obj.bounds.width, 4),
                        round(obj.bounds.height, 4),
                    )
                    for obj in slide.objects
                ),
            )
        )
    return fingerprint


def _attach_style_prior(
    slide_ir: SlideIRDocument,
    *,
    design_system: DesignSystem | None,
    canonical_generation_profile,
    style_prior_provider: StylePriorProvider | None,
) -> SlideIRDocument:
    provider = style_prior_provider or DEFAULT_STYLE_PRIOR_PROVIDER
    style_prior = provider.build(
        build_style_prior_context(
            slide_ir=slide_ir,
            design_system=design_system,
            canonical_generation_profile=canonical_generation_profile,
        )
    )
    styled_slide_ir = slide_ir.model_copy(update={"style_prior": style_prior})
    if _slide_ir_structural_fingerprint(slide_ir) != _slide_ir_structural_fingerprint(styled_slide_ir):
        raise ValueError("Style-prior integration must not alter structural SlideIR geometry.")
    return styled_slide_ir


def _build_slide_ir_slide(
    *,
    slide: BlueprintSlide,
    entry: SlideLedgerEntry,
    style: CompilerStyle,
    family: str,
    layout_warnings: list[str],
    appendix_start: int | None,
) -> SlideIRSlide:
    objects = _build_preview_objects(slide, entry, style, family)
    return SlideIRSlide(
        slide_number=slide.slide_number,
        slide_id=entry.slide_id,
        title=_normalize_text(slide.title),
        section=slide.section,
        one_line_takeaway=_normalize_text(slide.one_line_takeaway),
        main_message=_normalize_text(slide.main_message),
        deck_mode=slide.deck_mode,
        slide_role=slide.slide_role,
        visual_type=slide.visual_type,
        slide_archetype=slide.slide_archetype,
        layout_pattern_id=entry.layout_pattern_id,
        layout_family=family,
        layout_warnings=layout_warnings,
        numbering_label=_build_number_label(slide, appendix_start),
        objects=objects,
        notes=entry.change_note and [entry.change_note] if entry.change_note else [],
        primary_claim=_normalize_text(slide.primary_claim),
        audience_intent=_normalize_text(slide.audience_intent),
        must_keep_text=list(slide.must_keep_text),
        optional_text=list(slide.optional_text),
        core_content=list(slide.core_content),
        supporting_evidence=list(slide.supporting_evidence),
        required_evidence_assets=list(slide.required_evidence_assets),
        layout_slot_map=dict(slide.layout_slot_map),
        density_budget=slide.density_budget,
        evidence_class=slide.evidence_class,
        presenter_notes=slide.presenter_notes,
        production_bridge=slide.production_bridge,
        authoring_payload=dict(slide.authoring_payload),
        visual_source_preference=entry.visual_source_preference,
        production_mode=entry.production_mode,
        asset_dependency_kinds=list(entry.asset_dependency_kinds),
        batch_id=entry.batch_id,
        change_note=entry.change_note,
        unresolved_blockers=list(entry.unresolved_blockers or []),
    )


def _rect_overlaps(first: IRRect, second: IRRect) -> bool:
    if first.right <= second.left or second.right <= first.left:
        return False
    if first.bottom <= second.top or second.bottom <= first.top:
        return False
    return True


def _findings_for_ir_bounds(
    slide: SlideIRSlide,
    slide_width_in: float,
    slide_height_in: float,
) -> list[SlideIRValidationFinding]:
    findings: list[SlideIRValidationFinding] = []
    for obj in slide.objects:
        if obj.bounds.left < 0 or obj.bounds.top < 0:
            findings.append(
                SlideIRValidationFinding(
                    severity="error",
                    code="out_of_bounds",
                    message="Object placement starts outside slide bounds.",
                    slide_number=slide.slide_number,
                    object_id=obj.object_id,
                    details={"bounds": obj.bounds.model_dump(mode="json")},
                )
            )
        if obj.bounds.right > slide_width_in or obj.bounds.bottom > slide_height_in:
            findings.append(
                SlideIRValidationFinding(
                    severity="error",
                    code="out_of_bounds",
                    message="Object placement exceeds slide bounds.",
                    slide_number=slide.slide_number,
                    object_id=obj.object_id,
                    details={"bounds": obj.bounds.model_dump(mode="json")},
                )
            )

    for index, left_obj in enumerate(slide.objects):
        for right_obj in slide.objects[index + 1 :]:
            if _rect_overlaps(left_obj.bounds, right_obj.bounds):
                findings.append(
                    SlideIRValidationFinding(
                        severity="warning",
                        code="overlap",
                        message="Objects overlap in preview geometry.",
                        slide_number=slide.slide_number,
                        object_id=f"{left_obj.object_id}|{right_obj.object_id}",
                        details={
                            "first_object": left_obj.object_id,
                            "second_object": right_obj.object_id,
                        },
                    )
                )

    return findings


def adapt_blueprint_to_slide_ir(
    *,
    blueprint: Blueprint,
    design_system: DesignSystem | None = None,
    deck_constitution: DeckConstitution | None = None,
    layout_library: LayoutLibrary | None = None,
    slide_ledger: SlideLedger | None = None,
    asset_manifest: AssetManifest | None = None,
    viz_manifest: VizManifest | None = None,
    enable_layout_critic: bool = True,
    style_prior_provider: StylePriorProvider | None = None,
) -> SlideIRDocument:
    """Create a stable SlideIR payload from legacy compile-time schemas.

    The adapter is the only boundary where public compile-time schemas are read directly.
    Downstream preview and compile flows consume the returned SlideIR document only.
    """
    source_ledger = slide_ledger or SlideLedger(entries=[], deck_title=blueprint.deck_title)
    layout_patterns = _layout_pattern_index(layout_library) if layout_library is not None else {}
    slide_width_in, slide_height_in = _slide_width_height_from_ratio(blueprint.slide_ratio)
    spacing_scale = design_system.spacing_scale if design_system is not None else []
    style = CompilerStyle(
        slide_width=slide_width_in,
        slide_height=slide_height_in,
        margin=max(DEFAULT_MARGIN_IN, _spacing(spacing_scale, 4, DEFAULT_MARGIN_IN)),
        gap=max(DEFAULT_GAP_IN, _spacing(spacing_scale, 3, DEFAULT_GAP_IN)),
        colors={},
        fonts={},
        root=Path.cwd(),
        raster_dir=Path.cwd(),
        style_prior=None,
    )
    slide_map = {entry.slide_number: entry for entry in source_ledger.entries}
    deck_title = blueprint.deck_title
    ir_slides: list[SlideIRSlide] = []
    slide_specs: list[AdaptedSlideSpec] = []

    for slide in sorted(blueprint.slides, key=lambda item: item.slide_number):
        entry = slide_map.get(slide.slide_number)
        if entry is None:
            continue
        compile_ready_slide = _compile_ready_blueprint_slide(slide)
        pattern = _layout_pattern_lookup(layout_patterns, entry.layout_pattern_id)
        family, layout_warnings = _resolved_layout_family(compile_ready_slide, pattern)
        candidate_families = _layout_candidate_families(compile_ready_slide, family)
        slide_specs.append(
            AdaptedSlideSpec(
                slide=compile_ready_slide,
                entry=entry,
                layout_warnings=list(layout_warnings),
                candidate_families=candidate_families,
            )
        )
        ir_slides.append(
            _build_slide_ir_slide(
                slide=compile_ready_slide,
                entry=entry,
                style=style,
                family=family,
                layout_warnings=list(layout_warnings),
                appendix_start=blueprint.appendix_start,
            )
        )

    continuity = _build_slide_ir_continuity_metadata(
        deck_title=deck_title,
        slide_ratio=blueprint.slide_ratio,
        design_system=design_system,
        slides=ir_slides,
        style=style,
    )

    # Preserve the approved ledger snapshot as-is so downstream QA can classify
    # inconsistent source state instead of this bridge rejecting it prematurely.
    compile_context = SlideIRCompileContext.model_construct(
        design_system=design_system,
        layout_patterns=list(layout_patterns.values()),
        assets=list(asset_manifest.assets) if asset_manifest is not None else [],
        visuals=list(viz_manifest.visuals) if viz_manifest is not None else [],
        slide_ledger=source_ledger,
        ledger_entries=list(source_ledger.entries),
        appendix_start=blueprint.appendix_start,
        appendix_boundary_rule=deck_constitution.appendix_boundary_rule if deck_constitution is not None else None,
        layout_library_slide_ratio=layout_library.slide_ratio if layout_library is not None else None,
    )
    canonical_generation_profile = blueprint.canonical_generation_profile
    canonical_generation_profile_payload = (
        canonical_generation_profile.model_dump(mode="json", exclude_none=True)
        if canonical_generation_profile is not None
        else None
    )

    slide_ir = SlideIRDocument(
        deck_title=deck_title,
        slide_ratio=blueprint.slide_ratio,
        slide_width_in=slide_width_in,
        slide_height_in=slide_height_in,
        slides=ir_slides,
        generation_inputs={
            "slide_count": len(ir_slides),
            "ledger_matched": len(ir_slides) == len(source_ledger.entries),
            "layout_family_mode": "rule-based",
            "generation_mode": (
                canonical_generation_profile.mode.value
                if canonical_generation_profile is not None
                else "unspecified"
            ),
            "canonical_generation_profile": canonical_generation_profile_payload,
        },
        warnings=[],
        continuity=continuity,
        compile_context=compile_context,
    )
    if not enable_layout_critic:
        disabled_report = SlideIRLayoutCriticReport(
            deck_title=slide_ir.deck_title,
            slide_ratio=slide_ir.slide_ratio,
            critic_enabled=False,
            deterministic_fallback_available=True,
            selected_candidate_ids=[
                _layout_candidate_id(slide.slide_id, slide.layout_family, 1)
                for slide in slide_ir.slides
            ],
            summary={
                "slide_count": len(slide_ir.slides),
                "candidate_count": 0,
                "patched_candidate_count": 0,
                "selected_layout_families": [slide.layout_family for slide in slide_ir.slides],
                "selection_mode": "rule-based",
            },
        )
        generation_inputs = dict(slide_ir.generation_inputs)
        generation_inputs[LAYOUT_CRITIC_REPORT_KEY] = disabled_report.model_dump(mode="json", exclude_none=True)
        return _attach_style_prior(
            slide_ir.model_copy(update={"generation_inputs": generation_inputs}),
            design_system=design_system,
            canonical_generation_profile=canonical_generation_profile,
            style_prior_provider=style_prior_provider,
        )

    optimized_slide_ir, critic_report = _optimize_slide_ir_layouts(slide_ir, slide_specs)
    generation_inputs = dict(optimized_slide_ir.generation_inputs)
    generation_inputs["layout_family_mode"] = "critic-selected"
    generation_inputs[LAYOUT_CRITIC_REPORT_KEY] = critic_report.model_dump(mode="json", exclude_none=True)
    return _attach_style_prior(
        optimized_slide_ir.model_copy(update={"generation_inputs": generation_inputs}),
        design_system=design_system,
        canonical_generation_profile=canonical_generation_profile,
        style_prior_provider=style_prior_provider,
    )

def _append_unique(values: list[str], candidate: str) -> None:
    if candidate and candidate not in values:
        values.append(candidate)


def _layout_candidate_id(slide_id: str, family: str, index: int, *, patched: bool = False) -> str:
    candidate_id = f"{slide_id}:{family}:c{index:02d}"
    if patched:
        return f"{candidate_id}:patch1"
    return candidate_id


def _layout_candidate_families(slide: BlueprintSlide, resolved_family: str) -> tuple[str, ...]:
    families: list[str] = []
    _append_unique(families, resolved_family)
    _append_unique(families, _default_layout_family(slide))
    if slide.slide_role in {SlideRole.TITLE, SlideRole.SECTION_DIVIDER}:
        return tuple(families[:1])
    if slide.deck_mode == DeckMode.APPENDIX:
        _append_unique(families, "appendix-reference")
    if slide.visual_type in {VisualType.DOCUMENT_CROP, VisualType.CHART, VisualType.PHOTO}:
        _append_unique(families, "worked-example")
        _append_unique(families, "comparison")
    elif slide.visual_type in {VisualType.COMPARISON, VisualType.TABLE}:
        _append_unique(families, "comparison")
        _append_unique(families, "worked-example")
    elif slide.visual_type in {VisualType.PROCESS, VisualType.TIMELINE, VisualType.DECISION_PATH}:
        _append_unique(families, "process-flow")
    elif slide.visual_type in {VisualType.FRAMEWORK, VisualType.HIERARCHY, VisualType.INFOGRAPHIC, VisualType.METRIC_SUMMARY}:
        _append_unique(families, "concept-explainer")
        _append_unique(families, "definition-theorem")
    else:
        if slide.slide_role in {SlideRole.RECOMMENDATION, SlideRole.EXECUTIVE_SUMMARY}:
            _append_unique(families, "summary")
        _append_unique(families, "definition-theorem")
        _append_unique(families, "concept-explainer")
    return tuple(families[:MAX_LAYOUT_CANDIDATES])


def _rebuild_slide_ir_document(
    template: SlideIRDocument,
    slides: list[SlideIRSlide],
    *,
    generation_inputs: dict[str, Any] | None = None,
) -> SlideIRDocument:
    compile_context = template.compile_context
    design_system = compile_context.design_system if compile_context is not None else None
    style = _preview_style_for_slide_ir(
        slide_width_in=template.slide_width_in,
        slide_height_in=template.slide_height_in,
        design_system=design_system,
    )
    continuity = _build_slide_ir_continuity_metadata(
        deck_title=template.deck_title,
        slide_ratio=template.slide_ratio,
        design_system=design_system,
        slides=slides,
        style=style,
    )
    return template.model_copy(
        update={
            "slides": slides,
            "continuity": continuity,
            "generation_inputs": dict(generation_inputs) if generation_inputs is not None else dict(template.generation_inputs),
        }
    )


def _candidate_preview_summary(candidate_doc: SlideIRDocument, slide_number: int) -> dict[str, Any]:
    # Layout scoring consumes only preview slide/object structure. Continuity
    # fields on the preview payload are optional emit-only handoff data, and
    # this helper must keep working even when live preview continuity mirrors
    # are retired from the supported payload contract.
    preview_payload = render_slide_preview(candidate_doc)
    preview_slide = next(item for item in preview_payload["slides"] if item["slide_number"] == slide_number)
    return {
        "object_count": len(preview_slide["objects"]),
        "slots": [obj["slot"] for obj in preview_slide["objects"]],
        "layout_family": preview_slide["layout_family"],
    }


def _layout_semantic_fit_bonus(slide: BlueprintSlide, family: str) -> float:
    if slide.deck_mode == DeckMode.APPENDIX:
        return 16.0 if family == "appendix-reference" else 0.0
    if slide.visual_type in {VisualType.DOCUMENT_CROP, VisualType.CHART, VisualType.PHOTO}:
        if family == "worked-example":
            return 16.0
        if family == "comparison":
            return 10.0
        if family == "summary":
            return -12.0
    if slide.visual_type in {VisualType.COMPARISON, VisualType.TABLE}:
        if family == "comparison":
            return 16.0
        if family == "worked-example":
            return 8.0
    if slide.visual_type in {VisualType.PROCESS, VisualType.TIMELINE, VisualType.DECISION_PATH}:
        return 16.0 if family == "process-flow" else 0.0
    if slide.visual_type in {VisualType.FRAMEWORK, VisualType.HIERARCHY, VisualType.INFOGRAPHIC, VisualType.METRIC_SUMMARY}:
        return 12.0 if family in {"concept-explainer", "definition-theorem"} else 0.0
    if slide.slide_role in {SlideRole.RECOMMENDATION, SlideRole.EXECUTIVE_SUMMARY}:
        return 10.0 if family == "summary" else 0.0
    return 0.0


def _score_layout_candidate(
    *,
    slide: BlueprintSlide,
    candidate_slide: SlideIRSlide,
    candidate_doc: SlideIRDocument,
    candidate_id: str,
    patched: bool = False,
) -> SlideIRLayoutCandidateScore:
    validation = validate_slide_ir_geometry(candidate_doc)
    continuity = validate_slide_ir_continuity(candidate_doc)
    slide_findings = [finding for finding in validation.findings if finding.slide_number == candidate_slide.slide_number]
    slide_validation = SlideIRValidationReport(
        is_valid=not any(finding.severity == "error" for finding in slide_findings),
        object_count=len(candidate_slide.objects),
        findings=slide_findings,
    )
    preview_summary = _candidate_preview_summary(candidate_doc, candidate_slide.slide_number)
    slide_area = candidate_doc.slide_width_in * candidate_doc.slide_height_in
    content_objects = [obj for obj in candidate_slide.objects if obj.slot != "footer"]
    content_area = sum(obj.bounds.width * obj.bounds.height for obj in content_objects)
    content_bottom = max((obj.bounds.bottom for obj in content_objects), default=0.0)
    coverage_ratio = round((content_area / slide_area), 4) if slide_area else 0.0
    bottom_padding = round(max(candidate_doc.slide_height_in - content_bottom, 0.0), 4)
    visual_objects = [obj for obj in content_objects if obj.kind == "visual"]
    visual_width_ratio = round((visual_objects[0].bounds.width / candidate_doc.slide_width_in), 4) if visual_objects else 0.0
    uses_two_column_layout = bool(visual_objects) and visual_width_ratio < 0.58
    alignment_columns = sorted({round(obj.bounds.left, 2) for obj in content_objects})
    geometry_errors = slide_validation.error_count
    geometry_warnings = slide_validation.warning_count
    continuity_warning_count = continuity.warning_count
    geometry_score = max(0.0, 100.0 - (geometry_errors * 70.0) - (geometry_warnings * 22.0))
    continuity_score = max(0.0, 100.0 - (continuity_warning_count * 9.0))
    heuristic_score = 62.0
    if 0.18 <= coverage_ratio <= 0.58:
        heuristic_score += 10.0
    else:
        heuristic_score -= min(abs(coverage_ratio - 0.36) * 60.0, 18.0)
    if 0.4 <= bottom_padding <= 1.2:
        heuristic_score += 8.0
    else:
        heuristic_score -= min(abs(bottom_padding - 0.72) * 20.0, 14.0)
    heuristic_score += _layout_semantic_fit_bonus(slide, candidate_slide.layout_family)
    if visual_objects and uses_two_column_layout:
        heuristic_score += 6.0
    if visual_objects and not uses_two_column_layout and candidate_slide.supporting_evidence:
        heuristic_score -= 12.0
    if len(alignment_columns) <= 3:
        heuristic_score += 4.0
    else:
        heuristic_score -= (len(alignment_columns) - 3) * 3.0
    heuristic_score = max(0.0, min(100.0, heuristic_score))
    aesthetic_score = 70.0
    aesthetic_score -= min(abs(coverage_ratio - 0.32) * 120.0, 18.0)
    aesthetic_score -= max(len(alignment_columns) - 3, 0) * 6.0
    if bottom_padding < 0.25:
        aesthetic_score -= 12.0
    aesthetic_score = max(0.0, min(100.0, aesthetic_score))
    total_score = round(
        (geometry_score * 0.5)
        + (continuity_score * 0.2)
        + (heuristic_score * 0.25)
        + (aesthetic_score * 0.05),
        3,
    )
    warnings = _dedupe([*candidate_slide.layout_warnings, *_slide_ir_warnings(slide_validation), *continuity.as_guidance_lines()])
    preview_summary["uses_two_column_layout"] = uses_two_column_layout
    return SlideIRLayoutCandidateScore(
        slide_number=candidate_slide.slide_number,
        slide_id=candidate_slide.slide_id,
        candidate_id=candidate_id,
        layout_family=candidate_slide.layout_family,
        patched=patched,
        total_score=total_score,
        geometry_score=round(geometry_score, 3),
        continuity_score=round(continuity_score, 3),
        heuristic_score=round(heuristic_score, 3),
        aesthetic_score=round(aesthetic_score, 3),
        geometry_error_count=geometry_errors,
        geometry_warning_count=geometry_warnings,
        continuity_warning_count=continuity_warning_count,
        quality_signals={
            "coverage_ratio": coverage_ratio,
            "bottom_padding_in": bottom_padding,
            "alignment_columns": len(alignment_columns),
            "supporting_evidence_count": len(candidate_slide.supporting_evidence),
            "visual_width_ratio": visual_width_ratio,
        },
        preview_summary=preview_summary,
        warnings=warnings,
    )


def _candidate_needs_patch(score: SlideIRLayoutCandidateScore) -> bool:
    return (
        score.geometry_error_count > 0
        or score.geometry_warning_count > 0
        or float(score.quality_signals.get("bottom_padding_in", 0.0)) < 0.22
    )


def _patch_preview_geometry(slide: SlideIRSlide, *, slide_width_in: float, slide_height_in: float) -> tuple[SlideIRSlide, str] | None:
    visual = next((obj for obj in slide.objects if obj.slot == "primary-visual"), None)
    evidence_cursor = visual.bounds.bottom + 0.14 if visual is not None else 0.0
    patched_objects: list[SlideIRObject] = []
    changed = False
    for obj in slide.objects:
        left = obj.bounds.left
        top = obj.bounds.top
        if obj.slot.startswith("evidence-") and visual is not None:
            top = max(top, evidence_cursor)
            if _rect_overlaps(IRRect(left=left, top=top, width=obj.bounds.width, height=obj.bounds.height), visual.bounds):
                top = visual.bounds.bottom + 0.14
            evidence_cursor = top + obj.bounds.height + 0.08
        max_left = max(slide_width_in - obj.bounds.width, 0.0)
        max_top = max(slide_height_in - obj.bounds.height - (0.18 if obj.slot == "footer" else 0.62), 0.0)
        clamped_left = min(max(left, 0.0), max_left)
        clamped_top = min(max(top, 0.0), max_top)
        if round(clamped_left, 4) != round(obj.bounds.left, 4) or round(clamped_top, 4) != round(obj.bounds.top, 4):
            changed = True
            patched_objects.append(
                obj.model_copy(
                    update={
                        "bounds": IRRect(
                            left=round(clamped_left, 4),
                            top=round(clamped_top, 4),
                            width=obj.bounds.width,
                            height=obj.bounds.height,
                        )
                    }
                )
            )
        else:
            patched_objects.append(obj)
    if not changed:
        return None
    return (
        slide.model_copy(update={"objects": patched_objects}),
        "Applied one bounded preview-geometry patch pass to clamp bounds and separate overlapping evidence.",
    )


def _optimize_slide_ir_layouts(
    base_slide_ir: SlideIRDocument,
    slide_specs: list[AdaptedSlideSpec],
) -> tuple[SlideIRDocument, SlideIRLayoutCriticReport]:
    selected_slides = list(base_slide_ir.slides)
    selected_candidate_ids: list[str] = []
    all_scores: list[SlideIRLayoutCandidateScore] = []
    patched_candidate_ids: list[str] = []
    appendix_start = base_slide_ir.compile_context.appendix_start if base_slide_ir.compile_context is not None else None
    style = _preview_style_for_slide_ir(
        slide_width_in=base_slide_ir.slide_width_in,
        slide_height_in=base_slide_ir.slide_height_in,
        design_system=base_slide_ir.compile_context.design_system if base_slide_ir.compile_context is not None else None,
    )

    for index, spec in enumerate(slide_specs):
        candidate_runs: list[ScoredSlideCandidate] = []
        for candidate_index, family in enumerate(spec.candidate_families, start=1):
            candidate_slide = _build_slide_ir_slide(
                slide=spec.slide,
                entry=spec.entry,
                style=style,
                family=family,
                layout_warnings=list(spec.layout_warnings),
                appendix_start=appendix_start,
            )
            candidate_slides = list(selected_slides)
            candidate_slides[index] = candidate_slide
            candidate_doc = _rebuild_slide_ir_document(base_slide_ir, candidate_slides)
            candidate_runs.append(
                ScoredSlideCandidate(
                    slide=candidate_slide,
                    score=_score_layout_candidate(
                        slide=spec.slide,
                        candidate_slide=candidate_slide,
                        candidate_doc=candidate_doc,
                        candidate_id=_layout_candidate_id(spec.entry.slide_id, family, candidate_index),
                    ),
                )
            )

        best_candidate = max(
            enumerate(candidate_runs),
            key=lambda item: (item[1].score.total_score, -item[0]),
        )[1]
        if _candidate_needs_patch(best_candidate.score):
            patched = _patch_preview_geometry(
                best_candidate.slide,
                slide_width_in=base_slide_ir.slide_width_in,
                slide_height_in=base_slide_ir.slide_height_in,
            )
            if patched is not None:
                patched_slide, patch_warning = patched
                patched_slides = list(selected_slides)
                patched_slides[index] = patched_slide
                patched_doc = _rebuild_slide_ir_document(base_slide_ir, patched_slides)
                patched_score = _score_layout_candidate(
                    slide=spec.slide,
                    candidate_slide=patched_slide,
                    candidate_doc=patched_doc,
                    candidate_id=_layout_candidate_id(spec.entry.slide_id, best_candidate.slide.layout_family, len(candidate_runs), patched=True),
                    patched=True,
                )
                patched_score = patched_score.model_copy(update={"warnings": _dedupe([patch_warning, *patched_score.warnings])})
                candidate_runs.append(ScoredSlideCandidate(slide=patched_slide, score=patched_score))
                if (
                    patched_score.geometry_error_count < best_candidate.score.geometry_error_count
                    or patched_score.geometry_warning_count < best_candidate.score.geometry_warning_count
                    or patched_score.total_score >= best_candidate.score.total_score
                ):
                    best_candidate = candidate_runs[-1]

        selected_slides[index] = best_candidate.slide
        selected_candidate_ids.append(best_candidate.score.candidate_id)
        if best_candidate.score.patched:
            patched_candidate_ids.append(best_candidate.score.candidate_id)
        for candidate_run in candidate_runs:
            all_scores.append(
                candidate_run.score.model_copy(update={"selected": candidate_run.score.candidate_id == best_candidate.score.candidate_id})
            )

    optimized_slide_ir = _rebuild_slide_ir_document(base_slide_ir, selected_slides)
    geometry_report = validate_slide_ir_geometry(optimized_slide_ir)
    continuity_report = validate_slide_ir_continuity(optimized_slide_ir)
    critic_report = SlideIRLayoutCriticReport(
        deck_title=optimized_slide_ir.deck_title,
        slide_ratio=optimized_slide_ir.slide_ratio,
        critic_enabled=True,
        deterministic_fallback_available=True,
        selected_candidate_ids=selected_candidate_ids,
        scores=all_scores,
        summary={
            "slide_count": len(optimized_slide_ir.slides),
            "candidate_count": len(all_scores),
            "patched_candidate_count": len(patched_candidate_ids),
            "selected_layout_families": [slide.layout_family for slide in optimized_slide_ir.slides],
            "geometry_error_count": geometry_report.error_count,
            "geometry_warning_count": geometry_report.warning_count,
            "continuity_warning_count": continuity_report.warning_count,
        },
    )
    return optimized_slide_ir, critic_report


def _component_typography_tokens(slot: str) -> list[str]:
    lowered = slot.lower()
    if "title" in lowered:
        return ["title"]
    if "caption" in lowered or "note" in lowered or "source" in lowered:
        return ["caption"]
    return ["body"]


def _section_label(slide: SlideIRSlide) -> str:
    return slide.section.strip() or slide.deck_mode.value


def _section_structure(slides: list[SlideIRSlide]) -> list[SlideIRSectionStructure]:
    sections: dict[str, SlideIRSectionStructure] = {}
    for slide in sorted(slides, key=lambda item: item.slide_number):
        section = _section_label(slide)
        summary = sections.get(section)
        if summary is None:
            summary = SlideIRSectionStructure(section=section)
            sections[section] = summary
        summary.slide_numbers.append(slide.slide_number)
        summary.slide_ids.append(slide.slide_id)
        _append_unique(summary.deck_modes, slide.deck_mode.value)
        _append_unique(summary.layout_families, slide.layout_family)
    return list(sections.values())


def _typography_scale_refs(design_system: DesignSystem | None) -> list[SlideIRTypographyScaleReference]:
    if design_system is None:
        return []
    return [
        SlideIRTypographyScaleReference(
            token=token.token,
            font_family=token.font_family,
            size_pt=token.size_pt,
            weight=token.weight,
            usage=token.usage,
        )
        for token in design_system.typography_tokens
    ]


def _spacing_rhythm_refs(
    design_system: DesignSystem | None,
    *,
    style: CompilerStyle,
) -> list[SlideIRSpacingRhythmReference]:
    if design_system is None:
        return []
    return [
        SlideIRSpacingRhythmReference(
            scale_name="design-system-scale",
            scale_steps_pt=[float(step) for step in design_system.spacing_scale],
            margin_in=round(style.margin, 4),
            gap_in=round(style.gap, 4),
        )
    ]


def _component_style_refs(
    slides: list[SlideIRSlide],
    design_system: DesignSystem | None,
) -> list[SlideIRComponentStyleReference]:
    slot_usage: dict[str, dict[str, list[str] | int]] = {}
    color_tokens = [token.token for token in design_system.color_tokens] if design_system is not None else []
    for slide in slides:
        for obj in slide.objects:
            usage = slot_usage.setdefault(obj.slot, {"count": 0, "kinds": [], "layout_families": []})
            usage["count"] = int(usage["count"]) + 1
            _append_unique(usage["kinds"], obj.kind)
            _append_unique(usage["layout_families"], slide.layout_family)
    refs: list[SlideIRComponentStyleReference] = []
    for slot in sorted(slot_usage):
        usage = slot_usage[slot]
        if int(usage["count"]) < 2:
            continue
        kinds = list(usage["kinds"])
        refs.append(
            SlideIRComponentStyleReference(
                component_id=slot,
                style_family=kinds[0] if len(kinds) == 1 else "mixed",
                typography_tokens=_component_typography_tokens(slot),
                color_tokens=color_tokens[:2],
                spacing_refs=["design-system-scale"],
                layout_families=list(usage["layout_families"]),
            )
        )
    return refs


def _slide_density_budget_reference(slides: list[SlideIRSlide]) -> SlideIRDensityBudgetReference | None:
    for slide in slides:
        if slide.density_budget is None:
            continue
        return SlideIRDensityBudgetReference.model_validate(slide.density_budget.model_dump(mode="json", exclude_none=True))
    return None


def _continuity_anchor_type(previous_slide: SlideIRSlide, slide: SlideIRSlide) -> str:
    if previous_slide.deck_mode != slide.deck_mode and slide.deck_mode == DeckMode.APPENDIX:
        return "appendix_transition"
    if _section_label(previous_slide) != _section_label(slide):
        return "section_transition"
    if _section_label(previous_slide) == _section_label(slide):
        return "section"
    return "sequential"


def _continuity_anchor_label(previous_slide: SlideIRSlide, slide: SlideIRSlide) -> str:
    if previous_slide.deck_mode != slide.deck_mode and slide.deck_mode == DeckMode.APPENDIX:
        return "Transition to appendix"
    if _section_label(previous_slide) == _section_label(slide):
        return f"Continue {_section_label(slide)}"
    return f"{_section_label(previous_slide)} -> {_section_label(slide)}"


def _continuity_anchors(slides: list[SlideIRSlide]) -> list[SlideIRContinuityAnchor]:
    ordered = sorted(slides, key=lambda item: item.slide_number)
    anchors: list[SlideIRContinuityAnchor] = []
    for previous_slide, slide in zip(ordered, ordered[1:]):
        anchors.append(
            SlideIRContinuityAnchor(
                from_slide_number=previous_slide.slide_number,
                from_slide_id=previous_slide.slide_id,
                to_slide_number=slide.slide_number,
                to_slide_id=slide.slide_id,
                anchor_type=_continuity_anchor_type(previous_slide, slide),
                label=_continuity_anchor_label(previous_slide, slide),
            )
        )
    return anchors


def _build_slide_ir_continuity_metadata(
    *,
    deck_title: str,
    slide_ratio: str,
    design_system: DesignSystem | None,
    slides: list[SlideIRSlide],
    style: CompilerStyle,
) -> SlideIRDeckContinuityMetadata:
    return SlideIRDeckContinuityMetadata(
        deck_identity=SlideIRDeckIdentity(deck_title=deck_title, slide_ratio=slide_ratio),
        theme_identity=SlideIRThemeIdentity(
            theme_name=design_system.theme_name if design_system is not None else None,
            brand_name=design_system.brand_name if design_system is not None else None,
            visual_route_id=design_system.visual_route_id if design_system is not None else None,
            reference_source_family=design_system.reference_source_family if design_system is not None else None,
        ),
        section_structure=_section_structure(slides),
        typography_scale_refs=_typography_scale_refs(design_system),
        spacing_rhythm_refs=_spacing_rhythm_refs(design_system, style=style),
        repeated_component_style_refs=_component_style_refs(slides, design_system),
        slide_density_budget=_slide_density_budget_reference(slides),
        continuity_anchors=_continuity_anchors(slides),
    )


def _slide_ir_warnings(report: SlideIRValidationReport) -> list[str]:
    return [
        f"Slide IR {finding.code}: slide {finding.slide_number or '?'} object {finding.object_id or 'n/a'} - {finding.message}"
        for finding in report.findings
        if finding.severity == "warning"
    ]


def validate_slide_ir_geometry(slide_ir: SlideIRDocument) -> SlideIRValidationReport:
    findings: list[SlideIRValidationFinding] = []
    for slide in slide_ir.slides:
        findings.extend(_findings_for_ir_bounds(slide, slide_ir.slide_width_in, slide_ir.slide_height_in))
    return SlideIRValidationReport(
        is_valid=not any(finding.severity == "error" for finding in findings),
        object_count=sum(len(slide.objects) for slide in slide_ir.slides),
        findings=findings,
    )


def _title_hierarchy_findings(slides: list[SlideIRSlide]) -> list[SlideIRContinuityFinding]:
    findings: list[SlideIRContinuityFinding] = []
    for slide in sorted(slides, key=lambda item: item.slide_number):
        title = slide.title.strip()
        if len(title) <= 96:
            continue
        findings.append(
            SlideIRContinuityFinding(
                code="title_hierarchy_drift",
                message="Title length exceeds the stable deck heading range and may blur hierarchy.",
                slide_number=slide.slide_number,
                slide_id=slide.slide_id,
                details={"title_length": len(title)},
            )
        )
    return findings


def _typography_scale_findings(continuity: SlideIRDeckContinuityMetadata | None) -> list[SlideIRContinuityFinding]:
    if continuity is None:
        return []
    findings: list[SlideIRContinuityFinding] = []
    token_map: dict[str, SlideIRTypographyScaleReference] = {}
    duplicate_tokens: list[str] = []
    for reference in continuity.typography_scale_refs:
        if reference.token in token_map:
            duplicate_tokens.append(reference.token)
        else:
            token_map[reference.token] = reference
    if duplicate_tokens:
        findings.append(
            SlideIRContinuityFinding(
                code="typography_scale_drift",
                message=f"Duplicate typography tokens detected in continuity metadata: {', '.join(sorted(set(duplicate_tokens)))}.",
                details={"duplicate_tokens": sorted(set(duplicate_tokens))},
            )
        )
    title_ref = token_map.get("title")
    body_ref = token_map.get("body")
    caption_ref = token_map.get("caption")
    if title_ref is not None and body_ref is not None and title_ref.size_pt <= body_ref.size_pt:
        findings.append(
            SlideIRContinuityFinding(
                code="typography_scale_drift",
                message="Title scale is not larger than body scale in continuity metadata.",
                details={"title_size_pt": title_ref.size_pt, "body_size_pt": body_ref.size_pt},
            )
        )
    if body_ref is not None and caption_ref is not None and body_ref.size_pt <= caption_ref.size_pt:
        findings.append(
            SlideIRContinuityFinding(
                code="typography_scale_drift",
                message="Body scale is not larger than caption scale in continuity metadata.",
                details={"body_size_pt": body_ref.size_pt, "caption_size_pt": caption_ref.size_pt},
            )
        )
    return findings


def _spacing_rhythm_findings(continuity: SlideIRDeckContinuityMetadata | None) -> list[SlideIRContinuityFinding]:
    if continuity is None:
        return []
    findings: list[SlideIRContinuityFinding] = []
    for reference in continuity.spacing_rhythm_refs:
        steps = reference.scale_steps_pt
        if any(step <= 0 for step in steps) or any(current >= following for current, following in zip(steps, steps[1:])):
            findings.append(
                SlideIRContinuityFinding(
                    code="spacing_rhythm_drift",
                    message=f"Spacing rhythm `{reference.scale_name}` is not strictly increasing.",
                    details={"scale_steps_pt": steps},
                )
            )
        if (reference.margin_in is not None and reference.margin_in <= 0) or (reference.gap_in is not None and reference.gap_in <= 0):
            findings.append(
                SlideIRContinuityFinding(
                    code="spacing_rhythm_drift",
                    message=f"Spacing rhythm `{reference.scale_name}` contains non-positive margin or gap values.",
                    details={"margin_in": reference.margin_in, "gap_in": reference.gap_in},
                )
            )
    return findings


def _component_style_findings(continuity: SlideIRDeckContinuityMetadata | None) -> list[SlideIRContinuityFinding]:
    if continuity is None:
        return []
    grouped: dict[str, list[SlideIRComponentStyleReference]] = {}
    for reference in continuity.repeated_component_style_refs:
        grouped.setdefault(reference.component_id, []).append(reference)
    findings: list[SlideIRContinuityFinding] = []
    for component_id, references in sorted(grouped.items()):
        signatures = {
            json.dumps(reference.model_dump(mode="json", exclude={"component_id"}), sort_keys=True) for reference in references
        }
        if len(signatures) <= 1:
            continue
        findings.append(
            SlideIRContinuityFinding(
                code="component_style_drift",
                message=f"Repeated component `{component_id}` maps to multiple style signatures.",
                details={"component_id": component_id, "variant_count": len(signatures)},
            )
        )
    return findings


def _section_continuity_findings(slides: list[SlideIRSlide]) -> list[SlideIRContinuityFinding]:
    findings: list[SlideIRContinuityFinding] = []
    seen_sections: set[str] = set()
    previous_section: str | None = None
    for slide in sorted(slides, key=lambda item: item.slide_number):
        section = _section_label(slide)
        if section != previous_section and section in seen_sections:
            findings.append(
                SlideIRContinuityFinding(
                    code="section_continuity_break",
                    message=f"Section `{section}` resumes after another section intervened.",
                    slide_number=slide.slide_number,
                    slide_id=slide.slide_id,
                    details={"section": section},
                )
            )
        seen_sections.add(section)
        previous_section = section
    return findings


def _slide_text_char_count(slide: SlideIRSlide) -> int:
    text_parts = [
        slide.title,
        slide.one_line_takeaway,
        slide.main_message,
        slide.primary_claim,
        slide.audience_intent,
        *slide.must_keep_text,
        *slide.optional_text,
        *slide.core_content,
        *slide.supporting_evidence,
    ]
    return sum(len(part.strip()) for part in text_parts if part)


def _slide_density_score(slide: SlideIRSlide) -> int:
    text_chars = _slide_text_char_count(slide)
    evidence_items = len(slide.supporting_evidence) + len(slide.required_evidence_assets)
    slot_count = max(len(slide.layout_slot_map), len(slide.objects))
    return text_chars + (len(slide.objects) * 40) + (evidence_items * 30) + (slot_count * 20)


def _density_budget_text_ceiling(slide: SlideIRSlide, continuity: SlideIRDeckContinuityMetadata | None) -> int | None:
    if slide.density_budget is not None and slide.density_budget.text_char_ceiling is not None:
        return slide.density_budget.text_char_ceiling
    if continuity is not None and continuity.slide_density_budget is not None:
        return continuity.slide_density_budget.text_char_ceiling
    return None


def _density_findings(slide_ir: SlideIRDocument) -> list[SlideIRContinuityFinding]:
    ordered = sorted(slide_ir.slides, key=lambda item: item.slide_number)
    if not ordered:
        return []
    scores = [_slide_density_score(slide) for slide in ordered]
    baseline = median(scores)
    findings: list[SlideIRContinuityFinding] = []
    for slide, score in zip(ordered, scores):
        text_chars = _slide_text_char_count(slide)
        ceiling = _density_budget_text_ceiling(slide, slide_ir.continuity)
        over_ceiling = ceiling is not None and text_chars > int(ceiling * 1.5)
        is_outlier = baseline > 0 and score > (baseline * 1.9) and (score - baseline) >= 250
        if not over_ceiling and not is_outlier:
            continue
        findings.append(
            SlideIRContinuityFinding(
                code="density_outlier",
                message="Slide density materially exceeds the deck baseline and should be reviewed for continuity drift.",
                slide_number=slide.slide_number,
                slide_id=slide.slide_id,
                details={"density_score": score, "baseline_density_score": baseline, "text_chars": text_chars, "text_char_ceiling": ceiling},
            )
        )
    return findings


def validate_slide_ir_continuity(slide_ir: SlideIRDocument) -> SlideIRContinuityReport:
    findings = [
        *_title_hierarchy_findings(slide_ir.slides),
        *_typography_scale_findings(slide_ir.continuity),
        *_spacing_rhythm_findings(slide_ir.continuity),
        *_component_style_findings(slide_ir.continuity),
        *_section_continuity_findings(slide_ir.slides),
        *_density_findings(slide_ir),
    ]
    return SlideIRContinuityReport(is_valid=True, findings=findings)


def _preview_slide_payload(slide: SlideIRSlide) -> dict[str, Any]:
    return {
        "slide_number": slide.slide_number,
        "slide_id": slide.slide_id,
        "title": slide.title,
        "deck_mode": slide.deck_mode.value,
        "slide_role": slide.slide_role.value,
        "visual_type": slide.visual_type.value,
        "slide_archetype": slide.slide_archetype.value if slide.slide_archetype is not None else None,
        "layout_pattern_id": slide.layout_pattern_id,
        "layout_family": slide.layout_family,
        "objects": [obj.model_dump(mode="json", exclude_none=True) for obj in slide.objects],
        "notes": list(slide.notes),
    }


def _preview_metadata(slide_ir: SlideIRDocument) -> SlideIRPreviewMetadata:
    return SlideIRPreviewMetadata(
        slide_count=len(slide_ir.slides),
        object_count=sum(len(slide.objects) for slide in slide_ir.slides),
        slide_ids=[slide.slide_id for slide in slide_ir.slides],
        layout_families=_dedupe([slide.layout_family for slide in slide_ir.slides]),
    )


def _round_geometry(value: float) -> float:
    return round(value, 4)


def _geometry_summaries(slide_ir: SlideIRDocument) -> list[SlideIRGeometrySlideSummary]:
    summaries: list[SlideIRGeometrySlideSummary] = []
    for slide in slide_ir.slides:
        summaries.append(
            SlideIRGeometrySlideSummary(
                slide_number=slide.slide_number,
                slide_id=slide.slide_id,
                layout_family=slide.layout_family,
                object_count=len(slide.objects),
                objects=[
                    SlideIRGeometryObjectSummary(
                        object_id=obj.object_id,
                        slot=obj.slot,
                        kind=obj.kind,
                        left=_round_geometry(obj.bounds.left),
                        top=_round_geometry(obj.bounds.top),
                        width=_round_geometry(obj.bounds.width),
                        height=_round_geometry(obj.bounds.height),
                        right=_round_geometry(obj.bounds.right),
                        bottom=_round_geometry(obj.bounds.bottom),
                    )
                    for obj in slide.objects
                ],
            )
        )
    return summaries


def build_slide_ir_compile_report(
    slide_ir: SlideIRDocument,
    validation: SlideIRValidationReport,
    warnings: list[str],
    continuity: SlideIRContinuityReport | None = None,
) -> SlideIRCompileReport:
    continuity_report = continuity if continuity is not None else validate_slide_ir_continuity(slide_ir)
    continuity_guidance = continuity_report.as_guidance_lines()
    # The supported compile-report surface is now guidance-first for both new
    # in-memory objects and newly serialized artifacts. Older warning-mirror
    # payloads remain accepted only at the schema/load normalization boundary,
    # where they are upgraded onto canonical `continuity_guidance`.
    return SlideIRCompileReport(
        deck_title=slide_ir.deck_title,
        slide_ratio=slide_ir.slide_ratio,
        validation=validation,
        continuity=continuity_report,
        continuity_guidance=continuity_guidance,
        warnings=_dedupe(warnings),
        preview_metadata=_preview_metadata(slide_ir),
        geometry_summaries=_geometry_summaries(slide_ir),
    )


def _preview_continuity_guidance(continuity_report: SlideIRContinuityReport) -> list[str]:
    """Keep preview continuity guidance canonical for the live preview payload.

    Supported direct preview consumers stop at the in-repo emit-time helper
    path. Repo-outside direct consumption of preview continuity mirrors is
    outside the supported contract, so the live preview payload now exposes
    only canonical `continuity_guidance`. New in-repo preview integrations
    must not reintroduce policy dependence on a retired preview
    `continuity_warnings` mirror.
    """
    return _dedupe(continuity_report.as_guidance_lines())


def render_slide_preview(slide_ir: SlideIRDocument) -> dict[str, Any]:
    """Render a preview-ready JSON payload from a SlideIR document.

    The preview continuity contract is intentionally narrow: supported direct
    consumers are the in-repo emit-time helpers only, while repo-outside direct
    consumption of raw preview continuity mirrors is unsupported/out-of-scope.
    It exposes canonical operator-facing `continuity_guidance` and leaves
    heavier policy surfaces such as `continuity_alerts`,
    `compatibility_warning_codes`, and `compile_eligibility` to compile-report
    or persisted state artifacts.
    """
    continuity_report = validate_slide_ir_continuity(slide_ir)
    continuity_guidance = _preview_continuity_guidance(continuity_report)
    return {
        "deck_title": slide_ir.deck_title,
        "slide_ratio": slide_ir.slide_ratio,
        "slide_width_in": slide_ir.slide_width_in,
        "slide_height_in": slide_ir.slide_height_in,
        "slides": [_preview_slide_payload(slide) for slide in slide_ir.slides],
        "generation_inputs": dict(slide_ir.generation_inputs),
        "style_prior": slide_ir.style_prior.model_dump(mode="json", exclude_none=True) if slide_ir.style_prior is not None else None,
        "warnings": _dedupe([*slide_ir.warnings, *continuity_guidance]),
        "continuity": slide_ir.continuity.model_dump(mode="json", exclude_none=True) if slide_ir.continuity is not None else None,
        # Freeze the downstream preview contract around one canonical guidance
        # field. Heavier structured policy surfaces stay on compile-report or
        # persisted state artifacts instead of the emit-only preview payload.
        "continuity_guidance": continuity_guidance,
    }


def _parse_ratio(slide_ratio: str) -> tuple[float, float]:
    try:
        width_text, height_text = slide_ratio.split(":", 1)
        width_value = float(width_text)
        height_value = float(height_text)
    except Exception as exc:  # pragma: no cover - defensive
        raise ValueError(f"unsupported slide ratio {slide_ratio!r}") from exc
    if width_value <= 0 or height_value <= 0:
        raise ValueError("slide ratio must be positive")
    base_height = SLIDE_HEIGHT_IN
    return (base_height * (width_value / height_value), base_height)


def _rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    if len(value) == 3:
        value = "".join(char * 2 for char in value)
    return RGBColor.from_string(value.upper())


def _color_map(design_system: DesignSystem, style_prior: SlideIRStylePrior | None = None) -> dict[str, RGBColor]:
    colors = {token.token: _rgb(token.hex) for token in design_system.color_tokens}
    if style_prior is not None:
        for seed in style_prior.palette_seeds:
            colors.setdefault(seed.token, _rgb(seed.hex))
    colors.setdefault("ink", RGBColor(31, 41, 55))
    colors.setdefault("signal", RGBColor(194, 65, 12))
    colors.setdefault("canvas", RGBColor(248, 250, 252))
    colors.setdefault("muted", colors["canvas"])
    return colors


def _font_map(design_system: DesignSystem) -> dict[str, Any]:
    defaults = {
        "title": SimpleNamespace(font_family="Aptos Display", size_pt=24.0),
        "body": SimpleNamespace(font_family="Aptos", size_pt=11.0),
        "caption": SimpleNamespace(font_family="Aptos", size_pt=9.0),
    }
    mapped = {token.token: token for token in design_system.typography_tokens}
    for token, default in defaults.items():
        mapped.setdefault(token, default)
    return mapped


def _spacing(design_system: DesignSystem | list[float], index: int, fallback: float) -> float:
    spacing_scale = design_system.spacing_scale if isinstance(design_system, DesignSystem) else design_system
    if 0 <= index < len(spacing_scale):
        return spacing_scale[index] / 72.0
    return fallback


def _require_compile_context(slide_ir: SlideIRDocument) -> SlideIRCompileContext:
    if slide_ir.compile_context is None:
        raise ValueError("Slide IR compile context is required for PPTX compilation.")
    return slide_ir.compile_context


def _dedupe(items: list[str]) -> list[str]:
    seen: set[str] = set()
    ordered: list[str] = []
    for item in items:
        if item and item not in seen:
            seen.add(item)
            ordered.append(item)
    return ordered


def _ensure_blueprint_approved(blueprint: Blueprint, state_capsule: StateCapsule | None) -> None:
    approved = blueprint.approval_status == StageStatus.APPROVED or (state_capsule is not None and state_capsule.blueprint_approved)
    if not approved:
        raise ValueError("PPTX compilation requires blueprint approval before production.")


def _build_number_label(slide: BlueprintSlide, appendix_start: int | None) -> str:
    if slide.deck_mode == DeckMode.APPENDIX or (appendix_start is not None and slide.slide_number >= appendix_start):
        return f"Appendix {slide.slide_number}"
    return str(slide.slide_number)


def _layout_pattern_index(layout_library: LayoutLibrary) -> dict[str, LayoutPattern]:
    return {pattern.pattern_id: pattern for pattern in layout_library.patterns}


_LEGACY_LAYOUT_PATTERN_ALIASES: dict[str, str] = {
    "comparison": "evidence-table",
}


def _canonical_layout_pattern_id(pattern_id: str) -> str:
    return _LEGACY_LAYOUT_PATTERN_ALIASES.get(pattern_id, pattern_id)


def _layout_pattern_lookup(
    layout_index: dict[str, LayoutPattern],
    pattern_id: str,
) -> LayoutPattern | None:
    pattern = layout_index.get(pattern_id)
    if pattern is not None:
        return pattern
    return layout_index.get(_canonical_layout_pattern_id(pattern_id))


def _normalized_text(text: str | None) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _planned_text_char_count(slide: BlueprintSlide) -> int:
    items: list[str] = []
    seen: set[str] = set()
    for text in [
        slide.title,
        slide.one_line_takeaway,
        slide.main_message,
        *slide.core_content,
        *slide.must_keep_text,
        *slide.optional_text,
    ]:
        normalized = _normalized_text(text)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        items.append(normalized)
    return sum(len(item) for item in items)


def _planned_visual_node_count(slide: BlueprintSlide) -> int:
    payload = slide.authoring_payload if isinstance(slide.authoring_payload, dict) else {}
    counts: list[int] = []
    for value in payload.values():
        if isinstance(value, list):
            counts.append(len(value))
    return max(counts, default=0)


def _body_slot_ids(slide: BlueprintSlide) -> set[str]:
    return {slot for slot in slide.layout_slot_map.values() if slot != "header"}


def _compile_contract_errors(
    slide_ir: SlideIRDocument,
) -> list[str]:
    context = _require_compile_context(slide_ir)
    errors: list[str] = []
    layout_index = {pattern.pattern_id: pattern for pattern in context.layout_patterns}
    slides_by_number = {slide.slide_number: slide for slide in slide_ir.slides}
    ledger_by_number = {entry.slide_number: entry for entry in context.ledger_entries}

    for slide in slide_ir.slides:
        pattern = _layout_pattern_lookup(layout_index, slide.layout_pattern_id)
        if pattern is None:
            approved_ids = ", ".join(sorted(layout_index))
            errors.append(
                f"Slide {slide.slide_number} blueprint defect: layout_pattern_id `{slide.layout_pattern_id}` is not in the approved layout-library ids ({approved_ids})."
            )
            continue
        if not _normalized_text(slide.primary_claim):
            errors.append(f"Slide {slide.slide_number} blueprint defect: primary_claim is blank.")
        if not _normalized_text(slide.audience_intent):
            errors.append(f"Slide {slide.slide_number} blueprint defect: audience_intent is blank.")
        if not slide.must_keep_text:
            errors.append(f"Slide {slide.slide_number} blueprint defect: must_keep_text is empty.")
        if not slide.layout_slot_map:
            errors.append(f"Slide {slide.slide_number} blueprint defect: layout_slot_map is empty.")
        else:
            if "title" not in slide.layout_slot_map or "claim" not in slide.layout_slot_map:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: layout_slot_map must reserve title and claim slots."
                )
            if slide.visual_type not in {VisualType.TEXT, VisualType.QUOTE} and "primary_visual" not in slide.layout_slot_map:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: visual slides must map a `primary_visual` slot."
                )
            if (
                slide.evidence_class != SlideEvidenceClass.MESSAGE_ONLY
                and not {"evidence_marker", "supporting_text"} & set(slide.layout_slot_map)
            ):
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: evidence-backed slides must map either `evidence_marker` or `supporting_text`."
                )
            body_slots = _body_slot_ids(slide)
            if body_slots and len(body_slots) > pattern.body_slots:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: layout_slot_map requests {len(body_slots)} body slots but layout `{pattern.pattern_id}` allows {pattern.body_slots}."
                )
        if slide.density_budget is None:
            errors.append(f"Slide {slide.slide_number} blueprint defect: density_budget is missing from SlideIR.")
        else:
            slot_count = len(set(slide.layout_slot_map.values()))
            if slot_count > slide.density_budget.layout_slot_count:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: layout_slot_map uses {slot_count} slots but density_budget.layout_slot_count allows {slide.density_budget.layout_slot_count}."
                )
            planned_text_char_count = _planned_text_char_count(slide)
            if planned_text_char_count > slide.density_budget.text_char_ceiling:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: planned text density {planned_text_char_count} exceeds density_budget.text_char_ceiling={slide.density_budget.text_char_ceiling}."
                )
            if len(slide.core_content) > slide.density_budget.bullet_count_ceiling:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: core_content count {len(slide.core_content)} exceeds density_budget.bullet_count_ceiling={slide.density_budget.bullet_count_ceiling}."
                )
            if len(slide.supporting_evidence) > slide.density_budget.evidence_item_ceiling:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: supporting_evidence count {len(slide.supporting_evidence)} exceeds density_budget.evidence_item_ceiling={slide.density_budget.evidence_item_ceiling}."
                )
            planned_visual_nodes = _planned_visual_node_count(slide)
            if planned_visual_nodes > slide.density_budget.visual_node_ceiling:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: authored visual node count {planned_visual_nodes} exceeds density_budget.visual_node_ceiling={slide.density_budget.visual_node_ceiling}."
                )
        if context.appendix_start is not None:
            if slide.slide_number >= context.appendix_start and slide.deck_mode != DeckMode.APPENDIX:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: slides at or after appendix_start must be appendix."
                )
            if slide.slide_number < context.appendix_start and slide.deck_mode == DeckMode.APPENDIX:
                errors.append(
                    f"Slide {slide.slide_number} blueprint defect: slides before appendix_start cannot be appendix."
                )

    for entry in context.ledger_entries:
        slide = slides_by_number.get(entry.slide_number)
        if slide is None:
            continue
        if _canonical_layout_pattern_id(entry.layout_pattern_id) != _canonical_layout_pattern_id(slide.layout_pattern_id):
            errors.append(
                f"Slide {entry.slide_number} blueprint defect: ledger layout_pattern_id `{entry.layout_pattern_id}` does not match blueprint `{slide.layout_pattern_id}`."
            )
        if entry.slide_role != slide.slide_role:
            errors.append(
                f"Slide {entry.slide_number} blueprint defect: ledger slide_role `{entry.slide_role.value}` does not match blueprint `{slide.slide_role.value}`."
            )
        if entry.visual_type != slide.visual_type:
            errors.append(
                f"Slide {entry.slide_number} blueprint defect: ledger visual_type `{entry.visual_type.value}` does not match blueprint `{slide.visual_type.value}`."
            )
        if entry.deck_mode != slide.deck_mode:
            errors.append(
                f"Slide {entry.slide_number} blueprint defect: ledger deck_mode `{entry.deck_mode.value}` does not match blueprint `{slide.deck_mode.value}`."
            )
        if _normalized_text(entry.main_message) != _normalized_text(slide.main_message):
            errors.append(
                f"Slide {entry.slide_number} blueprint defect: ledger main_message drifted from the approved blueprint."
            )

    for slide_number in sorted(set(slides_by_number) - set(ledger_by_number)):
        errors.append(f"Slide {slide_number} blueprint defect: blueprint slide is missing from the slide ledger.")

    return errors


def _default_layout_family(slide: BlueprintSlide) -> str:
    if slide.slide_role == slide.slide_role.TITLE:
        return "cover"
    if slide.slide_role == slide.slide_role.SECTION_DIVIDER:
        return "section-divider"
    if slide.slide_role in {slide.slide_role.REFERENCES, slide.slide_role.APPENDIX_EVIDENCE}:
        return "appendix-reference"
    if slide.slide_role in {slide.slide_role.RECOMMENDATION, slide.slide_role.EXECUTIVE_SUMMARY}:
        return "summary"
    if slide.visual_type in {VisualType.PROCESS, VisualType.TIMELINE, VisualType.DECISION_PATH}:
        return "process-flow"
    if slide.visual_type in {VisualType.COMPARISON, VisualType.TABLE}:
        return "comparison"
    if slide.visual_type in {VisualType.DOCUMENT_CROP, VisualType.CHART, VisualType.PHOTO}:
        return "worked-example"
    if slide.visual_type in {VisualType.FRAMEWORK, VisualType.HIERARCHY, VisualType.INFOGRAPHIC, VisualType.METRIC_SUMMARY}:
        return "concept-explainer"
    return "definition-theorem"


def _layout_family(slide: BlueprintSlide, pattern: LayoutPattern | None) -> str:
    if slide.slide_role == slide.slide_role.TITLE:
        return "cover"
    if slide.slide_role == slide.slide_role.SECTION_DIVIDER:
        return "section-divider"
    pattern_id = pattern.pattern_id if pattern is not None else slide.layout_pattern_id
    explicit_map = {
        "cover": "cover",
        "cover-signal": "cover",
        "section-divider": "section-divider",
        "section-divider-band": "section-divider",
        "concept-explainer": "concept-explainer",
        "framework-grid": "concept-explainer",
        "definition-theorem": "definition-theorem",
        "title-thesis-body": "definition-theorem",
        "comparison": "comparison",
        "evidence-table": "comparison",
        "process-flow": "process-flow",
        "agenda-roadmap": "process-flow",
        "worked-example": "worked-example",
        "title-chart-insight": "worked-example",
        "title-visual-caption": "worked-example",
        "summary": "summary",
        "headline-evidence": "summary",
        "appendix-reference": "appendix-reference",
    }
    return explicit_map.get(pattern_id, _default_layout_family(slide))


def _resolved_layout_family(slide: BlueprintSlide, pattern: LayoutPattern | None) -> tuple[str, list[str]]:
    fallback_family = _default_layout_family(slide)
    if pattern is None:
        return fallback_family, [f"Slide {slide.slide_number}: unknown layout pattern `{slide.layout_pattern_id}`; compiled with `{fallback_family}` fallback."]
    warnings: list[str] = []
    if slide.slide_role not in pattern.slide_roles:
        warnings.append(
            f"Slide {slide.slide_number}: layout pattern `{pattern.pattern_id}` does not approve role `{slide.slide_role.value}`; used `{fallback_family}` fallback."
        )
    if slide.visual_type not in pattern.supported_visual_types:
        warnings.append(
            f"Slide {slide.slide_number}: layout pattern `{pattern.pattern_id}` does not approve visual `{slide.visual_type.value}`; used `{fallback_family}` fallback."
        )
    if warnings:
        return fallback_family, warnings
    return _layout_family(slide, pattern), []


def _rasterize_for_pptx(source_path: Path, raster_dir: Path, root: Path) -> str:
    if source_path.suffix.lower() != ".svg":
        return _display_path(source_path, root)
    raster_dir.mkdir(parents=True, exist_ok=True)
    raster_path = raster_dir / f"{source_path.stem}.png"
    if not raster_path.is_file():
        cairosvg.svg2png(bytestring=source_path.read_bytes(), write_to=str(raster_path))
    return _display_path(raster_path, root)


def _resolve_resources(
    slide_data: SlideIRSlide,
    compile_context: SlideIRCompileContext,
    root: Path,
    raster_dir: Path,
) -> ResolvedResources:
    assets = [
        asset
        for asset in compile_context.assets
        if (asset.slide_id == slide_data.slide_id or (asset.slide_id is None and asset.slide_number == slide_data.slide_number))
        and asset.status in {AssetStatus.APPROVED, AssetStatus.READY}
    ]
    visuals: list[ResolvedVisual] = []
    missing: list[str] = []
    allow_structured_visuals = (
        slide_data.visual_source_preference != VisualSourcePreference.DOCUMENT_CROP
        or slide_data.production_mode in {ProductionMode.STRUCTURED_VISUAL, ProductionMode.HYBRID}
        or AssetKind.STRUCTURED_VISUAL in slide_data.asset_dependency_kinds
    )

    if allow_structured_visuals:
        for record in compile_context.visuals:
            if record.spec.slide_id != slide_data.slide_id and record.spec.slide_number != slide_data.slide_number:
                continue
            resolved_output = _resolve_path(record.output_path, root)
            resolved_data = _resolve_path(record.data_output_path, root)
            output_text = None
            data_text = None
            if resolved_output is not None and resolved_output.is_file():
                output_text = _rasterize_for_pptx(resolved_output, raster_dir, root)
            elif record.output_path is not None:
                fallback_path = _resolve_path(record.fallback_output_path, root)
                if fallback_path is not None and fallback_path.is_file():
                    output_text = _rasterize_for_pptx(fallback_path, raster_dir, root)
                else:
                    missing.append(f"Missing visual output for {record.spec.spec_id}")
            if resolved_data is not None and resolved_data.is_file():
                data_text = _display_path(resolved_data, root)
            elif record.data_output_path is not None:
                missing.append(f"Missing visual data output for {record.spec.spec_id}")
            visuals.append(ResolvedVisual(record=record, output_path=output_text, data_output_path=data_text))

    for asset in assets:
        resolved = _resolve_path(asset.local_path, root)
        if resolved is None or not resolved.is_file():
            missing.append(f"Missing asset file for {asset.asset_id}")
    return ResolvedResources(assets=assets, visuals=visuals, missing_dependencies=missing)


def _source_lane_text(asset: AssetRecord | None, slide: SlideIRSlide) -> str:
    if asset is not None and asset.source_material_refs:
        refs: list[str] = []
        for ref in asset.source_material_refs[:2]:
            label = ref.label
            if ref.page is not None:
                label = f"{label} p.{ref.page}"
            refs.append(label)
        return "Source: " + "; ".join(refs)
    if slide.production_bridge is not None and slide.production_bridge.source_material_refs:
        refs = []
        for ref in slide.production_bridge.source_material_refs[:2]:
            label = ref.label
            if ref.page is not None:
                label = f"{label} p.{ref.page}"
            refs.append(label)
        return "Source: " + "; ".join(refs)
    return "Source: slide-native composition"


def _normalize_visible_text(text: str | None) -> str:
    return re.sub(r"\s+", " ", (text or "")).strip()


def _forbidden_visible_text_reason(text: str | None) -> str | None:
    normalized = _normalize_visible_text(text)
    if not normalized:
        return None
    for pattern, reason in FORBIDDEN_VISIBLE_TEXT_PATTERNS:
        if pattern.search(normalized):
            return reason
    return None


def _sanitize_visible_text(text: str | None) -> str:
    normalized = _normalize_visible_text(text)
    reason = _forbidden_visible_text_reason(normalized)
    if reason is not None:
        raise ValueError(f"Forbidden visible slide text detected ({reason}): {normalized}")
    return normalized


def _optional_visible_text(text: str | None, *, max_chars: int = 140) -> str | None:
    normalized = _normalize_visible_text(text)
    if not normalized:
        return None
    if _forbidden_visible_text_reason(normalized) is not None:
        return None
    if len(normalized) > max_chars:
        return None
    return normalized


def _visible_support_items(slide: BlueprintSlide, limit: int = 3) -> list[str]:
    candidates = slide.core_content or slide.required_evidence_assets or [slide.one_line_takeaway]
    items: list[str] = []
    title_text = _normalize_visible_text(slide.title).lower()
    main_message = _normalize_visible_text(slide.main_message).lower()
    takeaway = _normalize_visible_text(slide.one_line_takeaway).lower()
    for item in candidates:
        text = _normalize_visible_text(item)
        normalized = text.lower()
        if not text or normalized in {title_text, main_message, takeaway}:
            continue
        if text not in items:
            items.append(text)
        if len(items) >= limit:
            break
    fallback = _normalize_visible_text(slide.one_line_takeaway)
    if fallback and fallback.lower() not in {title_text, main_message}:
        return items or [fallback]
    return items


def _same_visible_text(left: str | None, right: str | None) -> bool:
    return _normalize_visible_text(left).lower() == _normalize_visible_text(right).lower()


def _text_duplicates_reference(text: str | None, reference: str | None) -> bool:
    candidate = _normalize_visible_text(text).lower()
    source = _normalize_visible_text(reference).lower()
    if not candidate or not source:
        return False
    if candidate == source:
        return True
    if min(len(candidate), len(source)) < 20:
        return False
    return candidate in source or source in candidate


def _deduped_subtitle(slide_data: BlueprintSlide, *, max_chars: int = 120) -> str | None:
    for candidate in (slide_data.one_line_takeaway, slide_data.main_message):
        text = _optional_visible_text(candidate, max_chars=max_chars)
        if text is None:
            continue
        if _text_duplicates_reference(text, slide_data.title):
            continue
        return text
    return None


def _compact_support_note(slide_data: BlueprintSlide, *, max_chars: int = 120) -> str | None:
    meta_phrases = (
        "use one compact concept card",
        "mapping slides should expose",
        "use the repeated slot",
        "the repeated slot should",
        "use appendix slides",
        "appendix support should cluster",
        "one slide should carry the whole loop",
        "use the arc as structure",
        "mechanism slides should show",
        "use the flow to connect",
        "the second pass should connect",
        "use the bridge to",
        "use the concept pass",
        "keep the limitation anchored",
    )
    for candidate in slide_data.core_content:
        text = _optional_visible_text(candidate, max_chars=max_chars)
        if text is None:
            continue
        if any(phrase in text.lower() for phrase in meta_phrases):
            continue
        if _text_duplicates_reference(text, slide_data.title):
            continue
        return text
    for candidate in (slide_data.main_message, slide_data.one_line_takeaway):
        text = _optional_visible_text(candidate, max_chars=max_chars)
        if text is None:
            continue
        if any(phrase in text.lower() for phrase in meta_phrases):
            continue
        if _text_duplicates_reference(text, slide_data.title):
            continue
        if _same_visible_text(text, slide_data.main_message) and _same_visible_text(slide_data.main_message, slide_data.one_line_takeaway):
            continue
        return text
    return None


def _strip_redundant_lead(text: str, reference: str | None) -> str:
    normalized = _normalize_visible_text(text)
    source = _normalize_visible_text(reference)
    if not normalized or not source:
        return normalized
    pattern = re.compile(rf"^{re.escape(source)}[\s:,-]+", re.IGNORECASE)
    stripped = pattern.sub("", normalized).strip()
    if not stripped:
        return normalized
    return stripped[:1].upper() + stripped[1:]


def _compact_process_body(text: str) -> str:
    normalized = _normalize_visible_text(text)
    replacements = {
        "Write each candidate as a chromosome.": "Encode each candidate.",
        "Decode candidates and score fitness.": "Score fitness.",
        "Bias reproduction toward better candidates.": "Choose fitter parents.",
        "Crossover and mutation create the next population.": "Create offspring.",
        "Stop when improvement or budget runs out.": "Repeat until stop.",
        "Genes, alleles, genotype, phenotype, and natural selection": "Core biology terms.",
        "Translate those roles into representation, fitness, and operators": "Map to encoding and fitness.",
        "Follow selection, crossover, and mutation through one loop": "Trace one operator loop.",
        "Use a toy population to see one generation change state": "Test on one toy population.",
        "End with limitations, tradeoffs, and practical use cases": "Close on limits and use cases.",
    }
    if normalized in replacements:
        return replacements[normalized]
    return normalized


def _tailored_process_steps(slide_data: BlueprintSlide, steps: list[dict[str, str]]) -> list[dict[str, str]]:
    title = _normalize_visible_text(slide_data.title).lower()
    if slide_data.slide_archetype == SlideArchetype.STEP_BY_STEP_MECHANISM:
        if "selection" in title:
            return [
                {"label": "Filter", "body": "Keep fitter parents in play."},
                {"label": "Reweight", "body": "Increase strong strings in the pool."},
                {"label": "Preserve", "body": "Avoid collapsing diversity too early."},
                {"label": "Shift state", "body": "The candidate pool changes next round."},
            ]
        if "crossover" in title:
            return [
                {"label": "Align", "body": "Choose compatible parent strings."},
                {"label": "Cut", "body": "Swap one segment across parents."},
                {"label": "Preserve", "body": "Keep useful partial structure intact."},
                {"label": "Test child", "body": "Check the offspring against fitness."},
            ]
        if "mutation" in title:
            return [
                {"label": "Pick locus", "body": "Target one local choice."},
                {"label": "Flip value", "body": "Change a small part of the string."},
                {"label": "Reopen search", "body": "Recover options crossover missed."},
                {"label": "Re-score", "body": "Compare the changed candidate again."},
            ]
        if "how one toy population changes" in title:
            return [
                {"label": "Start pool", "body": "Read the candidate strings first."},
                {"label": "Rank", "body": "Score and order the current pool."},
                {"label": "Spawn", "body": "Create offspring from the winners."},
                {"label": "Compare", "body": "Check how the next state differs."},
            ]
        if "decoding links representation to selection" in title:
            return [
                {"label": "Decode", "body": "Turn the string into behavior."},
                {"label": "Score", "body": "Measure the behavior against fitness."},
                {"label": "Compare", "body": "Rank candidates by the scores."},
                {"label": "Select", "body": "Use the scores to bias reproduction."},
            ]
    tailored: list[dict[str, str]] = []
    for step in steps:
        tailored.append(
            {
                "label": _normalize_visible_text(step["label"]),
                "body": _compact_process_body(step["body"]),
            }
        )
    if slide_data.slide_archetype == SlideArchetype.PROCESS_FLOW and len(tailored) > 4:
        return tailored[:4]
    return tailored


def _add_compact_support_note(
    slide,
    slide_data: BlueprintSlide,
    style: CompilerStyle,
    *,
    left: float,
    top: float,
    width: float,
    height: float = 0.42,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> bool:
    note = _compact_support_note(slide_data)
    if note is None:
        return False
    _add_textbox(
        slide,
        left,
        top,
        width,
        height,
        note,
        font_name=style.fonts["caption"].font_family,
        font_size=style.fonts["caption"].size_pt,
        color=style.colors["ink"],
        align=align,
    )
    return True


def _note_text(slide: SlideIRSlide, notes_map: dict[str, str]) -> str | None:
    values = [
        notes_map.get(slide.slide_id),
        notes_map.get(str(slide.slide_number)),
        slide.presenter_notes,
    ]
    merged: list[str] = []
    for value in values:
        if value and value not in merged:
            merged.append(value)
    if not merged:
        return None
    return "\n\n".join(merged)


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_name: str,
    font_size: float,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    visible_text = _sanitize_visible_text(text)
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = vertical_anchor
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = visible_text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color


def _fill_shape(shape, color: RGBColor, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()


def _style_prior_mask(style_prior: SlideIRStylePrior | None, mask_id: str | None, *, applies_to: str) -> SlideIRSafeAreaMask | None:
    if style_prior is None:
        return None
    if mask_id is not None:
        for mask in style_prior.safe_area_masks:
            if mask.mask_id == mask_id:
                return mask
    return next((mask for mask in style_prior.safe_area_masks if mask.applies_to == applies_to), None)


def _style_prior_color(style: CompilerStyle, layer: SlideIRBackgroundLayer) -> RGBColor:
    return style.colors.get(layer.color_token, style.colors.get("signal", RGBColor(194, 65, 12)))


def _style_prior_transparency(layer: SlideIRBackgroundLayer) -> float:
    return max(0.0, min(1.0, 1.0 - layer.opacity))


def _add_style_prior_background_layer(slide, style: CompilerStyle, layer: SlideIRBackgroundLayer) -> None:
    mask = _style_prior_mask(style.style_prior, layer.safe_area_mask_id, applies_to="background" if layer.layer_type != "motif" else "motif")
    top_in = mask.top_in if mask is not None else 0.5
    right_in = mask.right_in if mask is not None else 0.5
    bottom_in = mask.bottom_in if mask is not None else 0.45
    left_in = mask.left_in if mask is not None else 0.5

    shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
    left = 0.0
    top = 0.0
    width = style.slide_width
    height = style.slide_height
    if layer.shape_hint == "top-band":
        top = 0.0
        height = min(max(top_in * 0.72, 0.16), style.slide_height * 0.26)
    elif layer.shape_hint == "right-edge":
        width = min(max(right_in * 0.9, 0.22), style.slide_width * 0.18)
        left = style.slide_width - width
    elif layer.shape_hint == "bottom-band":
        height = min(max(bottom_in * 0.72, 0.16), style.slide_height * 0.26)
        top = style.slide_height - height
    elif layer.shape_hint == "corner-accent":
        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
        width = min(max((left_in + right_in) * 0.42, 0.8), style.slide_width * 0.22)
        height = min(max((top_in + bottom_in) * 0.38, 0.62), style.slide_height * 0.2)
        left = style.slide_width - width - (right_in * 0.16)
        top = style.slide_height - height - (bottom_in * 0.18)

    decorative_shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    _fill_shape(decorative_shape, _style_prior_color(style, layer), _style_prior_transparency(layer))


def _add_style_prior_background_layers(slide, style: CompilerStyle, *, appendix: bool) -> None:
    if style.style_prior is None:
        return
    for layer in style.style_prior.background_layers:
        if layer.scope == "appendix" and not appendix:
            continue
        _add_style_prior_background_layer(slide, style, layer)


def _add_background(slide, style: CompilerStyle, *, appendix: bool = False) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(style.slide_width),
        Inches(style.slide_height),
    )
    _fill_shape(panel, style.colors["canvas"], 0.0)
    _add_style_prior_background_layers(slide, style, appendix=appendix)
    if appendix:
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(style.slide_width),
            Inches(0.18),
        )
        _fill_shape(band, style.colors["signal"], 0.1)


def _add_title_block(
    slide,
    slide_data: BlueprintSlide,
    style: CompilerStyle,
    *,
    show_default_subtitle: bool = True,
    subtitle: str | None = None,
    eyebrow: str | None = None,
) -> None:
    title_font = style.fonts["title"]
    caption_font = style.fonts["caption"]
    title_top = style.margin * 0.72
    title_height = 0.82
    if eyebrow:
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(style.margin),
            Inches(style.margin * 0.38),
            Inches(2.15),
            Inches(0.26),
        )
        _fill_shape(pill, style.colors["signal"], 0.0)
        _add_textbox(
            slide,
            style.margin + 0.08,
            style.margin * 0.39,
            1.99,
            0.2,
            eyebrow,
            font_name=caption_font.font_family,
            font_size=caption_font.size_pt - 1,
            color=RGBColor(255, 255, 255),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        title_top = style.margin * 0.82
    _add_textbox(
        slide,
        style.margin,
        title_top,
        style.slide_width - (style.margin * 2),
        title_height,
        slide_data.title,
        font_name=title_font.font_family,
        font_size=title_font.size_pt,
        color=style.colors["ink"],
        bold=True,
    )
    subtitle_text = subtitle if subtitle is not None else (_deduped_subtitle(slide_data) if show_default_subtitle else None)
    if subtitle_text:
        _add_textbox(
            slide,
            style.margin,
            style.margin + 0.48,
            style.slide_width - (style.margin * 2),
            0.34,
            subtitle_text,
            font_name=caption_font.font_family,
            font_size=caption_font.size_pt,
            color=style.colors["ink"],
        )


def _add_footer(slide, number_label: str, deck_mode: DeckMode, style: CompilerStyle) -> None:
    caption_font = style.fonts["caption"]
    if deck_mode == DeckMode.APPENDIX:
        tag = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(style.margin),
            Inches(style.slide_height - 0.48),
            Inches(0.95),
            Inches(0.28),
        )
        _fill_shape(tag, style.colors["signal"], 0.0)
        _add_textbox(
            slide,
            style.margin + 0.06,
            style.slide_height - 0.455,
            0.82,
            0.22,
            "Appendix",
            font_name=caption_font.font_family,
            font_size=caption_font.size_pt - 1,
            color=RGBColor(255, 255, 255),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _add_textbox(
        slide,
        style.slide_width - 0.9,
        style.slide_height - 0.46,
        0.5,
        0.22,
        number_label,
        font_name=caption_font.font_family,
        font_size=caption_font.size_pt,
        color=style.colors["ink"],
        align=PP_ALIGN.RIGHT,
    )


def _set_notes(slide, text: str | None) -> bool:
    if not text:
        return False
    slide.notes_slide.notes_text_frame.text = text
    return True


def _add_placeholder_panel(slide, message: str, style: CompilerStyle, left: float, top: float, width: float, height: float) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _fill_shape(panel, style.colors["canvas"], 0.0)
    panel.line.color.rgb = style.colors["ink"]
    panel.line.transparency = 0.25
    _add_textbox(
        slide,
        left + 0.22,
        top + 0.18,
        width - 0.44,
        height - 0.36,
        message,
        font_name=style.fonts["body"].font_family,
        font_size=style.fonts["body"].size_pt,
        color=style.colors["ink"],
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def _add_support_panel(
    slide,
    slide_data: BlueprintSlide,
    style: CompilerStyle,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    heading: str = "Teaching cues",
) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _fill_shape(panel, style.colors["canvas"], 0.0)
    panel.line.color.rgb = style.colors["ink"]
    panel.line.transparency = 0.15
    _add_textbox(
        slide,
        left + 0.2,
        top + 0.16,
        width - 0.4,
        0.28,
        heading,
        font_name=style.fonts["caption"].font_family,
        font_size=style.fonts["caption"].size_pt,
        color=style.colors["signal"],
        bold=True,
    )
    _add_bullets(slide, _visible_support_items(slide_data, limit=3), style, left + 0.2, top + 0.48, width - 0.4, height - 0.64)


def _add_message_panel(slide, heading: str, body: str, style: CompilerStyle, left: float, top: float, width: float, height: float) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _fill_shape(panel, style.colors["canvas"], 0.0)
    panel.line.color.rgb = style.colors["signal"]
    panel.line.width = Pt(1.3)
    _add_textbox(
        slide,
        left + 0.2,
        top + 0.16,
        width - 0.4,
        0.26,
        heading,
        font_name=style.fonts["caption"].font_family,
        font_size=style.fonts["caption"].size_pt,
        color=style.colors["signal"],
        bold=True,
    )
    _add_textbox(
        slide,
        left + 0.2,
        top + 0.45,
        width - 0.4,
        height - 0.58,
        body,
        font_name=style.fonts["body"].font_family,
        font_size=style.fonts["body"].size_pt,
        color=style.colors["ink"],
    )


def _add_bullets(slide, items: list[str], style: CompilerStyle, left: float, top: float, width: float, height: float) -> None:
    if not items:
        return
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, item in enumerate(items[:3]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = f"- {_sanitize_visible_text(item)}"
        run.font.name = style.fonts["body"].font_family
        run.font.size = Pt(style.fonts["body"].size_pt)
        run.font.color.rgb = style.colors["ink"]


def _fit_picture(picture_path: Path, width: float, height: float) -> tuple[float, float]:
    try:
        with Image.open(picture_path) as image:
            image_width, image_height = image.size
    except Exception:
        return (width, height)
    if image_width <= 0 or image_height <= 0:
        return (width, height)
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio >= box_ratio:
        fitted_width = width
        fitted_height = width / image_ratio
    else:
        fitted_height = height
        fitted_width = height * image_ratio
    return (fitted_width, fitted_height)


def _add_picture(slide, picture_path: Path, style: CompilerStyle, left: float, top: float, width: float, height: float) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _fill_shape(frame, style.colors["canvas"], 0.0)
    frame.line.color.rgb = style.colors["ink"]
    frame.line.transparency = 0.55
    fitted_width, fitted_height = _fit_picture(picture_path, width, height)
    offset_left = left + ((width - fitted_width) / 2)
    offset_top = top + ((height - fitted_height) / 2)
    slide.shapes.add_picture(
        str(picture_path),
        Inches(offset_left),
        Inches(offset_top),
        width=Inches(fitted_width),
        height=Inches(fitted_height),
    )


def _assert_no_forbidden_visible_text(slide, slide_number: int) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _normalize_visible_text(shape.text_frame.text)
        reason = _forbidden_visible_text_reason(text)
        if reason is not None:
            raise ValueError(f"Slide {slide_number} contains forbidden visible text ({reason}): {text}")


def _parse_tsv(path: Path) -> list[list[str]]:
    rows = [line.rstrip("\n").split("\t") for line in path.read_text(encoding="utf-8").splitlines() if line.strip()]
    if not rows:
        return []
    width = max(len(row) for row in rows)
    return [row + [""] * (width - len(row)) for row in rows]


def _render_native_table(slide, table_path: Path, style: CompilerStyle, left: float, top: float, width: float, height: float) -> bool:
    rows = _parse_tsv(table_path)
    if len(rows) < 2:
        return False
    if len(rows) > 8:
        header = rows[0]
        body = rows[1:7] + [["...", *[""] * (len(header) - 1)]]
        rows = [header, *body]
    row_count = len(rows)
    column_count = len(rows[0])
    shape = slide.shapes.add_table(row_count, column_count, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for column_index in range(column_count):
        table.columns[column_index].width = Inches(width / column_count)
    for row_index in range(row_count):
        table.rows[row_index].height = Inches(height / row_count)
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            fill_color = style.colors["signal"] if row_index == 0 else style.colors["canvas"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            frame = cell.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = _sanitize_visible_text(value)
            run.font.name = style.fonts["body"].font_family
            run.font.size = Pt(style.fonts["body"].size_pt)
            run.font.bold = row_index == 0
            run.font.color.rgb = RGBColor(255, 255, 255) if row_index == 0 else style.colors["ink"]
    return True


def _authoring_rows(slide_data: BlueprintSlide) -> list[list[str]]:
    payload_rows = slide_data.authoring_payload.get("rows", [])
    if not isinstance(payload_rows, list):
        return []
    rows: list[list[str]] = []
    for row in payload_rows:
        if isinstance(row, list):
            rows.append([_sanitize_visible_text(str(value)) for value in row])
    return rows


def _authoring_columns(slide_data: BlueprintSlide) -> list[str]:
    payload_columns = slide_data.authoring_payload.get("columns", [])
    if not isinstance(payload_columns, list):
        return []
    return [_sanitize_visible_text(str(value)) for value in payload_columns]


def _authoring_cards(slide_data: BlueprintSlide) -> list[dict[str, str]]:
    payload_cards = slide_data.authoring_payload.get("cards", [])
    if not isinstance(payload_cards, list):
        return []
    cards: list[dict[str, str]] = []
    for card in payload_cards:
        if isinstance(card, dict):
            label = _sanitize_visible_text(str(card.get("label", "")))
            body = _sanitize_visible_text(str(card.get("body", "")))
            if label or body:
                cards.append({"label": label, "body": body})
    return cards


def _authoring_steps(slide_data: BlueprintSlide) -> list[dict[str, str]]:
    payload_steps = slide_data.authoring_payload.get("steps", [])
    if not isinstance(payload_steps, list):
        return []
    steps: list[dict[str, str]] = []
    for step in payload_steps:
        if isinstance(step, dict):
            label = _sanitize_visible_text(str(step.get("label", "")))
            body = _sanitize_visible_text(str(step.get("body", "")))
            if label or body:
                steps.append({"label": label, "body": body})
    return steps


def _render_authored_table(
    slide,
    *,
    columns: list[str],
    rows: list[list[str]],
    style: CompilerStyle,
    left: float,
    top: float,
    width: float,
    height: float,
    column_widths: list[float] | None = None,
    highlight_first_column: bool = False,
) -> bool:
    if not columns or not rows:
        return False
    row_count = len(rows) + 1
    column_count = len(columns)
    shape = slide.shapes.add_table(row_count, column_count, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    if column_widths is not None and len(column_widths) == column_count and sum(column_widths) > 0:
        total = sum(column_widths)
        resolved_widths = [width * (value / total) for value in column_widths]
    else:
        resolved_widths = [width / column_count for _ in range(column_count)]
    for column_index, column_width in enumerate(resolved_widths):
        table.columns[column_index].width = Inches(column_width)
    for row_index in range(row_count):
        table.rows[row_index].height = Inches(height / row_count)
    for row_index, row_values in enumerate([columns, *rows]):
        for column_index, value in enumerate(row_values[:column_count]):
            cell = table.cell(row_index, column_index)
            highlight_column = highlight_first_column and row_index > 0 and column_index == 0
            fill_color = style.colors["signal"] if row_index == 0 or highlight_column else style.colors["canvas"]
            text_color = RGBColor(255, 255, 255) if row_index == 0 or highlight_column else style.colors["ink"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.text = ""
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = _sanitize_visible_text(str(value))
            run.font.name = style.fonts["body"].font_family
            run.font.size = Pt(style.fonts["body"].size_pt - (1 if column_count >= 4 else 0))
            run.font.bold = row_index == 0
            run.font.color.rgb = text_color
    return True


def _required_table_payload(
    slide_data: BlueprintSlide,
    *,
    archetype_label: str,
    min_columns: int,
    min_rows: int,
) -> tuple[list[str], list[list[str]]]:
    columns = _authoring_columns(slide_data)
    rows = _authoring_rows(slide_data)
    if len(columns) < min_columns or len(rows) < min_rows:
        raise ValueError(
            f"Slide {slide_data.slide_number} requires a realized {archetype_label} with at least "
            f"{min_columns} columns and {min_rows} rows."
        )
    return columns, rows


def _require_labeled_table_axes(
    columns: list[str],
    rows: list[list[str]],
    *,
    slide_number: int,
    archetype_label: str,
) -> None:
    normalized_columns = [_normalize_visible_text(value) for value in columns]
    if any(not value for value in normalized_columns):
        raise ValueError(f"Slide {slide_number} {archetype_label} requires non-empty column labels.")
    lowered_columns = [value.lower() for value in normalized_columns]
    if len(lowered_columns) != len(set(lowered_columns)):
        raise ValueError(f"Slide {slide_number} {archetype_label} requires distinct comparison or mapping axes.")
    for row in rows:
        if not row:
            raise ValueError(f"Slide {slide_number} {archetype_label} requires labeled correspondence rows.")
        if not _normalize_visible_text(row[0]):
            raise ValueError(f"Slide {slide_number} {archetype_label} requires a labeled first-column row header.")


def _required_process_steps(slide_data: BlueprintSlide, *, min_steps: int = 3) -> list[dict[str, str]]:
    steps = _authoring_steps(slide_data)
    if len(steps) < min_steps:
        raise ValueError(f"Slide {slide_data.slide_number} requires at least {min_steps} ordered process steps.")
    return steps


def _require_worked_example_progression(rows: list[list[str]], slide_number: int) -> None:
    joined_rows = [" ".join(row).lower() for row in rows]
    required_groups = {
        "initial population state": ("start", "initial", "population"),
        "evaluation or selection state": ("fitness", "evaluate", "selection", "score"),
        "offspring or next-generation state": ("next", "offspring", "mutation", "new state"),
    }
    missing = [
        label
        for label, keywords in required_groups.items()
        if not any(any(keyword in row_text for keyword in keywords) for row_text in joined_rows)
    ]
    if missing:
        raise ValueError(f"Slide {slide_number} worked example is missing {', '.join(missing)}.")


def _worked_example_state_cards(rows: list[list[str]]) -> list[dict[str, str]]:
    visible_rows = rows[:4]
    cards: list[dict[str, str]] = []
    for index, row in enumerate(visible_rows):
        label = _sanitize_visible_text(row[0]) if row else f"State {index + 1}"
        label_lower = label.lower()
        if index == 0:
            body = "Start with a readable pool."
        elif any(keyword in label_lower for keyword in ("fitness", "evaluate", "score", "rank")):
            body = "Score and rank before selection."
        elif any(keyword in label_lower for keyword in ("selection", "crossover", "offspring", "recombine", "spawn")):
            body = "Create offspring from strong parents."
        elif index == len(visible_rows) - 1 or any(keyword in label_lower for keyword in ("mutation", "next", "update")):
            body = "Carry the changed pool forward."
        else:
            body = "Track how the pool changes here."
        cards.append({"label": label, "body": body})
    return cards


def _render_sequence_cards(
    slide,
    *,
    cards: list[dict[str, str]],
    style: CompilerStyle,
    left: float,
    top: float,
    width: float,
    height: float,
    numbered: bool = False,
) -> None:
    visible_cards = cards[:5]
    if not visible_cards:
        return
    connector_space = 0.28 if len(visible_cards) > 1 else 0.0
    card_width = (width - ((style.gap + connector_space) * max(len(visible_cards) - 1, 0))) / len(visible_cards)
    for index, card in enumerate(visible_cards):
        card_left = left + (index * (card_width + style.gap + connector_space))
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(card_left),
            Inches(top),
            Inches(card_width),
            Inches(height),
        )
        _fill_shape(panel, style.colors["canvas"], 0.0)
        panel.line.color.rgb = style.colors["signal"]
        panel.line.width = Pt(1.2)
        label_left = card_left + 0.18
        label_width = card_width - 0.36
        if numbered:
            badge = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(card_left + 0.14),
                Inches(top + 0.14),
                Inches(0.34),
                Inches(0.34),
            )
            _fill_shape(badge, style.colors["signal"], 0.0)
            _add_textbox(
                slide,
                card_left + 0.14,
                top + 0.16,
                0.34,
                0.26,
                str(index + 1),
                font_name=style.fonts["caption"].font_family,
                font_size=style.fonts["caption"].size_pt - 1,
                color=RGBColor(255, 255, 255),
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            label_left = card_left + 0.56
            label_width = card_width - 0.72
        _add_textbox(
            slide,
            label_left,
            top + 0.16,
            label_width,
            0.32,
            card["label"],
            font_name=style.fonts["caption"].font_family,
            font_size=style.fonts["caption"].size_pt,
            color=style.colors["signal"],
            bold=True,
            align=PP_ALIGN.LEFT,
        )
        _add_textbox(
            slide,
            card_left + 0.18,
            top + 0.54,
            card_width - 0.36,
            height - 0.72,
            card["body"],
            font_name=style.fonts["body"].font_family,
            font_size=style.fonts["body"].size_pt - 1,
            color=style.colors["ink"],
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        if index >= len(visible_cards) - 1:
            continue
        connector_left = card_left + card_width + 0.02
        connector_y = top + (height / 2)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(connector_left),
            Inches(connector_y),
            Inches(connector_left + connector_space + style.gap - 0.04),
            Inches(connector_y),
        )
        connector.line.color.rgb = style.colors["signal"]
        connector.line.width = Pt(1.8)
        chevron = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(connector_left + 0.06),
            Inches(connector_y - 0.18),
            Inches(max(connector_space + style.gap - 0.12, 0.12)),
            Inches(0.36),
        )
        _fill_shape(chevron, style.colors["signal"], 0.0)
        chevron.line.fill.background()


def _render_stacked_cards(
    slide,
    *,
    cards: list[dict[str, str]],
    style: CompilerStyle,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    visible_cards = cards[:3]
    if not visible_cards:
        return
    card_height = (height - (style.gap * max(len(visible_cards) - 1, 0))) / len(visible_cards)
    for index, card in enumerate(visible_cards):
        card_top = top + (index * (card_height + style.gap))
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(card_top),
            Inches(width),
            Inches(card_height),
        )
        _fill_shape(panel, style.colors["canvas"], 0.0)
        panel.line.color.rgb = style.colors["ink"]
        panel.line.transparency = 0.18
        _add_textbox(
            slide,
            left + 0.14,
            card_top + 0.16,
            width - 0.28,
            0.28,
            card["label"],
            font_name=style.fonts["caption"].font_family,
            font_size=style.fonts["caption"].size_pt,
            color=style.colors["signal"],
            bold=True,
        )
        _add_textbox(
            slide,
            left + 0.14,
            card_top + 0.5,
            width - 0.28,
            card_height - 0.66,
            card["body"],
            font_name=style.fonts["body"].font_family,
            font_size=style.fonts["body"].size_pt - 1,
            color=style.colors["ink"],
        )


def _compose_correspondence_matrix(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="correspondence matrix",
        min_columns=3,
        min_rows=2,
    )
    _require_labeled_table_axes(columns, rows, slide_number=slide_data.slide_number, archetype_label="correspondence matrix")
    _render_authored_table(
        slide,
        columns=columns,
        rows=rows,
        style=style,
        left=style.margin,
        top=1.48,
        width=style.slide_width - (style.margin * 2),
        height=4.9,
        column_widths=[0.22, 0.24, 0.54],
        highlight_first_column=True,
    )
    _add_compact_support_note(slide, slide_data, style, left=style.margin, top=6.52, width=style.slide_width - (style.margin * 2))
    return []


def _compose_two_column_mapping_table(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="mapping table",
        min_columns=2,
        min_rows=2,
    )
    _require_labeled_table_axes(columns, rows, slide_number=slide_data.slide_number, archetype_label="mapping table")
    _render_authored_table(
        slide,
        columns=columns,
        rows=rows,
        style=style,
        left=style.margin,
        top=1.48,
        width=style.slide_width - (style.margin * 2),
        height=4.85,
        column_widths=[0.2, 0.24, 0.56] if len(columns) == 3 else None,
        highlight_first_column=True,
    )
    _add_compact_support_note(slide, slide_data, style, left=style.margin, top=6.52, width=style.slide_width - (style.margin * 2))
    return []


def _compose_comparison_matrix(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=slide_data.deck_mode == DeckMode.APPENDIX)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="comparison matrix",
        min_columns=3,
        min_rows=2,
    )
    _require_labeled_table_axes(columns, rows, slide_number=slide_data.slide_number, archetype_label="comparison matrix")
    _render_authored_table(
        slide,
        columns=columns,
        rows=rows,
        style=style,
        left=style.margin,
        top=1.48,
        width=style.slide_width - (style.margin * 2),
        height=4.7,
        column_widths=[0.18, 0.23, 0.31, 0.28] if len(columns) == 4 else None,
        highlight_first_column=True,
    )
    _add_compact_support_note(slide, slide_data, style, left=style.margin, top=6.36, width=style.slide_width - (style.margin * 2))
    return []


def _compose_worked_example_state_table(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="worked example state table",
        min_columns=3,
        min_rows=3,
    )
    _require_labeled_table_axes(columns, rows, slide_number=slide_data.slide_number, archetype_label="worked example state table")
    _require_worked_example_progression(rows, slide_data.slide_number)
    state_cards = _worked_example_state_cards(rows)
    _render_sequence_cards(
        slide,
        cards=state_cards,
        style=style,
        left=style.margin,
        top=1.56,
        width=style.slide_width - (style.margin * 2),
        height=1.5,
        numbered=True,
    )
    _render_authored_table(
        slide,
        columns=columns,
        rows=rows,
        style=style,
        left=style.margin,
        top=3.38,
        width=style.slide_width - (style.margin * 2),
        height=2.7,
        column_widths=[0.22, 0.4, 0.38],
        highlight_first_column=True,
    )
    if len(rows) <= 3:
        _add_compact_support_note(slide, slide_data, style, left=style.margin, top=6.34, width=style.slide_width - (style.margin * 2))
    return []


def _compose_authored_process(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    steps = _tailored_process_steps(slide_data, _required_process_steps(slide_data))
    _render_sequence_cards(
        slide,
        cards=steps,
        style=style,
        left=style.margin,
        top=1.86,
        width=style.slide_width - (style.margin * 2),
        height=1.7,
        numbered=True,
    )
    _add_compact_support_note(slide, slide_data, style, left=style.margin, top=5.1, width=style.slide_width - (style.margin * 2), align=PP_ALIGN.CENTER)
    return []


def _compose_limitation_pitfall_callout(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    lead_body = _optional_visible_text(slide_data.main_message, max_chars=150) or _compact_support_note(slide_data, max_chars=150)
    lead = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(style.margin),
        Inches(1.72),
        Inches(7.25),
        Inches(3.25),
    )
    _fill_shape(lead, style.colors["canvas"], 0.0)
    lead.line.color.rgb = style.colors["signal"]
    lead.line.width = Pt(1.6)
    _add_textbox(
        slide,
        style.margin + 0.2,
        1.92,
        6.85,
        0.34,
        "Main limitation",
        font_name=style.fonts["caption"].font_family,
        font_size=style.fonts["caption"].size_pt,
        color=style.colors["signal"],
        bold=True,
    )
    _add_textbox(
        slide,
        style.margin + 0.2,
        2.36,
        6.85,
        2.18,
        lead_body,
        font_name=style.fonts["body"].font_family,
        font_size=style.fonts["body"].size_pt + 1,
        color=style.colors["ink"],
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    side_cards = _authoring_cards(slide_data) or [{"label": f"Guardrail {index + 1}", "body": item} for index, item in enumerate(_visible_support_items(slide_data, limit=2))]
    _render_stacked_cards(slide, cards=side_cards[:2], style=style, left=8.12, top=1.9, width=4.66, height=3.05)
    return []


def _compose_appendix_evidence_cluster(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=True)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="appendix evidence cluster",
        min_columns=2,
        min_rows=1,
    )
    source_cards = [
        {
            "label": row[0] if row else f"Source {index + 1}",
            "body": " ".join(
                part
                for part in [
                    row[1] if len(row) > 1 else "",
                    row[2] if len(row) > 2 else "",
                ]
                if part
            ),
        }
        for index, row in enumerate(rows[:3])
    ]
    if len(rows) == 1:
        _add_message_panel(
            slide,
            columns[0] if columns else "Evidence anchor",
            rows[0][1] if len(rows[0]) > 1 else rows[0][0],
            style,
            style.margin,
            1.5,
            7.85,
            4.85,
        )
    else:
        table_columns = columns[:2] if len(columns) >= 2 else columns
        table_rows = [row[:2] for row in rows]
        _render_authored_table(
            slide,
            columns=table_columns,
            rows=table_rows,
            style=style,
            left=style.margin,
            top=1.5,
            width=7.85,
            height=4.85,
            column_widths=[0.28, 0.72] if len(table_columns) == 2 else None,
            highlight_first_column=True,
        )
    _render_stacked_cards(slide, cards=source_cards, style=style, left=8.72, top=1.64, width=4.06, height=4.45)
    return []


def _compose_appendix_themed_evidence_cluster(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=True)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    cards = _authoring_cards(slide_data)
    if not cards:
        cards = [{"label": f"Anchor {index + 1}", "body": item} for index, item in enumerate(_visible_support_items(slide_data))]
    _add_message_panel(
        slide,
        "Theme anchor",
        _strip_redundant_lead(slide_data.main_message, slide_data.title),
        style,
        style.margin,
        1.56,
        4.05,
        4.88,
    )
    _render_sequence_cards(
        slide,
        cards=cards,
        style=style,
        left=4.86,
        top=1.64,
        width=style.slide_width - 4.86 - style.margin,
        height=4.7,
        numbered=False,
    )
    return []


def _compose_appendix_source_location_matrix(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=True)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="appendix source-location matrix",
        min_columns=3,
        min_rows=1,
    )
    _render_authored_table(
        slide,
        columns=columns[:3],
        rows=[row[:3] for row in rows],
        style=style,
        left=style.margin,
        top=1.5,
        width=8.2,
        height=4.9,
        column_widths=[0.24, 0.26, 0.5],
        highlight_first_column=True,
    )
    _add_message_panel(
        slide,
        "Source trail",
        _strip_redundant_lead(slide_data.main_message, slide_data.title),
        style,
        8.98,
        1.58,
        style.slide_width - 8.98 - style.margin,
        4.78,
    )
    return []


def _compose_appendix_annotated_excerpt_cluster(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=True)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    cards = _authoring_cards(slide_data)
    if not cards:
        cards = [{"label": "Source anchor", "body": item} for item in _visible_support_items(slide_data)]
    visible_cards = cards[:3]
    card_width = (style.slide_width - (style.margin * 2) - (style.gap * max(len(visible_cards) - 1, 0))) / max(len(visible_cards), 1)
    for index, card in enumerate(visible_cards):
        left = style.margin + (index * (card_width + style.gap))
        _add_message_panel(
            slide,
            card["label"],
            card["body"],
            style,
            left,
            1.72,
            card_width,
            4.6,
        )
    return []


def _compose_appendix_comparison_evidence_cluster(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=True)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="appendix comparison evidence cluster",
        min_columns=3,
        min_rows=2,
    )
    _render_authored_table(
        slide,
        columns=columns[:3],
        rows=[row[:3] for row in rows],
        style=style,
        left=style.margin,
        top=1.5,
        width=style.slide_width - (style.margin * 2),
        height=3.75,
        column_widths=[0.18, 0.41, 0.41],
        highlight_first_column=True,
    )
    _add_message_panel(
        slide,
        "Comparison reading",
        _strip_redundant_lead(slide_data.main_message, slide_data.title),
        style,
        style.margin,
        5.52,
        style.slide_width - (style.margin * 2),
        1.02,
    )
    return []


def _compose_appendix_source_map(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=True)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    columns, rows = _required_table_payload(
        slide_data,
        archetype_label="appendix source map",
        min_columns=3,
        min_rows=1,
    )
    _render_authored_table(
        slide,
        columns=columns[:3],
        rows=[row[:3] for row in rows],
        style=style,
        left=style.margin,
        top=1.5,
        width=style.slide_width - (style.margin * 2),
        height=3.95,
        column_widths=[0.24, 0.28, 0.48],
        highlight_first_column=True,
    )
    cards = [
        {
            "label": row[0] if row else f"Anchor {index + 1}",
            "body": " ".join(part for part in row[1:3] if part),
        }
        for index, row in enumerate(rows[:3])
    ]
    _render_sequence_cards(
        slide,
        cards=cards,
        style=style,
        left=style.margin,
        top=5.7,
        width=style.slide_width - (style.margin * 2),
        height=0.94,
        numbered=False,
    )
    return []


def _compose_authored_cards(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style, appendix=slide_data.deck_mode == DeckMode.APPENDIX)
    _add_title_block(slide, slide_data, style, show_default_subtitle=False)
    cards = _authoring_cards(slide_data)
    if not cards:
        cards = [{"label": f"Point {index + 1}", "body": item} for index, item in enumerate(_visible_support_items(slide_data))]
    card_count = min(len(cards), 3)
    card_width = (style.slide_width - (style.margin * 2) - (style.gap * max(card_count - 1, 0))) / max(card_count, 1)
    for index, card in enumerate(cards[:3]):
        left = style.margin + (index * (card_width + style.gap))
        label = card["label"]
        if _text_duplicates_reference(label, slide_data.title):
            label = "Core definition" if index == 0 else f"Point {index + 1}"
        body = _strip_redundant_lead(card["body"], slide_data.title)
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.92), Inches(card_width), Inches(2.75))
        _fill_shape(panel, style.colors["canvas"], 0.0)
        panel.line.color.rgb = style.colors["ink"]
        panel.line.transparency = 0.2
        _add_textbox(
            slide,
            left + 0.16,
            2.1,
            card_width - 0.32,
            0.28,
            label,
            font_name=style.fonts["caption"].font_family,
            font_size=style.fonts["caption"].size_pt,
            color=style.colors["signal"],
            bold=True,
        )
        _add_textbox(
            slide,
            left + 0.16,
            2.46,
            card_width - 0.32,
            1.85,
            body,
            font_name=style.fonts["body"].font_family,
            font_size=style.fonts["body"].size_pt,
            color=style.colors["ink"],
        )
    _add_compact_support_note(slide, slide_data, style, left=style.margin, top=5.28, width=style.slide_width - (style.margin * 2))
    return []


def _resolve_asset_picture_path(asset: AssetRecord, style: CompilerStyle) -> Path | None:
    resolved = _resolve_path(asset.local_path, style.root)
    if resolved is None or not resolved.is_file():
        return None
    if resolved.suffix.lower() == ".svg":
        rasterized = _rasterize_for_pptx(resolved, style.raster_dir, style.root)
        return _resolve_path(rasterized, style.root)
    return resolved


def _first_asset_by_kind(resources: ResolvedResources, style: CompilerStyle, kinds: set[AssetKind]) -> tuple[AssetRecord | None, Path | None]:
    for asset in resources.assets:
        if asset.asset_kind not in kinds:
            continue
        resolved = _resolve_asset_picture_path(asset, style)
        if resolved is not None and resolved.is_file():
            return (asset, resolved)
    return (None, None)


def _first_existing_asset(resources: ResolvedResources, style: CompilerStyle) -> tuple[AssetRecord | None, Path | None]:
    return _first_asset_by_kind(resources, style, {AssetKind.DOCUMENT_CROP, AssetKind.IMAGE, AssetKind.STRUCTURED_VISUAL, AssetKind.LOGO, AssetKind.ICON})


def _first_crop_asset(resources: ResolvedResources, style: CompilerStyle) -> tuple[AssetRecord | None, Path | None]:
    return _first_asset_by_kind(resources, style, {AssetKind.DOCUMENT_CROP, AssetKind.IMAGE})


def _first_structured_asset(resources: ResolvedResources, style: CompilerStyle) -> tuple[AssetRecord | None, Path | None]:
    return _first_asset_by_kind(resources, style, {AssetKind.STRUCTURED_VISUAL})


def _first_visual_image(resources: ResolvedResources, root: Path) -> tuple[ResolvedVisual | None, Path | None]:
    for visual in resources.visuals:
        resolved = _resolve_path(visual.output_path, root)
        if resolved is not None and resolved.is_file():
            return (visual, resolved)
    return (None, None)


def _first_visual_table(resources: ResolvedResources, root: Path) -> tuple[ResolvedVisual | None, Path | None]:
    for visual in resources.visuals:
        resolved = _resolve_path(visual.data_output_path, root)
        if resolved is not None and resolved.is_file():
            return (visual, resolved)
    return (None, None)


def _preferred_image_resource(
    slide_data: SlideIRSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> tuple[AssetRecord | ResolvedVisual | None, Path | None]:
    crop_asset, crop_path = _first_crop_asset(resources, style)
    structured_visual, visual_path = _first_visual_image(resources, style.root)
    structured_asset, structured_asset_path = _first_structured_asset(resources, style)
    bridge = slide_data.production_bridge
    if (
        bridge is not None
        and (
            bridge.visual_source_preference == VisualSourcePreference.DOCUMENT_CROP
            or bridge.production_mode == ProductionMode.SOURCE_REUSE
        )
    ):
        if crop_path is not None:
            return (crop_asset, crop_path)
        return (None, None)
    if visual_path is not None:
        return (structured_visual, visual_path)
    if structured_asset_path is not None:
        return (structured_asset, structured_asset_path)
    if crop_path is not None and bridge is not None and bridge.visual_source_preference == VisualSourcePreference.EITHER:
        return (crop_asset, crop_path)
    return (None, None)


def _compose_cover(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(style.margin),
        Inches(1.05),
        Inches(style.slide_width - (style.margin * 2)),
        Inches(0.28),
    )
    _fill_shape(band, style.colors["signal"], 0.0)
    _add_textbox(
        slide,
        style.margin,
        1.8,
        style.slide_width - (style.margin * 2),
        1.15,
        slide_data.title,
        font_name=style.fonts["title"].font_family,
        font_size=style.fonts["title"].size_pt + 8,
        color=style.colors["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        style.margin + 1.0,
        3.1,
        style.slide_width - ((style.margin + 1.0) * 2),
        0.75,
        slide_data.one_line_takeaway,
        font_name=style.fonts["body"].font_family,
        font_size=style.fonts["body"].size_pt + 1,
        color=style.colors["ink"],
        align=PP_ALIGN.CENTER,
    )
    _add_bullets(slide, _visible_support_items(slide_data), style, style.margin + 2.05, 4.1, style.slide_width - 4.1, 1.4)
    return []


def _compose_section_divider(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(style.margin),
        Inches(2.55),
        Inches(style.slide_width - (style.margin * 2)),
        Inches(1.35),
    )
    _fill_shape(band, style.colors["signal"], 0.1)
    _add_textbox(
        slide,
        style.margin + 0.3,
        2.82,
        style.slide_width - ((style.margin + 0.3) * 2),
        0.55,
        slide_data.section.upper(),
        font_name=style.fonts["caption"].font_family,
        font_size=style.fonts["caption"].size_pt + 1,
        color=style.colors["signal"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        style.margin + 0.8,
        3.15,
        style.slide_width - ((style.margin + 0.8) * 2),
        0.6,
        slide_data.title,
        font_name=style.fonts["title"].font_family,
        font_size=style.fonts["title"].size_pt + 2,
        color=style.colors["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return []


def _compose_agenda_roadmap(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style)
    steps = _visible_support_items(slide_data)
    _add_message_panel(slide, "Roadmap", slide_data.main_message, style, style.margin, 1.6, 3.5, 1.25)
    _add_bullets(slide, steps, style, 4.45, 1.72, 3.9, 1.65)
    step_width = (style.slide_width - (style.margin * 2) - (style.gap * max(len(steps) - 1, 0))) / max(len(steps), 1)
    for index, step in enumerate(steps):
        left = style.margin + (index * (step_width + style.gap))
        _add_message_panel(slide, f"Phase {index + 1}", step, style, left, 3.35, step_width, 1.5)
    return []


def _compose_concept_explainer(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style)
    _add_message_panel(slide, "Core point", slide_data.main_message, style, style.margin, 1.55, 4.3, 2.05)
    _add_bullets(slide, _visible_support_items(slide_data), style, 5.05, 1.72, 3.95, 2.0)
    _add_support_panel(slide, slide_data, style, 5.05, 3.95, 3.95, 1.95, heading="Interpretation")
    return []


def _compose_summary_next_step(slide, slide_data: BlueprintSlide, style: CompilerStyle) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style)
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(style.margin),
        Inches(1.65),
        Inches(style.slide_width - (style.margin * 2)),
        Inches(1.5),
    )
    _fill_shape(card, style.colors["signal"], 0.0)
    _add_textbox(
        slide,
        style.margin + 0.25,
        1.92,
        style.slide_width - ((style.margin + 0.25) * 2),
        0.9,
        slide_data.main_message,
        font_name=style.fonts["title"].font_family,
        font_size=style.fonts["title"].size_pt + 2,
        color=RGBColor(255, 255, 255),
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_message_panel(slide, "Why now", slide_data.one_line_takeaway, style, style.margin, 3.5, 4.0, 1.4)
    _add_bullets(slide, _visible_support_items(slide_data), style, 4.8, 3.68, 4.2, 1.45)
    return []


def _compose_case_study(
    slide,
    slide_data: BlueprintSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style)
    chosen_resource, chosen_path = _preferred_image_resource(slide_data, resources, style)
    linked_paths: list[str] = []
    if chosen_path is not None:
        _add_picture(slide, chosen_path, style, style.margin, 1.65, 5.6, 3.8)
        linked_paths.append(_display_path(chosen_path, style.root))
    else:
        _add_support_panel(slide, slide_data, style, style.margin, 1.65, 5.6, 3.8, heading="Worked example")
    _add_message_panel(slide, "Case insight", slide_data.main_message, style, 6.4, 1.65, 2.95, 1.4)
    _add_bullets(slide, _visible_support_items(slide_data), style, 6.4, 3.2, 2.95, 1.55)
    if isinstance(chosen_resource, AssetRecord):
        _add_textbox(
            slide,
            6.4,
            5.0,
            2.95,
            0.55,
            _source_lane_text(chosen_resource, slide_data),
            font_name=style.fonts["caption"].font_family,
            font_size=style.fonts["caption"].size_pt,
            color=style.colors["ink"],
        )
    return linked_paths


def _compose_evidence_chart(
    slide,
    slide_data: BlueprintSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    _add_background(slide, style, appendix=slide_data.deck_mode == DeckMode.APPENDIX)
    _add_title_block(slide, slide_data, style)
    chosen_resource, chosen_path = _preferred_image_resource(slide_data, resources, style)
    linked_paths: list[str] = []
    if chosen_path is not None:
        _add_picture(slide, chosen_path, style, style.margin, 1.65, 8.6, 4.55)
        linked_paths.append(_display_path(chosen_path, style.root))
    else:
        _add_support_panel(slide, slide_data, style, style.margin, 1.65, 8.6, 4.55, heading="Worked visual")
    caption_asset = chosen_resource if isinstance(chosen_resource, AssetRecord) else None
    caption_text = _source_lane_text(caption_asset, slide_data)
    _add_message_panel(slide, "Why it matters", slide_data.main_message, style, style.margin, 6.35, 4.35, 0.7)
    _add_textbox(
        slide,
        5.1,
        6.42,
        4.05,
        0.55,
        caption_text,
        font_name=style.fonts["caption"].font_family,
        font_size=style.fonts["caption"].size_pt,
        color=style.colors["ink"],
    )
    return linked_paths


def _compose_comparison(
    slide,
    slide_data: BlueprintSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    _add_background(slide, style, appendix=slide_data.deck_mode == DeckMode.APPENDIX)
    _add_title_block(slide, slide_data, style)
    table_visual, table_path = _first_visual_table(resources, style.root)
    image_visual, image_path = _first_visual_image(resources, style.root)
    structured_asset, structured_asset_path = _first_structured_asset(resources, style)
    linked_paths: list[str] = []
    if table_path is not None and _render_native_table(slide, table_path, style, style.margin, 1.72, 6.9, 3.8):
        linked_paths.append(_display_path(table_path, style.root))
    elif image_path is not None:
        _add_picture(slide, image_path, style, style.margin, 1.72, 6.9, 3.8)
        linked_paths.append(_display_path(image_path, style.root))
    elif structured_asset_path is not None:
        _add_picture(slide, structured_asset_path, style, style.margin, 1.72, 6.9, 3.8)
        linked_paths.append(_display_path(structured_asset_path, style.root))
    else:
        _add_support_panel(slide, slide_data, style, style.margin, 1.72, 6.9, 3.8, heading="Comparison cues")
    _add_message_panel(slide, "Comparison takeaway", slide_data.main_message, style, 7.0, 1.72, 2.3, 1.3)
    _add_bullets(slide, _visible_support_items(slide_data), style, 7.0, 3.25, 2.3, 1.6)
    if table_visual is not None or image_visual is not None or structured_asset is not None:
        source = table_visual.record.notes if table_visual is not None else image_visual.record.notes if image_visual is not None else structured_asset.notes
        source = _optional_visible_text(source)
        if source:
            _add_textbox(
                slide,
                7.0,
                5.3,
                2.3,
                0.5,
                source,
                font_name=style.fonts["caption"].font_family,
                font_size=style.fonts["caption"].size_pt,
                color=style.colors["ink"],
            )
    return linked_paths


def _compose_process_flow(
    slide,
    slide_data: BlueprintSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style)
    _visual, visual_path = _first_visual_image(resources, style.root)
    _asset, asset_path = _first_structured_asset(resources, style)
    linked_paths: list[str] = []
    chosen_path = visual_path or asset_path
    if chosen_path is not None:
        _add_picture(slide, chosen_path, style, style.margin, 1.65, style.slide_width - (style.margin * 2), 3.55)
        linked_paths.append(_display_path(chosen_path, style.root))
    else:
        steps = slide_data.required_evidence_assets[:4] or ["Step 1", "Step 2", "Step 3"]
        step_width = (style.slide_width - (style.margin * 2) - (style.gap * (len(steps) - 1))) / len(steps)
        for index, step in enumerate(steps):
            left = style.margin + (index * (step_width + style.gap))
            _add_message_panel(slide, f"Step {index + 1}", step, style, left, 2.25, step_width, 1.4)
    _add_message_panel(slide, "Takeaway", slide_data.main_message, style, style.margin, 5.45, style.slide_width - (style.margin * 2), 0.8)
    return linked_paths


def _compose_timeline(
    slide,
    slide_data: BlueprintSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style)
    _visual, visual_path = _first_visual_image(resources, style.root)
    _asset, asset_path = _first_structured_asset(resources, style)
    linked_paths: list[str] = []
    chosen_path = visual_path or asset_path
    if chosen_path is not None:
        _add_picture(slide, chosen_path, style, style.margin, 1.8, style.slide_width - (style.margin * 2), 3.25)
        linked_paths.append(_display_path(chosen_path, style.root))
    else:
        axis = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(style.margin + 0.3),
            Inches(3.2),
            Inches(style.slide_width - ((style.margin + 0.3) * 2)),
            Inches(0.06),
        )
        _fill_shape(axis, style.colors["signal"], 0.0)
        milestones = slide_data.required_evidence_assets[:4] or ["Start", "Decision", "Launch"]
        spacing = (style.slide_width - (style.margin * 2) - 0.7) / max(len(milestones), 1)
        for index, label in enumerate(milestones):
            left = style.margin + (index * spacing)
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + 0.15), Inches(3.05), Inches(0.18), Inches(0.18))
            _fill_shape(dot, style.colors["signal"], 0.0)
            _add_textbox(
                slide,
                left,
                3.35,
                spacing - 0.2,
                0.5,
                label,
                font_name=style.fonts["caption"].font_family,
                font_size=style.fonts["caption"].size_pt,
                color=style.colors["ink"],
                align=PP_ALIGN.CENTER,
            )
    _add_message_panel(slide, "Timeline point", slide_data.main_message, style, style.margin, 4.95, style.slide_width - (style.margin * 2), 0.75)
    return linked_paths


def _compose_framework_model(
    slide,
    slide_data: BlueprintSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    _add_background(slide, style)
    _add_title_block(slide, slide_data, style)
    _visual, visual_path = _first_visual_image(resources, style.root)
    _asset, asset_path = _first_structured_asset(resources, style)
    linked_paths: list[str] = []
    chosen_path = visual_path or asset_path
    if chosen_path is not None:
        _add_picture(slide, chosen_path, style, style.margin, 1.7, style.slide_width - (style.margin * 2), 3.65)
        linked_paths.append(_display_path(chosen_path, style.root))
    else:
        nodes = _visible_support_items(slide_data, limit=4)
        node_width = (style.slide_width - (style.margin * 2) - style.gap) / 2
        node_height = 1.35
        for index, node in enumerate(nodes[:4]):
            left = style.margin + ((index % 2) * (node_width + style.gap))
            top = 2.05 + ((index // 2) * (node_height + style.gap))
            _add_message_panel(slide, f"Concept {index + 1}", node, style, left, top, node_width, node_height)
    _add_textbox(
        slide,
        style.margin,
        5.75,
        style.slide_width - (style.margin * 2),
        0.45,
        slide_data.main_message,
        font_name=style.fonts["body"].font_family,
        font_size=style.fonts["body"].size_pt,
        color=style.colors["ink"],
        align=PP_ALIGN.CENTER,
    )
    return linked_paths


def _compose_appendix_reference(
    slide,
    slide_data: BlueprintSlide,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    if slide_data.visual_type in {VisualType.DOCUMENT_CROP, VisualType.PHOTO, VisualType.CHART}:
        linked_paths = _compose_evidence_chart(slide, slide_data, resources, style)
    else:
        linked_paths = _compose_comparison(slide, slide_data, resources, style)
    _add_textbox(
        slide,
        style.margin,
        6.7,
        style.slide_width - (style.margin * 2),
        0.26,
        "Supporting source material for the approved main story.",
        font_name=style.fonts["caption"].font_family,
        font_size=style.fonts["caption"].size_pt,
        color=style.colors["signal"],
        align=PP_ALIGN.CENTER,
    )
    return linked_paths


def _compose_slide(
    slide,
    slide_data: BlueprintSlide,
    entry: SlideLedgerEntry,
    family: str,
    resources: ResolvedResources,
    style: CompilerStyle,
) -> list[str]:
    if family == "cover":
        return _compose_cover(slide, slide_data, style)
    if family == "section-divider":
        return _compose_section_divider(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.CORRESPONDENCE_MATRIX:
        return _compose_correspondence_matrix(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.TWO_COLUMN_MAPPING_TABLE:
        return _compose_two_column_mapping_table(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.WORKED_EXAMPLE_STATE_TABLE:
        return _compose_worked_example_state_table(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.COMPARISON_MATRIX:
        return _compose_comparison_matrix(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.APPENDIX_THEMED_EVIDENCE_CLUSTER:
        return _compose_appendix_themed_evidence_cluster(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.APPENDIX_SOURCE_LOCATION_MATRIX:
        return _compose_appendix_source_location_matrix(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.APPENDIX_ANNOTATED_EXCERPT_CLUSTER:
        return _compose_appendix_annotated_excerpt_cluster(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.APPENDIX_COMPARISON_EVIDENCE_CLUSTER:
        return _compose_appendix_comparison_evidence_cluster(slide, slide_data, style)
    if slide_data.slide_archetype == SlideArchetype.APPENDIX_SOURCE_MAP:
        return _compose_appendix_source_map(slide, slide_data, style)
    if slide_data.slide_archetype in {
        SlideArchetype.APPENDIX_EVIDENCE_CLUSTER,
    }:
        return _compose_appendix_evidence_cluster(slide, slide_data, style)
    if slide_data.slide_archetype in {
        SlideArchetype.PROCESS_FLOW,
        SlideArchetype.STEP_BY_STEP_MECHANISM,
    }:
        return _compose_authored_process(slide, slide_data, style)
    if slide_data.slide_archetype in {
        SlideArchetype.LIMITATION_PITFALL_CALLOUT,
    }:
        return _compose_limitation_pitfall_callout(slide, slide_data, style)
    if slide_data.slide_archetype in {
        SlideArchetype.ANCHOR_CONCEPT_CARD,
        SlideArchetype.APPLICATION_VIGNETTE,
        SlideArchetype.SYNTHESIS_INTEGRATION,
    }:
        return _compose_authored_cards(slide, slide_data, style)
    if family == "summary":
        return _compose_summary_next_step(slide, slide_data, style)
    if family == "comparison":
        return _compose_comparison(slide, slide_data, resources, style)
    if family == "process-flow":
        if slide_data.visual_type == VisualType.TIMELINE:
            return _compose_timeline(slide, slide_data, resources, style)
        if slide_data.slide_role == SlideRole.PROCESS and slide_data.visual_type in {VisualType.PROCESS, VisualType.DECISION_PATH}:
            return _compose_agenda_roadmap(slide, slide_data, style)
        return _compose_process_flow(slide, slide_data, resources, style)
    if family == "worked-example":
        if slide_data.visual_type == VisualType.PHOTO:
            return _compose_case_study(slide, slide_data, resources, style)
        return _compose_evidence_chart(slide, slide_data, resources, style)
    if family in {"concept-explainer", "definition-theorem"}:
        return _compose_framework_model(slide, slide_data, resources, style)
    if family == "appendix-reference":
        return _compose_appendix_reference(slide, slide_data, resources, style)
    _ = entry
    return _compose_concept_explainer(slide, slide_data, style)


def _make_placeholder_slide(slide, entry: SlideLedgerEntry, style: CompilerStyle) -> None:
    _ = slide, entry, style
    raise ValueError("Placeholder slide emission is disabled for export-safe compilation.")


def _updated_batch_manifest(batch_manifest: BatchManifest | None, slide_ledger: SlideLedger) -> BatchManifest | None:
    if batch_manifest is None:
        return None
    updated_batches = []
    entries = {entry.slide_number: entry for entry in slide_ledger.entries}
    for batch in batch_manifest.batches:
        start = batch.slide_range.start
        end = batch.slide_range.end
        slice_entries = [entries[number] for number in range(start, end + 1) if number in entries]
        if slice_entries and all(entry.compile_status == StageStatus.COMPLETE for entry in slice_entries):
            status = StageStatus.COMPLETE
        elif slice_entries and any(entry.compile_status == StageStatus.BLOCKED for entry in slice_entries):
            status = StageStatus.BLOCKED
        else:
            status = batch.status
        updated_batches.append(batch.model_copy(update={"status": status}))
    return batch_manifest.model_copy(update={"batches": updated_batches})


def _updated_state_capsule(state_capsule: StateCapsule | None, warnings: list[str]) -> StateCapsule | None:
    if state_capsule is None:
        return None
    pending_actions = [action for action in state_capsule.pending_actions if not action.startswith("Compile PPTX")]
    if warnings:
        pending_actions.append("Resolve compiler warnings before QA sign-off.")
    else:
        pending_actions.append("Run deck QA on the compiled PPTX.")
    return state_capsule.model_copy(
        update={
            "active_gate": WorkflowGate.PRODUCTION_AND_QA,
            "blueprint_approved": True,
            "pending_actions": _dedupe(pending_actions),
        }
    )


def _merge_change_note(existing: str | None, note: str) -> str:
    if not existing:
        return note
    if note in existing:
        return existing
    return f"{existing} {note}"


def _compiler_style_from_slide_ir(slide_ir: SlideIRDocument, *, root: Path, raster_dir: Path) -> CompilerStyle:
    context = _require_compile_context(slide_ir)
    design_system = context.design_system
    if design_system is None:
        raise ValueError("Slide IR compile context is missing design_system for PPTX compilation.")
    slide_width, slide_height = _parse_ratio(slide_ir.slide_ratio)
    return CompilerStyle(
        slide_width=slide_width,
        slide_height=slide_height,
        margin=max(DEFAULT_MARGIN_IN, _spacing(design_system, 4, DEFAULT_MARGIN_IN)),
        gap=max(DEFAULT_GAP_IN, _spacing(design_system, 3, DEFAULT_GAP_IN)),
        colors=_color_map(design_system, slide_ir.style_prior),
        fonts=_font_map(design_system),
        root=root,
        raster_dir=raster_dir,
        style_prior=slide_ir.style_prior,
    )


def _compile_warnings_from_slide_ir(
    slide_ir: SlideIRDocument,
    validation: SlideIRValidationReport,
    continuity: SlideIRContinuityReport | None = None,
) -> list[str]:
    context = _require_compile_context(slide_ir)
    continuity_report = continuity if continuity is not None else validate_slide_ir_continuity(slide_ir)
    continuity_guidance = continuity_report.as_guidance_lines()
    warnings = list(slide_ir.warnings)
    if not validation.is_valid:
        warnings.append(
            f"Slide IR geometry has {validation.error_count} error(s); review out-of-bounds or overlap findings before relying on preview."
        )
    warnings.extend(_slide_ir_warnings(validation))
    warnings.extend(continuity_guidance)
    if context.design_system is not None and context.design_system.slide_ratio != slide_ir.slide_ratio:
        warnings.append(f"Design system ratio {context.design_system.slide_ratio} differs from blueprint ratio {slide_ir.slide_ratio}.")
    if context.layout_library_slide_ratio is not None and context.layout_library_slide_ratio != slide_ir.slide_ratio:
        warnings.append(f"Layout library ratio {context.layout_library_slide_ratio} differs from blueprint ratio {slide_ir.slide_ratio}.")
    if context.appendix_start is not None:
        for slide in slide_ir.slides:
            if slide.slide_number >= context.appendix_start and slide.deck_mode != DeckMode.APPENDIX:
                warnings.append(
                    f"Slide {slide.slide_number} falls after appendix_start but is not marked as appendix. Check constitution rule: {context.appendix_boundary_rule}"
                )
            if slide.slide_number < context.appendix_start and slide.deck_mode == DeckMode.APPENDIX:
                warnings.append(
                    f"Slide {slide.slide_number} is marked as appendix before appendix_start. Check constitution rule: {context.appendix_boundary_rule}"
                )
    for slide in slide_ir.slides:
        warnings.extend(slide.layout_warnings)
    return warnings


def _layout_critic_report_from_slide_ir(slide_ir: SlideIRDocument) -> SlideIRLayoutCriticReport | None:
    payload = slide_ir.generation_inputs.get(LAYOUT_CRITIC_REPORT_KEY)
    if not isinstance(payload, dict):
        return None
    try:
        return SlideIRLayoutCriticReport.model_validate(payload)
    except Exception:
        return None


def compile_pptx_from_slide_ir(
    slide_ir: SlideIRDocument,
    output_dir: str | Path,
    *,
    batch_manifest: BatchManifest | None = None,
    state_capsule: StateCapsule | None = None,
    notes_path: str | Path | None = None,
    pptx_name: str = "deck.pptx",
    root: str | Path | None = None,
) -> PptxCompileOutputs:
    context = _require_compile_context(slide_ir)
    if context.slide_ledger is None:
        raise ValueError("Slide IR compile context is missing slide_ledger for compatibility outputs.")
    output_root = Path(output_dir)
    output_root.mkdir(parents=True, exist_ok=True)
    root_path = Path(root).resolve() if root is not None else Path.cwd().resolve()
    notes_map = _load_notes(notes_path)
    raster_dir = output_root / "rasterized-visuals"
    slide_ir_report = validate_slide_ir_geometry(slide_ir)
    continuity_report = validate_slide_ir_continuity(slide_ir)
    warnings = _compile_warnings_from_slide_ir(slide_ir, slide_ir_report, continuity_report)
    contract_errors = _compile_contract_errors(slide_ir)
    if contract_errors:
        raise CompileContractError("PPTX compile contract failed:\n- " + "\n- ".join(contract_errors))

    style = _compiler_style_from_slide_ir(slide_ir, root=root_path, raster_dir=raster_dir)
    presentation = Presentation()
    presentation.slide_width = Inches(style.slide_width)
    presentation.slide_height = Inches(style.slide_height)
    blank_layout = presentation.slide_layouts[6]

    slides_by_number = {slide.slide_number: slide for slide in slide_ir.slides}
    ledger_by_number = {entry.slide_number: entry for entry in context.ledger_entries}
    for entry in sorted(context.ledger_entries, key=lambda item: item.slide_number):
        if entry.slide_number not in slides_by_number:
            raise ValueError(f"Ledger entry {entry.slide_id} has no matching blueprint slide.")
    updated_entries: list[SlideLedgerEntry] = []
    slide_links: list[SlideBuildLink] = []

    for slide_data in sorted(slide_ir.slides, key=lambda item: item.slide_number):
        entry = ledger_by_number.get(slide_data.slide_number)
        if entry is None:
            raise ValueError(f"Slide IR slide {slide_data.slide_id} has no matching ledger entry.")

        slide = presentation.slides.add_slide(blank_layout)
        resources = _resolve_resources(slide_data, context, root_path, raster_dir)
        family = slide_data.layout_family
        linked_paths = _compose_slide(slide, slide_data, entry, family, resources, style)
        number_label = slide_data.numbering_label or _build_number_label(slide_data, context.appendix_start)
        _add_footer(slide, number_label, slide_data.deck_mode, style)
        _assert_no_forbidden_visible_text(slide, slide_data.slide_number)
        notes_present = _set_notes(slide, _note_text(slide_data, notes_map))

        if entry.deck_mode != slide_data.deck_mode:
            warnings.append(
                f"Slide {slide_data.slide_number} deck_mode differs between ledger ({entry.deck_mode.value}) and SlideIR ({slide_data.deck_mode.value})."
            )

        missing_dependencies = list(resources.missing_dependencies)
        warnings.extend(f"Slide {slide_data.slide_number}: {message}" for message in missing_dependencies)
        compile_status = StageStatus.COMPLETE if not missing_dependencies else StageStatus.BLOCKED
        compile_note = (
            f"Compiled with layout family {family}."
            if not missing_dependencies
            else f"Compile blocked for layout family {family}: {'; '.join(missing_dependencies)}"
        )
        updated_entries.append(
            entry.model_copy(
                update={
                    "compile_status": compile_status,
                    "change_note": _merge_change_note(entry.change_note, compile_note),
                    "unresolved_blockers": _dedupe([*(entry.unresolved_blockers or []), *missing_dependencies]) or None,
                }
            )
        )

        slide_links.append(
            SlideBuildLink(
                slide_number=slide_data.slide_number,
                slide_id=slide_data.slide_id,
                title=slide_data.title,
                pptx_index=len(presentation.slides),
                deck_mode=slide_data.deck_mode,
                batch_id=slide_data.batch_id,
                layout_pattern_id=slide_data.layout_pattern_id,
                layout_family=family,
                visual_type=slide_data.visual_type,
                numbering_label=number_label,
                asset_ids=[asset.asset_id for asset in resources.assets],
                viz_spec_ids=[visual.record.spec.spec_id for visual in resources.visuals],
                linked_paths=_dedupe(linked_paths),
                notes_present=notes_present,
                compile_status=compile_status,
                missing_dependencies=missing_dependencies,
            )
        )

    pptx_path = output_root / pptx_name
    presentation.save(pptx_path)
    linkage_path = output_root / "slide-build-linkage.json"

    updated_slide_ledger = context.slide_ledger.model_copy(update={"entries": updated_entries})
    updated_batch_manifest = _updated_batch_manifest(batch_manifest, updated_slide_ledger)
    deduped_warnings = _dedupe(warnings)
    updated_state_capsule = _updated_state_capsule(state_capsule, deduped_warnings)
    compile_report = build_slide_ir_compile_report(slide_ir, slide_ir_report, deduped_warnings, continuity_report)

    linkage = SlideBuildLinkage(
        deck_title=slide_ir.deck_title,
        pptx_path=_display_path(pptx_path, root_path),
        slides=slide_links,
    )
    manifest = BuildManifest(
        deck_title=slide_ir.deck_title,
        pptx_path=_display_path(pptx_path, root_path),
        slide_ratio=slide_ir.slide_ratio,
        slide_count=len(slide_links),
        compiled_layout_patterns=_dedupe([link.layout_pattern_id for link in slide_links]),
        warnings=deduped_warnings,
        linkage_path=_display_path(linkage_path, root_path),
        batch_ids=_dedupe(
            [entry.batch_id for entry in updated_slide_ledger.entries if entry.batch_id]
            + ([batch.batch_id for batch in updated_batch_manifest.batches] if updated_batch_manifest is not None else [])
        ),
    )
    return PptxCompileOutputs(
        slide_ir=slide_ir,
        compile_report=compile_report,
        build_manifest=manifest,
        slide_build_linkage=linkage,
        slide_ledger=updated_slide_ledger,
        layout_critic_report=_layout_critic_report_from_slide_ir(slide_ir),
        batch_manifest=updated_batch_manifest,
        state_capsule=updated_state_capsule,
        pptx_path=pptx_path,
    )


def compile_pptx(
    blueprint: Blueprint,
    design_system: DesignSystem,
    deck_constitution: DeckConstitution,
    layout_library: LayoutLibrary,
    slide_ledger: SlideLedger,
    asset_manifest: AssetManifest,
    viz_manifest: VizManifest,
    output_dir: str | Path,
    *,
    batch_manifest: BatchManifest | None = None,
    state_capsule: StateCapsule | None = None,
    notes_path: str | Path | None = None,
    pptx_name: str = "deck.pptx",
    root: str | Path | None = None,
    enable_layout_critic: bool = True,
    style_prior_provider: StylePriorProvider | None = None,
) -> PptxCompileOutputs:
    _ensure_blueprint_approved(blueprint, state_capsule)
    slide_ir = adapt_blueprint_to_slide_ir(
        blueprint=blueprint,
        design_system=design_system,
        deck_constitution=deck_constitution,
        layout_library=layout_library,
        slide_ledger=slide_ledger,
        asset_manifest=asset_manifest,
        viz_manifest=viz_manifest,
        enable_layout_critic=enable_layout_critic,
        style_prior_provider=style_prior_provider,
    )
    return compile_pptx_from_slide_ir(
        slide_ir=slide_ir,
        output_dir=output_dir,
        batch_manifest=batch_manifest,
        state_capsule=state_capsule,
        notes_path=notes_path,
        pptx_name=pptx_name,
        root=root,
    )


def load_pptx_compile_file(
    path: str | Path,
) -> BuildManifest | SlideBuildLinkage | SlideIRCompileReport | SlideIRLayoutCriticReport:
    payload = json.loads(Path(path).read_text(encoding="utf-8"))
    # Preview payload snapshots are intentionally outside the persisted
    # pptx-compiler artifact contract; only schema-tagged compiler artifacts
    # are supported loader inputs here. For compile-report artifacts, the
    # supported serialized contract stops at canonical `continuity_guidance`
    # plus in-repo load normalization. Legacy serialized
    # `continuity_warnings` remains accepted only as a normalization input and
    # is not restored onto the loaded model surface.
    schema_name = payload.get("schema_name")
    if schema_name == BuildManifest.SCHEMA_NAME:
        return BuildManifest.model_validate(payload)
    if schema_name == SlideBuildLinkage.SCHEMA_NAME:
        return SlideBuildLinkage.model_validate(payload)
    if schema_name == SlideIRCompileReport.SCHEMA_NAME:
        return SlideIRCompileReport.model_validate(payload)
    if schema_name == SlideIRLayoutCriticReport.SCHEMA_NAME:
        return SlideIRLayoutCriticReport.model_validate(payload)
    raise ValueError(f"unsupported pptx-compiler schema {schema_name!r}")


def compile_pptx_from_files(
    blueprint_path: str | Path,
    design_system_path: str | Path,
    deck_constitution_path: str | Path,
    layout_library_path: str | Path,
    slide_ledger_path: str | Path,
    asset_manifest_path: str | Path,
    viz_manifest_path: str | Path,
    output_dir: str | Path,
    *,
    batch_manifest_path: str | Path | None = None,
    state_capsule_path: str | Path | None = None,
    notes_path: str | Path | None = None,
    pptx_name: str = "deck.pptx",
    root: str | Path | None = None,
    enable_layout_critic: bool = True,
    style_prior_provider: StylePriorProvider | None = None,
) -> PptxCompileOutputs:
    blueprint = load_state_file(blueprint_path)
    if blueprint.schema_name != "blueprint":
        raise TypeError(f"expected blueprint, found {blueprint.schema_name}")
    design_system = load_state_file(design_system_path)
    if design_system.schema_name != "design_system":
        raise TypeError(f"expected design_system, found {design_system.schema_name}")
    deck_constitution = load_state_file(deck_constitution_path)
    if deck_constitution.schema_name != "deck_constitution":
        raise TypeError(f"expected deck_constitution, found {deck_constitution.schema_name}")
    layout_library = load_state_file(layout_library_path)
    if layout_library.schema_name != "layout_library":
        raise TypeError(f"expected layout_library, found {layout_library.schema_name}")
    slide_ledger = load_state_file(slide_ledger_path)
    if slide_ledger.schema_name != "slide_ledger":
        raise TypeError(f"expected slide_ledger, found {slide_ledger.schema_name}")
    asset_manifest = load_state_file(asset_manifest_path)
    if asset_manifest.schema_name != "asset_manifest":
        raise TypeError(f"expected asset_manifest, found {asset_manifest.schema_name}")
    viz_manifest = load_state_file(viz_manifest_path)
    if viz_manifest.schema_name != "viz_manifest":
        raise TypeError(f"expected viz_manifest, found {viz_manifest.schema_name}")

    batch_manifest = None
    if batch_manifest_path is not None:
        batch_manifest = load_state_file(batch_manifest_path)
        if batch_manifest.schema_name != "batch_manifest":
            raise TypeError(f"expected batch_manifest, found {batch_manifest.schema_name}")
    state_capsule = None
    if state_capsule_path is not None:
        state_capsule = load_state_file(state_capsule_path)
        if state_capsule.schema_name != "state_capsule":
            raise TypeError(f"expected state_capsule, found {state_capsule.schema_name}")

    return compile_pptx(
        blueprint=blueprint,
        design_system=design_system,
        deck_constitution=deck_constitution,
        layout_library=layout_library,
        slide_ledger=slide_ledger,
        asset_manifest=asset_manifest,
        viz_manifest=viz_manifest,
        output_dir=output_dir,
        batch_manifest=batch_manifest,
        state_capsule=state_capsule,
        notes_path=notes_path,
        pptx_name=pptx_name,
        root=root,
        enable_layout_critic=enable_layout_critic,
        style_prior_provider=style_prior_provider,
    )


def write_pptx_compile_outputs(outputs: PptxCompileOutputs, output_dir: str | Path) -> dict[str, Path]:
    root = Path(output_dir)
    root.mkdir(parents=True, exist_ok=True)
    # Compiler preview payloads are emit-time handoff data from
    # render_slide_preview(...). Even after retiring the live preview warning
    # mirror, they remain out of persisted compile outputs. Gate 2
    # `authoring-preview.json` remains a separate artifact family and is not a
    # compiler preview persistence path.
    written = {
        "build_manifest": save_state_file(outputs.build_manifest, root / "build-manifest.json"),
        "slide_build_linkage": save_state_file(outputs.slide_build_linkage, root / "slide-build-linkage.json"),
        "slide_ledger": save_state_file(outputs.slide_ledger, root / "slide-ledger.json"),
    }
    if outputs.layout_critic_report is not None:
        written["layout_critic_report"] = save_state_file(outputs.layout_critic_report, root / "layout-critic-report.json")
    if outputs.batch_manifest is not None:
        written["batch_manifest"] = save_state_file(outputs.batch_manifest, root / "batch-manifest.json")
    if outputs.state_capsule is not None:
        written["state_capsule"] = save_state_file(outputs.state_capsule, root / "state-capsule.json")
    return written
