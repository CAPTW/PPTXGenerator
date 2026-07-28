"""Audit protected text/content zones against decorative and chrome intrusion."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_CONTRACTS_DIR = Path("outputs/template_contracts")
DEFAULT_GOLDEN_PPTX = Path("outputs/golden_template_masters.pptx")
DEFAULT_GOLDEN_REPORT = Path("outputs/golden_template_masters_report.json")
DEFAULT_FINAL_PPTX = Path("outputs/final_deck_large_premium.pptx")
DEFAULT_FINAL_MANIFEST = Path("outputs/final_deck_manifest.json")
DEFAULT_JSON_REPORT = Path("outputs/protected_zone_audit_report.json")
DEFAULT_MD_REPORT = Path("outputs/protected_zone_audit_report.md")

EMU_PER_INCH = 914400
FOOTER_Y = 6.45

ARCHETYPE_ALIASES = {
    "closing": "creative_cover",
    "concept_relationship": "concept_relationship_venn",
    "cover": "creative_cover",
    "data_dashboard": "kpi_donut_chart",
    "sequence": "work_support_sequence",
    "visual_toc": "visual_table_of_contents",
}

PRIMARY_ZONE_TYPES = {"title", "subtitle", "body", "card", "table", "chart"}


def build_protected_zone_audit_from_files(
    *,
    contracts_dir: str | Path = DEFAULT_CONTRACTS_DIR,
    golden_pptx_path: str | Path = DEFAULT_GOLDEN_PPTX,
    golden_report_path: str | Path = DEFAULT_GOLDEN_REPORT,
    final_pptx_path: str | Path = DEFAULT_FINAL_PPTX,
    final_manifest_path: str | Path = DEFAULT_FINAL_MANIFEST,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
) -> dict[str, Any]:
    report = build_protected_zone_audit(
        contracts_dir=contracts_dir,
        golden_pptx_path=golden_pptx_path,
        golden_report_path=golden_report_path,
        final_pptx_path=final_pptx_path,
        final_manifest_path=final_manifest_path,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return report


def build_protected_zone_audit(
    *,
    contracts_dir: str | Path,
    golden_pptx_path: str | Path,
    golden_report_path: str | Path,
    final_pptx_path: str | Path,
    final_manifest_path: str | Path,
) -> dict[str, Any]:
    contracts = _load_contracts(Path(contracts_dir))
    decks: list[dict[str, Any]] = []
    findings: list[dict[str, Any]] = []

    golden_records = _golden_records(Path(golden_report_path))
    golden_report = _inspect_deck("golden_template_masters", Path(golden_pptx_path), golden_records, contracts)
    decks.append(golden_report)
    findings.extend(golden_report["findings"])

    final_pptx = Path(final_pptx_path)
    if final_pptx.exists():
        final_records = _final_records(Path(final_manifest_path))
        final_report = _inspect_deck("final_deck_large_premium", final_pptx, final_records, contracts)
        decks.append(final_report)
        findings.extend(final_report["findings"])
    else:
        decks.append(
            {
                "scope": "final_deck_large_premium",
                "pptx_path": _display_path(final_pptx),
                "status": "skipped",
                "slide_count": 0,
                "slides": [],
                "findings": [],
                "findings_summary": {"total": 0, "severe": 0, "warning": 0, "minor": 0},
                "text_safe_archetypes": [],
            }
        )

    severe = sum(1 for finding in findings if finding["severity"] == "severe")
    warning = sum(1 for finding in findings if finding["severity"] == "warning")
    minor = sum(1 for finding in findings if finding["severity"] == "minor")
    return {
        "schema_name": "protected_zone_audit_report",
        "schema_version": "1.0",
        "status": "failed" if severe else "issues_reported" if warning or minor else "passed",
        "contracts_dir": _display_path(Path(contracts_dir)),
        "golden_pptx_path": _display_path(Path(golden_pptx_path)),
        "golden_report_path": _display_path(Path(golden_report_path)),
        "final_pptx_path": _display_path(Path(final_pptx_path)) if Path(final_pptx_path).exists() else None,
        "final_manifest_path": _display_path(Path(final_manifest_path)) if Path(final_manifest_path).exists() else None,
        "findings_summary": {"total": len(findings), "severe": severe, "warning": warning, "minor": minor},
        "text_safe_archetypes": _text_safe_archetypes(decks),
        "decks": decks,
        "findings": findings,
        "audit_policy": {
            "decorative_ornament_over_primary_text_zone": "severe",
            "footer_source_strip_over_content_zone": "severe",
            "card_chrome_leaves_insufficient_text_area": "severe",
            "table_chart_chrome_covers_labels_data": "severe",
            "index_rail_overlaps_primary_reading_path": "severe",
        },
    }


def _inspect_deck(scope: str, pptx_path: Path, records: dict[int, dict[str, Any]], contracts: dict[str, dict[str, Any]]) -> dict[str, Any]:
    if not pptx_path.exists():
        return {
            "scope": scope,
            "pptx_path": _display_path(pptx_path),
            "status": "skipped",
            "slide_count": 0,
            "slides": [],
            "findings": [],
            "findings_summary": {"total": 0, "severe": 0, "warning": 0, "minor": 0},
            "text_safe_archetypes": [],
        }

    presentation = Presentation(str(pptx_path))
    slides: list[dict[str, Any]] = []
    findings: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        record = records.get(slide_number, {})
        archetype_id = _canonical_archetype_id(str(record.get("archetype_id") or "unknown"))
        contract = contracts.get(archetype_id, {})
        shape_records = _shape_records(slide, presentation)
        zones = _zones_for_record(record)
        zones.extend(_image_zones(shape_records, zones))
        zone_reports = [
            _audit_zone(scope, slide_number, archetype_id, zone, shape_records)
            for zone in zones
        ]
        slide_findings = [
            _finding_from_zone(scope, slide_number, archetype_id, zone_report)
            for zone_report in zone_reports
            if zone_report["risk_level"] in {"warning", "severe"}
        ]
        slide_findings = [finding for finding in slide_findings if finding is not None]
        findings.extend(slide_findings)
        slides.append(
            {
                "slide_number": slide_number,
                "archetype_id": archetype_id,
                "layout_id": record.get("layout_id"),
                "contract_used": record.get("contract_used") or contract.get("_contract_path"),
                "zone_count": len(zone_reports),
                "zones": zone_reports,
                "text_safe": all(zone["risk_level"] != "severe" for zone in zone_reports),
                "finding_count": len(slide_findings),
                "finding_codes": [finding["code"] for finding in slide_findings],
            }
        )

    severe = sum(1 for finding in findings if finding["severity"] == "severe")
    warning = sum(1 for finding in findings if finding["severity"] == "warning")
    minor = sum(1 for slide in slides for zone in slide["zones"] if zone["risk_level"] == "minor")
    return {
        "scope": scope,
        "pptx_path": _display_path(pptx_path),
        "status": "failed" if severe else "issues_reported" if warning or minor else "passed",
        "slide_count": len(slides),
        "slides": slides,
        "findings": findings,
        "findings_summary": {"total": len(findings) + minor, "severe": severe, "warning": warning, "minor": minor},
        "text_safe_archetypes": sorted({slide["archetype_id"] for slide in slides if slide["text_safe"]}),
    }


def _audit_zone(scope: str, slide_number: int, archetype_id: str, zone: dict[str, Any], shapes: list[dict[str, Any]]) -> dict[str, Any]:
    bbox = zone["zone_bbox"]
    zone_area = max(0.001, bbox["w"] * bbox["h"])
    intersecting: list[dict[str, Any]] = []
    for shape in shapes:
        overlap = _intersection_area(bbox, shape["bbox"])
        if overlap <= 0.005:
            continue
        intersecting.append({**shape, "overlap_area": overlap, "overlap_ratio_of_zone": overlap / zone_area})

    text_objects = [item for item in intersecting if item["category"] == "text"]
    decorative_objects = [item for item in intersecting if item["category"] == "decorative"]
    chrome_objects = [item for item in intersecting if item["category"] == "chrome"]
    table_chart_objects = [item for item in intersecting if item["category"] in {"table", "chart"}]
    picture_objects = [item for item in intersecting if item["category"] == "picture"]

    decorative_area = sum(item["overlap_area"] for item in decorative_objects)
    chrome_area = sum(item["overlap_area"] for item in chrome_objects)
    text_area = sum(item["overlap_area"] for item in text_objects)
    overlap_ratio = min(1.0, (decorative_area + chrome_area) / zone_area)
    chrome_to_text_ratio = round(chrome_area / max(0.01, text_area), 3) if text_objects else None
    decorative_overlap_ratio = decorative_area / zone_area
    risk_level, code, action = _zone_risk(
        zone,
        decorative_objects,
        chrome_objects,
        text_objects,
        table_chart_objects,
        overlap_ratio,
        chrome_to_text_ratio,
        decorative_overlap_ratio,
    )
    return {
        "scope": scope,
        "slide_number": slide_number,
        "archetype_id": archetype_id,
        "zone_id": zone["zone_id"],
        "zone_type": zone["zone_type"],
        "slot_id": zone.get("slot_id"),
        "component_id": zone.get("component_id"),
        "zone_bbox": bbox,
        "text_object_count": len(text_objects),
        "decorative_object_count_inside": len(decorative_objects),
        "chrome_object_count_inside": len(chrome_objects),
        "table_chart_object_count_inside": len(table_chart_objects),
        "picture_object_count_inside": len(picture_objects),
        "overlap_ratio": round(overlap_ratio, 4),
        "chrome_to_text_area_ratio": chrome_to_text_ratio,
        "risk_level": risk_level,
        "finding_code": code,
        "recommended_action": action,
        "examples": [
            {
                "shape_index": item["shape_index"],
                "category": item["category"],
                "bbox": item["bbox"],
                "overlap_ratio_of_zone": round(item["overlap_ratio_of_zone"], 4),
            }
            for item in (decorative_objects + chrome_objects)[:5]
        ],
    }


def _zone_risk(
    zone: dict[str, Any],
    decorative_objects: list[dict[str, Any]],
    chrome_objects: list[dict[str, Any]],
    text_objects: list[dict[str, Any]],
    table_chart_objects: list[dict[str, Any]],
    overlap_ratio: float,
    chrome_to_text_ratio: float | None,
    decorative_overlap_ratio: float = 0.0,
) -> tuple[str, str | None, str]:
    zone_type = str(zone.get("zone_type") or "body")
    has_decorative = bool(decorative_objects)
    has_chrome = bool(chrome_objects)

    if zone_type in {"title", "body"} and has_decorative:
        return "severe", "DECORATIVE_ORNAMENT_OVER_PRIMARY_TEXT_ZONE", "Move ornaments outside the protected title/body reading area."
    if zone_type in PRIMARY_ZONE_TYPES and _footer_object_overlaps_zone(decorative_objects + chrome_objects):
        return "severe", "FOOTER_SOURCE_STRIP_OVER_CONTENT_ZONE", "Keep footer/source chrome below the protected content zone."
    if zone_type == "card" and has_chrome and not text_objects:
        return "warning", "CARD_TEXT_ZONE_EMPTY_OR_METADATA_GAP", "Verify this is a declared visual module frame; add text-zone metadata if it is intended to hold card text."
    if zone_type in {"table", "chart"} and decorative_objects and table_chart_objects and decorative_overlap_ratio > 0.08:
        return "severe", "TABLE_CHART_CHROME_COVERS_LABELS_DATA", "Move chart/table ornaments away from labels and data marks."
    if zone_type in {"table", "chart"} and has_chrome and table_chart_objects and overlap_ratio > 0.6:
        return "warning", "TABLE_CHART_CHROME_PRESSURE", "Review table/chart module chrome against label and data readability."
    if zone_type in PRIMARY_ZONE_TYPES and _index_rail_overlaps_zone(chrome_objects + decorative_objects):
        return "severe", "INDEX_RAIL_OVERLAPS_PRIMARY_READING_PATH", "Relocate index rail outside the primary reading path."

    if zone_type in {"footer", "image"}:
        if overlap_ratio > 0.5 and zone_type == "image" and has_chrome:
            return "warning", "PHOTO_FRAME_CHROME_HEAVY", "Keep photo-frame chrome outside the image focal area where possible."
        return "clear" if not has_decorative and not has_chrome else "minor", None, "No action; chrome is inside its declared zone."

    if has_decorative and zone_type in PRIMARY_ZONE_TYPES:
        return "warning", "DECORATIVE_OBJECT_IN_PROTECTED_ZONE", "Relocate nonessential decorative objects to the margin."
    if chrome_to_text_ratio is not None and chrome_to_text_ratio > 3.0:
        return "warning", "CHROME_PRESSURE_OVER_TEXT_ZONE", "Simplify chrome or expand text area within the protected zone."
    if overlap_ratio > 0.35 and has_chrome:
        return "warning", "CHROME_INSIDE_PROTECTED_ZONE", "Review whether chrome can be reduced without losing layout identity."
    if overlap_ratio > 0.04 and (has_decorative or has_chrome):
        return "minor", None, "Monitor; low-level chrome intersects the protected zone."
    return "clear", None, "No action."


def _finding_from_zone(scope: str, slide_number: int, archetype_id: str, zone_report: dict[str, Any]) -> dict[str, Any] | None:
    code = zone_report.get("finding_code")
    if not code:
        return None
    return {
        "code": code,
        "severity": zone_report["risk_level"],
        "scope": scope,
        "slide_number": slide_number,
        "archetype_id": archetype_id,
        "zone_id": zone_report["zone_id"],
        "zone_type": zone_report["zone_type"],
        "message": zone_report["recommended_action"],
        "details": {
            "zone_bbox": zone_report["zone_bbox"],
            "overlap_ratio": zone_report["overlap_ratio"],
            "text_object_count": zone_report["text_object_count"],
            "decorative_object_count_inside": zone_report["decorative_object_count_inside"],
            "chrome_object_count_inside": zone_report["chrome_object_count_inside"],
            "examples": zone_report["examples"],
        },
    }


def _zones_for_record(record: dict[str, Any]) -> list[dict[str, Any]]:
    plan = record.get("decorative_budget_plan") if isinstance(record.get("decorative_budget_plan"), dict) else {}
    raw_zones = list(plan.get("protected_text_zones") or []) + list(plan.get("protected_table_chart_zones") or [])
    zones: list[dict[str, Any]] = []
    for index, item in enumerate(raw_zones, start=1):
        if not isinstance(item, dict) or not isinstance(item.get("bounds"), dict):
            continue
        bbox = _normalize_bbox(item["bounds"])
        if bbox is None:
            continue
        zone_type = "footer" if bbox["y"] >= FOOTER_Y - 0.2 else _classify_zone_type(item)
        zones.append(
            {
                "zone_id": f"{zone_type}_{index}",
                "zone_type": zone_type,
                "slot_id": item.get("slot_id"),
                "slot_type": item.get("slot_type"),
                "component_id": item.get("component_id"),
                "zone_bbox": bbox,
            }
        )
    return zones


def _image_zones(shapes: list[dict[str, Any]], existing_zones: list[dict[str, Any]]) -> list[dict[str, Any]]:
    zones: list[dict[str, Any]] = []
    for index, shape in enumerate([item for item in shapes if item["category"] == "picture"], start=1):
        if any(_intersection_area(shape["bbox"], zone["zone_bbox"]) / max(0.001, shape["bbox"]["w"] * shape["bbox"]["h"]) > 0.7 for zone in existing_zones):
            continue
        zones.append(
            {
                "zone_id": f"image_{index}",
                "zone_type": "image",
                "slot_id": "photo_frame",
                "slot_type": "image",
                "component_id": "photo_frame_image",
                "zone_bbox": shape["bbox"],
            }
        )
    return zones


def _shape_records(slide: Any, presentation: Presentation) -> list[dict[str, Any]]:
    slide_area = max(0.001, (presentation.slide_width / EMU_PER_INCH) * (presentation.slide_height / EMU_PER_INCH))
    records: list[dict[str, Any]] = []
    for index, shape in enumerate(slide.shapes, start=1):
        bbox = _shape_bounds(shape)
        area_ratio = (bbox["w"] * bbox["h"]) / slide_area
        has_text = bool(getattr(shape, "has_text_frame", False) and str(shape.text or "").strip())
        if not has_text and getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE and area_ratio >= 0.82:
            continue
        records.append(
            {
                "shape_index": index,
                "category": _shape_category(shape, bbox),
                "bbox": bbox,
                "shape_type": str(getattr(shape, "shape_type", "")),
                "text_excerpt": str(getattr(shape, "text", "") or "")[:80] if has_text else "",
            }
        )
    return records


def _shape_category(shape: Any, bbox: dict[str, float]) -> str:
    if getattr(shape, "has_text_frame", False) and str(shape.text or "").strip():
        return "text"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE or _looks_like_small_ornament(bbox):
        return "decorative"
    return "chrome"


def _classify_zone_type(zone: dict[str, Any]) -> str:
    slot_id = str(zone.get("slot_id") or "").lower()
    slot_type = str(zone.get("slot_type") or "").lower()
    component_id = str(zone.get("component_id") or "").lower()
    joined = f"{slot_id} {slot_type} {component_id}"
    if slot_id == "title":
        return "title"
    if slot_id == "subtitle":
        return "subtitle"
    if "footer" in joined or "citation" in joined or "source_strip" in joined:
        return "footer"
    if "table" in joined:
        return "table"
    if "chart" in joined or "kpi" in joined:
        return "chart"
    if "image" in joined or "photo" in joined:
        return "image"
    if "card" in joined or slot_id in {"cards", "annotations", "index_navigation", "progress_markers", "metric_panels"}:
        return "card"
    return "body"


def _footer_object_overlaps_zone(objects: list[dict[str, Any]]) -> bool:
    return any(item["bbox"]["y"] >= FOOTER_Y - 0.1 for item in objects)


def _index_rail_overlaps_zone(objects: list[dict[str, Any]]) -> bool:
    for item in objects:
        bbox = item["bbox"]
        if bbox["x"] <= 0.35 and bbox["h"] >= 0.5:
            return True
    return False


def _load_contracts(directory: Path) -> dict[str, dict[str, Any]]:
    contracts: dict[str, dict[str, Any]] = {}
    if not directory.exists():
        return contracts
    for path in sorted(directory.glob("*.contract.json")):
        payload = _load_json(path)
        archetype_id = payload.get("archetype_id")
        if archetype_id:
            payload["_contract_path"] = _display_path(path)
            contracts[str(archetype_id)] = payload
    return contracts


def _golden_records(path: Path) -> dict[int, dict[str, Any]]:
    payload = _load_json(path)
    result: dict[int, dict[str, Any]] = {}
    for item in payload.get("compiled_layouts") or []:
        if isinstance(item, dict) and item.get("slide_number"):
            result[int(item["slide_number"])] = dict(item)
    return result


def _final_records(path: Path) -> dict[int, dict[str, Any]]:
    payload = _load_json(path)
    result: dict[int, dict[str, Any]] = {}
    for item in payload.get("compiled_slides") or []:
        if not isinstance(item, dict):
            continue
        slide_number = int(item.get("slide_number") or int(item.get("pptx_index") or 0) + 1)
        result[slide_number] = dict(item)
    return result


def _text_safe_archetypes(decks: list[dict[str, Any]]) -> dict[str, list[str]]:
    return {
        str(deck.get("scope")): sorted({str(slide.get("archetype_id")) for slide in deck.get("slides") or [] if slide.get("text_safe")})
        for deck in decks
    }


def _shape_bounds(shape: Any) -> dict[str, float]:
    return {
        "x": round(float(shape.left) / EMU_PER_INCH, 4),
        "y": round(float(shape.top) / EMU_PER_INCH, 4),
        "w": round(max(0.001, float(shape.width) / EMU_PER_INCH), 4),
        "h": round(max(0.001, float(shape.height) / EMU_PER_INCH), 4),
    }


def _normalize_bbox(raw: dict[str, Any]) -> dict[str, float] | None:
    try:
        return {
            "x": round(float(raw["x"]), 4),
            "y": round(float(raw["y"]), 4),
            "w": round(max(0.001, float(raw["w"])), 4),
            "h": round(max(0.001, float(raw["h"])), 4),
        }
    except (KeyError, TypeError, ValueError):
        return None


def _looks_like_small_ornament(bounds: dict[str, float]) -> bool:
    area = bounds["w"] * bounds["h"]
    return area <= 0.18 or bounds["w"] <= 0.12 or bounds["h"] <= 0.12


def _intersection_area(a: dict[str, float], b: dict[str, float]) -> float:
    x0 = max(a["x"], b["x"])
    y0 = max(a["y"], b["y"])
    x1 = min(a["x"] + a["w"], b["x"] + b["w"])
    y1 = min(a["y"] + a["h"], b["y"] + b["h"])
    if x1 <= x0 or y1 <= y0:
        return 0.0
    return round((x1 - x0) * (y1 - y0), 6)


def _canonical_archetype_id(archetype_id: str) -> str:
    return ARCHETYPE_ALIASES.get(str(archetype_id), str(archetype_id))


def _load_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    payload = json.loads(path.read_text(encoding="utf-8"))
    return payload if isinstance(payload, dict) else {}


def _display_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(Path.cwd().resolve())).replace("\\", "/")
    except ValueError:
        return str(path).replace("\\", "/")


def _markdown_report(report: dict[str, Any]) -> str:
    summary = report["findings_summary"]
    lines = [
        "# Protected Zone Audit Report",
        "",
        f"Status: `{report['status']}`",
        f"Findings: `{summary['total']}` total, `{summary['severe']}` severe, `{summary['warning']}` warnings, `{summary['minor']}` minor",
        "",
        "## Text-Safe Archetypes",
        "",
    ]
    for scope, archetypes in (report.get("text_safe_archetypes") or {}).items():
        lines.append(f"- `{scope}`: {', '.join(f'`{item}`' for item in archetypes) if archetypes else 'none'}")
    lines.extend(["", "## Decks", ""])
    for deck in report.get("decks") or []:
        deck_summary = deck.get("findings_summary") or {}
        lines.append(
            f"- `{deck.get('scope')}`: `{deck.get('status')}`, slides `{deck.get('slide_count')}`, "
            f"severe `{deck_summary.get('severe')}`, warnings `{deck_summary.get('warning')}`, minor `{deck_summary.get('minor')}`"
        )
    lines.extend(["", "## Findings", ""])
    for finding in report.get("findings") or []:
        lines.append(
            f"- `{finding['severity']}` `{finding['code']}` slide `{finding['slide_number']}` "
            f"`{finding['scope']}` `{finding['archetype_id']}` zone `{finding['zone_id']}`: {finding['message']}"
        )
    return "\n".join(lines) + "\n"


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Audit protected text/content zones for decorative and chrome intrusion.")
    parser.add_argument("--contracts-dir", type=Path, default=DEFAULT_CONTRACTS_DIR)
    parser.add_argument("--golden-pptx", type=Path, default=DEFAULT_GOLDEN_PPTX)
    parser.add_argument("--golden-report", type=Path, default=DEFAULT_GOLDEN_REPORT)
    parser.add_argument("--final-pptx", type=Path, default=DEFAULT_FINAL_PPTX)
    parser.add_argument("--final-manifest", type=Path, default=DEFAULT_FINAL_MANIFEST)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        report = build_protected_zone_audit_from_files(
            contracts_dir=args.contracts_dir,
            golden_pptx_path=args.golden_pptx,
            golden_report_path=args.golden_report,
            final_pptx_path=args.final_pptx,
            final_manifest_path=args.final_manifest,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"PROTECTED_ZONE_AUDIT_FAILED {type(exc).__name__}: {exc}")
        return 1
    print(f"WROTE {args.json_report}")
    print(f"WROTE {args.md_report}")
    print(f"PROTECTED_ZONE_AUDIT {report['status']}")
    return 1 if report["status"] == "failed" else 0


if __name__ == "__main__":
    raise SystemExit(main())
