"""Compile the E01.2 render-fidelity patched benchmark candidate."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = 16.0
SLIDE_H = 9.0


def build_editable_candidate_spec_e01_2(
    *,
    checklist_spec: dict[str, Any],
    bottom_action_spec: dict[str, Any],
    svg_role_map: dict[str, Any],
) -> dict[str, Any]:
    return {
        "schema_name": "editable_candidate_spec_e01_2",
        "slide_size": {"width_in": SLIDE_W, "height_in": SLIDE_H},
        "checklist_card_count": checklist_spec["card_count"],
        "bottom_action_count": bottom_action_spec["action_count"],
        "svg_role_count": svg_role_map["role_count"],
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "allowed_bounded_raster_targets": ["hero_photo_field", "thumbnail_callout_images"],
        "canva_parity_claimed": False,
    }


def compile_e01_2_candidate(
    *,
    reference_image: Path,
    oracle: dict[str, Any],
    checklist_spec: dict[str, Any],
    bottom_action_spec: dict[str, Any],
    output_pptx: Path,
) -> dict[str, Any]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    assets_dir = output_pptx.parent / "assets"
    assets_dir.mkdir(exist_ok=True)
    crops = _create_crops(reference_image, assets_dir)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_background(slide)
    hero = slide.shapes.add_picture(str(crops["hero"]), Inches(0), Inches(0), Inches(9.72), Inches(7.02))
    hero.name = "hero_photo_field_bounded_replaceable_image"
    _add_hero_depth_overlays(slide)
    _add_technical_overlay(slide)
    _add_thumbnail_callouts(slide, crops, oracle["thumbnail_callouts"])
    _add_checklist_v3(slide, oracle, checklist_spec)
    _add_bottom_action_bar_v3(slide, oracle, bottom_action_spec)
    _add_source_footer(slide)

    prs.save(output_pptx)
    return {
        "schema_name": "e01_2_candidate_compile_report",
        "status": "passed" if output_pptx.exists() else "failed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 1,
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use_count": 0,
        "bounded_raster_layer_count": 4,
        "canva_parity_claimed": False,
    }


def _create_crops(reference_image: Path, assets_dir: Path) -> dict[str, Path]:
    with Image.open(reference_image) as image:
        image = image.convert("RGB")
        w, h = image.size
        crops = {
            "hero": image.crop((0, 0, int(w * 0.605), int(h * 0.745))),
            "thumb_1": image.crop((int(w * 0.225), int(h * 0.585), int(w * 0.335), int(h * 0.72))),
            "thumb_2": image.crop((int(w * 0.352), int(h * 0.585), int(w * 0.462), int(h * 0.72))),
            "thumb_3": image.crop((int(w * 0.48), int(h * 0.585), int(w * 0.59), int(h * 0.72))),
        }
    paths: dict[str, Path] = {}
    for key, crop in crops.items():
        path = assets_dir / f"e01_2_{key}.png"
        crop.save(path)
        paths[key] = path
    return paths


def _add_background(slide: Any) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    shape.name = "background_base_native_dark_panel"
    _fill(shape, "061526")
    shape.line.fill.background()


def _add_hero_depth_overlays(slide: Any) -> None:
    # Native dark gradient stand-ins near the bottom/right keep the hero bounded while preserving text zones.
    for idx, (x, y, w, h, transparency) in enumerate(
        [(0, 6.52, 9.72, 0.5, 20), (9.15, 0.0, 0.58, 7.02, 12), (0.0, 0.0, 9.72, 0.42, 30)], start=1
    ):
        overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        overlay.name = f"hero_native_depth_overlay_{idx}"
        _fill(overlay, "061526", transparency=transparency)
        overlay.line.fill.background()


def _add_checklist_v3(slide: Any, oracle: dict[str, Any], checklist_spec: dict[str, Any]) -> None:
    panel = checklist_spec["panel_bbox_in"]
    panel_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(panel["x"]),
        Inches(panel["y"]),
        Inches(panel["w"]),
        Inches(panel["h"]),
    )
    panel_shape.name = "checklist_system_panel_v3_native"
    _fill(panel_shape, "061B29", transparency=3)
    _line(panel_shape, "24C5D8", 1.3)
    _add_line(slide, panel["x"] + 0.18, panel["y"] + 0.68, panel["x"] + panel["w"] - 0.18, panel["y"] + 0.68, "176F80", 0.85)
    _add_text(
        slide,
        "checklist_title_text_v3",
        oracle["title"],
        checklist_spec["title_bbox_in"]["x"],
        checklist_spec["title_bbox_in"]["y"],
        checklist_spec["title_bbox_in"]["w"],
        checklist_spec["title_bbox_in"]["h"],
        22,
        "58DDE8",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    for row in checklist_spec["cards"]:
        bbox = row["bbox_in"]
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(bbox["x"]), Inches(bbox["y"]), Inches(bbox["w"]), Inches(bbox["h"]))
        card.name = f"{row['card_id']}_panel_v3_native"
        _fill(card, "0A2A35", transparency=1)
        _line(card, "166C7A", 0.8)
        _add_line(slide, bbox["x"] + 1.27, bbox["y"], bbox["x"] + 1.27, bbox["y"] + bbox["h"], "174F5C", 0.7)
        _add_line(slide, bbox["x"] + 2.2, bbox["y"] + 0.2, bbox["x"] + 2.2, bbox["y"] + bbox["h"] - 0.2, "2FB9C7", 0.85)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(bbox["x"] + 0.16), Inches(bbox["y"] + 0.16), Inches(0.78), Inches(0.78))
        circle.name = f"{row['card_id']}_icon_circle_v3_vector"
        _fill(circle, "092231", transparency=25)
        _line(circle, "3DDCE8", 1.35)
        _add_role_icon(slide, row["icon_role"], bbox["x"] + 0.34, bbox["y"] + 0.34, 0.42, "6DE7F3")
        _add_text(slide, f"{row['card_id']}_number_text_v3", row["number"], bbox["x"] + 1.46, bbox["y"] + 0.26, 0.58, 0.48, 27, "36C7D9", bold=True)
        _add_text(slide, f"{row['card_id']}_heading_text_v3", row["heading"], bbox["x"] + 2.42, bbox["y"] + 0.19, 2.6, 0.27, 13, "F8FAFC", bold=True)
        _add_text(slide, f"{row['card_id']}_body_text_v3", row["body"], bbox["x"] + 2.42, bbox["y"] + 0.5, 2.55, 0.42, 10, "F5F7FA")
        chevron = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, Inches(bbox["x"] + bbox["w"] - 0.38), Inches(bbox["y"] + 0.43), Inches(0.22), Inches(0.27))
        chevron.name = f"{row['card_id']}_chevron_v3_vector"
        _fill(chevron, "3DDCE8")
        chevron.line.fill.background()


def _add_thumbnail_callouts(slide: Any, crops: dict[str, Path], callouts: list[dict[str, Any]]) -> None:
    positions = [(3.62, 5.6), (5.56, 5.6), (7.49, 5.6)]
    for idx, callout in enumerate(callouts, start=1):
        x, y = positions[idx - 1]
        pic = slide.shapes.add_picture(str(crops[f"thumb_{idx}"]), Inches(x), Inches(y), Inches(1.58), Inches(1.16))
        pic.name = f"thumbnail_{idx}_bounded_replaceable_image_v3"
        ring_outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.08), Inches(y - 0.1), Inches(1.74), Inches(1.36))
        ring_outer.name = f"thumbnail_{idx}_outer_ring_v3_vector"
        ring_outer.fill.background()
        _line(ring_outer, "3DDCE8", 1.5)
        _add_text(slide, f"thumbnail_{idx}_caption_text_v3", callout["label"], x - 0.02, y + 1.34, 1.78, 0.22, 8, "F8FAFC", bold=True, align=PP_ALIGN.CENTER)


def _add_bottom_action_bar_v3(slide: Any, oracle: dict[str, Any], bottom_action_spec: dict[str, Any]) -> None:
    bar = bottom_action_spec["bar_bbox_in"]
    container = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(bar["x"]), Inches(bar["y"]), Inches(bar["w"]), Inches(bar["h"]))
    container.name = "bottom_action_bar_container_v3_native"
    _fill(container, "06111A", transparency=0)
    container.line.fill.background()
    _add_line(slide, 0.0, bar["y"], SLIDE_W, bar["y"], "F59E1B", 1.2)
    _add_line(slide, 0.45, 8.73, 14.95, 8.73, "B87518", 0.85)
    _add_diagonal_accents(slide)

    for idx, action in enumerate(bottom_action_spec["actions"]):
        bbox = action["bbox_in"]
        if action["separator"]:
            _add_line(slide, bbox["x"] - 0.28, bbox["y"] + 0.08, bbox["x"] - 0.28, bbox["y"] + 0.66, "1D4550", 0.75)
        _add_role_icon(slide, action["icon_role"], bbox["x"], bbox["y"] + 0.05, 0.58, "F5A623")
        _add_text(slide, f"{action['action_id']}_top_label_v3", action["label_top"], bbox["x"] + 0.75, bbox["y"] + 0.08, 1.6, 0.24, 10, "F5A623", bold=True)
        _add_text(slide, f"{action['action_id']}_bottom_label_v3", action["label_bottom"], bbox["x"] + 0.75, bbox["y"] + 0.38, 1.85, 0.24, 9, "F5A623", bold=True)


def _add_diagonal_accents(slide: Any) -> None:
    for idx, x in enumerate((0.18, 0.52, 15.48, 15.78), start=1):
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(x), Inches(7.58 if idx < 3 else 8.28), Inches(0.32), Inches(0.32))
        accent.name = f"bottom_bar_gold_diagonal_accent_{idx}"
        _fill(accent, "F5A623")
        accent.line.fill.background()


def _add_source_footer(slide: Any) -> None:
    _add_line(slide, 0.42, 8.89, 15.25, 8.89, "174B58", 0.65)
    _add_text(slide, "source_footer_text_v3_editable", "SOURCE / FOOTER SLOT", 0.45, 8.82, 2.35, 0.16, 6, "A9D8DE")


def _add_technical_overlay(slide: Any) -> None:
    cx, cy = 5.82, 1.25
    for idx, size in enumerate((1.72, 1.34, 0.98, 0.62, 0.26), start=1):
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - size / 2), Inches(cy - size / 2), Inches(size), Inches(size))
        ring.name = f"technical_overlay_radar_ring_v3_{idx}"
        ring.fill.background()
        _line(ring, "0DAEC3", 0.65)
    _add_line(slide, cx - 1.2, cy, cx + 1.2, cy, "0DAEC3", 0.55)
    _add_line(slide, cx, cy - 1.1, cx, cy + 1.1, "0DAEC3", 0.55)
    for x1, y1, x2, y2 in ((0.35, 0.5, 4.7, 0.5), (3.0, 6.24, 9.2, 6.24), (6.84, 1.25, 8.78, 1.25), (8.78, 1.25, 9.1, 0.96)):
        _add_line(slide, x1, y1, x2, y2, "0DAEC3", 0.55)
    for x, y in ((4.63, 1.36), (6.08, 0.78), (8.62, 6.32), (3.32, 6.24), (5.82, 1.25)):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.065), Inches(0.065))
        dot.name = "technical_overlay_dot_v3_vector"
        _fill(dot, "23CAE0")
        dot.line.fill.background()


def _add_role_icon(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    if "clipboard" in role or "record" in role:
        _icon_document(slide, role, x, y, size, color)
    elif "valve" in role:
        _icon_valve(slide, role, x, y, size, color)
    elif "gauge" in role:
        _icon_gauge(slide, role, x, y, size, color)
    elif "shield" in role or "barrier" in role:
        _icon_shield(slide, role, x, y, size, color)
    elif "warning" in role:
        _icon_warning(slide, role, x, y, size, color)
    elif "lock" in role:
        _icon_lock(slide, role, x, y, size, color)
    elif "chat" in role:
        _icon_chat(slide, role, x, y, size, color)
    elif "team" in role:
        _icon_team(slide, role, x, y, size, color)
    else:
        _icon_gauge(slide, role, x, y, size, color)


def _icon_document(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    doc = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(size * 0.72), Inches(size))
    doc.name = f"{role}_svg_role_document_vector"
    doc.fill.background()
    _line(doc, color, 1.2)
    for i in range(3):
        _add_line(slide, x + size * 0.16, y + size * (0.25 + i * 0.18), x + size * 0.56, y + size * (0.25 + i * 0.18), color, 0.9)
    _add_line(slide, x + size * 0.2, y + size * 0.77, x + size * 0.34, y + size * 0.9, color, 1.2)
    _add_line(slide, x + size * 0.34, y + size * 0.9, x + size * 0.62, y + size * 0.58, color, 1.2)


def _icon_valve(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    _add_line(slide, x, y + size * 0.72, x + size, y + size * 0.72, color, 1.6)
    _add_line(slide, x + size * 0.5, y + size * 0.2, x + size * 0.5, y + size * 0.72, color, 1.6)
    _add_line(slide, x + size * 0.25, y + size * 0.2, x + size * 0.75, y + size * 0.2, color, 1.6)
    for px in (x, x + size * 0.78):
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(px), Inches(y + size * 0.54), Inches(size * 0.22), Inches(size * 0.36))
        body.name = f"{role}_svg_role_valve_vector"
        body.fill.background()
        _line(body, color, 1.2)


def _icon_gauge(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    gauge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    gauge.name = f"{role}_svg_role_gauge_vector"
    gauge.fill.background()
    _line(gauge, color, 1.2)
    _add_line(slide, x + size * 0.5, y + size * 0.5, x + size * 0.75, y + size * 0.28, color, 1.4)


def _icon_shield(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    shield = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(x), Inches(y), Inches(size), Inches(size))
    shield.name = f"{role}_svg_role_shield_vector"
    shield.fill.background()
    _line(shield, color, 1.2)
    _add_line(slide, x + size * 0.25, y + size * 0.52, x + size * 0.43, y + size * 0.7, color, 1.4)
    _add_line(slide, x + size * 0.43, y + size * 0.7, x + size * 0.76, y + size * 0.32, color, 1.4)


def _icon_warning(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    tri.name = f"{role}_svg_role_warning_vector"
    tri.fill.background()
    _line(tri, color, 1.4)
    _add_line(slide, x + size * 0.5, y + size * 0.35, x + size * 0.5, y + size * 0.64, color, 1.2)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * 0.47), Inches(y + size * 0.74), Inches(size * 0.06), Inches(size * 0.06))
    dot.name = f"{role}_svg_role_warning_dot_vector"
    _fill(dot, color)
    dot.line.fill.background()


def _icon_lock(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.12), Inches(y + size * 0.44), Inches(size * 0.76), Inches(size * 0.48))
    body.name = f"{role}_svg_role_lock_vector"
    body.fill.background()
    _line(body, color, 1.3)
    _add_line(slide, x + size * 0.28, y + size * 0.44, x + size * 0.28, y + size * 0.24, color, 1.2)
    _add_line(slide, x + size * 0.72, y + size * 0.44, x + size * 0.72, y + size * 0.24, color, 1.2)
    _add_line(slide, x + size * 0.28, y + size * 0.24, x + size * 0.72, y + size * 0.24, color, 1.2)


def _icon_chat(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + size * 0.1), Inches(size), Inches(size * 0.68))
    bubble.name = f"{role}_svg_role_chat_vector"
    bubble.fill.background()
    _line(bubble, color, 1.3)
    for i in range(3):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * (0.25 + i * 0.2)), Inches(y + size * 0.42), Inches(size * 0.08), Inches(size * 0.08))
        dot.name = f"{role}_svg_role_chat_dot_vector"
        _fill(dot, color)
        dot.line.fill.background()


def _icon_team(slide: Any, role: str, x: float, y: float, size: float, color: str) -> None:
    for i, offset in enumerate((0.08, 0.38, 0.68), start=1):
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * offset), Inches(y + size * 0.12), Inches(size * 0.2), Inches(size * 0.2))
        head.name = f"{role}_svg_role_team_head_{i}_vector"
        _fill(head, color)
        head.line.fill.background()
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * (offset - 0.04)), Inches(y + size * 0.42), Inches(size * 0.28), Inches(size * 0.34))
        body.name = f"{role}_svg_role_team_body_{i}_vector"
        _fill(body, color)
        body.line.fill.background()


def _add_text(
    slide: Any,
    name: str,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: str,
    *,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.01)
    frame.margin_right = Inches(0.01)
    frame.margin_top = Inches(0.0)
    frame.word_wrap = True
    lines = text.splitlines() or [text]
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = ""
        if align is not None:
            paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)


def _add_line(slide: Any, x1: float, y1: float, x2: float, y2: float, color: str, width_pt: float) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = "native_connector_line_v3"
    line.line.color.rgb = RGBColor.from_string(color)
    line.line.width = Pt(width_pt)


def _fill(shape: Any, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.fill.transparency = transparency


def _line(shape: Any, color: str, width_pt: float) -> None:
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(width_pt)

