from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from ..slide_scene import Callout, TextRun
from .common import append_warning, stable_scene_shape_name
from .style import SceneStyleContext, resolve_fill_style, resolve_spacing, resolve_stroke_style
from .text import _apply_run_style, _clear_bullet_style, _paragraph_groups


def render_callout(
    slide: Any,
    callout: Callout,
    *,
    style_context: SceneStyleContext,
    used_shape_names: set[str],
    slide_number: int,
    slide_id: str,
    warnings: list[Any],
) -> str:
    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(callout.bounds.x),
        Inches(callout.bounds.y),
        Inches(callout.bounds.width),
        Inches(callout.bounds.height),
    )
    background.name = stable_scene_shape_name("callout", f"{callout.object_id}:background", used_shape_names)
    background.fill.solid()
    background.fill.fore_color.rgb = resolve_fill_style(
        style_context,
        None,
        "#FFFFFF",
        slide_number=slide_number,
        slide_id=slide_id,
        object_id=callout.object_id,
    ).color_rgb
    background.line.fill.solid()
    background.line.fill.fore_color.rgb = resolve_stroke_style(
        style_context,
        callout.accent,
        default_hex="#CBD5E1",
        slide_number=slide_number,
        slide_id=slide_id,
        object_id=callout.object_id,
    ).color_rgb

    accent_width = min(0.12, max(0.06, round(callout.bounds.width * 0.025, 6)))
    padding = resolve_spacing(style_context, "md", 0.22, slide_number=slide_number, slide_id=slide_id, object_id=callout.object_id)
    inner_gap = resolve_spacing(style_context, "xs", 0.08, slide_number=slide_number, slide_id=slide_id, object_id=callout.object_id)
    text_left = callout.bounds.x + padding
    content_width = max(0.4, round(callout.bounds.width - 0.44, 6))
    if callout.accent is not None:
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(callout.bounds.x + 0.12),
            Inches(callout.bounds.y + 0.12),
            Inches(accent_width),
            Inches(max(0.2, round(callout.bounds.height - 0.24, 6))),
        )
        accent.name = stable_scene_shape_name("callout", f"{callout.object_id}:accent", used_shape_names)
        accent.fill.solid()
        accent.fill.fore_color.rgb = resolve_fill_style(
            style_context,
            callout.accent,
            "#C2410C",
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=callout.object_id,
        ).color_rgb
        accent.line.fill.background()
        text_left += accent_width + inner_gap
        content_width = max(0.3, round((callout.bounds.x + callout.bounds.width - 0.22) - text_left, 6))

    if callout.fit.mode not in {"none", "wrap", "fail"}:
        append_warning(
            warnings,
            code="fit_policy_not_enforced",
            severity="warning",
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=callout.object_id,
            message=f"Callout fit policy {callout.fit.mode!r} is not enforced exactly in the scene compile path.",
            details={"fit_mode": callout.fit.mode},
        )

    top = callout.bounds.y + 0.16
    if callout.title is not None:
        title_height = min(0.55, max(0.32, round(callout.bounds.height * 0.2, 6)))
        title_shape = slide.shapes.add_textbox(
            Inches(text_left),
            Inches(top),
            Inches(content_width),
            Inches(title_height),
        )
        title_shape.name = stable_scene_shape_name("callout", f"{callout.object_id}:title", used_shape_names)
        title_shape.fill.background()
        title_shape.line.fill.background()
        _render_runs(
            title_shape.text_frame,
            [callout.title],
            role="title",
            style_context=style_context,
            wrap=True,
            slide_number=slide_number,
            slide_id=slide_id,
            object_id=callout.object_id,
        )
        top += title_height + 0.08

    body_height = max(0.3, round((callout.bounds.y + callout.bounds.height - 0.16) - top, 6))
    body_shape = slide.shapes.add_textbox(
        Inches(text_left),
        Inches(top),
        Inches(content_width),
        Inches(body_height),
    )
    body_shape.name = stable_scene_shape_name("callout", f"{callout.object_id}:body", used_shape_names)
    body_shape.fill.background()
    body_shape.line.fill.background()
    _render_runs(
        body_shape.text_frame,
        callout.body,
        role="body",
        style_context=style_context,
        wrap=callout.fit.mode != "none",
        slide_number=slide_number,
        slide_id=slide_id,
        object_id=callout.object_id,
    )
    return background.name


def _render_runs(
    text_frame: Any,
    runs: list[TextRun],
    *,
    role: str,
    style_context: SceneStyleContext,
    wrap: bool,
    slide_number: int,
    slide_id: str,
    object_id: str,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = wrap
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraphs = _paragraph_groups(runs)
    for paragraph_index, group in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if paragraph_index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        _clear_bullet_style(paragraph)
        for run_index, (source_run, text) in enumerate(group):
            run = paragraph.add_run() if run_index > 0 else paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            _apply_run_style(run, source_run.model_copy(update={"text": text}), role, style_context, slide_number, slide_id, object_id)
