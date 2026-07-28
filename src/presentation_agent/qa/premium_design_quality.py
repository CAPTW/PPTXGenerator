"""Premium design acceptance QA for Codex/GPT-Image-2 design-board decks."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..generator_contracts import (
    validateEditableTemplateSpec,
    validateExtractedComponentLibrary,
    validateExtractedDesignBoardCodex,
    validateExtractedSlideArchetypes,
    validateExtractedStyleTokens,
)
from .final_deck_image_policy import build_final_deck_image_policy_report


DEFAULT_FINAL_DECK = Path("outputs/final_deck.pptx")
DEFAULT_LARGE_DECK = Path("outputs/final_deck_large.pptx")
DEFAULT_ASSEMBLY_PLAN = Path("outputs/deck_assembly_plan.json")
DEFAULT_LARGE_PREMIUM_ASSEMBLY_PLAN = Path("outputs/deck_assembly_plan_large_premium.json")
DEFAULT_FINAL_SPEC = Path("outputs/editable_template_spec.final.json")
DEFAULT_BOARD_MANIFEST = Path("outputs/template_design_board/template_design_board_manifest.json")
DEFAULT_DESIGN_EXTRACTION_DIR = Path("outputs/design_extraction")
DEFAULT_TEMPLATE_VISUAL_DIFF = Path("outputs/template_visual_diff_report.json")
DEFAULT_TEMPLATE_IMAGE_MANIFEST = Path("outputs/template_images/template_image_manifest.json")
DEFAULT_SPEC_FROM_BOARD_REPORT = Path("outputs/template_spec_from_design_board_report.json")
DEFAULT_CROP_MANIFEST = Path("outputs/template_design_board/design_board_crop_manifest.json")
DEFAULT_JSON_REPORT = Path("outputs/premium_design_quality_report.json")
DEFAULT_MD_REPORT = Path("outputs/premium_design_quality_report.md")
CANONICAL_PROMPT_MANIFEST = Path("design_prompts/prompt_manifest.json")
SMOKE_PHRASES = (
    "deterministic slot filling",
    "large deck smoke synthetic source",
    "Large Editable Template Smoke Deck",
    "Image frame",
)
REQUIRED_EXTRACTION_ARTIFACTS: dict[str, tuple[str, Callable[[dict[str, Any]], dict[str, Any]]]] = {
    "extracted_design_board.codex.json": ("extracted_design_board", validateExtractedDesignBoardCodex),
    "extracted_slide_archetypes.json": ("extracted_slide_archetypes", validateExtractedSlideArchetypes),
    "extracted_component_library.json": ("extracted_component_library", validateExtractedComponentLibrary),
    "extracted_style_tokens.json": ("extracted_style_tokens", validateExtractedStyleTokens),
}


def build_premium_design_quality_report(
    *,
    pptx_path: str | Path | None = None,
    deck_assembly_plan_path: str | Path = DEFAULT_ASSEMBLY_PLAN,
    editable_template_spec_path: str | Path = DEFAULT_FINAL_SPEC,
    design_board_manifest_path: str | Path = DEFAULT_BOARD_MANIFEST,
    design_extraction_dir: str | Path = DEFAULT_DESIGN_EXTRACTION_DIR,
    template_visual_diff_report_path: str | Path = DEFAULT_TEMPLATE_VISUAL_DIFF,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    template_spec_from_board_report_path: str | Path = DEFAULT_SPEC_FROM_BOARD_REPORT,
    design_board_crop_manifest_path: str | Path = DEFAULT_CROP_MANIFEST,
    run_label: str = "premium",
) -> dict[str, Any]:
    pptx = Path(pptx_path) if pptx_path is not None else _default_pptx_path()
    assembly_path = _resolve_assembly_plan_path(pptx, Path(deck_assembly_plan_path))
    spec_path = Path(editable_template_spec_path)
    board_manifest_path = Path(design_board_manifest_path)
    extraction_dir = Path(design_extraction_dir)
    visual_diff_path = Path(template_visual_diff_report_path)
    template_manifest_path = Path(template_image_manifest_path)
    spec_from_board_report_path = Path(template_spec_from_board_report_path)
    crop_manifest_path = Path(design_board_crop_manifest_path)

    findings: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    pptx_stats = _pptx_stats(pptx, findings)
    visible_text = pptx_stats.pop("visible_text", [])
    smoke_phrase_hits = _smoke_phrase_hits(visible_text)
    for phrase in smoke_phrase_hits:
        findings.append(_finding("SMOKE_TEXT_VISIBLE", "severe", f"Smoke phrase is visible in final deck: {phrase}", details={"phrase": phrase}))

    assembly_plan = _load_optional_json(assembly_path, findings, "DECK_ASSEMBLY_PLAN_MISSING")
    final_spec = _load_optional_json(spec_path, findings, "EDITABLE_TEMPLATE_SPEC_FINAL_MISSING")
    board_manifest = _load_optional_json(board_manifest_path, findings, "DESIGN_BOARD_MANIFEST_MISSING")
    prompt_manifest = _load_optional_json(CANONICAL_PROMPT_MANIFEST, findings, "CANONICAL_PROMPT_MANIFEST_MISSING")
    visual_diff = _load_optional_json(visual_diff_path, warnings, "TEMPLATE_VISUAL_DIFF_REPORT_MISSING", severity="warning")
    spec_from_board_report = _load_optional_json(spec_from_board_report_path, findings, "TEMPLATE_SPEC_FROM_DESIGN_BOARD_REPORT_MISSING")
    crop_manifest = _load_optional_json(crop_manifest_path, findings, "DESIGN_BOARD_CROP_MANIFEST_MISSING")

    spec_valid = _validate_payload(final_spec, validateEditableTemplateSpec, findings, "EDITABLE_TEMPLATE_SPEC_FINAL_INVALID")
    extraction_status = _extraction_status(extraction_dir, findings)
    image_policy = _image_policy_status(pptx, spec_path, assembly_path, template_manifest_path, findings)
    board_status = _design_board_status(board_manifest, findings)
    spec_derivation_status = _spec_derivation_status(final_spec, spec_from_board_report, findings)
    crop_status = _crop_manifest_status(crop_manifest, findings)
    tone_status = _tone_expression_status(assembly_plan, final_spec, findings)
    scale_status = _deck_scale_status(assembly_plan, pptx_stats, findings)
    reference_status = _reference_fidelity_status(visual_diff, run_label, findings, warnings)
    component_status = _component_system_status(final_spec, findings, warnings)
    layout_status = _layout_quality_status(assembly_plan, final_spec, findings, warnings)
    functional_smoke_status = _functional_smoke_status(smoke_phrase_hits, assembly_plan)

    generation_mode = str(board_manifest.get("generation_mode") or "")
    prompt_manifest_status = _prompt_manifest_status(prompt_manifest, findings)
    if run_label in {"premium", "design-accepted"} and generation_mode == "mock_fixture":
        findings.append(_finding("MOCK_FIXTURE_LABELED_PREMIUM", "severe", "A mock fixture design run cannot be labeled premium."))

    if image_policy.get("reference_template_image_embedded_count", 0) > 0:
        findings.append(_finding("TEMPLATE_REFERENCE_IMAGE_EMBEDDED", "severe", "A template reference PNG is embedded directly in the final deck."))
    if image_policy.get("full_slide_picture_count", 0) > 0:
        findings.append(_finding("FULL_SLIDE_RASTER_IMAGE", "severe", "A full-slide raster image exists in the final deck."))

    severe_count = sum(1 for finding in findings if finding["severity"] == "severe")
    premium_design_status = "passed" if severe_count == 0 else "failed"
    report = {
        "schema_name": "premium_design_quality_report",
        "schema_version": "1.0",
        "run_label": run_label,
        "status": premium_design_status,
        "premium_design_status": premium_design_status,
        "functional_smoke_status": functional_smoke_status,
        "editability_status": _editability_status(pptx_stats, spec_valid),
        "image_policy_status": image_policy.get("status", "unknown"),
        "reference_fidelity_status": reference_status,
        "tone_expression_status": tone_status,
        "deck_scale_status": scale_status,
        "design_board_status": board_status,
        "codex_extraction_status": extraction_status,
        "template_spec_derivation_status": spec_derivation_status,
        "design_board_crop_status": crop_status,
        "component_system_status": component_status,
        "layout_quality_status": layout_status,
        "premium_design_run": premium_design_status == "passed",
        "premium_design_run_statement": _premium_statement(premium_design_status),
        "pptx_path": _display_path(pptx),
        "deck_assembly_plan_path": _display_path(assembly_path),
        "editable_template_spec_path": _display_path(spec_path),
        "template_design_board_manifest_path": _display_path(board_manifest_path),
        "design_board_crop_manifest_path": _display_path(crop_manifest_path) if crop_manifest_path.exists() else None,
        "canonical_prompt_manifest_path": _display_path(CANONICAL_PROMPT_MANIFEST),
        "template_visual_diff_report_path": _display_path(visual_diff_path) if visual_diff_path.exists() else None,
        "pptx_stats": pptx_stats,
        "smoke_phrase_hits": smoke_phrase_hits,
        "generation_mode": generation_mode or None,
        "prompt_manifest_present": bool(prompt_manifest),
        "prompt_manifest_status": prompt_manifest_status,
        "image_policy": {
            "picture_shape_count": image_policy.get("picture_shape_count"),
            "full_slide_picture_count": image_policy.get("full_slide_picture_count"),
            "reference_template_image_embedded_count": image_policy.get("reference_template_image_embedded_count"),
            "undeclared_picture_shape_count": image_policy.get("undeclared_picture_shape_count"),
            "allowed_photo_frame_picture_count": image_policy.get("allowed_photo_frame_picture_count"),
        },
        "findings_summary": {
            "total": len(findings) + len(warnings),
            "severe": severe_count,
            "warning": len(warnings),
        },
        "findings": findings,
        "warnings": warnings,
    }
    return report


def build_premium_design_quality_report_from_files(
    *,
    pptx_path: str | Path | None = None,
    deck_assembly_plan_path: str | Path = DEFAULT_ASSEMBLY_PLAN,
    editable_template_spec_path: str | Path = DEFAULT_FINAL_SPEC,
    design_board_manifest_path: str | Path = DEFAULT_BOARD_MANIFEST,
    design_extraction_dir: str | Path = DEFAULT_DESIGN_EXTRACTION_DIR,
    template_visual_diff_report_path: str | Path = DEFAULT_TEMPLATE_VISUAL_DIFF,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    template_spec_from_board_report_path: str | Path = DEFAULT_SPEC_FROM_BOARD_REPORT,
    design_board_crop_manifest_path: str | Path = DEFAULT_CROP_MANIFEST,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
    run_label: str = "premium",
) -> Path:
    report = build_premium_design_quality_report(
        pptx_path=pptx_path,
        deck_assembly_plan_path=deck_assembly_plan_path,
        editable_template_spec_path=editable_template_spec_path,
        design_board_manifest_path=design_board_manifest_path,
        design_extraction_dir=design_extraction_dir,
        template_visual_diff_report_path=template_visual_diff_report_path,
        template_image_manifest_path=template_image_manifest_path,
        template_spec_from_board_report_path=template_spec_from_board_report_path,
        design_board_crop_manifest_path=design_board_crop_manifest_path,
        run_label=run_label,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return json_path


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Fail premium design acceptance when a deck is only a functional smoke artifact.")
    parser.add_argument("--pptx", type=Path, default=None)
    parser.add_argument("--deck-assembly-plan", type=Path, default=DEFAULT_ASSEMBLY_PLAN)
    parser.add_argument("--editable-template-spec", type=Path, default=DEFAULT_FINAL_SPEC)
    parser.add_argument("--design-board-manifest", type=Path, default=DEFAULT_BOARD_MANIFEST)
    parser.add_argument("--design-extraction-dir", type=Path, default=DEFAULT_DESIGN_EXTRACTION_DIR)
    parser.add_argument("--template-visual-diff-report", type=Path, default=DEFAULT_TEMPLATE_VISUAL_DIFF)
    parser.add_argument("--template-image-manifest", type=Path, default=DEFAULT_TEMPLATE_IMAGE_MANIFEST)
    parser.add_argument("--template-spec-from-board-report", type=Path, default=DEFAULT_SPEC_FROM_BOARD_REPORT)
    parser.add_argument("--design-board-crop-manifest", type=Path, default=DEFAULT_CROP_MANIFEST)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    parser.add_argument("--run-label", choices=["premium", "design-accepted", "smoke"], default="premium")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        output = build_premium_design_quality_report_from_files(
            pptx_path=args.pptx,
            deck_assembly_plan_path=args.deck_assembly_plan,
            editable_template_spec_path=args.editable_template_spec,
            design_board_manifest_path=args.design_board_manifest,
            design_extraction_dir=args.design_extraction_dir,
            template_visual_diff_report_path=args.template_visual_diff_report,
            template_image_manifest_path=args.template_image_manifest,
            template_spec_from_board_report_path=args.template_spec_from_board_report,
            design_board_crop_manifest_path=args.design_board_crop_manifest,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
            run_label=args.run_label,
        )
        report = _load_json(output)
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"PREMIUM_DESIGN_QA_FAILED {exc}")
        return 1
    print(f"WROTE {output}")
    if report.get("premium_design_status") != "passed":
        print("PREMIUM_DESIGN_QA failed")
        for finding in report.get("findings", []):
            if finding.get("severity") == "severe":
                print(f"PREMIUM_DESIGN_QA_FAILURE {finding.get('code')}: {finding.get('message')}")
        return 1
    print("PREMIUM_DESIGN_QA passed")
    return 0


def _default_pptx_path() -> Path:
    return DEFAULT_FINAL_DECK if DEFAULT_FINAL_DECK.exists() else DEFAULT_LARGE_DECK


def _resolve_assembly_plan_path(pptx: Path, assembly_path: Path) -> Path:
    if (
        assembly_path == DEFAULT_ASSEMBLY_PLAN
        and pptx.name == "final_deck_large_premium.pptx"
        and DEFAULT_LARGE_PREMIUM_ASSEMBLY_PLAN.exists()
    ):
        return DEFAULT_LARGE_PREMIUM_ASSEMBLY_PLAN
    return assembly_path


def _pptx_stats(path: Path, findings: list[dict[str, Any]]) -> dict[str, Any]:
    if not path.exists():
        findings.append(_finding("FINAL_DECK_MISSING", "severe", f"Final deck is missing: {_display_path(path)}"))
        return {"slide_count": 0, "shape_count": 0, "text_shape_count": 0, "table_count": 0, "chart_count": 0, "picture_shape_count": 0, "visible_text": []}
    deck = Presentation(path)
    stats = {"slide_count": len(deck.slides), "shape_count": 0, "text_shape_count": 0, "table_count": 0, "chart_count": 0, "picture_shape_count": 0, "visible_text": []}
    for slide in deck.slides:
        for shape in slide.shapes:
            stats["shape_count"] += 1
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    stats["text_shape_count"] += 1
                    stats["visible_text"].append(text)
            if getattr(shape, "has_table", False):
                stats["table_count"] += 1
            if getattr(shape, "has_chart", False):
                stats["chart_count"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                stats["picture_shape_count"] += 1
    return stats


def _smoke_phrase_hits(visible_text: list[str]) -> list[str]:
    deck_text = "\n".join(visible_text).lower()
    return [phrase for phrase in SMOKE_PHRASES if phrase.lower() in deck_text]


def _image_policy_status(pptx: Path, spec_path: Path, assembly_path: Path, template_manifest_path: Path, findings: list[dict[str, Any]]) -> dict[str, Any]:
    if not pptx.exists():
        return {"status": "skipped"}
    try:
        report = build_final_deck_image_policy_report(
            pptx_path=pptx,
            template_spec_path=spec_path,
            deck_assembly_plan_path=assembly_path,
            template_image_manifest_path=template_manifest_path,
        )
    except Exception as exc:
        findings.append(_finding("IMAGE_POLICY_QA_FAILED", "severe", f"Image policy QA could not run: {type(exc).__name__}: {exc}"))
        return {"status": "failed"}
    if report.get("status") != "passed":
        findings.append(_finding("IMAGE_POLICY_FAILED", "severe", "Final deck image policy failed.", details={"findings": report.get("findings", [])}))
    return report


def _design_board_status(board_manifest: dict[str, Any], findings: list[dict[str, Any]]) -> str:
    if not board_manifest:
        return "failed"
    if board_manifest.get("generation_mode") != "manual_codex":
        findings.append(_finding("DESIGN_BOARD_NOT_MANUAL_CODEX", "severe", "Premium design requires generation_mode=manual_codex."))
    if board_manifest.get("prompt_id") != "creative_academic_template_board_v1":
        findings.append(_finding("DESIGN_BOARD_PROMPT_ID_MISMATCH", "severe", "Premium design requires prompt_id=creative_academic_template_board_v1."))
    if board_manifest.get("premium_design_candidate") is not True:
        findings.append(_finding("DESIGN_BOARD_NOT_PREMIUM_CANDIDATE", "severe", "Template design board manifest must mark premium_design_candidate=true."))
    image_path = Path(str(board_manifest.get("board_image_path") or ""))
    if not image_path.exists():
        findings.append(_finding("DESIGN_BOARD_IMAGE_MISSING", "severe", "Design board image is missing.", details={"board_image_path": str(board_manifest.get("board_image_path"))}))
    if board_manifest.get("board_mode") != "single_design_board":
        findings.append(_finding("DESIGN_BOARD_MODE_INVALID", "severe", "Premium design requires the single design board workflow."))
    return "passed" if board_manifest.get("generation_mode") == "manual_codex" and board_manifest.get("prompt_id") == "creative_academic_template_board_v1" and board_manifest.get("premium_design_candidate") is True and image_path.exists() else "failed"


def _extraction_status(extraction_dir: Path, findings: list[dict[str, Any]]) -> str:
    passed = True
    for filename, (label, validator) in REQUIRED_EXTRACTION_ARTIFACTS.items():
        path = extraction_dir / filename
        if not path.exists():
            findings.append(_finding("CODEX_EXTRACTION_ARTIFACT_MISSING", "severe", f"Codex-assisted extraction artifact is missing: {path.as_posix()}", details={"artifact": label, "path": path.as_posix()}))
            passed = False
            continue
        payload = _load_json(path)
        if not _validate_payload(payload, validator, findings, f"{label.upper()}_INVALID"):
            passed = False
        elif label == "extracted_slide_archetypes" and not _archetypes_have_numeric_geometry(payload):
            findings.append(_finding("EXTRACTED_ARCHETYPES_NUMERIC_GEOMETRY_MISSING", "severe", "Extracted slide archetypes must include normalized_geometry and slot_geometry for every archetype."))
            passed = False
    return "passed" if passed else "failed"


def _archetypes_have_numeric_geometry(payload: dict[str, Any]) -> bool:
    for archetype in payload.get("archetypes") or []:
        if not isinstance(archetype, dict):
            return False
        normalized = archetype.get("normalized_geometry")
        slots = archetype.get("slot_geometry")
        if not isinstance(normalized, dict) or not isinstance(slots, list) or not slots:
            return False
        for zone in ("title_zone", "navigation_zone", "content_zone", "visual_zone", "footer_zone"):
            if not _is_numeric_box(normalized.get(zone)):
                return False
        if not isinstance(normalized.get("safe_margin"), dict):
            return False
        if any(not _is_numeric_box(slot) or "role" not in slot or "priority" not in slot for slot in slots):
            return False
    return True


def _is_numeric_box(payload: Any) -> bool:
    return isinstance(payload, dict) and all(isinstance(payload.get(key), (int, float)) for key in ("x", "y", "w", "h"))


def _crop_manifest_status(crop_manifest: dict[str, Any], findings: list[dict[str, Any]]) -> str:
    if not crop_manifest:
        return "failed"
    if crop_manifest.get("prompt_id") != "creative_academic_template_board_v1":
        findings.append(_finding("DESIGN_BOARD_CROP_PROMPT_ID_MISMATCH", "severe", "Crop manifest prompt_id must match creative_academic_template_board_v1."))
    if int(crop_manifest.get("crop_count") or 0) < 24:
        findings.append(_finding("DESIGN_BOARD_CROPS_INCOMPLETE", "severe", "Crop manifest must include all hero, thumbnail, component, layout, storytelling, and style-token crops."))
    missing_paths = [
        crop.get("path")
        for crop in crop_manifest.get("crops") or []
        if isinstance(crop, dict) and crop.get("path") and not Path(str(crop["path"])).exists()
    ]
    if missing_paths:
        findings.append(_finding("DESIGN_BOARD_CROP_IMAGE_MISSING", "severe", "One or more design board crop images are missing.", details={"missing_paths": missing_paths[:8]}))
    return "passed" if crop_manifest.get("prompt_id") == "creative_academic_template_board_v1" and int(crop_manifest.get("crop_count") or 0) >= 24 and not missing_paths else "failed"


def _spec_derivation_status(final_spec: dict[str, Any], report: dict[str, Any], findings: list[dict[str, Any]]) -> str:
    design_id = str(final_spec.get("design_id") or "")
    provenance = final_spec.get("provenance") or {}
    warning_codes = {str(warning.get("code") or "") for warning in report.get("warnings") or [] if isinstance(warning, dict)}
    if not design_id.startswith("design-board-"):
        findings.append(_finding("FINAL_SPEC_NOT_DESIGN_BOARD_DERIVED", "severe", "editable_template_spec.final.json is not design-board-derived."))
    if provenance.get("source") != "design_board_production_plan":
        findings.append(_finding("FINAL_SPEC_PROVENANCE_SOURCE_INVALID", "severe", "editable_template_spec.final.json provenance.source must be design_board_production_plan."))
    if provenance.get("production_plan_used") is not True:
        findings.append(_finding("FINAL_SPEC_PRODUCTION_PLAN_NOT_USED", "severe", "editable_template_spec.final.json must be built from the design production plan."))
    if provenance.get("extraction_source") != "actual":
        findings.append(_finding("FINAL_SPEC_EXTRACTION_SOURCE_INVALID", "severe", "Final spec provenance must record extraction_source=actual."))
    if provenance.get("codex_extraction_used") is not True:
        findings.append(_finding("FINAL_SPEC_CODEX_EXTRACTION_NOT_USED", "severe", "editable_template_spec.final.json provenance.codex_extraction_used must be true."))
    if provenance.get("prompt_id") != "creative_academic_template_board_v1":
        findings.append(_finding("FINAL_SPEC_PROMPT_ID_MISMATCH", "severe", "Final spec provenance prompt_id must match creative_academic_template_board_v1."))
    if any(code in warning_codes for code in {"MISSING_DESIGN_BOARD_EXTRACTION", "MISSING_EXTRACTION_ARTIFACT"}):
        findings.append(_finding("FINAL_SPEC_USED_SCHEMA_SAMPLE_EXTRACTION", "severe", "Final spec was built with schema sample extraction fallback, not real Codex design-board extraction."))
    if int(report.get("extracted_layout_count") or 0) < 18:
        findings.append(_finding("FINAL_SPEC_TOO_FEW_EXTRACTED_LAYOUTS", "severe", "Final spec does not contain at least 18 extracted layout archetypes."))
    ok = (
        design_id.startswith("design-board-")
        and provenance.get("source") == "design_board_production_plan"
        and provenance.get("codex_extraction_used") is True
        and provenance.get("production_plan_used") is True
        and provenance.get("extraction_source") == "actual"
        and provenance.get("prompt_id") == "creative_academic_template_board_v1"
        and not ({"MISSING_DESIGN_BOARD_EXTRACTION", "MISSING_EXTRACTION_ARTIFACT"} & warning_codes)
    )
    return "passed" if ok else "failed"


def _prompt_manifest_status(prompt_manifest: dict[str, Any], findings: list[dict[str, Any]]) -> str:
    prompts = prompt_manifest.get("prompts") if isinstance(prompt_manifest, dict) else None
    if not isinstance(prompts, list):
        return "failed"
    for entry in prompts:
        if not isinstance(entry, dict):
            continue
        if entry.get("prompt_id") == "creative_academic_template_board_v1":
            if entry.get("generation_mode") != "manual_codex":
                findings.append(_finding("PROMPT_MANIFEST_GENERATION_MODE_INVALID", "severe", "Canonical prompt manifest must use generation_mode=manual_codex."))
                return "failed"
            return "passed"
    findings.append(_finding("PROMPT_MANIFEST_CANONICAL_PROMPT_MISSING", "severe", "Canonical prompt manifest is missing creative_academic_template_board_v1."))
    return "failed"


def _tone_expression_status(assembly_plan: dict[str, Any], final_spec: dict[str, Any], findings: list[dict[str, Any]]) -> str:
    selected = assembly_plan.get("selected_tone_variant")
    binding_tones = {binding.get("selected_tone_variant") for binding in assembly_plan.get("slide_layout_bindings") or [] if isinstance(binding, dict)}
    tone_variants = ((final_spec.get("tokens") or {}).get("typography") or {}).get("tone_variants") or {}
    if not selected and not any(binding_tones):
        findings.append(_finding("NO_TONE_VARIANT_SELECTED", "severe", "Deck assembly plan does not select a tone variant."))
    if not tone_variants:
        findings.append(_finding("NO_TONE_VARIANTS_IN_SPEC", "severe", "Final template spec does not include academic/professional/creative tone variants."))
    return "passed" if (selected or any(binding_tones)) and tone_variants else "failed"


def _deck_scale_status(assembly_plan: dict[str, Any], pptx_stats: dict[str, Any], findings: list[dict[str, Any]]) -> str:
    scale = assembly_plan.get("deck_scale")
    slide_count = int(pptx_stats.get("slide_count") or 0)
    if not scale:
        findings.append(_finding("DECK_SCALE_MISSING", "severe", "Deck assembly plan does not record deck_scale."))
        return "failed"
    expected = "small" if slide_count <= 12 else "medium" if slide_count <= 30 else "large" if slide_count <= 80 else "very_large"
    if scale != expected:
        findings.append(_finding("DECK_SCALE_MISMATCH", "severe", "Deck scale does not match slide count.", details={"deck_scale": scale, "expected": expected, "slide_count": slide_count}))
        return "failed"
    return "passed"


def _reference_fidelity_status(visual_diff: dict[str, Any], run_label: str, findings: list[dict[str, Any]], warnings: list[dict[str, Any]]) -> str:
    if not visual_diff:
        warnings.append(_finding("TEMPLATE_VISUAL_DIFF_MISSING", "warning", "Template visual diff report is missing."))
        return "warning"
    if visual_diff.get("render_status") == "skipped" and run_label in {"premium", "design-accepted"}:
        findings.append(_finding("VISUAL_DIFF_SKIPPED_FOR_ACCEPTED_DESIGN", "severe", "Visual diff is skipped while the run is labeled premium/design-accepted."))
        return "failed"
    if int((visual_diff.get("findings_summary") or {}).get("severe") or 0) > 0:
        findings.append(_finding("TEMPLATE_VISUAL_DIFF_SEVERE", "severe", "Template visual diff contains severe violations."))
        return "failed"
    return "passed" if visual_diff.get("render_status") == "rendered" else "warning"


def _component_system_status(final_spec: dict[str, Any], findings: list[dict[str, Any]], warnings: list[dict[str, Any]]) -> str:
    required_components = {"dense_footer", "layered_card", "premium_kpi_card", "thin_grid_table", "diagonal_image_frame", "background_topology_grid"}
    components = [component for component in final_spec.get("components") or [] if isinstance(component, dict)]
    component_ids = {str(component.get("component_id") or component.get("id") or "") for component in components}
    missing = sorted(required_components - component_ids)
    if missing:
        findings.append(_finding("COMPONENT_SYSTEM_NOT_PREMIUM", "severe", "Final deck has no meaningful premium component system.", details={"missing_components": missing}))
        return "failed"
    provenance = final_spec.get("provenance") or {}
    if provenance.get("source") != "design_board_production_plan" or provenance.get("production_plan_used") is not True:
        findings.append(_finding("COMPONENT_SYSTEM_NOT_PRODUCTION_PLAN_DERIVED", "severe", "Premium component systems must be derived from the design production plan."))
        return "failed"
    translated_components = [
        component
        for component in components
        if ((component.get("default_tokens") or {}).get("source") == "component_translation_plan")
    ]
    if len(translated_components) < 12:
        findings.append(
            _finding(
                "COMPONENT_TRANSLATION_PLAN_TOO_THIN",
                "severe",
                "Final spec does not include enough component families translated from the production plan.",
                details={"translated_component_count": len(translated_components)},
            )
        )
        return "failed"
    layout_family_count = len([family for family in final_spec.get("layout_families") or [] if isinstance(family, dict)])
    if layout_family_count < 8:
        findings.append(
            _finding(
                "LAYOUT_FAMILY_PLAN_TOO_THIN",
                "severe",
                "Final spec does not include enough layout families from the production plan.",
                details={"layout_family_count": layout_family_count},
            )
        )
        return "failed"
    visual_fidelity_targets = final_spec.get("visual_fidelity_targets") or {}
    if not visual_fidelity_targets or len(visual_fidelity_targets) < 8:
        findings.append(_finding("VISUAL_FIDELITY_TARGETS_MISSING", "severe", "Final spec must carry machine-readable visual fidelity targets."))
        return "failed"
    generic_layouts = [
        str(layout.get("layout_id") or "")
        for layout in final_spec.get("layouts") or []
        if isinstance(layout, dict) and (str(layout.get("layout_id") or "").endswith("-mvp") or "standard-content-mvp" in str(layout.get("layout_id") or ""))
    ]
    if generic_layouts:
        findings.append(_finding("GENERIC_FALLBACK_LAYOUTS_IN_SPEC", "severe", "Premium final specs may not include generic MVP fallback layouts.", details={"layout_ids": generic_layouts[:8]}))
        return "failed"
    colors = ((final_spec.get("tokens") or {}).get("colors") or {})
    if _gray_card_palette(colors):
        findings.append(_finding("GENERIC_GRAY_CARD_LAYOUTS", "severe", "Template appears to rely on generic gray-card layouts."))
        return "failed"
    if len(colors) < 5:
        warnings.append(_finding("LOW_TOKEN_COLOR_COUNT", "warning", "Color token count is low for a premium design system."))
    return "passed"


def _layout_quality_status(assembly_plan: dict[str, Any], final_spec: dict[str, Any], findings: list[dict[str, Any]], warnings: list[dict[str, Any]]) -> str:
    layout_ids = {str(layout.get("layout_id") or "") for layout in final_spec.get("layouts") or [] if isinstance(layout, dict)}
    bindings = assembly_plan.get("slide_layout_bindings") or []
    selected = [str(binding.get("selected_layout_id") or binding.get("layout_id") or "") for binding in bindings if isinstance(binding, dict)]
    board_layouts = [layout_id for layout_id in selected if "layout-board-" in layout_id or layout_id in layout_ids and "board" in layout_id]
    if selected and not board_layouts:
        findings.append(_finding("DESIGN_BOARD_LAYOUTS_NOT_USED", "severe", "Deck assembly plan does not use design-board-derived layouts."))
        return "failed"
    generic_selected = [layout_id for layout_id in selected if layout_id.endswith("-mvp") or "standard-content-mvp" in layout_id]
    if generic_selected:
        findings.append(_finding("GENERIC_SMOKE_LAYOUT_USED", "severe", "Premium decks may not fall back to generic smoke/MVP layouts.", details={"layout_ids": sorted(set(generic_selected))}))
        return "failed"
    board_spec_layouts = [layout for layout in final_spec.get("layouts") or [] if isinstance(layout, dict) and str(layout.get("layout_id") or "").startswith("layout-board-")]
    missing_geometry_source = [layout.get("layout_id") for layout in board_spec_layouts if layout.get("extraction_geometry_source") != "extracted_slide_archetypes.slot_geometry"]
    if missing_geometry_source:
        findings.append(_finding("FINAL_SPEC_LAYOUTS_NOT_GEOMETRY_DERIVED", "severe", "Design-board layouts must be derived from extracted numeric slot geometry.", details={"layout_ids": missing_geometry_source[:8]}))
        return "failed"
    if len(set(selected)) <= 2 and len(selected) > 8:
        warnings.append(_finding("LOW_LAYOUT_VARIETY", "warning", "Deck uses very few layouts for its length."))
    return "passed"


def _functional_smoke_status(smoke_phrase_hits: list[str], assembly_plan: dict[str, Any]) -> str:
    if smoke_phrase_hits:
        return "smoke_artifact_detected"
    source = str(assembly_plan.get("source_slide_blueprint_path") or "")
    return "functional_artifact" if "large" in source else "not_detected"


def _editability_status(pptx_stats: dict[str, Any], spec_valid: bool) -> str:
    if not spec_valid:
        return "failed"
    if int(pptx_stats.get("text_shape_count") or 0) <= 0:
        return "failed"
    return "passed"


def _gray_card_palette(colors: dict[str, Any]) -> bool:
    if not colors:
        return True
    values = [str(value).lower() for value in colors.values() if isinstance(value, str)]
    gray_tokens = sum(1 for value in values if value in {"#ffffff", "#f8fafc", "#f1f5f9", "#e2e8f0", "#cbd5e1", "#64748b", "#475569", "#334155", "#1e293b", "#111827"})
    return len(values) >= 4 and gray_tokens / max(1, len(values)) > 0.72


def _validate_payload(payload: dict[str, Any], validator: Callable[[dict[str, Any]], dict[str, Any]], findings: list[dict[str, Any]], code: str) -> bool:
    if not payload:
        return False
    try:
        validator(payload)
    except Exception as exc:
        findings.append(_finding(code, "severe", f"Artifact validation failed: {type(exc).__name__}: {exc}"))
        return False
    return True


def _load_optional_json(path: Path, findings: list[dict[str, Any]], code: str, *, severity: str = "severe") -> dict[str, Any]:
    if not path.exists():
        findings.append(_finding(code, severity, f"Required artifact is missing: {path.as_posix()}"))
        return {}
    try:
        return _load_json(path)
    except (OSError, json.JSONDecodeError) as exc:
        findings.append(_finding(code, severity, f"Artifact could not be read: {type(exc).__name__}: {exc}", details={"path": path.as_posix()}))
        return {}


def _premium_statement(status: str) -> str:
    if status == "passed":
        return "This is a premium manual_codex design-board run, not a smoke fixture."
    return "This deck is not accepted as a premium design run; treat it as functional or incomplete until all severe findings are resolved."


def _finding(code: str, severity: str, message: str, *, details: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = {"code": code, "severity": severity, "message": message}
    if details:
        payload["details"] = details
    return payload


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Premium Design Quality Report",
        "",
        f"Status: `{report['status']}`",
        f"Premium design status: `{report['premium_design_status']}`",
        f"Run label: `{report['run_label']}`",
        f"Deck: `{report['pptx_path']}`",
        f"Functional smoke status: `{report['functional_smoke_status']}`",
        f"Editability status: `{report['editability_status']}`",
        f"Image policy status: `{report['image_policy_status']}`",
        f"Reference fidelity status: `{report['reference_fidelity_status']}`",
        f"Tone expression status: `{report['tone_expression_status']}`",
        f"Deck scale status: `{report['deck_scale_status']}`",
        f"Design board status: `{report['design_board_status']}`",
        f"Design board crop status: `{report['design_board_crop_status']}`",
        f"Codex extraction status: `{report['codex_extraction_status']}`",
        f"Spec derivation status: `{report['template_spec_derivation_status']}`",
        f"Component system status: `{report['component_system_status']}`",
        f"Layout quality status: `{report['layout_quality_status']}`",
        "",
        report["premium_design_run_statement"],
        "",
        "## Severe Findings",
        "",
    ]
    severe = [finding for finding in report["findings"] if finding["severity"] == "severe"]
    if severe:
        for finding in severe:
            lines.append(f"- `{finding['code']}`: {finding['message']}")
    else:
        lines.append("- None")
    lines.extend(["", "## Warnings", ""])
    if report["warnings"]:
        for warning in report["warnings"]:
            lines.append(f"- `{warning['code']}`: {warning['message']}")
    else:
        lines.append("- None")
    return "\n".join(lines) + "\n"


def _load_json(path: str | Path) -> Any:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _display_path(path: Path) -> str:
    return path.as_posix()


if __name__ == "__main__":
    raise SystemExit(main())
