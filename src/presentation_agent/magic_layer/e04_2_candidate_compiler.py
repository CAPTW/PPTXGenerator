"""Compile and render E04.2 patched deck."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation

from src.presentation_agent.magic_layer.e04_2_comparison_matrix_patch import patch_comparison_matrix_slide
from src.presentation_agent.magic_layer.e04_2_risk_register_patch import patch_risk_register_slide
from src.presentation_agent.magic_layer.e04_2_table_heavy_patch import patch_table_heavy_slide


def compile_e04_2_patched_deck(input_pptx: Path, output_pptx: Path) -> dict[str, Any]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(input_pptx, output_pptx)
    prs = Presentation(output_pptx)
    reports = {
        9: patch_comparison_matrix_slide(prs.slides[8]),
        11: patch_table_heavy_slide(prs.slides[10]),
        14: patch_risk_register_slide(prs.slides[13]),
    }
    prs.save(output_pptx)
    return {
        "schema_name": "e04_2_patch_application_report",
        "status": "passed" if all(report["status"] == "passed" for report in reports.values()) else "failed",
        "patched_deck_path": output_pptx.as_posix(),
        "slides_patched": [9, 11, 14],
        "patch_queue_items_resolved": 11,
        "updated_text_run_count": sum(report.get("updated_text_run_count", 0) for report in reports.values()),
        "slide_reports": reports,
    }


def render_e04_2_deck(pptx_path: Path, output_root: Path) -> dict[str, Any]:
    from src.presentation_agent.qa.render_pptx_preview import render_pptx_preview

    raw_dir = output_root / "renders" / "_raw"
    raw_dir.mkdir(parents=True, exist_ok=True)
    report = render_pptx_preview(pptx_path=pptx_path, output_dir=raw_dir, manifest_path=output_root / "renders" / "render_manifest.json", backend="auto", dpi=144)
    rendered = []
    for idx, row in enumerate(report.get("slides", []), start=1):
        source = Path(row.get("rendered_image_path") or "")
        if source.exists():
            target = output_root / "renders" / f"slide-{idx:03d}.png"
            shutil.copy2(source, target)
            row["rendered_image_path"] = target.as_posix()
            rendered.append(target)
    report["rendered_slide_count"] = len(rendered)
    report["expected_slide_count"] = 16
    report["rendered_paths"] = [path.as_posix() for path in rendered]
    return report

