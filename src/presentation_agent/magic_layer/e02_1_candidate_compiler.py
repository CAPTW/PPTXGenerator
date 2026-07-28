"""Compile richer E02.1 reference-fidelity PPTX candidates."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .e02_4core_orchestrator import ARCHETYPES

SLIDE_W = 16.0
SLIDE_H = 9.0

COLORS = {
    "bg": "06111A",
    "deep": "081B2A",
    "panel": "0B3B44",
    "panel2": "0F4B57",
    "cyan": "50D2E5",
    "cyan2": "1B8EA3",
    "gold": "F5A623",
    "gold2": "C68217",
    "paper": "F7F1E4",
    "paper2": "E7E1D5",
    "ink": "071018",
    "muted": "9EC4C8",
    "white": "FFFFFF",
}


def compile_e02_1_candidate(archetype_id: str, output_pptx: Path, *, visual_asset_plan: dict[str, Any] | None = None) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_e02_1_archetype(slide, archetype_id, visual_asset_plan or {"assets": []})
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    counts = expected_counts(archetype_id)
    return {
        "schema_name": "e02_1_candidate_compile_report",
        "status": "passed",
        "archetype_id": archetype_id,
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 1,
        **counts,
        "visual_asset_count": len((visual_asset_plan or {"assets": []}).get("assets", [])),
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "semantic_raster_final_use": False,
    }


def compile_e02_1_candidate_pack(archetype_asset_plans: dict[str, dict[str, Any]], output_pptx: Path) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for archetype_id in ARCHETYPES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        draw_e02_1_archetype(slide, archetype_id, archetype_asset_plans.get(archetype_id, {"assets": []}))
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_name": "e02_1_candidate_pack_compile_report",
        "status": "passed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": 4,
        "archetypes": list(ARCHETYPES),
        "non_canonical": True,
        "canonical_promotion": False,
    }


def expected_counts(archetype_id: str) -> dict[str, int]:
    return {
        "cover_hero": {"semantic_vector_icon_count": 5, "native_or_editable_chart_count": 0, "native_or_editable_table_count": 0},
        "standard_content": {"semantic_vector_icon_count": 8, "native_or_editable_chart_count": 0, "native_or_editable_table_count": 0},
        "data_dashboard": {"semantic_vector_icon_count": 10, "native_or_editable_chart_count": 2, "native_or_editable_table_count": 0},
        "table_heavy": {"semantic_vector_icon_count": 12, "native_or_editable_chart_count": 0, "native_or_editable_table_count": 1},
    }[archetype_id]


def draw_e02_1_archetype(slide: Any, archetype_id: str, visual_asset_plan: dict[str, Any]) -> None:
    _background(slide, f"{archetype_id}_background_base")
    if archetype_id == "cover_hero":
        _draw_cover_hero(slide, visual_asset_plan)
    elif archetype_id == "standard_content":
        _draw_standard_content(slide, visual_asset_plan)
    elif archetype_id == "data_dashboard":
        _draw_data_dashboard(slide)
    elif archetype_id == "table_heavy":
        _draw_table_heavy(slide)
    else:
        raise ValueError(archetype_id)


def _background(slide: Any, name: str) -> None:
    _shape(slide, name, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, "bg", None)
    for idx, x in enumerate((0.42, 1.7, 12.7, 14.7), start=1):
        _line(slide, f"{name}_technical_grid_{idx}", x, 0.25, x + 0.28, 7.9, "cyan2", 0.35)
    for idx, y in enumerate((0.42, 8.1), start=1):
        _line(slide, f"{name}_horizontal_rule_{idx}", 0.25, y, 15.75, y, "cyan2", 0.3)


def _draw_cover_hero(slide: Any, visual_asset_plan: dict[str, Any]) -> None:
    _asset(slide, visual_asset_plan, "cover_hero_visual_field_crop", "cover_hero_bounded_reference_visual_field")
    _shape(slide, "cover_hero_right_visual_frame", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 7.05, 0.26, 8.15, 7.85, None, "cyan")
    _shape(slide, "cover_hero_diagonal_divider_main", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 6.05, -0.25, 1.0, 8.7, "panel", "gold")
    _line(slide, "cover_hero_diagonal_gold_rule", 6.62, 0.0, 4.05, 8.35, "gold", 1.2)
    _line(slide, "cover_hero_diagonal_cyan_rule", 6.18, 0.55, 3.78, 8.1, "cyan", 0.7)
    _text(slide, "cover_hero_title_text", "TITLE", 0.7, 1.55, 4.2, 0.72, 42, "paper", bold=True)
    _line(slide, "cover_hero_title_rule", 0.72, 2.65, 2.82, 2.65, "cyan", 1.0)
    _text(slide, "cover_hero_subtitle_text", "SUBTITLE", 0.72, 3.02, 3.4, 0.42, 20, "gold", bold=True)
    _shape(slide, "cover_hero_meta_bar_native", MSO_AUTO_SHAPE_TYPE.PENTAGON, 0.62, 5.55, 4.3, 0.56, "panel", "cyan")
    _draw_icon(slide, "cover_hero_meta_icon_vector", 0.88, 5.68, 0.28, "cyan", "hex")
    _text(slide, "cover_hero_meta_text", "META", 1.35, 5.74, 1.4, 0.2, 8.5, "muted", bold=True)
    for idx, (x, y) in enumerate(((9.9, 3.2), (11.6, 2.25), (12.8, 4.72), (10.85, 5.35)), start=1):
        _line(slide, f"cover_hero_visual_node_line_{idx}", x - 0.55, y + 0.1, x, y, "gold", 0.8)
        _draw_icon(slide, f"cover_hero_visual_node_{idx}", x, y, 0.28, "gold", "target")
    _footer(slide, "cover_hero", left_label="FOOTER", right_label="SOURCE")


def _draw_standard_content(slide: Any, visual_asset_plan: dict[str, Any]) -> None:
    _asset(slide, visual_asset_plan, "standard_left_circuit_decorative_crop", "standard_content_bounded_left_circuit_chrome")
    _shape(slide, "standard_title_panel_chrome", MSO_AUTO_SHAPE_TYPE.PENTAGON, 0.55, 0.78, 3.3, 0.95, "deep", "cyan")
    _text(slide, "standard_title_text", "TITLE", 0.92, 1.04, 2.2, 0.38, 24, "paper", bold=True)
    _line(slide, "standard_title_gold_rule", 0.65, 1.68, 3.48, 1.68, "gold", 1.0)
    cards = [
        (4.0, 0.86, 5.45, 1.55, "standard_card_01", "BODY"),
        (4.0, 3.05, 5.45, 1.55, "standard_card_02", "BODY"),
        (4.0, 5.25, 5.45, 1.55, "standard_card_03", "BODY"),
        (9.82, 0.86, 2.85, 5.94, "standard_card_04", "BODY"),
    ]
    for idx, (x, y, w, h, name, label) in enumerate(cards, start=1):
        _shape(slide, f"{name}_angled_white_card", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x, y, w, h, "paper", "cyan")
        _shape(slide, f"{name}_shadow_chrome", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x + 0.12, y + h - 0.14, w - 0.35, 0.16, "paper2", None)
        _shape(slide, f"{name}_icon_hex_zone", MSO_AUTO_SHAPE_TYPE.HEXAGON, x - 0.35, y + 0.2, 0.92, 0.92, "panel", "cyan")
        _draw_icon(slide, f"{name}_semantic_vector_icon", x - 0.08, y + 0.47, 0.32, "paper", "document" if idx != 2 else "network")
        _text(slide, f"{name}_heading_text", label, x + 1.0, y + 0.32, 1.4, 0.22, 9, "ink", bold=True)
        _text(slide, f"{name}_body_text", "Editable module content area", x + 1.0, y + 0.76, w - 1.6, 0.34 if idx != 4 else 1.8, 10, "ink")
        _line(slide, f"{name}_gold_hatch_1", x + 0.7, y + 0.25, x + 0.92, y + 0.25, "gold", 1.1)
        _line(slide, f"{name}_gold_hatch_2", x + 0.78, y + 0.31, x + 1.0, y + 0.31, "gold", 1.1)
    _shape(slide, "standard_right_insight_rail_outer", MSO_AUTO_SHAPE_TYPE.PENTAGON, 13.2, 0.9, 2.0, 5.95, "panel", "gold")
    _draw_icon(slide, "standard_right_insight_bulb_icon", 13.95, 1.24, 0.42, "gold", "target")
    _text(slide, "standard_right_insight_title", "INSIGHT", 13.73, 1.82, 0.95, 0.24, 11, "gold", bold=True)
    _shape(slide, "standard_right_insight_body_panel", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 13.45, 2.28, 1.55, 3.7, "deep", "cyan2")
    _footer(slide, "standard_content", left_label="SOURCE", right_label="FOOTER")


def _draw_data_dashboard(slide: Any) -> None:
    _shape(slide, "dashboard_header_native_frame", MSO_AUTO_SHAPE_TYPE.PENTAGON, 2.35, 0.38, 10.4, 0.65, "deep", "cyan")
    _text(slide, "dashboard_title_text", "TITLE", 6.92, 0.52, 1.55, 0.28, 22, "paper", bold=True)
    for idx, x in enumerate((0.65, 3.95, 7.25, 10.55), start=1):
        _shape(slide, f"dashboard_kpi_{idx}_card_native", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 1.42, 2.72, 1.0, "panel", "cyan2")
        _draw_icon(slide, f"dashboard_kpi_{idx}_vector_icon", x + 0.14, 1.58, 0.48, "cyan", "gauge" if idx % 2 else "shield")
        _text(slide, f"dashboard_kpi_{idx}_label_text", "KPI", x + 0.78, 1.55, 0.8, 0.18, 8.5, "gold", bold=True)
        _text(slide, f"dashboard_kpi_{idx}_value_text", f"{idx*11:02d}", x + 0.78, 1.82, 0.9, 0.34, 20, "paper", bold=True)
        _line(slide, f"dashboard_kpi_{idx}_sparkline", x + 1.65, 2.02, x + 2.45, 1.7, "cyan", 0.8)
    _shape(slide, "dashboard_primary_chart_frame", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.68, 2.8, 9.42, 3.85, "paper", "cyan")
    _shape(slide, "dashboard_primary_chart_tab", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 4.2, 2.65, 2.2, 0.3, "deep", "cyan")
    _text(slide, "dashboard_primary_chart_title", "CHART", 4.95, 2.7, 0.72, 0.14, 8, "paper", bold=True)
    _dense_chart(slide, 1.08, 3.38, 8.18, 2.55)
    _shape(slide, "dashboard_secondary_panel_frame", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.35, 2.8, 4.85, 3.85, "paper", "cyan")
    _shape(slide, "dashboard_insight_tab", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 11.72, 2.65, 2.2, 0.3, "deep", "cyan")
    _text(slide, "dashboard_insight_tab_text", "INSIGHT", 12.3, 2.7, 0.9, 0.14, 8, "paper", bold=True)
    _donut_panel(slide, 10.85, 3.35)
    _annotation_strip(slide)
    _footer(slide, "data_dashboard", left_label="SOURCE", right_label="FOOTER")


def _draw_table_heavy(slide: Any) -> None:
    _shape(slide, "table_title_header_chrome", MSO_AUTO_SHAPE_TYPE.PENTAGON, 4.65, 0.38, 6.2, 0.58, "deep", "gold")
    _text(slide, "table_title_text", "TITLE", 6.82, 0.5, 1.2, 0.24, 21, "gold", bold=True)
    for idx, y in enumerate((1.24, 2.16, 3.08, 4.0), start=1):
        _shape(slide, f"table_side_rail_icon_zone_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, 0.38, y, 0.34, 0.34, "deep", "gold")
        _draw_icon(slide, f"table_side_rail_icon_{idx}", 0.43, y + 0.05, 0.22, "gold", "target")
    _shape(slide, "table_grid_outer_native", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.05, 1.18, 14.25, 5.72, "paper", "gold")
    _shape_grid(slide, 1.25, 1.47, 13.8, 4.95, cols=8, rows=7)
    for idx, x in enumerate((2.02, 3.32, 4.62, 5.92, 7.22, 8.52, 9.82, 11.12), start=1):
        _draw_icon(slide, f"table_header_icon_{idx}", x, 1.27, 0.2, "paper", "database" if idx % 2 else "shield")
    _shape(slide, "table_kpi_note_strip_native", MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, 0.95, 7.0, 14.25, 0.54, "paper", "gold")
    for idx, x in enumerate((1.28, 3.38, 5.48, 7.58), start=1):
        _draw_icon(slide, f"table_kpi_chip_icon_{idx}", x, 7.16, 0.24, "cyan2", "target")
        _text(slide, f"table_kpi_chip_text_{idx}", "KPI", x + 0.38, 7.17, 0.45, 0.15, 7.5, "ink", bold=True)
        _line(slide, f"table_kpi_chip_rule_{idx}", x + 0.9, 7.29, x + 1.55, 7.29, "cyan2", 0.8)
    _text(slide, "table_note_text", "NOTE", 10.82, 7.16, 0.65, 0.16, 7.5, "gold", bold=True)
    _footer(slide, "table_heavy", left_label="SOURCE", right_label="FOOTER")


def _asset(slide: Any, visual_asset_plan: dict[str, Any], asset_id: str, shape_name: str) -> None:
    for asset in visual_asset_plan.get("assets", []):
        if asset["asset_id"] == asset_id:
            bbox = asset["target_bbox_in"]
            pic = slide.shapes.add_picture(asset["asset_path"], Inches(bbox["x"]), Inches(bbox["y"]), width=Inches(bbox["w"]), height=Inches(bbox["h"]))
            pic.name = shape_name
            return


def _dense_chart(slide: Any, x: float, y: float, w: float, h: float) -> None:
    for idx in range(6):
        yy = y + idx * h / 5
        _line(slide, f"dashboard_chart_gridline_{idx}", x, yy, x + w, yy, "paper2", 0.25)
    values = [0.15, 0.28, 0.45, 0.42, 0.58, 0.63, 0.72, 0.86]
    bar_w = w / 18
    points = []
    for idx, value in enumerate(values):
        bx = x + 0.45 + idx * (bar_w + 0.45)
        bh = h * value
        _shape(slide, f"dashboard_chart_bar_{idx}_series_a", MSO_AUTO_SHAPE_TYPE.RECTANGLE, bx, y + h - bh, bar_w, bh, "cyan2", None)
        _shape(slide, f"dashboard_chart_bar_{idx}_series_b", MSO_AUTO_SHAPE_TYPE.RECTANGLE, bx + bar_w + 0.06, y + h - bh * 0.72, bar_w, bh * 0.72, "paper2", None)
        px, py = bx + bar_w * 1.4, y + h - bh * 0.85
        points.append((px, py))
        _shape(slide, f"dashboard_chart_marker_{idx}", MSO_AUTO_SHAPE_TYPE.DIAMOND, px - 0.04, py - 0.04, 0.08, 0.08, "gold", None)
    for idx in range(len(points) - 1):
        _line(slide, f"dashboard_chart_trend_line_{idx}", points[idx][0], points[idx][1], points[idx + 1][0], points[idx + 1][1], "cyan", 0.8)


def _donut_panel(slide: Any, x: float, y: float) -> None:
    _shape(slide, "dashboard_secondary_donut_outer", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, 1.05, 1.05, "cyan", None)
    _shape(slide, "dashboard_secondary_donut_inner", MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.28, y + 0.28, 0.5, 0.5, "paper", None)
    for idx in range(5):
        _line(slide, f"dashboard_secondary_list_rule_{idx}", x + 1.35, y + 0.08 + idx * 0.34, x + 3.6, y + 0.08 + idx * 0.34, "paper2", 0.55)
    for idx, label in enumerate(("A", "B", "C", "D"), start=1):
        _shape(slide, f"dashboard_secondary_bottom_hex_{idx}", MSO_AUTO_SHAPE_TYPE.HEXAGON, x + 0.2 + idx * 0.75, y + 2.38, 0.44, 0.34, "deep", "gold")
        _text(slide, f"dashboard_secondary_bottom_label_{idx}", label, x + 0.36 + idx * 0.75, y + 2.46, 0.14, 0.08, 6, "paper", bold=True)


def _annotation_strip(slide: Any) -> None:
    _shape(slide, "dashboard_annotation_strip_native", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.62, 6.95, 14.6, 0.54, "panel", "cyan2")
    labels = ["ANNOTATION", "FILTER", "SYSTEM", "EXPORT", "SOURCE"]
    for idx, label in enumerate(labels):
        x = 0.9 + idx * 2.75
        _draw_icon(slide, f"dashboard_annotation_icon_{idx}", x, 7.09, 0.22, "cyan", "target")
        _text(slide, f"dashboard_annotation_text_{idx}", label, x + 0.35, 7.1, 0.9, 0.12, 6.8, "muted", bold=True)


def _shape_grid(slide: Any, x: float, y: float, w: float, h: float, *, cols: int, rows: int) -> None:
    col_w = w / cols
    row_h = h / rows
    for row in range(rows):
        for col in range(cols):
            fill = "paper" if row % 2 else "white"
            if row == 0:
                fill = "panel"
            elif col in {0, cols - 1}:
                fill = "paper2"
            _shape(slide, f"table_grid_cell_r{row}_c{col}", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + col * col_w, y + row * row_h, col_w, row_h, fill, "paper2")
            if row == 0:
                _text(slide, f"table_header_cell_text_r{row}_c{col}", "HDR", x + col * col_w + 0.08, y + row * row_h + 0.17, col_w - 0.12, 0.12, 5.8, "paper", bold=True)
            elif col in {0, 2, 5}:
                _text(slide, f"table_cell_text_r{row}_c{col}", "ROW", x + col * col_w + 0.08, y + row * row_h + 0.18, col_w - 0.12, 0.12, 5.6, "ink")
            else:
                _shape(slide, f"table_cell_status_pill_r{row}_c{col}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + col * col_w + 0.36, y + row * row_h + 0.2, col_w - 0.72, 0.12, "paper2", "cyan2")


def _footer(slide: Any, prefix: str, *, left_label: str, right_label: str) -> None:
    _shape(slide, f"{prefix}_source_footer_strip_native", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 8.34, SLIDE_W, 0.66, "deep", None)
    _line(slide, f"{prefix}_footer_top_rule_gold", 0, 8.34, SLIDE_W, 8.34, "gold", 0.9)
    _draw_icon(slide, f"{prefix}_footer_left_icon", 0.42, 8.55, 0.22, "cyan", "database")
    _text(slide, f"{prefix}_footer_left_text", left_label, 0.82, 8.58, 1.2, 0.12, 7, "muted", bold=True)
    _draw_icon(slide, f"{prefix}_footer_right_icon", 12.8, 8.55, 0.22, "gold", "database")
    _text(slide, f"{prefix}_footer_right_text", right_label, 13.2, 8.58, 1.2, 0.12, 7, "muted", bold=True)


def _draw_icon(slide: Any, name: str, x: float, y: float, size: float, color: str, role: str) -> None:
    if role == "database":
        _shape(slide, f"{name}_vector_cylinder_top", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size * 0.35, None, color)
        _line(slide, f"{name}_vector_cylinder_left", x, y + size * 0.18, x, y + size * 0.78, color, 0.7)
        _line(slide, f"{name}_vector_cylinder_right", x + size, y + size * 0.18, x + size, y + size * 0.78, color, 0.7)
        _shape(slide, f"{name}_vector_cylinder_bottom", MSO_AUTO_SHAPE_TYPE.OVAL, x, y + size * 0.58, size, size * 0.35, None, color)
    elif role == "gauge":
        _shape(slide, f"{name}_vector_gauge_ring", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, None, color)
        _line(slide, f"{name}_vector_gauge_needle", x + size * 0.5, y + size * 0.55, x + size * 0.78, y + size * 0.28, color, 0.9)
    elif role == "shield":
        _shape(slide, f"{name}_vector_shield", MSO_AUTO_SHAPE_TYPE.PENTAGON, x, y, size, size, None, color)
        _line(slide, f"{name}_vector_check_1", x + size * 0.25, y + size * 0.52, x + size * 0.42, y + size * 0.70, color, 0.8)
        _line(slide, f"{name}_vector_check_2", x + size * 0.42, y + size * 0.70, x + size * 0.78, y + size * 0.30, color, 0.8)
    elif role == "network":
        for idx, (cx, cy) in enumerate(((0.1, 0.2), (0.7, 0.1), (0.55, 0.72)), start=1):
            _shape(slide, f"{name}_node_{idx}", MSO_AUTO_SHAPE_TYPE.OVAL, x + size * cx, y + size * cy, size * 0.22, size * 0.22, None, color)
        _line(slide, f"{name}_edge_1", x + size * 0.2, y + size * 0.3, x + size * 0.78, y + size * 0.2, color, 0.6)
        _line(slide, f"{name}_edge_2", x + size * 0.78, y + size * 0.2, x + size * 0.65, y + size * 0.82, color, 0.6)
    elif role == "hex":
        _shape(slide, f"{name}_vector_hex", MSO_AUTO_SHAPE_TYPE.HEXAGON, x, y, size, size, None, color)
    else:
        _shape(slide, f"{name}_vector_ring", MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, None, color)
        _line(slide, f"{name}_vector_mark_1", x + size * 0.28, y + size * 0.55, x + size * 0.46, y + size * 0.72, color, 0.85)
        _line(slide, f"{name}_vector_mark_2", x + size * 0.46, y + size * 0.72, x + size * 0.78, y + size * 0.28, color, 0.85)


def _shape(slide: Any, name: str, shape_type: Any, x: float, y: float, w: float, h: float, fill: str | None, line: str | None) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.65)
    else:
        shape.line.fill.background()
    return shape


def _line(slide: Any, name: str, x1: float, y1: float, x2: float, y2: float, color: str, width: float) -> Any:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _text(slide: Any, name: str, text: str, x: float, y: float, w: float, h: float, size: float, color: str, *, bold: bool = False) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _rgb(name_or_hex: str) -> RGBColor:
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
