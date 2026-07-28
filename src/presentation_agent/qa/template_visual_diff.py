"""Approximate visual diff between template references and rendered preview slides."""

from __future__ import annotations

import argparse
import hashlib
import json
import math
from pathlib import Path
from typing import Any

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .render_pptx_preview import DEFAULT_MANIFEST_PATH as DEFAULT_RENDER_REPORT_PATH
from .render_pptx_preview import render_pptx_preview


DEFAULT_TEMPLATE_IMAGE_MANIFEST = Path("outputs/template_images/template_image_manifest.json")
DEFAULT_TEMPLATE_SPEC = Path("outputs/editable_template_spec.final.json")
DEFAULT_TEMPLATE_PREVIEW_PPTX = Path("outputs/template_preview.pptx")
DEFAULT_PREVIEW_MANIFEST = Path("outputs/template_preview_manifest.json")
DEFAULT_RENDER_DIR = Path("outputs/template_preview_png")
DEFAULT_JSON_REPORT = Path("outputs/template_visual_diff_report.json")
DEFAULT_MD_REPORT = Path("outputs/template_visual_diff_report.md")
EMU_PER_INCH = 914400
FULL_SLIDE_AREA_RATIO_THRESHOLD = 0.88


def build_template_visual_diff_report(
    *,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    template_preview_pptx_path: str | Path = DEFAULT_TEMPLATE_PREVIEW_PPTX,
    preview_manifest_path: str | Path = DEFAULT_PREVIEW_MANIFEST,
    render_dir: str | Path = DEFAULT_RENDER_DIR,
    render_report_path: str | Path = DEFAULT_RENDER_REPORT_PATH,
) -> dict[str, Any]:
    template_image_manifest = _load_json(template_image_manifest_path)
    template_spec = _load_json(template_spec_path)
    preview_manifest = _load_json(preview_manifest_path) if Path(preview_manifest_path).exists() else {}
    render_report = _load_json(render_report_path) if Path(render_report_path).exists() else {"render_status": "not_provided"}
    render_status = _render_status(render_report)
    render_root = Path(render_dir)

    reference_by_archetype = _reference_images_by_archetype(template_image_manifest)
    picture_findings_by_slide = _template_preview_picture_findings(template_preview_pptx_path, template_image_manifest)
    compiled_layouts = _compiled_layouts(template_spec, preview_manifest)
    slides: list[dict[str, Any]] = []
    findings: list[dict[str, Any]] = []

    for fallback_index, compiled in enumerate(compiled_layouts, start=1):
        slide_index = int(compiled.get("slide_number") or fallback_index)
        archetype_id = str(compiled.get("archetype_id") or "")
        layout_id = str(compiled.get("layout_id") or "")
        rendered_path = _rendered_path(render_root, render_report, slide_index)
        reference_path = reference_by_archetype.get(archetype_id)
        layout = _layout_by_id(template_spec).get(layout_id, {})
        slide_report = _slide_visual_report(
            slide_index=slide_index,
            layout_id=layout_id,
            archetype_id=archetype_id,
            rendered_path=rendered_path,
            reference_path=reference_path,
            layout=layout,
            render_status=render_status,
            pptx_picture_findings=picture_findings_by_slide.get(slide_index, []),
        )
        slides.append(slide_report)
        findings.extend(slide_report["warnings"])
        findings.extend(slide_report["severe_violations"])

    if render_status == "skipped":
        findings.append(
            _finding(
                "RENDER_SKIPPED",
                "warning",
                "Template preview rendering was skipped because no local renderer was available.",
            )
        )

    severe_count = sum(1 for finding in findings if finding["severity"] == "severe")
    warning_count = sum(1 for finding in findings if finding["severity"] == "warning")
    report = {
        "schema_name": "template_visual_diff_report",
        "schema_version": "1.0",
        "template_image_manifest_path": _display_path(Path(template_image_manifest_path)),
        "template_spec_path": _display_path(Path(template_spec_path)),
        "template_preview_pptx_path": _display_path(Path(template_preview_pptx_path)),
        "render_dir": _display_path(render_root),
        "render_status": render_status,
        "render_backend": render_report.get("backend") or (render_report.get("renderer") or {}).get("name"),
        "status": "failed" if severe_count else "issues_reported" if warning_count else "passed",
        "qa_blocks_deck_generation": bool(severe_count),
        "slide_count": len(slides),
        "findings_summary": {
            "total": len(findings),
            "severe": severe_count,
            "warning": warning_count,
            "info": sum(1 for finding in findings if finding["severity"] == "info"),
        },
        "slides": slides,
        "findings": findings,
    }
    return report


def build_template_visual_diff_report_from_files(
    *,
    template_image_manifest_path: str | Path = DEFAULT_TEMPLATE_IMAGE_MANIFEST,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    template_preview_pptx_path: str | Path = DEFAULT_TEMPLATE_PREVIEW_PPTX,
    preview_manifest_path: str | Path = DEFAULT_PREVIEW_MANIFEST,
    render_dir: str | Path = DEFAULT_RENDER_DIR,
    render_report_path: str | Path = DEFAULT_RENDER_REPORT_PATH,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
    renderer: str = "auto",
) -> Path:
    render_pptx_preview(
        pptx_path=template_preview_pptx_path,
        output_dir=render_dir,
        report_path=render_report_path,
        backend=renderer,
    )
    report = build_template_visual_diff_report(
        template_image_manifest_path=template_image_manifest_path,
        template_spec_path=template_spec_path,
        template_preview_pptx_path=template_preview_pptx_path,
        preview_manifest_path=preview_manifest_path,
        render_dir=render_dir,
        render_report_path=render_report_path,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return json_path


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Run approximate visual diff between template reference PNGs and rendered template preview slides.")
    parser.add_argument("--template-image-manifest", type=Path, default=DEFAULT_TEMPLATE_IMAGE_MANIFEST)
    parser.add_argument("--template-spec", type=Path, default=DEFAULT_TEMPLATE_SPEC)
    parser.add_argument("--pptx", type=Path, default=DEFAULT_TEMPLATE_PREVIEW_PPTX)
    parser.add_argument("--preview-manifest", type=Path, default=DEFAULT_PREVIEW_MANIFEST)
    parser.add_argument("--render-dir", type=Path, default=DEFAULT_RENDER_DIR)
    parser.add_argument("--render-report", type=Path, default=DEFAULT_RENDER_REPORT_PATH)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    parser.add_argument("--renderer", choices=("auto", "powerpoint_com", "libreoffice", "none"), default="auto")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        output = build_template_visual_diff_report_from_files(
            template_image_manifest_path=args.template_image_manifest,
            template_spec_path=args.template_spec,
            template_preview_pptx_path=args.pptx,
            preview_manifest_path=args.preview_manifest,
            render_dir=args.render_dir,
            render_report_path=args.render_report,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
            renderer=args.renderer,
        )
        report = _load_json(output)
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        print(f"TEMPLATE_VISUAL_DIFF_FAILED {exc}")
        return 1
    print(f"WROTE {output}")
    if report.get("qa_blocks_deck_generation"):
        print("TEMPLATE_VISUAL_DIFF failed")
        return 1
    print(f"TEMPLATE_VISUAL_DIFF {report['status']}")
    return 0


def _slide_visual_report(
    *,
    slide_index: int,
    layout_id: str,
    archetype_id: str,
    rendered_path: Path,
    reference_path: Path | None,
    layout: dict[str, Any],
    render_status: str,
    pptx_picture_findings: list[dict[str, Any]],
) -> dict[str, Any]:
    warnings: list[dict[str, Any]] = []
    severe_violations: list[dict[str, Any]] = list(pptx_picture_findings)
    reference_metrics = _image_metrics(reference_path) if reference_path and reference_path.exists() else None
    rendered_metrics = _image_metrics(rendered_path) if rendered_path.exists() else None

    if reference_path is None:
        warnings.append(_finding("REFERENCE_IMAGE_MISSING", "warning", "No template reference image was available for this archetype.", slide_index))
    if rendered_metrics is None:
        if render_status != "skipped":
            warnings.append(_finding("RENDERED_PREVIEW_IMAGE_MISSING", "warning", "Rendered preview PNG is missing for this slide.", slide_index))
        return {
            "slide_index": slide_index,
            "layout_id": layout_id,
            "archetype_id": archetype_id,
            "rendered_image_path": _display_path(rendered_path),
            "reference_image_path": _display_path(reference_path) if reference_path else None,
            "reference_match": {"archetype_id": archetype_id, "layout_id": layout_id},
            "metrics": None,
            "similarity_summary": {"status": "rendered_preview_missing", "label": "No rendered preview image was available."},
            "warnings": warnings,
            "severe_violations": severe_violations,
        }

    if reference_metrics is None:
        metrics = _single_image_metrics(rendered_metrics)
        warnings.extend(_single_image_warnings(metrics, slide_index, layout, archetype_id))
    else:
        metrics = _comparison_metrics(rendered_metrics, reference_metrics, layout)
        classification = _classify_metrics(metrics, slide_index, layout, archetype_id)
        warnings.extend(classification["warnings"])
        severe_violations.extend(classification["severe_violations"])

    return {
        "slide_index": slide_index,
        "layout_id": layout_id,
        "archetype_id": archetype_id,
        "rendered_image_path": _display_path(rendered_path),
        "reference_image_path": _display_path(reference_path) if reference_path else None,
        "reference_match": {"archetype_id": archetype_id, "layout_id": layout_id},
        "metrics": metrics,
        "similarity_summary": _similarity_summary(metrics, severe_violations, warnings),
        "warnings": warnings,
        "severe_violations": severe_violations,
    }


def _comparison_metrics(rendered: dict[str, Any], reference: dict[str, Any], layout: dict[str, Any]) -> dict[str, Any]:
    regions = _regions_for_layout(layout)
    rendered_regions = {name: _region_occupancy(rendered["image"], bounds) for name, bounds in regions.items()}
    reference_regions = {name: _region_occupancy(reference["image"], bounds) for name, bounds in regions.items()}
    return {
        "canvas_size_match": _canvas_size_match(rendered, reference),
        "rendered_width_px": rendered["width_px"],
        "rendered_height_px": rendered["height_px"],
        "reference_width_px": reference["width_px"],
        "reference_height_px": reference["height_px"],
        "dominant_palette_distance": _palette_distance(rendered["dominant_rgb"], reference["dominant_rgb"]),
        "dominant_palette_similarity": round(1.0 - _palette_distance(rendered["dominant_rgb"], reference["dominant_rgb"]), 6),
        "mean_palette_distance": _palette_distance(rendered["mean_rgb"], reference["mean_rgb"]),
        "mean_palette_similarity": round(1.0 - _palette_distance(rendered["mean_rgb"], reference["mean_rgb"]), 6),
        "edge_density_difference": round(abs(rendered["edge_density"] - reference["edge_density"]), 6),
        "edge_density_delta": round(rendered["edge_density"] - reference["edge_density"], 6),
        "blank_area_difference": round(abs(rendered["blank_area_ratio"] - reference["blank_area_ratio"]), 6),
        "blank_area_delta": round(rendered["blank_area_ratio"] - reference["blank_area_ratio"], 6),
        "dark_area_difference": round(abs(rendered["dark_area_ratio"] - reference["dark_area_ratio"]), 6),
        "dark_area_delta": round(rendered["dark_area_ratio"] - reference["dark_area_ratio"], 6),
        "light_area_difference": round(abs(rendered["light_area_ratio"] - reference["light_area_ratio"]), 6),
        "light_area_delta": round(rendered["light_area_ratio"] - reference["light_area_ratio"], 6),
        "footer_area_occupancy": rendered_regions["footer"],
        "reference_footer_area_occupancy": reference_regions["footer"],
        "title_area_occupancy": rendered_regions["title"],
        "reference_title_area_occupancy": reference_regions["title"],
        "card_grid_region_occupancy": rendered_regions["card_grid"],
        "reference_card_grid_region_occupancy": reference_regions["card_grid"],
        "table_chart_frame_occupancy": rendered_regions["table_chart"],
        "reference_table_chart_frame_occupancy": reference_regions["table_chart"],
        "diagonal_region_occupancy": rendered_regions["diagonal"],
        "reference_diagonal_region_occupancy": reference_regions["diagonal"],
        "rendered_blank_area_ratio": rendered["blank_area_ratio"],
        "reference_blank_area_ratio": reference["blank_area_ratio"],
        "rendered_edge_density": rendered["edge_density"],
        "reference_edge_density": reference["edge_density"],
    }


def _single_image_metrics(rendered: dict[str, Any]) -> dict[str, Any]:
    return {
        "canvas_size_match": True,
        "rendered_width_px": rendered["width_px"],
        "rendered_height_px": rendered["height_px"],
        "dominant_palette_distance": None,
        "dominant_palette_similarity": None,
        "mean_palette_distance": None,
        "mean_palette_similarity": None,
        "edge_density_difference": None,
        "edge_density_delta": None,
        "blank_area_difference": None,
        "blank_area_delta": None,
        "dark_area_difference": None,
        "dark_area_delta": None,
        "light_area_difference": None,
        "light_area_delta": None,
        "footer_area_occupancy": _region_occupancy(rendered["image"], _normalized_region("footer")),
        "title_area_occupancy": _region_occupancy(rendered["image"], _normalized_region("title")),
        "card_grid_region_occupancy": _region_occupancy(rendered["image"], _normalized_region("card_grid")),
        "table_chart_frame_occupancy": _region_occupancy(rendered["image"], _normalized_region("table_chart")),
        "diagonal_region_occupancy": _region_occupancy(rendered["image"], _normalized_region("diagonal")),
        "rendered_blank_area_ratio": rendered["blank_area_ratio"],
        "rendered_edge_density": rendered["edge_density"],
    }


def _classify_metrics(metrics: dict[str, Any], slide_index: int, layout: dict[str, Any], archetype_id: str) -> dict[str, list[dict[str, Any]]]:
    warnings: list[dict[str, Any]] = []
    severe_violations: list[dict[str, Any]] = []
    if not metrics["canvas_size_match"]:
        warnings.append(_finding("CANVAS_SIZE_MISMATCH", "warning", "Rendered preview and reference image dimensions/aspect differ.", slide_index))

    if metrics["rendered_blank_area_ratio"] > 0.985 and metrics["rendered_edge_density"] < 0.004:
        severe_violations.append(_finding("RENDERED_SLIDE_MOSTLY_BLANK", "severe", "Rendered slide is mostly blank.", slide_index))
    elif metrics["rendered_blank_area_ratio"] > min(0.92, metrics["reference_blank_area_ratio"] + 0.22):
        warnings.append(_finding("VISUAL_DENSITY_BELOW_TARGET", "warning", "Rendered preview appears substantially sparser than the reference.", slide_index))

    if metrics["edge_density_delta"] is not None and metrics["edge_density_delta"] < -0.08:
        warnings.append(_finding("LOWER_EDGE_DENSITY_THAN_REFERENCE", "warning", "Rendered preview has lower edge density than the reference.", slide_index))
    elif metrics["edge_density_delta"] is not None and metrics["edge_density_delta"] > 0.16:
        warnings.append(_finding("TEMPLATE_TOO_DENSE", "warning", "Rendered preview has much higher edge density than the reference.", slide_index))

    if _requires_title(layout, archetype_id) and metrics["title_area_occupancy"] < 0.002:
        severe_violations.append(_finding("MISSING_TITLE_ZONE", "severe", "Rendered preview is missing the expected title zone.", slide_index))
    elif _requires_title(layout, archetype_id) and metrics["title_area_occupancy"] < max(0.01, metrics["reference_title_area_occupancy"] * 0.35):
        warnings.append(_finding("MISSING_TITLE_AREA", "warning", "Rendered preview has little visible occupancy in the expected title area.", slide_index))

    if _requires_footer(layout) and metrics["footer_area_occupancy"] < 0.002:
        severe_violations.append(_finding("MISSING_FOOTER", "severe", "Rendered preview is missing a required footer zone.", slide_index))
    elif _requires_footer(layout) and metrics["footer_area_occupancy"] < max(0.006, metrics["reference_footer_area_occupancy"] * 0.35):
        warnings.append(_finding("FOOTER_TOO_EMPTY", "warning", "Rendered preview footer region appears too empty compared with the reference.", slide_index))
    if metrics["footer_area_occupancy"] > 0.82:
        warnings.append(_finding("FOOTER_TOO_THICK", "warning", "Footer region appears too visually heavy.", slide_index))

    if metrics["title_area_occupancy"] > 0.82:
        warnings.append(_finding("TITLE_AREA_OVERSIZED", "warning", "Title region appears visually oversized.", slide_index))
    if _has_slot(layout, {"cards", "metric_panels"}) and metrics["card_grid_region_occupancy"] < 0.02:
        warnings.append(_finding("CARD_REGION_TOO_SPARSE", "warning", "Card or KPI region appears too sparse.", slide_index))

    if (metrics["dominant_palette_distance"] or 0) > 0.72 and (metrics["mean_palette_distance"] or 0) > 0.6:
        severe_violations.append(_finding("PALETTE_COMPLETELY_MISMATCHED", "severe", "Rendered preview palette is completely mismatched against the reference.", slide_index))
    elif (metrics["dominant_palette_distance"] or 0) > 0.32 or (metrics["mean_palette_distance"] or 0) > 0.24:
        warnings.append(_finding("PALETTE_MISMATCH", "warning", "Rendered preview palette differs substantially from the reference image.", slide_index))

    mismatch_score = sum(
        [
            1 if (metrics["dominant_palette_distance"] or 0) > 0.25 else 0,
            1 if (metrics["blank_area_difference"] or 0) > 0.28 else 0,
            1 if (metrics["edge_density_difference"] or 0) > 0.1 else 0,
            1 if metrics["title_area_occupancy"] < 0.02 else 0,
            1 if metrics["footer_area_occupancy"] < 0.005 else 0,
        ]
    )
    if mismatch_score >= 3:
        warnings.append(_finding("REFERENCE_IMAGE_LIKELY_NOT_REFLECTED", "warning", "Approximate metrics suggest the rendered preview does not reflect the reference image.", slide_index))
    return {"warnings": warnings, "severe_violations": severe_violations}


def _single_image_warnings(metrics: dict[str, Any], slide_index: int, layout: dict[str, Any], archetype_id: str) -> list[dict[str, Any]]:
    warnings: list[dict[str, Any]] = []
    if metrics["rendered_blank_area_ratio"] > 0.92:
        warnings.append(_finding("VISUAL_DENSITY_BELOW_TARGET", "warning", "Rendered preview appears sparse without a reference image.", slide_index))
    if _requires_title(layout, archetype_id) and metrics["title_area_occupancy"] < 0.01:
        warnings.append(_finding("MISSING_TITLE_AREA", "warning", "Rendered preview has little visible occupancy in the expected title area.", slide_index))
    if _requires_footer(layout) and metrics["footer_area_occupancy"] < 0.006:
        warnings.append(_finding("FOOTER_TOO_EMPTY", "warning", "Rendered preview footer region appears too empty.", slide_index))
    return warnings


def _image_metrics(path: Path | None) -> dict[str, Any] | None:
    if path is None or not path.exists():
        return None
    with Image.open(path) as source:
        image = source.convert("RGB")
        sample = image.resize((160, 90))
        stat = ImageStat.Stat(sample)
        grayscale = sample.convert("L")
        histogram = grayscale.histogram()
        width, height = image.size
        total = max(1, sample.size[0] * sample.size[1])
        light_ratio = sum(histogram[230:]) / total
        dark_ratio = sum(histogram[:35]) / total
        blank_ratio = max(light_ratio, dark_ratio)
        return {
            "image": sample,
            "width_px": width,
            "height_px": height,
            "mean_rgb": tuple(float(value) for value in stat.mean),
            "dominant_rgb": _dominant_rgb(sample),
            "edge_density": _edge_density(sample),
            "blank_area_ratio": round(blank_ratio, 6),
            "dark_area_ratio": round(dark_ratio, 6),
            "light_area_ratio": round(light_ratio, 6),
        }


def _dominant_rgb(image: Image.Image) -> tuple[float, float, float]:
    colors = image.quantize(colors=6, method=Image.Quantize.MEDIANCUT).convert("RGB").getcolors(maxcolors=1_000_000) or []
    if not colors:
        stat = ImageStat.Stat(image)
        return tuple(float(value) for value in stat.mean)
    _, rgb = max(colors, key=lambda item: item[0])
    return tuple(float(value) for value in rgb)


def _edge_density(image: Image.Image) -> float:
    gray = image.convert("L")
    pixels = gray.load()
    width, height = gray.size
    edges = 0
    total = 0
    for y in range(height - 1):
        for x in range(width - 1):
            current = pixels[x, y]
            if abs(current - pixels[x + 1, y]) > 24 or abs(current - pixels[x, y + 1]) > 24:
                edges += 1
            total += 1
    return round(edges / max(1, total), 6)


def _region_occupancy(image: Image.Image, bounds: dict[str, float]) -> float:
    width, height = image.size
    left = max(0, min(width - 1, int(bounds["x"] * width)))
    top = max(0, min(height - 1, int(bounds["y"] * height)))
    right = max(left + 1, min(width, int((bounds["x"] + bounds["w"]) * width)))
    bottom = max(top + 1, min(height, int((bounds["y"] + bounds["h"]) * height)))
    crop = image.crop((left, top, right, bottom)).convert("RGB")
    get_pixels = getattr(crop, "get_flattened_data", crop.getdata)
    pixels = list(get_pixels())
    if not pixels:
        return 0.0
    background = _dominant_rgb(crop)
    occupied = sum(1 for pixel in pixels if _rgb_distance(pixel, background) > 28)
    return round(occupied / len(pixels), 6)


def _regions_for_layout(layout: dict[str, Any]) -> dict[str, dict[str, float]]:
    regions = {
        "title": _normalized_region("title"),
        "footer": _normalized_region("footer"),
        "card_grid": _normalized_region("card_grid"),
        "table_chart": _normalized_region("table_chart"),
        "diagonal": _normalized_region("diagonal"),
    }
    for slot in layout.get("slots") or []:
        if not isinstance(slot, dict) or not isinstance(slot.get("bounds"), dict):
            continue
        slot_id = str(slot.get("slot_id") or "")
        bounds = _normalize_bounds(slot["bounds"])
        if slot_id == "title":
            regions["title"] = bounds
        elif slot_id == "footer":
            regions["footer"] = bounds
        elif slot_id in {"cards", "metric_panels", "supporting_panel"}:
            regions["card_grid"] = bounds
        elif str(slot.get("slot_type") or "") in {"table", "chart"} or slot_id in {"table", "matrix", "primary_chart", "secondary_chart"}:
            regions["table_chart"] = bounds
        elif str(slot.get("component_id") or "") == "diagonal_photo_panel" or slot_id in {"hero_image", "photo_frame"}:
            regions["diagonal"] = bounds
    return regions


def _normalize_bounds(bounds: dict[str, Any]) -> dict[str, float]:
    return {
        "x": float(bounds["x"]) / 13.333,
        "y": float(bounds["y"]) / 7.5,
        "w": float(bounds["w"]) / 13.333,
        "h": float(bounds["h"]) / 7.5,
    }


def _normalized_region(name: str) -> dict[str, float]:
    if name == "title":
        return {"x": 0.04, "y": 0.04, "w": 0.64, "h": 0.16}
    if name == "footer":
        return {"x": 0.04, "y": 0.88, "w": 0.92, "h": 0.08}
    if name == "table_chart":
        return {"x": 0.08, "y": 0.24, "w": 0.84, "h": 0.56}
    if name == "diagonal":
        return {"x": 0.55, "y": 0.06, "w": 0.38, "h": 0.78}
    return {"x": 0.05, "y": 0.24, "w": 0.9, "h": 0.52}


def _canvas_size_match(rendered: dict[str, Any], reference: dict[str, Any]) -> bool:
    rendered_ratio = rendered["width_px"] / max(1, rendered["height_px"])
    reference_ratio = reference["width_px"] / max(1, reference["height_px"])
    return abs(rendered_ratio - reference_ratio) <= 0.03


def _palette_distance(first: tuple[float, float, float], second: tuple[float, float, float]) -> float:
    return round(math.sqrt(sum((a - b) ** 2 for a, b in zip(first, second))) / 441.67295593, 6)


def _rgb_distance(first: tuple[int, int, int], second: tuple[float, float, float]) -> float:
    return math.sqrt(sum((float(a) - float(b)) ** 2 for a, b in zip(first, second)))


def _compiled_layouts(template_spec: dict[str, Any], preview_manifest: dict[str, Any]) -> list[dict[str, Any]]:
    manifest_layouts = preview_manifest.get("compiled_layouts")
    if isinstance(manifest_layouts, list) and manifest_layouts:
        return [item for item in manifest_layouts if isinstance(item, dict)]
    return [
        {
            "slide_number": index,
            "layout_id": layout.get("layout_id"),
            "archetype_id": layout.get("archetype_id"),
        }
        for index, layout in enumerate(template_spec.get("layouts") or [], start=1)
        if isinstance(layout, dict)
    ]


def _layout_by_id(template_spec: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {str(layout.get("layout_id")): layout for layout in template_spec.get("layouts") or [] if isinstance(layout, dict)}


def _reference_images_by_archetype(template_image_manifest: dict[str, Any]) -> dict[str, Path]:
    result: dict[str, Path] = {}
    for record in template_image_manifest.get("images") or []:
        if isinstance(record, dict) and record.get("archetype_id") and record.get("image_output_path"):
            result[str(record["archetype_id"])] = Path(str(record["image_output_path"]))
    return result


def _template_preview_picture_findings(
    template_preview_pptx_path: str | Path,
    template_image_manifest: dict[str, Any],
) -> dict[int, list[dict[str, Any]]]:
    pptx_file = Path(template_preview_pptx_path)
    if not pptx_file.exists():
        return {}
    try:
        deck = Presentation(pptx_file)
    except Exception:
        return {}
    slide_w = deck.slide_width / EMU_PER_INCH
    slide_h = deck.slide_height / EMU_PER_INCH
    slide_area = max(0.01, slide_w * slide_h)
    reference_hashes = _template_reference_hashes(template_image_manifest)
    findings_by_slide: dict[int, list[dict[str, Any]]] = {}

    for slide_index, slide in enumerate(deck.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            bounds = _shape_bounds(shape)
            area_ratio = round((bounds["w"] * bounds["h"]) / slide_area, 6)
            details = {
                "shape_index": shape_index,
                "bounds": bounds,
                "area_ratio": area_ratio,
            }
            if _is_full_slide_picture(bounds, slide_w, slide_h, area_ratio):
                findings_by_slide.setdefault(slide_index, []).append(
                    _finding(
                        "FULL_SLIDE_RASTER_BACKGROUND_DETECTED",
                        "severe",
                        "A picture shape covers most of the template preview slide.",
                        slide_index,
                        details,
                    )
                )
            if _shape_image_sha256(shape) in reference_hashes:
                findings_by_slide.setdefault(slide_index, []).append(
                    _finding(
                        "REFERENCE_IMAGE_EMBEDDED_DIRECTLY",
                        "severe",
                        "A template reference PNG is embedded directly in the template preview PPTX.",
                        slide_index,
                        details,
                    )
                )
    return findings_by_slide


def _template_reference_hashes(template_image_manifest: dict[str, Any]) -> set[str]:
    hashes: set[str] = set()
    for path in _template_reference_paths(template_image_manifest):
        try:
            hashes.add(hashlib.sha256(path.read_bytes()).hexdigest())
        except OSError:
            continue
    return hashes


def _template_reference_paths(template_image_manifest: dict[str, Any]) -> list[Path]:
    paths = []
    for record in template_image_manifest.get("images") or []:
        if isinstance(record, dict) and isinstance(record.get("image_output_path"), str):
            paths.append(Path(record["image_output_path"]))
    return paths


def _shape_image_sha256(shape: Any) -> str:
    image = getattr(shape, "image", None)
    blob = getattr(image, "blob", b"") if image is not None else b""
    return hashlib.sha256(blob).hexdigest()


def _shape_bounds(shape: Any) -> dict[str, float]:
    return {
        "x": round(shape.left / EMU_PER_INCH, 4),
        "y": round(shape.top / EMU_PER_INCH, 4),
        "w": round(shape.width / EMU_PER_INCH, 4),
        "h": round(shape.height / EMU_PER_INCH, 4),
    }


def _is_full_slide_picture(bounds: dict[str, float], slide_w: float, slide_h: float, area_ratio: float) -> bool:
    nearly_origin = bounds["x"] <= 0.2 and bounds["y"] <= 0.2
    nearly_full_extent = bounds["w"] >= slide_w * 0.9 and bounds["h"] >= slide_h * 0.9
    return area_ratio >= FULL_SLIDE_AREA_RATIO_THRESHOLD and nearly_origin and nearly_full_extent


def _requires_title(layout: dict[str, Any], archetype_id: str) -> bool:
    if archetype_id == "section_divider" or layout.get("explicitly_omits_title"):
        return False
    return any(str(slot.get("slot_id") or "") == "title" for slot in layout.get("slots") or [] if isinstance(slot, dict))


def _requires_footer(layout: dict[str, Any]) -> bool:
    if layout.get("explicitly_omits_footer"):
        return False
    for slot in layout.get("slots") or []:
        if isinstance(slot, dict) and str(slot.get("slot_id") or "") == "footer":
            return True
    return False


def _has_slot(layout: dict[str, Any], slot_ids: set[str]) -> bool:
    return any(str(slot.get("slot_id") or "") in slot_ids for slot in layout.get("slots") or [] if isinstance(slot, dict))


def _similarity_summary(
    metrics: dict[str, Any],
    severe_violations: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> dict[str, Any]:
    if metrics is None:
        return {"status": "not_available", "label": "No rendered metrics available."}
    if severe_violations:
        status = "severe_violation"
    elif warnings:
        status = "needs_review"
    else:
        status = "close_match"
    palette_similarity = metrics.get("mean_palette_similarity")
    if palette_similarity is None:
        palette_similarity = metrics.get("dominant_palette_similarity")
    density_delta = metrics.get("blank_area_delta")
    return {
        "status": status,
        "palette_similarity": palette_similarity,
        "edge_density_delta": metrics.get("edge_density_delta"),
        "blank_area_delta": density_delta,
        "title_area_occupancy": metrics.get("title_area_occupancy"),
        "footer_area_occupancy": metrics.get("footer_area_occupancy"),
        "warning_count": len(warnings),
        "severe_count": len(severe_violations),
    }


def _rendered_path(render_dir: Path, render_report: dict[str, Any], slide_index: int) -> Path:
    for record in render_report.get("slides") or []:
        if isinstance(record, dict) and int(record.get("slide_index") or -1) == slide_index and record.get("rendered_image_path"):
            return Path(str(record["rendered_image_path"]))
    output_paths = render_report.get("output_paths") or []
    if slide_index - 1 < len(output_paths):
        return Path(str(output_paths[slide_index - 1]))
    return render_dir / f"slide-{slide_index:03d}.png"


def _render_status(render_report: dict[str, Any]) -> str:
    return str(render_report.get("render_status") or render_report.get("status") or "not_provided")


def _finding(
    code: str,
    severity: str,
    message: str,
    slide_index: int | None = None,
    details: dict[str, Any] | None = None,
) -> dict[str, Any]:
    payload = {"code": code, "severity": severity, "message": message}
    if slide_index is not None:
        payload["slide_index"] = slide_index
    if details:
        payload["details"] = details
    return payload


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Template Visual Diff Report",
        "",
        f"Status: `{report['status']}`",
        f"Render status: `{report['render_status']}`",
        f"Slides: `{report['slide_count']}`",
        f"Findings: `{report['findings_summary']['total']}` total, `{report['findings_summary']['severe']}` severe",
        f"QA blocks deck generation: `{str(report['qa_blocks_deck_generation']).lower()}`",
        "",
        "| Slide | Layout | Archetype | Similarity | Palette | Edge delta | Blank delta | Severe | Warnings |",
        "|---:|---|---|---|---:|---:|---:|---|---|",
    ]
    for slide in report["slides"]:
        metrics = slide.get("metrics") or {}
        summary = slide.get("similarity_summary") or {}
        severe = ", ".join(item["code"] for item in slide.get("severe_violations") or []) or "none"
        warnings = ", ".join(warning["code"] for warning in slide["warnings"]) or "none"
        lines.append(
            f"| {slide['slide_index']} | `{slide['layout_id']}` | `{slide['archetype_id']}` | "
            f"`{summary.get('status', 'not_available')}` | {_metric_text(metrics.get('mean_palette_similarity'))} | "
            f"{_metric_text(metrics.get('edge_density_delta'))} | {_metric_text(metrics.get('blank_area_delta'))} | {severe} | {warnings} |"
        )
    lines.append("")
    if report["findings"]:
        lines.extend(["## Findings", ""])
        for finding in report["findings"]:
            slide = f" slide {finding['slide_index']}" if "slide_index" in finding else ""
            lines.append(f"- `{finding['severity']}` `{finding['code']}`{slide}: {finding['message']}")
    return "\n".join(lines) + "\n"


def _metric_text(value: Any) -> str:
    return "" if value is None else str(value)


def _load_json(path: str | Path) -> Any:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _display_path(path: Path | None) -> str | None:
    if path is None:
        return None
    return str(path.as_posix())


if __name__ == "__main__":
    raise SystemExit(main())
