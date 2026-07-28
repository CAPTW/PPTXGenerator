"""Editable shape-grid chart/table placeholders for contract compilation."""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt

from src.presentation_agent.magic_layer.e06_2_contract_shape_factory import add_contract_shape


def add_contract_chart_table(slide: Any, obj: dict[str, Any]) -> Any:
    return add_contract_shape(slide, obj)


def _add_grid_lines(slide: Any, obj: dict[str, Any], rows: int, cols: int) -> None:
    bbox = obj["bbox_emu"]
    x, y, w, h = int(bbox["x"]), int(bbox["y"]), int(bbox["w"]), int(bbox["h"])
    for col in range(1, cols):
        xx = x + int(w * col / cols)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, xx, y, xx, y + h)
        line.name = f"contract_grid::{obj['object_id']}::v{col}"
        line.line.color.rgb = RGBColor(45, 61, 78)
        line.line.width = Pt(0.2)
    for row in range(1, rows):
        yy = y + int(h * row / rows)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, yy, x + w, yy)
        line.name = f"contract_grid::{obj['object_id']}::h{row}"
        line.line.color.rgb = RGBColor(45, 61, 78)
        line.line.width = Pt(0.2)
