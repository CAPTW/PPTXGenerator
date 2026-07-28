"""QA policy for local SVG icon library usage in editable PPTX outputs."""

from __future__ import annotations

import argparse
import json
import posixpath
import re
import xml.etree.ElementTree as ET
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_MANIFESTS_DIR = Path("assets/icons/manifests")
DEFAULT_NORMALIZED_DIR = Path("assets/icons/normalized")
DEFAULT_JSON_REPORT = Path("outputs/icon_usage_policy_report.json")
DEFAULT_MD_REPORT = Path("outputs/icon_usage_policy_report.md")
DEFAULT_PPTX_PATHS = (
    Path("outputs/golden_template_masters.pptx"),
    Path("outputs/final_deck_large_premium.pptx"),
    Path("design_runs/run_002/outputs/run_002_editable_master_pack.pptx"),
)
RUN_002_ICON_LINKAGE_REPORT_JSON = Path("design_runs/run_002/outputs/icon_linkage_audit_report.json")
RUN_002_ICON_LINKAGE_REPORT_MD = Path("design_runs/run_002/outputs/icon_linkage_audit_report.md")

ICON_LIBRARY_MANIFEST = "icon_library_manifest.json"
ICON_ROLE_MAP = "icon_role_map.json"
ICON_STYLE_TOKENS = "icon_style_tokens.json"
EMU_PER_INCH = 914400
DEFAULT_MAX_ICONS_PER_SLIDE = 8
DEFAULT_MAX_ICON_AREA_RATIO = 0.03
MAX_ALLOWED_SVG_ICON_AREA_RATIO = 0.03
PPTX_REL_TYPE_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
PPTX_PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PPTX_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": PPTX_PACKAGE_REL_NS,
}

REQUIRED_ICON_ROLES = [
    "academic",
    "source",
    "citation",
    "evidence",
    "research",
    "method",
    "insight",
    "risk",
    "governance",
    "data",
    "table",
    "chart",
    "kpi",
    "timeline",
    "process",
    "comparison",
    "decision",
    "user",
    "calendar",
    "presenter",
    "globe",
    "section",
    "index",
    "appendix",
    "warning",
    "recommendation",
]


def build_icon_usage_policy_report_from_files(
    *,
    manifests_dir: str | Path = DEFAULT_MANIFESTS_DIR,
    normalized_dir: str | Path = DEFAULT_NORMALIZED_DIR,
    pptx_paths: tuple[Path, ...] = DEFAULT_PPTX_PATHS,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
) -> dict[str, Any]:
    report = build_icon_usage_policy_report(
        manifests_dir=manifests_dir,
        normalized_dir=normalized_dir,
        pptx_paths=pptx_paths,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return report


def build_icon_usage_policy_report(
    *,
    manifests_dir: str | Path,
    normalized_dir: str | Path,
    pptx_paths: tuple[Path, ...] | list[Path],
) -> dict[str, Any]:
    manifests_root = Path(manifests_dir)
    normalized_root = Path(normalized_dir)
    findings: list[dict[str, Any]] = []

    manifest_scan = _scan_manifest_files(manifests_root, findings)
    library_manifest = manifest_scan["loaded"].get(ICON_LIBRARY_MANIFEST, {})
    role_map = manifest_scan["loaded"].get(ICON_ROLE_MAP, {})
    style_tokens = manifest_scan["loaded"].get(ICON_STYLE_TOKENS, {})

    icon_records = _icon_records(library_manifest)
    role_entries = _role_entries(role_map)
    normalized_scan = _scan_normalized_svgs(normalized_root, icon_records, findings)
    role_coverage = _scan_required_role_coverage(role_entries, icon_records, normalized_scan["normalized_paths"], findings)
    asset_file_scan = _scan_icon_asset_files(findings)
    family_scan = _scan_icon_family_consistency(library_manifest, icon_records, findings)
    pptx_scans = []
    for pptx_path in pptx_paths:
        path = Path(pptx_path)
        scan = _scan_pptx_icon_usage(path, style_tokens, findings)
        scan["svg_linkage_audit"] = _audit_pptx_svg_linkage(path, role_entries, findings)
        pptx_scans.append(scan)

    severe = sum(1 for finding in findings if finding["severity"] == "severe")
    warning = sum(1 for finding in findings if finding["severity"] == "warning")
    normalized_count = normalized_scan["normalized_svg_count"]
    missing_role_count = len(role_coverage["missing_roles"])
    status = "failed" if severe else "needs_icon_assets" if normalized_count == 0 or missing_role_count else "issues_reported" if warning else "passed"
    return {
        "schema_name": "icon_usage_policy_report",
        "schema_version": "1.0",
        "status": status,
        "qa_blocks_template_usage": bool(severe),
        "manifests_dir": _display_path(manifests_root),
        "normalized_dir": _display_path(normalized_root),
        "outputs": {
            "json_report": _display_path(DEFAULT_JSON_REPORT),
            "markdown_report": _display_path(DEFAULT_MD_REPORT),
        },
        "findings_summary": {
            "total": len(findings),
            "severe": severe,
            "warning": warning,
        },
        "findings_by_code": dict(sorted(Counter(finding["code"] for finding in findings).items())),
        "manifest_scan": {
            key: value
            for key, value in manifest_scan.items()
            if key != "loaded"
        },
        "normalized_svg_scan": normalized_scan,
        "asset_file_scan": asset_file_scan,
        "role_coverage": role_coverage,
        "family_consistency": family_scan,
        "pptx_scans": pptx_scans,
        "findings": findings,
        "policy": {
            "svg_icons_allowed_as_vector_assets": True,
            "png_icons_forbidden": True,
            "full_slide_raster_background_forbidden": True,
            "raster_svg_content_forbidden": True,
            "external_svg_references_forbidden": True,
            "icons_are_supporting_visuals_only": True,
            "icons_must_not_satisfy_required_content_slots": True,
            "max_icons_per_slide": _max_icons_per_slide(style_tokens),
            "gate_integration_target": "gate:template-usability",
        },
    }


def _scan_manifest_files(manifests_root: Path, findings: list[dict[str, Any]]) -> dict[str, Any]:
    required = [ICON_LIBRARY_MANIFEST, ICON_ROLE_MAP, ICON_STYLE_TOKENS]
    loaded: dict[str, dict[str, Any]] = {}
    files = sorted(manifests_root.glob("*.json")) if manifests_root.exists() else []
    if not manifests_root.exists():
        findings.append(_finding("ICON_MANIFESTS_DIR_MISSING", "severe", "Icon manifests directory is missing.", {"path": _display_path(manifests_root)}))
    for name in required:
        path = manifests_root / name
        if not path.exists():
            findings.append(_finding("ICON_REQUIRED_MANIFEST_MISSING", "severe", f"Required icon manifest `{name}` is missing.", {"path": _display_path(path)}))
            continue
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            findings.append(_finding("ICON_MANIFEST_INVALID_JSON", "severe", f"Icon manifest `{name}` is not valid JSON.", {"path": _display_path(path), "error": str(exc)}))
            continue
        if not isinstance(payload, dict):
            findings.append(_finding("ICON_MANIFEST_NOT_OBJECT", "severe", f"Icon manifest `{name}` must be a JSON object.", {"path": _display_path(path)}))
            continue
        loaded[name] = payload
    return {
        "manifests_dir_exists": manifests_root.exists(),
        "manifest_files": [_display_path(path) for path in files],
        "required_manifests": required,
        "required_manifests_present": [name for name in required if name in loaded],
        "loaded": loaded,
    }


def _icon_records(library_manifest: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {
        str(record.get("icon_id")): record
        for record in library_manifest.get("icons") or []
        if isinstance(record, dict) and record.get("icon_id")
    }


def _role_entries(role_map: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {
        str(record.get("role")): record
        for record in role_map.get("roles") or []
        if isinstance(record, dict) and record.get("role")
    }


def _scan_normalized_svgs(normalized_root: Path, icon_records: dict[str, dict[str, Any]], findings: list[dict[str, Any]]) -> dict[str, Any]:
    svg_paths = sorted(normalized_root.rglob("*.svg")) if normalized_root.exists() else []
    normalized_paths = {_canonical_rel(path): path for path in svg_paths}
    svg_reports = []
    if not normalized_root.exists():
        findings.append(_finding("ICON_NORMALIZED_DIR_MISSING", "warning", "Normalized SVG icon directory is missing.", {"path": _display_path(normalized_root)}))
    for path in svg_paths:
        svg_findings = _inspect_svg_file(path)
        for item in svg_findings:
            findings.append(item)
        matching_record = _record_for_normalized_path(path, icon_records)
        if matching_record is None:
            findings.append(_finding("ICON_NORMALIZED_SVG_NOT_IN_MANIFEST", "warning", "Normalized SVG is not listed in icon_library_manifest.json.", {"path": _display_path(path)}))
        elif not _has_license_source_metadata(matching_record):
            findings.append(
                _finding(
                    "ICON_LICENSE_SOURCE_METADATA_MISSING",
                    "warning",
                    "Icon manifest record is missing license or source metadata.",
                    {
                        "icon_id": matching_record.get("icon_id"),
                        "normalized_path": _display_path(path),
                        "license": matching_record.get("license"),
                        "source_path": matching_record.get("source_path"),
                    },
                )
            )
        svg_reports.append(
            {
                "path": _display_path(path),
                "manifest_icon_id": matching_record.get("icon_id") if matching_record else None,
                "has_manifest_record": matching_record is not None,
                "license_source_metadata_present": bool(matching_record and _has_license_source_metadata(matching_record)),
                "finding_codes": [item["code"] for item in svg_findings],
            }
        )
    return {
        "normalized_dir_exists": normalized_root.exists(),
        "normalized_svg_count": len(svg_paths),
        "normalized_paths": sorted(normalized_paths),
        "svgs": svg_reports,
    }


def _inspect_svg_file(path: Path) -> list[dict[str, Any]]:
    findings: list[dict[str, Any]] = []
    try:
        text = path.read_text(encoding="utf-8")
    except OSError as exc:
        return [_finding("ICON_SVG_UNREADABLE", "severe", "Normalized SVG cannot be read.", {"path": _display_path(path), "error": str(exc)})]
    try:
        root = ET.fromstring(text)
    except ET.ParseError as exc:
        return [_finding("ICON_SVG_INVALID_XML", "severe", "Normalized SVG is not valid XML.", {"path": _display_path(path), "error": str(exc)})]
    for element in root.iter():
        tag = _local_name(element.tag)
        if tag == "image":
            findings.append(_finding("ICON_SVG_RASTER_IMAGE_ELEMENT", "severe", "SVG icon contains an image element.", {"path": _display_path(path)}))
        if tag in {"script", "foreignObject"}:
            findings.append(_finding("ICON_SVG_FORBIDDEN_ELEMENT", "severe", f"SVG icon contains forbidden `{tag}` element.", {"path": _display_path(path), "element": tag}))
        if tag == "text":
            findings.append(_finding("ICON_SVG_TEXT_ELEMENT", "warning", "SVG icon contains text; icons should be path/vector primitives.", {"path": _display_path(path)}))
        for attr_name, attr_value in element.attrib.items():
            if _local_name(attr_name).lower() in {"href", "src"}:
                href = str(attr_value or "").strip()
                if href.lower().startswith("data:image"):
                    findings.append(_finding("ICON_SVG_EMBEDDED_RASTER", "severe", "SVG icon contains embedded raster data.", {"path": _display_path(path)}))
                elif _is_external_reference(href):
                    findings.append(_finding("ICON_SVG_EXTERNAL_REFERENCE", "severe", "SVG icon contains an external reference.", {"path": _display_path(path), "href": href}))
    if re.search(r"data:image/(png|jpe?g|gif|webp|bmp)", text, re.IGNORECASE):
        findings.append(_finding("ICON_SVG_EMBEDDED_RASTER", "severe", "SVG icon contains embedded raster data.", {"path": _display_path(path)}))
    return findings


def _scan_required_role_coverage(
    role_entries: dict[str, dict[str, Any]],
    icon_records: dict[str, dict[str, Any]],
    normalized_paths: list[str],
    findings: list[dict[str, Any]],
) -> dict[str, Any]:
    normalized_set = set(normalized_paths)
    roles = sorted(set(REQUIRED_ICON_ROLES) | set(role_entries))
    role_reports: list[dict[str, Any]] = []
    missing_roles: list[str] = []
    for role in roles:
        entry = role_entries.get(role)
        if entry is None:
            missing_roles.append(role)
            findings.append(_finding("ICON_REQUIRED_ROLE_MISSING_FROM_ROLE_MAP", "severe", f"Required icon role `{role}` is missing from icon_role_map.json.", {"role": role}))
            role_reports.append({"role": role, "status": "missing_role_map_entry", "candidate_icon_ids": [], "resolved_icon_id": None})
            continue
        candidates = _candidate_icon_ids(entry)
        resolved = _resolved_candidate(candidates, icon_records, normalized_set)
        status = "covered" if resolved else "missing_icon_asset"
        if not resolved:
            missing_roles.append(role)
            findings.append(_finding("ICON_REQUIRED_ROLE_MISSING_ASSET", "warning", f"Required icon role `{role}` has no allowed normalized SVG asset.", {"role": role, "candidate_icon_ids": candidates}))
        role_reports.append(
            {
                "role": role,
                "status": status,
                "candidate_icon_ids": candidates,
                "resolved_icon_id": resolved.get("icon_id") if resolved else None,
                "resolved_family": resolved.get("source_family") if resolved else None,
                "resolved_normalized_path": resolved.get("normalized_path") if resolved else None,
            }
        )
    return {
        "required_roles": roles,
        "covered_roles": [item["role"] for item in role_reports if item["status"] == "covered"],
        "missing_roles": missing_roles,
        "roles": role_reports,
    }


def _scan_icon_asset_files(findings: list[dict[str, Any]]) -> dict[str, Any]:
    icon_root = Path("assets/icons")
    raster_paths = sorted(
        [
            path
            for pattern in ("*.png", "*.jpg", "*.jpeg", "*.webp", "*.gif", "*.bmp")
            for path in (icon_root.rglob(pattern) if icon_root.exists() else [])
        ]
    )
    png_icon_paths = [path for path in raster_paths if _is_forbidden_png_icon_asset(path)]
    vendor_non_icon_rasters = [path for path in raster_paths if path not in png_icon_paths]
    for path in png_icon_paths:
        findings.append(_finding("PNG_ICON_FILE_PRESENT", "severe", "PNG icon asset file exists in a compiler icon path.", {"path": _display_path(path)}))
    return {
        "icon_asset_root": _display_path(icon_root),
        "png_icon_file_count": len(png_icon_paths),
        "png_icon_files": [_display_path(path) for path in png_icon_paths],
        "vendor_non_icon_raster_file_count": len(vendor_non_icon_rasters),
        "vendor_non_icon_raster_examples": [_display_path(path) for path in vendor_non_icon_rasters[:20]],
    }


def _is_forbidden_png_icon_asset(path: Path) -> bool:
    parts = [part.lower() for part in path.parts]
    if "normalized" in parts:
        return True
    if "vendor" in parts:
        vendor_index = parts.index("vendor")
        vendor_parts = parts[vendor_index + 1 :]
        # Vendored repositories can contain raster documentation and social
        # preview assets. Treat only source icon directories as forbidden PNG
        # icon assets.
        return len(vendor_parts) >= 2 and vendor_parts[1] == "icons"
    if "assets" in parts and "icons" in parts:
        return True
    return False


def _scan_icon_family_consistency(library_manifest: dict[str, Any], icon_records: dict[str, dict[str, Any]], findings: list[dict[str, Any]]) -> dict[str, Any]:
    preferred = library_manifest.get("preferred_icon_families") if isinstance(library_manifest.get("preferred_icon_families"), dict) else {}
    allowed_families = {str(preferred.get("primary") or "tabler"), str(preferred.get("fallback") or "lucide")}
    used_families = sorted({str(record.get("source_family") or "unknown") for record in icon_records.values() if record.get("allowed_for_template")})
    unexpected = sorted(family for family in used_families if family not in allowed_families)
    if unexpected:
        findings.append(_finding("ICON_FAMILY_UNEXPECTED", "warning", "Allowed icon records include families outside the preferred primary/fallback set.", {"unexpected_families": unexpected, "allowed_families": sorted(allowed_families)}))
    if len(used_families) > 1:
        findings.append(_finding("ICON_FAMILY_MIXED", "warning", "Allowed icon records mix multiple icon families; keep this intentional and role-mapped.", {"families": used_families}))
    return {
        "primary_family": preferred.get("primary") or "tabler",
        "fallback_family": preferred.get("fallback") or "lucide",
        "allowed_families": sorted(allowed_families),
        "families_in_allowed_icon_records": used_families,
        "unexpected_families": unexpected,
        "status": "mixed" if len(used_families) > 1 else "single_family" if used_families else "no_icon_assets",
    }


def _scan_pptx_icon_usage(pptx_path: Path, style_tokens: dict[str, Any], findings: list[dict[str, Any]]) -> dict[str, Any]:
    if not pptx_path.exists():
        return {
            "pptx_path": _display_path(pptx_path),
            "status": "skipped_missing_pptx",
            "slide_count": 0,
            "slides": [],
            "icon_shape_count": 0,
            "png_icon_asset_count": 0,
        }
    deck = Presentation(pptx_path)
    max_icons = _max_icons_per_slide(style_tokens)
    max_area_ratio = _max_icon_area_ratio(style_tokens)
    slides: list[dict[str, Any]] = []
    deck_icon_count = 0
    png_icon_count = 0
    for slide_index, slide in enumerate(deck.slides, start=1):
        slide_w = deck.slide_width / EMU_PER_INCH
        slide_h = deck.slide_height / EMU_PER_INCH
        slide_area = max(0.01, slide_w * slide_h)
        icon_shapes: list[dict[str, Any]] = []
        text_shape_count = 0
        table_count = 0
        chart_count = 0
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False) and str(shape.text or "").strip():
                text_shape_count += 1
            if getattr(shape, "has_table", False):
                table_count += 1
            if getattr(shape, "has_chart", False):
                chart_count += 1
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            image_info = _shape_image_info(shape)
            bounds = _shape_bounds(shape)
            area_ratio = round((bounds["w"] * bounds["h"]) / slide_area, 6)
            is_icon = _is_icon_picture_shape(shape, image_info, area_ratio)
            if not is_icon:
                continue
            content_type = str(image_info.get("content_type") or "")
            partname = str(image_info.get("partname") or "")
            is_png = content_type == "image/png" or partname.lower().endswith(".png")
            is_svg = content_type == "image/svg+xml" or partname.lower().endswith(".svg")
            if is_png:
                png_icon_count += 1
                findings.append(_finding("PPTX_PNG_ICON_ASSET", "severe", "PPTX contains a PNG icon asset.", {"pptx_path": _display_path(pptx_path), "slide_number": slide_index, "shape_index": shape_index, "partname": partname, "content_type": content_type}))
            if not is_svg:
                findings.append(_finding("PPTX_NON_SVG_ICON_ASSET", "severe", "Icon-like PPTX picture is not an SVG vector asset.", {"pptx_path": _display_path(pptx_path), "slide_number": slide_index, "shape_index": shape_index, "partname": partname, "content_type": content_type}))
            if area_ratio > max_area_ratio:
                findings.append(_finding("ICON_REPLACES_CONTENT_SLOT_RISK", "severe", "Icon shape is too large and may be replacing a content slot.", {"pptx_path": _display_path(pptx_path), "slide_number": slide_index, "shape_index": shape_index, "area_ratio": area_ratio, "limit": max_area_ratio}))
            icon_shapes.append(
                {
                    "shape_index": shape_index,
                    "name": str(getattr(shape, "name", "") or ""),
                    "description": _shape_description(shape),
                    "bounds": bounds,
                    "area_ratio": area_ratio,
                    "content_type": content_type,
                    "partname": partname,
                    "is_svg": is_svg,
                    "is_png": is_png,
                }
            )
        if icon_shapes and text_shape_count == 0 and table_count == 0 and chart_count == 0:
            findings.append(_finding("ICON_REPLACES_CONTENT_SLOT_RISK", "severe", "Slide contains icon assets but no editable content shapes.", {"pptx_path": _display_path(pptx_path), "slide_number": slide_index, "icon_count": len(icon_shapes)}))
        if len(icon_shapes) > max_icons:
            findings.append(_finding("ICON_COUNT_EXCEEDS_DECORATIVE_BUDGET", "severe", "Slide icon count exceeds decorative budget.", {"pptx_path": _display_path(pptx_path), "slide_number": slide_index, "icon_count": len(icon_shapes), "limit": max_icons}))
        deck_icon_count += len(icon_shapes)
        slides.append(
            {
                "slide_number": slide_index,
                "icon_count": len(icon_shapes),
                "text_shape_count": text_shape_count,
                "table_count": table_count,
                "chart_count": chart_count,
                "icons": icon_shapes,
            }
        )
    return {
        "pptx_path": _display_path(pptx_path),
        "status": "scanned",
        "slide_count": len(deck.slides),
        "icon_shape_count": deck_icon_count,
        "png_icon_asset_count": png_icon_count,
        "slides": slides,
    }


def _audit_pptx_svg_linkage(pptx_path: Path, role_entries: dict[str, dict[str, Any]], findings: list[dict[str, Any]]) -> dict[str, Any]:
    if not pptx_path.exists():
        return _empty_linkage_audit(pptx_path, "skipped_missing_pptx")

    compile_report_path = _compile_report_path_for_pptx(pptx_path)
    compile_report = _load_json(compile_report_path) if compile_report_path else {}
    expected_icons = _expected_icons_from_compile_report(compile_report)
    unresolved_icon_roles = list((compile_report.get("icon_report") or {}).get("unresolved_icon_roles") or [])

    try:
        with zipfile.ZipFile(pptx_path) as archive:
            names = set(archive.namelist())
            svg_media = sorted(name for name in names if name.startswith("ppt/media/") and name.lower().endswith(".svg"))
            png_media = sorted(name for name in names if name.startswith("ppt/media/") and name.lower().endswith(".png"))
            raster_media = sorted(
                name
                for name in names
                if name.startswith("ppt/media/")
                and name.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"))
            )
            rel_maps, svg_relationships = _pptx_relationship_maps(archive, names)
            picture_records, svg_blip_records = _pptx_picture_records(archive, names, rel_maps)
    except Exception as exc:  # noqa: BLE001
        findings.append(_finding("PPTX_ICON_LINKAGE_AUDIT_FAILED", "severe", "PPTX SVG icon linkage audit could not inspect the package.", {"pptx_path": _display_path(pptx_path), "error": str(exc)}))
        audit = _empty_linkage_audit(pptx_path, "failed")
        audit["errors"] = [str(exc)]
        return audit

    referenced_svg = sorted({rel["target_part"] for rel in svg_relationships})
    blipped_svg = sorted({record["target_part"] for record in svg_blip_records if record.get("target_part")})
    orphan_svg = sorted(set(svg_media) - set(referenced_svg))
    unblipped_svg_relationships = [
        rel
        for rel in svg_relationships
        if not any(record.get("source_part") == rel.get("source_part") and record.get("relationship_id") == rel.get("relationship_id") for record in svg_blip_records)
    ]
    visible_icon_objects = [
        record
        for record in picture_records
        if record["is_visible"] and record["is_svg"] and record["is_compiler_marked_icon"]
    ]
    rendered_icon_counts = Counter(record["icon_id"] for record in visible_icon_objects if record.get("icon_id"))
    missing_expected = _missing_expected_icons(expected_icons, rendered_icon_counts)
    role_render = _rendered_role_counts(expected_icons, visible_icon_objects, role_entries)

    png_icon_records = [record for record in picture_records if record["is_compiler_marked_icon"] and record.get("target_part", "").lower().endswith(".png")]
    raster_icon_records = [
        record
        for record in picture_records
        if record["is_compiler_marked_icon"]
        and record.get("target_part", "").lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"))
    ]

    if svg_media and orphan_svg:
        findings.append(
            _finding(
                "PPTX_ORPHAN_SVG_ICON_MEDIA",
                "severe",
                "SVG media exists in the PPTX package but is not referenced by any slide, layout, or master relationship.",
                {"pptx_path": _display_path(pptx_path), "orphan_svg_media": orphan_svg},
            )
        )
    if svg_relationships and unblipped_svg_relationships:
        findings.append(
            _finding(
                "PPTX_SVG_RELATIONSHIP_WITHOUT_VISIBLE_OBJECT",
                "severe",
                "SVG image relationships exist but are not used by any picture blip.",
                {"pptx_path": _display_path(pptx_path), "relationships": unblipped_svg_relationships},
            )
        )
    if expected_icons and missing_expected:
        findings.append(
            _finding(
                "ICON_USED_IN_COMPILE_REPORT_NOT_RENDERED",
                "severe",
                "One or more icon_role entries marked used in the compile report have no visible SVG slide object.",
                {"pptx_path": _display_path(pptx_path), "missing_expected_icons": missing_expected},
            )
        )
    for record in png_icon_records:
        findings.append(_finding("PPTX_PNG_ICON_ASSET", "severe", "PPTX contains a PNG icon asset.", {"pptx_path": _display_path(pptx_path), "picture": record}))
    for record in raster_icon_records:
        findings.append(_finding("PPTX_RASTER_ICON_ASSET", "severe", "PPTX contains a raster image used as an icon.", {"pptx_path": _display_path(pptx_path), "picture": record}))
    for record in visible_icon_objects:
        if record["area_ratio"] > MAX_ALLOWED_SVG_ICON_AREA_RATIO:
            findings.append(
                _finding(
                    "ICON_REPLACES_CONTENT_SLOT_RISK",
                    "severe",
                    "SVG icon object is too large and may be replacing a required content slot.",
                    {"pptx_path": _display_path(pptx_path), "picture": record, "limit": MAX_ALLOWED_SVG_ICON_AREA_RATIO},
                )
            )

    return {
        "pptx_path": _display_path(pptx_path),
        "status": "scanned",
        "compile_report_path": _display_path(compile_report_path) if compile_report_path else None,
        "svg_media_count": len(svg_media),
        "referenced_svg_media_count": len(referenced_svg),
        "orphan_svg_media_count": len(orphan_svg),
        "orphan_svg_media": orphan_svg,
        "slide_picture_object_count": len([record for record in picture_records if record["source_part"].startswith("ppt/slides/")]),
        "visible_icon_object_count": len(visible_icon_objects),
        "svg_relationship_count": len(svg_relationships),
        "svg_blip_count": len(svg_blip_records),
        "png_icon_count": len(png_icon_records),
        "raster_icon_count": len(raster_icon_records),
        "icon_roles_expected": sorted(set(item["icon_role"] for item in expected_icons if item.get("icon_role"))),
        "icon_role_expected_counts": dict(sorted(Counter(item["icon_role"] for item in expected_icons if item.get("icon_role")).items())),
        "icon_roles_rendered": sorted(role_render),
        "icon_role_rendered_counts": dict(sorted(role_render.items())),
        "expected_icon_count": len(expected_icons),
        "missing_expected_icon_count": len(missing_expected),
        "missing_expected_icons": missing_expected,
        "unresolved_icon_roles": unresolved_icon_roles,
        "visible_icon_objects": visible_icon_objects,
        "svg_relationships": svg_relationships,
        "svg_blips": svg_blip_records,
        "unblipped_svg_relationships": unblipped_svg_relationships,
        "remaining_icon_linkage_failures": [
            *([{"code": "PPTX_ORPHAN_SVG_ICON_MEDIA", "count": len(orphan_svg)}] if orphan_svg else []),
            *([{"code": "PPTX_SVG_RELATIONSHIP_WITHOUT_VISIBLE_OBJECT", "count": len(unblipped_svg_relationships)}] if unblipped_svg_relationships else []),
            *([{"code": "ICON_USED_IN_COMPILE_REPORT_NOT_RENDERED", "count": len(missing_expected)}] if missing_expected else []),
            *([{"code": "PPTX_PNG_ICON_ASSET", "count": len(png_icon_records)}] if png_icon_records else []),
            *([{"code": "PPTX_RASTER_ICON_ASSET", "count": len(raster_icon_records)}] if raster_icon_records else []),
        ],
    }


def _empty_linkage_audit(pptx_path: Path, status: str) -> dict[str, Any]:
    return {
        "pptx_path": _display_path(pptx_path),
        "status": status,
        "svg_media_count": 0,
        "referenced_svg_media_count": 0,
        "orphan_svg_media_count": 0,
        "slide_picture_object_count": 0,
        "visible_icon_object_count": 0,
        "svg_relationship_count": 0,
        "svg_blip_count": 0,
        "icon_roles_expected": [],
        "icon_roles_rendered": [],
        "unresolved_icon_roles": [],
        "png_icon_count": 0,
        "raster_icon_count": 0,
        "remaining_icon_linkage_failures": [],
    }


def _pptx_relationship_maps(archive: zipfile.ZipFile, names: set[str]) -> tuple[dict[str, dict[str, dict[str, Any]]], list[dict[str, Any]]]:
    rel_maps: dict[str, dict[str, dict[str, Any]]] = {}
    svg_relationships: list[dict[str, Any]] = []
    rel_paths = sorted(
        name
        for name in names
        if (
            name.startswith("ppt/slides/_rels/")
            or name.startswith("ppt/slideLayouts/_rels/")
            or name.startswith("ppt/slideMasters/_rels/")
        )
        and name.endswith(".rels")
    )
    for rel_path in rel_paths:
        source_part = _source_part_for_rels(rel_path)
        if source_part is None:
            continue
        try:
            root = ET.fromstring(archive.read(rel_path))
        except ET.ParseError:
            continue
        rel_map: dict[str, dict[str, Any]] = {}
        for rel in root.findall("rel:Relationship", PPTX_NS):
            rel_id = str(rel.get("Id") or "")
            target = str(rel.get("Target") or "")
            target_mode = str(rel.get("TargetMode") or "")
            rel_type = str(rel.get("Type") or "")
            target_part = _resolve_relationship_target(source_part, target) if target_mode.lower() != "external" else target
            payload = {
                "relationship_id": rel_id,
                "relationship_part": rel_path,
                "source_part": source_part,
                "type": rel_type,
                "target": target,
                "target_part": target_part,
                "target_mode": target_mode or None,
            }
            rel_map[rel_id] = payload
            if rel_type == PPTX_REL_TYPE_IMAGE and target_part.lower().endswith(".svg"):
                svg_relationships.append(payload)
        rel_maps[source_part] = rel_map
    return rel_maps, svg_relationships


def _pptx_picture_records(
    archive: zipfile.ZipFile,
    names: set[str],
    rel_maps: dict[str, dict[str, dict[str, Any]]],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    picture_records: list[dict[str, Any]] = []
    svg_blip_records: list[dict[str, Any]] = []
    xml_parts = sorted(
        name
        for name in names
        if (
            (name.startswith("ppt/slides/") and name.endswith(".xml"))
            or (name.startswith("ppt/slideLayouts/") and name.endswith(".xml"))
            or (name.startswith("ppt/slideMasters/") and name.endswith(".xml"))
        )
    )
    for source_part in xml_parts:
        try:
            root = ET.fromstring(archive.read(source_part))
        except ET.ParseError:
            continue
        rel_map = rel_maps.get(source_part, {})
        for pic in root.findall(".//p:pic", PPTX_NS):
            record = _picture_record(source_part, pic, rel_map)
            picture_records.append(record)
            if record["is_svg"]:
                svg_blip_records.append(record)
    return picture_records, svg_blip_records


def _picture_record(source_part: str, pic: ET.Element, rel_map: dict[str, dict[str, Any]]) -> dict[str, Any]:
    c_nv_pr = pic.find(".//p:cNvPr", PPTX_NS)
    blip = pic.find(".//a:blip", PPTX_NS)
    xfrm = pic.find(".//a:xfrm", PPTX_NS)
    ext = xfrm.find("a:ext", PPTX_NS) if xfrm is not None else None
    r_id = blip.get(f"{{{PPTX_NS['r']}}}embed") if blip is not None else None
    rel = rel_map.get(str(r_id or ""), {})
    target_part = str(rel.get("target_part") or "")
    name = str((c_nv_pr.get("name") if c_nv_pr is not None else "") or "")
    description = str((c_nv_pr.get("descr") if c_nv_pr is not None else "") or "")
    cx = _int_or_zero(ext.get("cx") if ext is not None else None)
    cy = _int_or_zero(ext.get("cy") if ext is not None else None)
    area_ratio = round((cx * cy) / max(1, int(13.333 * EMU_PER_INCH) * int(7.5 * EMU_PER_INCH)), 6)
    icon_id = _icon_id_from_picture(name, description)
    return {
        "source_part": source_part,
        "relationship_id": r_id,
        "target_part": target_part,
        "name": name,
        "description": description,
        "icon_id": icon_id,
        "cx": cx,
        "cy": cy,
        "area_ratio": area_ratio,
        "is_visible": cx > 0 and cy > 0,
        "is_svg": target_part.lower().endswith(".svg"),
        "is_compiler_marked_icon": bool(icon_id),
    }


def _source_part_for_rels(rel_path: str) -> str | None:
    if "/_rels/" not in rel_path or not rel_path.endswith(".rels"):
        return None
    prefix, rel_name = rel_path.split("/_rels/", 1)
    return f"{prefix}/{rel_name[:-5]}"


def _resolve_relationship_target(source_part: str, target: str) -> str:
    if target.startswith("/"):
        return posixpath.normpath(target.lstrip("/"))
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def _icon_id_from_picture(name: str, description: str) -> str | None:
    if description.startswith("svg-icon:"):
        return description.split(":", 1)[1].strip() or None
    if name.startswith("SVG Icon "):
        return name.removeprefix("SVG Icon ").strip() or None
    return None


def _compile_report_path_for_pptx(pptx_path: Path) -> Path | None:
    candidates = [
        pptx_path.with_name(f"{pptx_path.stem}_report.json"),
        pptx_path.with_name(f"{pptx_path.stem}_manifest.json"),
    ]
    if pptx_path.name == "run_002_editable_master_pack.pptx":
        candidates.insert(0, pptx_path.with_name("run_002_master_pack_report.json"))
    if pptx_path.name == "golden_template_masters.pptx":
        candidates.insert(0, Path("outputs/golden_template_masters_report.json"))
    if pptx_path.name == "final_deck_large_premium.pptx":
        candidates.extend([Path("outputs/final_deck_manifest.json"), Path("outputs/large_premium_input_manifest.json")])
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _expected_icons_from_compile_report(report: dict[str, Any]) -> list[dict[str, Any]]:
    icon_report = report.get("icon_report") if isinstance(report.get("icon_report"), dict) else {}
    icons = icon_report.get("icons_used") if isinstance(icon_report.get("icons_used"), list) else []
    normalized: list[dict[str, Any]] = []
    for item in icons:
        if not isinstance(item, dict):
            continue
        normalized.append(
            {
                "icon_id": str(item.get("icon_id") or ""),
                "icon_role": str(item.get("icon_role") or ""),
                "relationship_id": item.get("relationship_id"),
                "object_id": item.get("object_id"),
                "object_name": item.get("object_name"),
                "object_description": item.get("object_description"),
                "slide_part": item.get("slide_part"),
                "context": item.get("context") or {},
            }
        )
    return normalized


def _missing_expected_icons(expected_icons: list[dict[str, Any]], rendered_icon_counts: Counter[str]) -> list[dict[str, Any]]:
    available = Counter(rendered_icon_counts)
    missing: list[dict[str, Any]] = []
    for item in expected_icons:
        icon_id = item.get("icon_id")
        if icon_id and available[icon_id] > 0:
            available[icon_id] -= 1
            continue
        missing.append(item)
    return missing


def _rendered_role_counts(
    expected_icons: list[dict[str, Any]],
    visible_icon_objects: list[dict[str, Any]],
    role_entries: dict[str, dict[str, Any]],
) -> Counter[str]:
    rendered_counts = Counter(record["icon_id"] for record in visible_icon_objects if record.get("icon_id"))
    role_counts: Counter[str] = Counter()
    if expected_icons:
        for item in expected_icons:
            icon_id = item.get("icon_id")
            role = item.get("icon_role")
            if icon_id and role and rendered_counts[icon_id] > 0:
                rendered_counts[icon_id] -= 1
                role_counts[role] += 1
        return role_counts
    icon_to_roles: dict[str, list[str]] = {}
    for role, entry in role_entries.items():
        preferred = str(entry.get("preferred_icon_id") or "")
        if preferred:
            icon_to_roles.setdefault(preferred, []).append(role)
    for record in visible_icon_objects:
        roles = icon_to_roles.get(str(record.get("icon_id") or ""), [])
        for role in roles:
            role_counts[role] += 1
    return role_counts


def _load_json(path: Path | None) -> dict[str, Any]:
    if path is None or not path.exists():
        return {}
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return value if isinstance(value, dict) else {}


def _int_or_zero(value: str | None) -> int:
    try:
        return int(value or 0)
    except (TypeError, ValueError):
        return 0


def _candidate_icon_ids(role_entry: dict[str, Any]) -> list[str]:
    candidates = []
    preferred = role_entry.get("preferred_icon_id")
    if isinstance(preferred, str) and preferred.strip():
        candidates.append(preferred.strip())
    for item in role_entry.get("fallback_icon_ids") or []:
        if isinstance(item, str) and item.strip():
            candidates.append(item.strip())
    return _dedupe(candidates)


def _resolved_candidate(candidates: list[str], icon_records: dict[str, dict[str, Any]], normalized_paths: set[str]) -> dict[str, Any] | None:
    for candidate in candidates:
        record = icon_records.get(candidate)
        if not record or not record.get("allowed_for_template"):
            continue
        normalized = str(record.get("normalized_path") or "").replace("\\", "/")
        if normalized in normalized_paths:
            return record
    return None


def _record_for_normalized_path(path: Path, icon_records: dict[str, dict[str, Any]]) -> dict[str, Any] | None:
    canonical = _canonical_rel(path)
    for record in icon_records.values():
        if str(record.get("normalized_path") or "").replace("\\", "/") == canonical:
            return record
    return None


def _has_license_source_metadata(record: dict[str, Any]) -> bool:
    license_value = str(record.get("license") or "").strip().lower()
    source_path = str(record.get("source_path") or "").strip()
    return bool(source_path and license_value and license_value not in {"unknown", "none", "tbd", "todo"})


def _is_external_reference(href: str) -> bool:
    lowered = href.strip().lower()
    if not lowered or lowered.startswith("#"):
        return False
    return lowered.startswith(("http://", "https://", "file:", "//", "data:")) or "://" in lowered


def _is_icon_picture_shape(shape: Any, image_info: dict[str, Any], area_ratio: float) -> bool:
    name = str(getattr(shape, "name", "") or "").lower()
    description = _shape_description(shape).lower()
    partname = str(image_info.get("partname") or "").lower()
    content_type = str(image_info.get("content_type") or "").lower()
    named_icon = "icon" in name or description.startswith("svg-icon:") or "icon" in description
    svg_icon = content_type == "image/svg+xml" and (named_icon or area_ratio <= DEFAULT_MAX_ICON_AREA_RATIO)
    return named_icon or svg_icon or ("icon" in partname and content_type.startswith("image/"))


def _shape_image_info(shape: Any) -> dict[str, Any]:
    element = getattr(shape, "_element", None)
    r_id = getattr(element, "blip_rId", None)
    if not r_id:
        return {}
    try:
        part = shape.part.related_part(r_id)
    except (KeyError, AttributeError, ValueError):
        return {}
    return {
        "rId": r_id,
        "content_type": str(getattr(part, "content_type", "") or ""),
        "partname": str(getattr(part, "partname", "") or ""),
    }


def _shape_description(shape: Any) -> str:
    try:
        c_nv_pr = shape._element.nvPicPr.cNvPr
        return str(c_nv_pr.get("descr") or "")
    except AttributeError:
        return ""


def _shape_bounds(shape: Any) -> dict[str, float]:
    return {
        "x": round(shape.left / EMU_PER_INCH, 4),
        "y": round(shape.top / EMU_PER_INCH, 4),
        "w": round(shape.width / EMU_PER_INCH, 4),
        "h": round(shape.height / EMU_PER_INCH, 4),
    }


def _max_icons_per_slide(style_tokens: dict[str, Any]) -> int:
    policy = style_tokens.get("decorative_budget") if isinstance(style_tokens.get("decorative_budget"), dict) else {}
    try:
        return max(1, int(policy.get("max_icons_per_slide") or DEFAULT_MAX_ICONS_PER_SLIDE))
    except (TypeError, ValueError):
        return DEFAULT_MAX_ICONS_PER_SLIDE


def _max_icon_area_ratio(style_tokens: dict[str, Any]) -> float:
    size_policy = style_tokens.get("size_policy") if isinstance(style_tokens.get("size_policy"), dict) else {}
    try:
        max_size = float(size_policy.get("maximum_template_icon_size_in") or 0.32)
    except (TypeError, ValueError):
        max_size = 0.32
    return max(DEFAULT_MAX_ICON_AREA_RATIO, round((max_size * max_size) / (13.333 * 7.5) * 9, 5))


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _dedupe(values: list[str]) -> list[str]:
    seen: set[str] = set()
    deduped: list[str] = []
    for value in values:
        if value not in seen:
            seen.add(value)
            deduped.append(value)
    return deduped


def _canonical_rel(path: Path) -> str:
    try:
        return str(path.relative_to(Path.cwd()).as_posix())
    except ValueError:
        return str(path.as_posix())


def _display_path(path: Path) -> str:
    try:
        return str(path.relative_to(Path.cwd()).as_posix())
    except ValueError:
        return str(path).replace("\\", "/")


def _finding(code: str, severity: str, message: str, details: dict[str, Any]) -> dict[str, Any]:
    return {
        "code": code,
        "severity": severity,
        "message": message,
        "details": details,
    }


def _markdown_report(report: dict[str, Any]) -> str:
    summary = report["findings_summary"]
    lines = [
        "# Icon Usage Policy Report",
        "",
        f"Status: `{report['status']}`",
        f"Findings: `{summary['total']}` total, `{summary['severe']}` severe, `{summary['warning']}` warnings",
        f"Normalized SVGs: `{report['normalized_svg_scan']['normalized_svg_count']}`",
        f"Missing role assets: `{len(report['role_coverage']['missing_roles'])}`",
        f"PNG icon files: `{report['asset_file_scan']['png_icon_file_count']}`",
        "",
        "## Findings By Code",
        "",
    ]
    if report.get("findings_by_code"):
        for code, count in report["findings_by_code"].items():
            lines.append(f"- `{code}`: `{count}`")
    else:
        lines.append("- none")
    lines.extend(["", "## Role Coverage", ""])
    for role in report["role_coverage"]["roles"]:
        lines.append(f"- `{role['role']}`: `{role['status']}`")
    lines.extend(["", "## PPTX Scans", ""])
    for scan in report.get("pptx_scans") or []:
        linkage = scan.get("svg_linkage_audit") if isinstance(scan.get("svg_linkage_audit"), dict) else {}
        lines.append(
            f"- `{scan['pptx_path']}`: `{scan['status']}`, slides `{scan['slide_count']}`, icons `{scan['icon_shape_count']}`, PNG icons `{scan['png_icon_asset_count']}`"
        )
        if linkage:
            lines.append(
                f"  - SVG linkage: media `{linkage.get('svg_media_count', 0)}`, referenced `{linkage.get('referenced_svg_media_count', 0)}`, "
                f"orphan `{linkage.get('orphan_svg_media_count', 0)}`, visible icon objects `{linkage.get('visible_icon_object_count', 0)}`"
            )
    return "\n".join(lines) + "\n"


def _write_run_002_linkage_audit_report(report: dict[str, Any]) -> list[Path]:
    run_pptx = "design_runs/run_002/outputs/run_002_editable_master_pack.pptx"
    audit: dict[str, Any] | None = None
    for scan in report.get("pptx_scans") or []:
        if scan.get("pptx_path") == run_pptx:
            candidate = scan.get("svg_linkage_audit")
            if isinstance(candidate, dict):
                audit = candidate
            break
    if audit is None:
        return []

    previous = _load_json(RUN_002_ICON_LINKAGE_REPORT_JSON)
    before_svg = previous.get("after_svg_media_count", audit.get("svg_media_count", 0))
    before_referenced = previous.get("after_referenced_svg_count", audit.get("referenced_svg_media_count", 0))
    payload = {
        "schema_name": "run_002_icon_linkage_audit_report",
        "schema_version": "1.0",
        "status": "passed" if not audit.get("remaining_icon_linkage_failures") else "failed",
        "pptx_path": run_pptx,
        "before_svg_media_count": before_svg,
        "before_referenced_svg_count": before_referenced,
        "after_svg_media_count": audit.get("svg_media_count", 0),
        "after_referenced_svg_count": audit.get("referenced_svg_media_count", 0),
        "orphan_svg_count": audit.get("orphan_svg_media_count", 0),
        "visible_icon_object_count": audit.get("visible_icon_object_count", 0),
        "svg_relationship_count": audit.get("svg_relationship_count", 0),
        "svg_blip_count": audit.get("svg_blip_count", 0),
        "slide_picture_object_count": audit.get("slide_picture_object_count", 0),
        "icon_roles_expected": audit.get("icon_roles_expected", []),
        "icon_roles_rendered": audit.get("icon_roles_rendered", []),
        "unresolved_icon_roles": audit.get("unresolved_icon_roles", []),
        "png_icon_count": audit.get("png_icon_count", 0),
        "raster_icon_count": audit.get("raster_icon_count", 0),
        "remaining_icon_linkage_failures": audit.get("remaining_icon_linkage_failures", []),
        "visible_icon_objects": audit.get("visible_icon_objects", []),
        "orphan_svg_media": audit.get("orphan_svg_media", []),
        "production_artifacts_modified": False,
    }
    RUN_002_ICON_LINKAGE_REPORT_JSON.parent.mkdir(parents=True, exist_ok=True)
    RUN_002_ICON_LINKAGE_REPORT_JSON.write_text(json.dumps(payload, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    RUN_002_ICON_LINKAGE_REPORT_MD.write_text(_run_002_linkage_markdown(payload), encoding="utf-8")
    return [RUN_002_ICON_LINKAGE_REPORT_JSON, RUN_002_ICON_LINKAGE_REPORT_MD]


def _run_002_linkage_markdown(report: dict[str, Any]) -> str:
    lines = [
        "# run_002 Icon Linkage Audit Report",
        "",
        f"- Status: `{report['status']}`",
        f"- PPTX: `{report['pptx_path']}`",
        f"- Before SVG media count: `{report['before_svg_media_count']}`",
        f"- Before referenced SVG count: `{report['before_referenced_svg_count']}`",
        f"- After SVG media count: `{report['after_svg_media_count']}`",
        f"- After referenced SVG count: `{report['after_referenced_svg_count']}`",
        f"- Orphan SVG count: `{report['orphan_svg_count']}`",
        f"- Visible icon object count: `{report['visible_icon_object_count']}`",
        f"- SVG relationship count: `{report['svg_relationship_count']}`",
        f"- SVG blip count: `{report['svg_blip_count']}`",
        f"- PNG icon count: `{report['png_icon_count']}`",
        f"- Raster icon count: `{report['raster_icon_count']}`",
        "",
        "## Roles",
        "",
        f"- Expected: `{', '.join(report['icon_roles_expected']) or 'none'}`",
        f"- Rendered: `{', '.join(report['icon_roles_rendered']) or 'none'}`",
        "",
        "## Remaining Linkage Failures",
        "",
    ]
    if report["remaining_icon_linkage_failures"]:
        for item in report["remaining_icon_linkage_failures"]:
            lines.append(f"- `{item.get('code')}`: `{item.get('count')}`")
    else:
        lines.append("- None")
    lines.extend(
        [
            "",
            "## Policy",
            "",
            "- SVG media must be referenced by slide, layout, or master relationships.",
            "- Expected compile-report icons must correspond to visible slide picture objects.",
            "- PNG icons and raster icon substitutes are forbidden.",
            "- Icons remain supporting visual components and cannot satisfy required content slots.",
            "",
        ]
    )
    return "\n".join(lines)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Validate SVG icon library usage in manifests, normalized assets, and PPTX outputs.")
    parser.add_argument("--manifests-dir", type=Path, default=DEFAULT_MANIFESTS_DIR)
    parser.add_argument("--normalized-dir", type=Path, default=DEFAULT_NORMALIZED_DIR)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    parser.add_argument("--pptx", type=Path, action="append", default=None, help="Additional or replacement PPTX path to scan. Can be supplied multiple times.")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    pptx_paths = tuple(args.pptx) if args.pptx else DEFAULT_PPTX_PATHS
    try:
        report = build_icon_usage_policy_report_from_files(
            manifests_dir=args.manifests_dir,
            normalized_dir=args.normalized_dir,
            pptx_paths=pptx_paths,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
        )
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"ICON_USAGE_POLICY_FAILED {exc}")
        return 1
    linkage_paths = _write_run_002_linkage_audit_report(report)
    print(f"WROTE {args.json_report}")
    print(f"WROTE {args.md_report}")
    for path in linkage_paths:
        print(f"WROTE {path}")
    print(f"ICON_USAGE_POLICY {report['status']}")
    return 1 if report.get("qa_blocks_template_usage") else 0


if __name__ == "__main__":
    raise SystemExit(main())
