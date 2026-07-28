"""Assemble non-canonical E03.1 16-archetype packs."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from .e03_16_orchestrator import ARCHETYPES, CORE_ARCHETYPES
from .e03_1_reference_fidelity_patch import draw_e03_1_archetype
from .e03_archetype_conversion import draw_e03_archetype


def assemble_e03_1_pack(visual_asset_plans: dict[str, dict[str, Any]], output_pptx: Path) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = 14630400
    prs.slide_height = 8229600
    for archetype in ARCHETYPES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if archetype in CORE_ARCHETYPES:
            draw_e03_archetype(slide, archetype, visual_asset_plans.get(archetype, {"assets": []}))
        else:
            draw_e03_1_archetype(slide, archetype, visual_asset_plans.get(archetype, {"assets": []}))
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_name": "e03_1_pack_report",
        "status": "passed",
        "pptx_path": output_pptx.as_posix(),
        "slide_count": len(ARCHETYPES),
        "archetypes": list(ARCHETYPES),
        "non_canonical": True,
        "canonical_promotion": False,
    }
