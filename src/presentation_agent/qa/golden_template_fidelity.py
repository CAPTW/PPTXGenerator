"""Compare production golden template masters against design-board references."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..compiler.template_compiler import (
    FORBIDDEN_GOLDEN_VISIBLE_STRINGS,
    compile_golden_template_masters_from_files,
)
from .render_pptx_preview import render_pptx_preview
from .template_crop_render_diff import (
    SLIDE_TYPE_BY_CROP_ROLE,
    _comparison_metrics,
    _display_path,
    _image_metrics,
    _match_score,
)


DEFAULT_CROP_MANIFEST = Path("outputs/template_design_board/design_board_crop_manifest.json")
DEFAULT_LAYOUT_REF_DIR = Path("outputs/template_design_board/layout_refs")
DEFAULT_TEMPLATE_SPEC = Path("outputs/editable_template_spec.final.json")
DEFAULT_GOLDEN_PPTX = Path("outputs/golden_template_masters.pptx")
DEFAULT_GOLDEN_REPORT = Path("outputs/golden_template_masters_report.json")
DEFAULT_RENDER_DIR = Path("outputs/golden_template_masters_png")
DEFAULT_RENDER_MANIFEST = Path("outputs/golden_template_masters_render_manifest.json")
DEFAULT_JSON_REPORT = Path("outputs/golden_template_fidelity_report.json")
DEFAULT_MD_REPORT = Path("outputs/golden_template_fidelity_report.md")


def build_golden_template_fidelity_from_files(
    *,
    crop_manifest_path: str | Path = DEFAULT_CROP_MANIFEST,
    layout_ref_dir: str | Path = DEFAULT_LAYOUT_REF_DIR,
    template_spec_path: str | Path = DEFAULT_TEMPLATE_SPEC,
    golden_pptx_path: str | Path = DEFAULT_GOLDEN_PPTX,
    golden_report_path: str | Path = DEFAULT_GOLDEN_REPORT,
    render_dir: str | Path = DEFAULT_RENDER_DIR,
    render_manifest_path: str | Path = DEFAULT_RENDER_MANIFEST,
    json_report_path: str | Path = DEFAULT_JSON_REPORT,
    md_report_path: str | Path = DEFAULT_MD_REPORT,
    renderer: str = "auto",
) -> dict[str, Any]:
    golden_pptx = Path(golden_pptx_path)
    golden_report_file = Path(golden_report_path)
    if not golden_pptx.exists() or not golden_report_file.exists():
        compile_golden_template_masters_from_files(
            spec_path=template_spec_path,
            output_path=golden_pptx,
            report_json_path=golden_report_file,
            report_md_path=Path("outputs/golden_template_masters_report.md"),
        )

    render_report = render_pptx_preview(
        pptx_path=golden_pptx,
        output_dir=render_dir,
        report_path=render_manifest_path,
        backend=renderer,
    )
    report = build_golden_template_fidelity_report(
        crop_manifest_path=crop_manifest_path,
        layout_ref_dir=layout_ref_dir,
        template_spec_path=template_spec_path,
        golden_pptx_path=golden_pptx,
        golden_report_path=golden_report_file,
        render_dir=render_dir,
        render_report=render_report,
    )
    json_path = Path(json_report_path)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(md_report_path).write_text(_markdown_report(report), encoding="utf-8")
    return report


def build_golden_template_fidelity_report(
    *,
    crop_manifest_path: str | Path,
    layout_ref_dir: str | Path,
    template_spec_path: str | Path,
    golden_pptx_path: str | Path,
    golden_report_path: str | Path,
    render_dir: str | Path,
    render_report: dict[str, Any],
) -> dict[str, Any]:
    crop_manifest = _load_json(crop_manifest_path)
    golden_report = _load_json(golden_report_path)
    spec = _load_json(template_spec_path)
    rendered_by_slide = _rendered_paths_by_slide(render_report, Path(render_dir))
    references = _references_by_archetype(crop_manifest, Path(layout_ref_dir))
    debug = _debug_text_scan(Path(golden_pptx_path))
    raster = _picture_policy_scan(Path(golden_pptx_path), spec)

    comparisons: list[dict[str, Any]] = []
    for record in golden_report.get("compiled_layouts") or []:
        if not isinstance(record, dict):
            continue
        archetype_id = str(record.get("archetype_id") or "")
        slide_number = int(record.get("slide_number") or 0)
        rendered_path = rendered_by_slide.get(slide_number, Path(render_dir) / f"slide-{slide_number:03d}.png")
        reference = references.get(archetype_id)
        comparisons.append(_compare_reference(archetype_id, record, reference, rendered_path))

    severe_findings: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    if debug["debug_text_violation_count"]:
        severe_findings.append(_finding("DEBUG_TEXT_VISIBLE", "Debug or internal slot text is visible in golden template masters.", "severe", debug))
    if raster["full_slide_picture_count"]:
        severe_findings.append(_finding("FULL_SLIDE_RASTER_BACKGROUND", "A full-slide raster picture exists in golden template masters.", "severe", raster))

    for comparison in comparisons:
        for finding in comparison.get("findings") or []:
            if finding.get("severity") == "severe":
                severe_findings.append(finding)
            else:
                warnings.append(finding)

    status = "failed" if severe_findings else "passed" if not warnings else "issues_reported"
    return {
        "schema_name": "golden_template_fidelity_report",
        "schema_version": "1.0",
        "status": status,
        "golden_pptx_path": _display_path(Path(golden_pptx_path)),
        "golden_report_path": _display_path(Path(golden_report_path)),
        "template_spec_path": _display_path(Path(template_spec_path)),
        "crop_manifest_path": _display_path(Path(crop_manifest_path)),
        "layout_ref_dir": _display_path(Path(layout_ref_dir)),
        "render_status": render_report.get("render_status"),
        "render_backend": render_report.get("backend"),
        "render_dir": _display_path(Path(render_dir)),
        "slide_count": len(comparisons),
        "debug_text_scan": debug,
        "picture_policy_scan": raster,
        "findings_summary": {
            "severe": len(severe_findings),
            "warning": len(warnings),
            "total": len(severe_findings) + len(warnings),
        },
        "severe_findings": severe_findings,
        "warnings": warnings,
        "layout_comparisons": comparisons,
        "reference_mode": "layout_refs_preferred_with_crop_fallback",
    }


def _compare_reference(archetype_id: str, record: dict[str, Any], reference: dict[str, Any] | None, rendered_path: Path) -> dict[str, Any]:
    findings: list[dict[str, Any]] = []
    if reference is None:
        findings.append(_finding("REFERENCE_MISSING", f"No board crop or full-size reference is available for {archetype_id}.", "warning"))
        return _comparison_shell(archetype_id, record, reference, rendered_path, None, 0.0, findings)
    ref_path = Path(reference["path"])
    if not ref_path.exists():
        findings.append(_finding("REFERENCE_IMAGE_MISSING", f"Reference image is missing for {archetype_id}.", "warning", {"path": reference["path"]}))
        return _comparison_shell(archetype_id, record, reference, rendered_path, None, 0.0, findings)
    if not rendered_path.exists():
        findings.append(_finding("RENDERED_MASTER_MISSING", f"Rendered golden master image is missing for {archetype_id}.", "severe", {"path": _display_path(rendered_path)}))
        return _comparison_shell(archetype_id, record, reference, rendered_path, None, 0.0, findings)

    metrics = _comparison_metrics(_image_metrics(ref_path), _image_metrics(rendered_path))
    score = _match_score(metrics)
    layout_family = str(record.get("layout_family_id") or "")
    if archetype_id == "creative_cover":
        rendered = metrics["rendered_metrics"]
        if rendered["dark_area_ratio"] < 0.45:
            findings.append(_finding("COVER_DARK_HERO_FIELD_MISSING", "Cover lacks the expected dark navy/teal hero field.", "severe", {"dark_area_ratio": rendered["dark_area_ratio"]}))
        if metrics["ornament_density_delta"] < -0.15:
            findings.append(_finding("COVER_TOPOLOGY_LAYER_WEAK", "Cover topology/ornament layer is too weak against the reference.", "severe", {"ornament_density_delta": metrics["ornament_density_delta"]}))
        if not _layout_has_image_slot(record):
            findings.append(_finding("COVER_PHOTO_MASK_MISSING", "Cover layout lacks a declared image/photo mask zone.", "severe"))
    if archetype_id == "section_divider" and not _layout_has_section_number(record):
        findings.append(_finding("SECTION_NUMBER_MISSING", "Section divider lacks oversized section number component.", "severe"))
    if layout_family in {"evidence_overview", "table_appendix", "kpi_dashboard"} and metrics["chart_table_module_presence_delta"] < -0.14:
        findings.append(_finding("DATA_LAYOUT_TOO_SPARSE", "Evidence/table/chart layout is too sparse against the reference.", "severe", {"chart_table_module_presence_delta": metrics["chart_table_module_presence_delta"]}))
    if metrics["palette_similarity"] < 0.68:
        findings.append(_finding("PALETTE_DIVERGENCE", "Golden master palette diverges from reference.", "warning", {"palette_similarity": metrics["palette_similarity"]}))
    if metrics["footer_occupancy_delta"] < -0.12:
        findings.append(_finding("FOOTER_OCCUPANCY_LOW", "Footer/citation occupancy is weaker than reference.", "warning", {"footer_occupancy_delta": metrics["footer_occupancy_delta"]}))
    return _comparison_shell(archetype_id, record, reference, rendered_path, metrics, score, findings)


def _comparison_shell(
    archetype_id: str,
    record: dict[str, Any],
    reference: dict[str, Any] | None,
    rendered_path: Path,
    metrics: dict[str, Any] | None,
    score: float,
    findings: list[dict[str, Any]],
) -> dict[str, Any]:
    return {
        "archetype_id": archetype_id,
        "layout_id": record.get("layout_id"),
        "layout_family_id": record.get("layout_family_id"),
        "slide_number": record.get("slide_number"),
        "reference_image_path": reference.get("path") if reference else None,
        "reference_source": reference.get("source") if reference else None,
        "rendered_image_path": _display_path(rendered_path),
        "similarity_score": round(score, 6),
        "metrics": metrics,
        "findings": findings,
    }


def _references_by_archetype(crop_manifest: dict[str, Any], layout_ref_dir: Path) -> dict[str, dict[str, Any]]:
    result: dict[str, dict[str, Any]] = {}
    for archetype in set(SLIDE_TYPE_BY_CROP_ROLE.values()):
        layout_ref = layout_ref_dir / f"{archetype}.png"
        if layout_ref.exists():
            result[archetype] = {"path": _display_path(layout_ref), "source": "layout_ref"}
    for crop in crop_manifest.get("crops") or []:
        if not isinstance(crop, dict):
            continue
        role = str(crop.get("crop_role") or crop.get("crop_id") or "")
        archetype = SLIDE_TYPE_BY_CROP_ROLE.get(role)
        if archetype and archetype not in result:
            result[archetype] = {"path": str(crop.get("path") or ""), "source": "design_board_crop", "crop_role": role}
    return result


def _debug_text_scan(pptx_path: Path) -> dict[str, Any]:
    hits: list[dict[str, Any]] = []
    presentation = Presentation(pptx_path)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = str(shape.text or "")
            for forbidden in FORBIDDEN_GOLDEN_VISIBLE_STRINGS:
                if _contains_forbidden(text, forbidden):
                    hits.append({"slide_index": slide_index, "forbidden_string": forbidden, "text": text[:160]})
    return {"debug_text_violation_count": len(hits), "hits": hits}


def _picture_policy_scan(pptx_path: Path, spec: dict[str, Any]) -> dict[str, Any]:
    presentation = Presentation(pptx_path)
    slide_w = float(spec.get("canvas", {}).get("width") or 13.333)
    slide_h = float(spec.get("canvas", {}).get("height") or 7.5)
    slide_area = slide_w * slide_h
    picture_count = 0
    full_slide_count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                area = (float(shape.width) / 914400.0) * (float(shape.height) / 914400.0)
                if area / max(0.01, slide_area) >= 0.88:
                    full_slide_count += 1
    return {"picture_shape_count": picture_count, "full_slide_picture_count": full_slide_count}


def _rendered_paths_by_slide(render_report: dict[str, Any], render_dir: Path) -> dict[int, Path]:
    result: dict[int, Path] = {}
    for record in render_report.get("slides") or []:
        if isinstance(record, dict) and record.get("rendered_image_path"):
            result[int(record.get("slide_index") or len(result) + 1)] = Path(str(record["rendered_image_path"]))
    for index, path in enumerate(render_report.get("output_paths") or [], start=1):
        result.setdefault(index, Path(str(path)))
    return result


def _layout_has_image_slot(record: dict[str, Any]) -> bool:
    layout_id = str(record.get("layout_id") or "")
    return layout_id.endswith("creative_cover")


def _layout_has_section_number(record: dict[str, Any]) -> bool:
    layout_id = str(record.get("layout_id") or "")
    return "section_divider" in layout_id


def _contains_forbidden(text: str, forbidden: str) -> bool:
    if forbidden in {"cards"}:
        return any(part.strip(" .,:;|/\\").lower() == forbidden for part in text.split())
    return forbidden.lower() in text.lower()


def _finding(code: str, message: str, severity: str, details: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = {"code": code, "severity": severity, "message": message}
    if details is not None:
        payload["details"] = details
    return payload


def _markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Golden Template Fidelity Report",
        "",
        f"Status: `{report['status']}`",
        f"Render: `{report.get('render_status')}` via `{report.get('render_backend')}`",
        f"Slide count: `{report.get('slide_count')}`",
        f"Debug text violations: `{report.get('debug_text_scan', {}).get('debug_text_violation_count')}`",
        f"Findings: `{report['findings_summary']['total']}` total, `{report['findings_summary']['severe']}` severe",
        "",
        "This QA compares production golden template masters against full-size layout references when present, otherwise against design-board crops. It also blocks debug text and full-slide raster backgrounds.",
        "",
        "## Layout Comparisons",
        "",
        "| Archetype | Layout | Reference | Score | Severe | Warnings |",
        "|---|---|---|---:|---:|---:|",
    ]
    for item in report.get("layout_comparisons") or []:
        severe = sum(1 for f in item.get("findings") or [] if f.get("severity") == "severe")
        warning = sum(1 for f in item.get("findings") or [] if f.get("severity") != "severe")
        lines.append(
            f"| `{item.get('archetype_id')}` | `{item.get('layout_id')}` | `{item.get('reference_source')}` | {item.get('similarity_score', 0):.3f} | {severe} | {warning} |"
        )
    lines.extend(["", "## Severe Findings", ""])
    if report.get("severe_findings"):
        for finding in report["severe_findings"]:
            lines.append(f"- `{finding['code']}`: {finding['message']}")
    else:
        lines.append("- None")
    lines.extend(["", "## Warnings", ""])
    if report.get("warnings"):
        for finding in report["warnings"]:
            lines.append(f"- `{finding['code']}`: {finding['message']}")
    else:
        lines.append("- None")
    return "\n".join(lines) + "\n"


def _load_json(path: str | Path) -> Any:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Compare golden template masters against design-board references.")
    parser.add_argument("--crop-manifest", type=Path, default=DEFAULT_CROP_MANIFEST)
    parser.add_argument("--layout-ref-dir", type=Path, default=DEFAULT_LAYOUT_REF_DIR)
    parser.add_argument("--template-spec", type=Path, default=DEFAULT_TEMPLATE_SPEC)
    parser.add_argument("--pptx", type=Path, default=DEFAULT_GOLDEN_PPTX)
    parser.add_argument("--golden-report", type=Path, default=DEFAULT_GOLDEN_REPORT)
    parser.add_argument("--render-dir", type=Path, default=DEFAULT_RENDER_DIR)
    parser.add_argument("--render-manifest", type=Path, default=DEFAULT_RENDER_MANIFEST)
    parser.add_argument("--json-report", type=Path, default=DEFAULT_JSON_REPORT)
    parser.add_argument("--md-report", type=Path, default=DEFAULT_MD_REPORT)
    parser.add_argument("--renderer", choices=("auto", "powerpoint_com", "libreoffice", "none"), default="auto")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        report = build_golden_template_fidelity_from_files(
            crop_manifest_path=args.crop_manifest,
            layout_ref_dir=args.layout_ref_dir,
            template_spec_path=args.template_spec,
            golden_pptx_path=args.pptx,
            golden_report_path=args.golden_report,
            render_dir=args.render_dir,
            render_manifest_path=args.render_manifest,
            json_report_path=args.json_report,
            md_report_path=args.md_report,
            renderer=args.renderer,
        )
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        print(f"GOLDEN_TEMPLATE_FIDELITY_FAILED {exc}")
        return 1
    print(f"WROTE {DEFAULT_JSON_REPORT}")
    print(f"GOLDEN_TEMPLATE_FIDELITY {report['status']}")
    return 1 if report["findings_summary"]["severe"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
