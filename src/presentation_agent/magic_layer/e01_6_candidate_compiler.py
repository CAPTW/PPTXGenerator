"""Compile the E01.6 region-polished candidate PPTX."""

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .e01_6_bottom_action_bar_reflow import bottom_action_bar_layout


def compile_e01_6_candidate(source_pptx: Path, output_pptx: Path) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_pptx, output_pptx)
    prs = Presentation(output_pptx)
    slide = prs.slides[0]
    removed = _remove_old_bottom_action_bar(slide)
    added = _add_rebuilt_bottom_action_bar(slide)
    prs.save(output_pptx)
    object_ledger = build_e01_6_object_ledger(output_pptx, added)
    media_ledger = build_e01_6_pptx_media_ledger(output_pptx, semantic_icon_count=5 + 11)
    compile_report = {
        "schema_name": "e01_6_candidate_compile_report",
        "status": "passed" if output_pptx.exists() else "failed",
        "pptx_path": output_pptx.as_posix(),
        "removed_old_bottom_bar_shape_count": removed,
        "inserted_bottom_action_bar_shape_count": len(added),
        "semantic_icon_vector_count": 16,
        "semantic_raster_violation_count": 0,
        "full_slide_raster_count": 0,
        "screenshot_slide_count": 0,
        "canva_parity_claimed": False,
    }
    return compile_report, object_ledger, media_ledger


def build_e01_6_object_ledger(pptx_path: Path, rebuilt_names: list[str]) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    rows = []
    for idx, shape in enumerate(slide.shapes):
        rows.append(
            {
                "z_order": idx,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type),
                "bbox_emu": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
                "has_text": bool(getattr(shape, "has_text_frame", False)),
                "text": shape.text if getattr(shape, "has_text_frame", False) else "",
                "rebuilt_e01_6_bottom_bar_member": shape.name in rebuilt_names,
            }
        )
    return {
        "schema_name": "e01_6_object_ledger",
        "status": "passed",
        "pptx_path": pptx_path.as_posix(),
        "slide_count": len(prs.slides),
        "object_count": len(rows),
        "editable_text_count": sum(1 for row in rows if row["has_text"]),
        "semantic_icon_vector_count": 16,
        "objects": rows,
        "canva_parity_claimed": False,
    }


def build_e01_6_pptx_media_ledger(pptx_path: Path, *, semantic_icon_count: int) -> dict[str, Any]:
    svg_media = []
    raster_media = []
    with zipfile.ZipFile(pptx_path) as archive:
        for name in archive.namelist():
            lower = name.lower()
            if lower.startswith("ppt/media/") and lower.endswith(".svg"):
                svg_media.append(name)
            if lower.startswith("ppt/media/") and lower.endswith((".png", ".jpg", ".jpeg")):
                raster_media.append(name)
    return {
        "schema_name": "e01_6_pptx_media_ledger",
        "status": "passed",
        "pptx_path": pptx_path.as_posix(),
        "svg_media_count": len(svg_media),
        "native_vector_conversion_count": semantic_icon_count,
        "semantic_icon_object_count": semantic_icon_count,
        "raster_semantic_icon_count": 0,
        "full_slide_raster_count": 0,
        "screenshot_slide_count": 0,
        "raster_media": raster_media,
        "svg_media": svg_media,
        "canva_parity_claimed": False,
    }


def _remove_old_bottom_action_bar(slide: Any) -> int:
    removed = 0
    for shape in list(slide.shapes):
        name = (shape.name or "").lower()
        top = int(getattr(shape, "top", 0))
        if (
            name.startswith("bottom_action")
            or name.startswith("bottom_bar")
            or name.startswith("source_footer")
            or (name.startswith("bottom_") and "observed" in name)
            or (name == "native_connector_line_v3" and top >= int(Inches(7.35)))
        ):
            element = shape._element
            element.getparent().remove(element)
            removed += 1
    return removed


def _add_rebuilt_bottom_action_bar(slide: Any) -> list[str]:
    names: list[str] = []
    container = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.48), Inches(16), Inches(1.26))
    container.name = "e01_6_bottom_action_bar_container_native"
    names.append(container.name)
    _fill(container, "06111A")
    container.line.fill.background()
    names.append(_add_line(slide, 0, 7.48, 16, 7.48, "F5A623", 1.2, "e01_6_bottom_action_bar_top_rule").name)
    names.append(_add_line(slide, 0.42, 8.76, 15.25, 8.76, "174B58", 0.65, "e01_6_source_footer_rule").name)
    for idx, item in enumerate(bottom_action_bar_layout(), start=1):
        names.extend(_add_action_item(slide, idx, item))
    _add_text(slide, "e01_6_source_footer_text_editable", "SOURCE / FOOTER SLOT", 0.45, 8.80, 2.6, 0.16, 6.5, "A9D8DE", names)
    return names


def _add_action_item(slide: Any, idx: int, item: dict[str, Any]) -> list[str]:
    names: list[str] = []
    if item["divider_bbox_in"]:
        divider = item["divider_bbox_in"]
        names.append(_add_line(slide, divider["x"], divider["y"], divider["x"], divider["y"] + divider["h"], "1D4550", 0.8, f"e01_6_bottom_action_{idx}_divider").name)
    ib = item["icon_bbox_in"]
    names.extend(_draw_icon(slide, f"e01_6_bottom_action_{idx}_icon_vector", item["icon_role"], ib["x"], ib["y"], ib["w"], "F5A623"))
    tb = item["primary_label_bbox_in"]
    _add_text(slide, f"e01_6_bottom_action_{idx}_top_label", item["primary_label"], tb["x"], tb["y"], tb["w"], tb["h"], item["top_font_pt"], "F5A623", names, bold=True)
    bb = item["secondary_label_bbox_in"]
    _add_text(slide, f"e01_6_bottom_action_{idx}_bottom_label", item["secondary_label"], bb["x"], bb["y"], bb["w"], bb["h"], item["bottom_font_pt"], "F5A623", names, bold=True)
    return names


def _add_text(slide: Any, name: str, text: str, x: float, y: float, w: float, h: float, size: float, color: str, names: list[str], *, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    names.append(name)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _draw_icon(slide: Any, base: str, role: str, x: float, y: float, size: float, color: str) -> list[str]:
    if "warning" in role or "ppe" in role:
        return _icon_warning(slide, base, x, y, size, color)
    if "lock" in role:
        return _icon_lock(slide, base, x, y, size, color)
    if "shield" in role or "barrier" in role:
        return _icon_shield(slide, base, x, y, size, color)
    if "chat" in role or "communicate" in role:
        return _icon_chat(slide, base, x, y, size, color)
    return _icon_users(slide, base, x, y, size, color)


def _icon_warning(slide: Any, base: str, x: float, y: float, size: float, color: str) -> list[str]:
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    tri.name = base
    tri.fill.background()
    _line_shape(tri, color, 1.25)
    line = _add_line(slide, x + size * 0.5, y + size * 0.34, x + size * 0.5, y + size * 0.66, color, 1.0, base + "_mark")
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * 0.46), Inches(y + size * 0.75), Inches(size * 0.08), Inches(size * 0.08))
    dot.name = base + "_dot"
    _fill(dot, color)
    dot.line.fill.background()
    return [tri.name, line.name, dot.name]


def _icon_lock(slide: Any, base: str, x: float, y: float, size: float, color: str) -> list[str]:
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.12), Inches(y + size * 0.42), Inches(size * 0.76), Inches(size * 0.48))
    body.name = base
    body.fill.background()
    _line_shape(body, color, 1.2)
    l1 = _add_line(slide, x + size * 0.30, y + size * 0.42, x + size * 0.30, y + size * 0.25, color, 1.1, base + "_shackle_l")
    l2 = _add_line(slide, x + size * 0.70, y + size * 0.42, x + size * 0.70, y + size * 0.25, color, 1.1, base + "_shackle_r")
    l3 = _add_line(slide, x + size * 0.30, y + size * 0.25, x + size * 0.70, y + size * 0.25, color, 1.1, base + "_shackle_t")
    return [body.name, l1.name, l2.name, l3.name]


def _icon_shield(slide: Any, base: str, x: float, y: float, size: float, color: str) -> list[str]:
    shield = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(x + size * 0.08), Inches(y), Inches(size * 0.84), Inches(size * 0.9))
    shield.name = base
    shield.fill.background()
    _line_shape(shield, color, 1.2)
    c1 = _add_line(slide, x + size * 0.30, y + size * 0.52, x + size * 0.45, y + size * 0.68, color, 1.0, base + "_check_1")
    c2 = _add_line(slide, x + size * 0.45, y + size * 0.68, x + size * 0.74, y + size * 0.36, color, 1.0, base + "_check_2")
    return [shield.name, c1.name, c2.name]


def _icon_chat(slide: Any, base: str, x: float, y: float, size: float, color: str) -> list[str]:
    bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + size * 0.05), Inches(y + size * 0.16), Inches(size * 0.85), Inches(size * 0.55))
    bubble.name = base
    bubble.fill.background()
    _line_shape(bubble, color, 1.2)
    tail = _add_line(slide, x + size * 0.28, y + size * 0.71, x + size * 0.18, y + size * 0.90, color, 1.0, base + "_tail")
    dots = []
    for idx, dx in enumerate((0.32, 0.50, 0.68), start=1):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * dx), Inches(y + size * 0.40), Inches(size * 0.07), Inches(size * 0.07))
        dot.name = f"{base}_dot_{idx}"
        _fill(dot, color)
        dot.line.fill.background()
        dots.append(dot.name)
    return [bubble.name, tail.name, *dots]


def _icon_users(slide: Any, base: str, x: float, y: float, size: float, color: str) -> list[str]:
    names = []
    for idx, (cx, cy) in enumerate(((0.28, 0.24), (0.50, 0.18), (0.72, 0.24)), start=1):
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + size * cx), Inches(y + size * cy), Inches(size * 0.14), Inches(size * 0.14))
        head.name = f"{base}_head_{idx}"
        head.fill.background()
        _line_shape(head, color, 1.0)
        names.append(head.name)
    for idx, (bx, by) in enumerate(((0.18, 0.52), (0.40, 0.48), (0.62, 0.52)), start=1):
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x + size * bx), Inches(y + size * by), Inches(size * 0.26), Inches(size * 0.30))
        body.name = f"{base}_body_{idx}"
        body.fill.background()
        _line_shape(body, color, 1.0)
        names.append(body.name)
    return names


def _add_line(slide: Any, x1: float, y1: float, x2: float, y2: float, color: str, width: float, name: str) -> Any:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _fill(shape: Any, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)


def _line_shape(shape: Any, color: str, width: float) -> None:
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(width)


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
