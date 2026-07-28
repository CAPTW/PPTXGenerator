"""Compile accepted E03H-V2 references into an editable pack PPTX."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from src.presentation_agent.magic_layer.e01h_v2_r1_report import write_json


def compile_e03h_v2_reference_pack(case_results: list[dict[str, Any]], output_dir: str | Path) -> dict[str, Any]:
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    accepted = [case for case in case_results if case.get("case_gate", {}).get("status") == "passed"]
    pptx_path = output / "editable_hybrid_reference_pack_v2.pptx"
    contact_sheet = output / "editable_hybrid_reference_pack_v2_contact_sheet.png"
    _build_pack_pptx(accepted, pptx_path)
    _contact_sheet([_render_path(case) for case in accepted], contact_sheet, "E03H-V2 editable hybrid reference pack")
    manifest = {
        "schema_name": "editable_hybrid_reference_pack_v2_render_manifest",
        "status": "passed" if accepted else "failed",
        "slide_count": len(accepted),
        "accepted_reference_ids": [case["reference_id"] for case in accepted],
        "pptx_path": pptx_path.as_posix(),
        "contact_sheet": contact_sheet.as_posix(),
        "full_slide_reference_background_count": 0,
        "screenshot_slide_count": 0,
        "semantic_raster_violation_count": 0,
        "canva_parity_claimed": False,
    }
    write_json(output / "editable_hybrid_reference_pack_v2_render_manifest.json", manifest)
    return {"status": manifest["status"], "pack_pptx": pptx_path.as_posix(), "contact_sheet": contact_sheet.as_posix(), "slide_count": len(accepted), "manifest": manifest}


def _render_path(case: dict[str, Any]) -> Path:
    if case.get("rendered_candidate"):
        return Path(case["rendered_candidate"])
    return Path(case["reference_dir"]) / "rendered_candidate.png"


def _build_pack_pptx(cases: list[dict[str, Any]], pptx_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for case in cases:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        dark = "dark" in case["reference_id"] or case["reference_id"] in {"maritime_checklist_hero", "process_workflow_infographic", "methodology_framework_layered"}
        bg = RGBColor(8, 27, 39) if dark else RGBColor(242, 238, 228)
        accent = RGBColor(42, 202, 218) if dark else RGBColor(36, 122, 148)
        title = RGBColor(255, 255, 255) if dark else RGBColor(31, 45, 55)
        _rect(slide, 0, 0, 13.333, 7.5, bg, None, "native_background_substrate")
        _text(slide, case["reference_id"].replace("_", " ").title(), 0.62, 0.42, 8.8, 0.36, 18, True, "semantic_text::title", title)
        _rect(slide, 0.82, 2.1, 3.25, 1.1, RGBColor(255, 252, 244) if not dark else RGBColor(11, 55, 68), accent, "native_card_panel::primary")
        _rect(slide, 4.35, 2.1, 3.25, 1.1, RGBColor(255, 252, 244) if not dark else RGBColor(11, 55, 68), accent, "native_card_panel::secondary")
        _rect(slide, 8.75, 1.78, 2.6, 3.1, RGBColor(236, 241, 239) if not dark else RGBColor(13, 43, 54), accent, "segmented_visual_field::bounded_nonsemantic")
        _text(slide, "Primary content", 1.0, 2.42, 2.4, 0.22, 8, False, "semantic_text::card_primary", title)
        _text(slide, "Bounded visual field", 4.55, 2.42, 2.4, 0.22, 8, False, "semantic_text::card_secondary", title)
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(10.76), Inches(0.68), Inches(0.18), Inches(0.18))
        icon.name = "sem_icon::generic_check::svg01_source_generic_check"
        icon.fill.solid()
        icon.fill.fore_color.rgb = accent
        icon.line.fill.background()
    prs.save(pptx_path)


def _rect(slide: Any, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None, name: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()


def _text(slide: Any, text: str, x: float, y: float, w: float, h: float, size: float, bold: bool, name: str, color: RGBColor) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _contact_sheet(images: list[Path], output: Path, title: str) -> None:
    cols = 4
    rows = max(1, (len(images) + cols - 1) // cols)
    canvas = Image.new("RGB", (cols * 390 + 20, rows * 245 + 54), (246, 247, 244))
    draw = ImageDraw.Draw(canvas)
    draw.text((18, 18), title, fill=(32, 45, 55), font=_font(22))
    for idx, path in enumerate(images):
        thumb = _load_thumb(path, (360, 203))
        x = 18 + (idx % cols) * 390
        y = 54 + (idx // cols) * 245
        canvas.paste(thumb, (x, y))
    output.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(output)


def _load_thumb(path: Path, size: tuple[int, int]) -> Image.Image:
    if not path.exists():
        return Image.new("RGB", size, (230, 230, 230))
    image = Image.open(path).convert("RGB")
    image.thumbnail(size)
    canvas = Image.new("RGB", size, (250, 250, 248))
    canvas.paste(image, ((size[0] - image.width) // 2, (size[1] - image.height) // 2))
    return canvas


def _font(size: int) -> ImageFont.ImageFont:
    for font_name in ("arial.ttf", "calibri.ttf", "segoeui.ttf"):
        try:
            return ImageFont.truetype(font_name, size)
        except OSError:
            continue
    return ImageFont.load_default()
