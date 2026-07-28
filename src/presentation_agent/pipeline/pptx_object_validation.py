"""Deterministic PPTX object-tree validation and SceneDeck comparison."""

from __future__ import annotations

from dataclasses import dataclass
import hashlib
import json
import math
import zipfile
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, ConfigDict, Field

from ..slide_scene import SceneDeck, TextBox
from ..pptx_renderers.common import scene_text_role_defaults
from .render_validation import validate_local_pptx


OBJECT_VALIDATION_REPORT_VERSION = "0.2"
ValidationMode = Literal["inspect", "warn", "enforce"]
ValidationProfile = Literal["basic", "scene-strict"]
FindingSeverity = Literal["info", "warning", "error"]
FindingCategory = Literal["semantic_mismatch", "editability", "native_object", "unsupported_shape", "text_fit"]
SemanticType = Literal[
    "text_object",
    "image_object",
    "native_table",
    "native_chart",
    "group",
    "placeholder",
    "auto_shape",
    "other",
]
FindingCode = Literal[
    "expected_text_object_missing",
    "expected_image_object_missing",
    "expected_native_table_missing",
    "expected_native_chart_missing",
    "expected_shape_missing",
    "expected_divider_missing",
    "expected_callout_missing",
    "expected_background_motif_missing",
    "picture_used_for_expected_table",
    "picture_used_for_expected_chart",
    "native_table_dimension_mismatch",
    "native_chart_series_count_mismatch",
    "native_chart_category_count_mismatch",
    "native_chart_type_mismatch",
    "unsupported_background_motif",
    "background_motif_trace_duplicate",
    "unsupported_shape_type",
    "duplicate_shape_name",
    "scene_slide_missing_in_pptx",
    "pptx_extra_slide_without_scene",
    "ambiguous_text_mapping",
    "image_object_missing_trace",
    "native_table_missing_trace",
    "native_chart_missing_trace",
    "scene_trace_duplicate",
    "scene_trace_type_mismatch",
    "text_box_empty",
    "text_box_overflow_risk",
    "text_box_line_count_exceeds_bounds",
    "text_box_font_below_min_fit_policy",
    "text_box_fit_policy_unsupported",
    "text_box_truncated_without_policy",
    "text_box_bullet_level_unsupported",
    "text_box_missing_trace",
]
TextFitDiagnosticCode = Literal[
    "text_box_empty",
    "text_box_overflow_risk",
    "text_box_line_count_exceeds_bounds",
    "text_box_font_below_min_fit_policy",
    "text_box_fit_policy_unsupported",
    "text_box_truncated_without_policy",
    "text_box_bullet_level_unsupported",
    "text_box_missing_trace",
]

_SEMANTIC_TYPE_ORDER: tuple[SemanticType, ...] = (
    "text_object",
    "image_object",
    "native_table",
    "native_chart",
    "group",
    "placeholder",
    "auto_shape",
    "other",
)
_SUPPORTED_BACKGROUND_SHAPE_TYPES = {"rect", "rectangle", "rounded_rect", "rounded-rect", "ellipse", "oval", "line"}
_SUPPORTED_TEXT_FIT_MODES = {"none", "wrap", "fail"}
_BULLET_LEVEL_HARD_LIMIT = 8
_SCENE_TRACE_DUPLICATE_CODES = {"scene_trace_duplicate", "background_motif_trace_duplicate"}
_SCENE_TRACE_MISSING_CODES = {
    "text_box_missing_trace",
    "image_object_missing_trace",
    "native_table_missing_trace",
    "native_chart_missing_trace",
    "expected_shape_missing",
    "expected_divider_missing",
    "expected_callout_missing",
    "expected_background_motif_missing",
}


@dataclass(frozen=True)
class _InspectedShapeDetail:
    name: str
    trace_kind: str | None
    trace_object_id: str | None
    trace_part: str | None
    semantic_types: tuple[SemanticType, ...]
    normalized_text: str | None
    paragraph_levels: tuple[int, ...]
    min_font_size_pt: float | None
    max_font_size_pt: float | None
    left_in: float | None
    top_in: float | None
    width_in: float | None
    height_in: float | None


class ValidationModel(BaseModel):
    model_config = ConfigDict(extra="forbid", populate_by_name=True)


class SemanticObjectCounts(ValidationModel):
    shape_count: int = 0
    text_object_count: int = 0
    image_object_count: int = 0
    native_table_count: int = 0
    native_chart_count: int = 0
    group_count: int = 0
    placeholder_count: int = 0
    auto_shape_count: int = 0
    other_shape_count: int = 0

    def increment(self, semantic_type: SemanticType) -> None:
        if semantic_type == "text_object":
            self.text_object_count += 1
        elif semantic_type == "image_object":
            self.image_object_count += 1
        elif semantic_type == "native_table":
            self.native_table_count += 1
        elif semantic_type == "native_chart":
            self.native_chart_count += 1
        elif semantic_type == "group":
            self.group_count += 1
        elif semantic_type == "placeholder":
            self.placeholder_count += 1
        elif semantic_type == "auto_shape":
            self.auto_shape_count += 1
        elif semantic_type == "other":
            self.other_shape_count += 1


class PptxShapeSummary(ValidationModel):
    shape_index: int
    shape_id: int | None = None
    name: str
    shape_type: str
    trace_kind: str | None = None
    trace_object_id: str | None = None
    trace_part: str | None = None
    semantic_types: list[SemanticType] = Field(default_factory=list)
    has_text: bool = False
    text_excerpt: str | None = None
    text_length: int | None = None
    paragraph_count: int | None = None
    max_paragraph_level: int | None = None
    min_font_size_pt: float | None = None
    max_font_size_pt: float | None = None
    left_in: float | None = None
    top_in: float | None = None
    width_in: float | None = None
    height_in: float | None = None
    child_count: int | None = None
    table_row_count: int | None = None
    table_column_count: int | None = None
    chart_type: str | None = None
    chart_series_count: int | None = None
    chart_category_count: int | None = None


class SceneExpectationSummary(ValidationModel):
    comparison_basis: Literal["none", "ordinal"] = "none"
    scene_slide_id: str | None = None
    layout_family: str | None = None
    expected_counts: SemanticObjectCounts = Field(default_factory=SemanticObjectCounts)
    actual_counts: SemanticObjectCounts = Field(default_factory=SemanticObjectCounts)


class PptxObjectFinding(ValidationModel):
    code: FindingCode
    severity: FindingSeverity
    category: FindingCategory
    message: str
    slide_number: int | None = None
    scene_slide_id: str | None = None
    shape_name: str | None = None
    details: dict[str, Any] = Field(default_factory=dict)


class PptxTextObjectDiagnostic(ValidationModel):
    code: TextFitDiagnosticCode
    severity: FindingSeverity
    slide_number: int
    scene_slide_id: str | None = None
    object_id: str
    shape_name: str | None = None
    details: dict[str, Any] = Field(default_factory=dict)


class PptxSlideTextFitSummary(ValidationModel):
    diagnostic_count: int = 0
    overflow_risk_count: int = 0
    unsupported_fit_policy_count: int = 0
    missing_text_trace_count: int = 0
    diagnostics: list[PptxTextObjectDiagnostic] = Field(default_factory=list)


class PptxSlideObjectInventory(ValidationModel):
    slide_number: int
    pptx_slide_id: int | None = None
    title_excerpt: str | None = None
    has_notes_slide: bool = False
    chart_relationship_count: int = 0
    counts: SemanticObjectCounts = Field(default_factory=SemanticObjectCounts)
    shapes: list[PptxShapeSummary] = Field(default_factory=list)


class PptxSlideValidationResult(ValidationModel):
    slide_number: int
    pptx_slide_id: int | None = None
    inventory: PptxSlideObjectInventory
    scene_comparison: SceneExpectationSummary | None = None
    text_fit_summary: PptxSlideTextFitSummary = Field(default_factory=PptxSlideTextFitSummary)
    findings: list[PptxObjectFinding] = Field(default_factory=list)


class PptxObjectFindingsSummary(ValidationModel):
    total_findings: int
    info_count: int
    warning_count: int
    error_count: int
    enforceable_count: int
    semantic_mismatch_count: int
    editability_count: int
    native_object_count: int
    unsupported_shape_count: int
    text_fit_count: int


class PptxDeckSummary(ValidationModel):
    file_size_bytes: int
    checksum: str
    slide_count: int
    chart_part_count: int
    zip_readable: bool
    presentation_xml_present: bool
    scene_slide_count: int | None = None


class PptxObjectValidationReport(ValidationModel):
    report_version: str = OBJECT_VALIDATION_REPORT_VERSION
    mode: ValidationMode
    profile: ValidationProfile = "basic"
    mode_result: Literal["passed", "issues_reported", "failed"]
    pptx_path: str
    scene_deck_path: str | None = None
    structural_hash: str
    deck_summary: PptxDeckSummary
    slide_count: int
    findings_summary: PptxObjectFindingsSummary
    overflow_risk_count: int = 0
    unsupported_fit_policy_count: int = 0
    missing_text_trace_count: int = 0
    missing_trace_count: int = 0
    duplicate_trace_count: int = 0
    slides: list[PptxSlideValidationResult] = Field(default_factory=list)
    findings: list[PptxObjectFinding] = Field(default_factory=list)

    def to_stable_payload(self) -> dict[str, Any]:
        return validation_report_to_stable_payload(self)

    def to_stable_json(self) -> str:
        return validation_report_to_stable_json(self)


def validate_pptx_objects(
    pptx_path: str | Path,
    *,
    scene_deck: SceneDeck | None = None,
    scene_deck_path: str | Path | None = None,
    mode: ValidationMode = "inspect",
    profile: ValidationProfile = "basic",
) -> PptxObjectValidationReport:
    pptx_file = Path(pptx_path).resolve()
    pptx_summary = validate_local_pptx(pptx_file)
    chart_part_count = _chart_part_count(pptx_file)
    presentation = Presentation(pptx_file)
    scene_slides = list(scene_deck.slides) if scene_deck is not None else []

    slide_results: list[PptxSlideValidationResult] = []
    deck_findings: list[PptxObjectFinding] = []
    deck_text_diagnostics: list[PptxTextObjectDiagnostic] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        inventory, inventory_findings, shape_details = _inspect_slide(slide_number, slide)
        scene_comparison = None
        comparison_findings: list[PptxObjectFinding] = []
        text_fit_summary = PptxSlideTextFitSummary()
        if scene_deck is not None and slide_number <= len(scene_slides):
            scene_slide = scene_slides[slide_number - 1]
            scene_comparison, comparison_findings, text_fit_summary = _compare_slide_to_scene(
                slide_number,
                scene_slide,
                inventory,
                shape_details,
                profile=profile,
            )
            deck_text_diagnostics.extend(text_fit_summary.diagnostics)
        slide_findings = _sorted_findings([*inventory_findings, *comparison_findings])
        slide_results.append(
            PptxSlideValidationResult(
                slide_number=slide_number,
                pptx_slide_id=getattr(slide, "slide_id", None),
                inventory=inventory,
                scene_comparison=scene_comparison,
                text_fit_summary=text_fit_summary,
                findings=slide_findings,
            )
        )
        deck_findings.extend(slide_findings)

    if scene_deck is not None and len(scene_slides) > len(slide_results):
        for scene_slide in scene_slides[len(slide_results) :]:
            deck_findings.append(
                PptxObjectFinding(
                    code="scene_slide_missing_in_pptx",
                    severity="error",
                    category="semantic_mismatch",
                    slide_number=scene_slide.slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    message=f"SceneDeck slide {scene_slide.slide_id} has no corresponding PPTX slide.",
                    details={"expected_slide_number": scene_slide.slide_number},
                )
            )
    if scene_deck is not None and len(slide_results) > len(scene_slides):
        for slide_result in slide_results[len(scene_slides) :]:
            deck_findings.append(
                PptxObjectFinding(
                    code="pptx_extra_slide_without_scene",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_result.slide_number,
                    message="PPTX contains an extra slide with no corresponding SceneDeck slide.",
                    details={"pptx_slide_number": slide_result.slide_number},
                )
            )

    findings = _sorted_findings(deck_findings)
    findings_summary = _summarize_findings(findings)
    overflow_risk_count = sum(1 for item in deck_text_diagnostics if item.code == "text_box_overflow_risk")
    unsupported_fit_policy_count = sum(1 for item in deck_text_diagnostics if item.code == "text_box_fit_policy_unsupported")
    missing_text_trace_count = sum(1 for item in deck_text_diagnostics if item.code == "text_box_missing_trace")
    missing_trace_count = sum(1 for finding in findings if finding.code in _SCENE_TRACE_MISSING_CODES)
    duplicate_trace_count = sum(1 for finding in findings if finding.code in _SCENE_TRACE_DUPLICATE_CODES)
    report = PptxObjectValidationReport(
        mode=mode,
        profile=profile,
        mode_result=_mode_result(mode, findings_summary),
        pptx_path=str(pptx_file),
        scene_deck_path=str(Path(scene_deck_path).resolve()) if scene_deck_path is not None else None,
        structural_hash="",
        deck_summary=PptxDeckSummary(
            file_size_bytes=int(pptx_summary["file_size_bytes"]),
            checksum=str(pptx_summary["checksum"]),
            slide_count=int(pptx_summary["slide_count"]),
            chart_part_count=chart_part_count,
            zip_readable=bool(pptx_summary["zip_readable"]),
            presentation_xml_present=bool(pptx_summary["presentation_xml_present"]),
            scene_slide_count=len(scene_slides) if scene_deck is not None else None,
        ),
        slide_count=len(slide_results),
        findings_summary=findings_summary,
        overflow_risk_count=overflow_risk_count,
        unsupported_fit_policy_count=unsupported_fit_policy_count,
        missing_text_trace_count=missing_text_trace_count,
        missing_trace_count=missing_trace_count,
        duplicate_trace_count=duplicate_trace_count,
        slides=slide_results,
        findings=findings,
    )
    structural_hash = validation_report_structural_hash(report)
    return report.model_copy(update={"structural_hash": structural_hash})


def validate_pptx_objects_from_files(
    pptx_path: str | Path,
    *,
    scene_deck_path: str | Path | None = None,
    mode: ValidationMode = "inspect",
    profile: ValidationProfile = "basic",
) -> PptxObjectValidationReport:
    scene_deck = None
    if scene_deck_path is not None:
        scene_deck = SceneDeck.model_validate_json(Path(scene_deck_path).read_text(encoding="utf-8"))
    return validate_pptx_objects(
        pptx_path,
        scene_deck=scene_deck,
        scene_deck_path=scene_deck_path,
        mode=mode,
        profile=profile,
    )


def write_pptx_object_validation_report(report: PptxObjectValidationReport, output_path: str | Path) -> Path:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(validation_report_to_stable_json(report) + "\n", encoding="utf-8")
    return output


def summarize_pptx_object_validation(report: PptxObjectValidationReport) -> list[str]:
    lines = [
        "PPTX_OBJECT_VALIDATION "
        f"mode={report.mode} "
        f"profile={report.profile} "
        f"result={report.mode_result} "
        f"slides={report.slide_count} "
        f"findings={report.findings_summary.total_findings} "
        f"warnings={report.findings_summary.warning_count} "
        f"errors={report.findings_summary.error_count} "
        f"enforceable={report.findings_summary.enforceable_count} "
        f"text_overflow_risk={report.overflow_risk_count} "
        f"trace_missing={report.missing_trace_count} "
        f"duplicate_traces={report.duplicate_trace_count}"
    ]
    for finding in report.findings:
        location = f"slide={finding.slide_number}" if finding.slide_number is not None else "slide=deck"
        lines.append(f"FINDING {location} code={finding.code} severity={finding.severity} {finding.message}")
    return lines


def validation_report_to_stable_payload(
    report: PptxObjectValidationReport,
    *,
    include_paths: bool = True,
) -> dict[str, Any]:
    payload = report.model_dump(mode="json", exclude_none=True)
    if not include_paths:
        payload.pop("pptx_path", None)
        payload.pop("scene_deck_path", None)
    return _normalize_for_stable_json(payload)


def validation_report_to_stable_json(report: PptxObjectValidationReport) -> str:
    return json.dumps(validation_report_to_stable_payload(report), sort_keys=True, separators=(",", ":"), ensure_ascii=True)


def validation_report_structural_hash(report: PptxObjectValidationReport) -> str:
    payload = validation_report_to_stable_payload(report, include_paths=False)
    payload.pop("structural_hash", None)
    text = json.dumps(payload, sort_keys=True, separators=(",", ":"), ensure_ascii=True)
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _inspect_slide(
    slide_number: int,
    slide: Any,
) -> tuple[PptxSlideObjectInventory, list[PptxObjectFinding], dict[str, _InspectedShapeDetail]]:
    counts = SemanticObjectCounts()
    findings: list[PptxObjectFinding] = []
    shape_summaries: list[PptxShapeSummary] = []
    shape_details: dict[str, _InspectedShapeDetail] = {}
    counts.shape_count = len(slide.shapes)
    duplicate_names: dict[str, int] = {}
    title_excerpt: str | None = None

    for index, shape in enumerate(slide.shapes, start=1):
        summary, shape_findings, detail = _inspect_shape(index, shape, slide_number)
        shape_summaries.append(summary)
        shape_details[summary.name] = detail
        for semantic_type in summary.semantic_types:
            counts.increment(semantic_type)
        findings.extend(shape_findings)
        duplicate_names[summary.name] = duplicate_names.get(summary.name, 0) + 1
        if title_excerpt is None and summary.text_excerpt:
            title_excerpt = summary.text_excerpt

    for shape_name, occurrences in sorted(duplicate_names.items()):
        if occurrences > 1:
            findings.append(
                PptxObjectFinding(
                    code="duplicate_shape_name",
                    severity="warning",
                    category="unsupported_shape",
                    slide_number=slide_number,
                    shape_name=shape_name,
                    message=f"Slide contains duplicate shape name {shape_name!r}.",
                    details={"occurrences": occurrences},
                )
            )

    inventory = PptxSlideObjectInventory(
        slide_number=slide_number,
        pptx_slide_id=getattr(slide, "slide_id", None),
        title_excerpt=title_excerpt,
        has_notes_slide=bool(getattr(slide, "has_notes_slide", False)),
        chart_relationship_count=_slide_chart_relationship_count(slide),
        counts=counts,
        shapes=shape_summaries,
    )
    return inventory, _sorted_findings(findings), shape_details


def _inspect_shape(index: int, shape: Any, slide_number: int) -> tuple[PptxShapeSummary, list[PptxObjectFinding], _InspectedShapeDetail]:
    shape_type = getattr(shape, "shape_type", None)
    semantic_types: list[SemanticType] = []
    findings: list[PptxObjectFinding] = []
    shape_name = str(getattr(shape, "name", f"shape-{index}"))
    trace_kind, trace_object_id, trace_part = _parse_scene_trace_name(shape_name)

    if bool(getattr(shape, "is_placeholder", False)):
        semantic_types.append("placeholder")
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        semantic_types.append("group")
    if bool(getattr(shape, "has_chart", False)):
        semantic_types.append("native_chart")
    if bool(getattr(shape, "has_table", False)):
        semantic_types.append("native_table")
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        semantic_types.append("image_object")
    if shape_type == MSO_SHAPE_TYPE.LINE:
        semantic_types.append("other")
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        semantic_types.append("auto_shape")

    normalized_text, paragraph_levels, min_font_size_pt, max_font_size_pt = _shape_text_details(shape)
    text_excerpt = _truncate_text_excerpt(normalized_text)
    has_text = bool(normalized_text)
    if has_text:
        semantic_types.append("text_object")

    if not semantic_types:
        semantic_types.append("other")
        findings.append(
            PptxObjectFinding(
                code="unsupported_shape_type",
                severity="warning",
                category="unsupported_shape",
                slide_number=slide_number,
                shape_name=shape_name,
                message=f"Unsupported or unknown shape type {_shape_type_name(shape_type)!r} encountered.",
                details={"shape_type": _shape_type_name(shape_type)},
            )
        )

    child_count = len(getattr(shape, "shapes", [])) if shape_type == MSO_SHAPE_TYPE.GROUP else None
    table_row_count = _table_row_count(shape)
    table_column_count = _table_column_count(shape)
    chart_type = _chart_type(shape)
    chart_series_count = _chart_series_count(shape)
    chart_category_count = _chart_category_count(shape)
    left_in, top_in, width_in, height_in = _shape_bounds_in(shape)
    summary = PptxShapeSummary(
        shape_index=index,
        shape_id=getattr(shape, "shape_id", None),
        name=shape_name,
        shape_type=_shape_type_name(shape_type),
        trace_kind=trace_kind,
        trace_object_id=trace_object_id,
        trace_part=trace_part,
        semantic_types=sorted(set(semantic_types), key=_SEMANTIC_TYPE_ORDER.index),
        has_text=has_text,
        text_excerpt=text_excerpt,
        text_length=len(normalized_text) if normalized_text is not None else None,
        paragraph_count=len(paragraph_levels) if paragraph_levels else 0,
        max_paragraph_level=max(paragraph_levels) if paragraph_levels else None,
        min_font_size_pt=min_font_size_pt,
        max_font_size_pt=max_font_size_pt,
        left_in=left_in,
        top_in=top_in,
        width_in=width_in,
        height_in=height_in,
        child_count=child_count,
        table_row_count=table_row_count,
        table_column_count=table_column_count,
        chart_type=chart_type,
        chart_series_count=chart_series_count,
        chart_category_count=chart_category_count,
    )
    detail = _InspectedShapeDetail(
        name=shape_name,
        trace_kind=trace_kind,
        trace_object_id=trace_object_id,
        trace_part=trace_part,
        semantic_types=tuple(summary.semantic_types),
        normalized_text=normalized_text,
        paragraph_levels=paragraph_levels,
        min_font_size_pt=min_font_size_pt,
        max_font_size_pt=max_font_size_pt,
        left_in=left_in,
        top_in=top_in,
        width_in=width_in,
        height_in=height_in,
    )
    return summary, findings, detail


def _compare_slide_to_scene(
    slide_number: int,
    scene_slide: Any,
    inventory: PptxSlideObjectInventory,
    shape_details: dict[str, _InspectedShapeDetail],
    *,
    profile: ValidationProfile,
) -> tuple[SceneExpectationSummary, list[PptxObjectFinding], PptxSlideTextFitSummary]:
    expected_counts = _scene_expected_counts(scene_slide)
    comparison = SceneExpectationSummary(
        comparison_basis="ordinal",
        scene_slide_id=scene_slide.slide_id,
        layout_family=scene_slide.layout_family,
        expected_counts=expected_counts,
        actual_counts=inventory.counts,
    )
    findings: list[PptxObjectFinding] = []
    text_diagnostics: list[PptxTextObjectDiagnostic] = []
    traced_shape_index = _index_traced_shapes(inventory.shapes)

    if inventory.counts.text_object_count < expected_counts.text_object_count:
        findings.append(
            _finding(
                code="expected_text_object_missing",
                severity="warning",
                category="semantic_mismatch",
                slide_number=slide_number,
                scene_slide_id=scene_slide.slide_id,
                message="PPTX contains fewer text-capable objects than SceneDeck expects.",
                expected=expected_counts.text_object_count,
                actual=inventory.counts.text_object_count,
            )
        )
    elif inventory.counts.text_object_count != expected_counts.text_object_count:
        findings.append(
            _finding(
                code="ambiguous_text_mapping",
                severity="info",
                category="semantic_mismatch",
                slide_number=slide_number,
                scene_slide_id=scene_slide.slide_id,
                message="Text object counts differ, so one-to-one text mapping remains ambiguous.",
                expected=expected_counts.text_object_count,
                actual=inventory.counts.text_object_count,
            )
        )

    if inventory.counts.image_object_count < expected_counts.image_object_count:
        findings.append(
            _finding(
                code="expected_image_object_missing",
                severity="warning",
                category="semantic_mismatch",
                slide_number=slide_number,
                scene_slide_id=scene_slide.slide_id,
                message="PPTX contains fewer image objects than SceneDeck expects.",
                expected=expected_counts.image_object_count,
                actual=inventory.counts.image_object_count,
                )
            )

    traced_shapes = set(_trace_object_ids(traced_shape_index, "shape"))
    traced_dividers = set(_trace_object_ids(traced_shape_index, "divider"))
    traced_callouts = _trace_callout_parts(traced_shape_index)
    traced_background_shapes = set(_trace_object_ids(traced_shape_index, "background_shape"))
    traced_background_dividers = set(_trace_object_ids(traced_shape_index, "background_divider"))

    traced_background_duplicates = _trace_duplicate_counts(
        inventory.shapes,
        kinds={"background_shape", "background_divider"},
        include_parts=False,
    )
    for (trace_kind, object_id), occurrences in sorted(traced_background_duplicates.items()):
        if occurrences > 1:
            findings.append(
                PptxObjectFinding(
                    code="background_motif_trace_duplicate",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    message="PPTX contains duplicate traced background motif objects.",
                    details={"trace_kind": trace_kind, "object_id": object_id, "occurrences": occurrences},
                )
                )

    for item in scene_slide.objects:
        kind = getattr(item, "kind", None)
        if kind == "shape" and item.object_id not in traced_shapes:
            findings.append(
                PptxObjectFinding(
                    code="expected_shape_missing",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    message="PPTX is missing an expected scene shape object.",
                    details={"object_id": item.object_id},
                )
            )
        elif kind == "divider" and item.object_id not in traced_dividers:
            findings.append(
                PptxObjectFinding(
                    code="expected_divider_missing",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    message="PPTX is missing an expected divider line object.",
                    details={"object_id": item.object_id},
                )
            )
        elif kind == "callout":
            parts = traced_callouts.get(item.object_id, set())
            required_parts = {"background", "body"}
            if getattr(item, "title", None) is not None:
                required_parts.add("title")
            if getattr(item, "accent", None) is not None:
                required_parts.add("accent")
            if not required_parts.issubset(parts):
                findings.append(
                    PptxObjectFinding(
                        code="expected_callout_missing",
                        severity="warning",
                        category="semantic_mismatch",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        message="PPTX is missing one or more expected callout parts.",
                        details={
                            "object_id": item.object_id,
                            "expected_parts": sorted(required_parts),
                            "actual_parts": sorted(parts),
                        },
                    )
                )
        elif profile == "scene-strict" and kind == "image":
            findings.extend(
                _strict_trace_findings_for_object(
                    slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    trace_kind="image",
                    expected_code="image_object_missing_trace",
                    expected_message="PPTX is missing a traced image object for a SceneDeck image.",
                    expected_semantic_type="image_object",
                    traced_shape_index=traced_shape_index,
                )
            )
        elif profile == "scene-strict" and kind == "native_table":
            findings.extend(
                _strict_trace_findings_for_object(
                    slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    trace_kind="native_table",
                    expected_code="native_table_missing_trace",
                    expected_message="PPTX is missing a traced native table object for a SceneDeck table.",
                    expected_semantic_type="native_table",
                    traced_shape_index=traced_shape_index,
                )
            )
        elif profile == "scene-strict" and kind == "native_chart":
            findings.extend(
                _strict_trace_findings_for_object(
                    slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    trace_kind="native_chart",
                    expected_code="native_chart_missing_trace",
                    expected_message="PPTX is missing a traced native chart object for a SceneDeck chart.",
                    expected_semantic_type="native_chart",
                    traced_shape_index=traced_shape_index,
                )
            )

    for motif in scene_slide.background.motifs:
        kind = getattr(motif, "kind", None)
        if kind == "shape":
            shape_type = str(getattr(motif, "shape_type", "")).strip()
            if shape_type not in _SUPPORTED_BACKGROUND_SHAPE_TYPES:
                findings.append(
                    PptxObjectFinding(
                        code="unsupported_background_motif",
                        severity="warning",
                        category="unsupported_shape",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        message="SceneDeck background motif shape type is not supported by the scene validator contract.",
                        details={"object_id": motif.object_id, "shape_type": shape_type},
                    )
                )
                continue
            if motif.object_id not in traced_background_shapes:
                findings.append(
                    PptxObjectFinding(
                        code="expected_background_motif_missing",
                        severity="warning",
                        category="semantic_mismatch",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        message="PPTX is missing an expected background motif shape.",
                        details={"object_id": motif.object_id, "motif_kind": kind},
                    )
                )
        elif kind == "divider":
            if motif.object_id not in traced_background_dividers:
                findings.append(
                    PptxObjectFinding(
                        code="expected_background_motif_missing",
                        severity="warning",
                        category="semantic_mismatch",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        message="PPTX is missing an expected background motif divider.",
                        details={"object_id": motif.object_id, "motif_kind": kind},
                    )
                )

    if inventory.counts.native_table_count < expected_counts.native_table_count:
        findings.append(
            _finding(
                code="expected_native_table_missing",
                severity="warning",
                category="native_object",
                slide_number=slide_number,
                scene_slide_id=scene_slide.slide_id,
                message="PPTX is missing an expected native table object.",
                expected=expected_counts.native_table_count,
                actual=inventory.counts.native_table_count,
            )
        )
        if inventory.counts.image_object_count > expected_counts.image_object_count:
            findings.append(
                _finding(
                    code="picture_used_for_expected_table",
                    severity="warning",
                    category="editability",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    message="A picture may be standing in for an expected native table.",
                    expected=expected_counts.native_table_count,
                    actual=inventory.counts.image_object_count,
                )
            )

    traced_tables = {
        object_id: shapes[0]
        for object_id, shapes in _trace_kind_index(traced_shape_index, "native_table").items()
        if shapes and "native_table" in shapes[0].semantic_types
    }
    for item in scene_slide.objects:
        if getattr(item, "kind", None) != "native_table":
            continue
        traced_shape = traced_tables.get(item.object_id)
        if traced_shape is None:
            continue
        expected_rows = 1 + len(item.rows)
        expected_columns = len(item.headers)
        actual_rows = traced_shape.table_row_count
        actual_columns = traced_shape.table_column_count
        if actual_rows != expected_rows or actual_columns != expected_columns:
            findings.append(
                PptxObjectFinding(
                    code="native_table_dimension_mismatch",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    shape_name=traced_shape.name,
                    message="Scene-rendered native table dimensions differ from SceneDeck expectations.",
                    details={
                        "object_id": item.object_id,
                        "expected_rows": expected_rows,
                        "actual_rows": actual_rows,
                        "expected_columns": expected_columns,
                        "actual_columns": actual_columns,
                    },
                )
            )

    if inventory.counts.native_chart_count < expected_counts.native_chart_count:
        findings.append(
            _finding(
                code="expected_native_chart_missing",
                severity="warning",
                category="native_object",
                slide_number=slide_number,
                scene_slide_id=scene_slide.slide_id,
                message="PPTX is missing an expected native chart object.",
                expected=expected_counts.native_chart_count,
                actual=inventory.counts.native_chart_count,
            )
        )
        if inventory.counts.image_object_count > expected_counts.image_object_count:
            findings.append(
                _finding(
                    code="picture_used_for_expected_chart",
                    severity="warning",
                    category="editability",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    message="A picture may be standing in for an expected native chart.",
                    expected=expected_counts.native_chart_count,
                    actual=inventory.counts.image_object_count,
                )
            )

    traced_charts = {
        object_id: shapes[0]
        for object_id, shapes in _trace_kind_index(traced_shape_index, "native_chart").items()
        if shapes and "native_chart" in shapes[0].semantic_types
    }
    for item in scene_slide.objects:
        if getattr(item, "kind", None) != "native_chart":
            continue
        traced_shape = traced_charts.get(item.object_id)
        if traced_shape is None:
            continue
        expected_series_count = len(item.series)
        expected_category_count = len(item.categories)
        if traced_shape.chart_series_count != expected_series_count:
            findings.append(
                PptxObjectFinding(
                    code="native_chart_series_count_mismatch",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    shape_name=traced_shape.name,
                    message="Scene-rendered native chart series count differs from SceneDeck expectations.",
                    details={
                        "object_id": item.object_id,
                        "expected_series_count": expected_series_count,
                        "actual_series_count": traced_shape.chart_series_count,
                    },
                )
            )
        if traced_shape.chart_category_count != expected_category_count:
            findings.append(
                PptxObjectFinding(
                    code="native_chart_category_count_mismatch",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    shape_name=traced_shape.name,
                    message="Scene-rendered native chart category count differs from SceneDeck expectations.",
                    details={
                        "object_id": item.object_id,
                        "expected_category_count": expected_category_count,
                        "actual_category_count": traced_shape.chart_category_count,
                    },
                )
            )
        if traced_shape.chart_type is not None and traced_shape.chart_type != item.chart_type:
            findings.append(
                PptxObjectFinding(
                    code="native_chart_type_mismatch",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    shape_name=traced_shape.name,
                    message="Scene-rendered native chart type differs from SceneDeck expectations.",
                    details={
                        "object_id": item.object_id,
                        "expected_chart_type": item.chart_type,
                        "actual_chart_type": traced_shape.chart_type,
                    },
                )
            )

    if profile == "scene-strict":
        findings.extend(_scene_trace_duplicate_findings(slide_number, scene_slide.slide_id, inventory.shapes))
        findings.extend(_scene_trace_type_findings(slide_number, scene_slide, traced_shape_index))

    text_diagnostics = _text_fit_diagnostics_for_scene_slide(
        slide_number,
        scene_slide,
        traced_shape_index,
        shape_details,
    )
    if profile == "scene-strict":
        findings.extend(_text_diagnostics_to_findings(text_diagnostics))

    return comparison, _sorted_findings(findings), _summarize_text_fit_diagnostics(text_diagnostics)


def _index_traced_shapes(shapes: list[PptxShapeSummary]) -> dict[tuple[str, str], list[PptxShapeSummary]]:
    traced: dict[tuple[str, str], list[PptxShapeSummary]] = {}
    for shape in shapes:
        if shape.trace_kind is None or shape.trace_object_id is None:
            continue
        traced.setdefault((shape.trace_kind, shape.trace_object_id), []).append(shape)
    for key in traced:
        traced[key] = sorted(traced[key], key=lambda item: (item.shape_index, item.name))
    return traced


def _trace_kind_index(
    traced_shape_index: dict[tuple[str, str], list[PptxShapeSummary]],
    trace_kind: str,
) -> dict[str, list[PptxShapeSummary]]:
    return {
        object_id: shapes
        for (kind, object_id), shapes in traced_shape_index.items()
        if kind == trace_kind
    }


def _trace_object_ids(traced_shape_index: dict[tuple[str, str], list[PptxShapeSummary]], trace_kind: str) -> list[str]:
    return sorted(_trace_kind_index(traced_shape_index, trace_kind))


def _trace_callout_parts(traced_shape_index: dict[tuple[str, str], list[PptxShapeSummary]]) -> dict[str, set[str]]:
    parts: dict[str, set[str]] = {}
    for object_id, shapes in _trace_kind_index(traced_shape_index, "callout").items():
        parts[object_id] = {shape.trace_part or "" for shape in shapes}
    return parts


def _trace_duplicate_counts(
    shapes: list[PptxShapeSummary],
    *,
    kinds: set[str],
    include_parts: bool,
) -> dict[tuple[str, str], int]:
    counts: dict[tuple[str, str], int] = {}
    for shape in shapes:
        if shape.trace_kind not in kinds or shape.trace_object_id is None:
            continue
        object_id = shape.trace_object_id if not include_parts else f"{shape.trace_object_id}:{shape.trace_part or ''}"
        key = (shape.trace_kind, object_id)
        counts[key] = counts.get(key, 0) + 1
    return counts


def _scene_trace_duplicate_findings(
    slide_number: int,
    scene_slide_id: str,
    shapes: list[PptxShapeSummary],
) -> list[PptxObjectFinding]:
    findings: list[PptxObjectFinding] = []
    trace_name_counts: dict[str, int] = {}
    generic_counts: dict[tuple[str, str], int] = {}
    callout_counts: dict[tuple[str, str, str], int] = {}
    for shape in shapes:
        if not shape.name.startswith("scene:"):
            continue
        trace_name_counts[shape.name] = trace_name_counts.get(shape.name, 0) + 1
        if shape.trace_kind is None or shape.trace_object_id is None:
            continue
        if shape.trace_kind in {"background_shape", "background_divider"}:
            continue
        if shape.trace_kind == "callout":
            key = (shape.trace_kind, shape.trace_object_id, shape.trace_part or "")
            callout_counts[key] = callout_counts.get(key, 0) + 1
            continue
        key = (shape.trace_kind, shape.trace_object_id)
        generic_counts[key] = generic_counts.get(key, 0) + 1
    for trace_name, occurrences in sorted(trace_name_counts.items()):
        if occurrences > 1:
            findings.append(
                PptxObjectFinding(
                    code="scene_trace_duplicate",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide_id,
                    shape_name=trace_name,
                    message="PPTX contains a duplicate scene trace name.",
                    details={"trace_name": trace_name, "occurrences": occurrences},
                )
            )
    for (trace_kind, object_id), occurrences in sorted(generic_counts.items()):
        if occurrences > 1:
            findings.append(
                PptxObjectFinding(
                    code="scene_trace_duplicate",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide_id,
                    message="PPTX contains duplicate traced scene objects.",
                    details={"trace_kind": trace_kind, "object_id": object_id, "occurrences": occurrences},
                )
            )
    for (trace_kind, object_id, trace_part), occurrences in sorted(callout_counts.items()):
        if occurrences > 1:
            findings.append(
                PptxObjectFinding(
                    code="scene_trace_duplicate",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide_id,
                    message="PPTX contains duplicate traced callout parts.",
                    details={
                        "trace_kind": trace_kind,
                        "object_id": object_id,
                        "trace_part": trace_part,
                        "occurrences": occurrences,
                    },
                )
            )
    return findings


def _strict_trace_findings_for_object(
    slide_number: int,
    *,
    scene_slide_id: str,
    object_id: str,
    trace_kind: str,
    expected_code: FindingCode,
    expected_message: str,
    expected_semantic_type: SemanticType,
    traced_shape_index: dict[tuple[str, str], list[PptxShapeSummary]],
) -> list[PptxObjectFinding]:
    traced_shapes = traced_shape_index.get((trace_kind, object_id), [])
    if not traced_shapes:
        return [
            PptxObjectFinding(
                code=expected_code,
                severity="warning",
                category="semantic_mismatch",
                slide_number=slide_number,
                scene_slide_id=scene_slide_id,
                message=expected_message,
                details={"object_id": object_id, "trace_kind": trace_kind},
            )
        ]
    findings: list[PptxObjectFinding] = []
    traced_shape = traced_shapes[0]
    if expected_semantic_type not in traced_shape.semantic_types:
        findings.append(
            PptxObjectFinding(
                code="scene_trace_type_mismatch",
                severity="warning",
                category="semantic_mismatch",
                slide_number=slide_number,
                scene_slide_id=scene_slide_id,
                shape_name=traced_shape.name,
                message="Traced scene object does not map to the expected PPTX semantic type.",
                details={
                    "object_id": object_id,
                    "trace_kind": trace_kind,
                    "expected_semantic_type": expected_semantic_type,
                    "actual_semantic_types": list(traced_shape.semantic_types),
                },
            )
        )
    return findings


def _scene_trace_type_findings(
    slide_number: int,
    scene_slide: Any,
    traced_shape_index: dict[tuple[str, str], list[PptxShapeSummary]],
) -> list[PptxObjectFinding]:
    findings: list[PptxObjectFinding] = []
    for item in scene_slide.objects:
        kind = getattr(item, "kind", None)
        object_id = getattr(item, "object_id", None)
        if not object_id:
            continue
        if kind == "text_box":
            traced_shapes = traced_shape_index.get(("text_box", object_id), [])
            if traced_shapes and "text_object" not in traced_shapes[0].semantic_types:
                findings.append(
                    PptxObjectFinding(
                        code="scene_trace_type_mismatch",
                        severity="warning",
                        category="semantic_mismatch",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        shape_name=traced_shapes[0].name,
                        message="TextBox trace does not map to a text-capable PPTX shape.",
                        details={
                            "object_id": object_id,
                            "trace_kind": "text_box",
                            "actual_semantic_types": list(traced_shapes[0].semantic_types),
                        },
                    )
                )
        elif kind == "shape":
            traced_shapes = traced_shape_index.get(("shape", object_id), [])
            if traced_shapes and not _is_editable_shape_trace(traced_shapes[0]):
                findings.append(
                    PptxObjectFinding(
                        code="scene_trace_type_mismatch",
                        severity="warning",
                        category="semantic_mismatch",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        shape_name=traced_shapes[0].name,
                        message="Scene shape trace does not map to an editable PPTX shape.",
                        details={"object_id": object_id, "trace_kind": "shape", "shape_type": traced_shapes[0].shape_type},
                    )
                )
        elif kind == "divider":
            traced_shapes = traced_shape_index.get(("divider", object_id), [])
            if traced_shapes and not _is_line_trace(traced_shapes[0]):
                findings.append(
                    PptxObjectFinding(
                        code="scene_trace_type_mismatch",
                        severity="warning",
                        category="semantic_mismatch",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        shape_name=traced_shapes[0].name,
                        message="Divider trace does not map to a native PPTX line shape.",
                        details={"object_id": object_id, "trace_kind": "divider", "shape_type": traced_shapes[0].shape_type},
                    )
                )
    for motif in scene_slide.background.motifs:
        kind = getattr(motif, "kind", None)
        object_id = getattr(motif, "object_id", None)
        if not object_id:
            continue
        trace_kind = "background_shape" if kind == "shape" else "background_divider" if kind == "divider" else None
        if trace_kind is None:
            continue
        traced_shapes = traced_shape_index.get((trace_kind, object_id), [])
        if not traced_shapes:
            continue
        valid = _is_editable_shape_trace(traced_shapes[0]) if trace_kind == "background_shape" else _is_line_trace(traced_shapes[0])
        if not valid:
            findings.append(
                PptxObjectFinding(
                    code="scene_trace_type_mismatch",
                    severity="warning",
                    category="semantic_mismatch",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    shape_name=traced_shapes[0].name,
                    message="Background motif trace does not map to the expected editable PPTX shape type.",
                    details={"object_id": object_id, "trace_kind": trace_kind, "shape_type": traced_shapes[0].shape_type},
                )
            )
    return findings


def _is_editable_shape_trace(shape: PptxShapeSummary) -> bool:
    return "auto_shape" in shape.semantic_types or _is_line_trace(shape)


def _is_line_trace(shape: PptxShapeSummary) -> bool:
    return shape.shape_type == "LINE"


def _text_fit_diagnostics_for_scene_slide(
    slide_number: int,
    scene_slide: Any,
    traced_shape_index: dict[tuple[str, str], list[PptxShapeSummary]],
    shape_details: dict[str, _InspectedShapeDetail],
) -> list[PptxTextObjectDiagnostic]:
    diagnostics: list[PptxTextObjectDiagnostic] = []
    for item in scene_slide.objects:
        if not isinstance(item, TextBox):
            continue
        traced_shapes = traced_shape_index.get(("text_box", item.object_id), [])
        if not traced_shapes:
            diagnostics.append(
                PptxTextObjectDiagnostic(
                    code="text_box_missing_trace",
                    severity="warning",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    details={"fit_mode": item.fit.mode},
                )
            )
            continue
        traced_shape = traced_shapes[0]
        detail = shape_details.get(traced_shape.name)
        expected_text = _scene_text_normalized(item)
        if not expected_text:
            diagnostics.append(
                PptxTextObjectDiagnostic(
                    code="text_box_empty",
                    severity="warning",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    shape_name=traced_shape.name,
                    details={"source": "scene"},
                )
            )
            continue
        if detail is None or not detail.normalized_text:
            diagnostics.append(
                PptxTextObjectDiagnostic(
                    code="text_box_empty",
                    severity="warning",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    shape_name=traced_shape.name,
                    details={"source": "pptx"},
                )
            )
        if item.fit.mode not in _SUPPORTED_TEXT_FIT_MODES:
            diagnostics.append(
                PptxTextObjectDiagnostic(
                    code="text_box_fit_policy_unsupported",
                    severity="warning",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    shape_name=traced_shape.name,
                    details={"fit_mode": item.fit.mode, "overflow_action": item.fit.overflow_action},
                )
            )
        if any(bullet.level > _BULLET_LEVEL_HARD_LIMIT for bullet in item.bullet_list):
            diagnostics.append(
                PptxTextObjectDiagnostic(
                    code="text_box_bullet_level_unsupported",
                    severity="warning",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    shape_name=traced_shape.name,
                    details={"max_level": max(bullet.level for bullet in item.bullet_list)},
                )
            )
        elif item.bullet_list and detail is not None and detail.paragraph_levels:
            expected_max_level = max(bullet.level for bullet in item.bullet_list)
            actual_max_level = max(detail.paragraph_levels)
            if actual_max_level < expected_max_level:
                diagnostics.append(
                    PptxTextObjectDiagnostic(
                        code="text_box_bullet_level_unsupported",
                        severity="warning",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        object_id=item.object_id,
                        shape_name=traced_shape.name,
                        details={"expected_max_level": expected_max_level, "actual_max_level": actual_max_level},
                    )
                )
        if item.fit.min_font_size_pt is not None:
            actual_min_font = detail.min_font_size_pt if detail is not None else traced_shape.min_font_size_pt
            estimated_min_font = actual_min_font or _scene_text_min_font_size(item)
            if estimated_min_font is not None and estimated_min_font < item.fit.min_font_size_pt:
                diagnostics.append(
                    PptxTextObjectDiagnostic(
                        code="text_box_font_below_min_fit_policy",
                        severity="warning",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        object_id=item.object_id,
                        shape_name=traced_shape.name,
                        details={
                            "actual_min_font_size_pt": round(estimated_min_font, 3),
                            "required_min_font_size_pt": item.fit.min_font_size_pt,
                        },
                    )
                )
        if detail is not None and detail.normalized_text:
            actual_text = detail.normalized_text
            if (
                item.fit.mode != "truncate"
                and actual_text != expected_text
                and len(actual_text) < len(expected_text)
                and expected_text.startswith(actual_text)
            ):
                diagnostics.append(
                    PptxTextObjectDiagnostic(
                        code="text_box_truncated_without_policy",
                        severity="warning",
                        slide_number=slide_number,
                        scene_slide_id=scene_slide.slide_id,
                        object_id=item.object_id,
                        shape_name=traced_shape.name,
                        details={"expected_text_length": len(expected_text), "actual_text_length": len(actual_text)},
                    )
                )
        line_estimate = _estimate_text_box_lines(item, detail)
        if line_estimate["estimated_line_count"] > line_estimate["line_capacity"]:
            diagnostics.append(
                PptxTextObjectDiagnostic(
                    code="text_box_line_count_exceeds_bounds",
                    severity="warning",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    shape_name=traced_shape.name,
                    details=line_estimate,
                )
            )
            diagnostics.append(
                PptxTextObjectDiagnostic(
                    code="text_box_overflow_risk",
                    severity="warning",
                    slide_number=slide_number,
                    scene_slide_id=scene_slide.slide_id,
                    object_id=item.object_id,
                    shape_name=traced_shape.name,
                    details=line_estimate,
                )
            )
    return sorted(diagnostics, key=lambda item: (item.slide_number, item.object_id, item.code, item.shape_name or ""))


def _text_diagnostics_to_findings(diagnostics: list[PptxTextObjectDiagnostic]) -> list[PptxObjectFinding]:
    return [
        PptxObjectFinding(
            code=diagnostic.code,
            severity=diagnostic.severity,
            category="text_fit",
            slide_number=diagnostic.slide_number,
            scene_slide_id=diagnostic.scene_slide_id,
            shape_name=diagnostic.shape_name,
            message=_text_fit_message(diagnostic.code),
            details={"object_id": diagnostic.object_id, **diagnostic.details},
        )
        for diagnostic in diagnostics
    ]


def _summarize_text_fit_diagnostics(diagnostics: list[PptxTextObjectDiagnostic]) -> PptxSlideTextFitSummary:
    return PptxSlideTextFitSummary(
        diagnostic_count=len(diagnostics),
        overflow_risk_count=sum(1 for item in diagnostics if item.code == "text_box_overflow_risk"),
        unsupported_fit_policy_count=sum(1 for item in diagnostics if item.code == "text_box_fit_policy_unsupported"),
        missing_text_trace_count=sum(1 for item in diagnostics if item.code == "text_box_missing_trace"),
        diagnostics=diagnostics,
    )


def _scene_text_normalized(text_box: TextBox) -> str:
    paragraphs: list[str] = []
    if text_box.bullet_list:
        for bullet in text_box.bullet_list:
            text = "".join(run.text for run in bullet.runs)
            normalized = " ".join(text.split()).strip()
            if normalized:
                paragraphs.append(normalized)
    else:
        for paragraph in _scene_text_paragraphs(text_box):
            if paragraph["text"]:
                paragraphs.append(paragraph["text"])
    return "\n".join(paragraphs)


def _scene_text_paragraphs(text_box: TextBox) -> list[dict[str, Any]]:
    paragraphs: list[dict[str, Any]] = []
    if text_box.bullet_list:
        for bullet in text_box.bullet_list:
            text = "".join(run.text for run in bullet.runs)
            normalized = " ".join(text.split()).strip()
            if not normalized:
                continue
            paragraphs.append(
                {
                    "text": normalized,
                    "level": bullet.level,
                    "font_size_pt": max((_scene_run_font_size(run, text_box.role) for run in bullet.runs), default=_scene_default_font_size(text_box.role)),
                }
            )
        return paragraphs
    current_runs: list[tuple[str, float]] = []
    for run in text_box.runs:
        parts = run.text.split("\n")
        for index, part in enumerate(parts):
            normalized = " ".join(part.split()).strip()
            if normalized:
                current_runs.append((normalized, _scene_run_font_size(run, text_box.role)))
            if index < len(parts) - 1:
                if current_runs:
                    paragraphs.append(
                        {
                            "text": " ".join(text for text, _ in current_runs).strip(),
                            "level": 0,
                            "font_size_pt": max(size for _, size in current_runs),
                        }
                    )
                current_runs = []
    if current_runs:
        paragraphs.append(
            {
                "text": " ".join(text for text, _ in current_runs).strip(),
                "level": 0,
                "font_size_pt": max(size for _, size in current_runs),
            }
        )
    return [item for item in paragraphs if item["text"]]


def _scene_default_font_size(role: str) -> float:
    _, size_pt, _ = scene_text_role_defaults(role)
    return size_pt


def _scene_run_font_size(run: Any, role: str) -> float:
    return float(getattr(run, "size_pt", None) or _scene_default_font_size(role))


def _scene_text_min_font_size(text_box: TextBox) -> float | None:
    sizes = []
    for paragraph in _scene_text_paragraphs(text_box):
        sizes.append(paragraph["font_size_pt"])
    return min(sizes) if sizes else None


def _estimate_text_box_lines(text_box: TextBox, detail: _InspectedShapeDetail | None) -> dict[str, Any]:
    width_in = detail.width_in if detail is not None and detail.width_in is not None else text_box.bounds.width
    height_in = detail.height_in if detail is not None and detail.height_in is not None else text_box.bounds.height
    paragraphs = _scene_text_paragraphs(text_box)
    estimated_line_count = 0
    used_height_in = 0.0
    for paragraph in paragraphs or [{"text": _scene_text_normalized(text_box), "level": 0, "font_size_pt": _scene_default_font_size(text_box.role)}]:
        paragraph_text = str(paragraph["text"])
        font_size_pt = float(paragraph["font_size_pt"])
        indent_in = min(float(paragraph["level"]) * 0.25, max(width_in - 0.1, 0.0))
        effective_width_in = max(width_in - indent_in, 0.25)
        chars_per_line = max(1, int((effective_width_in * 72.0) / max(font_size_pt * 0.55, 1.0)))
        line_count = max(1, math.ceil(len(paragraph_text) / chars_per_line))
        estimated_line_count += line_count
        used_height_in += line_count * ((font_size_pt * 1.2) / 72.0)
    average_font_size_pt = sum(paragraph["font_size_pt"] for paragraph in paragraphs) / len(paragraphs) if paragraphs else _scene_default_font_size(text_box.role)
    line_capacity = max(1, int(height_in / max((average_font_size_pt * 1.2) / 72.0, 0.01)))
    return {
        "estimated_line_count": estimated_line_count,
        "line_capacity": line_capacity,
        "bounds_width_in": round(width_in, 6),
        "bounds_height_in": round(height_in, 6),
        "estimated_used_height_in": round(used_height_in, 6),
    }


def _text_fit_message(code: TextFitDiagnosticCode) -> str:
    return {
        "text_box_empty": "Traced text box content is empty.",
        "text_box_overflow_risk": "Text box has a conservative overflow risk estimate.",
        "text_box_line_count_exceeds_bounds": "Estimated line count exceeds available text box bounds.",
        "text_box_font_below_min_fit_policy": "Traced text box font size falls below the SceneDeck fit policy minimum.",
        "text_box_fit_policy_unsupported": "SceneDeck text fit policy is not fully enforced on the scene compile path.",
        "text_box_truncated_without_policy": "Traced text appears truncated without a truncate fit policy.",
        "text_box_bullet_level_unsupported": "Bullet nesting depth is not fully supported by the scene validator contract.",
        "text_box_missing_trace": "PPTX is missing a traced text box for a SceneDeck TextBox.",
    }[code]


def _scene_expected_counts(scene_slide: Any) -> SemanticObjectCounts:
    counts = SemanticObjectCounts(shape_count=len(scene_slide.objects))
    for item in scene_slide.objects:
        kind = getattr(item, "kind", None)
        if kind == "text_box":
            counts.text_object_count += 1
        elif kind == "image":
            counts.image_object_count += 1
        elif kind == "native_table":
            counts.native_table_count += 1
        elif kind == "native_chart":
            counts.native_chart_count += 1
        elif kind == "group":
            counts.group_count += 1
        elif kind == "shape":
            counts.other_shape_count += 1
        elif kind == "divider":
            counts.other_shape_count += 1
        elif kind == "callout":
            counts.text_object_count += 1 + (1 if getattr(item, "title", None) is not None else 0)
    return counts


def _finding(
    *,
    code: FindingCode,
    severity: FindingSeverity,
    category: FindingCategory,
    slide_number: int | None,
    scene_slide_id: str | None = None,
    message: str,
    expected: int | None = None,
    actual: int | None = None,
) -> PptxObjectFinding:
    details: dict[str, Any] = {}
    if expected is not None:
        details["expected"] = expected
    if actual is not None:
        details["actual"] = actual
    return PptxObjectFinding(
        code=code,
        severity=severity,
        category=category,
        slide_number=slide_number,
        scene_slide_id=scene_slide_id,
        message=message,
        details=details,
    )


def _chart_part_count(pptx_path: Path) -> int:
    with zipfile.ZipFile(pptx_path) as archive:
        return len([name for name in archive.namelist() if name.startswith("ppt/charts/chart") and name.endswith(".xml")])


def _slide_chart_relationship_count(slide: Any) -> int:
    rels = getattr(getattr(slide, "part", None), "rels", {})
    return sum(1 for rel in rels.values() if str(getattr(rel, "reltype", "")).endswith("/chart"))


def _shape_text_details(shape: Any) -> tuple[str | None, tuple[int, ...], float | None, float | None]:
    if not bool(getattr(shape, "has_text_frame", False)):
        return (None, tuple(), None, None)
    paragraphs: list[str] = []
    paragraph_levels: list[int] = []
    font_sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        paragraph_text = "".join(str(getattr(run, "text", "")) for run in paragraph.runs)
        if not paragraph_text:
            paragraph_text = str(getattr(paragraph, "text", ""))
        normalized = " ".join(paragraph_text.split()).strip()
        if not normalized:
            continue
        paragraphs.append(normalized)
        paragraph_levels.append(int(getattr(paragraph, "level", 0) or 0))
        for run in paragraph.runs:
            run_size = getattr(getattr(run, "font", None), "size", None)
            if run_size is None:
                continue
            try:
                font_sizes.append(float(run_size.pt))
            except Exception:
                continue
    normalized_text = "\n".join(paragraphs) if paragraphs else None
    min_font = min(font_sizes) if font_sizes else None
    max_font = max(font_sizes) if font_sizes else None
    return (normalized_text, tuple(paragraph_levels), min_font, max_font)


def _truncate_text_excerpt(text: str | None) -> str | None:
    if not text:
        return None
    return text[:120]


def _shape_bounds_in(shape: Any) -> tuple[float | None, float | None, float | None, float | None]:
    return (
        _emu_to_inches(getattr(shape, "left", None)),
        _emu_to_inches(getattr(shape, "top", None)),
        _emu_to_inches(getattr(shape, "width", None)),
        _emu_to_inches(getattr(shape, "height", None)),
    )


def _emu_to_inches(value: Any) -> float | None:
    if value is None:
        return None
    try:
        inches = float(value) / 914400.0
    except Exception:
        return None
    return round(inches, 6)


def _shape_type_name(shape_type: Any) -> str:
    return str(getattr(shape_type, "name", shape_type))


def _parse_scene_trace_name(name: str) -> tuple[str | None, str | None, str | None]:
    if not name.startswith("scene:"):
        return (None, None, None)
    if name.startswith("scene:callout:"):
        remainder = name[len("scene:callout:") :]
        if ":" not in remainder:
            return (None, None, None)
        object_id, trace_part = remainder.rsplit(":", 1)
        if not object_id or not trace_part:
            return (None, None, None)
        return ("callout", object_id, trace_part)
    parts = name.split(":", 2)
    if len(parts) != 3 or not parts[1] or not parts[2]:
        return (None, None, None)
    return (parts[1], parts[2], None)


def _table_row_count(shape: Any) -> int | None:
    if not bool(getattr(shape, "has_table", False)):
        return None
    try:
        return len(shape.table.rows)
    except Exception:
        return None


def _table_column_count(shape: Any) -> int | None:
    if not bool(getattr(shape, "has_table", False)):
        return None
    try:
        return len(shape.table.columns)
    except Exception:
        return None


def _chart_type(shape: Any) -> str | None:
    if not bool(getattr(shape, "has_chart", False)):
        return None
    try:
        raw_type = getattr(shape.chart.chart_type, "name", shape.chart.chart_type)
    except Exception:
        return None
    chart_type = str(raw_type).upper()
    if chart_type.startswith("BAR_"):
        return "bar"
    if chart_type.startswith("COLUMN_"):
        return "column"
    if chart_type.startswith("LINE_"):
        return "line"
    if chart_type.startswith("SCATTER_"):
        return "scatter"
    return str(raw_type).lower()


def _chart_series_count(shape: Any) -> int | None:
    if not bool(getattr(shape, "has_chart", False)):
        return None
    try:
        return len(shape.chart.series)
    except Exception:
        return None


def _chart_category_count(shape: Any) -> int | None:
    if not bool(getattr(shape, "has_chart", False)):
        return None
    try:
        plots = getattr(shape.chart, "plots", None)
        if plots is None or not plots:
            return None
        return len(plots[0].categories)
    except Exception:
        return None


def _sorted_findings(findings: list[PptxObjectFinding]) -> list[PptxObjectFinding]:
    return sorted(
        findings,
        key=lambda item: (
            -1 if item.slide_number is None else item.slide_number,
            item.code,
            item.severity,
            item.message,
            item.shape_name or "",
            item.scene_slide_id or "",
        ),
    )


def _summarize_findings(findings: list[PptxObjectFinding]) -> PptxObjectFindingsSummary:
    info_count = sum(1 for finding in findings if finding.severity == "info")
    warning_count = sum(1 for finding in findings if finding.severity == "warning")
    error_count = sum(1 for finding in findings if finding.severity == "error")
    return PptxObjectFindingsSummary(
        total_findings=len(findings),
        info_count=info_count,
        warning_count=warning_count,
        error_count=error_count,
        enforceable_count=warning_count + error_count,
        semantic_mismatch_count=sum(1 for finding in findings if finding.category == "semantic_mismatch"),
        editability_count=sum(1 for finding in findings if finding.category == "editability"),
        native_object_count=sum(1 for finding in findings if finding.category == "native_object"),
        unsupported_shape_count=sum(1 for finding in findings if finding.category == "unsupported_shape"),
        text_fit_count=sum(1 for finding in findings if finding.category == "text_fit"),
    )


def _mode_result(mode: ValidationMode, summary: PptxObjectFindingsSummary) -> Literal["passed", "issues_reported", "failed"]:
    if mode == "enforce" and summary.enforceable_count > 0:
        return "failed"
    if summary.total_findings > 0:
        return "issues_reported"
    return "passed"


def _normalize_for_stable_json(value: Any) -> Any:
    if isinstance(value, dict):
        return {key: _normalize_for_stable_json(value[key]) for key in sorted(value)}
    if isinstance(value, list):
        return [_normalize_for_stable_json(item) for item in value]
    if isinstance(value, float):
        if not math.isfinite(value):
            raise ValueError("validation payload cannot serialize non-finite floats")
        normalized = round(value, 6)
        return 0.0 if normalized == 0 else normalized
    return value


__all__ = [
    "OBJECT_VALIDATION_REPORT_VERSION",
    "PptxDeckSummary",
    "PptxObjectFinding",
    "PptxObjectFindingsSummary",
    "PptxObjectValidationReport",
    "PptxShapeSummary",
    "PptxSlideObjectInventory",
    "PptxSlideValidationResult",
    "SceneExpectationSummary",
    "SemanticObjectCounts",
    "summarize_pptx_object_validation",
    "validate_pptx_objects",
    "validate_pptx_objects_from_files",
    "validation_report_structural_hash",
    "validation_report_to_stable_json",
    "validation_report_to_stable_payload",
    "write_pptx_object_validation_report",
]
