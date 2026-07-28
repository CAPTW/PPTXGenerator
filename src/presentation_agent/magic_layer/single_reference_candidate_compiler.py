"""Compile one E01 editable PPT candidate from object graph primitives."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

SLIDE_W = 16.0
SLIDE_H = 9.0


def compile_single_reference_candidate(
    *,
    object_graph: dict[str, Any],
    output_pptx: Path,
    reference_image: Path,
) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Non-full-slide, bounded visual crop for the hero/photo field only.
    hero_crop = _create_hero_crop(reference_image, output_pptx.parent / "assets" / "hero_visual_field_crop.png")

    for node in object_graph["nodes"]:
        role = node["semantic_role"]
        x, y, w, h = _inches(node["bbox_norm"])
        if role == "background_base":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
            shape.name = node["object_id"]
            _fill(shape, "08111F")
            shape.line.fill.background()
        elif role == "hero_visual_field":
            pic = slide.shapes.add_picture(str(hero_crop), Inches(x), Inches(y), Inches(w), Inches(h))
            pic.name = node["object_id"]
        elif role in {"title_text", "subtitle_text", "step_number_text", "step_heading_text", "step_body_text", "badge_text", "source_footer_text"}:
            _add_text(slide, node, text_for_role(role, node["object_id"]))
        elif role in {"checklist_panel", "card_panel", "source_footer_strip"}:
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.name = node["object_id"]
            _fill(shape, "0F2A38" if role != "source_footer_strip" else "061526", transparency=8)
            _line(shape, "29C7E8" if role != "source_footer_strip" else "1C7FA0", width_pt=1.2)
        elif role == "semantic_icon":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.name = node["object_id"]
            _fill(shape, "2DD4FF")
            _line(shape, "D5F7FF", width_pt=0.8)
        elif role == "accent_line":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.name = node["object_id"]
            _fill(shape, "F4B43F")
            shape.line.fill.background()
        elif role == "technical_overlay":
            _add_overlay(slide)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "schema_name": "single_reference_candidate_compile_report",
        "status": "passed" if output_pptx.exists() else "failed",
        "editable_candidate_pptx": output_pptx.as_posix(),
        "slide_count": 1,
        "full_slide_reference_background": False,
        "screenshot_slide": False,
        "bounded_reference_crop_used": True,
        "semantic_raster_final_use_count": 0,
        "canva_parity_claimed": False,
    }


def text_for_role(role: str, object_id: str) -> str:
    if role == "title_text":
        return "CHECKLIST TITLE SLOT"
    if role == "subtitle_text":
        return "SUBTITLE / CONTEXT SLOT"
    if role == "step_number_text":
        suffix = object_id.rsplit("_", 1)[-1]
        return f"{int(suffix):02d}" if suffix.isdigit() else "##"
    if role == "step_heading_text":
        return "STEP HEADING SLOT"
    if role == "step_body_text":
        return "Step body slot"
    if role == "badge_text":
        return "BADGE SLOT"
    if role == "source_footer_text":
        return "SOURCE / FOOTER SLOT"
    return "TEXT SLOT"


def _create_hero_crop(reference_image: Path, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(reference_image) as image:
        image = image.convert("RGB")
        w, h = image.size
        crop = image.crop((0, 0, int(w * 0.36), int(h * 0.64)))
        crop.save(output_path)
    return output_path


def _inches(bbox_norm: list[float]) -> tuple[float, float, float, float]:
    x, y, w, h = bbox_norm
    return x * SLIDE_W, y * SLIDE_H, w * SLIDE_W, h * SLIDE_H


def _fill(shape: Any, color: str, *, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.fill.transparency = transparency


def _line(shape: Any, color: str, *, width_pt: float = 1.0) -> None:
    shape.line.color.rgb = RGBColor.from_string(color)
    shape.line.width = Pt(width_pt)


def _add_text(slide: Any, node: dict[str, Any], text: str) -> None:
    x, y, w, h = _inches(node["bbox_norm"])
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = node["object_id"]
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.name = "Aptos"
    run.font.bold = node["semantic_role"] in {"title_text", "step_heading_text", "badge_text"}
    run.font.size = Pt(_font_size(node["semantic_role"]))
    run.font.color.rgb = RGBColor.from_string("F8FAFC")


def _font_size(role: str) -> int:
    return {
        "title_text": 22,
        "subtitle_text": 10,
        "step_number_text": 16,
        "step_heading_text": 8,
        "step_body_text": 6,
        "badge_text": 6,
        "source_footer_text": 7,
    }.get(role, 8)


def _add_overlay(slide: Any) -> None:
    for i in range(8):
        x = 0.04 + i * 0.11
        y = 0.04 + (i % 2) * 0.025
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x * SLIDE_W), Inches(y * SLIDE_H), Inches(0.09), Inches(0.09))
        shape.name = f"technical_overlay_dot_{i+1}"
        _fill(shape, "2DD4FF", transparency=10)
        shape.line.fill.background()
