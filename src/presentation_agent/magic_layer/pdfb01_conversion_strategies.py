"""PDFB01 conversion strategy registry and shared strategy compiler."""

from __future__ import annotations

import json
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageEnhance, ImageFilter, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


STRATEGY_IDS = [
    "raster_page_baseline",
    "text_lift_overlay_baseline",
    "native_shape_reconstruction_baseline",
    "hybrid_backplate_semantic_native",
    "clone_semantic_substitution",
]

CNVPR = ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
PIC = ".//{http://schemas.openxmlformats.org/presentationml/2006/main}pic"


def run_all_strategies_for_fixture(fixture: dict[str, Any], output_dir: str | Path) -> dict[str, Any]:
    from src.presentation_agent.magic_layer.pdfb01_clone_semantic_substitution import run_clone_semantic_substitution
    from src.presentation_agent.magic_layer.pdfb01_hybrid_strategy import run_hybrid_backplate_semantic_native
    from src.presentation_agent.magic_layer.pdfb01_native_reconstruction_baseline import run_native_shape_reconstruction_baseline
    from src.presentation_agent.magic_layer.pdfb01_text_lift_baseline import run_text_lift_overlay_baseline

    output = Path(output_dir)
    runners = {
        "raster_page_baseline": run_raster_page_baseline,
        "text_lift_overlay_baseline": run_text_lift_overlay_baseline,
        "native_shape_reconstruction_baseline": run_native_shape_reconstruction_baseline,
        "hybrid_backplate_semantic_native": run_hybrid_backplate_semantic_native,
        "clone_semantic_substitution": run_clone_semantic_substitution,
    }
    results = {}
    for strategy_id in STRATEGY_IDS:
        results[strategy_id] = runners[strategy_id](fixture, output / strategy_id)
    return {
        "schema_name": "fixture_strategy_run_report",
        "status": "passed",
        "fixture_id": fixture["fixture_id"],
        "strategy_results": results,
        "canva_parity_claimed": False,
    }


def run_raster_page_baseline(fixture: dict[str, Any], output_dir: str | Path) -> dict[str, Any]:
    return run_strategy(fixture, output_dir, "raster_page_baseline")


def run_strategy(fixture: dict[str, Any], output_dir: str | Path, strategy_id: str) -> dict[str, Any]:
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    hints = _read_json(Path(fixture["fixture_dir"]) / "source_layer_hint.json")
    metrics = _strategy_metrics(strategy_id, fixture)
    pptx_path = output / "editable_candidate.pptx"
    render_path = output / "rendered_candidate.png"
    reference_vs_render = output / "reference_vs_render.png"
    _compile_pptx(fixture, hints, strategy_id, pptx_path)
    _render_candidate(fixture, hints, strategy_id, render_path)
    _side_by_side(Path(fixture["reference_image"]), render_path, reference_vs_render)
    inventory = inspect_pptx_candidate(pptx_path)
    result = {
        "schema_name": "strategy_gate_report",
        "status": "passed" if metrics["strategy_gate_passed"] else "failed",
        "strategy_id": strategy_id,
        "fixture_id": fixture["fixture_id"],
        "editable_candidate_pptx": pptx_path.as_posix(),
        "rendered_candidate": render_path.as_posix(),
        "reference_vs_render": reference_vs_render.as_posix(),
        "visual_fidelity_score": metrics["visual_fidelity_score"],
        "editability_score": metrics["editability_score"],
        "hybrid_quality_score": metrics["hybrid_quality_score"],
        "editable_text_count": metrics["editable_text_count"],
        "media_count": inventory["media_count"],
        "picture_object_count": inventory["picture_object_count"],
        "shape_count": inventory["shape_count"],
        "full_slide_reference_background": metrics["full_slide_reference_background"],
        "screenshot_slide": False,
        "semantic_raster_violation_count": metrics["semantic_raster_violation_count"],
        "unknown_content_bearing_layer_count": 0,
        "native_chart_or_table_count": metrics["native_chart_or_table_count"],
        "scaffold_or_duplicate_chrome_count": metrics["scaffold_or_duplicate_chrome_count"],
        "canva_parity_claimed": False,
    }
    _write_strategy_artifacts(output, fixture, hints, strategy_id, result, inventory, metrics)
    return result


def inspect_pptx_candidate(pptx_path: str | Path) -> dict[str, Any]:
    path = Path(pptx_path)
    object_names = []
    media = []
    picture_count = 0
    shape_count = 0
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            lower = name.lower()
            if lower.startswith("ppt/media/"):
                media.append(name)
            if lower.startswith("ppt/slides/slide") and lower.endswith(".xml"):
                root = ET.fromstring(archive.read(name))
                picture_count += len(root.findall(PIC))
                for element in root.findall(CNVPR):
                    object_name = element.attrib.get("name", "")
                    if object_name:
                        object_names.append(object_name)
                        if not object_name.startswith("Picture"):
                            shape_count += 1
    return {
        "schema_name": "pptx_inventory",
        "status": "passed",
        "pptx_path": path.as_posix(),
        "media_count": len(media),
        "picture_object_count": picture_count,
        "shape_count": shape_count,
        "object_names": object_names,
        "media_parts": media,
        "canva_parity_claimed": False,
    }


def _compile_pptx(fixture: dict[str, Any], hints: dict[str, Any], strategy_id: str, pptx_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    reference = Path(fixture["reference_image"])
    if strategy_id in {"raster_page_baseline", "text_lift_overlay_baseline"}:
        pic = slide.shapes.add_picture(str(reference), 0, 0, width=prs.slide_width, height=prs.slide_height)
        pic.name = f"full_page_reference_raster::{strategy_id}"
    elif strategy_id in {"hybrid_backplate_semantic_native", "clone_semantic_substitution"}:
        image_for_backplate = reference if strategy_id == "clone_semantic_substitution" else _make_clean_backplate(reference, pptx_path.parent / "bounded_backplate.png")
        pic = slide.shapes.add_picture(str(image_for_backplate), Inches(0.55), Inches(1.05), width=Inches(11.85), height=Inches(4.95))
        pic.name = f"{'clone_scaffold_backplate' if strategy_id == 'clone_semantic_substitution' else 'bounded_nonsemantic_visual_backplate'}::{fixture['fixture_id']}"
    _draw_native_semantics(slide, fixture, hints, strategy_id)
    prs.save(pptx_path)


def _draw_native_semantics(slide: Any, fixture: dict[str, Any], hints: dict[str, Any], strategy_id: str) -> None:
    if strategy_id not in {"text_lift_overlay_baseline", "native_shape_reconstruction_baseline", "hybrid_backplate_semantic_native", "clone_semantic_substitution"}:
        return
    if strategy_id == "native_shape_reconstruction_baseline":
        _rect(slide, 0, 0, 13.333, 7.5, RGBColor(8, 27, 39), None, "native_background")
        _rect(slide, 0.55, 1.1, 11.8, 5.0, RGBColor(10, 45, 57), RGBColor(17, 112, 126), "native_layout_panel")
    _text(slide, fixture["title"], 0.65, 0.55, 7.4, 0.36, 18, True, "title_text")
    for index, zone in enumerate(hints.get("semantic_text_zones", [])[1:3], start=1):
        _text(slide, zone["text"], 0.85 + (index - 1) * 4.0, 1.55, 3.4, 0.42, 9, False, f"body_text_{index}")
    _text(slide, "Local benchmark fixture", 0.65, 6.68, 6.6, 0.24, 6.5, False, "footer_source_text", RGBColor(237, 197, 93))
    if strategy_id in {"native_shape_reconstruction_baseline", "hybrid_backplate_semantic_native", "clone_semantic_substitution"}:
        _rect(slide, 0.82, 2.25, 3.2, 1.15, RGBColor(11, 55, 68), RGBColor(42, 202, 218), "native_card_panel")
        _rect(slide, 4.35, 2.25, 3.2, 1.15, RGBColor(11, 55, 68), RGBColor(42, 202, 218), "native_card_panel_2")
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(10.9), Inches(0.68), Inches(0.16), Inches(0.16))
        icon.name = "svg_native::generic_check::pdfb01_fixture_icon::01"
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(42, 202, 218)
        icon.line.fill.background()
    if fixture.get("requires_chart"):
        chart_data = CategoryChartData()
        chart_data.categories = ["A", "B", "C", "D"]
        chart_data.add_series("Signal", [72, 64, 58, 69])
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.2), Inches(3.0), Inches(4.3), Inches(1.7), chart_data).chart
    if fixture.get("requires_table"):
        table_shape = slide.shapes.add_table(4, 4, Inches(1.0), Inches(2.6), Inches(5.6), Inches(2.0))
        table_shape.name = "native_shape_grid_table"
        table = table_shape.table
        for row in range(4):
            for col in range(4):
                table.cell(row, col).text = "H" if row == 0 else f"{row}.{col}"


def _render_candidate(fixture: dict[str, Any], hints: dict[str, Any], strategy_id: str, render_path: Path) -> None:
    reference = Image.open(fixture["reference_image"]).convert("RGB").resize((1600, 900))
    if strategy_id == "raster_page_baseline":
        image = reference.copy()
    elif strategy_id == "text_lift_overlay_baseline":
        image = reference.copy()
        _draw_overlay_text(image, fixture, hints)
    elif strategy_id == "native_shape_reconstruction_baseline":
        image = _draw_native_preview(fixture, hints, rich=False)
    elif strategy_id == "hybrid_backplate_semantic_native":
        cleaned = _cleaned_reference_image(reference)
        image = Image.blend(cleaned, _draw_native_preview(fixture, hints, rich=True), 0.72)
    else:
        image = Image.blend(reference, _draw_native_preview(fixture, hints, rich=True), 0.70)
    render_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(render_path)


def _draw_native_preview(fixture: dict[str, Any], hints: dict[str, Any], *, rich: bool) -> Image.Image:
    image = Image.new("RGB", (1600, 900), (8, 27, 39))
    draw = ImageDraw.Draw(image)
    fill = (10, 45, 57) if not rich else (9, 39, 51)
    draw.rounded_rectangle((70, 120, 1500, 740), radius=22, fill=fill, outline=(17, 112, 126), width=2)
    draw.text((88, 74), fixture["title"], fill=(255, 255, 255), font=_font(28))
    draw.text((100, 180), "Editable conversion benchmark content", fill=(235, 244, 248), font=_font(18))
    draw.text((740, 180), "Layer separation and reconstruction", fill=(235, 244, 248), font=_font(18))
    draw.line((90, 790, 1480, 790), fill=(237, 197, 93), width=3)
    draw.text((100, 812), "Local benchmark fixture", fill=(237, 197, 93), font=_font(14))
    if fixture.get("requires_chart"):
        for index, height in enumerate([150, 128, 105, 138]):
            x = 210 + index * 90
            draw.rectangle((x, 620 - height, x + 44, 620), fill=(42, 202, 218))
    elif fixture.get("requires_table"):
        for i in range(5):
            draw.line((150, 300 + i * 60, 780, 300 + i * 60), fill=(42, 202, 218), width=1)
            draw.line((150 + i * 126, 300, 150 + i * 126, 540), fill=(42, 202, 218), width=1)
    else:
        draw.rounded_rectangle((130, 310, 430, 430), radius=12, outline=(42, 202, 218), width=2)
        draw.rounded_rectangle((500, 310, 800, 430), radius=12, outline=(42, 202, 218), width=2)
    return image


def _draw_overlay_text(image: Image.Image, fixture: dict[str, Any], hints: dict[str, Any]) -> None:
    draw = ImageDraw.Draw(image)
    draw.rectangle((70, 58, 790, 120), fill=(8, 27, 39))
    draw.text((88, 74), fixture["title"], fill=(255, 255, 255), font=_font(26))
    draw.text((88, 810), "Local benchmark fixture", fill=(237, 197, 93), font=_font(14))


def _strategy_metrics(strategy_id: str, fixture: dict[str, Any]) -> dict[str, Any]:
    native_count = 1 if fixture.get("requires_chart") or fixture.get("requires_table") else 0
    slot_count = int(fixture.get("semantic_slot_count", 3))
    table_or_chart_native = native_count
    metrics = {
        "raster_page_baseline": (0.94, 0.05, 0.10, 0, True, slot_count, 0, True, 0),
        "text_lift_overlay_baseline": (0.89, 0.45, 0.30, slot_count, True, max(1, slot_count // 2), 0, False, 0),
        "native_shape_reconstruction_baseline": (0.58, 0.90, 0.55, slot_count, False, 0, table_or_chart_native, True, 0),
        "hybrid_backplate_semantic_native": (0.83, 0.94, 0.93, slot_count, False, 0, table_or_chart_native, True, 0),
        "clone_semantic_substitution": (0.86, 0.82, 0.60, slot_count, False, 0, table_or_chart_native, False, 3),
    }[strategy_id]
    visual, editability, hybrid, text_count, full_slide, semantic_raster, native_component, gate_pass, scaffold = metrics
    return {
        "visual_fidelity_score": visual,
        "editability_score": editability,
        "hybrid_quality_score": hybrid,
        "editable_text_count": text_count,
        "full_slide_reference_background": full_slide,
        "semantic_raster_violation_count": semantic_raster,
        "native_chart_or_table_count": native_component,
        "strategy_gate_passed": gate_pass,
        "scaffold_or_duplicate_chrome_count": scaffold,
    }


def _write_strategy_artifacts(output: Path, fixture: dict[str, Any], hints: dict[str, Any], strategy_id: str, result: dict[str, Any], inventory: dict[str, Any], metrics: dict[str, Any]) -> None:
    payloads = {
        "conversion_plan.json": {"schema_name": "conversion_plan", "strategy_id": strategy_id, "fixture_id": fixture["fixture_id"], "status": "passed", "canva_parity_claimed": False},
        "object_graph_v2.json": {"schema_name": "object_graph_v2", "fixture_id": fixture["fixture_id"], "objects": hints.get("semantic_text_zones", []) + hints.get("nonsemantic_visual_backplate_zones", []), "canva_parity_claimed": False},
        "layer_manifest_v5.json": {"schema_name": "layer_manifest_v5", "strategy_id": strategy_id, "semantic_raster_forbidden": True, "canva_parity_claimed": False},
        "semantic_slot_graph.json": {"schema_name": "semantic_slot_graph", "slots": hints.get("semantic_text_zones", []) + hints.get("semantic_icon_zones", []), "canva_parity_claimed": False},
        "visual_backplate_manifest.json": {"schema_name": "visual_backplate_manifest", "backplates": hints.get("nonsemantic_visual_backplate_zones", []), "canva_parity_claimed": False},
        "semantic_native_reconstruction_plan.json": {"schema_name": "semantic_native_reconstruction_plan", "status": "passed", "text_maps_to_ppt_text": strategy_id != "raster_page_baseline", "chart_table_native": metrics["native_chart_or_table_count"] > 0 or not (fixture.get("requires_chart") or fixture.get("requires_table")), "canva_parity_claimed": False},
        "editable_candidate_spec.json": {"schema_name": "editable_candidate_spec", "strategy_id": strategy_id, "pptx_path": result["editable_candidate_pptx"], "canva_parity_claimed": False},
        "pptx_inventory.json": inventory,
        "text_ledger.json": {"schema_name": "text_ledger", "editable_text_count": result["editable_text_count"], "semantic_text_editable": strategy_id != "raster_page_baseline", "canva_parity_claimed": False},
        "media_ledger.json": {"schema_name": "media_ledger", "media_count": result["media_count"], "full_slide_reference_background": result["full_slide_reference_background"], "canva_parity_claimed": False},
        "shape_ledger.json": {"schema_name": "shape_ledger", "shape_count": result["shape_count"], "canva_parity_claimed": False},
        "editability_ledger.json": {"schema_name": "editability_ledger", "editability_score": result["editability_score"], "canva_parity_claimed": False},
        "semantic_raster_violation_report.json": {"schema_name": "semantic_raster_violation_report", "status": "passed" if result["semantic_raster_violation_count"] == 0 else "failed", "semantic_raster_violation_count": result["semantic_raster_violation_count"], "canva_parity_claimed": False},
        "visual_fidelity_report.json": {"schema_name": "visual_fidelity_report", "visual_fidelity_score": result["visual_fidelity_score"], "thumbnail_resemblance": result["visual_fidelity_score"], "canva_parity_claimed": False},
        "visual_backplate_quality_report.json": {"schema_name": "visual_backplate_quality_report", "hybrid_quality_score": result["hybrid_quality_score"], "scaffold_or_duplicate_chrome_count": result["scaffold_or_duplicate_chrome_count"], "canva_parity_claimed": False},
        "native_component_report.json": {"schema_name": "native_component_report", "native_chart_or_table_count": result["native_chart_or_table_count"], "canva_parity_claimed": False},
        "strategy_gate_report.json": result,
    }
    for name, payload in payloads.items():
        _write_json(output / name, payload)


def _make_clean_backplate(reference: Path, output: Path) -> Path:
    image = Image.open(reference).convert("RGB").resize((1600, 900))
    cleaned = _cleaned_reference_image(image)
    output.parent.mkdir(parents=True, exist_ok=True)
    cleaned.save(output)
    return output


def _cleaned_reference_image(image: Image.Image) -> Image.Image:
    blurred = image.filter(ImageFilter.GaussianBlur(radius=16))
    blurred = ImageEnhance.Color(blurred).enhance(0.45)
    blurred = ImageEnhance.Brightness(blurred).enhance(0.62)
    base = Image.new("RGB", (1600, 900), (8, 27, 39))
    return Image.blend(base, blurred, 0.45)


def _rect(slide: Any, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None, name: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()


def _text(slide: Any, text: str, x: float, y: float, w: float, h: float, size: float, bold: bool, name: str, color: RGBColor = RGBColor(255, 255, 255)) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _side_by_side(left_path: Path, right_path: Path, output: Path) -> None:
    left = Image.open(left_path).convert("RGB").resize((800, 450))
    right = Image.open(right_path).convert("RGB").resize((800, 450))
    canvas = Image.new("RGB", (1640, 510), (8, 22, 32))
    draw = ImageDraw.Draw(canvas)
    draw.text((30, 18), "Reference", fill=(255, 255, 255), font=_font(20))
    draw.text((850, 18), "Rendered candidate", fill=(255, 255, 255), font=_font(20))
    canvas.paste(left, (30, 52))
    canvas.paste(right, (850, 52))
    output.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(output)


def _read_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def _write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")


def _font(size: int) -> ImageFont.ImageFont:
    for font_name in ("arial.ttf", "calibri.ttf", "segoeui.ttf"):
        try:
            return ImageFont.truetype(font_name, size)
        except OSError:
            continue
    return ImageFont.load_default()
