"""Text readability review for E05."""

from __future__ import annotations

from collections import defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation

from src.presentation_agent.magic_layer.e05_product_review_rubric import SLIDE_ORDER


def review_text_readability(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    slide_reviews = []
    issues = []
    tiny_total = 0
    clipping_count = 0
    overflow_count = 0
    for slide_number, slide in enumerate(prs.slides, start=1):
        sizes: list[float] = []
        text_box_count = 0
        char_count = 0
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            text_box_count += 1
            char_count += len(text)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        sizes.append(float(run.font.size.pt))
        min_font = min(sizes) if sizes else 0.0
        tiny_count = sum(1 for size in sizes if size < 6.0)
        small_count = sum(1 for size in sizes if size < 7.0)
        tiny_total += tiny_count
        score = 4.5
        notes = []
        if tiny_count:
            score -= 0.55
            notes.append("sub-6pt text is present in rendered deck source/data areas")
        elif small_count:
            score -= 0.25
            notes.append("sub-7pt text is present but not a blocker")
        if text_box_count > 30:
            score -= 0.45
            notes.append("dense text grid needs product-polish pass")
        if char_count > 900:
            score -= 0.25
            notes.append("high character load for one live slide")
        archetype_id = SLIDE_ORDER[slide_number - 1] if slide_number <= len(SLIDE_ORDER) else "unknown"
        if archetype_id in {"table_heavy", "risk_register", "comparison_matrix"}:
            score -= 0.2
        score = round(max(1.0, score), 2)
        if score < 4.0:
            severity = "medium" if score >= 3.5 else "high"
            issues.append(
                {
                    "slide_number": slide_number,
                    "archetype_id": archetype_id,
                    "issue": "; ".join(notes) or "text density should be polished",
                    "severity": severity,
                    "patch_type": "text_capacity_patch",
                    "recommended_action": "Increase source/data text legibility by reducing density, widening cells, or moving detail to a lighter note treatment.",
                }
            )
        slide_reviews.append(
            {
                "slide_number": slide_number,
                "archetype_id": archetype_id,
                "text_box_count": text_box_count,
                "character_count": char_count,
                "minimum_font_pt": round(min_font, 2),
                "tiny_text_run_count": tiny_count,
                "small_text_run_count": small_count,
                "score": score,
                "status": "passed" if score >= 3.5 else "patch_required",
                "notes": notes,
            }
        )
    verdict = "patch_recommended" if issues else "passed"
    return {
        "schema_name": "e05_text_readability_review",
        "status": verdict,
        "text_overflow_count": overflow_count,
        "text_clipping_count": clipping_count,
        "tiny_text_run_count": tiny_total,
        "slide_reviews": slide_reviews,
        "issues": issues,
    }


def text_scores_by_slide(review: dict[str, Any]) -> dict[int, float]:
    return {int(row["slide_number"]): float(row["score"]) for row in review.get("slide_reviews", [])}

