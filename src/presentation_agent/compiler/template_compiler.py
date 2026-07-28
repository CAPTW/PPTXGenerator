"""Compile editable_template_spec.json into an editable PowerPoint preview deck."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..generator_contracts import validateEditableTemplateSpec
from .deck_compiler import _DeckStyle as _SharedDeckStyle
from .deck_compiler import _render_slide as _render_shared_deck_slide
from .decorative_budget import resolve_decorative_budget, update_budget_decisions_after_render
from .icon_resolver import IconResolver


DEFAULT_SPEC_PATH = Path("outputs/editable_template_spec.json")
DEFAULT_OUTPUT_PATH = Path("outputs/template_preview.pptx")
DEFAULT_MANIFEST_PATH = Path("outputs/template_preview_manifest.json")
DEFAULT_GOLDEN_OUTPUT_PATH = Path("outputs/golden_template_masters.pptx")
DEFAULT_GOLDEN_REPORT_JSON = Path("outputs/golden_template_masters_report.json")
DEFAULT_GOLDEN_REPORT_MD = Path("outputs/golden_template_masters_report.md")
DEFAULT_GOLDEN_MASTER_SPECS_DIR = Path("outputs/golden_master_specs")
DEFAULT_GOLDEN_CONTENT_FIXTURES_DIR = Path("fixtures/golden_master_content")
DEFAULT_TEMPLATE_CONTRACTS_DIR = Path("outputs/template_contracts")
DEFAULT_PHOTO_PLACEHOLDER_MANIFEST = Path("assets/photo_placeholders/photo_placeholder_manifest.json")
SAFE_MARGIN_IN = 0.4

FORBIDDEN_GOLDEN_VISIBLE_STRINGS = [
    "SAFE MARGINS",
    "Editable placeholder block",
    "title | text | title_block",
    "ARCHETYPE:",
    "DENSITY:",
    "FOOTER SYSTEM",
    "IMAGE FRAME ONLY",
    "PREVIEW WARNINGS",
    "layout-board-",
    "cards",
    "timeline_steps",
    "source_notes",
    "Source-backed evidence",
    "Golden template master sample content",
    "Hierarchy cue",
    "Editable motif",
    "Decision-ready placeholder",
    "Master Fixture",
    "Template Spec",
    "A canonical master",
    "Creative Academic Systems",
]


def compile_template_pack(editable_template_spec: dict[str, Any], output_path: str | Path) -> dict[str, Any]:
    """Write a template preview PPTX with one editable preview slide per layout."""

    validateEditableTemplateSpec(editable_template_spec)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    canvas = editable_template_spec["canvas"]
    presentation.slide_width = Inches(float(canvas["width"]))
    presentation.slide_height = Inches(float(canvas["height"]))
    blank_layout = presentation.slide_layouts[6]
    style = _Style(editable_template_spec)
    icon_resolver = IconResolver()

    compiled_layouts: list[dict[str, Any]] = []
    icon_reports: list[dict[str, Any]] = []
    for index, layout in enumerate(editable_template_spec["layouts"], start=1):
        slide = presentation.slides.add_slide(blank_layout)
        icon_report = icon_resolver.empty_report()
        _draw_preview_slide(slide, layout, style, index, icon_resolver=icon_resolver, icon_report=icon_report)
        icon_reports.append(icon_report)
        preview_style = _preview_deck_style(editable_template_spec, layout, index)
        compiled_layouts.append(
            {
                "slide_number": index,
                "layout_id": layout["layout_id"],
                "archetype_id": layout["archetype_id"],
                "layout_family_id": layout.get("layout_family_id"),
                "slot_count": len(layout["slots"]),
                "ornament_density_mode": preview_style.ornament_density,
                "rendering_path": "deck_compiler_component_primitives",
                "icons_used": icon_report["icons_used"],
                "missing_icons": icon_report["missing_icons"],
                "unresolved_icon_roles": icon_report["unresolved_icon_roles"],
                "icon_asset_paths": icon_report["icon_asset_paths"],
                "icon_family": icon_report["icon_family"],
            }
        )

    presentation.save(output)
    return {
        "schema_name": "template_preview_manifest",
        "schema_version": "1.0",
        "source_spec_path": _display_path(DEFAULT_SPEC_PATH),
        "pptx_path": _display_path(output),
        "slide_count": len(editable_template_spec["layouts"]),
        "rendering_path": "deck_compiler_component_primitives",
        "icon_report": icon_resolver.merge_reports(icon_reports),
        "compiled_layouts": compiled_layouts,
        "editable_preview_policy": {
            "text_slots_are_text_boxes": True,
            "cards_panels_rules_are_shapes": True,
            "tables_are_editable_preview_tables_or_shape_grids": True,
            "chart_slots_are_editable_chart_frames": True,
            "photo_slots_are_placeholders_inside_defined_frames": True,
            "no_full_slide_raster_background": True,
        },
    }


def compile_template_pack_from_files(
    *,
    spec_path: str | Path = DEFAULT_SPEC_PATH,
    output_path: str | Path = DEFAULT_OUTPUT_PATH,
    manifest_path: str | Path = DEFAULT_MANIFEST_PATH,
) -> Path:
    spec_file = Path(spec_path)
    spec = json.loads(spec_file.read_text(encoding="utf-8"))
    manifest = compile_template_pack(spec, output_path)
    manifest["source_spec_path"] = _display_path(spec_file)
    manifest_file = Path(manifest_path)
    manifest_file.parent.mkdir(parents=True, exist_ok=True)
    manifest_file.write_text(json.dumps(manifest, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    return Path(output_path)


def compile_golden_template_masters(
    editable_template_spec: dict[str, Any],
    output_path: str | Path,
    *,
    master_specs_dir: str | Path = DEFAULT_GOLDEN_MASTER_SPECS_DIR,
    content_fixtures_dir: str | Path = DEFAULT_GOLDEN_CONTENT_FIXTURES_DIR,
    template_contracts_dir: str | Path = DEFAULT_TEMPLATE_CONTRACTS_DIR,
) -> dict[str, Any]:
    """Write production-looking editable template master slides without debug overlays."""

    validateEditableTemplateSpec(editable_template_spec)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    master_specs = _load_golden_master_specs(master_specs_dir)
    content_fixtures = _load_golden_content_fixtures(content_fixtures_dir, editable_template_spec["layouts"])
    template_contracts = _load_template_contracts(template_contracts_dir, editable_template_spec["layouts"])

    presentation = Presentation()
    canvas = editable_template_spec["canvas"]
    presentation.slide_width = Inches(float(canvas["width"]))
    presentation.slide_height = Inches(float(canvas["height"]))
    blank_layout = presentation.slide_layouts[6]
    icon_resolver = IconResolver()

    compiled_layouts: list[dict[str, Any]] = []
    icon_reports: list[dict[str, Any]] = []
    for index, layout in enumerate(editable_template_spec["layouts"], start=1):
        slide = presentation.slides.add_slide(blank_layout)
        archetype_id = str(layout.get("archetype_id") or "")
        master_spec = master_specs.get(archetype_id) or _fallback_golden_master_spec(layout)
        fixture = content_fixtures[archetype_id]
        contract = template_contracts[archetype_id]
        contract_plan = _contract_aware_fixture(layout, fixture, contract)
        binding = _golden_binding(layout, index, master_spec, contract, contract_plan)
        deck_style = _SharedDeckStyle(editable_template_spec).for_tone(str(binding["selected_tone_variant"]))
        deck_style = deck_style.for_layout(layout, binding)
        warnings: list[dict[str, Any]] = list(contract_plan["warnings"])
        icon_report = icon_resolver.empty_report()
        _render_shared_deck_slide(
            slide,
            _golden_source(layout, index, master_spec, contract_plan["fixture"]),
            layout,
            binding,
            deck_style,
            warnings,
            icon_resolver=icon_resolver,
            icon_report=icon_report,
        )
        icon_reports.append(icon_report)
        render_stats = _slide_render_stats(slide)
        decorative_budget_plan = update_budget_decisions_after_render(
            binding.get("decorative_budget_plan") or {},
            shape_count=render_stats["shape_count"],
        )
        shape_budget_status = _shape_budget_status(render_stats, contract)
        if shape_budget_status["status"] != "ok":
            warnings.append(
                _golden_renderer_warning(
                    "SHAPE_BUDGET_EXCEEDED",
                    "Shape count exceeds the template contract budget.",
                    {"shape_count": render_stats["shape_count"], "limit": shape_budget_status["limit"]},
                )
            )
        decoration_budget_status = _decoration_budget_status(render_stats, contract)
        if decoration_budget_status["status"] != "ok":
            warnings.append(
                _golden_renderer_warning(
                    "DECORATION_BUDGET_PRESSURE",
                    "Decorative object volume may compete with text zones.",
                    {
                        "decoration_to_text_ratio": decoration_budget_status["decoration_to_text_ratio"],
                        "limit": decoration_budget_status["limit"],
                    },
                )
            )
        layout_reference = master_spec.get("layout_reference") if isinstance(master_spec.get("layout_reference"), dict) else {}
        priority_ref_required = bool(layout_reference.get("priority_ref_required") or layout_reference.get("priority_archetype"))
        priority_ref_available = bool(layout_reference.get("priority_ref_available"))
        reference_source = _golden_reference_source(layout_reference)
        compiled_layouts.append(
            {
                "slide_number": index,
                "layout_id": layout["layout_id"],
                "archetype_id": layout["archetype_id"],
                "layout_family_id": layout.get("layout_family_id"),
                "master_spec_path": master_spec.get("_spec_path"),
                "contract_used": contract.get("_contract_path"),
                "content_fixture_path": fixture.get("_fixture_path"),
                "fixture_used": fixture.get("_fixture_path"),
                "text_capacity_status": contract_plan["text_capacity_status"],
                "reading_order_status": contract_plan["reading_order_status"],
                "decoration_budget_status": decoration_budget_status,
                "shape_budget_status": shape_budget_status,
                "render_object_stats": render_stats,
                "decorative_budget_plan": decorative_budget_plan,
                "identity_priority": master_spec.get("identity_priority"),
                "identity_lock": master_spec.get("identity_lock"),
                "required_visual_signature": master_spec.get("required_visual_signature"),
                "forbidden_drift": master_spec.get("forbidden_drift"),
                "minimum_visual_features": master_spec.get("minimum_visual_features"),
                "editable_translation": master_spec.get("editable_translation"),
                "refinement_patch_directives": master_spec.get("refinement_patch_directives") or [],
                "layout_reference": layout_reference,
                "priority_ref_required": priority_ref_required,
                "priority_ref_available": priority_ref_available,
                "reference_source": reference_source,
                "expected_density": master_spec.get("expected_density"),
                "palette_mode": master_spec.get("palette_mode"),
                "selected_tone_variant": deck_style.tone_variant,
                "ornament_density_mode": deck_style.ornament_density,
                "rendering_path": "deck_compiler_component_primitives",
                "icons_used": icon_report["icons_used"],
                "missing_icons": icon_report["missing_icons"],
                "unresolved_icon_roles": icon_report["unresolved_icon_roles"],
                "icon_asset_paths": icon_report["icon_asset_paths"],
                "icon_family": icon_report["icon_family"],
                "warnings": warnings,
                "remaining_warnings": warnings,
            }
        )

    presentation.save(output)
    return {
        "schema_name": "golden_template_masters_report",
        "schema_version": "1.0",
        "status": "created",
        "source_spec_path": _display_path(DEFAULT_SPEC_PATH),
        "golden_master_specs_dir": _display_path(Path(master_specs_dir)),
        "golden_content_fixtures_dir": _display_path(Path(content_fixtures_dir)),
        "template_contracts_dir": _display_path(Path(template_contracts_dir)),
        "photo_placeholder_manifest_path": _display_path(DEFAULT_PHOTO_PLACEHOLDER_MANIFEST),
        "canonical_master_specs_loaded": len(master_specs),
        "content_fixtures_loaded": len(content_fixtures),
        "template_contracts_loaded": len(template_contracts),
        "pptx_path": _display_path(output),
        "slide_count": len(compiled_layouts),
        "rendering_path": "deck_compiler_component_primitives",
        "icon_report": icon_resolver.merge_reports(icon_reports),
        "debug_overlay_visible": False,
        "forbidden_visible_strings": FORBIDDEN_GOLDEN_VISIBLE_STRINGS,
        "priority_ref_summary": _priority_ref_summary(compiled_layouts),
        "compiled_layouts": compiled_layouts,
        "editable_master_policy": {
            "text_slots_are_text_boxes": True,
            "cards_panels_rules_are_shapes": True,
            "tables_are_editable_preview_tables_or_shape_grids": True,
            "chart_slots_are_editable_chart_frames": True,
            "photo_slots_are_placeholders_inside_defined_frames": True,
            "photo_placeholder_provenance_recorded": DEFAULT_PHOTO_PLACEHOLDER_MANIFEST.exists(),
            "no_full_slide_raster_background": True,
            "design_board_image_not_inserted": True,
        },
    }


def compile_golden_template_masters_from_files(
    *,
    spec_path: str | Path = DEFAULT_SPEC_PATH,
    output_path: str | Path = DEFAULT_GOLDEN_OUTPUT_PATH,
    report_json_path: str | Path = DEFAULT_GOLDEN_REPORT_JSON,
    report_md_path: str | Path = DEFAULT_GOLDEN_REPORT_MD,
    master_specs_dir: str | Path = DEFAULT_GOLDEN_MASTER_SPECS_DIR,
    content_fixtures_dir: str | Path = DEFAULT_GOLDEN_CONTENT_FIXTURES_DIR,
    template_contracts_dir: str | Path = DEFAULT_TEMPLATE_CONTRACTS_DIR,
) -> Path:
    spec_file = Path(spec_path)
    spec = json.loads(spec_file.read_text(encoding="utf-8"))
    report = compile_golden_template_masters(
        spec,
        output_path,
        master_specs_dir=master_specs_dir,
        content_fixtures_dir=content_fixtures_dir,
        template_contracts_dir=template_contracts_dir,
    )
    report["source_spec_path"] = _display_path(spec_file)
    report_file = Path(report_json_path)
    report_file.parent.mkdir(parents=True, exist_ok=True)
    report_file.write_text(json.dumps(report, indent=2, sort_keys=True, ensure_ascii=True) + "\n", encoding="utf-8")
    Path(report_md_path).write_text(_golden_markdown_report(report), encoding="utf-8")
    return Path(output_path)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Compile editable_template_spec.json into outputs/template_preview.pptx.")
    parser.add_argument("--spec", type=Path, default=DEFAULT_SPEC_PATH)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT_PATH)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST_PATH)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        output = compile_template_pack_from_files(spec_path=args.spec, output_path=args.output, manifest_path=args.manifest)
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        print(f"COMPILE_TEMPLATE_PACK_FAILED {exc}")
        return 1
    print(f"WROTE {output}")
    return 0


def _golden_source(layout: dict[str, Any], slide_number: int, master_spec: dict[str, Any], fixture: dict[str, Any]) -> dict[str, Any]:
    archetype_id = str(layout.get("archetype_id") or "template")
    title, subtitle = _golden_title_pair(archetype_id, fixture)
    content_blocks = _golden_content_blocks(layout, archetype_id, fixture)
    return {
        "slide_id": f"golden-master-{slide_number:03d}",
        "section_id": f"S{max(1, min(slide_number, 18)):02d}",
        "slide_type": archetype_id,
        "title": title,
        "subtitle": subtitle,
        "content_density": str(layout.get("density") or "medium"),
        "strict_slot_content": True,
        "content_blocks": content_blocks,
        "chart_data": fixture.get("chart_data") or {
            "categories": ["Govern", "Measure", "Adapt", "Scale"],
            "series": [
                {"name": "Readiness", "values": [82, 67, 74, 88]},
                {"name": "Evidence", "values": [58, 71, 64, 79]},
            ],
        },
        "table_data": fixture.get("table_data") or {
            "headers": ["Signal", "Evidence", "Action"],
            "rows": [
                ["Pattern", "Review cue", "Align"],
                ["Hierarchy", "Clear focal area", "Prioritize"],
                ["Rhythm", "Measured cadence", "Stabilize"],
                ["Closure", "Visible decision", "Align"],
            ],
        },
        "image_needs": [{"slot": "image", "usage": "declared_photo_frame_only"}],
        "speaker_notes": str(fixture.get("speaker_notes") or "Synthetic sample content tests editable template capacity."),
        "citations": fixture.get("citations") or [{"label": "Sample Basis", "source": "Synthetic capacity sample"}],
        "design_intent": {
            "golden_template_master": True,
            "layout_family_id": layout.get("layout_family_id"),
            "canonical_archetype_identity": master_spec.get("identity_priority"),
        },
    }


def _golden_binding(
    layout: dict[str, Any],
    slide_number: int,
    master_spec: dict[str, Any],
    contract: dict[str, Any] | None = None,
    contract_plan: dict[str, Any] | None = None,
) -> dict[str, Any]:
    archetype_id = str(layout.get("archetype_id") or "template")
    decorative_budget = (contract or {}).get("decorative_budget") if isinstance((contract or {}).get("decorative_budget"), dict) else {}
    budget_plan = resolve_decorative_budget(
        layout,
        contract,
        binding={
            "slide_type": archetype_id,
            "layout_family_id": layout.get("layout_family_id"),
            "component_density": str(master_spec.get("expected_density") or layout.get("density") or "medium"),
        },
        slot_capacity_status=(contract_plan or {}).get("text_capacity_status") if isinstance((contract_plan or {}).get("text_capacity_status"), dict) else {},
    )
    slot_bindings = {
        str(slot.get("slot_id")): _slot_binding(slot)
        for slot in layout.get("slots") or []
        if isinstance(slot, dict) and slot.get("slot_id")
    }
    return {
        "slide_id": f"golden-master-{slide_number:03d}",
        "slide_type": archetype_id,
        "selected_layout_id": layout["layout_id"],
        "layout_family_id": layout.get("layout_family_id"),
        "selected_tone_variant": str(master_spec.get("palette_mode") or _preview_tone(layout)),
        "ornament_density": str(master_spec.get("ornament_system", {}).get("density_mode") or layout.get("density") or "medium"),
        "component_density": str(master_spec.get("expected_density") or layout.get("density") or "medium"),
        "layout_reference_source": _golden_reference_source(master_spec.get("layout_reference") if isinstance(master_spec.get("layout_reference"), dict) else {}),
        "identity_lock": str(master_spec.get("identity_lock") or "standard"),
        "required_visual_signature": master_spec.get("required_visual_signature") or {},
        "forbidden_drift": master_spec.get("forbidden_drift") or [],
        "minimum_visual_features": master_spec.get("minimum_visual_features") or {},
        "editable_translation": master_spec.get("editable_translation") or {},
        "reference_geometry": master_spec.get("reference_geometry") or {},
        "component_grammar": master_spec.get("component_grammar") or {},
        "refinement_patch_directives": master_spec.get("refinement_patch_directives") or [],
        "layout_ref_extraction": master_spec.get("layout_ref_extraction") or {},
        "deck_scale": "golden_template_master",
        "text_first_contract": True,
        "presentation_role": (contract or {}).get("presentation_role"),
        "reading_order": (contract or {}).get("reading_order") or [],
        "content_capacity": (contract or {}).get("content_capacity") or {},
        "decorative_budget": {
            **decorative_budget,
            "max_ornament_density": budget_plan["allowed_ornament_density"],
            "max_shape_count_target": budget_plan["max_shape_count_target"],
            "max_background_coverage": budget_plan["max_background_coverage"],
        },
        "decorative_budget_plan": budget_plan,
        "contract_plan": contract_plan or {},
        "section_rhythm_role": "divider" if archetype_id == "section_divider" else "overview",
        "slot_bindings": slot_bindings,
        "component_bindings": {
            str(slot.get("slot_id")): slot.get("component_id")
            for slot in layout.get("slots") or []
            if isinstance(slot, dict) and slot.get("slot_id")
        },
        "fallback_used": False,
        "fallback_reason": None,
        "production_master_rendering_path": "deck_compiler_component_primitives",
    }


def _golden_content_blocks(layout: dict[str, Any], archetype_id: str, fixture: dict[str, Any]) -> list[dict[str, Any]]:
    fixture_blocks = fixture.get("content_blocks")
    if isinstance(fixture_blocks, list) and all(isinstance(block, dict) for block in fixture_blocks):
        return [dict(block) for block in fixture_blocks]
    labels = {
        "creative_cover": ["Design thesis", "Audience signal", "Evidence promise"],
        "visual_table_of_contents": ["Context", "Method", "Evidence", "Decision"],
        "section_divider": ["Section thesis", "Research direction"],
        "research_overview": ["Context signal", "Evidence base", "Strategic question"],
        "problem_statement": ["Constraint", "Impact", "Decision point"],
        "research_gap": ["Known", "Unknown", "Opportunity"],
        "literature_map": ["Theme cluster", "Evidence lineage", "Interpretive gap"],
        "methodology_framework": ["Frame", "Measure", "Synthesize", "Govern"],
        "technical_flow_chart": ["Input", "Transform", "Validate", "Output"],
        "work_support_sequence": ["Prepare", "Execute", "Review", "Improve"],
        "photo_caption_grid": ["Observation", "Context", "Source cue"],
        "comparison_matrix": ["Option A", "Option B", "Trade-off"],
        "concept_relationship_venn": ["Overlap", "Boundary", "Dependency"],
        "three_level_explanation": ["Principle", "Mechanism", "Implication"],
        "circular_process": ["Sense", "Decide", "Act", "Learn"],
        "kpi_donut_chart": ["Leading signal", "Current state", "Target band"],
        "timeline_roadmap": ["Now", "Next", "Later"],
        "data_table_appendix": ["Finding", "Basis", "Owner"],
    }
    default_items = labels.get(archetype_id, ["Insight", "Evidence", "Action"])
    slots = [slot for slot in layout.get("slots") or [] if isinstance(slot, dict) and slot.get("slot_type") == "content"]
    blocks: list[dict[str, Any]] = []
    for index, slot in enumerate(slots[:8]):
        slot_id = str(slot.get("slot_id") or f"module_{index + 1}")
        title = default_items[index % len(default_items)]
        if _golden_slot_needs_multiple_items(slot_id, str(slot.get("component_id") or "")):
            for item_index in range(6):
                blocks.append(
                    {
                        "block_id": f"golden-{archetype_id}-{index + 1}-{item_index + 1}",
                        "slot": slot_id,
                        "title": default_items[item_index % len(default_items)],
                "content": [
                    "Evidence signal",
                    "Analytic implication",
                    "Decision signal",
                ],
                    }
                )
            continue
        blocks.append(
            {
                "block_id": f"golden-{archetype_id}-{index + 1}",
                "slot": slot_id,
                "title": title,
                "content": [
                    "Evidence remains visible",
                    "Rationale stays inspectable",
                    "Actionable synthesis",
                ],
            }
        )
    if not blocks:
        blocks.append(
            {
                "block_id": f"golden-{archetype_id}-summary",
                "slot": "body",
                "title": "Core message",
                "content": ["Evidence-led narrative", "Inspectable decision system"],
            }
        )
    return blocks


def _golden_slot_needs_multiple_items(slot_id: str, component_id: str) -> bool:
    normalized = slot_id.lower()
    component = component_id.lower()
    return (
        normalized in {"cards", "caption_cards", "index_navigation", "progress_markers", "step_cards"}
        or "card" in normalized
        or component in {"card", "cards", "layered_card", "glass_overlay_card", "index_navigation"}
    )


def _golden_title_pair(archetype_id: str, fixture: dict[str, Any] | None = None) -> tuple[str, str]:
    if isinstance(fixture, dict) and (fixture.get("title") or fixture.get("subtitle")):
        return str(fixture.get("title") or archetype_id.replace("_", " ").title()), str(fixture.get("subtitle") or "Editable master layout")
    titles = {
        "creative_cover": ("Creative Academic Systems", "Evidence-led strategy for complex decisions"),
        "visual_table_of_contents": ("Navigation Map", "A structured path from context to action"),
        "section_divider": ("Evidence Architecture", "A new section begins with a clear thesis"),
        "research_overview": ("Research Context", "Signals, sources, and implications in one view"),
        "problem_statement": ("Problem Frame", "Where complexity blocks reliable execution"),
        "research_gap": ("Research Gap", "What the evidence does not yet resolve"),
        "literature_map": ("Literature Map", "Themes, claims, and source lineage"),
        "methodology_framework": ("Method Framework", "How the analysis is structured"),
        "technical_flow_chart": ("Technical Flow", "System logic from input to decision"),
        "work_support_sequence": ("Work Sequence", "Operating rhythm and support model"),
        "photo_caption_grid": ("Field Evidence", "Visual observations with source context"),
        "comparison_matrix": ("Comparison Matrix", "Trade-offs across decision criteria"),
        "concept_relationship_venn": ("Concept Relationships", "Shared boundaries and dependencies"),
        "three_level_explanation": ("Three-Level Explanation", "Principle, mechanism, implication"),
        "circular_process": ("Learning Loop", "Sense, decide, act, and adapt"),
        "kpi_donut_chart": ("KPI Dashboard", "Directional signals and target bands"),
        "timeline_roadmap": ("Timeline Roadmap", "Near-term sequence and milestones"),
        "data_table_appendix": ("Data Appendix", "Compact evidence for detailed review"),
    }
    return titles.get(archetype_id, (archetype_id.replace("_", " ").title(), "Editable master layout"))


def _load_golden_master_specs(master_specs_dir: str | Path) -> dict[str, dict[str, Any]]:
    directory = Path(master_specs_dir)
    if not directory.exists():
        return {}
    specs: dict[str, dict[str, Any]] = {}
    for path in sorted(directory.glob("*.json")):
        payload = json.loads(path.read_text(encoding="utf-8"))
        if not isinstance(payload, dict):
            raise ValueError(f"golden master spec {path} must be a JSON object")
        archetype_id = str(payload.get("archetype_id") or "")
        if not archetype_id:
            raise ValueError(f"golden master spec {path} is missing archetype_id")
        missing = sorted(_required_golden_master_spec_fields() - set(payload))
        if missing:
            raise ValueError(f"golden master spec {path} missing fields: {', '.join(missing)}")
        payload = dict(payload)
        payload["_spec_path"] = _display_path(path)
        specs[archetype_id] = payload
    return specs


def _load_golden_content_fixtures(content_fixtures_dir: str | Path, layouts: list[dict[str, Any]]) -> dict[str, dict[str, Any]]:
    directory = Path(content_fixtures_dir)
    if not directory.exists():
        raise ValueError(f"golden content fixtures directory is missing: {directory}")
    fixtures: dict[str, dict[str, Any]] = {}
    for layout in layouts:
        archetype_id = str(layout.get("archetype_id") or "")
        if not archetype_id:
            continue
        path = directory / f"{archetype_id}.json"
        if not path.exists():
            raise ValueError(f"golden content fixture is missing for {archetype_id}: {path}")
        payload = json.loads(path.read_text(encoding="utf-8"))
        if not isinstance(payload, dict):
            raise ValueError(f"golden content fixture {path} must be a JSON object")
        if str(payload.get("archetype_id") or "") != archetype_id:
            raise ValueError(f"golden content fixture {path} has archetype_id {payload.get('archetype_id')!r}, expected {archetype_id!r}")
        missing = [field for field in ("title", "content_blocks", "citations", "speaker_notes") if field not in payload]
        if missing:
            raise ValueError(f"golden content fixture {path} missing fields: {', '.join(missing)}")
        blocks = payload.get("content_blocks")
        if not isinstance(blocks, list) or not all(isinstance(block, dict) for block in blocks):
            raise ValueError(f"golden content fixture {path} must define content_blocks as a list of objects")
        visible_text = "\n".join(_fixture_visible_strings(payload))
        hits = [phrase for phrase in FORBIDDEN_GOLDEN_VISIBLE_STRINGS if _contains_forbidden(visible_text, phrase)]
        if hits:
            raise ValueError(f"golden content fixture {path} contains forbidden visible phrase(s): {', '.join(sorted(hits))}")
        payload = dict(payload)
        payload["_fixture_path"] = _display_path(path)
        fixtures[archetype_id] = payload
    return fixtures


def _load_template_contracts(template_contracts_dir: str | Path, layouts: list[dict[str, Any]]) -> dict[str, dict[str, Any]]:
    directory = Path(template_contracts_dir)
    if not directory.exists():
        raise ValueError(f"template contracts directory is missing: {directory}")
    contracts: dict[str, dict[str, Any]] = {}
    for layout in layouts:
        archetype_id = str(layout.get("archetype_id") or "")
        if not archetype_id:
            continue
        path = directory / f"{archetype_id}.contract.json"
        if not path.exists():
            raise ValueError(f"template usability contract is missing for {archetype_id}: {path}")
        payload = json.loads(path.read_text(encoding="utf-8"))
        if not isinstance(payload, dict):
            raise ValueError(f"template usability contract {path} must be a JSON object")
        if str(payload.get("archetype_id") or "") != archetype_id:
            raise ValueError(f"template usability contract {path} has archetype_id {payload.get('archetype_id')!r}, expected {archetype_id!r}")
        for field in ("presentation_role", "reading_order", "content_capacity", "decorative_budget"):
            if field not in payload:
                raise ValueError(f"template usability contract {path} missing field: {field}")
        payload = dict(payload)
        payload["_contract_path"] = _display_path(path)
        contracts[archetype_id] = payload
    return contracts


def _contract_aware_fixture(layout: dict[str, Any], fixture: dict[str, Any], contract: dict[str, Any]) -> dict[str, Any]:
    shaped = json.loads(json.dumps(fixture))
    warnings: list[dict[str, Any]] = []
    capacity = contract.get("content_capacity") if isinstance(contract.get("content_capacity"), dict) else {}
    overflow = contract.get("overflow_policy") if isinstance(contract.get("overflow_policy"), dict) else {}
    title_status = _text_field_capacity_status("title", shaped.get("title"), int(capacity.get("title_max_chars") or 0), overflow, warnings)
    subtitle_status = _text_field_capacity_status("subtitle", shaped.get("subtitle"), int(capacity.get("subtitle_max_chars") or 0), overflow, warnings)
    body_limit = int(capacity.get("body_bullet_count") or 0)
    body_line_limit = int(capacity.get("body_bullet_max_chars") or 0)
    card_limit = int(capacity.get("card_count") or 0)
    card_title_limit = int(capacity.get("card_title_max_chars") or 0)
    card_body_limit = int(capacity.get("card_body_max_chars") or 0)
    ordered_blocks = _blocks_in_contract_order(shaped.get("content_blocks") or [], contract)
    shaped_blocks, block_status = _shape_blocks_to_capacity(
        ordered_blocks,
        layout,
        body_limit=body_limit,
        body_line_limit=body_line_limit,
        card_limit=card_limit,
        card_title_limit=card_title_limit,
        card_body_limit=card_body_limit,
        subtitle=shaped.get("subtitle"),
        warnings=warnings,
    )
    shaped["content_blocks"] = shaped_blocks
    table_status = _shape_table_to_capacity(shaped, capacity, warnings)
    chart_status = _shape_chart_to_capacity(shaped, capacity, warnings)
    return {
        "fixture": shaped,
        "warnings": warnings,
        "text_capacity_status": {
            "status": "ok" if not [w for w in warnings if str(w.get("code", "")).startswith(("TITLE_", "SUBTITLE_", "BODY_", "CARD_", "TABLE_", "CHART_"))] else "adjusted",
            "title": title_status,
            "subtitle": subtitle_status,
            "body": block_status,
            "table": table_status,
            "chart": chart_status,
            "presentation_role": contract.get("presentation_role"),
        },
        "reading_order_status": _reading_order_contract_status(layout, contract),
    }


def _text_field_capacity_status(field: str, value: Any, limit: int, overflow: dict[str, Any], warnings: list[dict[str, Any]]) -> dict[str, Any]:
    text = str(value or "")
    status = {"status": "ok", "chars": len(text), "limit": limit}
    if limit and len(text) > limit:
        status["status"] = "warning"
        warnings.append(
            _golden_renderer_warning(
                f"{field.upper()}_ZONE_CAPACITY_WARNING",
                f"{field.title()} text exceeds the template contract capacity.",
                {"chars": len(text), "limit": limit, "overflow_policy": overflow},
            )
        )
    return status


def _blocks_in_contract_order(blocks: list[Any], contract: dict[str, Any]) -> list[dict[str, Any]]:
    order = [str(item) for item in (contract.get("reading_order") or [])]
    rank = {slot_id: index for index, slot_id in enumerate(order)}
    normalized = [dict(block) for block in blocks if isinstance(block, dict)]
    return sorted(normalized, key=lambda block: (rank.get(str(block.get("slot") or ""), 999), normalized.index(block) if block in normalized else 999))


def _shape_blocks_to_capacity(
    blocks: list[dict[str, Any]],
    layout: dict[str, Any],
    *,
    body_limit: int,
    body_line_limit: int,
    card_limit: int,
    card_title_limit: int,
    card_body_limit: int,
    subtitle: Any,
    warnings: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    slot_components = {
        str(slot.get("slot_id")): str(slot.get("component_id") or "")
        for slot in layout.get("slots") or []
        if isinstance(slot, dict)
    }
    archetype_id = str(layout.get("archetype_id") or "")
    effective_card_limit = max(card_limit, 6) if archetype_id == "research_overview" else card_limit
    budget = 999 if archetype_id == "research_overview" else body_limit if body_limit > 0 else 999
    used_lines = 1 if body_limit > 0 and len(str(subtitle or "").strip()) >= 8 else 0
    shaped: list[dict[str, Any]] = []
    card_count = 0
    compressed_blocks = 0
    dropped_blocks = 0
    for block in blocks:
        slot_id = str(block.get("slot") or "")
        component_id = slot_components.get(slot_id, "")
        card_like = _is_card_like_slot(slot_id, component_id)
        if card_like and effective_card_limit > 0 and card_count >= effective_card_limit:
            dropped_blocks += 1
            continue
        next_block = dict(block)
        title_limit = card_title_limit if card_like else body_line_limit
        if title_limit and isinstance(next_block.get("title"), str) and len(str(next_block["title"])) > title_limit:
            next_block["title"] = _compress_visible_text(str(next_block["title"]), title_limit)
            compressed_blocks += 1
        title_line_count = 1 if len(str(next_block.get("title") or "").strip()) >= 8 else 0
        remaining = budget - used_lines - title_line_count
        content = next_block.get("content")
        content_items = [str(item) for item in content] if isinstance(content, list) else ([str(content)] if isinstance(content, str) and content else [])
        max_body_chars = card_body_limit if card_like and card_body_limit else body_line_limit
        compact_items: list[str] = []
        if card_like and archetype_id == "research_overview":
            compact_items = [_compress_visible_text(item, max_body_chars) for item in content_items[:1] if item]
        elif body_limit <= 0:
            compact_items = [_compress_visible_text(item, max_body_chars) for item in content_items if item][:1 if card_like else len(content_items)]
        elif not card_like and remaining > 0:
            compact_items = [_compress_visible_text(item, max_body_chars) for item in content_items[:remaining] if item]
        elif card_like and remaining > 1 and len(shaped) < max(1, body_limit // 2):
            compact_items = [_compress_visible_text(item, max_body_chars) for item in content_items[:1] if item]
        if len(compact_items) < len(content_items):
            compressed_blocks += 1
        next_block["content"] = compact_items
        block_lines = title_line_count + sum(1 for item in compact_items if len(item.strip()) >= 8)
        if body_limit > 0 and used_lines + block_lines > budget and title_line_count:
            next_block["title"] = _short_capacity_label(str(next_block.get("title") or ""))
            block_lines = sum(1 for item in compact_items if len(item.strip()) >= 8)
            compressed_blocks += 1
        if body_limit > 0 and used_lines + block_lines > budget:
            dropped_blocks += 1
            continue
        used_lines += block_lines
        if card_like:
            card_count += 1
        shaped.append(next_block)
    if compressed_blocks:
        warnings.append(
            _golden_renderer_warning(
                "BODY_CONTENT_COMPRESSED_TO_CONTRACT",
                "Fixture body text was compressed before rendering to respect the template contract.",
                {"compressed_blocks": compressed_blocks, "body_line_budget": body_limit, "planned_body_lines": used_lines},
            )
        )
    if dropped_blocks:
        warnings.append(
            _golden_renderer_warning(
                "BODY_CONTENT_MOVED_TO_NOTES",
                "Fixture blocks were omitted from visible golden master text and retained as note-level sample detail.",
                {"dropped_blocks": dropped_blocks, "body_line_budget": body_limit},
            )
        )
    if effective_card_limit != card_limit:
        warnings.append(
            _golden_renderer_warning(
                "CARD_DENSITY_FIDELITY_EXCEPTION",
                "One additional concise evidence module was retained to preserve the reference visual family.",
                {"contract_card_budget": card_limit, "effective_card_budget": effective_card_limit},
            )
        )
    return shaped, {
        "status": "ok" if not compressed_blocks and not dropped_blocks else "adjusted",
        "body_line_budget": body_limit,
        "planned_body_lines": used_lines,
        "visible_block_count": len(shaped),
        "card_budget": card_limit,
        "effective_card_budget": effective_card_limit,
        "planned_card_count": card_count,
        "compressed_blocks": compressed_blocks,
        "dropped_blocks": dropped_blocks,
    }


def _shape_table_to_capacity(fixture: dict[str, Any], capacity: dict[str, Any], warnings: list[dict[str, Any]]) -> dict[str, Any]:
    data = fixture.get("table_data")
    if not isinstance(data, dict):
        return {"status": "not_applicable"}
    max_rows = int(capacity.get("table_max_rows") or 0)
    max_cols = int(capacity.get("table_max_columns") or 0)
    headers = list(data.get("headers") or [])
    rows = [list(row) for row in data.get("rows") or [] if isinstance(row, list)]
    original = {"rows": len(rows), "columns": len(headers)}
    if max_cols > 0 and len(headers) > max_cols:
        headers = headers[:max_cols]
        rows = [row[:max_cols] for row in rows]
    if max_rows > 0 and len(rows) > max_rows:
        rows = rows[:max_rows]
    data["headers"] = headers
    data["rows"] = rows
    adjusted = original["rows"] != len(rows) or original["columns"] != len(headers)
    if adjusted:
        warnings.append(
            _golden_renderer_warning(
                "TABLE_CAPACITY_ADJUSTED",
                "Table sample data was trimmed to fit the template contract.",
                {"before": original, "after": {"rows": len(rows), "columns": len(headers)}},
            )
        )
    return {"status": "adjusted" if adjusted else "ok", "rows": len(rows), "columns": len(headers), "limits": {"rows": max_rows, "columns": max_cols}}


def _shape_chart_to_capacity(fixture: dict[str, Any], capacity: dict[str, Any], warnings: list[dict[str, Any]]) -> dict[str, Any]:
    data = fixture.get("chart_data")
    if not isinstance(data, dict):
        return {"status": "not_applicable"}
    label_limit = int(capacity.get("chart_label_max_chars") or 0)
    if label_limit <= 0:
        return {"status": "ok", "label_limit": label_limit}
    adjusted = False
    categories = []
    for item in data.get("categories") or []:
        text = str(item)
        short = _compress_visible_text(text, label_limit)
        adjusted = adjusted or short != text
        categories.append(short)
    data["categories"] = categories
    for series in data.get("series") or []:
        if isinstance(series, dict) and isinstance(series.get("name"), str):
            short = _compress_visible_text(series["name"], label_limit)
            adjusted = adjusted or short != series["name"]
            series["name"] = short
    if adjusted:
        warnings.append(
            _golden_renderer_warning(
                "CHART_LABEL_CAPACITY_ADJUSTED",
                "Chart labels were shortened to fit the template contract.",
                {"label_limit": label_limit},
            )
        )
    return {"status": "adjusted" if adjusted else "ok", "label_limit": label_limit}


def _reading_order_contract_status(layout: dict[str, Any], contract: dict[str, Any]) -> dict[str, Any]:
    contract_order = [str(item) for item in contract.get("reading_order") or []]
    layout_order = [str(slot.get("slot_id")) for slot in layout.get("slots") or [] if isinstance(slot, dict)]
    missing = [slot_id for slot_id in contract_order if slot_id not in layout_order]
    ordered_indexes = [layout_order.index(slot_id) for slot_id in contract_order if slot_id in layout_order]
    monotonic = ordered_indexes == sorted(ordered_indexes)
    return {
        "status": "ok" if not missing and monotonic else "warning",
        "contract_order": contract_order,
        "layout_order": layout_order,
        "missing_slots": missing,
        "monotonic_in_layout": monotonic,
        "enforced_in_fixture_order": True,
    }


def _shape_budget_status(render_stats: dict[str, Any], contract: dict[str, Any]) -> dict[str, Any]:
    budget = contract.get("decorative_budget") if isinstance(contract.get("decorative_budget"), dict) else {}
    limit = int(budget.get("max_shape_count_target") or 0)
    count = int(render_stats.get("shape_count") or 0)
    return {"status": "ok" if not limit or count <= limit else "warning", "shape_count": count, "limit": limit}


def _decoration_budget_status(render_stats: dict[str, Any], contract: dict[str, Any]) -> dict[str, Any]:
    text_count = max(1, int(render_stats.get("text_shape_count") or 0))
    decorative = max(0, int(render_stats.get("shape_count") or 0) - text_count - int(render_stats.get("table_count") or 0) - int(render_stats.get("chart_count") or 0) - int(render_stats.get("picture_count") or 0))
    ratio = round(decorative / text_count, 3)
    density = str(((contract.get("decorative_budget") or {}).get("max_ornament_density") or "medium")).lower()
    limit = 12.0 if density == "high" else 9.0 if density == "medium" else 7.0
    return {"status": "ok" if ratio <= limit else "warning", "decorative_shape_count": decorative, "text_shape_count": text_count, "decoration_to_text_ratio": ratio, "limit": limit}


def _slide_render_stats(slide: Any) -> dict[str, int]:
    stats = {"shape_count": 0, "text_shape_count": 0, "table_count": 0, "chart_count": 0, "picture_count": 0}
    for shape in slide.shapes:
        stats["shape_count"] += 1
        if getattr(shape, "has_text_frame", False) and str(shape.text or "").strip():
            stats["text_shape_count"] += 1
        if getattr(shape, "has_table", False):
            stats["table_count"] += 1
        if getattr(shape, "has_chart", False):
            stats["chart_count"] += 1
        if int(getattr(shape, "shape_type", 0)) == 13:
            stats["picture_count"] += 1
    return stats


def _is_card_like_slot(slot_id: str, component_id: str) -> bool:
    normalized = slot_id.lower()
    component = component_id.lower()
    return normalized in {"cards", "caption_cards", "index_navigation", "progress_markers", "step_cards", "metric_panels"} or "card" in normalized or component in {"card", "cards", "layered_card", "glass_overlay_card", "index_navigation", "kpi_card", "premium_kpi_card"}


def _compress_visible_text(text: str, limit: int) -> str:
    value = str(text or "").strip()
    if not limit or len(value) <= limit:
        return value
    if limit <= 8:
        return value[:limit].rstrip()
    return value[: limit - 1].rstrip() + "."


def _short_capacity_label(text: str) -> str:
    words = [word for word in str(text or "").replace(":", " ").split() if word]
    if not words:
        return "Key"
    return words[0][:7]


def _golden_renderer_warning(code: str, message: str, details: dict[str, Any]) -> dict[str, Any]:
    return {"code": code, "severity": "warning", "message": message, "details": details}


def _fixture_visible_strings(payload: dict[str, Any]) -> list[str]:
    visible: list[str] = []
    for field in ("title", "subtitle", "footer_source"):
        value = payload.get(field)
        if isinstance(value, str):
            visible.append(value)
    for block in payload.get("content_blocks") or []:
        if not isinstance(block, dict):
            continue
        if isinstance(block.get("title"), str):
            visible.append(str(block["title"]))
        content = block.get("content")
        if isinstance(content, str):
            visible.append(content)
        elif isinstance(content, list):
            visible.extend(str(item) for item in content if isinstance(item, (str, int, float)))
    for card in payload.get("cards") or []:
        if isinstance(card, dict):
            visible.extend(str(card.get(field)) for field in ("title", "body") if isinstance(card.get(field), str))
    for citation in payload.get("citations") or []:
        if isinstance(citation, dict):
            visible.extend(str(citation.get(field)) for field in ("label", "source") if isinstance(citation.get(field), str))
    table_data = payload.get("table_data")
    if isinstance(table_data, dict):
        visible.extend(str(item) for item in table_data.get("headers") or [] if isinstance(item, (str, int, float)))
        for row in table_data.get("rows") or []:
            if isinstance(row, list):
                visible.extend(str(item) for item in row if isinstance(item, (str, int, float)))
    chart_data = payload.get("chart_data")
    if isinstance(chart_data, dict):
        visible.extend(str(item) for item in chart_data.get("categories") or [] if isinstance(item, (str, int, float)))
        for series in chart_data.get("series") or []:
            if isinstance(series, dict) and isinstance(series.get("name"), str):
                visible.append(str(series["name"]))
    return visible


def _contains_forbidden(text: str, phrase: str) -> bool:
    if phrase == "cards":
        return f" {phrase} " in f" {text.lower()} "
    return phrase.lower() in text.lower()


def _required_golden_master_spec_fields() -> set[str]:
    return {
        "archetype_id",
        "identity_priority",
        "required_visual_motifs",
        "forbidden_substitutions",
        "title_zone_geometry",
        "image_photo_zone_geometry",
        "footer_source_geometry",
        "ornament_system",
        "card_chart_table_chrome_style",
        "expected_density",
        "palette_mode",
        "identity_lock",
        "required_visual_signature",
        "forbidden_drift",
        "minimum_visual_features",
        "reference_source",
        "editable_translation",
    }


def _fallback_golden_master_spec(layout: dict[str, Any]) -> dict[str, Any]:
    archetype_id = str(layout.get("archetype_id") or "template")
    return {
        "archetype_id": archetype_id,
        "identity_priority": ["editable production master", archetype_id.replace("_", " ")],
        "identity_lock": "standard",
        "required_visual_signature": {
            "summary": f"{archetype_id.replace('_', ' ')} editable production master",
            "analysis_method": "fallback_golden_master_spec",
        },
        "required_visual_motifs": ["visible title hierarchy", "editable layout chrome", "footer rule"],
        "forbidden_substitutions": ["debug overlay", "full-slide raster background", "source-backed planning content"],
        "forbidden_drift": ["debug overlay", "full-slide raster background", "source-backed planning content"],
        "title_zone_geometry": _slot_geometry(layout, {"title", "headline"}),
        "image_photo_zone_geometry": _slot_geometry(layout, set(), slot_type="image"),
        "footer_source_geometry": _slot_geometry(layout, {"footer"}, slot_type="footer"),
        "reference_geometry": {},
        "ornament_system": {"density_mode": str(layout.get("density") or "medium"), "motifs": ["grid", "rule", "micro nodes"]},
        "card_chart_table_chrome_style": {"mode": "editable component chrome"},
        "component_grammar": {},
        "expected_density": str(layout.get("density") or "medium"),
        "palette_mode": _preview_tone(layout),
        "minimum_visual_features": {
            "dark_area_ratio": 0.05,
            "footer_presence": True,
            "index_presence": False,
            "ornament_presence": True,
            "image_mask_presence": False,
            "card_density": str(layout.get("density") or "medium"),
            "table_density": "none",
        },
        "reference_source": "board_crop",
        "editable_translation": {
            "ppt_text": ["title", "body", "footer"],
            "ppt_shape": ["layout chrome", "cards", "rules"],
            "ppt_table": ["editable table when declared"],
            "ppt_chart": ["editable chart when declared"],
            "svg_ornament": ["simple rule ornaments"],
            "photo_frame_image": ["declared frames only"],
        },
    }


def _golden_reference_source(layout_reference: dict[str, Any]) -> str:
    source = str(layout_reference.get("reference_source") or "")
    if source in {"layout_ref", "board_crop"}:
        return source
    preferred = str(layout_reference.get("preferred_source") or "")
    if preferred == "layout_ref_16_9":
        return "layout_ref"
    return "board_crop"


def _priority_ref_summary(compiled_layouts: list[dict[str, Any]]) -> dict[str, Any]:
    required = [item for item in compiled_layouts if item.get("priority_ref_required")]
    missing = [str(item.get("archetype_id")) for item in required if not item.get("priority_ref_available")]
    return {
        "status": "ready" if not missing else "NEEDS_REFERENCE_ASSET",
        "required_count": len(required),
        "available_count": len(required) - len(missing),
        "missing_count": len(missing),
        "missing_archetypes": missing,
        "normal_fallback_allowed": True,
    }


def _slot_geometry(layout: dict[str, Any], slot_ids: set[str], slot_type: str | None = None) -> dict[str, Any]:
    matches = []
    for slot in layout.get("slots") or []:
        if not isinstance(slot, dict):
            continue
        if slot.get("slot_id") in slot_ids or (slot_type and slot.get("slot_type") == slot_type):
            matches.append(
                {
                    "slot_id": slot.get("slot_id"),
                    "slot_type": slot.get("slot_type"),
                    "component_id": slot.get("component_id"),
                    "bounds": slot.get("bounds"),
                }
            )
    return {"slots": matches, "coordinate_system": "inches_16_9"}


def _golden_markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# Golden Template Masters Report",
        "",
        f"Status: `{report.get('status')}`",
        f"PPTX: `{report.get('pptx_path')}`",
        f"Source spec: `{report.get('source_spec_path')}`",
        f"Slide count: `{report.get('slide_count')}`",
        f"Rendering path: `{report.get('rendering_path')}`",
        "",
        "These master slides are production-looking editable layout references. Debug overlays, slot labels, safe margin labels, and internal layout identifiers are forbidden here and remain limited to `outputs/template_preview.pptx`.",
        "",
        "## Layouts",
        "",
    ]
    for item in report.get("compiled_layouts") or []:
        text_status = (item.get("text_capacity_status") or {}).get("status")
        decoration_status = (item.get("decoration_budget_status") or {}).get("status")
        shape_status = (item.get("shape_budget_status") or {}).get("status")
        warning_count = len(item.get("remaining_warnings") or item.get("warnings") or [])
        lines.append(
            f"- `{item.get('archetype_id')}`: slide `{item.get('slide_number')}`, role `{((item.get('text_capacity_status') or {}).get('presentation_role'))}`, text `{text_status}`, decoration `{decoration_status}`, shapes `{shape_status}`, warnings `{warning_count}`, fixture `{item.get('fixture_used')}`, contract `{item.get('contract_used')}`"
        )
    return "\n".join(lines) + "\n"


class _Style:
    def __init__(self, spec: dict[str, Any]) -> None:
        self.spec = spec
        self.canvas_width = float(spec["canvas"]["width"])
        self.canvas_height = float(spec["canvas"]["height"])
        tokens = spec["tokens"]
        self.colors = tokens["colors"]
        self.typography = tokens["typography"]
        self.line_weights = tokens.get("line_weights", {})
        self.spacing = tokens.get("spacing", {})

    def color(self, token: str, fallback: str = "#111827") -> RGBColor:
        return _hex_to_rgb(self.colors.get(token, fallback))

    def font_size(self, token: str, fallback: float) -> Pt:
        typography = self.typography.get(token, {})
        return Pt(float(typography.get("size_pt", fallback)))

    def font_family(self, token: str) -> str:
        typography = self.typography.get(token, {})
        return str(typography.get("font_family", "Aptos"))


def _draw_preview_slide(
    slide: Any,
    layout: dict[str, Any],
    style: _Style,
    slide_number: int,
    *,
    icon_resolver: IconResolver | None = None,
    icon_report: dict[str, Any] | None = None,
) -> None:
    """Render preview using the final-deck component path plus lightweight labels."""

    source = _preview_source(layout, slide_number)
    binding = _preview_binding(layout, slide_number)
    deck_style = _preview_deck_style(style.spec, layout, slide_number)
    warnings: list[dict[str, Any]] = []
    _render_shared_deck_slide(
        slide,
        source,
        layout,
        binding,
        deck_style,
        warnings,
        icon_resolver=icon_resolver,
        icon_report=icon_report,
    )
    _add_preview_overlay(slide, layout, style, slide_number, warnings)


def _preview_deck_style(spec: dict[str, Any], layout: dict[str, Any], slide_number: int) -> _SharedDeckStyle:
    binding = _preview_binding(layout, slide_number)
    deck_style = _SharedDeckStyle(spec).for_tone(str(binding["selected_tone_variant"]))
    return deck_style.for_layout(layout, binding)


def _preview_source(layout: dict[str, Any], slide_number: int) -> dict[str, Any]:
    archetype_id = str(layout.get("archetype_id") or "template")
    title = _preview_title(archetype_id)
    slots = [slot for slot in layout.get("slots") or [] if isinstance(slot, dict)]
    content_blocks = [
        {
            "block_id": f"preview-{slot.get('slot_id', index)}",
            "slot": str(slot.get("slot_id") or f"slot_{index}"),
            "title": _slot_label(slot),
            "content": _slot_content(slot),
        }
        for index, slot in enumerate(slots, start=1)
        if slot.get("slot_type") in {"content", "text"}
    ]
    if not content_blocks:
        content_blocks = [
            {"block_id": "preview-message", "slot": "body", "title": "CONTENT", "content": "Editable template component"},
            {"block_id": "preview-evidence", "slot": "cards", "title": "EVIDENCE", "content": "Source-backed module"},
        ]
    return {
        "slide_id": f"template-preview-{slide_number:03d}",
        "section_id": f"S{max(1, slide_number):02d}",
        "slide_type": archetype_id,
        "title": title,
        "subtitle": "SUBTITLE AREA" if archetype_id == "creative_cover" else "Editable layout reference",
        "content_density": str(layout.get("density") or "medium"),
        "content_blocks": content_blocks,
        "chart_data": {
            "categories": ["A", "B", "C", "D"],
            "series": [
                {"name": "KPI", "values": [78, 64, 86, 72]},
                {"name": "Delta", "values": [42, 51, 47, 58]},
            ],
        },
        "table_data": {
            "headers": ["Signal", "Evidence", "Implication"],
            "rows": [
                ["Trace", "Source-backed", "Prioritize"],
                ["Method", "Repeatable", "Standardize"],
                ["Risk", "Bounded", "Monitor"],
            ],
        },
        "image_needs": [{"slot": "image", "usage": "declared_photo_frame_only"}],
        "speaker_notes": "Template preview generated from editable component primitives.",
        "citations": [{"label": "FOOTER", "source": "SOURCE"}],
        "design_intent": {
            "preview": True,
            "layout_family_id": layout.get("layout_family_id"),
        },
    }


def _preview_binding(layout: dict[str, Any], slide_number: int) -> dict[str, Any]:
    archetype_id = str(layout.get("archetype_id") or "template")
    slot_bindings = {
        str(slot.get("slot_id")): _slot_binding(slot)
        for slot in layout.get("slots") or []
        if isinstance(slot, dict) and slot.get("slot_id")
    }
    return {
        "slide_id": f"template-preview-{slide_number:03d}",
        "slide_type": archetype_id,
        "selected_layout_id": layout["layout_id"],
        "layout_family_id": layout.get("layout_family_id"),
        "selected_tone_variant": _preview_tone(layout),
        "deck_scale": "template_preview",
        "section_rhythm_role": "divider" if archetype_id == "section_divider" else "overview",
        "slot_bindings": slot_bindings,
        "component_bindings": {
            str(slot.get("slot_id")): slot.get("component_id")
            for slot in layout.get("slots") or []
            if isinstance(slot, dict) and slot.get("slot_id")
        },
        "fallback_used": False,
        "fallback_reason": None,
        "preview_rendering_path": "deck_compiler_component_primitives",
    }


def _slot_binding(slot: dict[str, Any]) -> str:
    slot_id = str(slot.get("slot_id") or "")
    slot_type = str(slot.get("slot_type") or "")
    if slot_id in {"title", "headline"}:
        return "title"
    if slot_id == "subtitle":
        return "subtitle"
    if slot_id == "footer" or slot_type == "footer":
        return "citations"
    if slot_type == "table":
        return "table_data"
    if slot_type == "chart":
        return "chart_data"
    if slot_type == "image":
        return "image_needs"
    return "content_blocks"


def _preview_tone(layout: dict[str, Any]) -> str:
    archetype_id = str(layout.get("archetype_id") or "")
    family_id = str(layout.get("layout_family_id") or "")
    if archetype_id in {"creative_cover", "section_divider"} or family_id == "expressive_cover_divider":
        return "creative"
    if family_id in {"evidence_overview", "table_appendix", "methodology_framework"}:
        return "academic"
    return "professional"


def _preview_title(archetype_id: str) -> str:
    if archetype_id == "creative_cover":
        return "TITLE"
    if archetype_id == "section_divider":
        return "SECTION"
    return archetype_id.replace("_", " ").title()


def _slot_label(slot: dict[str, Any]) -> str:
    slot_id = str(slot.get("slot_id") or "slot").replace("_", " ").upper()
    component = str(slot.get("component_id") or "").replace("_", " ").upper()
    return f"{slot_id} {component}".strip()


def _slot_content(slot: dict[str, Any]) -> str:
    component = str(slot.get("component_id") or "")
    slot_id = str(slot.get("slot_id") or "")
    if slot_id == "cards" or component in {"card", "layered_card"}:
        return "Card title\nEvidence note\nSource tag"
    if component in {"radial_process", "curved_timeline", "connector_style"} or slot_id in {"diagram", "process", "timeline"}:
        return "Frame\nReview\nDecide\nLearn"
    if slot_id == "metric_panels" or component in {"kpi_card", "premium_kpi_card"}:
        return "KPI\nTrend\nDelta"
    return "Editable placeholder block"


def _add_preview_overlay(
    slide: Any,
    layout: dict[str, Any],
    style: _Style,
    slide_number: int,
    warnings: list[dict[str, Any]],
) -> None:
    _draw_safe_margins(slide, style)
    _add_header(slide, layout, style, slide_number)
    for slot in layout["slots"]:
        _draw_slot_boundary(slide, slot, style)
    _add_footer_system_label(slide, layout, style)
    if warnings:
        _text_box(
            slide,
            f"PREVIEW WARNINGS: {len(warnings)}",
            10.9,
            0.34,
            1.65,
            0.18,
            style,
            "label",
            "accent",
            align=PP_ALIGN.RIGHT,
        )


def _draw_slot_boundary(slide: Any, slot: dict[str, Any], style: _Style) -> None:
    bounds = slot["bounds"]
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.background()
    rect.line.color.rgb = style.color("accent_secondary", "#0F766E")
    rect.line.width = Pt(0.45)
    rect.line.dash_style = 4
    label = f"{slot['slot_id']} | {slot.get('slot_type', '')} | {slot.get('component_id', '')}"
    _text_box(slide, label, x + 0.06, y + 0.04, max(0.4, min(w - 0.12, 2.7)), 0.14, style, "label", "accent_secondary")
    if str(slot.get("slot_type") or "") == "image":
        _text_box(slide, "IMAGE FRAME ONLY", x + 0.08, y + h - 0.22, max(0.5, min(w - 0.16, 2.2)), 0.14, style, "footer", "accent_secondary")


def _add_footer_system_label(slide: Any, layout: dict[str, Any], style: _Style) -> None:
    _text_box(
        slide,
        f"FOOTER SYSTEM | editable source note, slide index, citation strip | {layout['layout_id']}",
        0.72,
        7.03,
        10.6,
        0.16,
        style,
        "footer",
        "muted_text",
    )


def _draw_background(slide: Any, style: _Style) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(style.canvas_width),
        Inches(style.canvas_height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style.color("background", "#F8FAFC")
    shape.line.fill.background()

    grid_color = style.color("grid", "#E2E8F0")
    step = 0.5
    x = step
    while x < style.canvas_width:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(0.0), Inches(x), Inches(style.canvas_height))
        line.line.color.rgb = grid_color
        line.line.width = Pt(0.25)
        x += step
    y = step
    while y < style.canvas_height:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.0), Inches(y), Inches(style.canvas_width), Inches(y))
        line.line.color.rgb = grid_color
        line.line.width = Pt(0.25)
        y += step
    for x, y in ((0.72, 0.72), (12.45, 0.72), (0.72, 6.55), (12.45, 6.55)):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.05), Inches(0.05))
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.color("accent_secondary", "#0F766E")
        dot.line.fill.background()


def _draw_safe_margins(slide: Any, style: _Style) -> None:
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(SAFE_MARGIN_IN),
        Inches(SAFE_MARGIN_IN),
        Inches(style.canvas_width - SAFE_MARGIN_IN * 2),
        Inches(style.canvas_height - SAFE_MARGIN_IN * 2),
    )
    rect.fill.background()
    rect.line.color.rgb = style.color("accent_secondary", "#0F766E")
    rect.line.width = Pt(1.1)
    rect.line.dash_style = 4
    _text_box(slide, "SAFE MARGINS", SAFE_MARGIN_IN + 0.05, SAFE_MARGIN_IN + 0.04, 1.25, 0.18, style, "label", "accent_secondary")


def _add_header(slide: Any, layout: dict[str, Any], style: _Style, slide_number: int) -> None:
    _text_box(
        slide,
        f"{slide_number:02d}  {layout['layout_id']}",
        0.55,
        0.08,
        7.1,
        0.25,
        style,
        "label",
        "muted_text",
    )
    _text_box(
        slide,
        f"ARCHETYPE: {layout['archetype_id']}    DENSITY: {layout.get('density', 'unspecified')}",
        7.85,
        0.08,
        4.9,
        0.25,
        style,
        "label",
        "muted_text",
        align=PP_ALIGN.RIGHT,
    )


def _add_footer_system(slide: Any, layout: dict[str, Any], style: _Style) -> None:
    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(6.78),
        Inches(12.22),
        Inches(0.45),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = style.color("surface", "#FFFFFF")
    footer.line.color.rgb = style.color("line", "#CBD5E1")
    footer.line.width = Pt(0.8)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(6.78), Inches(0.12), Inches(0.45))
    band.fill.solid()
    band.fill.fore_color.rgb = style.color("accent", "#2563EB")
    band.line.fill.background()
    rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.72), Inches(6.84), Inches(12.55), Inches(6.84))
    rule.line.color.rgb = style.color("grid", "#E2E8F0")
    rule.line.width = Pt(0.5)
    _text_box(
        slide,
        f"FOOTER SYSTEM | editable slide number, source note, section marker | {layout['layout_id']}",
        0.72,
        6.88,
        10.7,
        0.18,
        style,
        "footer",
        "muted_text",
    )


def _draw_slot(slide: Any, slot: dict[str, Any], style: _Style) -> None:
    bounds = slot["bounds"]
    slot_type = slot["slot_type"]
    component_id = slot.get("component_id", "")
    label = f"{slot['slot_id']} | {slot_type} | {component_id}"
    if slot_type == "table":
        _draw_table_slot(slide, bounds, label, style)
    elif slot_type == "chart":
        _draw_chart_slot(slide, bounds, label, style)
    elif slot_type == "image":
        _draw_image_slot(slide, bounds, label, style, component_id)
    elif component_id == "section_marker":
        _draw_section_marker_slot(slide, bounds, label, style)
    elif component_id == "kpi_card":
        _draw_kpi_slot(slide, bounds, label, style)
    elif component_id == "card":
        _draw_card_slot(slide, bounds, label, style)
    else:
        _draw_shape_slot(slide, bounds, label, style, slot_type)


def _draw_card_slot(slide: Any, bounds: dict[str, Any], label: str, style: _Style) -> None:
    _draw_shape_slot(slide, bounds, label, style, "content")
    x, y, w, _h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(float(bounds["h"])))
    rail.fill.solid()
    rail.fill.fore_color.rgb = style.color("accent_secondary", "#0F766E")
    rail.line.fill.background()
    top_rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.18), Inches(y + 0.18), Inches(x + w - 0.18), Inches(y + 0.18))
    top_rule.line.color.rgb = style.color("grid", "#E2E8F0")
    top_rule.line.width = Pt(0.6)


def _draw_kpi_slot(slide: Any, bounds: dict[str, Any], label: str, style: _Style) -> None:
    _draw_shape_slot(slide, bounds, label, style, "content")
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    for index in range(3):
        card_y = y + 0.25 + index * max(0.38, (h - 0.5) / 3)
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.24), Inches(card_y), Inches(0.12), Inches(0.12))
        marker.fill.solid()
        marker.fill.fore_color.rgb = style.color("accent", "#2563EB")
        marker.line.fill.background()
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.48), Inches(card_y + 0.06), Inches(x + w - 0.24), Inches(card_y + 0.06))
        rule.line.color.rgb = style.color("line", "#CBD5E1")
        rule.line.width = Pt(0.5)


def _draw_shape_slot(slide: Any, bounds: dict[str, Any], label: str, style: _Style, slot_type: str) -> None:
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if slot_type in {"content", "text"} else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(float(bounds["x"])),
        Inches(float(bounds["y"])),
        Inches(float(bounds["w"])),
        Inches(float(bounds["h"])),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = style.color("surface", "#FFFFFF")
    shape.line.color.rgb = style.color("accent", "#2563EB")
    shape.line.width = Pt(1.4)
    shape.text_frame.clear()
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = label
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.name = style.font_family("label")
    run.font.size = style.font_size("label", 9)
    run.font.color.rgb = style.color("text", "#111827")
    shape.text_frame.margin_left = Inches(0.08)
    shape.text_frame.margin_right = Inches(0.08)
    shape.text_frame.margin_top = Inches(0.05)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def _draw_table_slot(slide: Any, bounds: dict[str, Any], label: str, style: _Style) -> None:
    rows, cols = 4, 4
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(float(bounds["x"])),
        Inches(float(bounds["y"])),
        Inches(float(bounds["w"])),
        Inches(float(bounds["h"])),
    )
    table = table_shape.table
    for row_index in range(rows):
        for col_index in range(cols):
            cell = table.cell(row_index, col_index)
            cell.text = "TABLE" if row_index == 0 and col_index == 0 else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = style.color("surface", "#FFFFFF") if row_index else style.color("surface_alt", "#EEF2F7")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = style.font_family("label")
                    run.font.size = style.font_size("label", 9)
                    run.font.color.rgb = style.color("text", "#111827")
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    for index in range(1, 4):
        line_y = y + h * index / 4
        rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(line_y), Inches(x + w), Inches(line_y))
        rule.line.color.rgb = style.color("grid", "#E2E8F0")
        rule.line.width = Pt(0.35)
    _slot_label_overlay(slide, bounds, label, style)


def _draw_chart_slot(slide: Any, bounds: dict[str, Any], label: str, style: _Style) -> None:
    _draw_shape_slot(slide, bounds, label, style, "shape")
    x = float(bounds["x"])
    y = float(bounds["y"])
    w = float(bounds["w"])
    h = float(bounds["h"])
    for index, height_ratio in enumerate((0.35, 0.62, 0.48)):
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x + 0.35 + index * 0.5),
            Inches(y + h - 0.35 - h * height_ratio * 0.55),
            Inches(0.25),
            Inches(h * height_ratio * 0.55),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = style.color("accent", "#2563EB")
        bar.line.fill.background()
    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.28), Inches(y + h - 0.32), Inches(x + w - 0.28), Inches(y + h - 0.32))
    axis.line.color.rgb = style.color("muted_text", "#475569")
    axis.line.width = Pt(0.7)


def _draw_image_slot(slide: Any, bounds: dict[str, Any], label: str, style: _Style, component_id: str) -> None:
    _draw_shape_slot(slide, bounds, label, style, "shape")
    x = float(bounds["x"])
    y = float(bounds["y"])
    w = float(bounds["w"])
    h = float(bounds["h"])
    line_color = style.color("accent_secondary", "#0F766E")
    first = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y + h))
    second = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + w), Inches(y), Inches(x), Inches(y + h))
    for line in (first, second):
        line.line.color.rgb = line_color
        line.line.width = Pt(1.0)
    if component_id == "diagonal_photo_panel":
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(x + w * 0.55), Inches(y), Inches(w * 0.32), Inches(h))
        band.fill.solid()
        band.fill.fore_color.rgb = style.color("surface_alt", "#EEF2F7")
        band.line.color.rgb = style.color("accent_secondary", "#0F766E")
        band.line.width = Pt(0.8)
    _text_box(slide, "IMAGE FRAME ONLY", x + 0.14, y + h - 0.34, min(w - 0.28, 2.0), 0.18, style, "label", "accent_secondary")


def _draw_section_marker_slot(slide: Any, bounds: dict[str, Any], label: str, style: _Style) -> None:
    x, y, w, h = (float(bounds[key]) for key in ("x", "y", "w", "h"))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    band.fill.solid()
    band.fill.fore_color.rgb = style.color("text", "#111827")
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(min(0.18, w)), Inches(h))
    accent.fill.solid()
    accent.fill.fore_color.rgb = style.color("accent", "#2563EB")
    accent.line.fill.background()
    _text_box(slide, label, x + 0.24, y + 0.04, max(0.4, w - 0.3), max(0.12, h - 0.08), style, "label", "surface")


def _slot_label_overlay(slide: Any, bounds: dict[str, Any], label: str, style: _Style) -> None:
    _text_box(
        slide,
        label,
        float(bounds["x"]) + 0.08,
        float(bounds["y"]) + 0.06,
        max(0.5, float(bounds["w"]) - 0.16),
        0.18,
        style,
        "label",
        "text",
    )


def _text_box(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    style: _Style,
    typography_token: str,
    color_token: str,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = style.font_family(typography_token)
    run.font.size = style.font_size(typography_token, 9)
    run.font.color.rgb = style.color(color_token, "#111827")
    return box


def _hex_to_rgb(value: str) -> RGBColor:
    text = value.strip().lstrip("#")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _display_path(path: Path) -> str:
    return str(path.as_posix())


if __name__ == "__main__":
    raise SystemExit(main())
